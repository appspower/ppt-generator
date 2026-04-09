"""컴포넌트 렌더러 - 회사 컬러/스타일 정책 반영.

Badge: pill형태 (roundRect)
Card: 3단계 스타일 (default/dark/accent)
SubMarker: 소제목 옆 오렌지 수직 바
AccentBox: 풀너비 강조 배너
"""

from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from .styles import (
    CL_WHITE, CL_BLACK, CL_BG, CL_BORDER,
    CL_BODY_TEXT, CL_DARK, CL_GREY, CL_GREY_MID, CL_GREY_LIGHT,
    CL_ACCENT, CL_ACCENT_MID, CL_ACCENT_LIGHT,
    CL_TABLE_HEADER, CL_TABLE_ROW_ALT,
    FONT_BODY, FONT_TITLE,
    FONT_SIZE_BODY, FONT_SIZE_SMALL, FONT_SIZE_BULLET,
    FONT_SIZE_HEADER, FONT_SIZE_BADGE, FONT_SIZE_KICKER,
    FONT_SIZE_TABLE_HEADER, FONT_SIZE_TABLE_BODY, FONT_SIZE_KPI,
    COL_GAP,
)
from ..models.schema import (
    CardComponent, TextBlockComponent, BadgeComponent,
    KickerComponent, BulletComponent, TableComponent,
    ChartComponent, ProcessFlowComponent, DividerComponent,
    ImageComponent, ProcessStep,
)


# ============================================================
# Card (3 styles: default, dark, accent)
# ============================================================

def render_card(slide: Slide, comp: CardComponent, x, y, w, h):
    from .styles import estimate_text_height

    style = comp.style
    # 방안 2: accent는 상단 바만 오렌지, 본문은 흰색 (색상 절제)
    if style == "dark":
        top_bar_c, fill_c, text_c, header_c = CL_DARK, CL_DARK, CL_WHITE, CL_WHITE
        border_c = None
    elif style == "accent":
        top_bar_c, fill_c, text_c, header_c = CL_ACCENT, CL_WHITE, CL_BODY_TEXT, CL_ACCENT
        border_c = CL_BORDER
    else:
        top_bar_c, fill_c, text_c, header_c = CL_GREY, CL_WHITE, CL_BODY_TEXT, CL_BLACK
        border_c = CL_BORDER

    pad_inch = 0.08  # ④ 패딩 축소 0.15→0.08
    inner_w = w / 914400 - 2 * pad_inch

    # 텍스트 실제 높이 추정
    est_h = pad_inch * 2  # 상하 패딩
    if comp.header:
        est_h += estimate_text_height(comp.header, 11, inner_w, bold=True) + 0.05
    if comp.subtitle:
        est_h += estimate_text_height(comp.subtitle, 9, inner_w) + 0.08
    for block in comp.content:
        if block.type == "kpi":
            est_h += 0.45  # KPI 값 + 설명
        else:
            est_h += estimate_text_height(block.text, 10, inner_w) + 0.05
    if comp.content and comp.bullets:
        est_h += 0.1  # 구분 간격
    for bullet in comp.bullets:
        est_h += estimate_text_height(bullet, 9, inner_w - 0.15) + 0.03

    # 카드 높이: 추정치보다 작지 않게, 단 할당 h는 절대 넘지 않음
    # (Top 정렬이라 텍스트 적으면 하단 여백 발생 → 추정 높이로 카드를 줄임)
    est_card_h = Inches(est_h + 0.2)
    card_h = max(Inches(1.2), min(h, est_card_h))

    # ① 카드 상단 악센트 라인 (2px — 디자이너 스타일)
    top_bar_h = Inches(0.02)
    top_bar = slide.shapes.add_shape(1, x, y, w, top_bar_h)
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = top_bar_c
    top_bar.line.fill.background()

    # ③ 카드 본문 — roundRect (디자이너 스타일)
    body_y = y + top_bar_h
    body_h = card_h - top_bar_h
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, body_y, w, body_h)
    try:
        box.adjustments[0] = 0.03  # 미묘한 둥근 모서리
    except (IndexError, TypeError):
        pass
    box.fill.solid()
    box.fill.fore_color.rgb = fill_c
    if border_c:
        box.line.color.rgb = border_c
        box.line.width = Emu(6350)
    else:
        box.line.fill.background()

    # 과제3: 드롭 쉐도우 추가
    from .renderers.base import add_shadow
    add_shadow(box, blur=3, dist=2, color="C0C0C0", alpha=30000)

    pad = Inches(pad_inch)
    tx = slide.shapes.add_textbox(x + pad, body_y + Inches(0.08), w - 2 * pad, body_h - Inches(0.16))
    tf = tx.text_frame
    tf.word_wrap = True
    # 텍스트가 적으면 카드 중앙 정렬 (하단 빈공간 방지)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True

    # 과제3: 헤더 + 아이콘 앵커
    if comp.header:
        # 아이콘 심볼 (스타일별)
        icon_map = {"accent": "●", "dark": "◆", "default": "▸"}
        icon = icon_map.get(style, "▸")
        p = tf.paragraphs[0]
        p.text = f"{icon}  {comp.header}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = header_c
        p.font.name = FONT_BODY
        p.space_after = Pt(2)
        first = False

    # 서브타이틀
    if comp.subtitle:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = comp.subtitle
        p.font.size = Pt(9)
        p.font.color.rgb = CL_GREY if style == "default" else CL_ACCENT_LIGHT
        p.font.name = FONT_BODY
        p.space_after = Pt(6)
        first = False

    # 콘텐츠 블록
    for block in comp.content:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if block.type == "kpi":
            p.text = block.value
            p.font.size = FONT_SIZE_KPI
            p.font.bold = True
            p.font.color.rgb = CL_ACCENT if style == "default" else CL_WHITE
            p.font.name = FONT_BODY
            p.space_after = Pt(1)
            if block.text:
                p2 = tf.add_paragraph()
                p2.text = block.text
                p2.font.size = Pt(8)
                p2.font.color.rgb = CL_GREY if style == "default" else CL_ACCENT_LIGHT
                p2.font.name = FONT_BODY
                p2.space_after = Pt(8)
        else:
            p.text = block.text
            p.font.size = FONT_SIZE_BODY
            p.font.bold = block.bold
            p.font.color.rgb = text_c
            p.font.name = FONT_BODY
            p.space_after = Pt(3)
        first = False

    # KPI와 불릿 사이 구분선
    if comp.content and comp.bullets:
        sep = tf.add_paragraph()
        sep.text = "─" * 15
        sep.font.size = Pt(5)
        sep.font.color.rgb = CL_BORDER
        sep.space_after = Pt(3)

    # 불릿 (빈 문자열은 간격으로 처리)
    for bullet in comp.bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if bullet.strip():
            p.text = f"\u2022  {bullet}"
            p.font.size = Pt(9)
            p.font.color.rgb = text_c
            p.font.name = FONT_BODY
            p.space_after = Pt(3)
        else:
            # 빈 줄은 작은 간격만 — 빈 불릿 표시 방지
            p.text = ""
            p.font.size = Pt(4)
            p.space_after = Pt(2)
        first = False


# ============================================================
# Text Block
# ============================================================

def render_text_block(slide: Slide, comp: TextBlockComponent, x, y, w, h):
    size_map = {"small": FONT_SIZE_SMALL, "body": FONT_SIZE_BODY, "large": FONT_SIZE_HEADER}
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = comp.text
    p.font.size = size_map.get(comp.size, FONT_SIZE_BODY)
    p.font.bold = comp.bold
    p.font.color.rgb = CL_BODY_TEXT
    p.font.name = FONT_BODY


# ============================================================
# Badge (Pill 형태 - roundRect)
# ============================================================

def render_badge(slide: Slide, comp: BadgeComponent, x, y, w, h):
    badge_h = min(h, Inches(0.28))
    badge_w = min(w, Inches(0.13) * len(comp.label) + Inches(0.4))

    if comp.style == "accent":
        fill, text_c = CL_ACCENT, CL_WHITE
    elif comp.style == "dark":
        fill, text_c = CL_DARK, CL_WHITE
    else:
        fill, text_c = CL_BG, CL_BODY_TEXT

    # roundRect (pill shape)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, badge_w, badge_h,
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    # adj=30000 for pill radius (components.md 기준)
    try:
        box.adjustments[0] = 0.5
    except (IndexError, TypeError):
        pass

    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = comp.label
    p.font.size = FONT_SIZE_BADGE
    p.font.bold = True
    p.font.color.rgb = text_c
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


# ============================================================
# Kicker (상단 작은 강조 텍스트)
# ============================================================

def render_kicker(slide: Slide, comp: KickerComponent, x, y, w, h):
    # 오렌지 수직 바 (SubMarker 스타일)
    bar_w = Inches(0.04)
    bar = slide.shapes.add_shape(1, x, y, bar_w, min(h, Inches(0.22)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CL_ACCENT
    bar.line.fill.background()

    # 텍스트
    tx = slide.shapes.add_textbox(x + Inches(0.1), y, w - Inches(0.1), min(h, Inches(0.25)))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = comp.text.upper()
    p.font.size = FONT_SIZE_KICKER
    p.font.bold = True
    p.font.color.rgb = CL_ACCENT
    p.font.name = FONT_BODY


# ============================================================
# Bullet
# ============================================================

def render_bullet(slide: Slide, comp: BulletComponent, x, y, w, h):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(comp.items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if comp.style == "kpi":
            p.text = f"\u25A0  {item}"
        else:
            p.text = f"\u2022  {item}"
        p.font.size = FONT_SIZE_BODY
        p.font.color.rgb = CL_BODY_TEXT
        p.font.name = FONT_BODY
        p.space_after = Pt(4)


# ============================================================
# Table (검정 헤더 + 교대행)
# ============================================================

def render_table(slide: Slide, comp: TableComponent, x, y, w, h):
    data = comp.data
    n_rows = len(data.rows) + 1
    n_cols = len(data.headers)
    # 사용 가능한 높이를 행 수로 균등 분배 (오버플로 방지)
    row_h = int(h / n_rows)
    table_h = row_h * n_rows
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, table_h).table

    # 행 높이 명시적 설정 — 자동 확장 방지
    for r in range(n_rows):
        tbl.rows[r].height = row_h

    col_w = int(w / n_cols)
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    # 헤더 (검정 배경 + 흰 텍스트)
    for j, hdr in enumerate(data.headers):
        cell = tbl.cell(0, j)
        cell.text = hdr
        _style_cell(cell, is_header=True)

    # 데이터 (교대행)
    for i, row in enumerate(data.rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            _style_cell(cell, is_header=False, alt_row=(i % 2 == 1))


def _style_cell(cell, is_header: bool, alt_row: bool = False):
    if is_header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = CL_TABLE_HEADER
    elif alt_row:
        cell.fill.solid()
        cell.fill.fore_color.rgb = CL_TABLE_ROW_ALT
    else:
        cell.fill.background()

    for p in cell.text_frame.paragraphs:
        p.font.name = FONT_BODY
        if is_header:
            p.font.size = FONT_SIZE_TABLE_HEADER
            p.font.bold = True
            p.font.color.rgb = CL_WHITE
        else:
            p.font.size = FONT_SIZE_TABLE_BODY
            p.font.color.rgb = CL_BODY_TEXT

    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)


# ============================================================
# Process Flow (오렌지 그라데이션 단계별)
# ============================================================

def render_process_flow(slide: Slide, comp: ProcessFlowComponent, x, y, w, h):
    n = len(comp.steps)
    gap_w = Inches(0.08)
    step_w = int((w - gap_w * (n - 1)) / n)
    step_h = min(h, Inches(1.8))
    step_y = y + int((h - step_h) / 2)

    # 단계별 색상 그라데이션 (Orange → Grey)
    step_colors = _gradient_colors(n)

    for i, step in enumerate(comp.steps):
        sx = x + i * (step_w + int(gap_w))

        # 번호 영역 (상단)
        num_h = Inches(0.5)
        num_box = slide.shapes.add_shape(1, sx, step_y, step_w, num_h)
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = step_colors[i]
        num_box.line.fill.background()

        tf = num_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{i + 1:02d}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = CL_WHITE if i < n // 2 + 1 else CL_BLACK
        p.font.name = FONT_TITLE
        p.alignment = PP_ALIGN.LEFT

        # 내용 영역 (하단)
        body_y = step_y + num_h
        body_h = step_h - num_h
        tx = slide.shapes.add_textbox(sx + Inches(0.08), body_y + Inches(0.08),
                                       step_w - Inches(0.16), body_h - Inches(0.16))
        tf = tx.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = step.label
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = CL_BLACK
        p.font.name = FONT_BODY
        p.space_after = Pt(4)

        if step.description:
            p2 = tf.add_paragraph()
            p2.text = step.description
            p2.font.size = FONT_SIZE_SMALL
            p2.font.color.rgb = CL_GREY
            p2.font.name = FONT_BODY


def _gradient_colors(n: int) -> list[RGBColor]:
    """단계별 Orange → Grey 그라데이션 색상을 생성."""
    if n <= 1:
        return [CL_ACCENT]
    colors = [CL_ACCENT, CL_ACCENT_MID, CL_ACCENT_LIGHT, CL_GREY_LIGHT, CL_GREY_MID, CL_GREY]
    if n <= len(colors):
        return colors[:n]
    # 반복
    return [colors[i % len(colors)] for i in range(n)]


# ============================================================
# Chevron Process (쉐브론 화살표 프로세스)
# ============================================================

def render_chevron_process(slide: Slide, comp, x, y, w, h):
    """쉐브론 프로세스 — 라벨+설명을 쉐브론 안에 통합 (1행 구조)."""
    from pptx.enum.shapes import MSO_SHAPE

    n = len(comp.steps)
    overlap = Inches(0.05)
    step_w = int((w + overlap * (n - 1)) / n)
    step_h = min(h, Inches(0.65))
    step_y = y + int((h - step_h) / 2)  # 할당 영역 내 수직 중앙

    colors = _gradient_colors(n)

    for i, step in enumerate(comp.steps):
        sx = x + i * (step_w - int(overlap))
        text_c = CL_WHITE if i < (n + 1) // 2 else CL_BLACK

        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, sx, step_y, step_w, step_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 라벨 (굵게)
        p = tf.paragraphs[0]
        p.text = step.label
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = text_c
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)

        # 설명 (쉐브론 안에 2번째 줄)
        if step.description:
            p2 = tf.add_paragraph()
            p2.text = step.description
            p2.font.size = Pt(7)
            p2.font.bold = False
            p2.font.color.rgb = text_c
            p2.font.name = FONT_BODY
            p2.alignment = PP_ALIGN.CENTER


# ============================================================
# Framework Matrix (N×M 프레임워크 매트릭스)
# ============================================================

def render_framework_matrix(slide: Slide, comp, x, y, w, h):
    """N×M 프레임워크 매트릭스 — 행/열 헤더 + 셀 하이라이트."""
    from ..models.schema import FrameworkMatrixComponent

    n_rows = len(comp.cells)
    n_cols = len(comp.cells[0]) if n_rows > 0 else 0
    has_row_h = len(comp.row_headers) > 0
    has_col_h = len(comp.col_headers) > 0

    # 전체 행/열 수 (헤더 포함)
    total_rows = n_rows + (1 if has_col_h else 0)
    total_cols = n_cols + (1 if has_row_h else 0)

    # 행별 가변 높이: 전체 높이를 행 수로 균등 분배 (공간 채우기)
    row_h_val = int(h / total_rows)
    col_w_val = int(w / total_cols)

    # 열 헤더 행은 더 작게
    header_row_h = min(row_h_val, Inches(0.4))
    data_row_h = int((h - (header_row_h if has_col_h else 0)) / n_rows) if n_rows > 0 else row_h_val

    for ri in range(total_rows):
        is_header_row = has_col_h and ri == 0
        current_row_h = header_row_h if is_header_row else data_row_h
        for ci in range(total_cols):
            cx = x + ci * col_w_val
            # 행별 y 계산
            if is_header_row:
                cy = y
            else:
                data_ri = ri - (1 if has_col_h else 0)
                cy = y + (header_row_h if has_col_h else 0) + data_ri * data_row_h

            is_col_header = has_col_h and ri == 0
            is_row_header = has_row_h and ci == 0

            # 셀 텍스트 결정
            if is_col_header and is_row_header:
                text = ""
                fill = CL_WHITE
                text_c = CL_BLACK
                bold = False
            elif is_col_header:
                col_idx = ci - (1 if has_row_h else 0)
                text = comp.col_headers[col_idx] if col_idx < len(comp.col_headers) else ""
                fill = CL_DARK  # ② 오렌지→다크 (색상 절제)
                text_c = CL_WHITE
                bold = True
            elif is_row_header:
                row_idx = ri - (1 if has_col_h else 0)
                text = comp.row_headers[row_idx] if row_idx < len(comp.row_headers) else ""
                fill = CL_DARK
                text_c = CL_WHITE
                bold = True
            else:
                row_idx = ri - (1 if has_col_h else 0)
                col_idx = ci - (1 if has_row_h else 0)
                cell = comp.cells[row_idx][col_idx]
                text = cell.text
                if cell.highlight:
                    # 방안 2: highlight는 연한 배경 + 오렌지 텍스트 (절제)
                    fill = RGBColor(0xFF, 0xF3, 0xEB)  # 매우 연한 피치
                    text_c = CL_ACCENT
                elif cell.style == "accent":
                    # 방안 2: accent도 배경은 흰색, 텍스트만 오렌지 Bold
                    fill = CL_WHITE
                    text_c = CL_ACCENT
                    bold = True
                elif cell.style == "dark":
                    fill = CL_DARK
                    text_c = CL_WHITE
                else:
                    fill = CL_WHITE if (row_idx % 2 == 0) else CL_TABLE_ROW_ALT
                    text_c = CL_BODY_TEXT
                bold = False

            # 셀 렌더링
            box = slide.shapes.add_shape(1, cx, cy, col_w_val, current_row_h)
            box.fill.solid()
            box.fill.fore_color.rgb = fill
            box.line.color.rgb = CL_BORDER
            box.line.width = Emu(6350)

            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.06)
            tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.03)
            tf.margin_bottom = Inches(0.03)

            # 줄바꿈(\n) 처리: 여러 paragraph로 분리
            lines = text.split('\n') if text else [""]
            for li, line in enumerate(lines):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(8)
                p.font.bold = bold
                p.font.color.rgb = text_c
                p.font.name = FONT_BODY
                p.space_after = Pt(1)


# ============================================================
# Numbered Circle (번호 원형 뱃지)
# ============================================================

def render_numbered_circle(slide: Slide, comp, x, y, w, h):
    """번호 원형 뱃지 + 텍스트 리스트."""
    from pptx.enum.shapes import MSO_SHAPE

    n = len(comp.items)
    item_h = min(int(h / n), Inches(0.6))
    circle_d = min(item_h - Inches(0.05), Inches(0.4))

    style_colors = {
        "accent": CL_ACCENT,
        "dark": CL_DARK,
        "grey": CL_GREY,
    }
    circle_color = style_colors.get(comp.style, CL_ACCENT)

    for i, item in enumerate(comp.items):
        iy = y + i * item_h

        # 원형 번호
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, iy, circle_d, circle_d,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = circle_color
        circle.line.fill.background()

        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{i + 1:02d}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = CL_WHITE
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER

        # 텍스트
        tx = slide.shapes.add_textbox(
            x + circle_d + Inches(0.12), iy,
            w - circle_d - Inches(0.12), item_h,
        )
        tf = tx.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = FONT_SIZE_BODY
        p.font.color.rgb = CL_BODY_TEXT
        p.font.name = FONT_BODY


# ============================================================
# Takeaway Bar (하단 핵심 메시지 바)
# ============================================================

def render_takeaway_bar(slide: Slide, comp, x, y, w, h):
    """하단 풀너비 핵심 메시지 바 — 좌측 오렌지 마커 포함."""
    bar_h = Inches(0.45)

    style_map = {
        "accent": (CL_ACCENT, CL_WHITE),
        "dark": (CL_DARK, CL_WHITE),
    }
    fill_c, text_c = style_map.get(comp.style, (CL_DARK, CL_WHITE))

    # 메인 바
    bar = slide.shapes.add_shape(1, x, y, w, bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = fill_c
    bar.line.fill.background()

    # 좌측 오렌지 마커
    if comp.style == "dark":
        marker = slide.shapes.add_shape(1, x, y, Inches(0.05), bar_h)
        marker.fill.solid()
        marker.fill.fore_color.rgb = CL_ACCENT
        marker.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = comp.message
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = text_c
    p.font.name = FONT_BODY


# ============================================================
# Numbered Quadrant (2×2 번호 정보 블록)
# ============================================================

def render_numbered_quadrant(slide: Slide, comp, x, y, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    n = len(comp.items)
    cols = 2
    rows = (n + 1) // 2
    gap = Inches(0.1)
    cell_w = int((w - gap * (cols - 1)) / cols)
    cell_h = int((h - gap * (rows - 1)) / rows)

    style_colors = {"accent": CL_ACCENT, "dark": CL_DARK, "grey": CL_GREY, "default": CL_GREY_MID}

    for idx, item in enumerate(comp.items):
        r, c = divmod(idx, cols)
        cx = x + c * (cell_w + gap)
        cy = y + r * (cell_h + gap)
        bar_c = style_colors.get(item.style, CL_GREY_MID)

        # 상단 색상 바
        bar = slide.shapes.add_shape(1, cx, cy, cell_w, Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_c
        bar.line.fill.background()

        # 본문 박스
        box = slide.shapes.add_shape(1, cx, cy + Inches(0.05), cell_w, cell_h - Inches(0.05))
        box.fill.solid()
        box.fill.fore_color.rgb = CL_WHITE
        box.line.color.rgb = CL_BORDER
        box.line.width = Emu(6350)

        # 번호 원
        _circle_d = Inches(0.3)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.1), cy + Inches(0.15), _circle_d, _circle_d)
        circ.fill.solid()
        circ.fill.fore_color.rgb = bar_c
        circ.line.fill.background()
        tf = circ.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = item.number
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = CL_WHITE
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER

        # 제목
        tx = slide.shapes.add_textbox(cx + Inches(0.5), cy + Inches(0.15), cell_w - Inches(0.6), Inches(0.3))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = item.title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = CL_BLACK
        p.font.name = FONT_BODY

        # 불릿
        if item.bullets:
            btx = slide.shapes.add_textbox(cx + Inches(0.15), cy + Inches(0.5), cell_w - Inches(0.3), cell_h - Inches(0.6))
            tf = btx.text_frame
            tf.word_wrap = True
            for bi, bullet in enumerate(item.bullets):
                p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                p.text = f"\u2022  {bullet}"
                p.font.size = Pt(8)
                p.font.color.rgb = CL_BODY_TEXT
                p.font.name = FONT_BODY
                p.space_after = Pt(2)


# ============================================================
# Harvey Ball Matrix
# ============================================================

def render_harvey_ball_matrix(slide: Slide, comp, x, y, w, h):
    from pptx.enum.shapes import MSO_SHAPE

    n_rows = len(comp.row_headers)
    n_cols = len(comp.col_headers)
    total_rows = n_rows + 1  # +header
    total_cols = n_cols + 1  # +row label

    row_h_val = int(h / total_rows)
    col_w_val = int(w / total_cols)
    label_col_w = int(col_w_val * 1.5)
    data_col_w = int((w - label_col_w) / n_cols)

    for ri in range(total_rows):
        for ci in range(total_cols):
            is_header_row = ri == 0
            is_label_col = ci == 0

            cw = label_col_w if is_label_col else data_col_w
            cx = x + (label_col_w if ci > 0 else 0) + (ci - 1) * data_col_w if ci > 0 else x
            if ci == 0:
                cx = x
            else:
                cx = x + label_col_w + (ci - 1) * data_col_w
            cy = y + ri * row_h_val

            if is_header_row and is_label_col:
                fill, text, text_c, bold = CL_WHITE, "", CL_BLACK, False
            elif is_header_row:
                fill, text = CL_DARK, comp.col_headers[ci - 1]
                text_c, bold = CL_WHITE, True
            elif is_label_col:
                fill, text = CL_GREY_LIGHT, comp.row_headers[ri - 1]
                text_c, bold = CL_BLACK, True
            else:
                fill = CL_WHITE if (ri % 2 == 0) else RGBColor(0xFA, 0xFA, 0xFA)
                text, text_c, bold = "", CL_BLACK, False

            # 셀 박스
            box = slide.shapes.add_shape(1, cx, cy, cw, row_h_val)
            box.fill.solid()
            box.fill.fore_color.rgb = fill
            box.line.color.rgb = CL_BORDER
            box.line.width = Emu(6350)

            if text:
                tf = box.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Inches(0.06)
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(8)
                p.font.bold = bold
                p.font.color.rgb = text_c
                p.font.name = FONT_BODY
                p.alignment = PP_ALIGN.CENTER if is_header_row else PP_ALIGN.LEFT

            # Harvey Ball (데이터 셀)
            if not is_header_row and not is_label_col:
                score = comp.scores[ri - 1][ci - 1] if ri - 1 < len(comp.scores) and ci - 1 < len(comp.scores[ri - 1]) else 0
                ball_d = min(row_h_val - Inches(0.1), Inches(0.25))
                bx = cx + int((cw - ball_d) / 2)
                by = cy + int((row_h_val - ball_d) / 2)

                # 배경 원 (빈)
                bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, ball_d, ball_d)
                bg.fill.solid()
                bg.fill.fore_color.rgb = CL_BORDER
                bg.line.fill.background()

                # 채움 원 (score에 따라 크기 조절로 근사)
                if score > 0:
                    fill_ratio = score / 100.0
                    fill_d = int(ball_d * (fill_ratio ** 0.5))
                    offset = int((ball_d - fill_d) / 2)
                    fg = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx + offset, by + offset, fill_d, fill_d)
                    fg.fill.solid()
                    fg.fill.fore_color.rgb = CL_ACCENT
                    fg.line.fill.background()


# ============================================================
# Vertical Flow (수직 화살표 체인)
# ============================================================

def render_vertical_flow(slide: Slide, comp, x, y, w, h):
    n = len(comp.steps)
    step_h = int((h - Inches(0.08) * (n - 1)) / n)
    label_w = Inches(1.5)
    detail_w = w - label_w - Inches(3.2)
    step_box_w = Inches(3.0)

    for i, step in enumerate(comp.steps):
        sy = y + i * (step_h + Inches(0.08))
        bar_c = CL_ACCENT if step.style == "accent" else CL_DARK

        # 좌: 번호 + 라벨 (대비 강화 — 회색 대신 다크)
        lbl = slide.shapes.add_shape(1, x, sy, label_w, step_h)
        lbl.fill.solid()
        lbl.fill.fore_color.rgb = bar_c
        lbl.line.fill.background()
        tf = lbl.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        p = tf.paragraphs[0]
        p.text = f"{i + 1:02d}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = CL_WHITE
        p.font.name = FONT_BODY

        # 중앙: 스텝 박스
        mid_x = x + label_w + Inches(0.1)
        # 상단 바
        bar = slide.shapes.add_shape(1, mid_x, sy, step_box_w, Inches(0.04))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_c
        bar.line.fill.background()
        box = slide.shapes.add_shape(1, mid_x, sy + Inches(0.04), step_box_w, step_h - Inches(0.04))
        box.fill.solid()
        box.fill.fore_color.rgb = CL_WHITE
        box.line.color.rgb = CL_BORDER
        box.line.width = Emu(6350)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = step.label
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = CL_BLACK
        p.font.name = FONT_BODY

        # 우: 상세 텍스트
        if step.detail:
            dtx = slide.shapes.add_textbox(mid_x + step_box_w + Inches(0.15), sy, detail_w, step_h)
            tf = dtx.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for di, line in enumerate(step.detail.split('\n')):
                p = tf.paragraphs[0] if di == 0 else tf.add_paragraph()
                p.text = f"\u2022  {line}" if line.strip() else ""
                p.font.size = Pt(8)
                p.font.color.rgb = CL_BODY_TEXT
                p.font.name = FONT_BODY
                p.space_after = Pt(2)

        # 화살표 (마지막 제외)
        if i < n - 1:
            arr_x = x + label_w + Inches(0.1) + int(step_box_w / 2) - Inches(0.08)
            arr_y = sy + step_h
            arr = slide.shapes.add_shape(1, arr_x, arr_y, Inches(0.16), Inches(0.06))
            arr.fill.solid()
            arr.fill.fore_color.rgb = CL_GREY
            arr.line.fill.background()


# ============================================================
# Funnel (깔때기)
# ============================================================

def render_funnel(slide: Slide, comp, x, y, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    n = len(comp.stages)
    stage_h = int((h - Inches(0.05) * (n - 1)) / n)
    max_w = w
    min_w = int(w * 0.3)

    for i, stage in enumerate(comp.stages):
        ratio = 1.0 - (i / max(n - 1, 1)) * 0.7
        sw = int(max_w * ratio)
        sx = x + int((max_w - sw) / 2)
        sy = y + i * (stage_h + Inches(0.05))
        color = _gradient_colors(n)[i] if i < n else CL_GREY

        # 사다리꼴 대신 사각형 (width 변화로 깔때기 표현)
        box = slide.shapes.add_shape(1, sx, sy, sw, stage_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.1)
        p = tf.paragraphs[0]
        label = stage.label
        if stage.value:
            label += f"  |  {stage.value}"
        p.text = label
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = CL_WHITE if i < (n + 1) // 2 else CL_BLACK
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER

        # 전환율 (우측)
        if stage.conversion:
            cvt = slide.shapes.add_textbox(sx + sw + Inches(0.1), sy, Inches(0.8), stage_h)
            tf = cvt.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = f"→ {stage.conversion}"
            p.font.size = Pt(7)
            p.font.color.rgb = CL_GREY
            p.font.name = FONT_BODY


# ============================================================
# RAG Table (R/A/G 상태 테이블)
# ============================================================

def render_rag_table(slide: Slide, comp, x, y, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    n = len(comp.items)
    row_h = min(int(h / (n + 1)), Inches(0.45))

    # 헤더행
    headers = ["항목", "상태", "추세", "비고"]
    col_ws = [int(w * 0.35), int(w * 0.15), int(w * 0.15), int(w * 0.35)]

    cx = x
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        box = slide.shapes.add_shape(1, cx, y, cw, row_h)
        box.fill.solid()
        box.fill.fore_color.rgb = CL_DARK
        box.line.color.rgb = CL_BORDER
        box.line.width = Emu(6350)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.06)
        p = tf.paragraphs[0]
        p.text = hdr
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = CL_WHITE
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER
        cx += cw

    # 데이터행
    rag_colors = {"green": RGBColor(0x27, 0xAE, 0x60), "amber": RGBColor(0xF3, 0x9C, 0x12), "red": RGBColor(0xC0, 0x39, 0x2B)}
    trend_symbols = {"up": "▲", "flat": "●", "down": "▼"}

    for ri, item in enumerate(comp.items):
        ry = y + (ri + 1) * row_h
        fill = CL_WHITE if ri % 2 == 0 else RGBColor(0xFA, 0xFA, 0xFA)
        cx = x

        for ci, cw in enumerate(col_ws):
            box = slide.shapes.add_shape(1, cx, ry, cw, row_h)
            box.fill.solid()
            box.fill.fore_color.rgb = fill
            box.line.color.rgb = CL_BORDER
            box.line.width = Emu(6350)

            if ci == 0:  # 항목명
                tf = box.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Inches(0.08)
                p = tf.paragraphs[0]
                p.text = item.name
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = CL_BLACK
                p.font.name = FONT_BODY
            elif ci == 1:  # RAG 원형
                rag_c = rag_colors.get(item.status, CL_GREY)
                d = min(row_h - Inches(0.08), Inches(0.22))
                circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                    cx + int((cw - d) / 2), ry + int((row_h - d) / 2), d, d)
                circ.fill.solid()
                circ.fill.fore_color.rgb = rag_c
                circ.line.fill.background()
            elif ci == 2:  # 추세
                tf = box.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                sym = trend_symbols.get(item.trend, "●")
                p.text = sym
                trend_c = rag_colors.get("green" if item.trend == "up" else "red" if item.trend == "down" else "amber", CL_GREY)
                p.font.size = Pt(10)
                p.font.color.rgb = trend_c
                p.font.name = FONT_BODY
                p.alignment = PP_ALIGN.CENTER
            elif ci == 3:  # 비고
                tf = box.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Inches(0.06)
                p = tf.paragraphs[0]
                p.text = item.note
                p.font.size = Pt(8)
                p.font.color.rgb = CL_BODY_TEXT
                p.font.name = FONT_BODY

            cx += cw


# ============================================================
# Divider (구분선)
# ============================================================

def render_divider(slide: Slide, comp: DividerComponent, x, y, w, h):
    line_y = y + int(h / 2)
    line = slide.shapes.add_shape(1, x, line_y, w, Emu(9144))
    line.fill.solid()
    line.fill.fore_color.rgb = CL_BORDER
    line.line.fill.background()


# ============================================================
# Dispatcher
# ============================================================

COMPONENT_RENDERERS = {
    "card": render_card,
    "text_block": render_text_block,
    "badge": render_badge,
    "kicker": render_kicker,
    "bullet": render_bullet,
    "table": render_table,
    "chart": None,
    "process_flow": render_process_flow,
    "chevron_process": render_chevron_process,
    "framework_matrix": render_framework_matrix,
    "numbered_circle": render_numbered_circle,
    "takeaway_bar": render_takeaway_bar,
    "numbered_quadrant": render_numbered_quadrant,
    "harvey_ball_matrix": render_harvey_ball_matrix,
    "vertical_flow": render_vertical_flow,
    "funnel": render_funnel,
    "rag_table": render_rag_table,
    "divider": render_divider,
    "image": None,
}


def render_component(slide: Slide, comp, x, y, w, h) -> None:
    renderer = COMPONENT_RENDERERS.get(comp.type)
    if renderer is None:
        print(f"[WARN] 미구현 컴포넌트: {comp.type}")
        return
    renderer(slide, comp, x, y, w, h)
