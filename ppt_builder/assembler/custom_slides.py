"""커스텀 슬라이드 렌더러 — 완성본 샘플 1:1 재현용.

일반 컴포넌트 조합으로는 한계가 있는 복잡 슬라이드를
전용 함수로 정밀 렌더링한다.
"""

from pptx.slide import Slide
from pptx.presentation import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor
from lxml import etree

# 회사 컬러
CL_ACCENT = RGBColor(0xFD, 0x51, 0x08)
CL_ACCENT_MID = RGBColor(0xFE, 0x7C, 0x39)
CL_ACCENT_LIGHT = RGBColor(0xFF, 0xAA, 0x72)
CL_DARK = RGBColor(0x1A, 0x1A, 0x1A)
CL_BLACK = RGBColor(0x00, 0x00, 0x00)
CL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CL_GREY = RGBColor(0xA1, 0xA8, 0xB3)
CL_GREY_MID = RGBColor(0xB5, 0xBC, 0xC4)
CL_GREY_LIGHT = RGBColor(0xCB, 0xD1, 0xD6)
CL_BG = RGBColor(0xF2, 0xF2, 0xF2)
CL_BORDER = RGBColor(0xDD, 0xDD, 0xDE)
FONT = "Arial"
FONT_T = "Georgia"


def _add_shadow(shape, blur=4, dist=2, color="808080"):
    """shape에 드롭 쉐도우를 추가한다."""
    try:
        sp_pr = shape._element.spPr
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        eff_lst = etree.SubElement(sp_pr, f"{{{ns}}}effectLst")
        shdw = etree.SubElement(eff_lst, f"{{{ns}}}outerShdw")
        shdw.set("blurRad", str(blur * 12700))
        shdw.set("dist", str(dist * 12700))
        shdw.set("dir", "2700000")
        srgb = etree.SubElement(shdw, f"{{{ns}}}srgbClr")
        srgb.set("val", color)
        alpha = etree.SubElement(srgb, f"{{{ns}}}alpha")
        alpha.set("val", "40000")
    except Exception:
        pass


def _rect(slide, x, y, w, h, fill=CL_WHITE, border=None, rounded=False):
    """사각형 (선택적 둥근 모서리)."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else 1
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Emu(6350)
    else:
        s.line.fill.background()
    if rounded:
        try:
            s.adjustments[0] = 0.03
        except:
            pass
    return s


def _text(slide, x, y, w, h, text, sz=9, bold=False, color=CL_BLACK,
          font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """텍스트 박스."""
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    # 줄바꿈 처리
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(1)
    return tx


def _connector(slide, x1, y1, x2, y2, color=CL_GREY, width=1.0, arrow=False):
    """커넥터 라인 (선택적 화살표)."""
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if arrow:
        try:
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            ln = c._element.spPr.find(f"{{{ns}}}ln")
            if ln is None:
                ln = etree.SubElement(c._element.spPr, f"{{{ns}}}ln")
            tail = etree.SubElement(ln, f"{{{ns}}}tailEnd")
            tail.set("type", "arrow")
            tail.set("w", "med")
            tail.set("len", "med")
        except:
            pass
    return c


def _circle(slide, x, y, d, fill=CL_ACCENT, text="", text_c=CL_WHITE, sz=9):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    if text:
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = True
        p.font.color.rgb = text_c
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
    return c


def _footer(slide, source_text=""):
    """완성본 스타일 푸터."""
    # 구분선
    _rect(slide, 0.3, 7.1, 9.4, 0.005, CL_BORDER)
    # Confidential
    _text(slide, 0.3, 7.15, 2.5, 0.15, "Strictly Private and Confidential", 5, False, CL_GREY)
    # 출처
    if source_text:
        _text(slide, 3.0, 7.15, 4.5, 0.15, source_text, 5, False, CL_GREY)
    # 로고
    _text(slide, 0.3, 7.32, 0.8, 0.12, "pwc", 7, True, CL_ACCENT)
    _text(slide, 8.8, 7.32, 1.0, 0.12, "HD현대", 7, True, CL_DARK, align=PP_ALIGN.RIGHT)


# ============================================================
# 샘플 1 (09): HD현대 조선 사업 추진 현황 — 3컬럼 전략 프레임워크
# ============================================================

def render_sample_09(prs: Presentation) -> Slide:
    """완성본 09번 1:1 재현 — 3병렬 전략 프레임워크."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # === 상단 헤더 (완성본: 얇은 다크 바 없이 텍스트만) ===
    _text(slide, 0.3, 0.15, 5.0, 0.25, "HD현대의 조선 사업 추진 현황", 12, True, CL_BLACK, FONT_T)
    _text(slide, 7.0, 0.18, 2.7, 0.2, "2. 프로젝트 추진 전략 > 1. 사업에 대한 이해", 6, False, CL_GREY, align=PP_ALIGN.RIGHT)
    # 제목 아래 오렌지 라인
    _rect(slide, 0.3, 0.42, 9.4, 0.005, CL_ACCENT)

    # === 서브 헤더 ===
    _text(slide, 0.3, 0.5, 9.4, 0.35,
          "최근 HD현대 조선 사업은 운영 구조(글로벌 Multi-Yard), 사업의 범위(함정, MRO 등), 기술 실현(FoS 디지털 전환) 측면에서 확장이 맞물리며 운영 복잡도와 난이도가 구조적으로 상승함",
          8, False, CL_BLACK)

    # === 3컬럼 본문 (비균등: 다크/흰/연회) ===
    col_w = 3.05
    col_gap = 0.08
    col_y = 0.95
    col_h = 5.8

    cols = [
        {
            "bg": CL_DARK, "text_c": CL_WHITE, "title_c": CL_WHITE,
            "label": "Group-Wide Expansion",
            "subtitle": "'단일 Yard 중심 운영에서 Multi-Yard &\nGlobal Operation 체계로 확장'",
            "visual_type": "map",
            "bullets": [
                "Yard 간 물량 재배치가 Key",
                "필리핀 Subic / 루마니아 Mangalia",
                "중국 Qingdao / 국내 울산·거제",
                "설계/조달 표준화 → Global Template",
                "실적·원가 통합 관리 체계 필수",
            ]
        },
        {
            "bg": CL_WHITE, "text_c": CL_BLACK, "title_c": CL_BLACK,
            "label": "Biz Portfolio Diversification",
            "subtitle": "'단일 선박 건조 중심 탈피에서\n복합 사업 구조로 다각화'",
            "visual_type": "portfolio",
            "bullets": [
                "상선 (Merchant) — Core Business",
                "해양 (MRO/FPU) — Growth Engine",
                "함정 (Naval) — New Business",
                "특수선 (Offshore) — Niche",
                "사업별 프로세스·원가·일정 차별화",
            ]
        },
        {
            "bg": CL_BG, "text_c": CL_BLACK, "title_c": CL_BLACK,
            "label": "Path to FoS & Beyond",
            "subtitle": "'디지털/차세대 기반\n기술 실현 체계 구축'",
            "visual_type": "fos",
            "bullets": [
                "Step 1: 현재 보이는 조선소 (Twin FoS)",
                "FoS 전략 = 디지털 플랫폼 전환",
                "Siemens PLM → 3D 설계",
                "AI/ML → Palantir 분석",
                "MASS/자율운항 + Humanoid 야드 자동화",
            ]
        },
    ]

    for i, col in enumerate(cols):
        cx = 0.3 + i * (col_w + col_gap)
        border = CL_BORDER if col["bg"] == CL_WHITE else None

        # 컬럼 배경
        _rect(slide, cx, col_y, col_w, col_h, col["bg"], border, rounded=True)
        _add_shadow(_rect(slide, cx, col_y, col_w, col_h, col["bg"], border, rounded=True))

        # 컬럼 제목
        _text(slide, cx + 0.1, col_y + 0.1, col_w - 0.2, 0.25,
              col["label"], 10, True, col["title_c"])

        # 서브타이틀
        _text(slide, cx + 0.1, col_y + 0.38, col_w - 0.2, 0.45,
              col["subtitle"], 7, False,
              CL_GREY_LIGHT if col["bg"] == CL_DARK else CL_GREY)

        # 시각 요소 영역 (중앙 — 다이어그램 placeholder)
        vis_y = col_y + 0.9
        vis_h = 2.0
        if col["visual_type"] == "map":
            # 세계지도 placeholder — 원형 + 마커 점
            _rect(slide, cx + 0.2, vis_y, col_w - 0.4, vis_h,
                  RGBColor(0x2D, 0x2D, 0x2D), rounded=True)
            _text(slide, cx + 0.3, vis_y + 0.3, col_w - 0.6, 0.3,
                  "Global Yard Network", 8, True, CL_WHITE, align=PP_ALIGN.CENTER)
            # 마커 점들
            markers = [(0.5, 0.8), (1.0, 0.6), (1.5, 1.0), (2.0, 0.7), (0.8, 1.3)]
            for mx, my in markers:
                _circle(slide, cx + mx, vis_y + my, 0.15, CL_ACCENT, "", CL_WHITE, 5)

        elif col["visual_type"] == "portfolio":
            # Product Portfolio 박스 다이어그램
            items = [
                ("상선", CL_ACCENT, 0.0),
                ("해양 MRO", CL_ACCENT_MID, 0.5),
                ("함정", CL_GREY, 1.0),
                ("특수선", CL_GREY_LIGHT, 1.5),
            ]
            for label, color, dy in items:
                bw = col_w - 0.4 - dy * 0.2
                _rect(slide, cx + 0.2 + dy * 0.1, vis_y + dy * 0.45 + 0.1, bw, 0.35, color, rounded=True)
                tc = CL_WHITE if color in (CL_ACCENT, CL_ACCENT_MID) else CL_BLACK
                _text(slide, cx + 0.3 + dy * 0.1, vis_y + dy * 0.45 + 0.12, bw - 0.2, 0.3,
                      label, 8, True, tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 화살표
            _text(slide, cx + col_w/2 - 0.15, vis_y + 0.05, 0.3, 0.2,
                  "▼", 10, False, CL_ACCENT, align=PP_ALIGN.CENTER)

        elif col["visual_type"] == "fos":
            # FoS 단계 다이어그램
            _rect(slide, cx + 0.2, vis_y, col_w - 0.4, 0.4, CL_ACCENT, rounded=True)
            _text(slide, cx + 0.3, vis_y + 0.05, col_w - 0.6, 0.3,
                  "FoS Strategy", 9, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 기술 아이콘 행
            techs = [("Siemens\nPLM", CL_DARK), ("AI\nPalantir", CL_ACCENT_MID), ("Humanoid", CL_GREY)]
            for j, (tech, tc) in enumerate(techs):
                tx = cx + 0.3 + j * 0.85
                _rect(slide, tx, vis_y + 0.55, 0.75, 0.55, tc, rounded=True)
                _text(slide, tx + 0.05, vis_y + 0.57, 0.65, 0.5,
                      tech, 6, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # Step 박스
            _rect(slide, cx + 0.2, vis_y + 1.25, col_w - 0.4, 0.6, CL_WHITE, CL_BORDER, rounded=True)
            _text(slide, cx + 0.3, vis_y + 1.3, col_w - 0.6, 0.5,
                  "Step 1: 현재 보이는 조선소 (Twin FoS)\nStep 2: FoS 전략 디지털 플랫폼 전환",
                  7, False, CL_BLACK)

        # 불릿 리스트 (하단)
        bullet_y = col_y + 3.1
        for bi, bullet in enumerate(col["bullets"]):
            _text(slide, cx + 0.15, bullet_y + bi * 0.28, col_w - 0.3, 0.25,
                  f"• {bullet}", 7, False, col["text_c"])

        # 하단 구분선
        _rect(slide, cx + 0.15, col_y + col_h - 0.5, col_w - 0.3, 0.003,
              CL_GREY_LIGHT if col["bg"] == CL_DARK else CL_BORDER)

        # 하단 참고 텍스트
        note_texts = [
            "MRO/서비스/해양·함정 포트폴리오",
            "Product Portfolio 구조 다각화",
            "Palantir 기반 실시간 의사결정"
        ]
        _text(slide, cx + 0.15, col_y + col_h - 0.45, col_w - 0.3, 0.35,
              note_texts[i], 6, False,
              CL_GREY_LIGHT if col["bg"] == CL_DARK else CL_GREY)

    # === 푸터 ===
    _footer(slide, "출처: HD현대 사업 현황 분석, PwC 전략컨설팅팀")

    return slide
