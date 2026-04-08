"""프레젠테이션 엔진 - 슬라이드 스키마를 .pptx로 변환."""

from pathlib import Path
from pptx import Presentation

from ..models.schema import PresentationSchema
from ..models.enums import SlideType
from .renderers.cover import CoverRenderer
from .renderers.section_divider import SectionDividerRenderer
from .renderers.conclusion import ConclusionRenderer
from .renderers.content import ContentSlideRenderer
from .renderers.template_inject import TemplateInjectRenderer


class PresentationEngine:
    """슬라이드 스키마를 받아서 .pptx를 생성하는 엔진."""

    RENDERERS = {
        SlideType.COVER: CoverRenderer(),
        SlideType.CONTENT: ContentSlideRenderer(),
        SlideType.SECTION_DIVIDER: SectionDividerRenderer(),
        SlideType.CONCLUSION: ConclusionRenderer(),
        SlideType.TEMPLATE: TemplateInjectRenderer(),
    }

    def __init__(self, template: Path | None = None):
        if template and template.exists():
            self.prs = Presentation(str(template))
        else:
            self.prs = Presentation()
            # 4:3 Standard
            self.prs.slide_width = 9144000    # 10.0 inches in EMU
            self.prs.slide_height = 6858000   # 7.5 inches in EMU

    def render(self, schema: PresentationSchema, output: Path) -> Path:
        """전체 프레젠테이션을 렌더링한다."""
        output.parent.mkdir(parents=True, exist_ok=True)

        for slide_def in schema.slides:
            renderer = self.RENDERERS.get(slide_def.type)
            if renderer is None:
                print(f"[WARN] 미구현 슬라이드 타입: {slide_def.type} — 건너뜀")
                continue
            renderer.render(self.prs, slide_def)

        self.prs.save(str(output))
        return output
