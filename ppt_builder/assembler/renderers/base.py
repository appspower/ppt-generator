"""렌더러 기본 클래스 — 헤더바 3종 + 그림자 + 아이콘 앵커."""

from abc import ABC, abstractmethod
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree

from ..styles import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    CL_WHITE, CL_BLACK, CL_BODY_TEXT, CL_ACCENT, CL_GREY, CL_GREY_LIGHT,
    CL_DARK, CL_BORDER,
    FONT_TITLE, FONT_BODY,
)


def add_shadow(shape, blur=3, dist=2, color="A0A0A0", alpha=35000):
    """shape에 드롭 쉐도우 추가 (과제3)."""
    try:
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = shape._element.spPr
        # 기존 effectLst 제거
        for old in sp_pr.findall(f"{{{ns}}}effectLst"):
            sp_pr.remove(old)
        eff = etree.SubElement(sp_pr, f"{{{ns}}}effectLst")
        shdw = etree.SubElement(eff, f"{{{ns}}}outerShdw")
        shdw.set("blurRad", str(blur * 12700))
        shdw.set("dist", str(dist * 12700))
        shdw.set("dir", "2700000")
        srgb = etree.SubElement(shdw, f"{{{ns}}}srgbClr")
        srgb.set("val", color)
        a = etree.SubElement(srgb, f"{{{ns}}}alpha")
        a.set("val", str(alpha))
    except Exception:
        pass


class BaseRenderer(ABC):

    @abstractmethod
    def render(self, prs: Presentation, slide_def) -> Slide:
        pass

    def add_blank_slide(self, prs: Presentation) -> Slide:
        layout = prs.slide_layouts[6]
        return prs.slides.add_slide(layout)

    # === 과제2: 헤더바 3종 ===

    def add_header_bar(self, slide: Slide, style: str = "standard") -> None:
        """헤더바 3종: standard(다크), minimal(라인만), accent(오렌지)."""
        if style == "minimal":
            # 헤더바 없이, 하단 1px 오렌지 라인만
            line = slide.shapes.add_shape(1, Inches(0.3), Inches(0.45), Inches(9.4), Emu(9525))
            line.fill.solid()
            line.fill.fore_color.rgb = CL_ACCENT
            line.line.fill.background()
        elif style == "accent":
            # 오렌지 배경 바
            bar = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Inches(0.55))
            bar.fill.solid()
            bar.fill.fore_color.rgb = CL_ACCENT
            bar.line.fill.background()
        else:  # standard
            bar = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Inches(0.55))
            bar.fill.solid()
            bar.fill.fore_color.rgb = CL_DARK
            bar.line.fill.background()
            # 좌측 오렌지 악센트 (2px)
            accent = slide.shapes.add_shape(1, 0, 0, Inches(0.03), Inches(0.55))
            accent.fill.solid()
            accent.fill.fore_color.rgb = CL_ACCENT
            accent.line.fill.background()

    def add_title(self, slide: Slide, text: str, style: str = "standard") -> None:
        """제목 — 헤더 스타일에 따라 색상 변경."""
        if style == "minimal":
            # 헤더바 없음 — minimal은 breadcrumb과 충돌 방지를 위해 폰트 조금 줄임
            txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(7.0), Inches(0.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = CL_BLACK
            p.font.name = FONT_TITLE
        else:
            # standard/accent — 흰 텍스트
            txBox = slide.shapes.add_textbox(Inches(0.15), Inches(0.08), Inches(6.5), Inches(0.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = CL_WHITE
            p.font.name = FONT_TITLE

    def add_breadcrumb(self, slide: Slide, text: str, style: str = "standard") -> None:
        if not text:
            return
        color = CL_GREY if style == "minimal" else CL_GREY_LIGHT
        bc = slide.shapes.add_textbox(Inches(7.0), Inches(0.12), Inches(2.7), Inches(0.25))
        tf = bc.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(6)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.RIGHT

    def add_header_message(self, slide: Slide, text: str) -> None:
        if not text:
            return
        txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(9.4), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(8)
        p.font.color.rgb = CL_BODY_TEXT
        p.font.name = FONT_BODY

    def add_footnote(self, slide: Slide, text: str) -> None:
        # 구분선
        line = slide.shapes.add_shape(1, Inches(0.3), Inches(7.1), Inches(9.4), Emu(4572))
        line.fill.solid()
        line.fill.fore_color.rgb = CL_BORDER
        line.line.fill.background()
        # Confidential
        conf = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(2.5), Inches(0.12))
        tf = conf.text_frame
        p = tf.paragraphs[0]
        p.text = "Strictly Private and Confidential"
        p.font.size = Pt(5)
        p.font.color.rgb = CL_GREY
        p.font.name = FONT_BODY
        # 출처
        if text:
            src = slide.shapes.add_textbox(Inches(3.0), Inches(7.15), Inches(4.5), Inches(0.12))
            tf = src.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(5)
            p.font.color.rgb = CL_GREY
            p.font.name = FONT_BODY
        # 로고
        logo_l = slide.shapes.add_textbox(Inches(0.3), Inches(7.3), Inches(0.8), Inches(0.12))
        tf = logo_l.text_frame
        p = tf.paragraphs[0]
        p.text = "pwc"
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = CL_ACCENT
        p.font.name = FONT_BODY
        logo_r = slide.shapes.add_textbox(Inches(8.8), Inches(7.3), Inches(1.0), Inches(0.12))
        tf = logo_r.text_frame
        p = tf.paragraphs[0]
        p.text = "HD\ud604\ub300"
        p.font.size = Pt(6)
        p.font.bold = True
        p.font.color.rgb = CL_DARK
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.RIGHT
