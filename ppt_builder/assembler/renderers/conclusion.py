"""종료/결론 슬라이드 렌더러 - 핑크 그라데이션 배경."""

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .base import BaseRenderer
from ..styles import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    CL_BLACK, CL_WHITE, CL_ACCENT, CL_ACCENT_LIGHT, CL_GREY,
    FONT_TITLE, FONT_BODY,
)
from ...models.schema import ConclusionSlide


class ConclusionRenderer(BaseRenderer):

    def render(self, prs: Presentation, slide_def: ConclusionSlide) -> Slide:
        slide = self.add_blank_slide(prs)

        # 그라데이션 배경 (연한 핑크/피치 → 흰)
        # python-pptx로 그라데이션이 제한적이므로 2개 박스로 근사
        # 상단: 연한 오렌지/피치 톤
        top_bg = slide.shapes.add_shape(
            1, 0, 0, SLIDE_WIDTH, Inches(4.0),
        )
        # 연한 피치/핑크 배경 (Light Orange를 배경으로)
        top_bg.fill.solid()
        top_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xD6)  # 매우 연한 피치
        top_bg.line.fill.background()

        # 하단: 흰색
        bot_bg = slide.shapes.add_shape(
            1, 0, Inches(4.0), SLIDE_WIDTH, Inches(3.5),
        )
        bot_bg.fill.solid()
        bot_bg.fill.fore_color.rgb = CL_WHITE
        bot_bg.line.fill.background()

        # 하단 악센트 라인
        line = slide.shapes.add_shape(
            1, 0, SLIDE_HEIGHT - Inches(0.06),
            SLIDE_WIDTH, Inches(0.06),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = CL_ACCENT
        line.line.fill.background()

        # "Thank you"
        title = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.2),
            Inches(8.0), Inches(1.5),
        )
        tf = title.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_def.title
        p.font.size = Pt(40)
        p.font.bold = False
        p.font.color.rgb = CL_BLACK
        p.font.name = FONT_TITLE

        # 부제/연락처
        if slide_def.subtitle:
            sub = slide.shapes.add_textbox(
                Inches(0.6), Inches(4.5),
                Inches(8.0), Inches(0.8),
            )
            tf = sub.text_frame
            p = tf.paragraphs[0]
            p.text = slide_def.subtitle
            p.font.size = Pt(11)
            p.font.color.rgb = CL_GREY
            p.font.name = FONT_BODY

        return slide
