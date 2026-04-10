"""본문(Content) 슬라이드 렌더러 — 공간 채우기 전략.

핵심 원칙: 슬라이드 전체를 채운다. 빈 공간을 남기지 않는다.
- 단일 레이아웃: 콘텐츠 영역 전체를 사용
- stacked: 비율대로 분배, TakeawayBar는 하단 고정
"""

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Emu

from .base import BaseRenderer
from ..layout_engine import calculate_positions_custom
from ..components import render_component
from ...models.schema import ContentSlide
from ...models.enums import ComponentType, LayoutType
from ..styles import CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H, FOOTER_Y


# TakeawayBar 고정 높이
TAKEAWAY_H = Inches(0.45)
TAKEAWAY_GAP = Inches(0.1)


class ContentSlideRenderer(BaseRenderer):

    def render(self, prs: Presentation, slide_def: ContentSlide) -> Slide:
        slide = self.add_blank_slide(prs)

        # 과제2: 헤더 스타일 전달
        hs = slide_def.header_style if hasattr(slide_def, 'header_style') else "standard"

        # 프레임 분기
        if slide_def.layout == LayoutType.FULLSCREEN:
            self._render_fullscreen(slide, slide_def)
        elif slide_def.layout == LayoutType.SIDEBAR:
            self.add_header_bar(slide, hs)
            self.add_title(slide, slide_def.title, hs)
            self.add_breadcrumb(slide, slide_def.breadcrumb, hs)
            self._render_sidebar(slide, slide_def)
        elif slide_def.layout == LayoutType.STACKED and slide_def.sections:
            self.add_header_bar(slide, hs)
            self.add_title(slide, slide_def.title, hs)
            self.add_breadcrumb(slide, slide_def.breadcrumb, hs)
            self.add_header_message(slide, slide_def.header_message)
            self._render_stacked(slide, slide_def)
        else:
            self.add_header_bar(slide, hs)
            self.add_title(slide, slide_def.title, hs)
            self.add_breadcrumb(slide, slide_def.breadcrumb, hs)
            self.add_header_message(slide, slide_def.header_message)
            self._render_single(slide, slide_def)

        self.add_footnote(slide, slide_def.footnote)
        return slide

    def _render_single(self, slide: Slide, slide_def: ContentSlide):
        """단일 레이아웃 — 전체 콘텐츠 영역을 채운다."""
        normal = []
        fullwidth = []
        for elem in slide_def.elements:
            if elem.type == ComponentType.TAKEAWAY_BAR:
                fullwidth.append(elem)
            else:
                normal.append(elem)

        # TakeawayBar 공간 확보 후, 나머지 전체를 일반 컴포넌트에 할당
        fw_total = (TAKEAWAY_H + TAKEAWAY_GAP) * len(fullwidth)
        normal_h = CONTENT_H - fw_total  # 축소하지 않음 — 전체를 채움

        if normal:
            # 과제1: col_ratios 전달
            cr = slide_def.col_ratios if hasattr(slide_def, 'col_ratios') else None
            positions = calculate_positions_custom(
                layout=slide_def.layout,
                n_elements=len(normal),
                n_cols=slide_def.n_cols,
                elements=normal,
                content_x=CONTENT_X, content_y=CONTENT_Y,
                content_w=CONTENT_W, content_h=normal_h,
                col_ratios=cr if cr else None,
            )
            for comp, (x, y, w, h) in zip(normal, positions):
                render_component(slide, comp, x, y, w, h)

        # TakeawayBar — 푸터 바로 위에 고정
        if fullwidth:
            fw_y = Inches(float(FOOTER_Y) / 914400) - TAKEAWAY_GAP - TAKEAWAY_H * len(fullwidth)
            for comp in fullwidth:
                render_component(slide, comp, CONTENT_X, fw_y, CONTENT_W, TAKEAWAY_H)
                fw_y += TAKEAWAY_H + TAKEAWAY_GAP

    def _render_stacked(self, slide: Slide, slide_def: ContentSlide):
        """복합(stacked) 레이아웃 — 비율대로 전체 영역을 채운다."""
        takeaways = []
        content_sections = []
        for section in slide_def.sections:
            normal = [e for e in section.elements if e.type != ComponentType.TAKEAWAY_BAR]
            taws = [e for e in section.elements if e.type == ComponentType.TAKEAWAY_BAR]
            takeaways.extend(taws)
            if normal:
                content_sections.append((section, normal))

        # 사용 가능한 전체 높이 (TakeawayBar + 갭 제외)
        fw_space = (TAKEAWAY_H + TAKEAWAY_GAP) * len(takeaways)
        gap = Inches(0.08)
        n_gaps = max(0, len(content_sections) - 1)
        available_h = CONTENT_H - fw_space - gap * n_gaps

        # 비율대로 전체 분배 (축소하지 않음)
        total_ratio = sum(s.height_ratio for s, _ in content_sections)
        current_y = CONTENT_Y

        for si, (section, normal_elems) in enumerate(content_sections):
            ratio = section.height_ratio / total_ratio
            section_h = int(available_h * ratio)

            positions = calculate_positions_custom(
                layout=section.layout,
                n_elements=len(normal_elems),
                n_cols=section.n_cols,
                elements=normal_elems,
                content_x=CONTENT_X, content_y=current_y,
                content_w=CONTENT_W, content_h=section_h,
            )
            for comp, (x, y, w, h) in zip(normal_elems, positions):
                render_component(slide, comp, x, y, w, h)

            current_y += section_h + gap

        # TakeawayBar — 푸터 바로 위에 고정
        if takeaways:
            fw_y = Inches(float(FOOTER_Y) / 914400) - TAKEAWAY_GAP - TAKEAWAY_H * len(takeaways)
            for comp in takeaways:
                render_component(slide, comp, CONTENT_X, fw_y, CONTENT_W, TAKEAWAY_H)
                fw_y += TAKEAWAY_H + TAKEAWAY_GAP

    def _render_sidebar(self, slide: Slide, slide_def: ContentSlide):
        """프레임 B: 좌측 사이드바 + 우측 콘텐츠."""
        from pptx.enum.text import MSO_ANCHOR
        from ..styles import CL_ACCENT, CL_DARK, CL_BG, CL_WHITE, CL_BLACK, CL_BORDER, CL_GREY

        sidebar_w = Inches(2.5)
        main_x = Inches(3.1)
        main_w = Inches(6.6)
        top_y = Inches(0.65)

        # 사이드바 항목
        items = slide_def.sidebar_items or ["항목 1", "항목 2", "항목 3"]
        active = slide_def.sidebar_active
        item_h = Inches(0.55)

        for i, item in enumerate(items):
            iy = top_y + Inches(i * 0.6)
            is_active = (i == active)
            fill = CL_ACCENT if is_active else CL_BG
            text_c = CL_WHITE if is_active else CL_BLACK
            border = None if is_active else CL_BORDER

            box = slide.shapes.add_shape(1, Inches(0.3), iy, sidebar_w, item_h)
            box.fill.solid()
            box.fill.fore_color.rgb = fill
            if border:
                box.line.color.rgb = border
                box.line.width = Emu(6350)
            else:
                box.line.fill.background()

            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.1)
            p = tf.paragraphs[0]
            p.text = item
            from pptx.util import Pt
            p.font.size = Pt(9)
            p.font.bold = is_active
            p.font.color.rgb = text_c
            p.font.name = "Arial"

        # 우측 콘텐츠 — 일반 컴포넌트 배치
        normal = [e for e in slide_def.elements if e.type != ComponentType.TAKEAWAY_BAR]
        takeaways = [e for e in slide_def.elements if e.type == ComponentType.TAKEAWAY_BAR]

        main_h = CONTENT_H - (TAKEAWAY_H + TAKEAWAY_GAP) * len(takeaways)

        if normal:
            positions = calculate_positions_custom(
                layout=LayoutType.FULL,
                n_elements=len(normal),
                n_cols=1,
                elements=normal,
                content_x=main_x, content_y=Inches(0.65),
                content_w=main_w, content_h=main_h,
            )
            for comp, (cx, cy, cw, ch) in zip(normal, positions):
                render_component(slide, comp, cx, cy, cw, ch)

        if takeaways:
            fw_y = Inches(float(FOOTER_Y) / 914400) - TAKEAWAY_GAP - TAKEAWAY_H
            for comp in takeaways:
                render_component(slide, comp, CONTENT_X, fw_y, CONTENT_W, TAKEAWAY_H)
                fw_y += TAKEAWAY_H + TAKEAWAY_GAP

    def _render_fullscreen(self, slide: Slide, slide_def: ContentSlide):
        """프레임 C: 헤더바 없이 전체화면 — 제목은 좌상단 텍스트."""
        from pptx.util import Pt
        from ..styles import CL_BLACK, CL_ACCENT, CL_GREY, CL_GREY_LIGHT, CL_BORDER

        # 제목 (좌상단, 헤더바 없이)
        tx = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(7.0), Inches(0.4))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = slide_def.title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = CL_BLACK
        p.font.name = "Georgia"

        # 브레드크럼
        if slide_def.breadcrumb:
            bc = slide.shapes.add_textbox(Inches(7.5), Inches(0.25), Inches(2.2), Inches(0.2))
            tf = bc.text_frame
            p = tf.paragraphs[0]
            p.text = slide_def.breadcrumb
            p.font.size = Pt(7)
            p.font.color.rgb = CL_GREY
            p.font.name = "Arial"
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.RIGHT

        # 제목 아래 구분선
        line = slide.shapes.add_shape(1, Inches(0.3), Inches(0.65), Inches(9.4), Emu(6350))
        line.fill.solid()
        line.fill.fore_color.rgb = CL_ACCENT
        line.line.fill.background()

        # 콘텐츠 — 전체 영역 (0.75" ~ 7.1")
        fs_y = Inches(0.75)
        fs_h = Inches(6.35)

        normal = [e for e in slide_def.elements if e.type != ComponentType.TAKEAWAY_BAR]
        takeaways = [e for e in slide_def.elements if e.type == ComponentType.TAKEAWAY_BAR]

        fw_space = (TAKEAWAY_H + TAKEAWAY_GAP) * len(takeaways)
        normal_h = fs_h - fw_space

        if normal:
            positions = calculate_positions_custom(
                layout=LayoutType.FULL if slide_def.n_cols <= 1 else LayoutType.COLUMNS,
                n_elements=len(normal),
                n_cols=slide_def.n_cols,
                elements=normal,
                content_x=CONTENT_X, content_y=fs_y,
                content_w=CONTENT_W, content_h=normal_h,
            )
            for comp, (cx, cy, cw, ch) in zip(normal, positions):
                render_component(slide, comp, cx, cy, cw, ch)

        if takeaways:
            fw_y = Inches(float(FOOTER_Y) / 914400) - TAKEAWAY_GAP - TAKEAWAY_H
            for comp in takeaways:
                render_component(slide, comp, CONTENT_X, fw_y, CONTENT_W, TAKEAWAY_H)
                fw_y += TAKEAWAY_H + TAKEAWAY_GAP
