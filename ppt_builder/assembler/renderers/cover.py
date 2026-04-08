"""표지(Cover) 슬라이드 렌더러 - 오렌지 평행사변형 장식."""

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

from .base import BaseRenderer
from ..styles import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    CL_BLACK, CL_BODY_TEXT, CL_GREY, CL_ACCENT, CL_ACCENT_MID, CL_WHITE,
    FONT_TITLE, FONT_BODY,
    FONT_SIZE_COVER_TITLE, FONT_SIZE_SUBTITLE,
)
from ...models.schema import CoverSlide


class CoverRenderer(BaseRenderer):

    def render(self, prs: Presentation, slide_def: CoverSlide) -> Slide:
        slide = self.add_blank_slide(prs)

        # --- 오렌지 평행사변형 장식 (2개) ---
        self._add_parallelogram(
            slide, Inches(3.2), Inches(3.0), Inches(3.5), Inches(1.2),
            CL_ACCENT, skew=-15,
        )
        self._add_parallelogram(
            slide, Inches(5.5), Inches(3.0), Inches(3.5), Inches(1.2),
            CL_ACCENT_MID, skew=-15,
        )

        # --- 하단 악센트 라인 ---
        line = slide.shapes.add_shape(
            1, 0, SLIDE_HEIGHT - Inches(0.08),
            SLIDE_WIDTH, Inches(0.08),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = CL_ACCENT
        line.line.fill.background()

        # --- 제목 ---
        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.5),
            Inches(6.0), Inches(1.5),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_def.title
        p.font.size = FONT_SIZE_COVER_TITLE
        p.font.bold = False
        p.font.color.rgb = CL_BLACK
        p.font.name = FONT_TITLE

        # --- 부제 ---
        if slide_def.subtitle:
            sub = slide.shapes.add_textbox(
                Inches(0.6), Inches(3.2),
                Inches(6.0), Inches(0.6),
            )
            tf = sub.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_def.subtitle
            p.font.size = FONT_SIZE_SUBTITLE
            p.font.color.rgb = CL_BODY_TEXT
            p.font.name = FONT_BODY

        # --- 작성자 & 날짜 (좌하단) ---
        info_parts = [p for p in [slide_def.author, slide_def.date] if p]
        if info_parts:
            for i, part in enumerate(info_parts):
                info = slide.shapes.add_textbox(
                    Inches(0.6), Inches(6.0 + i * 0.35),
                    Inches(5.0), Inches(0.3),
                )
                tf = info.text_frame
                p = tf.paragraphs[0]
                p.text = part
                p.font.size = Pt(11)
                p.font.color.rgb = CL_GREY
                p.font.name = FONT_BODY

        return slide

    def _add_parallelogram(self, slide, x, y, w, h, color, skew=-15):
        """평행사변형 도형을 추가한다."""
        # MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM = 7
        shape = slide.shapes.add_shape(7, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.rotation = 0.0
