"""Executive Summary 슬라이드 렌더러."""

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .base import BaseRenderer
from ..styles import (
    CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH,
    COLOR_TEXT, COLOR_ACCENT, COLOR_PRIMARY,
    FONT_BODY, FONT_SIZE_BODY, FONT_SIZE_BULLET,
)
from ...models.schema import ExecSummarySlide


class ExecSummaryRenderer(BaseRenderer):

    def render(self, prs: Presentation, slide_def: ExecSummarySlide) -> Slide:
        slide = self.add_blank_slide(prs)
        self.add_title(slide, slide_def.title)

        # Bullets
        bullet_top = CONTENT_TOP + Inches(0.2)
        txBox = slide.shapes.add_textbox(
            CONTENT_LEFT + Inches(0.3), bullet_top,
            CONTENT_WIDTH - Inches(0.6), Inches(4.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(slide_def.bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = FONT_SIZE_BODY
            p.font.color.rgb = COLOR_TEXT
            p.font.name = FONT_BODY
            p.space_after = Pt(8)
            p.level = 0
            # 불릿 마커
            p.text = f"\u25A0  {bullet}"  # ■ 마커
            run = p.runs[0]
            run.font.size = FONT_SIZE_BODY
            run.font.color.rgb = COLOR_TEXT
            run.font.name = FONT_BODY

        # Highlight box (선택적)
        if slide_def.highlight_box:
            box_top = bullet_top + Inches(4.3)
            from pptx.util import Emu
            box = slide.shapes.add_shape(
                1,  # Rectangle
                CONTENT_LEFT + Inches(0.3), box_top,
                CONTENT_WIDTH - Inches(0.6), Inches(0.8),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = COLOR_ACCENT
            box.line.fill.background()

            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_def.highlight_box
            p.font.size = Pt(13)
            p.font.bold = True
            from pptx.dml.color import RGBColor
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER

        return slide
