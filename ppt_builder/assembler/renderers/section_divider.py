"""섹션 구분 슬라이드 렌더러 - 좌측 오렌지 바 + 다크 배경."""

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .base import BaseRenderer
from ..styles import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    CL_BLACK, CL_DARK, CL_WHITE, CL_ACCENT,
    FONT_TITLE, FONT_SIZE_SECTION,
)
from ...models.schema import SectionDividerSlide


class SectionDividerRenderer(BaseRenderer):

    def render(self, prs: Presentation, slide_def: SectionDividerSlide) -> Slide:
        slide = self.add_blank_slide(prs)

        # 배경 (다크)
        bg = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
        bg.fill.solid()
        bg.fill.fore_color.rgb = CL_DARK
        bg.line.fill.background()

        # 좌측 오렌지 수직 바 (Sacred Rectangle)
        bar = slide.shapes.add_shape(
            1, 0, Inches(2.0),
            Inches(0.35), Inches(4.0),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = CL_ACCENT
        bar.line.fill.background()

        # 섹션 번호
        if slide_def.section_number is not None:
            num = slide.shapes.add_textbox(
                Inches(0.7), Inches(2.2), Inches(3), Inches(1),
            )
            tf = num.text_frame
            p = tf.paragraphs[0]
            p.text = f"{slide_def.section_number:02d}"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = CL_ACCENT
            p.font.name = FONT_TITLE

        # 섹션 제목
        title_top = Inches(3.2) if slide_def.section_number else Inches(2.8)
        title = slide.shapes.add_textbox(
            Inches(0.7), title_top, Inches(8.5), Inches(1.5),
        )
        tf = title.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_def.title
        p.font.size = FONT_SIZE_SECTION
        p.font.bold = True
        p.font.color.rgb = CL_WHITE
        p.font.name = FONT_TITLE

        return slide
