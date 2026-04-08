"""타이틀(표지) 슬라이드 렌더러."""

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .base import BaseRenderer
from ..styles import (
    SLIDE_WIDTH, SLIDE_HEIGHT,
    COLOR_PRIMARY, COLOR_TEXT, COLOR_TEXT_LIGHT,
    FONT_TITLE, FONT_BODY,
)
from ...models.schema import TitleSlide


class TitleRenderer(BaseRenderer):

    def render(self, prs: Presentation, slide_def: TitleSlide) -> Slide:
        slide = self.add_blank_slide(prs)

        # Main title - 중앙 배치
        title_top = Inches(2.5)
        txBox = slide.shapes.add_textbox(
            Inches(1.5), title_top,
            SLIDE_WIDTH - Inches(3), Inches(1.2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_def.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.font.name = FONT_TITLE
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if slide_def.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1.5), title_top + Inches(1.4),
                SLIDE_WIDTH - Inches(3), Inches(0.8),
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_def.subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_TEXT
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER

        # Author & Date - 하단
        if slide_def.author or slide_def.date:
            info_text = " | ".join(
                part for part in [slide_def.author, slide_def.date] if part
            )
            info_box = slide.shapes.add_textbox(
                Inches(1.5), SLIDE_HEIGHT - Inches(1.5),
                SLIDE_WIDTH - Inches(3), Inches(0.5),
            )
            tf = info_box.text_frame
            p = tf.paragraphs[0]
            p.text = info_text
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT_LIGHT
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER

        # 하단 구분선
        from pptx.util import Emu
        line = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(2), title_top + Inches(1.2),
            SLIDE_WIDTH - Inches(4), Emu(28000),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_PRIMARY
        line.line.fill.background()

        return slide
