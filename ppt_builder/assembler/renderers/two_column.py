"""2단 레이아웃 슬라이드 렌더러."""

from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .base import BaseRenderer
from ..styles import (
    CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH,
    COLOR_TEXT, COLOR_PRIMARY, COLOR_TEXT_LIGHT,
    FONT_BODY, FONT_TITLE, FONT_SIZE_BODY,
)
from ...models.schema import TwoColumnSlide, ColumnContent


class TwoColumnRenderer(BaseRenderer):

    def render(self, prs: Presentation, slide_def: TwoColumnSlide) -> Slide:
        slide = self.add_blank_slide(prs)
        self.add_title(slide, slide_def.title)

        col_width = (CONTENT_WIDTH - Inches(0.5)) / 2  # 0.5" gap
        col_height = Inches(4.5)

        # Left column
        self._render_column(
            slide, slide_def.left,
            CONTENT_LEFT, CONTENT_TOP + Inches(0.2),
            col_width, col_height,
        )

        # Right column
        self._render_column(
            slide, slide_def.right,
            CONTENT_LEFT + col_width + Inches(0.5), CONTENT_TOP + Inches(0.2),
            col_width, col_height,
        )

        # 중앙 구분선
        from pptx.util import Emu
        divider = slide.shapes.add_shape(
            1,  # Rectangle
            CONTENT_LEFT + col_width + Inches(0.22),
            CONTENT_TOP + Inches(0.2),
            Emu(12700),  # 1pt width
            col_height,
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = COLOR_TEXT_LIGHT
        divider.line.fill.background()

        return slide

    def _render_column(self, slide, content: ColumnContent, left, top, width, height):
        """한 컬럼의 내용을 렌더링."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        first_para = True

        # Heading
        if content.heading:
            p = tf.paragraphs[0]
            p.text = content.heading
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_PRIMARY
            p.font.name = FONT_TITLE
            p.space_after = Pt(8)
            first_para = False

        # Body text
        if content.text:
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            p.text = content.text
            p.font.size = FONT_SIZE_BODY
            p.font.color.rgb = COLOR_TEXT
            p.font.name = FONT_BODY
            p.space_after = Pt(6)
            first_para = False

        # Bullets
        for bullet in content.bullets:
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            p.text = f"\u25A0  {bullet}"
            p.font.size = FONT_SIZE_BODY
            p.font.color.rgb = COLOR_TEXT
            p.font.name = FONT_BODY
            p.space_after = Pt(4)
            first_para = False
