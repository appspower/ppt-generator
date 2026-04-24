"""Track 1 Stage 1 — 슬라이드 메타데이터 추출 + SimHash dedup.

입력: 마스터 템플릿 .pptx (1,251장)
출력: slide_meta.json (각 슬라이드의 파싱 메타)

핵심:
  - iter_leaf_shapes로 그룹 재귀 (Phase D 평탄화 재사용)
  - SimHash(텍스트) + 구조 시그니처로 dedup 후보 생성
  - 덱 경계 감지 시그니처 (title, page_number, footer)도 함께 추출
"""
from __future__ import annotations

import hashlib
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from ..template.edit_ops import iter_leaf_shapes
from .schemas import SlideMeta

# 페이지 번호 감지 (숫자만, 또는 "P.숫자", "/숫자")
_PAGE_NUM_RE = re.compile(r"^\s*(?:p\.?\s*)?(\d{1,3})\s*(?:/\s*\d+)?\s*$", re.IGNORECASE)
_KOREAN_RE = re.compile(r"[\uac00-\ud7a3]")
_ENGLISH_RE = re.compile(r"[a-zA-Z]")


def _simhash_64(text: str) -> str:
    """간단한 64-bit SimHash. 토큰 → md5 → 비트 투표.

    dedup 용이라 완벽할 필요는 없음 — near-duplicate 감지가 목표.
    """
    if not text.strip():
        return "0" * 16

    # 토큰화 (단어 + 한글 음절 단위)
    tokens = re.findall(r"\w+", text)
    if not tokens:
        return "0" * 16

    # 토큰별 가중치 1
    bit_votes = [0] * 64
    for tok in tokens:
        h = int(hashlib.md5(tok.encode("utf-8")).hexdigest()[:16], 16)
        for i in range(64):
            if (h >> i) & 1:
                bit_votes[i] += 1
            else:
                bit_votes[i] -= 1

    # 양수 = 1, 음수 = 0
    out = 0
    for i in range(64):
        if bit_votes[i] > 0:
            out |= (1 << i)
    return f"{out:016x}"


def _is_page_number(text: str) -> bool:
    t = text.strip()
    if len(t) > 12:
        return False
    return bool(_PAGE_NUM_RE.match(t))


def _extract_title(slide) -> str:
    """가장 큰 폰트 size + 위치 상단인 text shape 감지.

    Phase A1 관찰: 슬라이드 상단 2인치 내 가장 큰 텍스트가 실질 title.
    """
    best = ("", 0)  # (text, font_size_pt)
    for idx, sh in iter_leaf_shapes(slide):
        if not sh.has_text_frame:
            continue
        # 위치 필터 (상단 2인치 = 1828800 EMU)
        try:
            top = sh.top or 0
            if top > Emu(2 * 914400):  # 2 inches
                continue
        except Exception:
            continue

        tf = sh.text_frame
        if not tf.text.strip():
            continue

        # 최대 폰트 size
        max_size = 0
        for para in tf.paragraphs:
            for run in para.runs:
                sz = run.font.size
                if sz is not None:
                    pt = sz.pt
                    if pt > max_size:
                        max_size = pt

        if max_size > best[1]:
            best = (tf.text.strip()[:80], max_size)

    return best[0]


def _extract_page_number(slide) -> str:
    """하단 1.5인치 내에서 페이지 번호 패턴 탐색."""
    from pptx.util import Inches
    # slide.part.parent는 pres — 슬라이드 높이 얻기
    try:
        slide_h = slide.part.package.presentation_part.presentation.slide_height
    except Exception:
        slide_h = Inches(7.5)

    bottom_threshold = slide_h - Inches(1.5)

    for idx, sh in iter_leaf_shapes(slide):
        if not sh.has_text_frame:
            continue
        try:
            top = sh.top or 0
            if top < bottom_threshold:
                continue
        except Exception:
            continue

        txt = sh.text_frame.text.strip()
        if _is_page_number(txt):
            return txt
    return ""


def _extract_footer(slide) -> str:
    """하단 1인치 내 페이지 번호 아닌 긴 텍스트 = footer 후보."""
    from pptx.util import Inches
    try:
        slide_h = slide.part.package.presentation_part.presentation.slide_height
    except Exception:
        slide_h = Inches(7.5)
    bottom_threshold = slide_h - Inches(1)

    for idx, sh in iter_leaf_shapes(slide):
        if not sh.has_text_frame:
            continue
        try:
            top = sh.top or 0
            if top < bottom_threshold:
                continue
        except Exception:
            continue

        txt = sh.text_frame.text.strip()
        if txt and not _is_page_number(txt) and len(txt) > 5:
            return txt[:80]
    return ""


def _structure_signature(counts: dict) -> str:
    """구조 시그니처 — 클러스터링용."""
    parts = []
    # 주요 타입만 (많은 것부터)
    if counts.get("paragraph", 0) > 0:
        parts.append(f"para{counts['paragraph']}")
    if counts.get("picture", 0) > 0:
        parts.append(f"pic{counts['picture']}")
    if counts.get("table", 0) > 0:
        parts.append(f"tbl{counts['table']}")
    if counts.get("chart", 0) > 0:
        parts.append(f"chart{counts['chart']}")
    if counts.get("group", 0) > 0:
        parts.append(f"grp{counts['group']}")
    return "|".join(parts) if parts else "empty"


def extract_slide_meta(slide, slide_index: int) -> SlideMeta:
    """단일 슬라이드의 메타데이터 추출."""
    from pptx.shapes.picture import Picture

    # 카운트
    total_shapes = 0
    leaf_shapes = 0
    picture_count = 0
    chart_count = 0
    table_count = 0
    group_count = 0
    paragraph_count = 0
    placeholder_count = 0
    all_text_parts: list[str] = []

    # 슬라이드 레이아웃
    try:
        layout_name = slide.slide_layout.name or ""
    except Exception:
        layout_name = ""

    # 최상위 + 평탄화
    for sh in slide.shapes:
        total_shapes += 1
        if getattr(sh, "shape_type", None) == 6:  # GROUP
            group_count += 1

    for idx, sh in iter_leaf_shapes(slide):
        leaf_shapes += 1

        # 타입 분류
        if isinstance(sh, Picture):
            picture_count += 1
        elif getattr(sh, "has_chart", False):
            chart_count += 1
        elif getattr(sh, "has_table", False):
            table_count += 1

        # 텍스트
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                paragraph_count += 1
                txt = para.text
                if txt:
                    all_text_parts.append(txt)
                    # `~~`는 자체가 placeholder (둘러싸는 형식 아님). Phase A1 검수 기준.
                    placeholder_count += txt.count("~~")

    text_total = "\n".join(all_text_parts)
    korean_char_count = len(_KOREAN_RE.findall(text_total))
    english_char_count = len(_ENGLISH_RE.findall(text_total))
    text_hash_sim = _simhash_64(text_total)

    # 구조 시그니처
    structure_sig = _structure_signature({
        "paragraph": paragraph_count,
        "picture": picture_count,
        "table": table_count,
        "chart": chart_count,
        "group": group_count,
    })

    return SlideMeta(
        slide_index=slide_index,
        layout_name=layout_name,
        shape_count_total=total_shapes,
        shape_count_leaf=leaf_shapes,
        picture_count=picture_count,
        chart_count=chart_count,
        table_count=table_count,
        group_count=group_count,
        text_total=text_total,
        text_hash_sim=text_hash_sim,
        placeholder_count=placeholder_count,
        korean_char_count=korean_char_count,
        english_char_count=english_char_count,
        structure_sig=structure_sig,
        has_group=group_count > 0,
        has_smartart=False,  # Phase A1: 1251장 중 1장만. 무시
        title_text=_extract_title(slide),
        page_number_text=_extract_page_number(slide),
        footer_text=_extract_footer(slide),
    )


def run_extract(pptx_path: Path, output_path: Path, start: int = 0, end: int | None = None) -> list[SlideMeta]:
    """마스터 템플릿 전체를 순회하며 SlideMeta 추출."""
    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    end = end or total
    end = min(end, total)

    result: list[SlideMeta] = []
    for i in range(start, end):
        if i % 50 == 0:
            print(f"  [{i}/{end}] processing...")
        try:
            meta = extract_slide_meta(prs.slides[i], i)
            result.append(meta)
        except Exception as e:
            print(f"  [ERR] slide {i}: {type(e).__name__}: {e}")
            # 최소 SlideMeta (실패도 기록)
            result.append(SlideMeta(slide_index=i, structure_sig="error"))

    # 저장
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(
            [m.model_dump() for m in result],
            f, ensure_ascii=False, indent=2,
        )
    print(f"  saved {len(result)} metas -> {output_path}")
    return result


if __name__ == "__main__":
    ROOT = Path(__file__).resolve().parent.parent.parent
    src = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
    out = ROOT / "output" / "catalog" / "slide_meta.json"
    run_extract(src, out)
