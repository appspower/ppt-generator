"""Track 1 Stage 2d — 슬롯별 max_chars 측정.

근거
----
`ppt_builder/evaluate.py:149-172` overflow 규칙:
  텍스트 밀도 > 200 chars/in² → CRITICAL (−15점)

이 임계값을 역산해서 각 `~~` placeholder shape의 "안전 최대 글자 수"를
산출. Phase A2의 `SlotSchema.max_chars` 를 채운다.

공식
----
  max_chars(slot) = bbox_area_in2 × 200 × 0.95

(0.95 = safety margin — evaluate.py HIGH severity 80% 지점)

출력
----
  slot_schemas.json : list[dict]
    { slide_index, slot_id, shape_type, bbox_area_in2, max_chars,
      left_pct, top_pct, width_pct, height_pct,
      has_placeholder, layout_name }
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from ..template.edit_ops import iter_leaf_shapes


CHARS_PER_IN2_CRITICAL = 200.0     # evaluate.py:160
SAFETY_MARGIN = 0.95                # HIGH 임계의 80%


def _emu_to_in2(width_emu: int, height_emu: int) -> float:
    """EMU bbox → in² (1 inch = 914400 EMU)."""
    if width_emu <= 0 or height_emu <= 0:
        return 0.0
    return (width_emu / 914400.0) * (height_emu / 914400.0)


def extract_slots_for_slide(
    slide,
    slide_index: int,
    slide_w_emu: int,
    slide_h_emu: int,
) -> list[dict]:
    """단일 슬라이드 → 슬롯 리스트."""
    slots: list[dict] = []
    try:
        layout_name = slide.slide_layout.name or ""
    except Exception:
        layout_name = ""

    for flat_idx, sh in iter_leaf_shapes(slide):
        try:
            left = int(sh.left or 0)
            top = int(sh.top or 0)
            width = int(sh.width or 0)
            height = int(sh.height or 0)
        except Exception:
            continue
        if width <= 0 or height <= 0:
            continue

        has_tf = bool(getattr(sh, "has_text_frame", False))
        placeholder_count = 0
        raw_text = ""
        if has_tf:
            try:
                raw_text = sh.text_frame.text or ""
                placeholder_count = raw_text.count("~~")
            except Exception:
                pass

        area = _emu_to_in2(width, height)
        max_chars = int(area * CHARS_PER_IN2_CRITICAL * SAFETY_MARGIN)

        slots.append({
            "slide_index": slide_index,
            "slot_id": f"slide_{slide_index:04d}:{flat_idx}",
            "flat_idx": flat_idx,
            "shape_type": str(sh.shape_type) if hasattr(sh, "shape_type") else "",
            "bbox_area_in2": round(area, 3),
            "max_chars": max_chars,
            "left_pct": round(left / max(slide_w_emu, 1), 3),
            "top_pct": round(top / max(slide_h_emu, 1), 3),
            "width_pct": round(width / max(slide_w_emu, 1), 3),
            "height_pct": round(height / max(slide_h_emu, 1), 3),
            "has_text_frame": has_tf,
            "placeholder_count": placeholder_count,
            "has_placeholder": placeholder_count > 0,
            "layout_name": layout_name,
        })
    return slots


def run_stage_2d(
    pptx_path: Path,
    output_path: Path,
) -> dict:
    prs = Presentation(str(pptx_path))
    slide_w = int(prs.slide_width or Emu(9144000))
    slide_h = int(prs.slide_height or Emu(6858000))
    total = len(prs.slides)
    print(f"[2d] {total} slides, size={slide_w}x{slide_h} EMU", flush=True)

    all_slots: list[dict] = []
    for i in range(total):
        if i % 100 == 0:
            print(f"  [{i}/{total}]", flush=True)
        all_slots.extend(extract_slots_for_slide(prs.slides[i], i, slide_w, slide_h))

    # 집계
    text_slots = [s for s in all_slots if s["has_text_frame"]]
    ph_slots = [s for s in all_slots if s["has_placeholder"]]
    import numpy as np
    mc = np.array([s["max_chars"] for s in text_slots])
    summary = {
        "total_slots": len(all_slots),
        "text_slots": len(text_slots),
        "placeholder_slots": len(ph_slots),
        "max_chars_mean": float(mc.mean()) if len(mc) else 0,
        "max_chars_median": float(np.median(mc)) if len(mc) else 0,
        "max_chars_p95": float(np.percentile(mc, 95)) if len(mc) else 0,
        "total_placeholder_count": sum(s["placeholder_count"] for s in all_slots),
    }
    print(f"[2d] summary: {summary}", flush=True)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(all_slots, f, ensure_ascii=False, indent=2)
    with open(output_path.parent / "slot_schemas_summary.json", "w", encoding="utf-8") as f:
        json.dump(summary, f, ensure_ascii=False, indent=2)
    print(f"[2d] saved {len(all_slots)} slots -> {output_path}", flush=True)
    return summary


if __name__ == "__main__":
    ROOT = Path(__file__).resolve().parent.parent.parent
    pptx = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
    out = ROOT / "output" / "catalog" / "slot_schemas.json"
    run_stage_2d(pptx, out)
