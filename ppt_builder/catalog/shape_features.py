"""Track 1 Stage 2a — OOXML shape-based structural feature 추출.

텍스트 의존 0. 1,251장 마스터 템플릿이 placeholder(`~~`) 상태에서도
layout archetype별 분리가 가능하도록 설계.

Feature 구성 (총 68-dim)
------------------------
  F1  grid6x6 occupancy binary       36-dim  (shape center point)
  F2  shape_type histogram (norm L2)  12-dim  (Auto/Line/Text/Group/Pic/Table/Chart/Ph/Freeform/기타)
  F3  placeholder_grid (3x4)          12-dim  (`~~` shape bbox 점유)
  F4  structure hash signature         8-dim  (MD5(structure_sig) 앞 8bytes → float)

데이터 근거
-----------
scripts/deep_layout_analysis.py 실험 결과:
  grid6x6 unique 81.8%, Jaccard@0.7 ≈ 56 클러스터 (sweet spot 50-200)
  shape_type hist 단독은 degenerate (0.76 mean sim) — 보조용으로만 사용
  structure_sig 95.7% unique (1,197/1,251)
"""
from __future__ import annotations

import hashlib
import json
import re
from pathlib import Path

import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.picture import Picture
from pptx.util import Emu

from ..template.edit_ops import iter_leaf_shapes
from .schemas import SlideMeta


# --- Shape type 분류 (12 bucket) ---------------------------------------------
SHAPE_TYPES = [
    "auto_shape", "line", "text_box", "group", "picture",
    "table", "chart", "placeholder", "freeform", "connector",
    "ole", "other",
]
_TYPE_INDEX = {t: i for i, t in enumerate(SHAPE_TYPES)}


def _classify_shape(sh) -> str:
    """shape을 12 bucket 중 하나로 분류."""
    if isinstance(sh, Picture):
        return "picture"
    try:
        st = sh.shape_type
    except Exception:
        return "other"

    if getattr(sh, "has_chart", False):
        return "chart"
    if getattr(sh, "has_table", False):
        return "table"
    if getattr(sh, "is_placeholder", False):
        return "placeholder"

    # MSO_SHAPE_TYPE 매핑
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "auto_shape"
    if st == MSO_SHAPE_TYPE.LINE:
        return "line"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text_box"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if st == MSO_SHAPE_TYPE.FREEFORM:
        return "freeform"
    if hasattr(MSO_SHAPE_TYPE, "LINE_CONNECTOR") and st == MSO_SHAPE_TYPE.LINE_CONNECTOR:
        return "connector"
    if hasattr(MSO_SHAPE_TYPE, "EMBEDDED_OLE_OBJECT") and st == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        return "ole"
    return "other"


# --- Slide dimension 캐시 -----------------------------------------------------
def _slide_size(prs) -> tuple[int, int]:
    return int(prs.slide_width or Emu(9144000)), int(prs.slide_height or Emu(6858000))


# --- F1: grid6x6 occupancy ---------------------------------------------------
def _grid_occupancy(shapes_data: list[dict], w: int, h: int, gx: int = 6, gy: int = 6) -> np.ndarray:
    """shape center point를 gx×gy 그리드에 binary occupancy."""
    grid = np.zeros((gy, gx), dtype=np.float32)
    if w <= 0 or h <= 0:
        return grid.flatten()
    for s in shapes_data:
        cx = s["left"] + s["width"] / 2
        cy = s["top"] + s["height"] / 2
        ix = min(int(cx / w * gx), gx - 1)
        iy = min(int(cy / h * gy), gy - 1)
        ix = max(0, ix)
        iy = max(0, iy)
        grid[iy, ix] = 1.0
    return grid.flatten()


# --- F2: shape_type histogram -------------------------------------------------
def _type_hist(shapes_data: list[dict]) -> np.ndarray:
    h = np.zeros(len(SHAPE_TYPES), dtype=np.float32)
    for s in shapes_data:
        h[_TYPE_INDEX.get(s["type"], _TYPE_INDEX["other"])] += 1
    # L2 normalize (degenerate 방지)
    n = np.linalg.norm(h)
    if n > 0:
        h = h / n
    return h


# --- F3: placeholder grid (3x4) -----------------------------------------------
def _placeholder_grid(
    shapes_data: list[dict],
    w: int, h: int,
    gx: int = 4, gy: int = 3,
) -> np.ndarray:
    """`~~` 텍스트 포함 shape의 bbox 영역 점유율 (3x4 grid)."""
    grid = np.zeros((gy, gx), dtype=np.float32)
    if w <= 0 or h <= 0:
        return grid.flatten()
    for s in shapes_data:
        if not s.get("has_placeholder_text"):
            continue
        l, t, ww, hh = s["left"], s["top"], s["width"], s["height"]
        if ww <= 0 or hh <= 0:
            continue
        # bbox 를 grid에 soft-project (면적 비례)
        x0, y0 = l / w * gx, t / h * gy
        x1, y1 = (l + ww) / w * gx, (t + hh) / h * gy
        ix0, iy0 = int(max(0, x0)), int(max(0, y0))
        ix1, iy1 = min(gx - 1, int(x1)), min(gy - 1, int(y1))
        for iy in range(iy0, iy1 + 1):
            for ix in range(ix0, ix1 + 1):
                grid[iy, ix] += 1.0
    # 정규화
    mx = grid.max()
    if mx > 0:
        grid = grid / mx
    return grid.flatten()


# --- F4: structure_sig hash ---------------------------------------------------
def _structure_sig_hash(sig: str, dim: int = 8) -> np.ndarray:
    """structure_sig 문자열을 MD5 → uniform [0,1] 8-dim."""
    if not sig:
        return np.zeros(dim, dtype=np.float32)
    h = hashlib.md5(sig.encode("utf-8")).digest()
    arr = np.frombuffer(h[:dim], dtype=np.uint8).astype(np.float32) / 255.0
    return arr


# --- Shape data 추출 (공통) ---------------------------------------------------
_PLACEHOLDER_RE = re.compile(r"~~")


def _extract_shapes_data(slide) -> list[dict]:
    """각 shape의 type + bbox + placeholder 존재 여부를 평탄화 순회로 추출."""
    data: list[dict] = []
    for _, sh in iter_leaf_shapes(slide):
        try:
            left = int(sh.left or 0)
            top = int(sh.top or 0)
            width = int(sh.width or 0)
            height = int(sh.height or 0)
        except Exception:
            left = top = width = height = 0

        has_ph = False
        if sh.has_text_frame:
            try:
                txt = sh.text_frame.text or ""
                has_ph = bool(_PLACEHOLDER_RE.search(txt))
            except Exception:
                pass

        data.append({
            "type": _classify_shape(sh),
            "left": left, "top": top, "width": width, "height": height,
            "has_placeholder_text": has_ph,
        })
    return data


# --- 메인: 슬라이드 → 68-dim feature vector ----------------------------------
def extract_shape_feature(slide, meta: SlideMeta, slide_w: int, slide_h: int) -> np.ndarray:
    """단일 슬라이드 → 68-dim structural feature."""
    shapes_data = _extract_shapes_data(slide)
    f1 = _grid_occupancy(shapes_data, slide_w, slide_h, gx=6, gy=6)       # 36
    f2 = _type_hist(shapes_data)                                           # 12
    f3 = _placeholder_grid(shapes_data, slide_w, slide_h, gx=4, gy=3)      # 12
    f4 = _structure_sig_hash(meta.structure_sig, dim=8)                    # 8
    feat = np.concatenate([f1, f2, f3, f4]).astype(np.float32)
    # 각 블록 가중치 (F2는 degenerate 방지로 0.5)
    weights = np.concatenate([
        np.full(36, 1.0, dtype=np.float32),  # F1 grid
        np.full(12, 0.5, dtype=np.float32),  # F2 type hist (약화)
        np.full(12, 0.7, dtype=np.float32),  # F3 placeholder grid
        np.full(8, 0.3, dtype=np.float32),   # F4 sig hash (보조)
    ])
    feat = feat * weights
    # L2 normalize
    n = np.linalg.norm(feat)
    if n > 0:
        feat = feat / n
    return feat


def extract_all_features(
    pptx_path: Path,
    meta_path: Path,
    output_path: Path,
) -> np.ndarray:
    """pptx 전체 슬라이드 → (N, 68) feature matrix 저장."""
    with open(meta_path, "r", encoding="utf-8") as f:
        metas = [SlideMeta.model_validate(x) for x in json.load(f)]

    prs = Presentation(str(pptx_path))
    slide_w, slide_h = _slide_size(prs)
    total = len(prs.slides)
    print(f"[feat] {total} slides, slide_size={slide_w}x{slide_h}", flush=True)

    feats = np.zeros((total, 68), dtype=np.float32)
    for i in range(total):
        if i % 100 == 0:
            print(f"  [{i}/{total}]", flush=True)
        feats[i] = extract_shape_feature(prs.slides[i], metas[i], slide_w, slide_h)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    np.save(output_path, feats)
    print(f"[feat] saved shape={feats.shape} -> {output_path}", flush=True)
    return feats


if __name__ == "__main__":
    ROOT = Path(__file__).resolve().parent.parent.parent
    pptx = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
    meta = ROOT / "output" / "catalog" / "slide_meta.json"
    out = ROOT / "output" / "catalog" / "shape_features.npy"
    extract_all_features(pptx, meta, out)
