"""PPT 네이티브 차트 — python-pptx add_chart() 기반.

matplotlib PNG가 아닌 PPT 내장 차트로, 생성 후 PowerPoint에서 편집 가능.
Region-aware: (slide, region) → 차트 shape 반환.
"""

from __future__ import annotations

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu

from ppt_builder.primitives import Region


# 색상 상수
_ACCENT = RGBColor(0xFD, 0x51, 0x08)     # PwC 오렌지
_ACCENT_MID = RGBColor(0xFE, 0x7C, 0x39)
_GREY_800 = RGBColor(0x4A, 0x4F, 0x58)
_GREY_400 = RGBColor(0x9A, 0xA0, 0xA8)
_GREY_200 = RGBColor(0xE2, 0xE5, 0xE8)
_POSITIVE = RGBColor(0x27, 0xAE, 0x60)
_NEGATIVE = RGBColor(0xC0, 0x39, 0x2B)


def _style_chart(chart, hide_legend=False):
    """공통 차트 스타일링 — 미니멀 컨설팅 스타일."""
    chart.has_title = False
    if hide_legend:
        chart.has_legend = False
    else:
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(8)

    # 축 스타일링
    if hasattr(chart, 'category_axis'):
        cat_ax = chart.category_axis
        cat_ax.has_major_gridlines = False
        cat_ax.major_tick_mark = 2  # OUTSIDE -> NONE
        cat_ax.tick_labels.font.size = Pt(8)
        cat_ax.format.line.fill.background()

    if hasattr(chart, 'value_axis'):
        val_ax = chart.value_axis
        val_ax.has_major_gridlines = True
        val_ax.major_gridlines.format.line.color.rgb = _GREY_200
        val_ax.major_gridlines.format.line.width = Pt(0.5)
        val_ax.major_tick_mark = 2
        val_ax.tick_labels.font.size = Pt(8)
        val_ax.format.line.fill.background()


# ============================================================
# 1. Vertical Bar Chart (세로 막대)
# ============================================================

def chart_vertical_bar(
    slide: Slide,
    *,
    categories: list[str],
    values: list[float],
    highlight_idx: int = -1,
    region: Region,
    series_name: str = "",
) -> object:
    """세로 막대 차트 — 최대값 자동 강조."""
    data = CategoryChartData()
    data.categories = categories
    data.add_series(series_name, values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(region.x), Inches(region.y),
        Inches(region.w), Inches(region.h),
        data,
    )
    chart = chart_shape.chart
    _style_chart(chart, hide_legend=True)

    # 바 색상: highlight만 오렌지, 나머지 회색
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = _GREY_400

    if highlight_idx >= 0 and highlight_idx < len(values):
        pt = series.points[highlight_idx]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = _ACCENT

    # 데이터 라벨
    series.has_data_labels = True
    series.data_labels.font.size = Pt(9)
    series.data_labels.font.bold = True
    series.data_labels.number_format = '#,##0'

    return chart_shape


# ============================================================
# 2. Line Chart (꺾은선)
# ============================================================

def chart_line(
    slide: Slide,
    *,
    categories: list[str],
    series_data: list[dict],  # [{"name": "...", "values": [...]}]
    region: Region,
) -> object:
    """꺾은선 차트 — 다중 시리즈."""
    data = CategoryChartData()
    data.categories = categories
    for sd in series_data:
        data.add_series(sd["name"], sd["values"])

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(region.x), Inches(region.y),
        Inches(region.w), Inches(region.h),
        data,
    )
    chart = chart_shape.chart
    _style_chart(chart)

    colors = [_ACCENT, _GREY_800, _GREY_400, _ACCENT_MID]
    for i, series in enumerate(chart.series):
        series.format.line.color.rgb = colors[i % len(colors)]
        series.format.line.width = Pt(2.0)
        series.smooth = False

    return chart_shape


# ============================================================
# 3. Pie / Donut Chart
# ============================================================

def chart_donut(
    slide: Slide,
    *,
    categories: list[str],
    values: list[float],
    region: Region,
    center_text: str = "",
) -> object:
    """도넛 차트 — 중앙 텍스트 지원."""
    data = CategoryChartData()
    data.categories = categories
    data.add_series("", values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(region.x), Inches(region.y),
        Inches(region.w), Inches(region.h),
        data,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(8)

    # 색상
    colors = [_ACCENT, _GREY_800, _GREY_400, _ACCENT_MID, _GREY_200]
    plot = chart.plots[0]
    series = plot.series[0]
    for i in range(len(values)):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]

    # 데이터 라벨
    series.has_data_labels = True
    series.data_labels.font.size = Pt(9)
    series.data_labels.number_format = '0%'
    series.data_labels.show_percentage = True
    series.data_labels.show_value = False

    return chart_shape


# ============================================================
# 4. Stacked Bar Chart (누적 가로 막대)
# ============================================================

def chart_stacked_bar(
    slide: Slide,
    *,
    categories: list[str],
    series_data: list[dict],  # [{"name": "...", "values": [...]}]
    region: Region,
) -> object:
    """누적 가로 막대 차트 — 구성 비율 비교."""
    data = CategoryChartData()
    data.categories = categories
    for sd in series_data:
        data.add_series(sd["name"], sd["values"])

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED_100,
        Inches(region.x), Inches(region.y),
        Inches(region.w), Inches(region.h),
        data,
    )
    chart = chart_shape.chart
    _style_chart(chart)

    colors = [_ACCENT, _GREY_800, _GREY_400, _ACCENT_MID, _GREY_200]
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[i % len(colors)]

    return chart_shape


# ============================================================
# 5. Scatter / Bubble Chart
# ============================================================

def chart_scatter(
    slide: Slide,
    *,
    points: list[dict],  # [{"label": "...", "x": float, "y": float, "size": float}]
    x_label: str = "",
    y_label: str = "",
    region: Region,
) -> object:
    """산점도/버블 차트 — 포지셔닝 맵."""
    data = XyChartData()

    for pt in points:
        series = data.add_series(pt["label"])
        series.add_data_point(pt["x"], pt["y"])

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(region.x), Inches(region.y),
        Inches(region.w), Inches(region.h),
        data,
    )
    chart = chart_shape.chart
    _style_chart(chart)

    colors = [_ACCENT, _GREY_800, _GREY_400, _ACCENT_MID, _NEGATIVE, _POSITIVE]
    for i, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[i % len(colors)]
        marker = series.marker
        marker.size = 12
        marker.format.fill.solid()
        marker.format.fill.fore_color.rgb = colors[i % len(colors)]

    return chart_shape
