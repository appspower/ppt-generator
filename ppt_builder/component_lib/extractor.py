"""ComponentExtractor — 외부 .pptx에서 컴퍼넌트(GroupShape 등)를 추출한다.

Phase 1: 명시적 GroupShape만 추출.
Phase 2 (TODO): 자동 클러스터링으로 ungrouped shape도 추출.

추출 단위는 Component dataclass. 다음을 포함:
- shape XML (lxml deepcopy)
- 이미지 blob 데이터 (cross-presentation 복사를 위해)
- 메타데이터 (bbox, color palette, slot 후보 등)
"""

from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from .models import Component, Slot


# Namespace 헬퍼
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

A = lambda tag: f"{{{A_NS}}}{tag}"
P = lambda tag: f"{{{P_NS}}}{tag}"
R = lambda tag: f"{{{R_NS}}}{tag}"


class ComponentExtractor:
    """외부 .pptx에서 GroupShape를 추출하여 Component 객체로 변환한다."""

    def extract_groups(self, pptx_path: str | Path) -> list[Component]:
        """파일 내 모든 슬라이드에서 GroupShape를 추출한다.

        Args:
            pptx_path: 외부 .pptx 경로

        Returns:
            Component 리스트
        """
        path = Path(pptx_path)
        if not path.exists():
            raise FileNotFoundError(f"Template not found: {path}")

        prs = Presentation(str(path))
        components: list[Component] = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    component = self._extract_group(
                        group_shape=shape,
                        slide=slide,
                        source_file=path.name,
                        slide_index=slide_idx,
                    )
                    if component:
                        components.append(component)

        return components

    def _extract_group(
        self,
        group_shape,
        slide,
        source_file: str,
        slide_index: int,
    ) -> Component | None:
        """단일 GroupShape를 Component로 변환한다."""
        try:
            grp_elm = group_shape._element  # <p:grpSp>
        except AttributeError:
            return None

        # 1. shape XML deep copy
        sp_xml_copy = deepcopy(grp_elm)

        # 2. 이미지 blob 수집 (cross-presentation 복사용)
        image_blobs, image_content_types = self._collect_image_blobs(
            sp_xml_copy, slide
        )

        # 3. 메타데이터 분석
        text_density, fonts, has_chart, has_smartart, has_table, has_images, child_count = \
            self._analyze_xml(sp_xml_copy)
        has_images = has_images or len(image_blobs) > 0

        # 4. 색상 팔레트 추출
        color_palette = self._extract_colors(sp_xml_copy)

        # 5. 슬롯 후보 자동 추출 (단순 휴리스틱: 모든 텍스트 run을 슬롯으로)
        slots = self._extract_slot_candidates(sp_xml_copy)

        # 6. 카테고리 자동 분류 (간단 휴리스틱)
        category, subcategory = self._classify(slots, child_count, has_table)

        # 7. ID 생성
        component_id = self._generate_id(source_file, slide_index, group_shape.name)

        # 8. bbox (원본)
        bbox = (
            int(group_shape.left or 0),
            int(group_shape.top or 0),
            int(group_shape.width or 0),
            int(group_shape.height or 0),
        )

        return Component(
            id=component_id,
            category=category,
            subcategory=subcategory,
            name=group_shape.name or "Unnamed",
            source_file=source_file,
            source_slide_index=slide_index,
            bbox_emu=bbox,
            sp_xml_bytes=etree.tostring(sp_xml_copy),
            image_blobs=image_blobs,
            image_content_types=image_content_types,
            slots=slots,
            color_palette=color_palette,
            font_families=list(set(fonts)),
            text_density=text_density,
            shape_count=child_count,
            has_images=has_images,
            has_charts=has_chart,
            has_smartart=has_smartart,
            has_table=has_table,
        )

    def _collect_image_blobs(
        self, grp_elm, slide
    ) -> tuple[dict[str, bytes], dict[str, str]]:
        """그룹 내부의 모든 image rId를 수집한다.

        모던 PPTX는 한 이미지에 두 개의 rId를 가질 수 있다:
        - <a:blip r:embed="rId3"/> (PNG fallback)
        - <asvg:svgBlip r:embed="rId4"/> (SVG vector)
        둘 다 수집해야 PowerPoint가 빨간 X를 표시하지 않는다.
        """
        blobs: dict[str, bytes] = {}
        content_types: dict[str, str] = {}

        slide_part = slide.part
        slide_rels = slide_part.rels

        embed_attr_q = R("embed")
        link_attr_q = R("link")

        # 모든 element에서 r:embed/r:link 속성 수집
        used_rids: set[str] = set()
        for el in grp_elm.iter():
            rid = el.get(embed_attr_q) or el.get(link_attr_q)
            if rid:
                used_rids.add(rid)

        for rid in used_rids:
            rel = slide_rels.get(rid)
            if rel is None:
                continue
            if "image" not in rel.reltype:
                continue
            if rel.is_external:
                # 외부 링크는 Phase 1 범위 밖
                continue
            try:
                image_part = rel.target_part
                blobs[rid] = image_part.blob
                content_types[rid] = image_part.content_type
            except Exception:
                pass

        return blobs, content_types

    def _analyze_xml(self, grp_elm) -> tuple[int, list[str], bool, bool, bool, bool, int]:
        """XML 트리를 순회하며 메타데이터를 추출한다."""
        text_density = 0
        fonts: list[str] = []
        has_chart = False
        has_smartart = False
        has_table = False
        has_images = False
        child_count = 0

        # 자식 shape 개수 (1단계 자식만)
        for child in grp_elm:
            tag = etree.QName(child.tag).localname
            if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
                child_count += 1

        # 모든 텍스트 추출
        for t in grp_elm.iter(A("t")):
            if t.text:
                text_density += len(t.text)

        # 폰트 추출
        for rfont in grp_elm.iter(A("latin")):
            typeface = rfont.get("typeface")
            if typeface:
                fonts.append(typeface)
        for rfont in grp_elm.iter(A("ea")):
            typeface = rfont.get("typeface")
            if typeface:
                fonts.append(typeface)

        # 차트 감지
        for graphic_data in grp_elm.iter(A("graphicData")):
            uri = graphic_data.get("uri", "")
            if "chart" in uri:
                has_chart = True
            if "diagram" in uri or "smartart" in uri.lower():
                has_smartart = True
            if "table" in uri:
                has_table = True

        # 이미지 (pic 태그)
        for _ in grp_elm.iter(P("pic")):
            has_images = True
            break

        return text_density, fonts, has_chart, has_smartart, has_table, has_images, child_count

    def _extract_colors(self, grp_elm) -> list[str]:
        """sRGB 색상 코드를 모두 추출한다 (중복 제거)."""
        colors: set[str] = set()
        for clr in grp_elm.iter(A("srgbClr")):
            val = clr.get("val")
            if val:
                colors.add(f"#{val.upper()}")
        return sorted(colors)

    def _extract_slot_candidates(self, grp_elm) -> list[Slot]:
        """텍스트가 들어있는 모든 paragraph를 슬롯 후보로 추출한다.

        Phase 1에서는 단순히 paragraph 단위로 slot을 만든다.
        slot_id는 자동 생성 (slot_0, slot_1, ...).
        Phase 3에서 의미 분석으로 고도화 예정.
        """
        slots: list[Slot] = []
        slot_idx = 0

        # <p:sp> 단위로 순회
        for sp in grp_elm.iter(P("sp")):
            tx_body = sp.find(P("txBody"))
            if tx_body is None:
                continue

            for p in tx_body.iter(A("p")):
                # paragraph의 텍스트 합치기
                texts = []
                font_size = 0
                bold = False
                for r in p.iter(A("r")):
                    t = r.find(A("t"))
                    if t is not None and t.text:
                        texts.append(t.text)
                    rpr = r.find(A("rPr"))
                    if rpr is not None:
                        sz = rpr.get("sz")
                        if sz:
                            font_size = max(font_size, int(sz) / 100)
                        if rpr.get("b") == "1":
                            bold = True

                full_text = "".join(texts).strip()
                if not full_text:
                    continue

                slot = Slot(
                    slot_id=f"slot_{slot_idx}",
                    semantic_role="text",
                    original_text=full_text,
                    font_size_pt=font_size,
                    font_bold=bold,
                    max_chars=max(len(full_text) * 2, 50),
                )
                slots.append(slot)
                slot_idx += 1

        return slots

    def _classify(
        self, slots: list[Slot], child_count: int, has_table: bool
    ) -> tuple[str, str]:
        """카테고리 자동 분류 (Phase 1: 단순 휴리스틱).

        Phase 3에서 의미 분석으로 고도화 예정.
        """
        # 텍스트 키워드 기반
        all_text = " ".join(s.original_text.lower() for s in slots)

        if any(k in all_text for k in ["strength", "weakness", "swot"]):
            return "framework", "swot"
        if any(k in all_text for k in ["bcg", "growth-share", "stars", "cash cow"]):
            return "framework", "bcg_matrix"
        if any(k in all_text for k in ["porter", "five forces", "rivalry"]):
            return "framework", "porter"
        if any(k in all_text for k in ["pestel", "political", "economic"]):
            return "framework", "pestel"
        if any(k in all_text for k in ["value chain", "primary activities"]):
            return "framework", "value_chain"
        if any(k in all_text for k in ["timeline", "roadmap", "gantt"]):
            return "timeline", "generic"
        if has_table:
            return "table", "generic"

        # shape 개수로 추정
        if child_count >= 6:
            return "framework", "generic"
        if child_count >= 3:
            return "callout", "multi"
        return "callout", "generic"

    def _generate_id(self, source_file: str, slide_index: int, name: str) -> str:
        """안정적인 컴퍼넌트 ID를 생성한다."""
        base = source_file.replace(".pptx", "")
        # 영문/숫자/언더스코어만 남김
        clean_name = re.sub(r"[^a-zA-Z0-9]+", "_", (name or "group")).strip("_").lower()
        return f"{base}_s{slide_index}_{clean_name}"
