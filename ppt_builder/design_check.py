"""Phase B — 디자인 결정 체크리스트.

목적: Claude가 매번 같은 디테일 수준으로 슬라이드를 만들도록 강제한다.
패턴 라이브러리만으로는 "01 번호를 원형으로 할지 사각으로 할지" 같은
결정이 들쭉날쭉해진다. 이 모듈이 그런 결정을 명시화하고,
빌드 후 디자인 품질을 점검한다.

두 종류의 검사:
1. PRE-BUILD (decide_*) — 콘텐츠 spec을 받아 디자인 결정을 추천
   - "이 spec에 어울리는 도형은?" 같은 질문에 명확한 답을 내림
   - Claude는 이걸 호출해서 자기가 빌드할 때 결정 근거를 가짐
2. POST-BUILD (inspect_*) — 빌드된 캔버스를 분석해 문제 검출
   - validate_visual()이 잡지 못하는 디자인 품질 문제
   - 시각 위계 부족, 영역별 빈 공간, 색상 균형, 정보 밀도

기존 검증과의 분담:
- evaluate.py: 스코어 (텍스트 길이, 폰트 크기 등)
- visual_validate.py: 시각 잘림 / 오버플로 / 겹침
- design_check.py: 디자인 품질 (위계, 균형, 결정 일관성)
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Literal, Optional

from pptx import Presentation


# ============================================================
# Public dataclasses
# ============================================================


@dataclass
class DesignDecision:
    """디자인 결정 결과 — Claude가 빌드 전에 참조."""
    rationale: str
    recommendation: dict
    alternatives: list[dict] = field(default_factory=list)


@dataclass
class DesignIssue:
    severity: Literal["high", "medium", "low"]
    category: str  # hierarchy / balance / density / decision / contrast
    message: str
    suggestion: str = ""

    def __str__(self) -> str:
        return f"[{self.severity.upper()}/{self.category}] {self.message}"


@dataclass
class DesignReport:
    issues: list[DesignIssue] = field(default_factory=list)
    metrics: dict = field(default_factory=dict)

    @property
    def passed(self) -> bool:
        return not any(i.severity == "high" for i in self.issues)

    def summary(self) -> str:
        lines = [f"Design issues: {len(self.issues)}"]
        for i in self.issues[:10]:
            lines.append(f"  - {i}")
            if i.suggestion:
                lines.append(f"    → {i.suggestion}")
        return "\n".join(lines)


# ============================================================
# PRE-BUILD: 콘텐츠 → 디자인 결정 추천
# ============================================================


def decide_number_marker(
    *,
    item_count: int,
    has_sequence: bool,
    space_h: float,
) -> DesignDecision:
    """01/02/03 같은 번호 시퀀스에 어떤 도형을 쓸지 결정.

    Args:
        item_count: 항목 개수
        has_sequence: 순서가 의미를 갖는가? (시간/단계/우선순위)
        space_h: 항목당 사용 가능한 세로 공간 (인치)

    Returns:
        DesignDecision with shape: 'circle' | 'square' | 'chevron'
    """
    # 순서가 명확한 단계 (4개 이하, 가로로 펼쳐짐) → chevron
    if has_sequence and item_count <= 5 and space_h < 0.7:
        return DesignDecision(
            rationale=(
                f"{item_count}개 단계 항목, 단계 의미 명확 → "
                "chevron 화살표가 시퀀스를 시각적으로 가장 잘 표현"
            ),
            recommendation={
                "shape": "chevron",
                "primitive": "Canvas.chevron",
                "size_hint": "h=0.4~0.5, w=균등 분배",
            },
            alternatives=[
                {"shape": "arrow_chain", "reason": "단계가 6개 이상이면"},
            ],
        )

    # 세로 리스트, 항목별 디테일이 있음 → circle
    if space_h >= 0.5:
        return DesignDecision(
            rationale=(
                f"{item_count}개 항목, 항목당 {space_h:.1f}\" 공간 — "
                "원형 번호가 텍스트 옆에 작게 들어가 깔끔한 위계 형성"
            ),
            recommendation={
                "shape": "circle",
                "primitive": "Canvas.circle or Canvas.numbered_list",
                "size_hint": "d=0.32~0.40",
            },
            alternatives=[
                {"shape": "square", "reason": "더 격식 있는 톤이 필요하면"},
            ],
        )

    # 작은 공간 + 다수 항목 → 사각 박스
    return DesignDecision(
        rationale=(
            f"{item_count}개 항목, 좁은 공간 ({space_h:.1f}\") — "
            "사각 박스가 공간 효율 최고"
        ),
        recommendation={
            "shape": "square",
            "primitive": "Canvas.box with text",
            "size_hint": "h=0.3~0.4 정사각",
        },
        alternatives=[
            {"shape": "circle", "reason": "디자인 다양성을 위해"},
        ],
    )


def decide_emphasis_color(
    *,
    palette: Literal["grey", "grey_orange", "monochrome"] = "grey",
    needs_strong_emphasis: bool = False,
) -> DesignDecision:
    """강조 색상 결정. 사용자 정책: 회색 톤 우선, 진한 오렌지 회피."""
    if palette == "grey":
        primary = "grey_900" if needs_strong_emphasis else "grey_800"
        return DesignDecision(
            rationale="회색 위계 정책 — 강조는 다크 그레이로",
            recommendation={
                "primary": primary,
                "secondary": "grey_700",
                "tertiary": "grey_400",
                "background": "grey_100",
            },
        )
    if palette == "monochrome":
        return DesignDecision(
            rationale="순수 모노크롬",
            recommendation={
                "primary": "black",
                "secondary": "grey_700",
                "tertiary": "grey_400",
                "background": "white",
            },
        )
    # grey_orange — 오렌지 사용 허용 (현재 사용자는 비선호)
    return DesignDecision(
        rationale="회색+오렌지 (사용자 정책상 비선호 — 명시 요청 시만)",
        recommendation={
            "primary": "accent",
            "secondary": "grey_800",
            "tertiary": "grey_400",
            "background": "white",
        },
    )


def decide_density(
    *,
    available_area: float,  # sq-in
    intended_chars: int,
) -> DesignDecision:
    """공간 대비 텍스트 밀도가 적절한지 판단.

    가이드:
    - 너무 빈약: 100자 / sq-in 미만 → 콘텐츠 추가 필요
    - 적절: 100~300자 / sq-in
    - 과밀: 300자 / sq-in 초과 → 폰트 줄이거나 박스 키우기
    """
    if available_area <= 0:
        return DesignDecision(
            rationale="영역 없음",
            recommendation={"action": "n/a"},
        )
    density = intended_chars / available_area
    if density < 100:
        return DesignDecision(
            rationale=(
                f"밀도 {density:.0f} chars/sq-in — 너무 빈약. "
                "sub-bullet, mini-list, 또는 보조 정보로 채워야 함"
            ),
            recommendation={
                "action": "enrich",
                "suggestions": [
                    "sub-bullet 2~3개 추가",
                    "보조 KPI 또는 stat_block 추가",
                    "mini-callout으로 디테일 보강",
                ],
            },
        )
    if density > 300:
        return DesignDecision(
            rationale=(
                f"밀도 {density:.0f} chars/sq-in — 과밀. "
                "텍스트 압축 또는 영역 확장 필요"
            ),
            recommendation={
                "action": "reduce_or_expand",
                "suggestions": [
                    "텍스트 압축 (불필요 단어 제거)",
                    "영역 확장",
                    "두 영역으로 분할",
                ],
            },
        )
    return DesignDecision(
        rationale=f"밀도 {density:.0f} chars/sq-in — 적절",
        recommendation={"action": "ok"},
    )


# 한국어 컨설팅 단호 어미 화이트리스트 — 레퍼런스 분석 기반
KOREAN_DEFINITIVE_ENDINGS = (
    "임", "함", "됨", "음",
    "구축", "확보", "지원", "제공", "달성", "수립", "도출", "정의",
    "강화", "개선", "절감", "단축", "확대", "수행", "적용", "고도화",
    "활용", "이행", "추진", "운영", "실시", "마련", "체계화",
    "필요", "권장", "예상", "기대",
    "있음", "없음",
    "발견", "분류", "검증", "확인", "분석", "차별화", "우위",
    "압축", "자동화", "통합", "정합화", "최적화", "표준화",
)


def check_head_message_ending(text: str) -> Optional[DesignIssue]:
    """head message 텍스트가 컨설팅 톤 한국어 어미로 끝나는지 검사.

    한국어 텍스트가 아니거나 검증을 건너뛸 경우 None 반환.
    검사 통과: None. 미통과: LOW severity DesignIssue.
    """
    if not text or not text.strip():
        return None
    text = text.strip().rstrip(".")
    # 한국어 비율로 한국어 텍스트인지 판단
    kr_chars = sum(1 for ch in text if 0xAC00 <= ord(ch) <= 0xD7A3)
    if kr_chars / max(len(text), 1) < 0.3:
        # 영문/숫자 위주 → 검사 skip
        return None
    # 어미 매칭
    for ending in KOREAN_DEFINITIVE_ENDINGS:
        if text.endswith(ending):
            return None
    return DesignIssue(
        severity="low",
        category="tone",
        message=f"Head message가 컨설팅 단호 어미로 끝나지 않음: \"...{text[-10:]}\"",
        suggestion="\"...임\", \"...구축\", \"...지원\" 등 단호한 한국어 어미 권장",
    )


def decide_head_message_form(
    *,
    intent: Literal[
        "executive", "timeline", "comparison", "process", "quadrant", "data"
    ],
) -> DesignDecision:
    """슬라이드 의도별로 head message 어미 추천.

    레퍼런스 분석 기반: PwC 컨설팅 산출물의 head message 어미 패턴.
    """
    forms = {
        "executive": {
            "preferred_endings": ["임", "함", "구축", "확보"],
            "tone": "결론 단정, 임팩트 숫자 포함",
            "example": "Palantir 투입으로 SAP 전환 일정 14% 단축, 테스트 70%↓, DT 50%↓",
            "anti_pattern": "물음표나 슬로건 형태 (예: '어떻게 할 것인가?')",
        },
        "timeline": {
            "preferred_endings": ["수립", "체계화", "이행", "추진"],
            "tone": "단계적 실행 의지",
            "example": "3단계 Quick Win 전략으로 2~3주 내 가치 입증 후 점진 확대",
            "anti_pattern": "단계 나열만 (예: 'L1, L2, L3, L4')",
        },
        "comparison": {
            "preferred_endings": ["임", "차별화", "우위", "확보"],
            "tone": "비교 결론 단호",
            "example": "Palantir의 차별점은 Ontology+AIP+Workshop 단일 플랫폼 통합임",
            "anti_pattern": "모호한 평가 (예: '각자 장단점이 있음')",
        },
        "process": {
            "preferred_endings": ["수행", "구축", "자동화", "압축"],
            "tone": "프로세스 결과 명시",
            "example": "Blueprint→결함 등록 5단계 자동 파이프라인으로 공수 70% 절감",
            "anti_pattern": "단계 명사형만 (예: '5단계 프로세스')",
        },
        "quadrant": {
            "preferred_endings": ["분류", "우선순위화", "권장", "도출"],
            "tone": "사분면 → 의사결정 권고",
            "example": "8개 모듈을 ROI×난이도로 분류 — Quick Win 영역에서 시작 권장",
            "anti_pattern": "사분면 라벨만 (예: '4개 영역의 모듈')",
        },
        "data": {
            "preferred_endings": ["증명", "확인", "도출", "발견"],
            "tone": "데이터 → 인사이트",
            "example": "월별 결함 추이가 6개월 후 안정화되는 패턴을 확인",
            "anti_pattern": "차트 설명만 (예: '월별 결함 그래프')",
        },
    }
    rec = forms[intent]
    return DesignDecision(
        rationale=f"{intent} 슬라이드 head message 어미 권장: {', '.join(rec['preferred_endings'])}",
        recommendation=rec,
    )


def decide_layout_archetype(
    *,
    intent: Literal[
        "executive", "timeline", "comparison", "process", "quadrant", "data"
    ],
    item_count: int,
) -> DesignDecision:
    """슬라이드 의도(intent)와 항목 개수에 따라 적합한 레이아웃 패턴 추천."""
    rules = {
        "executive": {
            "pattern": "hero_with_metrics",
            "layout": "left hero (35-45%) + right metric grid + bottom takeaway",
            "rationale": "C-level은 결론+근거 숫자가 즉시 보여야 함",
        },
        "timeline": {
            "pattern": "horizontal_phases_with_deliverables",
            "layout": "top chevron sequence + bottom per-phase deliverable lists",
            "rationale": "시간 흐름은 가로축이 자연스럽고, 각 단계 산출물이 함께 보여야 신뢰",
        },
        "comparison": {
            "pattern": "side_by_side_or_matrix",
            "layout": (
                "2개: 좌우 callout / "
                "3-4개: 균등 columns / "
                "5+: framework matrix"
            ),
            "rationale": "비교는 동일 차원에서 나란히 보여야 함",
        },
        "process": {
            "pattern": "horizontal_arrow_chain_with_callouts",
            "layout": "arrow_chain 위 + 각 단계 옆 callout",
            "rationale": "프로세스는 흐름 방향이 명확해야 함",
        },
        "quadrant": {
            "pattern": "2x2_with_center_insight",
            "layout": "2x2 grid + 중앙 또는 하단 인사이트 박스",
            "rationale": "사분면은 두 축의 의미가 강조되어야 함",
        },
        "data": {
            "pattern": "chart_with_callouts",
            "layout": "큰 차트/표 + 옆 인사이트 stat_blocks",
            "rationale": "데이터는 숫자만 있으면 의미 전달 안 됨, 인사이트 동반",
        },
    }
    rec = rules[intent]
    return DesignDecision(
        rationale=rec["rationale"],
        recommendation=rec,
    )


# ============================================================
# POST-BUILD: 캔버스 분석으로 디자인 품질 점검
# ============================================================


def inspect_design(
    pptx_path: str,
    *,
    expect_dense: bool = True,
    forbidden_colors: Optional[list[str]] = None,
    pattern_kind: Optional[str] = None,
) -> DesignReport:
    """렌더된 .pptx를 디자인 관점에서 점검한다.

    검사 항목:
    1. 색상 균형 — 전체에서 어떤 색이 얼마나 쓰였는가
    2. 영역별 빈 공간 — 슬라이드를 4분면으로 나눠 콘텐츠 없는 영역 검출
    3. 시각 위계 — 폰트 크기 분포가 충분한 단계를 가지는가
    4. 정보 밀도 — 슬라이드 전체 밀도 (expect_dense=True일 때)
       pattern_kind로 패턴별 임계값 다르게 적용 가능
    5. 금지색 사용 — forbidden_colors에 들어 있는 색이 쓰였는가

    Args:
        pattern_kind: 'executive' | 'timeline' | 'comparison' | 'process' |
                      'quadrant' | 'data' — 패턴별 density 임계값 조정.
                      comparison/quadrant는 짧은 셀+여백이 본질이라 낮은
                      밀도가 정상.
    """
    # 패턴별 density 임계값
    # 기본: 12 (그 미만이면 너무 빈약)
    # comparison/quadrant: 6 (셀이 짧고 여백이 의미)
    # data: 8 (차트가 면적 대부분 차지)
    density_floors = {
        "executive": 12,
        "timeline": 12,
        "comparison": 6,
        "quadrant": 6,
        "process": 10,
        "data": 8,
    }
    density_floor = density_floors.get(pattern_kind or "", 12)
    issues: list[DesignIssue] = []
    metrics: dict = {}

    if forbidden_colors is None:
        # 사용자 정책: 진한 오렌지 회피
        forbidden_colors = ["FD5108", "FE7C39"]

    prs = Presentation(pptx_path)

    for si, slide in enumerate(prs.slides):
        sn = si + 1
        slide_w = prs.slide_width / 914400
        slide_h = prs.slide_height / 914400

        # ---------- 1. 색상 분포 ----------
        color_use: dict[str, float] = {}  # color → total area
        forbidden_hits: list[tuple[str, str]] = []  # (color, where)
        font_sizes: list[float] = []
        total_chars = 0

        # head message 후보 — 헤더 영역(y < 1.2")에서 가장 큰 폰트 텍스트
        head_candidates: list[tuple[float, str]] = []  # (font_size, text)

        # ---------- 2. 영역별 점유 ----------
        # 4분면 (TL, TR, BL, BR) 각 영역의 사용 면적
        quadrants = {"TL": 0.0, "TR": 0.0, "BL": 0.0, "BR": 0.0}
        slide_area = slide_w * slide_h

        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            l = shape.left / 914400
            t = shape.top / 914400
            w = (shape.width or 0) / 914400
            h = (shape.height or 0) / 914400
            if w <= 0 or h <= 0:
                continue
            area = w * h

            # 색상
            try:
                if shape.fill.type is not None:
                    rgb = str(shape.fill.fore_color.rgb)
                    color_use[rgb] = color_use.get(rgb, 0) + area
                    if rgb.upper() in (c.upper() for c in forbidden_colors):
                        forbidden_hits.append((rgb, f"shape@{l:.1f},{t:.1f}"))
            except Exception:
                pass

            # 4분면 — shape 중심점이 어디에 속하는가
            cx = l + w / 2
            cy = t + h / 2
            row = "T" if cy < slide_h / 2 else "B"
            col = "L" if cx < slide_w / 2 else "R"
            quadrants[row + col] += area

            # 텍스트 분석
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                total_chars += len(txt)
                # head message 후보: 헤더 영역(y < 1.2") + 폰트 ≥ 13pt + 길이 10자 이상
                shape_y = (shape.top or 0) / 914400
                max_font = 0.0
                for p in shape.text_frame.paragraphs:
                    if p.font.size:
                        sz = p.font.size / 12700
                        font_sizes.append(sz)
                        if sz > max_font:
                            max_font = sz
                if shape_y < 1.2 and max_font >= 13 and len(txt) >= 10:
                    head_candidates.append((max_font, txt))

        # ---------- 1. 금지색 ----------
        if forbidden_hits:
            issues.append(
                DesignIssue(
                    severity="high",
                    category="contrast",
                    message=f"Slide {sn}: 금지색 사용 ({len(forbidden_hits)}건) — {forbidden_hits[0][0]}",
                    suggestion="회색 위계로 교체 (grey_700/800/900)",
                )
            )

        # ---------- 2. 색상 균형 ----------
        # 가장 큰 비중 색이 70% 초과면 모노톤 과다
        if color_use:
            total_colored = sum(color_use.values())
            if total_colored > 0:
                top_color, top_area = max(color_use.items(), key=lambda kv: kv[1])
                top_ratio = top_area / total_colored
                metrics[f"slide_{sn}_top_color"] = (top_color, round(top_ratio, 2))
                if top_ratio > 0.85 and top_color.upper() not in ("FFFFFF",):
                    issues.append(
                        DesignIssue(
                            severity="medium",
                            category="balance",
                            message=f"Slide {sn}: 단일 색이 {top_ratio:.0%} 점유 — 시각 위계 약함",
                            suggestion="2~3 단계 회색 위계 추가 (grey_200/grey_700 등)",
                        )
                    )

        # ---------- 3. 영역 균형 ----------
        # 4분면 점유율 계산 — 임계값 3% (작은 컴포넌트도 있을 수 있음)
        for q, area in quadrants.items():
            ratio = area / slide_area if slide_area > 0 else 0
            metrics[f"slide_{sn}_{q}"] = round(ratio, 2)
            if ratio < 0.03:
                issues.append(
                    DesignIssue(
                        severity="medium",
                        category="balance",
                        message=f"Slide {sn}: {q} 영역이 거의 비어 있음 ({ratio:.0%})",
                        suggestion=f"{q}에 보조 콘텐츠 (stat_block, callout, badge) 추가 검토",
                    )
                )

        # ---------- 4. 시각 위계 ----------
        if font_sizes:
            unique_sizes = sorted(set(round(s) for s in font_sizes))
            metrics[f"slide_{sn}_font_steps"] = unique_sizes
            if len(unique_sizes) < 3:
                issues.append(
                    DesignIssue(
                        severity="medium",
                        category="hierarchy",
                        message=f"Slide {sn}: 폰트 위계 단계 {len(unique_sizes)}개 — 위계 부족",
                        suggestion="제목/부제/본문/디테일 4단계 권장",
                    )
                )
            elif max(unique_sizes) - min(unique_sizes) < 5:
                issues.append(
                    DesignIssue(
                        severity="low",
                        category="hierarchy",
                        message=f"Slide {sn}: 폰트 크기 범위 좁음 ({min(unique_sizes)}~{max(unique_sizes)}pt)",
                        suggestion="제목과 본문 차이를 6pt 이상으로 키워야 위계 강화",
                    )
                )

        # ---------- 5. 정보 밀도 (패턴별 임계값) ----------
        # 임계값 근거: v3 (사용자 OK한 밀도)가 ~19 chars/sq-in, v1 (빈약)이 ~10
        # comparison/quadrant는 짧은 셀이 정상 → density_floor=6
        if expect_dense:
            content_area = slide_area * 0.7  # 헤더/푸터 제외 추정
            density = total_chars / content_area if content_area > 0 else 0
            metrics[f"slide_{sn}_density"] = round(density, 1)
            if density < density_floor:
                issues.append(
                    DesignIssue(
                        severity="high",
                        category="density",
                        message=f"Slide {sn}: 정보 밀도 매우 낮음 ({density:.0f} chars/sq-in, 패턴 floor {density_floor})",
                        suggestion="콘텐츠 대폭 보강 필요 — sub-bullet, stat_block, mini-list",
                    )
                )
            elif density > 80:
                issues.append(
                    DesignIssue(
                        severity="medium",
                        category="density",
                        message=f"Slide {sn}: 정보 밀도 과다 ({density:.0f} chars/sq-in)",
                        suggestion="텍스트 압축 또는 영역 확장",
                    )
                )

        # ---------- 6. Head message 어미 검사 ----------
        if head_candidates:
            # 가장 큰 폰트 = head message로 추정
            head_candidates.sort(reverse=True)
            head_text = head_candidates[0][1]
            metrics[f"slide_{sn}_head"] = head_text[:50]
            ending_issue = check_head_message_ending(head_text)
            if ending_issue:
                # 슬라이드 번호 prefix
                ending_issue.message = f"Slide {sn}: " + ending_issue.message
                issues.append(ending_issue)

    return DesignReport(issues=issues, metrics=metrics)


# ============================================================
# 통합 워크플로 헬퍼
# ============================================================


def inspect_box_fill_ratios(pptx_path: str) -> list[dict]:
    """큰 박스(>= 2"×2") 안의 콘텐츠 채움 비율을 측정한다.

    각 박스에 대해:
    - 박스 좌표/크기
    - 박스 안에 위치한 텍스트 shape들의 총 면적
    - fill_ratio = 텍스트 면적 / 박스 면적 (0.0~1.0)
    - fill_ratio < 0.4이면 비어 보이는 박스

    Returns:
        list of {"slide": int, "box_idx": int, "box_rect": (l,t,w,h),
                 "fill_ratio": float, "sparse": bool}
    """
    prs = Presentation(pptx_path)
    results = []

    for si, slide in enumerate(prs.slides):
        # 1) 모든 shape의 좌표 수집
        shapes_data = []
        for idx, s in enumerate(slide.shapes):
            if s.left is None or s.top is None:
                continue
            l = s.left / 914400
            t = s.top / 914400
            w = (s.width or 0) / 914400
            h = (s.height or 0) / 914400
            if w <= 0 or h <= 0:
                continue

            has_fill = False
            try:
                if s.fill.type is not None:
                    has_fill = True
            except Exception:
                pass

            has_text = s.has_text_frame and bool(s.text_frame.text.strip())
            shapes_data.append({
                "idx": idx, "l": l, "t": t, "w": w, "h": h,
                "area": w * h, "has_fill": has_fill, "has_text": has_text,
            })

        # 2) 큰 박스 식별 (fill 있고, 2"×2" 이상)
        big_boxes = [
            s for s in shapes_data
            if s["has_fill"] and s["w"] >= 1.5 and s["h"] >= 1.5 and s["area"] >= 3.0
        ]

        for box in big_boxes:
            bl, bt, bw, bh = box["l"], box["t"], box["w"], box["h"]
            box_area = bw * bh

            # 3) 박스 안에 있는 텍스트 shape들의 면적 합산
            text_area = 0.0
            for s in shapes_data:
                if not s["has_text"]:
                    continue
                # shape 중심이 박스 안에 있는가?
                cx = s["l"] + s["w"] / 2
                cy = s["t"] + s["h"] / 2
                if bl <= cx <= bl + bw and bt <= cy <= bt + bh:
                    text_area += s["area"]

            fill_ratio = text_area / box_area if box_area > 0 else 0
            results.append({
                "slide": si + 1,
                "box_idx": box["idx"],
                "box_rect": (round(bl, 2), round(bt, 2), round(bw, 2), round(bh, 2)),
                "fill_ratio": round(fill_ratio, 2),
                "sparse": fill_ratio < 0.4,
            })

    return results


def decide_aux_content(
    *,
    pattern_kind: Literal[
        "executive", "timeline", "comparison", "process", "quadrant", "data"
    ],
    box_context: str = "",
    fill_ratio: float = 0.0,
) -> DesignDecision:
    """빈 박스에 어떤 보조 콘텐츠를 채울지 자동 결정한다.

    Claude가 이걸 호출하는 게 아니라, 패턴 함수가 빈 공간을 감지했을 때
    자동으로 적절한 aux 슬롯을 결정하는 데 사용한다.

    Args:
        pattern_kind: 패턴 종류
        box_context: 박스가 어떤 단계/영역에 해당하는지 (예: "L1 가시화", "Blueprint 수집")
        fill_ratio: 현재 채움 비율 (0.0~1.0)

    Returns:
        DesignDecision with recommended aux content types
    """
    # 패턴별 aux 콘텐츠 후보 — 우선순위 순서
    aux_pool = {
        "process": [
            {"type": "prerequisites", "label": "전제 조건",
             "rationale": "이 단계를 시작하기 위한 선행 조건 — 단계 간 의존성을 보여줌"},
            {"type": "risks", "label": "리스크 / 주의사항",
             "rationale": "이 단계에서 발생 가능한 위험 — 실무자가 즉시 참고"},
            {"type": "metrics", "label": "정량 효과",
             "rationale": "이 단계 완료 시 기대 효과 — 투자 근거 강화"},
            {"type": "example", "label": "실증 사례",
             "rationale": "유사 프로젝트에서의 적용 결과 — 신뢰성 확보"},
        ],
        "timeline": [
            {"type": "prerequisites", "label": "선행 조건",
             "rationale": "이전 단계 완료 기준 — 의존성 체인 시각화"},
            {"type": "gate_criteria", "label": "Gate 기준",
             "rationale": "다음 단계 착수 조건 — PMO 의사결정 근거"},
            {"type": "team", "label": "투입 인력",
             "rationale": "이 단계에 필요한 역할 — 리소스 계획 지원"},
            {"type": "risks", "label": "리스크",
             "rationale": "일정 지연 가능 요인 — 선제 대응 가능"},
        ],
        "executive": [
            {"type": "metrics", "label": "핵심 KPI",
             "rationale": "정량적 성과 수치 — 경영진 즉시 판단 가능"},
            {"type": "next_steps", "label": "후속 조치",
             "rationale": "의사결정 후 즉시 실행 사항 — action item"},
        ],
        "comparison": [
            {"type": "notes", "label": "비고",
             "rationale": "셀만으로 설명이 부족한 경우 보조 정보"},
            {"type": "recommendation", "label": "권장 사항",
             "rationale": "비교 결과 기반 권고 — 의사결정 지원"},
        ],
        "quadrant": [
            {"type": "action_items", "label": "실행 과제",
             "rationale": "각 사분면에 대한 구체적 후속 조치"},
            {"type": "criteria", "label": "분류 기준",
             "rationale": "축의 의미를 더 구체적으로 설명"},
        ],
        "data": [
            {"type": "insight", "label": "인사이트",
             "rationale": "데이터에서 도출한 의미 — 숫자만으로는 부족"},
            {"type": "methodology", "label": "분석 방법",
             "rationale": "어떻게 이 결론에 도달했는가 — 신뢰성"},
        ],
    }

    pool = aux_pool.get(pattern_kind, aux_pool["executive"])

    # fill_ratio에 따라 추천 개수 결정
    # 매우 빈약 (< 0.25) → 2개 추천, 약간 빈약 (0.25~0.4) → 1개 추천
    if fill_ratio < 0.25:
        n_recommend = min(2, len(pool))
    elif fill_ratio < 0.4:
        n_recommend = 1
    else:
        n_recommend = 0

    recommended = pool[:n_recommend]
    alternatives = pool[n_recommend:n_recommend + 2]

    return DesignDecision(
        rationale=(
            f"{pattern_kind} 패턴, fill_ratio={fill_ratio:.0%} — "
            f"{n_recommend}개 aux 콘텐츠 추천"
            + (f" ({', '.join(r['label'] for r in recommended)})" if recommended else "")
        ),
        recommendation={
            "n_aux": n_recommend,
            "aux_types": recommended,
            "box_context": box_context,
        },
        alternatives=[{"type": a["type"], "reason": a["rationale"]} for a in alternatives],
    )


def design_check_pipeline(
    pptx_path: str,
    *,
    expect_dense: bool = True,
    forbidden_colors: Optional[list[str]] = None,
    pattern_kind: Optional[str] = None,
) -> DesignReport:
    """convenience: render_validated 후 호출하여 디자인 점검까지 한 번에.

    디자인 점검 + 박스별 fill ratio 측정을 모두 수행.
    """
    report = inspect_design(
        pptx_path,
        expect_dense=expect_dense,
        forbidden_colors=forbidden_colors,
        pattern_kind=pattern_kind,
    )
    # 박스별 fill ratio 측정 결과를 metrics에 추가
    box_ratios = inspect_box_fill_ratios(pptx_path)
    sparse_boxes = [b for b in box_ratios if b["sparse"]]
    report.metrics["box_fill_ratios"] = box_ratios
    report.metrics["sparse_box_count"] = len(sparse_boxes)

    # sparse 박스가 있으면 MEDIUM issue 추가
    for b in sparse_boxes:
        report.issues.append(
            DesignIssue(
                severity="medium",
                category="density",
                message=(
                    f"Slide {b['slide']}: 박스 {b['box_rect']} fill_ratio={b['fill_ratio']:.0%} — "
                    "콘텐츠 부족으로 빈 공간 과다"
                ),
                suggestion="aux 콘텐츠(전제조건/리스크/KPI/사례) 추가 또는 shrink-to-fit 적용",
            )
        )

    return report
