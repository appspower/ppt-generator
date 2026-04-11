"""PPT 자동 품질 평가 — Step 5: EVALUATE.

사용법:
  from ppt_builder.evaluate import evaluate_pptx
  report = evaluate_pptx("output/my_slide.pptx")
  print(report)
"""

from pathlib import Path
from pptx import Presentation


def evaluate_pptx(path: str | Path) -> dict:
    """생성된 PPT의 품질을 자동 평가한다.

    Returns:
        {
            "score": int (0~100),
            "pass": bool,
            "issues": list[str],
            "metrics": dict,
        }
    """
    prs = Presentation(str(path))
    issues = []
    metrics = {}

    for si, slide in enumerate(prs.slides):
        slide_issues = _evaluate_slide(slide, si)
        issues.extend(slide_issues)

    # 덱 레벨 평가
    deck_issues = _evaluate_deck(list(prs.slides))
    issues.extend(deck_issues)

    # 점수 계산 (100점 - 이슈당 감점)
    deductions = {
        "CRITICAL": 15,
        "HIGH": 8,
        "MEDIUM": 4,
        "LOW": 2,
    }
    total_deduction = sum(deductions.get(iss.split(":")[0], 3) for iss in issues)
    score = max(0, 100 - total_deduction)

    return {
        "score": score,
        "pass": score >= 75,
        "issues": issues,
        "metrics": metrics,
        "slide_count": len(prs.slides),
    }


def _evaluate_slide(slide, slide_idx: int) -> list[str]:
    """단일 슬라이드 평가."""
    issues = []
    shapes = slide.shapes

    # 1. Shape 수
    n_shapes = len(shapes)
    if n_shapes < 5:
        issues.append(f"LOW: Slide {slide_idx+1}: shape 수 부족 ({n_shapes}개)")

    # 2. 텍스트 밀도
    total_chars = 0
    font_sizes = {}
    accent_shapes = 0
    total_colored = 0

    for s in shapes:
        # 색상 분석
        try:
            if s.fill.type is not None:
                rgb = str(s.fill.fore_color.rgb)
                total_colored += 1
                if rgb in ("FD5108", "FE7C39"):  # accent colors
                    accent_shapes += 1
        except:
            pass

        # 텍스트 분석
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                chars = len(p.text)
                total_chars += chars
                if p.font.size:
                    sz = round(p.font.size / 12700, 0)
                    font_sizes[sz] = font_sizes.get(sz, 0) + chars

    # 3. 텍스트 밀도 (차트 중심 슬라이드 완화: 300+, 텍스트 중심: 500+)
    has_chart = any(s.has_chart for s in shapes if hasattr(s, 'has_chart'))
    min_chars = 200 if has_chart else 300
    if total_chars < 80:
        issues.append(f"HIGH: Slide {slide_idx+1}: 텍스트 부족 ({total_chars}자)")
    elif total_chars < min_chars:
        issues.append(f"MEDIUM: Slide {slide_idx+1}: 텍스트 밀도 낮음 ({total_chars}자, 목표 {min_chars}+)")
    elif total_chars > 2000:
        issues.append(f"LOW: Slide {slide_idx+1}: 텍스트 과다 ({total_chars}자)")

    # 4. 폰트 위계 체크
    if font_sizes:
        size_range = max(font_sizes.keys()) - min(font_sizes.keys())
        if size_range < 4:
            issues.append(f"MEDIUM: Slide {slide_idx+1}: 폰트 위계 부족 (범위 {size_range}pt)")

    # 5. 색상 절제 체크
    if total_colored > 0:
        accent_ratio = accent_shapes / total_colored
        if accent_ratio > 0.25:
            issues.append(f"MEDIUM: Slide {slide_idx+1}: accent 과다 ({accent_ratio:.0%})")

    # 6. 빈 공간 체크
    max_bottom = 0
    for s in shapes:
        if s.top and s.height:
            bottom = (s.top + s.height) / 914400
            if bottom > max_bottom:
                max_bottom = bottom
    empty_bottom = 7.5 - max_bottom
    if empty_bottom > 1.5:
        issues.append(f"HIGH: Slide {slide_idx+1}: 하단 빈 공간 과다 ({empty_bottom:.1f}\")")

    # 7. 제목 체크 (첫 번째 텍스트가 인사이트 문장인지)
    title_text = ""
    for s in shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            title_text = s.text_frame.text.strip()
            break
    if title_text and len(title_text) < 10:
        issues.append(f"LOW: Slide {slide_idx+1}: 제목이 너무 짧음 (라벨형?)")

    # 8. Overflow 추정 — 텍스트 프레임 밀도 (Track C 베이스라인에서 발견한 사각지대 #1)
    # 텍스트가 컨테이너를 넘치는지 python-pptx로 정확히 알 수 없지만,
    # "작은 박스에 긴 텍스트"를 근사적으로 감지한다.
    for s in shapes:
        if s.has_text_frame and s.width and s.height:
            frame_text = s.text_frame.text.strip()
            if not frame_text:
                continue
            # 컨테이너 면적 (제곱인치)
            area_sq_in = (s.width / 914400) * (s.height / 914400)
            if area_sq_in < 0.5:
                continue  # 너무 작은 shape는 라벨이므로 무시
            # 글자수 / 면적 비율 (chars per sq inch)
            density = len(frame_text) / area_sq_in
            if density > 200:
                issues.append(
                    f"HIGH: Slide {slide_idx+1}: overflow 의심 — "
                    f"텍스트 밀도 {density:.0f}자/in² (shape '{frame_text[:20]}...')"
                )
            elif density > 140:
                issues.append(
                    f"MEDIUM: Slide {slide_idx+1}: overflow 위험 — "
                    f"텍스트 밀도 {density:.0f}자/in² (shape '{frame_text[:20]}...')"
                )

    # 9. 시각 다양성 — 컴포넌트 타입 다양성 (사각지대 #2: 평면성)
    # 단순 표만 있는 슬라이드(style_b_r8 유형)를 감지한다.
    shape_categories = set()
    for s in shapes:
        if s.has_table:
            shape_categories.add("table")
        elif s.has_chart:
            shape_categories.add("chart")
        elif hasattr(s, "image"):
            shape_categories.add("image")
        elif s.has_text_frame:
            text = s.text_frame.text.strip()
            if not text:
                continue
            # 큰 텍스트 (헤더/타이틀) vs 작은 텍스트 (본문)
            if s.height and s.height / 914400 > 0.6:
                shape_categories.add("text_large")
            else:
                shape_categories.add("text_small")
        else:
            shape_categories.add("shape")

    # 의미 있는 카테고리가 2종 미만이면 평면적
    if n_shapes >= 5 and len(shape_categories) < 2:
        issues.append(
            f"MEDIUM: Slide {slide_idx+1}: 시각 다양성 부족 — "
            f"카테고리 {len(shape_categories)}종 ({', '.join(shape_categories)}), 컨설팅 PPT는 3종+ 권장"
        )

    # 10. 공간 활용률 — shape 면적 합 / 슬라이드 면적
    total_shape_area = 0
    for s in shapes:
        if s.width and s.height:
            total_shape_area += (s.width / 914400) * (s.height / 914400)
    slide_area = 10 * 7.5  # 10x7.5 inches
    if slide_area > 0:
        coverage = total_shape_area / slide_area
        if coverage < 0.35:
            issues.append(
                f"HIGH: Slide {slide_idx+1}: 공간 활용률 낮음 ({coverage:.0%})"
            )
        elif coverage < 0.50:
            issues.append(
                f"MEDIUM: Slide {slide_idx+1}: 공간 활용률 부족 ({coverage:.0%}, 목표 50%+)"
            )

    return issues


def _evaluate_deck(slides) -> list[str]:
    """덱 레벨 평가 — 인접 슬라이드 유사도."""
    issues = []
    if len(slides) < 2:
        return issues

    prev_shape_count = 0
    for si, slide in enumerate(slides):
        n = len(slide.shapes)
        if si > 0 and n > 0 and prev_shape_count > 0:
            # shape 수가 ±2 이내이고 둘 다 5개 이상이면 유사 경고
            if abs(n - prev_shape_count) <= 2 and n >= 5 and prev_shape_count >= 5:
                # 추가 확인: 같은 타입의 shape가 많은지
                pass  # 향후 고도화
        prev_shape_count = n

    return issues


def print_report(report: dict):
    """평가 보고서를 출력한다."""
    import sys, io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    status = "PASS" if report["pass"] else "FAIL"
    print(f"\n{'='*50}")
    print(f"PPT Quality: {report['score']}/100 [{status}]")
    print(f"Slides: {report['slide_count']}")
    print(f"{'='*50}")

    if report["issues"]:
        print(f"\nIssues ({len(report['issues'])}):")
        for iss in report["issues"]:
            print(f"  - {iss}")
    else:
        print("\nNo issues")

    return report


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python -m ppt_builder.evaluate <pptx_file>")
        sys.exit(1)
    report = evaluate_pptx(sys.argv[1])
    print_report(report)
