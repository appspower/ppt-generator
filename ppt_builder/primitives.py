"""Layer 2 — 절대좌표 디자인 프리미티브.

매번 새로운 화면 구성을 만들 수 있도록, 19개 프리셋 컴포넌트가 아니라
"백지 + 펜"을 제공한다. Claude가 슬라이드별로 좌표를 직접 결정하여
배치하면, 같은 주제로도 매번 다른 레이아웃이 나온다.

설계 원칙:
- 모든 좌표는 인치 단위 (사용자 직관)
- 굵고 각진 도형 (둥근 모서리 ❌, 두꺼운 라인 ✓)
- 안전망: validate_visual()이 잡아주므로 자유도 높여도 됨
- 기존 컴포넌트 시스템과 독립 — 기존 코드는 손대지 않음

사용 예:
    from pptx import Presentation
    from ppt_builder.primitives import Canvas

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    c = Canvas(slide)
    c.title("핵심 메시지", x=0.3, y=0.2, w=9.4)
    c.box(x=0.3, y=1.0, w=4.5, h=3.0, fill="dark", border=2)
    c.text("내용", x=0.5, y=1.2, w=4.1, h=2.6, color="white", size=11)
    c.arrow(x1=4.9, y1=2.5, x2=5.5, y2=2.5)
    c.kpi(value="14%", label="일정 단축", x=6.0, y=1.0, w=3.7, h=1.5)
"""

from __future__ import annotations

from typing import Literal, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

from .assembler.styles import (
    CL_ACCENT,
    CL_ACCENT_LIGHT,
    CL_ACCENT_MID,
    CL_BLACK,
    CL_BORDER,
    CL_DARK,
    CL_GREY,
    CL_GREY_LIGHT,
    CL_GREY_MID,
    CL_NEGATIVE,
    CL_POSITIVE,
    CL_WHITE,
    FONT_BODY,
    FONT_TITLE,
)


# ============================================================
# Design tokens — 굵고 각진 스타일
# ============================================================

# 색상 별칭 (사용자 친화 이름)
COLORS = {
    "white": CL_WHITE,
    "black": CL_BLACK,
    "dark": CL_DARK,
    "accent": CL_ACCENT,
    "accent_mid": CL_ACCENT_MID,
    "accent_light": CL_ACCENT_LIGHT,
    "grey": CL_GREY,
    "grey_mid": CL_GREY_MID,
    "grey_light": CL_GREY_LIGHT,
    "border": CL_BORDER,
    "positive": CL_POSITIVE,
    "negative": CL_NEGATIVE,
    # 회색 위계 확장 (5단계) — 완전 블랙 회피용
    "grey_900": RGBColor(0x2E, 0x33, 0x3A),  # 거의 검정에 가까운 다크 그레이
    "grey_800": RGBColor(0x4A, 0x4F, 0x58),  # 다크 그레이 (Hero 배경용)
    "grey_700": RGBColor(0x6B, 0x71, 0x7B),  # 미디엄 다크 그레이
    "grey_400": RGBColor(0x9A, 0xA0, 0xA8),  # 미디엄 그레이
    "grey_200": RGBColor(0xE2, 0xE5, 0xE8),  # 라이트 그레이 (배경)
    "grey_100": RGBColor(0xF1, 0xF3, 0xF5),  # 매우 라이트 그레이
}

# 굵은 라인 두께 — 1pt = 12700 EMU
STROKE_THIN = Emu(6350)        # 0.5pt
STROKE_NORMAL = Emu(12700)     # 1pt
STROKE_BOLD = Emu(19050)       # 1.5pt
STROKE_HEAVY = Emu(25400)      # 2pt
STROKE_EXTRA = Emu(38100)      # 3pt


def color(name_or_rgb) -> RGBColor:
    """문자열 별칭 또는 RGBColor를 RGBColor로 정규화."""
    if isinstance(name_or_rgb, RGBColor):
        return name_or_rgb
    if isinstance(name_or_rgb, str):
        if name_or_rgb in COLORS:
            return COLORS[name_or_rgb]
        if name_or_rgb.startswith("#") and len(name_or_rgb) == 7:
            return RGBColor(
                int(name_or_rgb[1:3], 16),
                int(name_or_rgb[3:5], 16),
                int(name_or_rgb[5:7], 16),
            )
    raise ValueError(f"Unknown color: {name_or_rgb!r}")


# ============================================================
# Canvas — 백지 위에 도형/텍스트를 절대좌표로 그리는 펜
# ============================================================


class Canvas:
    """슬라이드를 백지로 보고 절대좌표로 도형을 그린다.

    좌표는 인치 단위 (10x7.5 슬라이드 기준 0~10, 0~7.5).
    원점은 좌상단.
    """

    def __init__(self, slide: Slide):
        self.slide = slide
        self._drawn: list[tuple] = []  # 디버깅/검증용 추적

    # --------------------------------------------------------
    # Atomic shapes
    # --------------------------------------------------------

    def box(
        self,
        *,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: str | RGBColor | None = "white",
        border: float | None = 0.75,
        border_color: str | RGBColor = "grey_mid",
        shape: Literal["rect", "rounded"] = "rect",
    ):
        """굵고 각진 사각형. 둥근 모서리는 명시적으로 요청해야 적용됨.

        Args:
            x, y, w, h: 인치 단위 위치/크기
            fill: 채움색 ("white"/"dark"/"accent"/None=투명)
            border: 테두리 두께 pt (None=없음, 1.5pt가 기본 굵기)
            border_color: 테두리 색
            shape: "rect" (기본, 각진) | "rounded" (둥근)
        """
        mso_shape = (
            MSO_SHAPE.ROUNDED_RECTANGLE if shape == "rounded" else MSO_SHAPE.RECTANGLE
        )
        box = self.slide.shapes.add_shape(
            mso_shape, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        # 둥근 모서리 강도 (선택 시에만)
        if shape == "rounded":
            try:
                box.adjustments[0] = 0.05  # 매우 미묘
            except (IndexError, TypeError):
                pass

        if fill is None:
            box.fill.background()
        else:
            box.fill.solid()
            box.fill.fore_color.rgb = color(fill)

        if border is None:
            box.line.fill.background()
        else:
            box.line.color.rgb = color(border_color)
            box.line.width = Pt(border)

        self._drawn.append(("box", x, y, w, h))
        return box

    def circle(
        self,
        *,
        x: float,
        y: float,
        d: float,
        fill: str | RGBColor | None = "white",
        border: float | None = 0.75,
        border_color: str | RGBColor = "grey_700",
        text: str = "",
        text_color: str | RGBColor = "black",
        text_size: float = 11,
        text_bold: bool = True,
    ):
        """원형 도형 (지름 d). 번호 뱃지/카테고리 마커 등에 사용."""
        oval = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
        )
        if fill is None:
            oval.fill.background()
        else:
            oval.fill.solid()
            oval.fill.fore_color.rgb = _color(fill)
        if border is None:
            oval.line.fill.background()
        else:
            oval.line.color.rgb = _color(border_color)
            oval.line.width = Pt(border)

        if text:
            tf = oval.text_frame
            tf.margin_left = Inches(0)
            tf.margin_right = Inches(0)
            tf.margin_top = Inches(0)
            tf.margin_bottom = Inches(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(text_size)
            p.font.bold = text_bold
            p.font.color.rgb = _color(text_color)
            p.font.name = FONT_BODY
            p.alignment = PP_ALIGN.CENTER

        self._drawn.append(("circle", x, y, d, d))
        return oval

    def chevron(
        self,
        *,
        x: float,
        y: float,
        w: float,
        h: float,
        fill: str | RGBColor = "grey_400",
        border: float | None = None,
        text: str = "",
        text_color: str | RGBColor = "white",
        text_size: float = 10,
        text_bold: bool = True,
    ):
        """5각형 화살표 (chevron). 단계/순서 표시에 사용."""
        chev = self.slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        chev.fill.solid()
        chev.fill.fore_color.rgb = _color(fill)
        if border is None:
            chev.line.fill.background()
        else:
            chev.line.color.rgb = _color("grey_mid")
            chev.line.width = Pt(border)

        if text:
            tf = chev.text_frame
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.15)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            for i, line in enumerate(text.split("\n")):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(text_size if i == 0 else text_size - 1)
                p.font.bold = text_bold and i == 0
                p.font.color.rgb = _color(text_color)
                p.font.name = FONT_BODY
                p.alignment = PP_ALIGN.CENTER
        self._drawn.append(("chevron", x, y, w, h))
        return chev

    def line(
        self,
        *,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        color: str | RGBColor = "black",
        width: float = 1.5,
    ):
        """직선 (수평/수직/대각)."""
        ln = self.slide.shapes.add_connector(
            1,  # straight line
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        ln.line.color.rgb = _color(color)
        ln.line.width = Pt(width)
        self._drawn.append(("line", x1, y1, x2, y2))
        return ln

    def arrow(
        self,
        *,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        color: str | RGBColor = "dark",
        width: float = 1.0,
    ):
        """화살표 — 보통 두께, 다크 기본."""
        ln = self.slide.shapes.add_connector(
            1,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        ln.line.color.rgb = _color(color)
        ln.line.width = Pt(width)
        # 화살표 머리 추가
        from pptx.oxml.ns import qn
        from copy import deepcopy
        line_elem = ln.line._get_or_add_ln()
        # 기존 head/tail 제거
        for tag in ("headEnd", "tailEnd"):
            for el in line_elem.findall(qn(f"a:{tag}")):
                line_elem.remove(el)
        from lxml import etree
        tail = etree.SubElement(line_elem, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
        self._drawn.append(("arrow", x1, y1, x2, y2))
        return ln

    def text(
        self,
        body: str,
        *,
        x: float,
        y: float,
        w: float,
        h: float,
        size: float = 10,
        bold: bool = False,
        color: str | RGBColor = "black",
        font: str = FONT_BODY,
        align: Literal["left", "center", "right"] = "left",
        anchor: Literal["top", "middle", "bottom"] = "top",
    ):
        """텍스트 박스 — 절대좌표 + 명시적 폰트 제어."""
        tx = self.slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        tf = tx.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }[anchor]

        align_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }
        for i, line in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = _color(color)
            p.font.name = font
            p.alignment = align_map[align]
        self._drawn.append(("text", x, y, w, h, len(body)))
        return tx

    # --------------------------------------------------------
    # Composite helpers — 자주 쓰는 조합 (그래도 자유도 유지)
    # --------------------------------------------------------

    def title(
        self,
        text: str,
        *,
        x: float = 0.3,
        y: float = 0.2,
        w: float = 9.4,
        h: float = 0.45,
        size: float = 16,
        underline: bool = True,
        underline_color: str = "dark",
        underline_thickness: float = 0.02,
    ):
        """슬라이드 제목 — 검정 굵은 + 하단 가는 라인 (선택)."""
        self.text(
            text,
            x=x,
            y=y,
            w=w,
            h=h,
            size=size,
            bold=True,
            color="black",
            font=FONT_TITLE,
            anchor="middle",
        )
        if underline:
            self.box(
                x=x,
                y=y + h + 0.02,
                w=w,
                h=underline_thickness,
                fill=underline_color,
                border=None,
            )

    def kpi(
        self,
        *,
        value: str,
        label: str,
        x: float,
        y: float,
        w: float,
        h: float,
        detail: str = "",
        stripe: bool = True,
    ):
        """대형 KPI 숫자 + 라벨 (+ optional detail) — 회색 톤."""
        self.box(
            x=x, y=y, w=w, h=h,
            fill="white", border=0.75, border_color="grey_mid",
        )
        if stripe:
            self.box(
                x=x, y=y, w=0.08, h=h,
                fill="grey_700", border=None,
            )
        # 큰 숫자
        v_size = 26 if len(value) <= 3 else 20
        self.text(
            value,
            x=x + 0.2, y=y + 0.12, w=w - 0.3, h=h * 0.5,
            size=v_size, bold=True, color="grey_900",
            font=FONT_TITLE, anchor="top",
        )
        # 라벨
        self.text(
            label,
            x=x + 0.2, y=y + h * 0.55, w=w - 0.3, h=0.25,
            size=9, bold=True, color="grey_900", anchor="top",
        )
        # 디테일 (선택)
        if detail:
            self.text(
                detail,
                x=x + 0.2, y=y + h * 0.55 + 0.27, w=w - 0.3, h=h * 0.3,
                size=7, color="grey_700", anchor="top",
            )

    def label_chip(
        self,
        text: str,
        *,
        x: float,
        y: float,
        w: float = 1.2,
        h: float = 0.28,
        fill: str = "grey_700",
        text_color: str = "white",
        size: float = 8,
    ):
        """카테고리/태그 — 작은 색 박스 + 흰 텍스트."""
        self.box(x=x, y=y, w=w, h=h, fill=fill, border=None)
        self.text(
            text,
            x=x, y=y, w=w, h=h,
            size=size, bold=True,
            color=text_color, align="center", anchor="middle",
        )

    def divider_h(
        self,
        *,
        x: float,
        y: float,
        w: float,
        color: str = "grey_mid",
        width: float = 0.75,
    ):
        """수평 구분선 — 보통 두께."""
        self.line(x1=x, y1=y, x2=x + w, y2=y, color=color, width=width)

    def divider_v(
        self,
        *,
        x: float,
        y: float,
        h: float,
        color: str = "grey_mid",
        width: float = 0.75,
    ):
        """수직 구분선 — 보통 두께."""
        self.line(x1=x, y1=y, x2=x, y2=y + h, color=color, width=width)

    # --------------------------------------------------------
    # Phase A — Rich composites
    # --------------------------------------------------------

    def badge(
        self,
        text: str,
        *,
        x: float,
        y: float,
        h: float = 0.24,
        fill: str = "grey_200",
        text_color: str = "grey_900",
        size: float = 8,
        rounded: bool = True,
    ):
        """작은 카테고리 라벨 — 자동 폭 계산.

        chip보다 작고, 자동으로 폭이 텍스트에 맞춰진다.
        """
        # 텍스트 길이로 폭 추정 (8pt 기준)
        char_w = 0.06 if size <= 8 else 0.07
        # 한글 비중 가산
        kr = sum(1 for ch in text if ord(ch) > 0x1100) / max(len(text), 1)
        char_w *= 1 + kr * 0.6
        w = max(0.4, len(text) * char_w + 0.18)
        shape_kind = "rounded" if rounded else "rect"
        self.box(
            x=x, y=y, w=w, h=h,
            fill=fill, border=None, shape=shape_kind,
        )
        self.text(
            text,
            x=x, y=y, w=w, h=h,
            size=size, bold=True,
            color=text_color, align="center", anchor="middle",
        )
        return w  # 자동 폭 반환 — 여러 뱃지 연달아 배치할 때 유용

    def callout_box(
        self,
        *,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str = "",
        body: str = "",
        bullets: list[str] | None = None,
        bar_color: str = "grey_700",
        bar_width: float = 0.08,
        fill: str = "white",
        bordered: bool = True,
        title_size: float = 11,
        body_size: float = 9,
    ):
        """좌측 컬러바 + 강조 박스. title/body/bullets 조합 가능."""
        if bordered:
            self.box(
                x=x, y=y, w=w, h=h,
                fill=fill, border=0.75, border_color="grey_mid",
            )
        else:
            self.box(x=x, y=y, w=w, h=h, fill=fill, border=None)
        # 좌측 컬러바
        self.box(
            x=x, y=y, w=bar_width, h=h,
            fill=bar_color, border=None,
        )
        pad = 0.12
        cx = x + bar_width + pad
        cw = w - bar_width - pad * 2
        cy = y + 0.1

        if title:
            self.text(
                title,
                x=cx, y=cy, w=cw, h=0.3,
                size=title_size, bold=True,
                color="grey_900", anchor="top",
            )
            cy += 0.32

        if body:
            # 본문 줄 수 추정
            body_lines = body.count("\n") + 1
            body_h = max(0.25, body_lines * 0.18)
            self.text(
                body,
                x=cx, y=cy, w=cw, h=body_h,
                size=body_size, color="grey_900", anchor="top",
            )
            cy += body_h + 0.05

        if bullets:
            for i, bul in enumerate(bullets):
                self.text(
                    f"▪  {bul}",
                    x=cx, y=cy + i * 0.22, w=cw, h=0.2,
                    size=body_size - 1, color="grey_700", anchor="top",
                )

    def arrow_chain(
        self,
        items: list[str],
        *,
        x: float,
        y: float,
        w: float,
        h: float,
        gap: float = 0.18,
        fill: str = "grey_200",
        text_color: str = "grey_900",
        text_size: float = 10,
        with_arrows: bool = True,
        arrow_color: str = "grey_700",
    ):
        """박스 N개 → 화살표 → 박스. 시퀀스/플로우 표현.

        gap 안에 화살표가 들어간다 (with_arrows=True일 때).
        """
        n = len(items)
        if n == 0:
            return
        # 박스 폭 계산: 전체 폭에서 (gap × n-1)을 뺀 후 균등 분배
        box_w = (w - gap * (n - 1)) / n
        for i, item in enumerate(items):
            bx = x + i * (box_w + gap)
            self.box(
                x=bx, y=y, w=box_w, h=h,
                fill=fill, border=0.75, border_color="grey_mid",
            )
            self.text(
                item,
                x=bx, y=y, w=box_w, h=h,
                size=text_size, bold=True,
                color=text_color, align="center", anchor="middle",
            )
            # 화살표 (마지막 박스 제외)
            if with_arrows and i < n - 1:
                ay = y + h / 2
                ax_start = bx + box_w + 0.02
                ax_end = bx + box_w + gap - 0.02
                self.arrow(
                    x1=ax_start, y1=ay, x2=ax_end, y2=ay,
                    color=arrow_color, width=1.0,
                )

    def dot_grid(
        self,
        *,
        x: float,
        y: float,
        filled: int,
        total: int = 5,
        d: float = 0.14,
        gap: float = 0.06,
        fill_on: str = "grey_800",
        fill_off: str = "grey_200",
    ):
        """N개 점 중 filled개가 채워진 진행률 표시. Harvey Ball 대안."""
        for i in range(total):
            cx = x + i * (d + gap)
            self.circle(
                x=cx, y=y, d=d,
                fill=fill_on if i < filled else fill_off,
                border=None,
                text="",
            )

    def mini_table(
        self,
        *,
        x: float,
        y: float,
        w: float,
        h: float,
        headers: list[str],
        rows: list[list[str]],
        col_ratios: list[float] | None = None,
        header_size: float = 9,
        body_size: float = 8,
    ):
        """매우 컴팩트한 표 — 강조 헤더 + 교대행. python-pptx table 대신
        절대좌표 박스로 그려서 행 높이 자동확장 문제 회피."""
        n_cols = len(headers)
        n_rows = len(rows)
        if col_ratios and len(col_ratios) == n_cols:
            total = sum(col_ratios)
            widths = [w * (r / total) for r in col_ratios]
        else:
            widths = [w / n_cols] * n_cols

        header_h = 0.28
        body_h = (h - header_h) / max(n_rows, 1)
        # 헤더
        cx = x
        for i, hdr in enumerate(headers):
            self.box(
                x=cx, y=y, w=widths[i], h=header_h,
                fill="grey_800", border=None,
            )
            self.text(
                hdr,
                x=cx + 0.05, y=y, w=widths[i] - 0.1, h=header_h,
                size=header_size, bold=True, color="white", anchor="middle",
            )
            cx += widths[i]
        # 바디
        for ri, row in enumerate(rows):
            ry = y + header_h + ri * body_h
            row_fill = "white" if ri % 2 == 0 else "grey_100"
            cx = x
            for ci, cell in enumerate(row):
                self.box(
                    x=cx, y=ry, w=widths[ci], h=body_h,
                    fill=row_fill, border=0.5, border_color="grey_mid",
                )
                self.text(
                    str(cell),
                    x=cx + 0.06, y=ry, w=widths[ci] - 0.12, h=body_h,
                    size=body_size, color="grey_900", anchor="middle",
                )
                cx += widths[ci]

    def stat_block(
        self,
        *,
        value: str,
        label: str,
        x: float,
        y: float,
        w: float,
        h: float = 0.9,
        align: Literal["left", "center", "right"] = "left",
        accent: bool = False,
    ):
        """작은 통계 블록 — KPI보다 가벼움. 라벨 + 큰 숫자만, 박스 없음.

        라벨이 위에, 숫자가 아래에 위치 (시각적 위계 강조).
        """
        # 라벨 (작게, 위)
        self.text(
            label,
            x=x, y=y, w=w, h=0.22,
            size=8, color="grey_700", align=align, anchor="top",
        )
        # 큰 숫자 (아래)
        v_color = "grey_900"
        v_size = 22 if len(value) <= 5 else 16
        self.text(
            value,
            x=x, y=y + 0.22, w=w, h=h - 0.22,
            size=v_size, bold=True, color=v_color,
            font=FONT_TITLE, align=align, anchor="top",
        )
        # 미세한 하단 라인 (선택)
        if accent:
            self.box(
                x=x, y=y + h - 0.04, w=min(0.5, w * 0.3), h=0.025,
                fill="grey_700", border=None,
            )

    def numbered_list(
        self,
        items: list[tuple[str, str]],
        *,
        x: float,
        y: float,
        w: float,
        item_h: float = 0.6,
        gap: float = 0.12,
        circle_d: float = 0.36,
        circle_fill: str = "grey_900",
        circle_text_color: str = "white",
        text_color: str = "grey_900",
        title_size: float = 11,
        detail_size: float = 8,
    ):
        """원형 번호 + 제목 + 디테일 리스트.

        items: [(title, detail), (title, detail), ...]
        """
        for i, (title, detail) in enumerate(items):
            iy = y + i * (item_h + gap)
            # 원형 번호
            self.circle(
                x=x, y=iy + 0.04, d=circle_d,
                fill=circle_fill, border=None,
                text=f"{i+1:02d}", text_color=circle_text_color,
                text_size=10,
            )
            # 텍스트 영역
            tx = x + circle_d + 0.15
            tw = w - circle_d - 0.15
            self.text(
                title,
                x=tx, y=iy, w=tw, h=0.28,
                size=title_size, bold=True, color=text_color, anchor="top",
            )
            if detail:
                self.text(
                    detail,
                    x=tx, y=iy + 0.28, w=tw, h=item_h - 0.28,
                    size=detail_size, color="grey_700", anchor="top",
                )

    def section_label(
        self,
        text: str,
        *,
        x: float,
        y: float,
        w: float,
        size: float = 10,
        bar_color: str = "grey_700",
    ):
        """섹션 헤더 — 좌측 미세 바 + 굵은 텍스트. 영역 구분용."""
        bar_w = 0.08
        self.box(
            x=x, y=y + 0.04, w=bar_w, h=0.18,
            fill=bar_color, border=None,
        )
        self.text(
            text,
            x=x + bar_w + 0.08, y=y, w=w - bar_w - 0.08, h=0.28,
            size=size, bold=True, color="grey_900", anchor="top",
        )


# ============================================================
# 내부 helper
# ============================================================


def _color(name_or_rgb) -> RGBColor:
    """color()의 별칭 — 메서드 인자명 'color'와 충돌 회피용."""
    return color(name_or_rgb)
