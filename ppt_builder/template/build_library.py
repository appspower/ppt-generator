"""template_library.pptx 빌더 — 완성본 사진 정밀 분석 기반 고품질 버전.

완성본 94장 분석 결과 TOP 5 패턴:
  0: 4-Column Reference (색상 헤더 + 밀도 높은 불릿)
  1: Framework Matrix + Sidebar Navigation
  2: 4-Step Decision Flow (아이콘 + 시스템 박스)
  3: Comparison (2안 비교 + 추천 하이라이트 + 하단 결론)
  4: Timeline Roadmap (수평 바 + 마일스톤)
  5: Before/After (AS-IS → TO-BE)
  6: Process Grid (번호 + 헤더 + 본문, 3×2 그리드)
  7: KPI Dashboard (4패널 + 번호 원형)
  8: SWOT 2×2
  9: Hub-Spoke Architecture

공통 요소:
  - 상단 헤더 바 (다크 or 악센트, h≈0.55")
  - 브레드크럼 (우상단)
  - 하단 푸터 (로고 + confidential)
  - 하단 시사점/결론 바
  - "Illustrative" 뱃지
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ============================================================
# Company Colors (정밀)
# ============================================================
CL_ACCENT = RGBColor(0xFD, 0x51, 0x08)
CL_ACCENT_MID = RGBColor(0xFE, 0x7C, 0x39)
CL_ACCENT_LIGHT = RGBColor(0xFF, 0xAA, 0x72)
CL_DARK = RGBColor(0x1A, 0x1A, 0x1A)
CL_BLACK = RGBColor(0x00, 0x00, 0x00)
CL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CL_GREY = RGBColor(0xA1, 0xA8, 0xB3)
CL_GREY_MID = RGBColor(0xB5, 0xBC, 0xC4)
CL_GREY_LIGHT = RGBColor(0xCB, 0xD1, 0xD6)
CL_BG = RGBColor(0xF2, 0xF2, 0xF2)
CL_BORDER = RGBColor(0xDD, 0xDD, 0xDE)
CL_RED_HIGHLIGHT = RGBColor(0xC0, 0x39, 0x2B)
FONT = "Arial"
FONT_T = "Georgia"

# ============================================================
# Helpers (정밀 배치)
# ============================================================
def _s(prs):
    """빈 슬라이드."""
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, x, y, w, h, fill=CL_WHITE, border=None, border_w=6350):
    """사각형."""
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Emu(border_w)
    else:
        s.line.fill.background()
    return s

def _text(slide, x, y, w, h, text, sz=9, bold=False, color=CL_BLACK,
          font=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """텍스트 박스."""
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tx

def _circle(slide, x, y, d, fill=CL_ACCENT, text="", text_c=CL_WHITE, sz=10):
    """원형."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = True
    p.font.color.rgb = text_c
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    return c

def _connector(slide, x1, y1, x2, y2, color=CL_GREY, width=9525):
    """직선 커넥터."""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    c.line.color.rgb = color
    c.line.width = Emu(width)
    return c

# ============================================================
# 공통: 헤더 + 푸터 + 브레드크럼 (모든 슬라이드)
# ============================================================
def _header(slide, title_text="{{title}}", breadcrumb="{{breadcrumb}}"):
    """완성본 스타일: 얇은 다크 헤더 바 + 오렌지 좌측 악센트 + 브레드크럼."""
    # 헤더 바 (다크, 더 얇게 - 완성본 기준 ~0.55")
    _rect(slide, 0, 0, 10, 0.55, CL_DARK)
    # 좌측 오렌지 악센트 바 (3px)
    _rect(slide, 0, 0, 0.04, 0.55, CL_ACCENT)
    # 타이틀
    _text(slide, 0.15, 0.08, 6.5, 0.4, title_text, 14, True, CL_WHITE, FONT_T)
    # 브레드크럼 (우상단)
    _text(slide, 6.8, 0.12, 3.0, 0.3, breadcrumb, 7, False, CL_GREY_LIGHT, FONT, PP_ALIGN.RIGHT)

def _footer(slide):
    """완성본 스타일: 하단 구분선 + confidential."""
    # 구분선
    _rect(slide, 0.3, 7.15, 9.4, 0.005, CL_BORDER)
    # Confidential 텍스트
    _text(slide, 0.3, 7.2, 5, 0.2, "Strictly Private and Confidential", 6, False, CL_GREY)

def _illustrative_badge(slide):
    """우상단 'Illustrative' 뱃지."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(0.65), Inches(1.2), Inches(0.22),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = CL_ACCENT_MID
    s.line.fill.background()
    try:
        s.adjustments[0] = 0.5
    except (IndexError, TypeError):
        pass
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Illustrative"
    p.font.size = Pt(7)
    p.font.bold = True
    p.font.color.rgb = CL_WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

def _takeaway(slide, text="{{takeaway}}", y=6.5):
    """하단 시사점/결론 바 (완성본의 가장 빈번한 패턴)."""
    _rect(slide, 0.3, y, 9.4, 0.5, CL_DARK)
    # 오렌지 좌측 마커
    _rect(slide, 0.3, y, 0.06, 0.5, CL_ACCENT)
    _text(slide, 0.5, y + 0.05, 9.0, 0.4, text, 9, True, CL_WHITE, FONT, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# Slide 0: 4-Column Reference (가장 빈번)
# ============================================================
def build_4col_reference(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 서브 헤더 (설명 텍스트)
    _text(slide, 0.3, 0.65, 9.4, 0.35, "{{subtitle}}", 9, False, CL_BLACK)

    # 4개 컬럼
    cols = [
        ("{{col1_header}}", CL_ACCENT),
        ("{{col2_header}}", CL_ACCENT_MID),
        ("{{col3_header}}", CL_GREY),
        ("{{col4_header}}", CL_DARK),
    ]
    col_w = 2.25
    gap = 0.1
    start_x = 0.3

    for i, (header, color) in enumerate(cols):
        cx = start_x + i * (col_w + gap)
        # 헤더 바
        _rect(slide, cx, 1.1, col_w, 0.4, color)
        _text(slide, cx, 1.12, col_w, 0.36, header, 9, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 본문 영역
        _rect(slide, cx, 1.5, col_w, 4.8, CL_WHITE, CL_BORDER)
        _text(slide, cx + 0.08, 1.6, col_w - 0.16, 4.6, f"{{{{col{i+1}_content}}}}", 8, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 1: Framework Matrix + Sidebar
# ============================================================
def build_framework_sidebar(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 좌측 사이드바 (프레임워크 네비게이션)
    sidebar_w = 2.8
    rows = ["{{nav_1}}", "{{nav_2}}", "{{nav_3}}", "{{nav_4}}"]
    for i, label in enumerate(rows):
        ry = 0.7 + i * 0.85
        fill = CL_ACCENT if i == 0 else CL_BG
        text_c = CL_WHITE if i == 0 else CL_BLACK
        border = None if i == 0 else CL_BORDER
        _rect(slide, 0.3, ry, sidebar_w, 0.75, fill, border)
        _text(slide, 0.4, ry + 0.05, sidebar_w - 0.2, 0.65, label, 8, True, text_c, anchor=MSO_ANCHOR.MIDDLE)

    # 우측 본문 영역
    _rect(slide, 3.3, 0.7, 6.4, 5.5, CL_WHITE, CL_BORDER)
    # 본문 헤더
    _rect(slide, 3.3, 0.7, 6.4, 0.4, CL_ACCENT)
    _text(slide, 3.5, 0.72, 6.0, 0.36, "{{content_header}}", 10, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 본문 내용
    _text(slide, 3.5, 1.2, 6.0, 4.8, "{{content_body}}", 9, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 2: 4-Step Decision Flow
# ============================================================
def build_decision_flow(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    _illustrative_badge(slide)

    # 서브 설명
    _text(slide, 0.3, 0.65, 9.4, 0.3, "{{subtitle}}", 9, False, CL_BLACK)

    # 4단계 (완성본: 번호원+아이콘+헤더+시스템박스 구조)
    steps = 4
    step_w = 2.15
    gap = 0.15
    sx = 0.3

    for i in range(steps):
        cx = sx + i * (step_w + gap)

        # 번호 원형
        _circle(slide, cx + 0.05, 1.1, 0.35, CL_ACCENT, f"{i+1:02d}", CL_WHITE, 11)

        # 단계 제목
        _text(slide, cx + 0.5, 1.1, step_w - 0.5, 0.35, f"{{{{step{i+1}_title}}}}", 10, True, CL_BLACK)

        # 내용 박스 (테두리 카드)
        _rect(slide, cx, 1.55, step_w, 3.2, CL_WHITE, CL_BORDER)
        _text(slide, cx + 0.08, 1.65, step_w - 0.16, 3.0, f"{{{{step{i+1}_content}}}}", 8, False, CL_BLACK)

        # 시스템 뱃지 (하단)
        _rect(slide, cx + 0.1, 4.5, step_w - 0.2, 0.3, CL_BG, CL_BORDER)
        _text(slide, cx + 0.1, 4.5, step_w - 0.2, 0.3, f"{{{{step{i+1}_system}}}}", 7, True, CL_GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 화살표 (마지막 제외)
        if i < steps - 1:
            ax = cx + step_w + 0.02
            _text(slide, ax - 0.02, 2.6, 0.18, 0.5, "▶", 14, False, CL_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 3: Comparison (2안 비교 + 추천 하이라이트)
# ============================================================
def build_comparison(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 검토 배경 (상단)
    _rect(slide, 0.3, 0.65, 9.4, 0.6, CL_BG, CL_BORDER)
    _text(slide, 0.4, 0.7, 9.2, 0.5, "{{background}}", 8, False, CL_BLACK)

    # Option A (좌측)
    _rect(slide, 0.3, 1.4, 4.5, 0.4, CL_DARK)
    _text(slide, 0.4, 1.42, 4.3, 0.36, "{{option_a_title}}", 10, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 0.3, 1.8, 4.5, 3.0, CL_WHITE, CL_BORDER)
    _text(slide, 0.4, 1.9, 4.3, 2.8, "{{option_a_content}}", 8, False, CL_BLACK)

    # Option B (우측, 추천 — 빨간 테두리)
    _rect(slide, 5.2, 1.4, 4.5, 0.4, CL_ACCENT)
    _text(slide, 5.3, 1.42, 4.3, 0.36, "{{option_b_title}}  ★ 추천", 10, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 5.2, 1.8, 4.5, 3.0, CL_WHITE, CL_RED_HIGHLIGHT, 19050)
    _text(slide, 5.3, 1.9, 4.3, 2.8, "{{option_b_content}}", 8, False, CL_BLACK)

    # 하단 결론 (2행)
    _rect(slide, 0.3, 5.0, 9.4, 0.4, CL_BG)
    _text(slide, 0.4, 5.02, 9.2, 0.36, "{{comparison_criteria}}", 8, True, CL_BLACK, anchor=MSO_ANCHOR.MIDDLE)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 4: Timeline Roadmap
# ============================================================
def build_timeline(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 타임라인 수평 바
    bar_y = 3.5
    _rect(slide, 0.5, bar_y, 9.0, 0.08, CL_GREY)

    # 5개 마일스톤
    for i in range(5):
        mx = 1.0 + i * 2.0
        # 마일스톤 원형 (타임라인 위)
        colors = [CL_ACCENT, CL_ACCENT, CL_ACCENT_MID, CL_GREY_MID, CL_GREY]
        _circle(slide, mx - 0.2, bar_y - 0.2, 0.45, colors[i], f"{{{{year_{i+1}}}}}", CL_WHITE if i < 3 else CL_BLACK, 8)

        # 콘텐츠 (상하 교대)
        if i % 2 == 0:
            by = bar_y - 2.3
            # 수직 연결선
            _connector(slide, mx, by + 1.6, mx, bar_y - 0.2, CL_BORDER)
        else:
            by = bar_y + 0.5
            _connector(slide, mx, bar_y + 0.25, mx, by, CL_BORDER)

        # 콘텐츠 카드
        _rect(slide, mx - 0.7, by, 1.6, 1.5, CL_WHITE, CL_BORDER)
        _text(slide, mx - 0.6, by + 0.05, 1.4, 0.25, f"{{{{milestone_{i+1}_title}}}}", 8, True, CL_ACCENT if i < 3 else CL_BLACK)
        _text(slide, mx - 0.6, by + 0.35, 1.4, 1.0, f"{{{{milestone_{i+1}_content}}}}", 7, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 5: Before/After
# ============================================================
def build_before_after(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # AS-IS (좌측)
    _rect(slide, 0.3, 0.7, 4.2, 0.4, CL_GREY)
    _text(slide, 0.4, 0.72, 4.0, 0.36, "AS-IS (현재)", 10, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 0.3, 1.1, 4.2, 3.8, CL_BG, CL_BORDER)
    _text(slide, 0.4, 1.2, 4.0, 3.6, "{{before_content}}", 8, False, CL_BLACK)

    # 화살표 (중앙)
    _circle(slide, 4.65, 2.6, 0.7, CL_ACCENT, "→", CL_WHITE, 20)

    # TO-BE (우측)
    _rect(slide, 5.5, 0.7, 4.2, 0.4, CL_ACCENT)
    _text(slide, 5.6, 0.72, 4.0, 0.36, "TO-BE (미래)", 10, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 5.5, 1.1, 4.2, 3.8, CL_WHITE, CL_ACCENT, 12700)
    _text(slide, 5.6, 1.2, 4.0, 3.6, "{{after_content}}", 8, False, CL_BLACK)

    # 기대효과 바
    _rect(slide, 5.5, 5.1, 4.2, 0.45, CL_ACCENT)
    _text(slide, 5.6, 5.12, 4.0, 0.41, "{{expected_effect}}", 8, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)

    _takeaway(slide, "{{takeaway}}", 5.8)


# ============================================================
# Slide 6: Process Grid (3×2, 번호+헤더+본문)
# ============================================================
def build_process_grid(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    rows, cols = 2, 3
    cw, ch = 3.0, 2.5
    gx, gy = 0.12, 0.12
    sx, sy = 0.3, 0.7

    colors = [CL_ACCENT, CL_ACCENT_MID, CL_ACCENT_LIGHT, CL_GREY_LIGHT, CL_GREY_MID, CL_GREY]

    for r in range(rows):
        for c in range(cols):
            idx = r * cols + c
            px = sx + c * (cw + gx)
            py = sy + r * (ch + gy)

            # 번호 바 (상단)
            _rect(slide, px, py, cw, 0.45, colors[idx])
            _text(slide, px + 0.1, py + 0.03, 0.6, 0.4, f"{idx:02d}", 18, True, CL_WHITE if idx < 3 else CL_BLACK, FONT_T)

            # 헤더
            _text(slide, px + 0.1, py + 0.55, cw - 0.2, 0.3, f"{{{{grid{idx}_header}}}}", 9, True, CL_BLACK)

            # 본문
            _rect(slide, px, py + 0.9, cw, ch - 0.9, CL_WHITE, CL_BORDER)
            _text(slide, px + 0.08, py + 0.95, cw - 0.16, ch - 1.1, f"{{{{grid{idx}_content}}}}", 7, False, CL_BLACK)


# ============================================================
# Slide 7: KPI Dashboard (4패널)
# ============================================================
def build_kpi_dashboard(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    panels = [
        (0.3, 0.75, "{{kpi1_value}}", "{{kpi1_label}}", "{{kpi1_detail}}"),
        (5.15, 0.75, "{{kpi2_value}}", "{{kpi2_label}}", "{{kpi2_detail}}"),
        (0.3, 3.85, "{{kpi3_value}}", "{{kpi3_label}}", "{{kpi3_detail}}"),
        (5.15, 3.85, "{{kpi4_value}}", "{{kpi4_label}}", "{{kpi4_detail}}"),
    ]
    pw, ph = 4.55, 2.8

    for i, (px, py, value, label, detail) in enumerate(panels):
        # 패널 배경
        _rect(slide, px, py, pw, ph, CL_WHITE, CL_BORDER)
        # 번호 원형
        _circle(slide, px + 0.15, py + 0.15, 0.35, CL_ACCENT if i < 2 else CL_GREY, f"{i+1:02d}", CL_WHITE, 10)
        # KPI 라벨
        _text(slide, px + 0.6, py + 0.18, pw - 0.8, 0.3, label, 9, True, CL_BLACK)
        # KPI 값 (대형)
        _text(slide, px + 0.3, py + 0.7, pw - 0.6, 0.8, value, 32, True, CL_ACCENT if i < 2 else CL_DARK, align=PP_ALIGN.CENTER)
        # 상세
        _text(slide, px + 0.2, py + 1.6, pw - 0.4, 1.0, detail, 8, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 8: SWOT 2×2
# ============================================================
def build_swot(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    quads = [
        (0.3, 0.7, CL_ACCENT, "{{q1_title}}", "{{q1_content}}"),
        (5.15, 0.7, CL_ACCENT_MID, "{{q2_title}}", "{{q2_content}}"),
        (0.3, 3.85, CL_GREY, "{{q3_title}}", "{{q3_content}}"),
        (5.15, 3.85, CL_DARK, "{{q4_title}}", "{{q4_content}}"),
    ]
    qw, qh = 4.55, 2.85

    for qx, qy, fill, title, content in quads:
        _rect(slide, qx, qy, qw, 0.4, fill)
        _text(slide, qx + 0.1, qy + 0.02, qw - 0.2, 0.36, title, 10, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        _rect(slide, qx, qy + 0.4, qw, qh - 0.4, CL_WHITE, CL_BORDER)
        _text(slide, qx + 0.1, qy + 0.5, qw - 0.2, qh - 0.6, content, 8, False, CL_BLACK)


# ============================================================
# Slide 9: Hub-Spoke Architecture
# ============================================================
def build_hub_spoke(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    _illustrative_badge(slide)

    # Central hub
    cx_c, cy_c = 5.0, 3.8
    _circle(slide, cx_c - 0.8, cy_c - 0.8, 1.6, CL_ACCENT, "{{hub_label}}", CL_WHITE, 11)

    # 6 spokes
    spoke_labels = ["{{spoke_1}}", "{{spoke_2}}", "{{spoke_3}}", "{{spoke_4}}", "{{spoke_5}}", "{{spoke_6}}"]
    radius = 2.3

    for i, label in enumerate(spoke_labels):
        angle = math.radians(60 * i - 90)
        nx = cx_c + radius * math.cos(angle)
        ny = cy_c + radius * math.sin(angle)

        # 스포크 박스
        bw, bh = 1.3, 0.55
        _rect(slide, nx - bw/2, ny - bh/2, bw, bh, CL_WHITE, CL_BORDER)
        _text(slide, nx - bw/2 + 0.05, ny - bh/2, bw - 0.1, bh, label, 8, True, CL_BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 연결선
        lx = cx_c + 0.8 * math.cos(angle)
        ly = cy_c + 0.8 * math.sin(angle)
        _connector(slide, lx, ly, nx, ny, CL_GREY_MID)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 10: Task/Image/Activity (3행 구조 — 완성본 최빈 패턴)
# ============================================================
def build_task_image_activity(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    rows = [
        ("Task", CL_DARK, CL_WHITE),
        ("Description", CL_BG, CL_BLACK),
        ("Activity", CL_WHITE, CL_BLACK),
    ]
    row_h = 1.7
    for ri, (label, fill, text_c) in enumerate(rows):
        ry = 0.7 + ri * (row_h + 0.08)
        # 행 라벨 (좌측)
        _rect(slide, 0.3, ry, 1.5, row_h, CL_ACCENT if ri == 0 else CL_GREY_LIGHT)
        _text(slide, 0.35, ry, 1.4, row_h, f"{{{{row{ri}_label}}}}", 9, True,
              CL_WHITE if ri == 0 else CL_BLACK, anchor=MSO_ANCHOR.MIDDLE)

        # 콘텐츠 셀 (3열)
        for ci in range(3):
            cx = 1.95 + ci * 2.55
            _rect(slide, cx, ry, 2.4, row_h, fill, CL_BORDER)
            _text(slide, cx + 0.08, ry + 0.08, 2.24, row_h - 0.16,
                  f"{{{{cell_{ri}_{ci}}}}}", 8, False, text_c)

    _takeaway(slide, "{{takeaway}}", 6.2)


# ============================================================
# Slide 11: Waterfall / Bridge (단계별 증감)
# ============================================================
def build_waterfall(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    n_steps = 6
    bar_w = 1.2
    gap = 0.25
    base_y = 5.5
    max_h = 3.5
    sx = 0.8

    labels = [f"{{{{step_{i}_label}}}}" for i in range(n_steps)]
    # 시뮬레이션 높이 (시작-증가-감소-증가-감소-최종)
    heights = [2.8, 0.8, -0.5, 1.0, -0.6, 3.5]
    cumulative = 0

    for i in range(n_steps):
        cx = sx + i * (bar_w + gap)

        if i == 0 or i == n_steps - 1:
            # 시작/최종 바
            bh = abs(heights[i]) / max_h * 3.0
            by = base_y - bh
            color = CL_DARK if i == 0 else CL_ACCENT
        elif heights[i] > 0:
            cumulative += heights[i]
            bh = heights[i] / max_h * 3.0
            by = base_y - (cumulative / max_h * 3.0)
            color = CL_ACCENT_MID
        else:
            bh = abs(heights[i]) / max_h * 3.0
            by = base_y - (cumulative / max_h * 3.0)
            cumulative += heights[i]
            color = CL_GREY

        _rect(slide, cx, by, bar_w, bh, color)
        _text(slide, cx, by - 0.25, bar_w, 0.25, f"{{{{step_{i}_value}}}}", 10, True,
              CL_BLACK, align=PP_ALIGN.CENTER)
        _text(slide, cx, base_y + 0.1, bar_w, 0.3, labels[i], 7, False,
              CL_BLACK, align=PP_ALIGN.CENTER)

    # 베이스 라인
    _rect(slide, 0.5, base_y, 9.0, 0.01, CL_BORDER)
    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 12: Center Diagram + Surrounding Text (중앙 포커스형)
# ============================================================
def build_center_focus(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 중앙 원
    _circle(slide, 3.7, 2.5, 2.6, CL_ACCENT, "{{center_label}}", CL_WHITE, 14)

    # 4방향 텍스트 블록
    positions = [
        (0.3, 1.0, 3.0, 1.2, "{{top_left}}"),      # 좌상
        (6.7, 1.0, 3.0, 1.2, "{{top_right}}"),      # 우상
        (0.3, 4.5, 3.0, 1.2, "{{bottom_left}}"),    # 좌하
        (6.7, 4.5, 3.0, 1.2, "{{bottom_right}}"),   # 우하
    ]
    for i, (px, py, pw, ph, placeholder) in enumerate(positions):
        # 카드 상단 바
        bar_c = CL_ACCENT if i < 2 else CL_GREY
        _rect(slide, px, py, pw, 0.05, bar_c)
        _rect(slide, px, py + 0.05, pw, ph - 0.05, CL_WHITE, CL_BORDER)
        _text(slide, px + 0.1, py + 0.15, pw - 0.2, 0.25, f"{{{{label_{i}}}}}", 9, True, CL_BLACK)
        _text(slide, px + 0.1, py + 0.45, pw - 0.2, ph - 0.55, placeholder, 8, False, CL_BLACK)

        # 연결선 (중앙으로)
        lx = 5.0 if px < 5 else 5.0
        ly = 3.8
        _connector(slide, px + pw/2, py + ph/2, lx, ly, CL_GREY_LIGHT)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 13: Dense Data Table (조건부 색상 — 완성본 빈출)
# ============================================================
def build_dense_table(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    _illustrative_badge(slide)

    # 8행 5열 밀도 높은 테이블
    n_rows, n_cols = 8, 5
    col_w = 1.8
    row_h = 0.55
    start_x, start_y = 0.3, 0.75

    # 열 헤더
    headers = [f"{{{{col{c}_header}}}}" for c in range(n_cols)]
    for c in range(n_cols):
        cx = start_x + c * col_w + (0.5 if c == 0 else 0)
        cw = col_w + 0.5 if c == 0 else col_w
        _rect(slide, cx if c > 0 else start_x, start_y, cw if c > 0 else col_w + 0.5, 0.4, CL_DARK)
        _text(slide, cx if c > 0 else start_x + 0.05, start_y, cw if c > 0 else col_w + 0.45, 0.4,
              headers[c], 8, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 데이터 행
    for r in range(n_rows):
        ry = start_y + 0.4 + r * row_h
        fill = CL_WHITE if r % 2 == 0 else CL_BG
        for c in range(n_cols):
            cx = start_x + c * col_w + (0.5 if c == 0 else 0)
            cw = col_w + 0.5 if c == 0 else col_w
            cell_fill = fill
            if c == 0:
                cell_fill = CL_GREY_LIGHT
            _rect(slide, cx if c > 0 else start_x, ry, cw if c > 0 else col_w + 0.5, row_h, cell_fill, CL_BORDER)
            _text(slide, (cx if c > 0 else start_x) + 0.05, ry, (cw if c > 0 else col_w + 0.45), row_h,
                  f"{{{{r{r}_c{c}}}}}", 7, c == 0, CL_BLACK, anchor=MSO_ANCHOR.MIDDLE)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 14: Two-Panel Deep Dive (좌:분석 우:시각화)
# ============================================================
def build_two_panel(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 좌측 패널: 텍스트 분석
    _rect(slide, 0.3, 0.7, 4.5, 0.35, CL_DARK)
    _text(slide, 0.4, 0.72, 4.3, 0.31, "{{left_header}}", 9, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # 좌측 불릿 영역
    _rect(slide, 0.3, 1.05, 4.5, 4.6, CL_WHITE, CL_BORDER)
    for i in range(6):
        iy = 1.15 + i * 0.72
        # 번호 원
        _circle(slide, 0.45, iy, 0.3, CL_ACCENT if i < 3 else CL_GREY, f"{i+1:02d}", CL_WHITE, 8)
        # 텍스트
        _text(slide, 0.85, iy, 3.8, 0.65, f"{{{{left_item_{i}}}}}", 8, False, CL_BLACK)

    # 우측 패널: 시각화
    _rect(slide, 5.2, 0.7, 4.5, 0.35, CL_ACCENT)
    _text(slide, 5.3, 0.72, 4.3, 0.31, "{{right_header}}", 9, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 5.2, 1.05, 4.5, 4.6, CL_BG, CL_BORDER)
    _text(slide, 5.4, 1.2, 4.1, 4.3, "{{right_content}}", 8, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Main
# ============================================================
def build_template_library(output_path: Path = None):
    if output_path is None:
        output_path = Path(__file__).parent.parent.parent / "templates" / "template_library.pptx"

    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000

    build_4col_reference(prs)      # 0
    build_framework_sidebar(prs)   # 1
    build_decision_flow(prs)       # 2
    build_comparison(prs)          # 3
    build_timeline(prs)            # 4
    build_before_after(prs)        # 5
    build_process_grid(prs)        # 6
    build_kpi_dashboard(prs)       # 7
    build_swot(prs)                # 8
    build_hub_spoke(prs)           # 9
    build_task_image_activity(prs) # 10
    build_waterfall(prs)           # 11
    build_center_focus(prs)        # 12
    build_dense_table(prs)         # 13
    build_two_panel(prs)           # 14
    build_raci(prs)                # 15 TIER1
    build_swimlane(prs)            # 16 TIER1
    build_pestel(prs)              # 17 TIER1
    build_scr(prs)                 # 18 TIER1
    build_left_right_split(prs)    # 19 TIER1
    build_porter_five_forces(prs)  # 20 TIER2
    build_value_chain(prs)         # 21 TIER2
    build_bcg_matrix(prs)          # 22 TIER2
    build_org_chart(prs)           # 23 TIER2
    build_gantt_roadmap(prs)       # 24 TIER2
    build_prioritization_2x2(prs)  # 25 TIER2
    build_tornado(prs)             # 26 TIER2
    build_decision_tree(prs)       # 27 TIER2
    build_revenue_tree(prs)        # 28 TIER2
    build_three_option(prs)        # 29 TIER2
    build_circular_loop(prs)       # 30 TIER3
    build_mckinsey_7s(prs)         # 31 TIER3
    build_three_horizons(prs)      # 32 TIER3
    build_mekko(prs)               # 33 TIER3
    build_table_with_bars(prs)     # 34 TIER3

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Template library built: {output_path} ({len(prs.slides)} slides)")
    return output_path


# ============================================================
# Slide 15: RACI Matrix
# ============================================================
def build_raci(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    n_rows, n_cols = 6, 5  # 6 tasks × 5 roles
    row_h = 0.7
    col_w = 1.5
    label_w = 2.5
    sx, sy = 0.3, 0.75

    raci_colors = {"R": CL_ACCENT, "A": CL_ACCENT_MID, "C": CL_GREY_LIGHT, "I": CL_WHITE}

    # 열 헤더 (역할)
    for c in range(n_cols):
        cx = sx + label_w + c * col_w
        _rect(slide, cx, sy, col_w, 0.4, CL_DARK)
        _text(slide, cx, sy, col_w, 0.4, f"{{{{role_{c}}}}}", 8, True, CL_WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 행 라벨 헤더
    _rect(slide, sx, sy, label_w, 0.4, CL_DARK)
    _text(slide, sx + 0.05, sy, label_w - 0.1, 0.4, "Activity / Task", 8, True, CL_WHITE,
          anchor=MSO_ANCHOR.MIDDLE)

    # 데이터 행
    for r in range(n_rows):
        ry = sy + 0.4 + r * row_h
        fill = CL_WHITE if r % 2 == 0 else CL_BG
        _rect(slide, sx, ry, label_w, row_h, fill, CL_BORDER)
        _text(slide, sx + 0.05, ry, label_w - 0.1, row_h, f"{{{{task_{r}}}}}", 8, True, CL_BLACK,
              anchor=MSO_ANCHOR.MIDDLE)
        for c in range(n_cols):
            cx = sx + label_w + c * col_w
            _rect(slide, cx, ry, col_w, row_h, fill, CL_BORDER)
            _text(slide, cx, ry, col_w, row_h, f"{{{{raci_{r}_{c}}}}}", 12, True, CL_ACCENT,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 16: Swimlane Process
# ============================================================
def build_swimlane(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    lanes = 3
    steps = 5
    lane_h = 1.6
    lane_label_w = 1.5
    step_w = 1.4
    gap = 0.08
    sx, sy = 0.3, 0.75

    for li in range(lanes):
        ly = sy + li * (lane_h + gap)
        fill = CL_ACCENT if li == 0 else CL_GREY_LIGHT if li == 1 else CL_BG
        text_c = CL_WHITE if li == 0 else CL_BLACK
        _rect(slide, sx, ly, lane_label_w, lane_h, fill)
        _text(slide, sx + 0.05, ly, lane_label_w - 0.1, lane_h, f"{{{{lane_{li}}}}}", 9, True, text_c,
              anchor=MSO_ANCHOR.MIDDLE)

        for si in range(steps):
            step_x = sx + lane_label_w + 0.1 + si * (step_w + 0.08)
            step_fill = CL_WHITE
            _rect(slide, step_x, ly + 0.15, step_w, lane_h - 0.3, step_fill, CL_BORDER)
            _text(slide, step_x + 0.05, ly + 0.15, step_w - 0.1, lane_h - 0.3,
                  f"{{{{step_{li}_{si}}}}}", 7, False, CL_BLACK)

    # 화살표 (상단 레인)
    for si in range(steps - 1):
        ax = sx + lane_label_w + 0.1 + si * (step_w + 0.08) + step_w
        _text(slide, ax, sy + 0.6, 0.08, 0.3, "→", 8, False, CL_ACCENT)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 17: PESTEL Analysis
# ============================================================
def build_pestel(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    factors = [
        ("P", "Political", CL_ACCENT),
        ("E", "Economic", CL_ACCENT_MID),
        ("S", "Social", CL_ACCENT_LIGHT),
        ("T", "Technological", CL_GREY),
        ("E", "Environmental", CL_GREY_MID),
        ("L", "Legal", CL_GREY_LIGHT),
    ]
    cols, rows = 3, 2
    cw, ch = 3.0, 2.7
    gap = 0.12
    sx, sy = 0.3, 0.75

    for idx, (letter, label, color) in enumerate(factors):
        r, c = divmod(idx, cols)
        px = sx + c * (cw + gap)
        py = sy + r * (ch + gap)

        # 상단 바
        _rect(slide, px, py, cw, 0.06, color)
        # 본문
        _rect(slide, px, py + 0.06, cw, ch - 0.06, CL_WHITE, CL_BORDER)
        # 레터 뱃지
        _circle(slide, px + 0.1, py + 0.15, 0.35, color, letter, CL_WHITE, 14)
        # 라벨
        _text(slide, px + 0.55, py + 0.15, cw - 0.65, 0.35, f"{{{{pestel_{idx}_title}}}}", 10, True, CL_BLACK)
        # 불릿 영역
        _text(slide, px + 0.1, py + 0.6, cw - 0.2, ch - 0.7, f"{{{{pestel_{idx}_content}}}}", 8, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}", 6.3)


# ============================================================
# Slide 18: SCR (Situation-Complication-Resolution)
# ============================================================
def build_scr(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    sections = [
        ("S", "Situation", CL_BG, CL_BLACK, 0.75),
        ("C", "Complication", RGBColor(0xFF, 0xF3, 0xEB), CL_ACCENT, 0.75),
        ("R", "Resolution", CL_WHITE, CL_BLACK, 3.5),
    ]
    sy = 0.75
    for letter, label, fill, text_c, height in sections:
        # 좌측 라벨
        label_fill = CL_ACCENT if letter == "R" else CL_GREY if letter == "C" else CL_DARK
        _rect(slide, 0.3, sy, 0.8, height, label_fill)
        _text(slide, 0.3, sy, 0.8, height, letter, 20, True, CL_WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 본문 영역
        _rect(slide, 1.15, sy, 8.55, height, fill, CL_BORDER)
        _text(slide, 1.25, sy + 0.05, 1.5, 0.3, label, 9, True, text_c)
        _text(slide, 1.25, sy + 0.35, 8.3, height - 0.45, f"{{{{{letter.lower()}_content}}}}", 9, False, CL_BLACK)

        sy += height + 0.08

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# Slide 19: Left-Right Split (개선 — 좌30% 맥락 + 우70% 상세)
# ============================================================
def build_left_right_split(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 좌측 패널 (Key Considerations)
    left_w = 2.8
    _rect(slide, 0.3, 0.75, left_w, 0.35, CL_DARK)
    _text(slide, 0.4, 0.77, left_w - 0.2, 0.31, "{{left_header}}", 9, True, CL_WHITE,
          anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 0.3, 1.1, left_w, 4.8, CL_BG, CL_BORDER)
    for i in range(5):
        iy = 1.2 + i * 0.9
        _circle(slide, 0.45, iy, 0.28, CL_ACCENT if i < 2 else CL_GREY, f"{i+1:02d}", CL_WHITE, 8)
        _text(slide, 0.85, iy, left_w - 0.65, 0.8, f"{{{{left_item_{i}}}}}", 8, False, CL_BLACK)

    # 우측 패널 (Main Content)
    right_x = 3.3
    right_w = 6.4
    _rect(slide, right_x, 0.75, right_w, 0.35, CL_ACCENT)
    _text(slide, right_x + 0.1, 0.77, right_w - 0.2, 0.31, "{{right_header}}", 9, True, CL_WHITE,
          anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, right_x, 1.1, right_w, 4.8, CL_WHITE, CL_BORDER)
    _text(slide, right_x + 0.1, 1.2, right_w - 0.2, 4.6, "{{right_content}}", 8, False, CL_BLACK)

    # 하단 번호 스트립
    for i in range(3):
        nx = 0.3 + i * 3.2
        _circle(slide, nx, 6.1, 0.3, CL_ACCENT, f"{i+1:02d}", CL_WHITE, 9)
        _text(slide, nx + 0.4, 6.1, 2.6, 0.3, f"{{{{bottom_item_{i}}}}}", 8, True, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# ============================================================
# TIER 2 TEMPLATES (10종)
# ============================================================

# Slide 20: Porter's Five Forces
def build_porter_five_forces(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    # 중앙: Industry Rivalry
    _rect(slide, 3.5, 2.8, 3.0, 1.2, CL_ACCENT)
    _text(slide, 3.6, 2.85, 2.8, 1.1, "{{center}}\nIndustry Rivalry", 10, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 4 방향 Force 박스
    forces = [
        (3.5, 0.85, "{{force_top}}", "Threat of\nNew Entrants"),    # 상
        (3.5, 4.5, "{{force_bottom}}", "Threat of\nSubstitutes"),    # 하
        (0.3, 2.8, "{{force_left}}", "Supplier\nPower"),             # 좌
        (6.7, 2.8, "{{force_right}}", "Buyer\nPower"),               # 우
    ]
    for fx, fy, placeholder, label in forces:
        _rect(slide, fx, fy, 3.0, 1.2, CL_WHITE, CL_BORDER)
        # 상단 바
        _rect(slide, fx, fy, 3.0, 0.05, CL_ACCENT_MID)
        _text(slide, fx + 0.1, fy + 0.1, 2.8, 0.35, label, 8, True, CL_BLACK, align=PP_ALIGN.CENTER)
        _text(slide, fx + 0.1, fy + 0.5, 2.8, 0.6, placeholder, 7, False, CL_BLACK)
    # 화살표 (텍스트)
    _text(slide, 4.7, 2.2, 0.5, 0.5, "▲", 14, False, CL_GREY, align=PP_ALIGN.CENTER)
    _text(slide, 4.7, 4.1, 0.5, 0.5, "▼", 14, False, CL_GREY, align=PP_ALIGN.CENTER)
    _text(slide, 2.8, 3.2, 0.5, 0.5, "◀", 14, False, CL_GREY, align=PP_ALIGN.CENTER)
    _text(slide, 6.5, 3.2, 0.5, 0.5, "▶", 14, False, CL_GREY, align=PP_ALIGN.CENTER)
    _takeaway(slide, "{{takeaway}}")


# Slide 21: Value Chain
def build_value_chain(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    # Primary Activities (쉐브론 밴드)
    activities = ["{{primary_1}}", "{{primary_2}}", "{{primary_3}}", "{{primary_4}}", "{{primary_5}}"]
    aw = 1.7
    for i, act in enumerate(activities):
        ax = 0.3 + i * (aw + 0.05)
        colors = [CL_ACCENT, CL_ACCENT_MID, CL_ACCENT_LIGHT, CL_GREY_LIGHT, CL_GREY]
        _rect(slide, ax, 0.75, aw, 2.0, colors[i])
        _text(slide, ax + 0.1, 0.85, aw - 0.2, 0.4, act, 8, True, CL_WHITE if i < 3 else CL_BLACK, align=PP_ALIGN.CENTER)
        _text(slide, ax + 0.1, 1.3, aw - 0.2, 1.3, f"{{{{detail_{i}}}}}", 7, False, CL_WHITE if i < 2 else CL_BLACK)
    # Margin 쐐기
    _rect(slide, 8.8, 0.75, 0.9, 2.0, CL_DARK)
    _text(slide, 8.85, 0.75, 0.8, 2.0, "Margin", 9, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Support Activities (하단 행)
    supports = ["{{support_1}}", "{{support_2}}", "{{support_3}}", "{{support_4}}"]
    sw = 2.3
    for i, sup in enumerate(supports):
        sx = 0.3 + i * (sw + 0.08)
        _rect(slide, sx, 3.0, sw, 0.8, CL_BG, CL_BORDER)
        _text(slide, sx + 0.1, 3.05, sw - 0.2, 0.7, sup, 8, False, CL_BLACK)
    _text(slide, 0.3, 3.85, 9.4, 0.25, "Support Activities: Infrastructure, HR, Technology, Procurement", 7, False, CL_GREY)
    _takeaway(slide, "{{takeaway}}", 4.3)


# Slide 22: BCG Growth-Share Matrix
def build_bcg_matrix(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    # 2×2 그리드
    qw, qh = 4.3, 2.6
    quads = [
        (0.3, 0.75, "★ Stars", CL_ACCENT, "{{stars}}"),
        (4.7, 0.75, "? Question Marks", CL_ACCENT_LIGHT, "{{questions}}"),
        (0.3, 3.5, "💰 Cash Cows", CL_GREY_LIGHT, "{{cashcows}}"),
        (4.7, 3.5, "🐕 Dogs", CL_BG, "{{dogs}}"),
    ]
    for qx, qy, label, fill, placeholder in quads:
        _rect(slide, qx, qy, qw, 0.4, fill)
        text_c = CL_WHITE if fill in (CL_ACCENT,) else CL_BLACK
        _text(slide, qx + 0.1, qy + 0.02, qw - 0.2, 0.36, label, 10, True, text_c, anchor=MSO_ANCHOR.MIDDLE)
        _rect(slide, qx, qy + 0.4, qw, qh - 0.4, CL_WHITE, CL_BORDER)
        _text(slide, qx + 0.1, qy + 0.5, qw - 0.2, qh - 0.6, placeholder, 8, False, CL_BLACK)
    # 축 라벨
    _text(slide, 0.3, 6.3, 9.0, 0.25, "← High Market Share                    Low Market Share →", 7, False, CL_GREY, align=PP_ALIGN.CENTER)
    _takeaway(slide, "{{takeaway}}")


# Slide 23: Org Chart (3-Level)
def build_org_chart(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    # Level 1 (CEO)
    _rect(slide, 3.5, 0.85, 3.0, 0.65, CL_ACCENT)
    _text(slide, 3.6, 0.87, 2.8, 0.61, "{{level1}}", 10, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 수직선
    _rect(slide, 4.95, 1.5, 0.01, 0.4, CL_GREY)
    # 수평선
    _rect(slide, 1.0, 1.9, 7.9, 0.01, CL_GREY)
    # Level 2 (4 Direct Reports)
    l2_labels = ["{{l2_1}}", "{{l2_2}}", "{{l2_3}}", "{{l2_4}}"]
    for i, lbl in enumerate(l2_labels):
        lx = 0.5 + i * 2.4
        _rect(slide, lx, 2.0, 2.1, 0.55, CL_DARK)
        _text(slide, lx + 0.05, 2.02, 2.0, 0.51, lbl, 8, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 수직선
        _rect(slide, lx + 1.0, 2.55, 0.01, 0.3, CL_GREY)
    # Level 3 (하위 — 각 L2 아래 2개)
    for i in range(4):
        lx = 0.5 + i * 2.4
        for j in range(2):
            bx = lx + j * 1.1
            _rect(slide, bx, 2.95, 1.0, 0.45, CL_BG, CL_BORDER)
            _text(slide, bx + 0.05, 2.97, 0.9, 0.41, f"{{{{l3_{i}_{j}}}}}", 7, False, CL_BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _takeaway(slide, "{{takeaway}}", 3.7)


# Slide 24: Gantt Roadmap (개선 — 트랙별 바 + 마일스톤)
def build_gantt_roadmap(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    # 헤더: 시간축
    phases = ["{{phase_1}}", "{{phase_2}}", "{{phase_3}}", "{{phase_4}}", "{{phase_5}}"]
    pw = 1.7
    for i, ph in enumerate(phases):
        px = 2.0 + i * (pw + 0.05)
        _rect(slide, px, 0.75, pw, 0.35, CL_DARK if i < 2 else CL_GREY)
        _text(slide, px, 0.77, pw, 0.31, ph, 8, True, CL_WHITE if i < 2 else CL_BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 트랙 행
    tracks = ["{{track_1}}", "{{track_2}}", "{{track_3}}", "{{track_4}}", "{{track_5}}"]
    bar_specs = [(0, 2), (1, 3), (0, 3), (2, 4), (3, 4)]  # (start_phase, end_phase)
    th = 0.65
    for ti, (track, (sp, ep)) in enumerate(zip(tracks, bar_specs)):
        ty = 1.2 + ti * (th + 0.08)
        # 트랙 라벨
        _rect(slide, 0.3, ty, 1.6, th, CL_BG, CL_BORDER)
        _text(slide, 0.35, ty, 1.5, th, track, 8, True, CL_BLACK, anchor=MSO_ANCHOR.MIDDLE)
        # Gantt 바
        bx = 2.0 + sp * (pw + 0.05)
        bw = (ep - sp) * (pw + 0.05) + pw
        colors = [CL_ACCENT, CL_ACCENT_MID, CL_ACCENT_LIGHT, CL_GREY_MID, CL_GREY]
        _rect(slide, bx, ty + 0.1, bw, th - 0.2, colors[ti])
        _text(slide, bx + 0.1, ty + 0.1, bw - 0.2, th - 0.2, f"{{{{bar_{ti}}}}}", 7, True,
              CL_WHITE if ti < 3 else CL_BLACK, anchor=MSO_ANCHOR.MIDDLE)
    # 마일스톤 (다이아몬드 대신 원형)
    for i in range(3):
        mx = 3.7 + i * 2.8
        _circle(slide, mx, 5.1, 0.25, CL_ACCENT, "◆", CL_WHITE, 8)
        _text(slide, mx - 0.3, 5.4, 0.85, 0.25, f"{{{{ms_{i}}}}}", 7, True, CL_ACCENT, align=PP_ALIGN.CENTER)
    _takeaway(slide, "{{takeaway}}", 5.8)


# Slide 25: 2×2 Prioritization Matrix
def build_prioritization_2x2(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    qw, qh = 4.3, 2.5
    quads = [
        (0.3, 0.85, "Quick Wins", CL_ACCENT, "{{q1}}", "High Impact / Low Effort"),
        (4.7, 0.85, "Big Bets", CL_ACCENT_MID, "{{q2}}", "High Impact / High Effort"),
        (0.3, 3.55, "Fill-ins", CL_GREY_LIGHT, "{{q3}}", "Low Impact / Low Effort"),
        (4.7, 3.55, "Deprioritize", CL_BG, "{{q4}}", "Low Impact / High Effort"),
    ]
    for qx, qy, label, fill, placeholder, desc in quads:
        _rect(slide, qx, qy, qw, 0.35, fill)
        text_c = CL_WHITE if fill in (CL_ACCENT, CL_ACCENT_MID) else CL_BLACK
        _text(slide, qx + 0.1, qy + 0.01, qw - 0.2, 0.33, f"{label} — {desc}", 8, True, text_c, anchor=MSO_ANCHOR.MIDDLE)
        _rect(slide, qx, qy + 0.35, qw, qh - 0.35, CL_WHITE, CL_BORDER)
        _text(slide, qx + 0.1, qy + 0.45, qw - 0.2, qh - 0.55, placeholder, 8, False, CL_BLACK)
    _text(slide, 0.3, 6.3, 4.3, 0.2, "← Low Effort        High Effort →", 7, False, CL_GREY, align=PP_ALIGN.CENTER)
    _text(slide, 4.7, 6.3, 4.3, 0.2, "← Low Effort        High Effort →", 7, False, CL_GREY, align=PP_ALIGN.CENTER)
    _takeaway(slide, "{{takeaway}}")


# Slide 26: Tornado Chart (Sensitivity)
def build_tornado(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    n_vars = 6
    bar_h = 0.55
    center_x = 5.0
    max_bar_w = 3.5
    sy = 0.85
    for i in range(n_vars):
        vy = sy + i * (bar_h + 0.1)
        # 변수명 (좌측)
        _text(slide, 0.3, vy, 2.0, bar_h, f"{{{{var_{i}}}}}", 8, True, CL_BLACK, anchor=MSO_ANCHOR.MIDDLE)
        # 음수 바 (좌로)
        neg_w = max_bar_w * (1 - i * 0.15)
        _rect(slide, center_x - neg_w, vy + 0.05, neg_w, bar_h - 0.1, CL_GREY)
        _text(slide, center_x - neg_w, vy + 0.05, neg_w, bar_h - 0.1, f"{{{{neg_{i}}}}}", 7, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 양수 바 (우로)
        pos_w = max_bar_w * (1 - i * 0.12)
        _rect(slide, center_x, vy + 0.05, pos_w, bar_h - 0.1, CL_ACCENT)
        _text(slide, center_x, vy + 0.05, pos_w, bar_h - 0.1, f"{{{{pos_{i}}}}}", 7, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 중앙선
    _rect(slide, center_x - 0.005, sy, 0.01, n_vars * (bar_h + 0.1), CL_DARK)
    _text(slide, center_x - 1.0, sy + n_vars * (bar_h + 0.1) + 0.1, 2.0, 0.2, "Base Case", 7, True, CL_BLACK, align=PP_ALIGN.CENTER)
    _takeaway(slide, "{{takeaway}}")


# Slide 27: Decision Tree
def build_decision_tree(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    # Root
    _rect(slide, 0.3, 2.5, 2.0, 0.8, CL_DARK)
    _text(slide, 0.4, 2.55, 1.8, 0.7, "{{root}}", 9, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Branch lines
    _connector(slide, 2.3, 2.9, 3.0, 1.8, CL_GREY)
    _connector(slide, 2.3, 2.9, 3.0, 4.0, CL_GREY)
    # L1 nodes
    _rect(slide, 3.0, 1.4, 2.0, 0.8, CL_ACCENT)
    _text(slide, 3.1, 1.45, 1.8, 0.7, "{{branch_a}}", 8, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 3.0, 3.6, 2.0, 0.8, CL_GREY)
    _text(slide, 3.1, 3.65, 1.8, 0.7, "{{branch_b}}", 8, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Branch A → 2 outcomes
    _connector(slide, 5.0, 1.8, 5.7, 1.2, CL_GREY)
    _connector(slide, 5.0, 1.8, 5.7, 2.4, CL_GREY)
    _rect(slide, 5.7, 0.8, 2.0, 0.7, CL_ACCENT_MID)
    _text(slide, 5.8, 0.85, 1.8, 0.6, "{{outcome_a1}}", 7, False, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 5.7, 2.1, 2.0, 0.7, CL_BG, CL_BORDER)
    _text(slide, 5.8, 2.15, 1.8, 0.6, "{{outcome_a2}}", 7, False, CL_BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Branch B → 2 outcomes
    _connector(slide, 5.0, 4.0, 5.7, 3.4, CL_GREY)
    _connector(slide, 5.0, 4.0, 5.7, 4.6, CL_GREY)
    _rect(slide, 5.7, 3.0, 2.0, 0.7, CL_BG, CL_BORDER)
    _text(slide, 5.8, 3.05, 1.8, 0.6, "{{outcome_b1}}", 7, False, CL_BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 5.7, 4.2, 2.0, 0.7, CL_BG, CL_BORDER)
    _text(slide, 5.8, 4.25, 1.8, 0.6, "{{outcome_b2}}", 7, False, CL_BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 결과 라벨
    for i, (oy, placeholder) in enumerate([(0.8, "{{value_a1}}"), (2.1, "{{value_a2}}"), (3.0, "{{value_b1}}"), (4.2, "{{value_b2}}")]):
        _text(slide, 7.8, oy + 0.1, 1.9, 0.5, placeholder, 8, True, CL_ACCENT if i == 0 else CL_BLACK)
    _takeaway(slide, "{{takeaway}}", 5.3)


# Slide 28: Revenue Decomposition Tree
def build_revenue_tree(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    # Root (Total Revenue)
    _rect(slide, 0.3, 2.3, 2.0, 1.0, CL_ACCENT)
    _text(slide, 0.4, 2.35, 1.8, 0.9, "{{total}}\nTotal Revenue", 9, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # L1: Price × Volume
    _connector(slide, 2.3, 2.8, 3.0, 1.6, CL_GREY)
    _connector(slide, 2.3, 2.8, 3.0, 4.0, CL_GREY)
    _rect(slide, 3.0, 1.2, 1.8, 0.8, CL_ACCENT_MID)
    _text(slide, 3.05, 1.25, 1.7, 0.7, "{{l1_price}}\nPrice", 8, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, 3.0, 3.6, 1.8, 0.8, CL_ACCENT_MID)
    _text(slide, 3.05, 3.65, 1.7, 0.7, "{{l1_volume}}\nVolume", 8, True, CL_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 연산자
    _text(slide, 3.5, 2.3, 0.8, 0.5, "×", 20, True, CL_GREY, align=PP_ALIGN.CENTER)
    # L2: Volume branches
    _connector(slide, 4.8, 4.0, 5.5, 3.0, CL_GREY)
    _connector(slide, 4.8, 4.0, 5.5, 5.0, CL_GREY)
    segments = [("{{seg_1}}", 2.6), ("{{seg_2}}", 3.6), ("{{seg_3}}", 4.6)]
    for label, sy in segments:
        _rect(slide, 5.5, sy, 1.8, 0.7, CL_GREY_LIGHT, CL_BORDER)
        _text(slide, 5.6, sy + 0.05, 1.6, 0.6, label, 8, False, CL_BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Growth annotations
    for i, sy in enumerate([2.6, 3.6, 4.6]):
        _text(slide, 7.5, sy + 0.1, 2.2, 0.5, f"{{{{growth_{i}}}}}", 8, False, CL_ACCENT if i == 0 else CL_BLACK)
    _takeaway(slide, "{{takeaway}}", 5.7)


# Slide 29: Three-Option Compare
def build_three_option(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    ow = 3.0
    gap = 0.1
    for i in range(3):
        ox = 0.3 + i * (ow + gap)
        is_selected = (i == 1)  # 가운데 추천
        header_fill = CL_ACCENT if is_selected else CL_GREY
        border = CL_ACCENT if is_selected else CL_BORDER
        border_w = 19050 if is_selected else 6350
        # 헤더
        _rect(slide, ox, 0.75, ow, 0.4, header_fill)
        suffix = "  ★ 추천" if is_selected else ""
        _text(slide, ox + 0.1, 0.77, ow - 0.2, 0.36, f"{{{{opt_{i}_title}}}}{suffix}", 9, True, CL_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # 본문
        _rect(slide, ox, 1.15, ow, 4.0, CL_WHITE, border, border_w)
        _text(slide, ox + 0.1, 1.25, ow - 0.2, 3.8, f"{{{{opt_{i}_content}}}}", 8, False, CL_BLACK)
        # 평가 뱃지 하단
        eval_fill = CL_ACCENT if is_selected else CL_BG
        eval_tc = CL_WHITE if is_selected else CL_BLACK
        _rect(slide, ox, 5.25, ow, 0.4, eval_fill)
        _text(slide, ox + 0.1, 5.27, ow - 0.2, 0.36, f"{{{{opt_{i}_eval}}}}", 8, True, eval_tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _takeaway(slide, "{{takeaway}}", 5.9)


# ============================================================
# TIER 3 TEMPLATES (5종 — 고난이도)
# ============================================================

# Slide 30: Circular/Loop Diagram (4-node cycle)
def build_circular_loop(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 4개 노드를 원형 배치 + 화살표
    cx_c, cy_c = 5.0, 3.3
    radius = 2.0
    nodes = [
        ("{{node_1}}", 0),    # 상 (12시)
        ("{{node_2}}", 90),   # 우 (3시)
        ("{{node_3}}", 180),  # 하 (6시)
        ("{{node_4}}", 270),  # 좌 (9시)
    ]
    node_colors = [CL_ACCENT, CL_ACCENT_MID, CL_GREY, CL_GREY_MID]

    for i, (label, angle_deg) in enumerate(nodes):
        angle = math.radians(angle_deg - 90)  # 12시 시작
        nx = cx_c + radius * math.cos(angle) - 0.7
        ny = cy_c + radius * math.sin(angle) - 0.35

        _rect(slide, nx, ny, 1.4, 0.7, node_colors[i])
        _text(slide, nx + 0.05, ny + 0.05, 1.3, 0.6, label, 9, True,
              CL_WHITE if i < 2 else CL_BLACK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 화살표 (노드 간 — 텍스트로 표현)
    arrows = [
        (cx_c + 0.8, cy_c - 1.5, "→"),   # 상→우
        (cx_c + 0.8, cy_c + 0.8, "↓"),   # 우→하
        (cx_c - 1.0, cy_c + 0.8, "←"),   # 하→좌
        (cx_c - 1.0, cy_c - 1.5, "↑"),   # 좌→상
    ]
    for ax, ay, sym in arrows:
        _text(slide, ax, ay, 0.4, 0.4, sym, 16, True, CL_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 중앙 라벨
    _circle(slide, cx_c - 0.6, cy_c - 0.6, 1.2, CL_WHITE, "{{center}}", CL_ACCENT, 10)

    # 4방향 상세 텍스트
    details = [
        (0.3, 0.85, 2.0, 1.0, "{{detail_1}}"),
        (7.5, 0.85, 2.2, 1.0, "{{detail_2}}"),
        (7.5, 4.5, 2.2, 1.0, "{{detail_3}}"),
        (0.3, 4.5, 2.0, 1.0, "{{detail_4}}"),
    ]
    for dx, dy, dw, dh, placeholder in details:
        _text(slide, dx, dy, dw, dh, placeholder, 7, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# Slide 31: McKinsey 7S Framework (7-node web)
def build_mckinsey_7s(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    cx_c, cy_c = 5.0, 3.5
    # Shared Values (중심)
    _circle(slide, cx_c - 0.7, cy_c - 0.7, 1.4, CL_ACCENT, "{{shared_values}}\nShared\nValues", CL_WHITE, 8)

    # 6개 외곽 노드 (60도 간격)
    outer_labels = [
        ("{{strategy}}", "Strategy", True),
        ("{{structure}}", "Structure", True),
        ("{{systems}}", "Systems", True),
        ("{{style}}", "Style", False),
        ("{{staff}}", "Staff", False),
        ("{{skills}}", "Skills", False),
    ]
    radius = 2.2
    for i, (placeholder, label, is_hard) in enumerate(outer_labels):
        angle = math.radians(60 * i - 90)
        nx = cx_c + radius * math.cos(angle)
        ny = cy_c + radius * math.sin(angle)

        fill = CL_DARK if is_hard else CL_GREY
        d = 0.9
        _circle(slide, nx - d/2, ny - d/2, d, fill, f"{label}", CL_WHITE, 7)

        # 연결선
        lx = cx_c + 0.7 * math.cos(angle)
        ly = cy_c + 0.7 * math.sin(angle)
        _connector(slide, lx, ly, nx, ny, CL_GREY_LIGHT)

        # 상세 텍스트 (외곽)
        tx = nx + (0.6 if math.cos(angle) > 0 else -2.2)
        ty = ny + (0.5 if math.sin(angle) > 0 else -0.8)
        _text(slide, tx, ty, 1.8, 0.6, placeholder, 7, False, CL_BLACK)

    # 범례
    _rect(slide, 0.3, 6.0, 0.2, 0.2, CL_DARK)
    _text(slide, 0.55, 6.0, 1.0, 0.2, "Hard S", 7, True, CL_BLACK)
    _rect(slide, 1.7, 6.0, 0.2, 0.2, CL_GREY)
    _text(slide, 1.95, 6.0, 1.0, 0.2, "Soft S", 7, True, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# Slide 32: Three Horizons Growth Map
def build_three_horizons(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)

    # 3개 S-curve 근사 (사각형 밴드로 표현)
    horizons = [
        ("H1", "{{h1_label}}", CL_ACCENT, 0.85, 5.5, 0.3, 4.5),       # 좌→중앙, 높이 감소
        ("H2", "{{h2_label}}", CL_ACCENT_MID, 2.5, 4.0, 0.5, 5.2),    # 중앙, 상승→하강
        ("H3", "{{h3_label}}", CL_GREY, 5.0, 2.5, 0.7, 5.8),          # 우측, 상승
    ]

    # 축
    _connector(slide, 0.5, 5.8, 9.5, 5.8, CL_BORDER)  # X축 (시간)
    _connector(slide, 0.5, 0.85, 0.5, 5.8, CL_BORDER)  # Y축 (수익)
    _text(slide, 0.3, 0.65, 1.0, 0.2, "Revenue", 7, True, CL_GREY)
    _text(slide, 8.5, 5.85, 1.0, 0.2, "Time →", 7, True, CL_GREY)

    for hz_label, placeholder, color, start_x, end_x, start_y, end_y in horizons:
        # 밴드 (사다리꼴 근사 — 3개 사각형으로 S-curve 흉내)
        mid_x = (start_x + end_x) / 2
        peak_y = min(start_y, end_y) - 0.5

        # 상승 구간
        _rect(slide, start_x, start_y, mid_x - start_x, 5.8 - start_y, color)
        # 정점 구간
        _rect(slide, mid_x, peak_y, (end_x - mid_x) * 0.6, 5.8 - peak_y, color)

        # 라벨
        _circle(slide, start_x + 0.1, start_y - 0.5, 0.4, color, hz_label, CL_WHITE, 10)
        _text(slide, start_x + 0.6, start_y - 0.5, 2.0, 0.4, placeholder, 8, True, CL_BLACK)

    # 상세 텍스트
    for i, (hx, hy, placeholder) in enumerate([
        (0.5, 0.85, "{{h1_detail}}"),
        (3.5, 2.0, "{{h2_detail}}"),
        (6.5, 1.5, "{{h3_detail}}"),
    ]):
        _text(slide, hx, hy, 2.5, 0.8, placeholder, 7, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}")


# Slide 33: Mekko/Marimekko Chart (가변 폭 스택드 바)
def build_mekko(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    _illustrative_badge(slide)

    # 5개 컬럼 (가변 폭) — 시장 규모에 비례
    widths = [2.5, 1.8, 1.5, 1.2, 1.0]  # 인치
    total_w = sum(widths)
    scale = 8.5 / total_w
    widths = [w * scale for w in widths]

    bar_h = 4.5
    base_y = 0.85
    sx = 0.5

    for ci, cw in enumerate(widths):
        # 3개 세그먼트 (합=100%)
        segments = [(0.5, CL_ACCENT), (0.3, CL_ACCENT_MID), (0.2, CL_GREY_LIGHT)]
        seg_y = base_y
        for si, (ratio, color) in enumerate(segments):
            seg_h = bar_h * ratio
            _rect(slide, sx, seg_y, cw - 0.05, seg_h, color)
            _text(slide, sx + 0.05, seg_y + 0.02, cw - 0.15, seg_h - 0.04,
                  f"{{{{seg_{ci}_{si}}}}}", 7, False,
                  CL_WHITE if si < 2 else CL_BLACK, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
            seg_y += seg_h

        # 컬럼 라벨 (하단)
        _text(slide, sx, base_y + bar_h + 0.05, cw - 0.05, 0.3,
              f"{{{{col_{ci}}}}}", 7, True, CL_BLACK, align=PP_ALIGN.CENTER)
        # 폭 라벨 (상단)
        _text(slide, sx, base_y - 0.25, cw - 0.05, 0.2,
              f"{{{{size_{ci}}}}}", 7, False, CL_GREY, align=PP_ALIGN.CENTER)

        sx += cw

    # 범례
    for i, (label, color) in enumerate([("{{legend_1}}", CL_ACCENT), ("{{legend_2}}", CL_ACCENT_MID), ("{{legend_3}}", CL_GREY_LIGHT)]):
        lx = 0.5 + i * 2.5
        _rect(slide, lx, 5.6, 0.2, 0.15, color)
        _text(slide, lx + 0.3, 5.58, 2.0, 0.2, label, 7, False, CL_BLACK)

    _takeaway(slide, "{{takeaway}}", 5.9)


# Slide 34: Dense Table with Inline Bar Charts
def build_table_with_bars(prs):
    slide = _s(prs)
    _header(slide, "{{title}}", "{{breadcrumb}}")
    _footer(slide)
    _illustrative_badge(slide)

    n_rows, n_cols = 6, 4
    col_ws = [2.5, 2.0, 3.5, 1.5]  # 항목, 값, 바, 상태
    row_h = 0.65
    sx, sy = 0.3, 0.75

    # 헤더
    headers = ["{{h_item}}", "{{h_value}}", "{{h_progress}}", "{{h_status}}"]
    cx = sx
    for ci, (hdr, cw) in enumerate(zip(headers, col_ws)):
        _rect(slide, cx, sy, cw, 0.35, CL_DARK)
        _text(slide, cx + 0.05, sy, cw - 0.1, 0.35, hdr, 8, True, CL_WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw

    # 데이터 행
    bar_percentages = [85, 60, 45, 90, 30, 75]
    for r in range(n_rows):
        ry = sy + 0.35 + r * row_h
        fill = CL_WHITE if r % 2 == 0 else CL_BG
        cx = sx
        for ci, cw in enumerate(col_ws):
            _rect(slide, cx, ry, cw, row_h, fill, CL_BORDER)

            if ci == 0:  # 항목
                _text(slide, cx + 0.08, ry, cw - 0.16, row_h, f"{{{{item_{r}}}}}", 8, True, CL_BLACK,
                      anchor=MSO_ANCHOR.MIDDLE)
            elif ci == 1:  # 값
                _text(slide, cx + 0.08, ry, cw - 0.16, row_h, f"{{{{value_{r}}}}}", 9, True, CL_ACCENT,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            elif ci == 2:  # 인라인 바
                pct = bar_percentages[r] / 100.0
                bar_max_w = cw - 0.3
                bar_w = bar_max_w * pct
                bar_y = ry + (row_h - 0.2) / 2
                # 배경 바
                _rect(slide, cx + 0.15, bar_y, bar_max_w, 0.2, CL_BG)
                # 채움 바
                bar_color = CL_ACCENT if pct > 0.7 else CL_ACCENT_MID if pct > 0.4 else CL_GREY
                _rect(slide, cx + 0.15, bar_y, bar_w, 0.2, bar_color)
                # % 라벨
                _text(slide, cx + 0.15 + bar_max_w + 0.02, bar_y - 0.02, 0.5, 0.25,
                      f"{bar_percentages[r]}%", 7, True, CL_BLACK)
            elif ci == 3:  # 상태
                _text(slide, cx + 0.08, ry, cw - 0.16, row_h, f"{{{{status_{r}}}}}", 8, False, CL_BLACK,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += cw

    _takeaway(slide, "{{takeaway}}")


if __name__ == "__main__":
    build_template_library()
