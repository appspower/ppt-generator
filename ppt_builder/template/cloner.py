"""슬라이드 클로너 - .pptx 템플릿에서 슬라이드를 복제한다.

핵심 원리:
  1. 템플릿 .pptx를 열어서 원하는 슬라이드의 XML을 가져옴
  2. 대상 프레젠테이션에 빈 슬라이드를 추가
  3. XML 내용과 관련 리소스(이미지, 관계)를 복사 — rId remap 필수
  4. 텍스트 placeholder를 실제 데이터로 치환

이미지 rId remap 패턴 (APryor6 / scanny):
  - src_part.rels에서 image rel을 찾아 blob/content_type을 꺼낸다
  - dst_part에 새 ImagePart를 생성하여 새 rId 발급
  - 새 sp_tree 내 모든 <a:blip r:embed="OLD_RID"/>를 새 rId로 치환
"""

import copy
from io import BytesIO
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.image import ImagePart
from pptx.slide import Slide

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _qn(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


def _guess_ext(content_type: str) -> str:
    mapping = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/jpg": "jpg",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/tiff": "tiff",
        "image/svg+xml": "svg",
        "image/x-emf": "emf",
        "image/x-wmf": "wmf",
    }
    return mapping.get(content_type.lower(), "png")


def _add_image_part(dst_part, blob: bytes, content_type: str) -> tuple:
    """Image blob을 dst_part에 추가하여 (image_part, rId)를 반환한다.

    SVG/EMF 등 PIL이 못 읽는 포맷도 지원하기 위해, raster 이미지(PNG/JPG 등)는
    SlidePart.get_or_add_image_part()를 사용하고, 벡터 이미지는 ImagePart를
    직접 생성한다.
    """
    raster_types = {
        "image/png", "image/jpeg", "image/jpg", "image/gif",
        "image/bmp", "image/tiff",
    }

    if content_type.lower() in raster_types:
        # PIL이 처리 가능 — 표준 경로
        return dst_part.get_or_add_image_part(BytesIO(blob))

    # 벡터/특수 포맷 — ImagePart를 직접 생성 (PIL 우회)
    package = dst_part.package
    ext = _guess_ext(content_type)
    partname = package.next_partname(f"/ppt/media/image%d.{ext}")
    image_part = ImagePart(partname, content_type, package, blob)
    rId = dst_part.relate_to(image_part, RT.IMAGE)
    return image_part, rId


class SlideCloner:
    """템플릿 .pptx에서 슬라이드를 복제하는 클래스."""

    def __init__(self, template_path: Path):
        """템플릿 파일을 로드한다."""
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        self.template_prs = Presentation(str(template_path))
        self._template_path = template_path

    @property
    def slide_count(self) -> int:
        return len(self.template_prs.slides)

    def get_slide_titles(self) -> list[str]:
        """템플릿의 모든 슬라이드 제목을 반환한다."""
        titles = []
        for slide in self.template_prs.slides:
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    title = shape.text_frame.text[:50]
                    break
            titles.append(title)
        return titles

    def clone_slide(self, target_prs: Presentation, slide_index: int) -> Slide:
        """템플릿의 slide_index번째 슬라이드를 target_prs에 복제한다.

        Args:
            target_prs: 대상 프레젠테이션
            slide_index: 복제할 슬라이드 인덱스 (0-based)

        Returns:
            복제된 Slide 객체
        """
        if slide_index >= len(self.template_prs.slides):
            raise IndexError(f"Slide index {slide_index} out of range (max: {len(self.template_prs.slides) - 1})")

        src_slide = self.template_prs.slides[slide_index]

        # 1. 대상에 빈 슬라이드 추가
        blank_layout = target_prs.slide_layouts[6]  # Blank
        new_slide = target_prs.slides.add_slide(blank_layout)

        # 2. 기존 빈 슬라이드의 shape 제거
        sp_tree = new_slide.shapes._spTree
        for sp in list(sp_tree):
            if sp.tag.endswith('}sp') or sp.tag.endswith('}pic') or sp.tag.endswith('}graphicFrame'):
                sp_tree.remove(sp)

        # 3. 소스 슬라이드의 shape을 deep copy + sp_tree에 append
        src_sp_tree = src_slide.shapes._spTree
        copied_shapes = []
        for sp in src_sp_tree:
            tag = sp.tag.split('}')[-1] if '}' in sp.tag else sp.tag
            if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
                new_sp = copy.deepcopy(sp)
                sp_tree.append(new_sp)
                copied_shapes.append(new_sp)

        # 4. 이미지 rId remap (소스 → 대상)
        self._remap_image_rels(src_slide, new_slide, copied_shapes)

        return new_slide

    def _remap_image_rels(self, src_slide, dst_slide, copied_shapes: list) -> None:
        """복사된 shape XML 내의 이미지 rId를 대상 슬라이드 기준으로 재매핑한다.

        모던 PPTX는 한 이미지에 대해 두 개의 rId를 가질 수 있다:
        - <a:blip r:embed="rId3"/> (PNG fallback, namespace=DrawingML)
        - <asvg:svgBlip r:embed="rId4"/> (SVG vector, namespace=Office 2016 SVG)

        a:blip만 처리하면 SVG가 깨져서 PowerPoint가 빨간 X를 표시한다.
        해법: 모든 element를 순회하며 r:embed/r:link 속성을 가진 노드를 모두 처리.
        """
        src_part = src_slide.part
        dst_part = dst_slide.part

        embed_attr = _qn(R_NS, "embed")
        link_attr = _qn(R_NS, "link")

        # 사용된 rId 모음 (모든 element에서)
        used_rids: set[str] = set()
        for sp in copied_shapes:
            for el in sp.iter():
                rid = el.get(embed_attr) or el.get(link_attr)
                if rid:
                    used_rids.add(rid)

        if not used_rids:
            return

        # 각 rId를 새 rId로 매핑
        rid_map: dict[str, str] = {}
        for old_rid in used_rids:
            try:
                rel = src_part.rels.get(old_rid)
                if rel is None:
                    continue
                if "image" not in rel.reltype:
                    continue
                if rel.is_external:
                    # 외부 링크 이미지는 다운로드 안 함 (Phase 1 범위 밖)
                    continue
                src_image_part = rel.target_part
                blob = src_image_part.blob
                content_type = src_image_part.content_type

                _, new_rid = _add_image_part(dst_part, blob, content_type)
                rid_map[old_rid] = new_rid
            except Exception as e:
                print(f"[WARN] image rId remap failed for {old_rid}: {e}")

        # XML 내 모든 r:embed/r:link 속성을 새 rId로 치환
        for sp in copied_shapes:
            for el in sp.iter():
                old_embed = el.get(embed_attr)
                if old_embed and old_embed in rid_map:
                    el.set(embed_attr, rid_map[old_embed])
                old_link = el.get(link_attr)
                if old_link and old_link in rid_map:
                    el.set(link_attr, rid_map[old_link])


def clone_slide_from_template(
    target_prs: Presentation,
    template_path: Path,
    slide_index: int,
) -> Slide:
    """편의 함수 - 템플릿에서 슬라이드 하나를 복제한다."""
    cloner = SlideCloner(template_path)
    return cloner.clone_slide(target_prs, slide_index)
