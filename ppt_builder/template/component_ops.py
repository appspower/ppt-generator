"""컴포넌트 추출/삽입 API — N1-Lite (single-component reuse) 인프라.

[docs/PROJECT_DIRECTION.md](../../docs/PROJECT_DIRECTION.md) 헌법 §1 N1-Lite의
"마스터에서 도형 그룹 1개 추출 → 빈 슬라이드에 단독 배치" 부분을 구현한다.

기술 근거: [docs/_research/component_extraction_research.md] T1 (lxml deepcopy + rId rewrite).
사양: [docs/N1_LITE_IMPLEMENTATION.md] §3 component_ops.py.

핵심 API
--------
- extract_group(src_slide, group_indices) -> ComponentXML
- insert_component(dst_slide, component, position=None) -> int
- create_blank_slide_with_master_theme(target_prs) -> Slide

extract → insert 패턴
--------------------
원본 슬라이드의 leaf shape 인덱스(iter_leaf_shapes 평탄 idx)를 받아 deepcopy.
공통 부모가 <p:grpSp>이고 그 그룹 전체가 선택되면 그룹 통째 복제(가장 안전),
아니면 합성 <p:grpSp>로 래핑한다. a16:creationId UUID 재생성으로 중복 손상 회피.

insert 시 dst_slide.part에 image/chart rel을 다시 묶어서(relate_to) rId를 재발급하고
복제된 XML 내 r:embed / r:link 속성을 새 rId로 일괄 치환한다.

v1 범위 외 (사양 §3.4)
---------------------
- chart 컴포넌트 (chart part + xlsx 동기 복사 별도 작업)
- cross-master theme 재작성
- transformation (resize/recolor)
- rel_remap override 매개변수 (필요 시 추가)
"""

from __future__ import annotations

import copy
import uuid
from dataclasses import dataclass, field
from typing import Iterable

from lxml import etree
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from .edit_ops import SlideEditError, iter_leaf_shapes


# --- 네임스페이스 / 태그 ---------------------------------------------------------

_TAG_GRPSP = qn("p:grpSp")
_TAG_SP = qn("p:sp")
_TAG_PIC = qn("p:pic")
_TAG_GRAPHICFRAME = qn("p:graphicFrame")
_TAG_CXNSP = qn("p:cxnSp")
_TAG_NVGRPSPPR = qn("p:nvGrpSpPr")
_TAG_GRPSPPR = qn("p:grpSpPr")
_TAG_SPPR = qn("p:spPr")
_TAG_XFRM_A = qn("a:xfrm")
_TAG_OFF = qn("a:off")
_TAG_EXT = qn("a:ext")
_TAG_CHOFF = qn("a:chOff")
_TAG_CHEXT = qn("a:chExt")
_TAG_EXTLST_P = qn("p:extLst")
# a16 네임스페이스는 python-pptx _nsmap에 없으므로 직접 URI 사용
_NS_A16 = "http://schemas.microsoft.com/office/drawing/2014/main"
_TAG_CREATIONID = "{%s}creationId" % _NS_A16

_ATTR_R_EMBED = qn("r:embed")
_ATTR_R_LINK = qn("r:link")

_SHAPE_TAGS = {_TAG_SP, _TAG_PIC, _TAG_GRAPHICFRAME, _TAG_GRPSP, _TAG_CXNSP}


# --- ComponentXML 데이터 컨테이너 -----------------------------------------------

@dataclass
class ComponentXML:
    """추출된 컴포넌트. insert_component의 입력.

    element        : deepcopy된 <p:grpSp> 또는 단일 shape XML.
    source_part    : 원본 SlidePart — rId 재바인딩 시 src 측 part 참조.
    source_rids    : element 내부에 등장하는 src rel ID 목록.
    bbox_emu       : {left, top, width, height} EMU 단위 (원본 위치).
    metadata       : 호출자가 자유롭게 채우는 부가 정보 (family, slot count 등).
    """

    element: etree._Element
    source_part: object  # SlidePart, but typing 'object' to avoid private import
    source_rids: set[str]
    bbox_emu: dict
    metadata: dict = field(default_factory=dict)


# --- 공통 헬퍼 -----------------------------------------------------------------

def _is_shape_element(el: etree._Element) -> bool:
    return el.tag in _SHAPE_TAGS


def _resolve_shapes(slide: Slide, indices: list[int]) -> list[BaseShape]:
    """flat_idx 리스트 → BaseShape 리스트. 입력 순서 보존."""
    if not indices:
        raise SlideEditError("group_indices is empty.")
    wanted = set(indices)
    by_idx: dict[int, BaseShape] = {}
    for fi, sh in iter_leaf_shapes(slide):
        if fi in wanted:
            by_idx[fi] = sh
        if len(by_idx) == len(wanted):
            break
    missing = wanted - by_idx.keys()
    if missing:
        raise SlideEditError(
            f"flat indices not found in slide: {sorted(missing)}"
        )
    return [by_idx[i] for i in indices]


def _bbox_from_shapes(shapes: list[BaseShape]) -> dict:
    """N개 shape의 외접 bbox (EMU). 그룹 wrap 시 사용."""
    lefts, tops, rights, bottoms = [], [], [], []
    for sh in shapes:
        left, top, w, h = sh.left, sh.top, sh.width, sh.height
        if left is None or top is None or w is None or h is None:
            # 그룹 멤버는 None이 나올 수 있음. 자체 element xfrm에서 직접 읽는다.
            xfrm = sh._element.find(".//" + _TAG_XFRM_A)
            if xfrm is None:
                continue
            off = xfrm.find(_TAG_OFF)
            ext = xfrm.find(_TAG_EXT)
            if off is None or ext is None:
                continue
            left = int(off.get("x", 0))
            top = int(off.get("y", 0))
            w = int(ext.get("cx", 0))
            h = int(ext.get("cy", 0))
        lefts.append(left)
        tops.append(top)
        rights.append(left + w)
        bottoms.append(top + h)
    if not lefts:
        return {"left": 0, "top": 0, "width": 0, "height": 0}
    L, T = min(lefts), min(tops)
    return {
        "left": L,
        "top": T,
        "width": max(rights) - L,
        "height": max(bottoms) - T,
    }


def _bbox_from_grpsp(grpsp: etree._Element) -> dict:
    """<p:grpSp>의 grpSpPr/xfrm/off+ext에서 bbox 읽기."""
    grpSpPr = grpsp.find(_TAG_GRPSPPR)
    if grpSpPr is None:
        return {"left": 0, "top": 0, "width": 0, "height": 0}
    xfrm = grpSpPr.find(_TAG_XFRM_A)
    if xfrm is None:
        return {"left": 0, "top": 0, "width": 0, "height": 0}
    off = xfrm.find(_TAG_OFF)
    ext = xfrm.find(_TAG_EXT)
    return {
        "left": int(off.get("x", 0)) if off is not None else 0,
        "top": int(off.get("y", 0)) if off is not None else 0,
        "width": int(ext.get("cx", 0)) if ext is not None else 0,
        "height": int(ext.get("cy", 0)) if ext is not None else 0,
    }


def _bbox_from_shape_element(el: etree._Element) -> dict:
    """단일 shape element의 spPr/xfrm 또는 grpSpPr/xfrm bbox."""
    if el.tag == _TAG_GRPSP:
        return _bbox_from_grpsp(el)
    sp_pr = el.find(_TAG_SPPR)
    if sp_pr is None:
        # graphicFrame은 <p:xfrm>를 직접 가짐 (a: 네임스페이스 아님)
        xfrm = el.find(qn("p:xfrm"))
    else:
        xfrm = sp_pr.find(_TAG_XFRM_A)
    if xfrm is None:
        return {"left": 0, "top": 0, "width": 0, "height": 0}
    off = xfrm.find(_TAG_OFF)
    ext = xfrm.find(_TAG_EXT)
    return {
        "left": int(off.get("x", 0)) if off is not None else 0,
        "top": int(off.get("y", 0)) if off is not None else 0,
        "width": int(ext.get("cx", 0)) if ext is not None else 0,
        "height": int(ext.get("cy", 0)) if ext is not None else 0,
    }


def _regen_creationids(root: etree._Element) -> int:
    """<a16:creationId id="{...}"/> 모두 새 GUID로 교체. 손상 방지(이슈 #961).

    반환: 교체된 element 개수.
    """
    n = 0
    for el in root.iter(_TAG_CREATIONID):
        el.set("id", "{" + str(uuid.uuid4()).upper() + "}")
        n += 1
    return n


def _collect_rids(root: etree._Element) -> set[str]:
    """element 내 r:embed / r:link 속성 값 수집."""
    rids: set[str] = set()
    for el in root.iter():
        v = el.get(_ATTR_R_EMBED)
        if v:
            rids.add(v)
        v = el.get(_ATTR_R_LINK)
        if v:
            rids.add(v)
    return rids


def _set_top_offset(element: etree._Element, left_emu: int, top_emu: int) -> None:
    """최상위 element의 xfrm/off를 (left, top)으로 설정.

    grpSp는 grpSpPr/xfrm, sp/cxnSp는 spPr/xfrm, graphicFrame은 p:xfrm.
    chOff/chExt는 그대로 둔다 (그룹 내부 좌표계는 그대로, 그룹 자체만 평행이동).
    """
    if element.tag == _TAG_GRPSP:
        container = element.find(_TAG_GRPSPPR)
        xfrm = container.find(_TAG_XFRM_A) if container is not None else None
    elif element.tag == _TAG_GRAPHICFRAME:
        xfrm = element.find(qn("p:xfrm"))
    else:
        container = element.find(_TAG_SPPR)
        xfrm = container.find(_TAG_XFRM_A) if container is not None else None

    if xfrm is None:
        # xfrm이 없으면 만든다 (sp/cxnSp의 경우 spPr 안에).
        if element.tag == _TAG_GRAPHICFRAME:
            xfrm = etree.SubElement(element, qn("p:xfrm"))
        elif element.tag == _TAG_GRPSP:
            container = element.find(_TAG_GRPSPPR)
            if container is None:
                container = etree.SubElement(element, _TAG_GRPSPPR)
            xfrm = etree.SubElement(container, _TAG_XFRM_A)
        else:
            container = element.find(_TAG_SPPR)
            if container is None:
                container = etree.SubElement(element, _TAG_SPPR)
            xfrm = etree.SubElement(container, _TAG_XFRM_A)

    off = xfrm.find(_TAG_OFF)
    if off is None:
        off = etree.SubElement(xfrm, _TAG_OFF)
    off.set("x", str(int(left_emu)))
    off.set("y", str(int(top_emu)))


def _wrap_in_grpsp(elements: list[etree._Element], bbox: dict) -> etree._Element:
    """N개 shape element를 새 <p:grpSp>로 감싼다. 자식 좌표계는 변형 없이 그대로."""
    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    grp = etree.Element(_TAG_GRPSP, nsmap=nsmap)

    # nvGrpSpPr (필수)
    nv = etree.SubElement(grp, _TAG_NVGRPSPPR)
    cnv_pr = etree.SubElement(nv, qn("p:cNvPr"))
    cnv_pr.set("id", "0")
    cnv_pr.set("name", "Group")
    etree.SubElement(nv, qn("p:cNvGrpSpPr"))
    etree.SubElement(nv, qn("p:nvPr"))

    # grpSpPr (xfrm)
    gpr = etree.SubElement(grp, _TAG_GRPSPPR)
    xfrm = etree.SubElement(gpr, _TAG_XFRM_A)
    L, T, W, H = bbox["left"], bbox["top"], bbox["width"], bbox["height"]
    off = etree.SubElement(xfrm, _TAG_OFF)
    off.set("x", str(int(L)))
    off.set("y", str(int(T)))
    ext = etree.SubElement(xfrm, _TAG_EXT)
    ext.set("cx", str(int(W)))
    ext.set("cy", str(int(H)))
    chOff = etree.SubElement(xfrm, _TAG_CHOFF)
    chOff.set("x", str(int(L)))
    chOff.set("y", str(int(T)))
    chExt = etree.SubElement(xfrm, _TAG_CHEXT)
    chExt.set("cx", str(int(W)))
    chExt.set("cy", str(int(H)))

    for el in elements:
        grp.append(el)
    return grp


def _common_xml_parent(elements: list[etree._Element]) -> etree._Element | None:
    """모든 element의 XML 부모가 동일하면 그 부모, 아니면 None."""
    parents = {el.getparent() for el in elements}
    if len(parents) == 1:
        return next(iter(parents))
    return None


# --- 핵심 API ------------------------------------------------------------------

def extract_group(src_slide: Slide, group_indices: list[int]) -> ComponentXML:
    """src_slide의 leaf shape 인덱스 리스트를 컴포넌트로 추출.

    선택 규칙:
      1) 1개 인덱스 → 그 shape 단독 deepcopy.
      2) 모든 인덱스가 동일 <p:grpSp> 부모 + 그 부모의 모든 shape 자식이 선택됨
         → 부모 grpSp 통째 deepcopy (가장 안전, 그룹 metadata 보존).
      3) 그 외 → 각 shape deepcopy 후 합성 <p:grpSp>로 래핑.

    공통:
      - a16:creationId UUID 재생성 (이슈 #961, PPT 손상 회피).
      - r:embed / r:link 수집 (insert_component에서 재바인딩).
      - bbox 계산.
    """
    shapes = _resolve_shapes(src_slide, group_indices)
    elements = [sh._element for sh in shapes]

    if len(group_indices) == 1:
        cloned = copy.deepcopy(elements[0])
        bbox = _bbox_from_shape_element(cloned)
    else:
        parent = _common_xml_parent(elements)
        if parent is not None and parent.tag == _TAG_GRPSP:
            siblings = [c for c in parent if _is_shape_element(c)]
            ids_selected = {id(e) for e in elements}
            ids_siblings = {id(s) for s in siblings}
            if ids_selected == ids_siblings:
                # 그룹 통째 추출
                cloned = copy.deepcopy(parent)
                bbox = _bbox_from_grpsp(cloned)
            else:
                # 그룹 일부만 → wrap
                cloned_children = [copy.deepcopy(e) for e in elements]
                bbox = _bbox_from_shapes(shapes)
                cloned = _wrap_in_grpsp(cloned_children, bbox)
        else:
            # 부모가 다르거나 spTree → wrap
            cloned_children = [copy.deepcopy(e) for e in elements]
            bbox = _bbox_from_shapes(shapes)
            cloned = _wrap_in_grpsp(cloned_children, bbox)

    _regen_creationids(cloned)
    rids = _collect_rids(cloned)
    return ComponentXML(
        element=cloned,
        source_part=src_slide.part,
        source_rids=rids,
        bbox_emu=bbox,
    )


def insert_component(
    dst_slide: Slide,
    component: ComponentXML,
    position: tuple[int, int] | None = None,
) -> int:
    """dst_slide에 component 삽입. rId 재바인딩 + creationId(이미 재생성됨) + 위치 옵션.

    Args:
        dst_slide: 대상 슬라이드
        component: extract_group 산출물. 같은 component를 N번 insert해도 안전하도록
                   매 호출마다 element를 deepcopy한다 (creationId도 재재생성).
        position : (left_emu, top_emu). None이면 원본 위치 유지.

    Returns:
        삽입된 컴포넌트 내 첫 leaf shape의 새 flat_idx.
    """
    # 다회 insert 안전 + 크로스-package 클래스 복원:
    # copy.deepcopy는 lxml._Element만 반환해 python-pptx가 has_ph_elm 등을 못 찾음.
    # tostring → parse_xml 라운드트립으로 oxml_parser의 커스텀 클래스(CT_GroupShape 등)를 재바인딩한다.
    new_element = parse_xml(etree.tostring(component.element))
    _regen_creationids(new_element)

    # rId 재바인딩
    src_part = component.source_part
    dst_part = dst_slide.part
    rid_map: dict[str, str] = {}
    for old_rid in component.source_rids:
        rel = src_part.rels.get(old_rid)
        if rel is None:
            continue
        if rel.is_external:
            # 외부 링크는 v1 범위 밖 — 원래 rId를 임시 보존하지만 깨질 수 있음
            continue
        target_part = rel.target_part
        new_rid = dst_part.relate_to(target_part, rel.reltype)
        rid_map[old_rid] = new_rid

    # XML 내 r:embed / r:link 치환
    if rid_map:
        for el in new_element.iter():
            v = el.get(_ATTR_R_EMBED)
            if v and v in rid_map:
                el.set(_ATTR_R_EMBED, rid_map[v])
            v = el.get(_ATTR_R_LINK)
            if v and v in rid_map:
                el.set(_ATTR_R_LINK, rid_map[v])

    # 위치 변경
    if position is not None:
        _set_top_offset(new_element, position[0], position[1])

    # 삽입 전 leaf 개수 기록 → 삽입 후 첫 새 leaf가 flat_idx[pre_count]
    pre_count = sum(1 for _ in iter_leaf_shapes(dst_slide))

    # spTree 삽입 (extLst 앞)
    sp_tree = dst_slide.shapes._spTree
    extlst = sp_tree.find(_TAG_EXTLST_P)
    if extlst is not None:
        extlst.addprevious(new_element)
    else:
        sp_tree.append(new_element)

    # 새 leaf가 정말 추가됐는지 확인 (방어적)
    post_count = sum(1 for _ in iter_leaf_shapes(dst_slide))
    if post_count <= pre_count:
        raise SlideEditError(
            f"Insert did not add leaves: pre={pre_count}, post={post_count}. "
            f"spTree insertion likely failed."
        )
    return pre_count


def create_blank_slide_with_master_theme(target_prs) -> Slide:
    """target_prs에 빈 레이아웃 슬라이드 1개 추가.

    target_prs를 마스터 .pptx에서 로드한 경우 (Presentation(master_path)),
    같은 theme/font/layout이 자동 상속된다 — N1-Lite의 "스타일 일치" 보장.
    blank layout은 마지막 레이아웃이거나 'Blank' 이름을 가진 것으로 추정.
    못 찾으면 인덱스 6 (python-pptx 기본 blank)을 시도, 그것도 없으면 마지막 레이아웃.
    """
    layouts = target_prs.slide_layouts
    chosen = None
    for layout in layouts:
        try:
            if (layout.name or "").strip().lower() == "blank":
                chosen = layout
                break
        except Exception:
            pass
    if chosen is None:
        try:
            chosen = layouts[6]
        except IndexError:
            chosen = layouts[len(layouts) - 1]

    slide = target_prs.slides.add_slide(chosen)

    # 레이아웃에 placeholder가 있으면 모두 제거 (단독 컴포넌트용 깨끗한 캔버스)
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree):
        tag = sp.tag
        if tag in _SHAPE_TAGS:
            # placeholder만 제거 (nvSpPr/nvPr/ph 존재 여부)
            ph = sp.find(".//" + qn("p:ph"))
            if ph is not None:
                sp_tree.remove(sp)
    return slide


def _apply_series_color(series, color_hex: str) -> None:
    """차트 series에 hex color 적용. fill + line 양쪽 시도 (chart type 무관).

    bar/column/pie: fill.fore_color.rgb 가 막대/조각 색을 바꾼다.
    line/scatter: line.color.rgb 가 선 색을 바꾼다.
    둘 다 try (실패 무시) — chart type별 분기 없이 안전.
    """
    from pptx.dml.color import RGBColor

    rgb = RGBColor.from_string(color_hex.lstrip("#").upper())
    try:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = rgb
    except Exception:
        pass
    try:
        series.format.line.color.rgb = rgb
    except Exception:
        pass


def replace_chart_data(
    slide: Slide,
    flat_idx: int,
    categories: list[str],
    series: list[tuple[str, list[float]]],
    *,
    series_colors: list[str | None] | None = None,
) -> None:
    """차트 더미 데이터 교체 (+ 옵션 색상 적용). python-pptx CategoryChartData 사용.

    Args:
        slide: 대상 슬라이드
        flat_idx: 차트 shape의 iter_leaf_shapes 평탄 인덱스
        categories: x축 카테고리 (예: ['Q1', 'Q2', 'Q3'])
        series: [(name, [val1, val2, val3]), ...] (예: [('Sales', [100, 200, 300])])
        series_colors: series 순서대로 hex '#RRGGBB' 또는 None. None인 시리즈는
            차트 원본 색 유지. series 길이와 맞아야 한다.

    Raises:
        SlideEditError: shape이 chart 아니거나 flat_idx 범위 밖,
            series_colors 길이가 series와 다른 경우

    예시:
        replace_chart_data(
            slide, flat_idx=2,
            categories=['Q1', 'Q2', 'Q3'],
            series=[('매출', [4200, 4500, 4800]), ('이익', [380, 410, 450])],
            series_colors=['#D04A02', None],  # 매출만 PwC 오렌지, 이익은 원본 색
        )
    """
    from pptx.chart.data import CategoryChartData
    from pptx.shapes.graphfrm import GraphicFrame

    target = None
    for fi, sh in iter_leaf_shapes(slide):
        if fi == flat_idx:
            target = sh
            break
    if target is None:
        raise SlideEditError(f"flat_idx={flat_idx} not found")
    if not isinstance(target, GraphicFrame) or not target.has_chart:
        raise SlideEditError(
            f"flat_idx={flat_idx} is not a chart "
            f"(type={type(target).__name__})"
        )
    if series_colors is not None and len(series_colors) != len(series):
        raise SlideEditError(
            f"series_colors length {len(series_colors)} != series length {len(series)}"
        )

    cd = CategoryChartData()
    cd.categories = categories
    for name, values in series:
        cd.add_series(name, values)
    target.chart.replace_data(cd)

    if series_colors:
        for i, color in enumerate(series_colors):
            if not color:
                continue
            try:
                _apply_series_color(target.chart.series[i], color)
            except Exception:
                # 차트 type이 series 색 변경을 지원하지 않는 경우 silent skip
                pass


def has_chart(slide: Slide) -> bool:
    """슬라이드에 차트가 있는지 빠른 확인."""
    from pptx.shapes.graphfrm import GraphicFrame
    for _, sh in iter_leaf_shapes(slide):
        if isinstance(sh, GraphicFrame) and sh.has_chart:
            return True
    return False


def chart_count(slide: Slide) -> int:
    """슬라이드 내 차트 개수."""
    from pptx.shapes.graphfrm import GraphicFrame
    n = 0
    for _, sh in iter_leaf_shapes(slide):
        if isinstance(sh, GraphicFrame) and sh.has_chart:
            n += 1
    return n


__all__ = [
    "ComponentXML",
    "extract_group",
    "insert_component",
    "create_blank_slide_with_master_theme",
    "replace_chart_data",
    "has_chart",
    "chart_count",
]
