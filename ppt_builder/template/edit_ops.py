"""편집 API 5종 — PPTAgent apis.py를 순수 python-pptx로 이식.

PPTAgent(EMNLP 2025)의 편집 API를 본 프로젝트 아키텍처에 맞게 재작성했다.
원본: C:/Users/y2kbo/.claude/projects/c--Users-y2kbo-Coding/PPTAgent/pptagent/apis.py

이식 함수
---------
- clone_paragraph(slide, div_id, paragraph_id) -> int
- replace_paragraph(slide, div_id, paragraph_id, text)
- del_paragraph(slide, div_id, paragraph_id)
- replace_image(slide, img_id, image_path)
- del_image(slide, img_id)

매핑 방침 (Option A, 평탄화 / flat)
-----------------------------------
- div_id: iter_leaf_shapes(slide)의 평탄 순회 인덱스 (0부터 시작)
- paragraph_id: shape.text_frame.paragraphs의 인덱스 (0부터 시작)
- 그룹은 재귀하여 내부 leaf shape만 인덱싱 (그룹 컨테이너 자체는 yield 안 함)
- PPTAgent의 schema_extractor.yaml / induct.py 파이프라인과 호환되도록 평탄화 선택

한글 East Asian 폰트 보존
------------------------
<a:ea typeface="맑은 고딕"> 등 한글 런의 타이포그래피는 parse_xml(_r.xml) XML 레벨
복제로 완벽 보존됨 (Phase A1 POC 슬라이드 1175에서 검증).

PPTAgent에서 제외한 부분
-----------------------
- CodeExecutor / SAFE_EVAL_GLOBALS: exec() 방식은 보안상 채택 안 함
- replace_image_with_table (table_XXXX.png 분기): Layer 2+ 후속 작업
- Closure 큐잉 패턴: 본 이식은 즉시 적용. 역순 인덱스 문제는 호출자가 관리
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterator

from bs4 import BeautifulSoup
from mistune import HTMLRenderer, create_markdown
from PIL import Image
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.text.text import _Paragraph, _Run


class SlideEditError(Exception):
    """편집 API 실패 시 발생."""


# --- 평탄화 순회 ---------------------------------------------------------------

def iter_leaf_shapes(slide: Slide) -> Iterator[tuple[int, BaseShape]]:
    """슬라이드 내 모든 leaf shape을 평탄화 순회하여 (flat_idx, shape)를 yield.

    그룹 shape은 재귀 진입하되 그룹 자체는 반환하지 않는다.
    순서: slide.shapes 선언 순 → 그룹 내부도 선언 순.
    """
    idx = [0]

    def _walk(shapes):
        for sh in shapes:
            if isinstance(sh, GroupShape):
                yield from _walk(sh.shapes)
            else:
                yield idx[0], sh
                idx[0] += 1

    yield from _walk(slide.shapes)


def _get_shape(slide: Slide, div_id: int) -> BaseShape:
    for fi, sh in iter_leaf_shapes(slide):
        if fi == div_id:
            return sh
    raise SlideEditError(
        f"Cannot find shape with div_id={div_id}. "
        f"Check range or prior deletions."
    )


def _get_paragraph(shape: BaseShape, paragraph_id: int) -> _Paragraph:
    if not shape.has_text_frame:
        raise SlideEditError(
            f"Shape (shape_id={shape.shape_id}) does not have a text frame."
        )
    paras = shape.text_frame.paragraphs
    if paragraph_id < 0 or paragraph_id >= len(paras):
        raise SlideEditError(
            f"paragraph_id={paragraph_id} out of range (0..{len(paras) - 1})."
        )
    return paras[paragraph_id]


def _runs_merge(para: _Paragraph) -> _Run:
    """paragraph의 첫 run만 남기고 나머지 제거. 빈 paragraph면 run을 추가."""
    runs = list(para.runs)
    if not runs:
        return para.add_run()
    first, *rest = runs
    for r in rest:
        r._r.getparent().remove(r._r)
    return first


# --- markdown → TextBlock -----------------------------------------------------

class _SlideRenderer(HTMLRenderer):
    """리스트를 래핑 태그로 감싸지 않는 렌더러 (PPTAgent 방식)."""

    def list(self, text: str, ordered: bool, **attrs) -> str:  # noqa: A002
        return text

    def list_item(self, text: str) -> str:
        return text


_markdown = create_markdown(renderer=_SlideRenderer(), plugins=["strikethrough"])

_MARKDOWN_STYLES = {
    "strong": "bold",
    "em": "italic",
    "code": "code",
    "del": "strikethrough",
}


@dataclass
class TextBlock:
    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False
    strikethrough: bool = False
    href: str | None = None

    def build_run(self, run: _Run) -> None:
        if self.bold:
            run.font.bold = True
        if self.italic:
            run.font.italic = True
        if self.code:
            run.font.name = "Consolas"
        if self.strikethrough:
            run.font.strikethrough = True
        if self.href is not None:
            run.hyperlink.address = self.href
        run.text = self.text


def _process_element(element, styles: dict | None = None) -> list[TextBlock]:
    if styles is None:
        styles = {}
    result: list[TextBlock] = []
    if isinstance(element, str):
        result.append(TextBlock(element, **styles))
        return result
    if element.name == "a":
        href = element.get("href")
        for child in element.children:
            blocks = _process_element(child, styles.copy())
            for b in blocks:
                b.href = href
            result.extend(blocks)
    elif _MARKDOWN_STYLES.get(element.name):
        new_styles = styles.copy()
        new_styles[_MARKDOWN_STYLES[element.name]] = True
        for child in element.children:
            result.extend(_process_element(child, new_styles))
    else:
        for child in element.children:
            result.extend(_process_element(child, styles))
    return result


def _markdown_to_blocks(text: str) -> list[TextBlock]:
    html = _markdown(text).strip()
    soup = BeautifulSoup(html, "html.parser")
    blocks = _process_element(soup)
    if not blocks:
        blocks = [TextBlock(text)]
    return blocks


# --- 5 API --------------------------------------------------------------------

def clone_paragraph(slide: Slide, div_id: int, paragraph_id: int) -> int:
    """paragraph를 XML 복제하여 같은 shape의 마지막 paragraph 뒤에 추가.

    반환값: 새로 생성된 paragraph의 인덱스 (len(paragraphs) - 1)

    한글 East Asian 폰트(<a:ea>), <a:latin>, font-size, vtab 등 모든 Run 속성 보존.
    """
    shape = _get_shape(slide, div_id)
    para = _get_paragraph(shape, paragraph_id)
    tf = shape.text_frame
    last_p = tf.paragraphs[-1]._p
    last_p.addnext(parse_xml(para._p.xml))
    return len(tf.paragraphs) - 1


def replace_paragraph(
    slide: Slide, div_id: int, paragraph_id: int, text: str
) -> None:
    """paragraph의 텍스트를 교체. markdown **bold**/*italic*/`code`/~~del~~ 지원.

    첫 run의 서식을 베이스로, markdown 블록별로 run을 XML 복제해 스타일을 덮어쓴다.
    """
    shape = _get_shape(slide, div_id)
    para = _get_paragraph(shape, paragraph_id)
    blocks = _markdown_to_blocks(text)
    first = _runs_merge(para)
    first.text = ""
    for _ in range(len(blocks) - 1):
        first._r.addnext(parse_xml(first._r.xml))
    for block, run in zip(blocks, para.runs):
        block.build_run(run)


def del_paragraph(slide: Slide, div_id: int, paragraph_id: int) -> None:
    """paragraph 삭제. 마지막 1개가 남으면 run 텍스트만 비워 서식을 보존한다."""
    shape = _get_shape(slide, div_id)
    para = _get_paragraph(shape, paragraph_id)
    tf = shape.text_frame
    if len(tf.paragraphs) == 1:
        for r in para.runs:
            r.text = ""
        return
    para._p.getparent().remove(para._p)


def replace_image(slide: Slide, img_id: int, image_path: str | Path) -> None:
    """Picture shape의 이미지를 교체. 비율 유지 + 세로 중앙 정렬."""
    shape = _get_shape(slide, img_id)
    if not isinstance(shape, Picture):
        raise SlideEditError(
            f"Shape at div_id={img_id} is not a Picture "
            f"(type={type(shape).__name__})."
        )
    image_path = str(image_path)

    image_part, new_rId = shape.part.get_or_add_image_part(image_path)
    blip = shape._element.find(".//" + qn("a:blip"))
    if blip is None:
        raise SlideEditError("Cannot find <a:blip> element inside Picture.")
    blip.set(qn("r:embed"), new_rId)

    img_w, img_h = Image.open(image_path).size
    if img_w <= 0 or img_h <= 0:
        return
    r = min(shape.width / img_w, shape.height / img_h)
    new_w = int(img_w * r)
    new_h = int(img_h * r)
    shape.top = int(shape.top + (shape.height - new_h) / 2)
    shape.width = new_w
    shape.height = new_h


def del_image(slide: Slide, img_id: int) -> None:
    """Picture shape을 슬라이드에서 삭제."""
    shape = _get_shape(slide, img_id)
    if not isinstance(shape, Picture):
        raise SlideEditError(
            f"Shape at div_id={img_id} is not a Picture "
            f"(type={type(shape).__name__})."
        )
    shape._element.getparent().remove(shape._element)


__all__ = [
    "SlideEditError",
    "TextBlock",
    "iter_leaf_shapes",
    "clone_paragraph",
    "replace_paragraph",
    "del_paragraph",
    "replace_image",
    "del_image",
]
