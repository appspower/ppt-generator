"""템플릿 에디터 - 원본 .pptx를 직접 편집하여 관계(relationships)를 100% 보존한다.

기존 cloner.py의 문제:
  - 슬라이드 XML을 새 프레젠테이션에 복제 → rId 관계 깨짐
  - 이미지, 차트, 하이퍼링크 등 모든 외부 참조 손실

해결 방식:
  - 원본 .pptx 복사본을 열어서 불필요한 슬라이드를 삭제
  - 남은 슬라이드에 TextSubstitutor로 텍스트 치환
  - 관계가 원본 그대로이므로 이미지/차트/서식 모두 보존

사용법:
  editor = TemplateEditor("templates/external/smartsheet_general.pptx")
  prs = editor.keep_slides([4, 7, 15])          # SWOT, Stages, Macro Plan만 남김
  editor.substitute(0, {"{{title}}": "실제 제목"})  # 첫 번째(=SWOT) 슬라이드 치환
  editor.save("output/result.pptx")
"""

import shutil
import tempfile
from pathlib import Path

from lxml import etree
from pptx import Presentation

from .substitutor import TextSubstitutor


class TemplateEditor:
    """원본 .pptx를 직접 편집하는 에디터. 관계 100% 보존."""

    def __init__(self, template_path: str | Path):
        """원본을 임시 파일로 복사한 후 열기.

        Args:
            template_path: 원본 .pptx 경로
        """
        self._src_path = Path(template_path)
        if not self._src_path.exists():
            raise FileNotFoundError(f"Template not found: {self._src_path}")

        # 원본 보호: 임시 복사본에서 작업
        self._tmp_dir = tempfile.mkdtemp(prefix="ppt_editor_")
        self._work_path = Path(self._tmp_dir) / self._src_path.name
        shutil.copy2(self._src_path, self._work_path)

        self._prs = Presentation(str(self._work_path))

    @property
    def prs(self) -> Presentation:
        """편집 중인 Presentation 객체."""
        return self._prs

    @property
    def slide_count(self) -> int:
        return len(self._prs.slides)

    def get_slide_info(self) -> list[dict]:
        """모든 슬라이드의 요약 정보를 반환한다."""
        info = []
        for i, slide in enumerate(self._prs.slides):
            title = ""
            text_len = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    text_len += len(txt)
                    if txt and not title:
                        title = txt[:80]
            info.append({"index": i, "title": title, "text_length": text_len})
        return info

    def keep_slides(self, indices: list[int]) -> "TemplateEditor":
        """지정한 인덱스의 슬라이드만 남기고 나머지를 삭제한다.

        Args:
            indices: 남길 슬라이드 인덱스 리스트 (0-based, 원본 기준)

        Returns:
            self (체이닝용)
        """
        if not indices:
            raise ValueError("At least one slide index must be specified")

        total = len(self._prs.slides)
        for idx in indices:
            if idx < 0 or idx >= total:
                raise IndexError(f"Slide index {idx} out of range (0-{total-1})")

        # 삭제할 인덱스 (역순으로 삭제해야 인덱스 안 밀림)
        to_remove = sorted(set(range(total)) - set(indices), reverse=True)

        for idx in to_remove:
            self._delete_slide(idx)

        return self

    def delete_slides(self, indices: list[int]) -> "TemplateEditor":
        """지정한 인덱스의 슬라이드를 삭제한다.

        Args:
            indices: 삭제할 슬라이드 인덱스 리스트 (0-based)

        Returns:
            self (체이닝용)
        """
        if not indices:
            return self

        total = len(self._prs.slides)
        for idx in sorted(set(indices), reverse=True):
            if idx < 0 or idx >= total:
                raise IndexError(f"Slide index {idx} out of range (0-{total-1})")
            self._delete_slide(idx)

        return self

    def substitute(self, slide_index: int, replacements: dict[str, str]) -> int:
        """특정 슬라이드의 텍스트를 치환한다.

        Args:
            slide_index: 현재 프레젠테이션 기준 슬라이드 인덱스
            replacements: {old_text: new_text} 딕셔너리

        Returns:
            치환된 횟수
        """
        if slide_index >= len(self._prs.slides):
            raise IndexError(f"Slide index {slide_index} out of range")

        slide = self._prs.slides[slide_index]
        sub = TextSubstitutor(slide)
        return sub.replace_all(replacements)

    def substitute_all(self, replacements: dict[str, str]) -> int:
        """모든 슬라이드에서 텍스트를 치환한다.

        Args:
            replacements: {old_text: new_text} 딕셔너리

        Returns:
            총 치환 횟수
        """
        total = 0
        for i in range(len(self._prs.slides)):
            total += self.substitute(i, replacements)
        return total

    def set_slide_size(self, width_inches: float, height_inches: float) -> "TemplateEditor":
        """슬라이드 크기를 변경한다 (EMU 단위 자동 변환).

        Args:
            width_inches: 너비 (인치)
            height_inches: 높이 (인치)

        Returns:
            self (체이닝용)
        """
        from pptx.util import Inches
        self._prs.slide_width = Inches(width_inches)
        self._prs.slide_height = Inches(height_inches)
        return self

    def save(self, output_path: str | Path) -> Path:
        """편집된 프레젠테이션을 저장한다.

        Args:
            output_path: 저장 경로

        Returns:
            저장된 파일 경로
        """
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(out))
        return out

    def merge_into(self, target_prs: Presentation) -> list:
        """편집된 슬라이드를 다른 프레젠테이션에 추가한다.

        원본 직접 편집 방식이므로, 이 메서드는 현재 편집 중인 프레젠테이션을
        먼저 임시 저장 → 다시 열어서 슬라이드를 복제하는 하이브리드 방식.
        단일 템플릿에서 여러 슬라이드를 가져올 때 사용.

        Args:
            target_prs: 대상 프레젠테이션

        Returns:
            추가된 슬라이드 리스트
        """
        # 임시 저장 후 SlideCloner로 복제 (편집 완료 상태이므로 관계 정리됨)
        tmp_path = Path(self._tmp_dir) / "_merge_tmp.pptx"
        self._prs.save(str(tmp_path))

        from .cloner import SlideCloner
        cloner = SlideCloner(tmp_path)
        slides = []
        for i in range(len(self._prs.slides)):
            slide = cloner.clone_slide(target_prs, i)
            slides.append(slide)
        return slides

    def _delete_slide(self, index: int) -> None:
        """python-pptx에서 공식 지원하지 않는 슬라이드 삭제를 XML 직접 조작으로 수행."""
        slide = self._prs.slides[index]
        rId = None

        # presentation.xml에서 슬라이드 참조 제거
        prs_part = self._prs.part
        for rel_key, rel in prs_part.rels.items():
            if rel.target_part is slide.part:
                rId = rel_key
                break

        if rId is None:
            raise RuntimeError(f"Cannot find relationship for slide at index {index}")

        # sldIdLst에서 해당 슬라이드 제거
        prs_elm = prs_part._element
        ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
        sldIdLst = prs_elm.find(f"{{{ns}}}sldIdLst")

        if sldIdLst is not None:
            for sldId in list(sldIdLst):
                r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                if sldId.get(f"{{{r_ns}}}id") == rId:
                    sldIdLst.remove(sldId)
                    break

        # 관계 제거
        prs_part.rels.pop(rId)

    def cleanup(self) -> None:
        """임시 파일 정리."""
        try:
            shutil.rmtree(self._tmp_dir, ignore_errors=True)
        except Exception:
            pass

    def __enter__(self):
        return self

    def __exit__(self, *args):
        self.cleanup()

    def __del__(self):
        self.cleanup()


def edit_template(
    template_path: str | Path,
    keep_indices: list[int],
    replacements: dict[str, str] | None = None,
    output_path: str | Path | None = None,
) -> Path:
    """편의 함수 - 템플릿에서 특정 슬라이드만 추출하고 텍스트를 치환한다.

    Args:
        template_path: 원본 .pptx 경로
        keep_indices: 남길 슬라이드 인덱스
        replacements: 전체 슬라이드에 적용할 텍스트 치환 (선택)
        output_path: 저장 경로 (없으면 output/ 하위에 자동 생성)

    Returns:
        저장된 파일 경로
    """
    if output_path is None:
        output_path = Path("output") / f"edited_{Path(template_path).stem}.pptx"

    with TemplateEditor(template_path) as editor:
        editor.keep_slides(keep_indices)
        if replacements:
            editor.substitute_all(replacements)
        return editor.save(output_path)
