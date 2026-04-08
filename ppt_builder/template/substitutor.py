"""텍스트 치환 엔진 - 복제된 슬라이드의 placeholder를 실제 데이터로 교체한다.

사용법:
  substitutor = TextSubstitutor(slide)
  substitutor.replace_all({
      "{{title}}": "실제 제목",
      "{{bullet_1}}": "첫 번째 항목",
      "Header": "새 헤더",
  })
"""

from pptx.slide import Slide
from pptx.util import Pt
from pptx.dml.color import RGBColor


class TextSubstitutor:
    """슬라이드 내 모든 텍스트 placeholder를 치환한다."""

    def __init__(self, slide: Slide):
        self.slide = slide

    def replace_all(self, replacements: dict[str, str]) -> int:
        """모든 placeholder를 치환한다.

        Args:
            replacements: {placeholder: replacement} 딕셔너리

        Returns:
            치환된 횟수
        """
        count = 0
        for shape in self.slide.shapes:
            count += self._replace_in_shape(shape, replacements)
        return count

    def replace_text(self, old: str, new: str) -> int:
        """단일 텍스트를 치환한다."""
        return self.replace_all({old: new})

    def get_all_texts(self) -> list[str]:
        """슬라이드 내 모든 텍스트를 추출한다."""
        texts = []
        for shape in self.slide.shapes:
            texts.extend(self._extract_texts(shape))
        return texts

    def _replace_in_shape(self, shape, replacements: dict) -> int:
        """shape 내 텍스트를 치환한다. \\n은 여러 paragraph로 분리."""
        count = 0

        # 텍스트 프레임
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            replaced = run.text.replace(old, new)
                            if '\n' in replaced:
                                # 첫 줄은 현재 run에, 나머지는 새 paragraph
                                lines = replaced.split('\n')
                                run.text = lines[0]
                                # 나머지 줄을 새 paragraph로 추가
                                tf = shape.text_frame
                                for line in lines[1:]:
                                    new_p = tf.add_paragraph()
                                    new_p.text = line
                                    # 원본 paragraph의 폰트 스타일 복사
                                    try:
                                        new_p.font.size = paragraph.font.size
                                        new_p.font.bold = paragraph.font.bold
                                        new_p.font.name = paragraph.font.name
                                        if paragraph.font.color and paragraph.font.color.rgb:
                                            new_p.font.color.rgb = paragraph.font.color.rgb
                                    except (AttributeError, TypeError):
                                        pass
                            else:
                                run.text = replaced
                            count += 1

        # 테이블
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for old, new in replacements.items():
                                if old in run.text:
                                    run.text = run.text.replace(old, new)
                                    count += 1

        # 그룹 shape (재귀)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                for child in shape.shapes:
                    count += self._replace_in_shape(child, replacements)
            except AttributeError:
                pass

        return count

    def _extract_texts(self, shape) -> list[str]:
        """shape에서 텍스트를 추출한다."""
        texts = []
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    texts.append(p.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        if p.text.strip():
                            texts.append(p.text)
        return texts


class CellStyler:
    """테이블 셀 스타일을 동적으로 변경한다."""

    @staticmethod
    def highlight_cell(cell, color: RGBColor):
        """셀 배경색을 변경한다."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    @staticmethod
    def set_text_color(cell, color: RGBColor):
        """셀 텍스트 색상을 변경한다."""
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = color
