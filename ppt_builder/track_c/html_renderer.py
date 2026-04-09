"""HTML/Tailwind → PNG 렌더러 (Playwright Chromium 기반).

복잡한 도식(2×2 매트릭스, BCG, 워터폴 등)을 python-pptx로 그리기 어려운 경우
HTML+CSS로 작성해 PNG로 렌더한 뒤 슬라이드에 이미지로 삽입한다.

Track A의 content.py 렌더러는 현재 ImageComponent를 처리하지 않으므로
Track C가 직접 python-pptx로 image-slide pptx를 빌드한다 (build_image_slide_pptx).

회사 컬러:
    Orange: #FD5108 / #FE7C39 / #FFAA72
    Grey:   #A1A8B3 / #B5BCC4 / #CBD1D6
    Body:   #000000
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# 회사 컬러 팔레트 (project_ppt_generator.md 메모리 기반)
COLORS = {
    "accent": "#FD5108",
    "accent_med": "#FE7C39",
    "accent_light": "#FFAA72",
    "grey_dark": "#A1A8B3",
    "grey_med": "#B5BCC4",
    "grey_light": "#CBD1D6",
    "body": "#000000",
}


def html_to_png(
    html: str,
    output_path: Path,
    width: int = 1600,
    height: int = 1000,
    device_scale_factor: int = 2,
) -> Path:
    """HTML 문자열을 PNG로 렌더한다.

    Args:
        html: 완전한 HTML 문서 (Tailwind CDN 포함 권장)
        output_path: PNG 출력 경로
        width, height: 뷰포트 픽셀 (기본 1600×1000, 4:3 적합)
        device_scale_factor: 2x retina 렌더링 (인쇄 품질)

    Returns:
        생성된 PNG 경로
    """
    output_path = Path(output_path).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with sync_playwright() as p:
        browser = p.chromium.launch()
        try:
            context = browser.new_context(
                viewport={"width": width, "height": height},
                device_scale_factor=device_scale_factor,
            )
            page = context.new_page()
            page.set_content(html, wait_until="networkidle")
            page.screenshot(path=str(output_path), full_page=False, omit_background=False)
        finally:
            browser.close()

    return output_path


# ============================================================
# 템플릿: 2×2 매트릭스 (BCG 스타일)
# ============================================================

_BCG_TEMPLATE = """<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8">
<title>2x2 Matrix</title>
<script src="https://cdn.tailwindcss.com"></script>
<style>
  body {{ font-family: 'Pretendard', 'Malgun Gothic', sans-serif; }}
  .quad {{ position: relative; }}
  .quad-label {{ position: absolute; top: 0.75rem; left: 0.75rem; font-weight: 700; font-size: 0.875rem; opacity: 0.7; }}
</style>
</head>
<body class="bg-white p-12">
  <h1 class="text-3xl font-bold text-black mb-2">{title}</h1>
  <p class="text-base text-gray-600 mb-8">{subtitle}</p>

  <div class="relative" style="height: 720px;">
    <!-- Y axis label -->
    <div class="absolute left-0 top-1/2 -translate-y-1/2 -translate-x-12 -rotate-90 origin-center text-sm font-semibold text-gray-700 whitespace-nowrap">{y_axis_label}</div>
    <!-- Y high/low -->
    <div class="absolute left-2 top-2 text-xs font-bold text-gray-500">↑ {y_high}</div>
    <div class="absolute left-2 bottom-2 text-xs font-bold text-gray-500">↓ {y_low}</div>

    <!-- X axis label -->
    <div class="absolute left-1/2 -bottom-10 -translate-x-1/2 text-sm font-semibold text-gray-700">{x_axis_label}</div>
    <!-- X low/high -->
    <div class="absolute -bottom-6 left-2 text-xs font-bold text-gray-500">← {x_low}</div>
    <div class="absolute -bottom-6 right-2 text-xs font-bold text-gray-500">{x_high} →</div>

    <!-- 2x2 grid -->
    <div class="grid grid-cols-2 grid-rows-2 gap-1 ml-12" style="height: 100%;">
      <!-- Top-left (low x, high y) -->
      <div class="quad p-6 border-2" style="background-color: #FFAA72; border-color: #FD5108;">
        <div class="quad-label">{q1_label}</div>
        <div class="mt-8">
          <h3 class="text-xl font-bold text-black mb-3">{q1_title}</h3>
          <ul class="text-sm text-black space-y-2">{q1_items}</ul>
        </div>
      </div>
      <!-- Top-right (high x, high y) -->
      <div class="quad p-6 border-2" style="background-color: #FD5108; color: white; border-color: #FD5108;">
        <div class="quad-label" style="color: white;">{q2_label}</div>
        <div class="mt-8">
          <h3 class="text-xl font-bold mb-3">{q2_title}</h3>
          <ul class="text-sm space-y-2">{q2_items}</ul>
        </div>
      </div>
      <!-- Bottom-left (low x, low y) -->
      <div class="quad p-6 border-2" style="background-color: #CBD1D6; border-color: #A1A8B3;">
        <div class="quad-label">{q3_label}</div>
        <div class="mt-8">
          <h3 class="text-xl font-bold text-black mb-3">{q3_title}</h3>
          <ul class="text-sm text-black space-y-2">{q3_items}</ul>
        </div>
      </div>
      <!-- Bottom-right (high x, low y) -->
      <div class="quad p-6 border-2" style="background-color: #B5BCC4; border-color: #A1A8B3;">
        <div class="quad-label">{q4_label}</div>
        <div class="mt-8">
          <h3 class="text-xl font-bold text-black mb-3">{q4_title}</h3>
          <ul class="text-sm text-black space-y-2">{q4_items}</ul>
        </div>
      </div>
    </div>
  </div>
</body>
</html>
"""


def render_2x2_matrix(
    output_path: Path,
    title: str,
    quadrants: dict[str, dict[str, Any]],
    subtitle: str = "",
    x_axis_label: str = "",
    x_low: str = "Low",
    x_high: str = "High",
    y_axis_label: str = "",
    y_low: str = "Low",
    y_high: str = "High",
    width: int = 1600,
    height: int = 1000,
) -> Path:
    """2×2 매트릭스(BCG 스타일)를 렌더링.

    Args:
        output_path: PNG 출력 경로
        title: 매트릭스 상단 제목 (액션 타이틀)
        quadrants: {
            "top_left":     {"label": "01", "title": "...", "items": ["..."]},
            "top_right":    {"label": "02", "title": "...", "items": ["..."]},
            "bottom_left":  {"label": "03", "title": "...", "items": ["..."]},
            "bottom_right": {"label": "04", "title": "...", "items": ["..."]},
        }
        subtitle: 부제목
        x/y_axis_label: 축 라벨
        x/y_low/high: 축 끝 라벨
    """

    def _items_html(items: list[str]) -> str:
        return "".join(f"<li>• {item}</li>" for item in items)

    def _quad(key: str) -> dict[str, Any]:
        q = quadrants.get(key, {})
        return {
            "label": q.get("label", ""),
            "title": q.get("title", ""),
            "items_html": _items_html(q.get("items", [])),
        }

    tl = _quad("top_left")
    tr = _quad("top_right")
    bl = _quad("bottom_left")
    br = _quad("bottom_right")

    html = _BCG_TEMPLATE.format(
        title=title,
        subtitle=subtitle,
        x_axis_label=x_axis_label,
        x_low=x_low,
        x_high=x_high,
        y_axis_label=y_axis_label,
        y_low=y_low,
        y_high=y_high,
        q1_label=tl["label"], q1_title=tl["title"], q1_items=tl["items_html"],
        q2_label=tr["label"], q2_title=tr["title"], q2_items=tr["items_html"],
        q3_label=bl["label"], q3_title=bl["title"], q3_items=bl["items_html"],
        q4_label=br["label"], q4_title=br["title"], q4_items=br["items_html"],
    )

    return html_to_png(html, output_path, width=width, height=height)


# ============================================================
# 템플릿: 워터폴 차트
# ============================================================

_WATERFALL_TEMPLATE = """<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8">
<title>Waterfall Chart</title>
<script src="https://cdn.tailwindcss.com"></script>
<style>
  body {{ font-family: 'Pretendard', 'Malgun Gothic', sans-serif; }}
</style>
</head>
<body class="bg-white p-12">
  <h1 class="text-3xl font-bold text-black mb-2">{title}</h1>
  <p class="text-base text-gray-600 mb-8">{subtitle}</p>

  <svg viewBox="0 0 1500 720" xmlns="http://www.w3.org/2000/svg" style="width: 100%; height: 720px;">
    <!-- baseline -->
    <line x1="80" y1="620" x2="1480" y2="620" stroke="#A1A8B3" stroke-width="2"/>
    <!-- y axis label -->
    <text x="40" y="320" font-size="14" font-weight="700" fill="#4b5563" text-anchor="middle" transform="rotate(-90, 40, 320)">{y_label}</text>
    {bars_svg}
  </svg>

  <p class="text-xs text-gray-500 mt-2">{footnote}</p>
</body>
</html>
"""


def render_waterfall(
    output_path: Path,
    title: str,
    bars: list[dict[str, Any]],
    subtitle: str = "",
    y_label: str = "값",
    footnote: str = "",
    width: int = 1600,
    height: int = 1000,
) -> Path:
    """워터폴 차트 렌더링.

    Args:
        bars: [
            {"label": "시작", "value": 100, "type": "total"},
            {"label": "신규 매출", "value": 30, "type": "increase"},
            {"label": "이탈", "value": -15, "type": "decrease"},
            {"label": "확장", "value": 20, "type": "increase"},
            {"label": "최종", "value": 135, "type": "total"},  # 자동 계산값과 같아야 함
        ]
        type: total (시작/최종), increase (+), decrease (-)
    """
    # 좌표 계산
    n = len(bars)
    x_pad = 80
    chart_w = 1500 - 2 * x_pad
    bar_w = chart_w / (n * 1.5)
    gap = bar_w * 0.5

    # 누적 합과 max 값
    cumulative = 0.0
    cumulatives: list[tuple[float, float]] = []  # (start, end) per bar
    for b in bars:
        if b["type"] == "total":
            cumulatives.append((0.0, float(b["value"])))
            cumulative = float(b["value"])
        else:
            new_cum = cumulative + float(b["value"])
            cumulatives.append((cumulative, new_cum))
            cumulative = new_cum

    max_val = max(max(s, e) for s, e in cumulatives)
    chart_h = 500
    y_base = 620

    def y_of(v: float) -> float:
        return y_base - (v / max_val) * chart_h

    # SVG bars
    svg_parts: list[str] = []
    color_map = {
        "total": "#000000",
        "increase": "#FD5108",
        "decrease": "#A1A8B3",
    }

    for i, b in enumerate(bars):
        start, end = cumulatives[i]
        x = x_pad + i * (bar_w + gap)
        top = y_of(max(start, end))
        bottom = y_of(min(start, end))
        h = max(bottom - top, 4)
        color = color_map.get(b["type"], "#A1A8B3")

        svg_parts.append(
            f'<rect x="{x}" y="{top}" width="{bar_w}" height="{h}" fill="{color}"/>'
        )
        # value label above bar
        sign = "+" if b["type"] == "increase" else ("-" if b["type"] == "decrease" else "")
        val_text = f"{sign}{abs(int(b['value'])):,}"
        svg_parts.append(
            f'<text x="{x + bar_w/2}" y="{top - 8}" font-size="14" font-weight="700" fill="#000" text-anchor="middle">{val_text}</text>'
        )
        # label below baseline
        svg_parts.append(
            f'<text x="{x + bar_w/2}" y="{y_base + 24}" font-size="13" font-weight="600" fill="#374151" text-anchor="middle">{b["label"]}</text>'
        )
        # cumulative label below name
        svg_parts.append(
            f'<text x="{x + bar_w/2}" y="{y_base + 42}" font-size="11" fill="#9ca3af" text-anchor="middle">누적 {int(end):,}</text>'
        )

        # connector line to next bar
        if i < n - 1:
            next_start, _ = cumulatives[i + 1]
            connector_y = y_of(end)
            x_end = x_pad + (i + 1) * (bar_w + gap)
            svg_parts.append(
                f'<line x1="{x + bar_w}" y1="{connector_y}" x2="{x_end}" y2="{connector_y}" stroke="#9ca3af" stroke-width="1.5" stroke-dasharray="3,3"/>'
            )

    html = _WATERFALL_TEMPLATE.format(
        title=title,
        subtitle=subtitle,
        y_label=y_label,
        footnote=footnote,
        bars_svg="\n    ".join(svg_parts),
    )

    return html_to_png(html, output_path, width=width, height=height)


# ============================================================
# 템플릿: McKinsey Horizon 3
# ============================================================

_HORIZON3_TEMPLATE = """<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8">
<title>McKinsey Horizon 3</title>
<script src="https://cdn.tailwindcss.com"></script>
<style>
  body {{ font-family: 'Pretendard', 'Malgun Gothic', sans-serif; }}
</style>
</head>
<body class="bg-white p-12">
  <h1 class="text-3xl font-bold text-black mb-2">{title}</h1>
  <p class="text-base text-gray-600 mb-6">{subtitle}</p>

  <svg viewBox="0 0 1500 700" xmlns="http://www.w3.org/2000/svg" style="width: 100%; height: 700px;">
    <!-- axes -->
    <line x1="80" y1="600" x2="1480" y2="600" stroke="#374151" stroke-width="2"/>
    <line x1="80" y1="80" x2="80" y2="600" stroke="#374151" stroke-width="2"/>
    <text x="780" y="650" font-size="14" font-weight="700" fill="#374151" text-anchor="middle">시간 →</text>
    <text x="40" y="340" font-size="14" font-weight="700" fill="#374151" text-anchor="middle" transform="rotate(-90, 40, 340)">가치/임팩트 →</text>

    <!-- Horizon 1: Core (수성) — opens immediately, plateaus then declines -->
    <path d="M 80 350 Q 250 200, 500 240 Q 750 280, 1000 360 Q 1250 440, 1480 520"
          stroke="#FD5108" stroke-width="6" fill="none" opacity="0.95"/>
    <text x="270" y="195" font-size="20" font-weight="800" fill="#FD5108">H1 — Core</text>
    <text x="270" y="220" font-size="13" font-weight="600" fill="#374151">{h1_label}</text>

    <!-- Horizon 2: Emerging — starts later, peaks higher -->
    <path d="M 350 580 Q 600 480, 850 320 Q 1100 200, 1350 180"
          stroke="#FE7C39" stroke-width="6" fill="none" opacity="0.85"/>
    <text x="700" y="305" font-size="20" font-weight="800" fill="#FE7C39">H2 — Emerging</text>
    <text x="700" y="330" font-size="13" font-weight="600" fill="#374151">{h2_label}</text>

    <!-- Horizon 3: Future — starts even later, highest potential -->
    <path d="M 700 590 Q 950 540, 1200 380 Q 1380 250, 1480 120"
          stroke="#A1A8B3" stroke-width="6" fill="none" stroke-dasharray="8,6" opacity="0.85"/>
    <text x="1080" y="425" font-size="20" font-weight="800" fill="#6b7280">H3 — Future</text>
    <text x="1080" y="450" font-size="13" font-weight="600" fill="#374151">{h3_label}</text>
  </svg>

  <div class="grid grid-cols-3 gap-4 mt-2">
    <div class="border-l-4 p-4" style="border-color: #FD5108;">
      <h3 class="font-bold text-lg" style="color: #FD5108;">H1 · {h1_period}</h3>
      <p class="text-xs text-gray-500 mb-2">{h1_subtitle}</p>
      <ul class="text-sm text-gray-800 space-y-1">{h1_items}</ul>
    </div>
    <div class="border-l-4 p-4" style="border-color: #FE7C39;">
      <h3 class="font-bold text-lg" style="color: #FE7C39;">H2 · {h2_period}</h3>
      <p class="text-xs text-gray-500 mb-2">{h2_subtitle}</p>
      <ul class="text-sm text-gray-800 space-y-1">{h2_items}</ul>
    </div>
    <div class="border-l-4 p-4" style="border-color: #A1A8B3;">
      <h3 class="font-bold text-lg text-gray-700">H3 · {h3_period}</h3>
      <p class="text-xs text-gray-500 mb-2">{h3_subtitle}</p>
      <ul class="text-sm text-gray-800 space-y-1">{h3_items}</ul>
    </div>
  </div>
</body>
</html>
"""


def render_horizon3(
    output_path: Path,
    title: str,
    horizons: dict[str, dict[str, Any]],
    subtitle: str = "",
    width: int = 1600,
    height: int = 1100,
) -> Path:
    """McKinsey Horizon 3 모델.

    horizons: {
        "h1": {"label": "축약 라벨", "period": "0-12개월", "subtitle": "...", "items": [...]},
        "h2": {"label": "...",     "period": "12-24개월", "subtitle": "...", "items": [...]},
        "h3": {"label": "...",     "period": "24개월+",   "subtitle": "...", "items": [...]},
    }
    """
    def _items_html(items: list[str]) -> str:
        return "".join(f"<li>• {x}</li>" for x in items)

    def _h(key: str) -> dict[str, Any]:
        h = horizons.get(key, {})
        return {
            "label": h.get("label", ""),
            "period": h.get("period", ""),
            "subtitle": h.get("subtitle", ""),
            "items_html": _items_html(h.get("items", [])),
        }

    h1 = _h("h1")
    h2 = _h("h2")
    h3 = _h("h3")

    html = _HORIZON3_TEMPLATE.format(
        title=title,
        subtitle=subtitle,
        h1_label=h1["label"], h1_period=h1["period"], h1_subtitle=h1["subtitle"], h1_items=h1["items_html"],
        h2_label=h2["label"], h2_period=h2["period"], h2_subtitle=h2["subtitle"], h2_items=h2["items_html"],
        h3_label=h3["label"], h3_period=h3["period"], h3_subtitle=h3["subtitle"], h3_items=h3["items_html"],
    )

    return html_to_png(html, output_path, width=width, height=height)


# ============================================================
# 템플릿: Porter 가치사슬
# ============================================================

_VALUE_CHAIN_TEMPLATE = """<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8">
<title>Porter Value Chain</title>
<script src="https://cdn.tailwindcss.com"></script>
<style>
  body {{ font-family: 'Pretendard', 'Malgun Gothic', sans-serif; }}
  .arrow {{
    clip-path: polygon(0 0, calc(100% - 30px) 0, 100% 50%, calc(100% - 30px) 100%, 0 100%);
  }}
</style>
</head>
<body class="bg-white p-12">
  <h1 class="text-3xl font-bold text-black mb-2">{title}</h1>
  <p class="text-base text-gray-600 mb-8">{subtitle}</p>

  <div class="flex" style="height: 720px;">
    <div class="flex-1 flex flex-col">
      <!-- Support activities -->
      <div class="mb-2">
        <div class="text-sm font-bold text-gray-700 mb-2">지원 활동 (Support Activities)</div>
        {support_rows}
      </div>

      <!-- Primary activities (arrow row) -->
      <div class="mt-2">
        <div class="text-sm font-bold text-gray-700 mb-2">본원적 활동 (Primary Activities)</div>
        <div class="flex gap-1">
          {primary_cols}
        </div>
      </div>

      <!-- Bottom: Margin highlight -->
      <div class="mt-6 p-4 border-2" style="background-color: #FFAA72; border-color: #FD5108;">
        <div class="text-xs font-bold uppercase mb-1" style="color: #FD5108;">차별화 원천 → 마진</div>
        <div class="text-base font-bold text-black">{margin_message}</div>
      </div>
    </div>
  </div>
</body>
</html>
"""


def render_value_chain(
    output_path: Path,
    title: str,
    support_activities: list[dict[str, Any]],
    primary_activities: list[dict[str, Any]],
    margin_message: str = "",
    subtitle: str = "",
    width: int = 1600,
    height: int = 1100,
) -> Path:
    """Porter 가치사슬.

    Args:
        support_activities: [{"name": "Firm Infrastructure", "items": ["..."], "highlight": False}, ...]
            보통 4개 (Firm Infra / HR / Technology / Procurement)
        primary_activities: [{"name": "Inbound Logistics", "items": ["..."], "highlight": False}, ...]
            보통 5개 (Inbound / Operations / Outbound / M&S / Service)
        margin_message: 마진 박스에 표시할 핵심 메시지
    """
    def _items_html(items: list[str]) -> str:
        return "".join(f'<li class="text-xs text-gray-700 leading-tight">• {x}</li>' for x in items)

    # Support rows (수평 박스 4개 stacked)
    support_html = []
    for s in support_activities:
        bg = "#FFAA72" if s.get("highlight") else "#F3F4F6"
        border = "#FD5108" if s.get("highlight") else "#D1D5DB"
        support_html.append(
            f'<div class="border-2 px-4 py-2 mb-1" style="background-color: {bg}; border-color: {border};">'
            f'  <div class="font-bold text-sm text-black mb-1">{s["name"]}</div>'
            f'  <ul class="grid grid-cols-3 gap-x-3">{_items_html(s.get("items", []))}</ul>'
            f'</div>'
        )

    # Primary cols (수평 화살표 형태)
    primary_html = []
    n_cols = len(primary_activities)
    for i, p in enumerate(primary_activities):
        is_last = (i == n_cols - 1)
        bg = "#FD5108" if p.get("highlight") else ("#FE7C39" if i % 2 == 0 else "#FFAA72")
        text_color = "white" if p.get("highlight") else "#000"
        arrow_class = "arrow" if is_last else ""
        # f-string 안 백슬래시 회피: 미리 컴파일
        items_li = "".join(
            '<li class="text-xs leading-tight" style="color: ' + text_color + ';">• ' + str(x) + '</li>'
            for x in p.get("items", [])
        )
        primary_html.append(
            f'<div class="flex-1 p-3 {arrow_class}" style="background-color: {bg}; min-height: 220px;">'
            f'  <div class="font-bold text-sm mb-2" style="color: {text_color};">{p["name"]}</div>'
            f'  <ul class="space-y-1">{items_li}</ul>'
            f'</div>'
        )

    html = _VALUE_CHAIN_TEMPLATE.format(
        title=title,
        subtitle=subtitle,
        support_rows="\n        ".join(support_html),
        primary_cols="\n          ".join(primary_html),
        margin_message=margin_message or "본원·지원 활동 전체에 걸친 차별화 → 시장 대비 우위 마진 확보",
    )

    return html_to_png(html, output_path, width=width, height=height)


# ============================================================
# Image Slide PPTX 빌더 (Track A ImageComponent 미구현 우회)
# ============================================================

def build_image_slide_pptx(
    output_pptx: Path,
    image_path: Path,
    title: str = "",
    breadcrumb: str = "",
    footnote: str = "",
) -> Path:
    """이미지 1장을 메인 콘텐츠로 하는 1-슬라이드 pptx를 빌드한다.

    4:3 (10×7.5인치) 슬라이드:
        - 상단 0.5인치: 액션 타이틀 (회사 컬러 검정 18pt 볼드)
        - 우상단: breadcrumb (작은 회색)
        - 본문 영역 0.7~6.7 인치: 이미지 (1:1.6 비율로 폭 자동 맞춤)
        - 하단 0.05인치: 회사 컬러 강조선
        - 최하단: footnote (작은 회색)

    Track A content 렌더러가 ImageComponent를 처리하지 않아 Track C가 직접 빌드.
    """
    output_pptx = Path(output_pptx).resolve()
    image_path = Path(image_path).resolve()
    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    if not image_path.exists():
        raise FileNotFoundError(f"image not found: {image_path}")

    prs = Presentation()
    prs.slide_width = Inches(10)   # 4:3
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # 빈 레이아웃
    slide = prs.slides.add_slide(blank_layout)

    accent = RGBColor(0xFD, 0x51, 0x08)
    grey_dark = RGBColor(0xA1, 0xA8, 0xB3)
    black = RGBColor(0x00, 0x00, 0x00)

    # 1. Title
    if title:
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(8.5), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = black

    # 2. Breadcrumb (우상단)
    if breadcrumb:
        bb = slide.shapes.add_textbox(Inches(7.5), Inches(0.2), Inches(2.4), Inches(0.3))
        bf = bb.text_frame
        bp = bf.paragraphs[0]
        bp.text = breadcrumb
        bp.font.size = Pt(9)
        bp.font.color.rgb = grey_dark
        bp.alignment = 2  # right

    # 3. Image (본문)
    img_left = Inches(0.5)
    img_top = Inches(0.95)
    img_width = Inches(9.0)  # 좌우 0.5인치 마진
    slide.shapes.add_picture(
        str(image_path),
        img_left,
        img_top,
        width=img_width,
    )

    # 4. Accent line (하단 구분)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4),
        Inches(7.05),
        Inches(9.2),
        Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()

    # 5. Footnote
    if footnote:
        fb = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(9.2), Inches(0.3))
        ff = fb.text_frame
        fp = ff.paragraphs[0]
        fp.text = footnote
        fp.font.size = Pt(8)
        fp.font.color.rgb = grey_dark

    prs.save(str(output_pptx))
    return output_pptx
