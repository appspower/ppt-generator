"""시각적 PPT 검증 — 정적 분석 + (가능 시) PDF 변환 검사.

evaluate.py가 잡지 못하는 시각 문제를 검출:
- 슬라이드 경계 밖으로 넘어간 shape (overflow)
- 테이블 자동 확장으로 인한 넘침 (행 합 > 할당 높이)
- 텍스트 박스 간 겹침 (overlap)
- 하단 빈 공간 과다
- 카드/박스 텍스트 추정 오버플로

PDF 변환은 Windows + PowerPoint COM이 있을 때만 시도하며, 없으면
정적 분석만 수행한다 — 환경에 무관하게 절대 예외를 던지지 않음.

기존 evaluate.py와 독립적으로 동작하여 어느 호출 흐름도 깨지 않는다.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation


# ============================================================
# Public API
# ============================================================

class VisualCheckUnavailable(Exception):
    """PDF 변환 환경이 갖춰지지 않은 경우 (Windows + PowerPoint 필요).

    호출자는 이 예외를 잡아서 정적 검사만으로 진행할지 결정한다.
    """


@dataclass
class VisualReport:
    pdf_path: Optional[Path] = None
    issues: list[str] = field(default_factory=list)
    metrics: dict = field(default_factory=dict)
    pdf_available: bool = False

    @property
    def passed(self) -> bool:
        return not self.issues

    def severity_count(self) -> dict[str, int]:
        counts = {"OVERFLOW": 0, "TABLE": 0, "OVERLAP": 0, "EMPTY": 0, "OTHER": 0}
        for iss in self.issues:
            tag = iss.split(":", 1)[0]
            if tag.startswith("OVERFLOW"):
                counts["OVERFLOW"] += 1
            elif tag.startswith("TABLE"):
                counts["TABLE"] += 1
            elif tag.startswith("TEXT_OVERLAP") or tag.startswith("OVERLAP"):
                counts["OVERLAP"] += 1
            elif tag.startswith("EMPTY"):
                counts["EMPTY"] += 1
            else:
                counts["OTHER"] += 1
        return counts


def validate_visual(
    pptx_path: str | Path,
    *,
    convert_pdf: bool = True,
    pdf_path: str | Path | None = None,
) -> VisualReport:
    """PPTX 파일의 시각 문제를 검출한다.

    Args:
        pptx_path: 검사할 PPTX 파일 경로.
        convert_pdf: True면 PDF 변환을 시도한다 (PowerPoint COM 필요).
                     실패해도 정적 분석은 항상 수행한다.
        pdf_path: 출력 PDF 경로. None이면 pptx 옆에 같은 이름으로 생성.

    Returns:
        VisualReport — 발견된 issues 리스트, 측정치, PDF 경로.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        return VisualReport(issues=[f"FILE_NOT_FOUND: {pptx_path}"])

    issues: list[str] = []
    metrics: dict = {}

    # 정적 분석 — 환경 무관, 항상 실행
    static_issues, static_metrics = _check_pptx_static(pptx_path)
    issues.extend(static_issues)
    metrics.update(static_metrics)

    # PDF 변환 — 가능하면 시도
    pdf_out: Optional[Path] = None
    pdf_available = False
    if convert_pdf:
        try:
            pdf_out = convert_pptx_to_pdf(pptx_path, pdf_path)
            pdf_available = True
            metrics["pdf_path"] = str(pdf_out)
        except VisualCheckUnavailable as e:
            metrics["pdf_skip_reason"] = str(e)
        except Exception as e:
            issues.append(f"PDF_CONVERSION_FAILED: {e}")

    return VisualReport(
        pdf_path=pdf_out,
        issues=issues,
        metrics=metrics,
        pdf_available=pdf_available,
    )


def convert_pptx_to_pdf(
    pptx_path: str | Path,
    pdf_path: str | Path | None = None,
) -> Path:
    """PowerPoint COM을 사용해 PPTX → PDF 변환.

    Windows + Microsoft PowerPoint가 설치되어 있어야 한다.
    그 외 환경에서는 VisualCheckUnavailable 예외를 던진다.
    """
    pptx_path = Path(pptx_path)
    if pdf_path is None:
        pdf_path = pptx_path.with_suffix(".pdf")
    pdf_path = Path(pdf_path)

    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore
    except ImportError as e:
        raise VisualCheckUnavailable(
            "pywin32가 필요합니다 (Windows + PowerPoint 환경에서만 동작)."
        ) from e

    pythoncom.CoInitialize()
    ppt = None
    try:
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        src = os.path.abspath(str(pptx_path))
        dst = os.path.abspath(str(pdf_path))
        pres = ppt.Presentations.Open(src, WithWindow=False)
        try:
            pres.SaveAs(dst, 32)  # 32 = ppSaveAsPDF
        finally:
            pres.Close()
    except Exception as e:
        raise VisualCheckUnavailable(
            f"PowerPoint COM 호출 실패 (PowerPoint 미설치?): {e}"
        ) from e
    finally:
        if ppt is not None:
            try:
                ppt.Quit()
            except Exception:
                pass

    return pdf_path


# ============================================================
# Static analysis (PDF 없이도 잡을 수 있는 시각 문제)
# ============================================================

def _check_pptx_static(pptx_path: Path) -> tuple[list[str], dict]:
    issues: list[str] = []
    metrics: dict = {}

    prs = Presentation(str(pptx_path))
    slide_w_in = prs.slide_width / 914400
    slide_h_in = prs.slide_height / 914400
    metrics["slide_size"] = (round(slide_w_in, 2), round(slide_h_in, 2))
    metrics["slide_count"] = len(prs.slides)

    for si, slide in enumerate(prs.slides):
        sn = si + 1
        slide_metrics: dict = {"shape_count": len(slide.shapes)}

        # 1. 슬라이드 경계 overflow
        max_bottom = 0.0
        for idx, shape in enumerate(slide.shapes):
            if shape.left is None or shape.top is None:
                continue
            l = shape.left / 914400
            t = shape.top / 914400
            w = (shape.width or 0) / 914400
            h = (shape.height or 0) / 914400
            r = l + w
            b = t + h
            if b > max_bottom:
                max_bottom = b

            if r > slide_w_in + 0.1:
                issues.append(
                    f"OVERFLOW_RIGHT: Slide {sn} shape {idx} "
                    f"(right={r:.2f}\" > slide_w {slide_w_in:.2f}\")"
                )
            if b > slide_h_in + 0.1:
                issues.append(
                    f"OVERFLOW_BOTTOM: Slide {sn} shape {idx} "
                    f"(bottom={b:.2f}\" > slide_h {slide_h_in:.2f}\")"
                )

        slide_metrics["max_bottom"] = round(max_bottom, 2)
        empty_bottom = slide_h_in - max_bottom
        slide_metrics["empty_bottom"] = round(empty_bottom, 2)
        if empty_bottom > 1.5:
            issues.append(
                f"EMPTY_BOTTOM: Slide {sn} has {empty_bottom:.2f}\" empty space at bottom"
            )

        # 2. 테이블 자동 확장 감지 (행 높이 합 > 할당 높이)
        for idx, shape in enumerate(slide.shapes):
            tbl = _get_table(shape)
            if tbl is None:
                continue
            try:
                row_total = sum(r.height for r in tbl.rows) / 914400
                alloc = (shape.height or 0) / 914400
                if alloc > 0 and row_total > alloc + 0.1:
                    issues.append(
                        f"TABLE_OVERFLOW: Slide {sn} table {idx} "
                        f"(rows={row_total:.2f}\" > alloc={alloc:.2f}\")"
                    )
            except Exception:
                pass

        # 3. 텍스트 박스 간 겹침
        text_shapes = []
        for idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt:
                continue
            if shape.left is None or shape.top is None:
                continue
            l = shape.left / 914400
            t = shape.top / 914400
            w = (shape.width or 0) / 914400
            h = (shape.height or 0) / 914400
            text_shapes.append((idx, l, t, w, h))

        for i, (i1, l1, t1, w1, h1) in enumerate(text_shapes):
            for i2, l2, t2, w2, h2 in text_shapes[i + 1 :]:
                ox = max(0, min(l1 + w1, l2 + w2) - max(l1, l2))
                oy = max(0, min(t1 + h1, t2 + h2) - max(t1, t2))
                area = ox * oy
                # 작은 겹침은 무시 (icon anchor 등 의도된 경우)
                if area > 0.5:
                    issues.append(
                        f"TEXT_OVERLAP: Slide {sn} shapes {i1}+{i2} "
                        f"overlap {area:.2f} sq-in"
                    )

        # 4. 텍스트 길이 vs 박스 크기 휴리스틱 (잘림 추정)
        for idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt or shape.width is None or shape.height is None:
                continue
            w_in = shape.width / 914400
            h_in = shape.height / 914400
            if w_in < 0.5 or h_in < 0.2:
                continue
            # 9pt 본문 기준 — 폭당 한글 문자수 ~ w_in * 11, 줄당 0.16"
            est_chars_per_line = max(1, int(w_in * 11))
            est_lines = max(1, len(txt) // est_chars_per_line + txt.count("\n"))
            est_height = est_lines * 0.18
            if est_height > h_in * 1.4 and len(txt) > 80:
                issues.append(
                    f"TEXT_TOO_DENSE: Slide {sn} shape {idx} "
                    f"(est {est_lines} lines, {len(txt)} chars in {w_in:.1f}x{h_in:.1f}\")"
                )

        metrics[f"slide_{sn}"] = slide_metrics

    return issues, metrics


def _get_table(shape):
    """python-pptx Shape에서 table 객체를 안전하게 가져온다."""
    try:
        if shape.has_table:
            return shape.table
    except Exception:
        pass
    try:
        # 일부 shape는 has_table 속성이 없음
        return shape.table
    except Exception:
        return None
