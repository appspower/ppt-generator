"""Phase A3 Step 1 — 5 벤치마크 시나리오 정밀 측정.

목적
----
Mode A 단독으로 5개 컨설팅 시나리오 PPT를 생성하고 정량 측정한다.
- A. role 매칭 성공률
- B. 슬롯 채움률 (~~ 잔존 비율)
- C. 텍스트 오버플로 빈도
- F. role 풀 고갈 (재사용 횟수)
(D PPTEval / E 시각 검증은 PNG 생성 후 별도 단계)

사용
----
python scripts/benchmark_5_scenarios.py [scenario_id]
  scenario_id 미지정 시 5개 모두 실행
"""
from __future__ import annotations

import json
import sys
import time
from collections import Counter
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation

from ppt_builder.template import edit_ops
from ppt_builder.template.editor import TemplateEditor


CATALOG_PATH = ROOT / "output" / "catalog" / "final_labels.json"
SKELETONS_PATH = ROOT / "output" / "catalog" / "skeletons.json"
SLOT_SCHEMAS_PATH = ROOT / "output" / "catalog" / "slot_schemas.json"
TEMPLATE_PATH = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
OUTPUT_ROOT = ROOT / "output" / "benchmark"


# ----------------------------------------------------------------------------
# Fallback rules (extended for all roles in skeletons)
# ----------------------------------------------------------------------------

ARCHETYPE_FALLBACK = {
    "opening": ["cover_divider"],
    "closing": ["cover_divider"],
    "agenda": ["vertical_list", "cards_3col", "left_title_right_body"],
    "divider": ["cover_divider"],
    "appendix": ["table_native", "vertical_list", "dense_grid"],
    "situation": ["cards_3col", "cards_2col", "left_title_right_body"],
    "complication": ["cards_3col", "left_title_right_body", "vertical_list"],
    "evidence": ["table_native", "chart_native", "dense_grid", "cards_3col"],
    "analysis": ["matrix_2x2", "table_native", "dense_grid"],
    "recommendation": ["cards_3col", "vertical_list", "flowchart"],
    "roadmap": ["roadmap", "timeline_h"],
    "benefit": ["cards_3col", "cards_2col", "vertical_list"],
    "risk": ["matrix_2x2", "table_native", "cards_3col"],
}


# ----------------------------------------------------------------------------
# Virtual content per scenario
# ----------------------------------------------------------------------------

SCENARIO_CONTENT = {
    "transformation_roadmap_10": {
        "scenario_name": "D철강 2030 넷제로 전환 로드맵",
        "skeleton_id": "transformation_roadmap_10",
        "content_by_role": {
            "opening": ["D철강 2030 넷제로 전환 로드맵"],
            "situation": ["2024년 배출량 1,200만톤 / 산업 평균 대비 +15%"],
            "complication": ["EU CBAM 2026 시행 / 비용 영향 연 800억원"],
            "analysis": ["현 배출 1,200만톤 vs 2030 목표 600만톤 / 갭 50%"],
            "recommendation": ["전기로 전환 + 수소환원제철 + 그린전력 PPA"],
            "roadmap": [
                "Phase 1 (2026-2027) 전기로 도입, 폐열회수 250만톤 감축",
                "Phase 2 (2028-2030) 수소환원 파일럿, 350만톤 추가 감축",
            ],
            "benefit": ["탄소비용 절감 4,200억원 / ESG 등급 A 진입"],
            "risk": ["수소 단가 변동 / 정책 후퇴 / 설비 투자 회수기간 8년"],
            "closing": ["2030 600만톤 / 2040 Net-Zero / 다음 단계 의사결정"],
        },
    },
    "consulting_proposal_30": {
        "scenario_name": "HD현대 SAP S/4HANA 전환 제안",
        "skeleton_id": "consulting_proposal_30",
        "content_by_role": {
            "opening": ["HD현대 SAP S/4HANA 전환 제안 — PwC Strategy"],
            "agenda": ["1. 현황 / 2. 과제 / 3. 솔루션 / 4. 로드맵 / 5. 효과 / 6. 위험"],
            "situation": [
                "ECC 6.0 EOS 2027년 / 7개 계열사 25개 인스턴스 분산 운영",
                "월결산 D+12일 / 데이터 정합성 이슈 분기 평균 3.2건",
            ],
            "complication": [
                "신성장(에너지/AI) 대응 ERP 한계 / 글로벌 RHQ 통합 불가",
                "Custom 개발 누적 8,400건 / 표준 회귀 어려움 ",
            ],
            "evidence": [
                "동종 그룹 SAP 전환 평균 ROI 18개월",
                "S/4HANA 도입사 결산속도 평균 47% 단축",
                "현업 인터뷰 142건 / Pain Point Top 5 식별",
                "글로벌 BMP 12개 영역 평가 (CFO/COO/CIO 90%)",
            ],
            "analysis": [
                "GAP 분석: 모듈 8/12 표준화 가능 / 4 모듈 재설계 필요",
                "데이터 거버넌스 성숙도 Level 2 → Level 4 전환",
                "조직역량 진단: SAP 전문가 사내 23명 (필요 80명)",
            ],
            "divider": ["Section 3. 솔루션 — Big Bang vs Phased"],
            "recommendation": [
                "Phased 접근: 본사 → 계열사 순 (24개월)",
                "Single Global Template 기반 7개사 일괄 도입",
                "RISE with SAP 패키지 + Cloud Edition 채택",
                "Fiori UX 표준화 + 모바일 결재 100%",
                "Master Data Governance 단일 ERP 통합",
                "Center of Excellence 30명 신설",
            ],
            "roadmap": [
                "Phase 1 (M1-M6): Discover + Design / Blueprint 확정",
                "Phase 2 (M7-M18): Build + Test / 본사 GoLive",
                "Phase 3 (M19-M24): Rollout 6개 계열사 + Stabilize",
            ],
            "benefit": [
                "결산 D+12 → D+5 / 운영비 연 240억원 절감",
                "표준 프로세스 80% / Custom 8,400 → 1,600건",
            ],
            "risk": [
                "데이터 마이그레이션 리스크 / Cleansing 사전 6개월 권장",
                "변화관리: 사용자 교육 시수 1인 평균 32시간",
            ],
            "closing": [
                "PwC + SAP Korea 통합 추진 / 24개월 단일 책임 운영",
                "다음 단계: Discovery Workshop 5일 / 5월 착수 권고",
            ],
            "appendix": ["참조: PwC SAP Reference Cases 2023-2025 (33건)"],
        },
    },
    "analysis_report_15": {
        "scenario_name": "2026 Q1 재무 성과 분석",
        "skeleton_id": "analysis_report_15",
        "content_by_role": {
            "opening": ["2026 Q1 재무 성과 분석 — Group CFO Office"],
            "agenda": ["1. 손익 / 2. 매출 분해 / 3. 비용 분석 / 4. 권고 / 5. 위험"],
            "situation": [
                "Q1 매출 4.2조원 (YoY +6.3%) / 영업이익 3,800억원 (YoY -8.1%)",
                "원가율 73.8% / 판관비율 16.2% / 영업이익률 9.0%",
            ],
            "evidence": [
                "사업부별 매출: 조선 +12% / 건설기계 +4% / 에너지 -2%",
                "원가 구성: 원자재 58% / 인건비 22% / 외주 11% / 기타 9%",
            ],
            "analysis": [
                "영업이익 감소 -8.1% 주요인: 원자재 +14% (철광석 +22%)",
                "환율 효과: 평균 1,335원 (전기 1,310원) → 매출 +210억원",
                "재고자산 회전율 6.8회 (전기 7.2회) → 재고 누적 신호",
                "EBITDA 4,950억원 / EBITDA 마진 11.8% (전년 12.4%)",
                "Cash Conversion Cycle 78일 (전기 71일) → 운전자본 압박",
            ],
            "recommendation": [
                "원가 절감: 원자재 장기 계약 비중 35% → 50%",
                "재고 관리: 안전재고 정책 재설정 / SKU 합리화 20%",
            ],
            "risk": ["철광석 가격 변동성 / 환율 약세 시 추가 손실 가능"],
            "closing": ["Q2 회복 시나리오 + 중기 가이던스 조정 권고"],
        },
    },
    "change_management_20": {
        "scenario_name": "그룹 조직개편 Change Management",
        "skeleton_id": "change_management_20",
        "content_by_role": {
            "opening": ["그룹 조직개편 Change Management 실행안"],
            "agenda": ["1. Why / 2. As-Is / 3. To-Be / 4. 실행 / 5. 효과 / 6. 위험"],
            "situation": [
                "조직 사일로 4개 BU / 의사결정 평균 32일 / 중복 기능 18건",
                "직원 만족도 6.4/10 (업계 평균 7.2) / 자발적 이직률 12.4%",
            ],
            "complication": [
                "신사업(AI/에너지전환) 기존 조직으로 대응 불가",
            ],
            "analysis": [
                "Best Practice 5개사 벤치마크: Matrix 3사 / Tribe 2사",
                "Span of Control 평균 4.2 (Best 7.1) → 관리층 과다",
            ],
            "recommendation": [
                "Function Pool + Squad 모델 도입 / 70:30 비율",
                "그룹 공통 기능 통합 (HR/Finance/IT) → CoE 3개",
                "임원 -22% / 팀장 -15% / 의사결정 단계 5→3",
                "성과 평가: 조직 KPI 60% + Squad KPI 40%",
            ],
            "roadmap": [
                "M1-M2: Vision 공유 + 리더 코칭 (전임원 80시간)",
                "M3-M4: 신조직 발족 + 인력 재배치 (자발 신청제)",
                "M5-M9: 안정화 + 성과 모니터링 / 분기 펄스 서베이",
            ],
            "evidence": [
                "Pulse 1차 조사 (n=2,418): 변화 수용도 64% / 우려 31%",
                "리더 워크숍 결과: 신모델 동의 78% / 보완 요청 영역 4개",
            ],
            "benefit": ["의사결정 32 → 14일 / 신사업 매출 비중 8 → 18% (3년)"],
            "risk": [
                "핵심인재 이탈 위험 / Stay Bonus + 경력 경로 명시 필수",
                "변화 피로 누적 / 6개월 차 펄스 서베이 결과 따라 조정",
            ],
            "closing": ["Day-1 5월 1일 / CEO 메시지 + 90일 Quick Win"],
        },
    },
    "executive_strategy_40": {
        "scenario_name": "그룹 2026-2030 중장기 전략",
        "skeleton_id": "executive_strategy_40",
        "content_by_role": {
            "opening": ["그룹 2026-2030 중장기 전략 — Board Review"],
            "agenda": ["1. 환경 / 2. 진단 / 3. 비전 / 4. 전략 / 5. 실행 / 6. 위험"],
            "situation": [
                "글로벌 매크로: 금리 5.2% / 인플레 둔화 / 지정학 변동성 ↑",
                "산업 변화: AI/에너지 전환 / 공급망 재편 / 규제 강화",
                "당사 위치: 매출 18조 / 영업이익 1.6조 / 순현금 3.2조",
            ],
            "complication": [
                "기존 사업 성장 한계: CAGR 3.1% (시장 2.8%)",
                "신성장 동력 부재: AI/에너지 매출 비중 4.8% (목표 25%)",
            ],
            "evidence": [
                "고객 NPS 42 (업계 38) / 브랜드 가치 10위권",
                "R&D 매출 비중 2.8% (글로벌 Top 5.4%)",
                "ESG 등급 A- / Scope 3 배출 미공개 → 개선 시급",
            ],
            "analysis": [
                "Portfolio 진단: Star 2 / Cash Cow 4 / Question 3 / Dog 2",
                "역량 갭: 디지털 인재 -40% / 글로벌 PM -55%",
            ],
            "divider": ["Vision 2030 — Sustainable Industrial Champion"],
            "recommendation": [
                "Vision: 매출 30조 / 영업이익률 12% / ESG AAA",
                "Pillar 1: Core 강화 (조선/건설기계/에너지)",
                "Pillar 2: AI 사업 신설 (산업용 AI Platform)",
                "Pillar 3: 에너지 전환 (수소/CCUS/그린전력)",
                "Pillar 4: Global Footprint (북미 30%/유럽 15%)",
                "Pillar 5: Talent (디지털 +2,000명 / 글로벌 +500명)",
                "Pillar 6: ESG Leadership (Scope 3 -40% by 2030)",
                "Pillar 7: M&A 5조원 규모 / 3건 이상 (Tech/Energy)",
                "Pillar 8: Capital Allocation (Buyback 5천억/년)",
            ],
            "roadmap": [
                "Phase 1 (2026): Foundation / Vision 공유 + 조직 개편",
                "Phase 2 (2027-2028): Acceleration / M&A 2건 + AI 출시",
                "Phase 3 (2029): Scale / 글로벌 확장 + 수소 상용화",
                "Phase 4 (2030): Lead / Vision KPI 달성 + 차기 수립",
            ],
            "benefit": [
                "매출 18 → 30조 / 영업이익 1.6 → 3.6조 (CAGR 9.0%)",
                "주주환원 누적 5조원 / 시가총액 2배 목표",
            ],
            "risk": [
                "M&A 실행 위험 / 통합 실패 시 손상차손 가능",
                "AI/수소 R&D 회수기간 7-10년 / 손익 변동성 ↑",
            ],
            "closing": [
                "Vision 2030 의결 요청 / Q3 KPI Dashboard 도입",
                "다음 단계: Pillar별 TF 발족 / 6월 Day-1 선포",
            ],
            "appendix": [
                "Reference: 글로벌 Industrial Champion 5개사 분석",
                "Reference: AI Platform 시장 규모 2030 1.4조 USD",
                "Reference: 수소 경제성 Break-even 2028 (수소 4,000원/kg)",
            ],
        },
    },
}


# ----------------------------------------------------------------------------
# Loaders
# ----------------------------------------------------------------------------

def load_catalog() -> list[dict]:
    with open(CATALOG_PATH, encoding="utf-8") as f:
        return json.load(f)["labels"]


def load_skeletons() -> dict:
    with open(SKELETONS_PATH, encoding="utf-8") as f:
        return json.load(f)


def load_slot_schemas() -> dict[str, dict]:
    """slide_index → flat_idx → schema dict."""
    with open(SLOT_SCHEMAS_PATH, encoding="utf-8") as f:
        rows = json.load(f)
    by_slide: dict[int, dict[int, dict]] = {}
    for r in rows:
        by_slide.setdefault(r["slide_index"], {})[r["flat_idx"]] = r
    return by_slide


# ----------------------------------------------------------------------------
# Retrieval
# ----------------------------------------------------------------------------

def candidates_for_role(
    labels: list[dict], role: str, used: set[int]
) -> tuple[list[dict], str]:
    direct = [
        l for l in labels
        if role in l.get("narrative_role", []) and l["slide_index"] not in used
    ]
    if direct:
        direct.sort(key=lambda l: l.get("overall_confidence", 0), reverse=True)
        return direct, "role"

    fallback_archs = ARCHETYPE_FALLBACK.get(role, [])
    fb = [
        l for l in labels
        if any(a in fallback_archs for a in l.get("archetype", []))
        and l["slide_index"] not in used
    ]
    if fb:
        fb.sort(key=lambda l: l.get("overall_confidence", 0), reverse=True)
        return fb, "archetype"

    return [], "none"


def select_deck(labels: list[dict], narrative: list[str]) -> list[dict]:
    """narrative 시퀀스대로 슬라이드 선정. 풀 고갈 시 reuse 허용."""
    used: set[int] = set()
    plan = []
    role_use_count: dict[str, int] = {}

    for role in narrative:
        cands, source = candidates_for_role(labels, role, used)

        if cands:
            chosen = cands[0]
            used.add(chosen["slide_index"])
            plan.append({
                "role": role,
                "slide_index": chosen["slide_index"],
                "source": source,
                "archetype": chosen.get("archetype", []),
                "macro": chosen.get("macro"),
                "confidence": chosen.get("overall_confidence", 0),
                "reuse_count": 0,
            })
        else:
            # Reuse fallback: 같은 role 또는 fallback archetype의 이미 쓰인 것 재사용
            reuse_pool = [
                l for l in labels
                if role in l.get("narrative_role", [])
            ]
            if not reuse_pool:
                reuse_pool = [
                    l for l in labels
                    if any(a in ARCHETYPE_FALLBACK.get(role, [])
                           for a in l.get("archetype", []))
                ]
            if reuse_pool:
                reuse_pool.sort(
                    key=lambda l: l.get("overall_confidence", 0), reverse=True
                )
                chosen = reuse_pool[0]
                role_use_count[role] = role_use_count.get(role, 0) + 1
                plan.append({
                    "role": role,
                    "slide_index": chosen["slide_index"],
                    "source": "reuse",
                    "archetype": chosen.get("archetype", []),
                    "macro": chosen.get("macro"),
                    "confidence": chosen.get("overall_confidence", 0),
                    "reuse_count": role_use_count[role],
                })
            else:
                plan.append({
                    "role": role,
                    "slide_index": None,
                    "source": "none",
                    "archetype": [],
                    "macro": None,
                    "confidence": 0,
                    "reuse_count": 0,
                })
    return plan


# ----------------------------------------------------------------------------
# Slide text capacity
# ----------------------------------------------------------------------------

def slide_text_summary(src_prs, slide_index: int) -> dict:
    slide = src_prs.slides[slide_index]
    para_records = []
    for fi, sh in edit_ops.iter_leaf_shapes(slide):
        if not getattr(sh, "has_text_frame", False):
            continue
        try:
            tf = sh.text_frame
        except Exception:
            continue
        for pi, p in enumerate(tf.paragraphs):
            text = p.text
            if text:
                para_records.append({
                    "div_id": fi,
                    "para_id": pi,
                    "text": text,
                    "len": len(text),
                })
    para_records.sort(key=lambda r: r["len"], reverse=True)
    return {
        "slide_index": slide_index,
        "n_paragraphs": len(para_records),
        "max_text_len": max((r["len"] for r in para_records), default=0),
        "total_chars": sum(r["len"] for r in para_records),
        "top3": para_records[:3],
        "all": para_records,
    }


# ----------------------------------------------------------------------------
# Build
# ----------------------------------------------------------------------------

def _reorder_sldIdLst(prs, original_indices_in_keep_order: list[int],
                      desired_order: list[int]) -> None:
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    pos_by_original = {orig: i for i, orig in enumerate(original_indices_in_keep_order)}
    new_order = [children[pos_by_original[orig]] for orig in desired_order]
    for c in children:
        sldIdLst.remove(c)
    for c in new_order:
        sldIdLst.append(c)


def build_pptx(plan: list[dict], content_by_role: dict, pptx_out: Path,
               slot_schemas: dict) -> dict:
    pptx_out.parent.mkdir(parents=True, exist_ok=True)

    # 0. 텍스트 capacity 측정 (편집 전)
    src_prs = Presentation(str(TEMPLATE_PATH))
    summaries: dict[int, dict] = {}
    for p in plan:
        if p["slide_index"] is None:
            continue
        if p["slide_index"] not in summaries:
            summaries[p["slide_index"]] = slide_text_summary(src_prs, p["slide_index"])
    del src_prs

    # 1. keep_slides + 재정렬
    desired_order_with_dup = [
        p["slide_index"] for p in plan if p["slide_index"] is not None
    ]
    keep_unique_sorted = sorted(set(desired_order_with_dup))

    # plan에 중복 slide_index 있으면 keep_slides는 1번만 유지하고
    # editor.prs에서 그 슬라이드를 N번 복제해야 한다.
    # 단순화: 본 시뮬레이션은 plan 단계에서 재사용도 가능하므로
    # 중복 slide_index가 있으면 마지막 기준으로 1번만 사용 (덱이 짧아질 수 있음).
    # 정확한 reuse 시뮬레이션은 SlideCloner를 직접 써야 함.
    # → 실용성을 위해 unique slide_index만 keep + plan을 unique sequence로 단순화
    seen = set()
    plan_unique = []
    plan_skipped = []
    for p in plan:
        if p["slide_index"] is None:
            plan_skipped.append({**p, "reason": "no_candidate"})
            continue
        if p["slide_index"] in seen:
            plan_skipped.append({**p, "reason": "duplicate_slide_skipped"})
            continue
        seen.add(p["slide_index"])
        plan_unique.append(p)

    desired_order = [p["slide_index"] for p in plan_unique]
    keep_unique_sorted = sorted(set(desired_order))

    editor = TemplateEditor(TEMPLATE_PATH)
    editor.keep_slides(keep_unique_sorted)
    _reorder_sldIdLst(editor.prs, keep_unique_sorted, desired_order)

    # 2. edit_ops 적용
    edit_results = []
    role_content_idx: dict[str, int] = {}

    for step_idx, item in enumerate(plan_unique):
        role = item["role"]
        sidx = item["slide_index"]
        step = step_idx + 1

        # 가상 컨텐츠
        contents = content_by_role.get(role, [role.upper()])
        idx = role_content_idx.get(role, 0)
        new_text = contents[min(idx, len(contents) - 1)]
        role_content_idx[role] = idx + 1

        summary = summaries[sidx]
        target_slide = editor.prs.slides[step_idx]

        edit_ok = False
        edit_reason = ""
        target_div_id = None
        target_para_id = None
        max_chars_limit = None
        overflow = False

        if summary["top3"]:
            cand = summary["top3"][0]
            target_div_id = cand["div_id"]
            target_para_id = cand["para_id"]

            # max_chars 룩업 (slot_schemas)
            slide_schemas = slot_schemas.get(sidx, {})
            slot = slide_schemas.get(cand["div_id"])
            max_chars_limit = slot.get("max_chars") if slot else None
            if max_chars_limit is not None and len(new_text) > max_chars_limit:
                overflow = True

            try:
                edit_ops.replace_paragraph(
                    target_slide, cand["div_id"], cand["para_id"], new_text
                )
                edit_ok = True
                edit_reason = (
                    f"div={cand['div_id']} para={cand['para_id']} "
                    f"orig_len={cand['len']} new_len={len(new_text)} "
                    f"max={max_chars_limit}"
                )
            except Exception as e:
                edit_reason = f"FAIL: {type(e).__name__}: {e}"
        else:
            edit_reason = "no_text_paragraph"

        edit_results.append({
            "step": step,
            "role": role,
            "slide_index": sidx,
            "source": item["source"],
            "new_text": new_text,
            "edit_ok": edit_ok,
            "edit_reason": edit_reason,
            "target_div_id": target_div_id,
            "target_para_id": target_para_id,
            "max_chars_limit": max_chars_limit,
            "overflow": overflow,
            "n_paragraphs_total": summary["n_paragraphs"],
            "n_paragraphs_filled": 1 if edit_ok else 0,
            "fill_ratio": (1 / summary["n_paragraphs"]) if (edit_ok and summary["n_paragraphs"]) else 0,
        })

    editor.save(pptx_out)
    editor.cleanup()

    return {
        "plan_unique": plan_unique,
        "plan_skipped": plan_skipped,
        "edits": edit_results,
        "pptx": str(pptx_out),
    }


# ----------------------------------------------------------------------------
# PNG render
# ----------------------------------------------------------------------------

def render_pngs(pptx_path: Path, png_dir: Path) -> list[Path]:
    import pythoncom
    import win32com.client

    png_dir.mkdir(parents=True, exist_ok=True)
    for old in png_dir.glob("*.png"):
        old.unlink()

    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            str(pptx_path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False,
        )
        total = presentation.Slides.Count
        out = []
        for i in range(1, total + 1):
            p = png_dir / f"step_{i:02d}.png"
            try:
                presentation.Slides(i).Export(str(p), "PNG", 1568, 1176)
                out.append(p)
            except Exception as e:
                print(f"  [err] slide {i}: {e}", flush=True)
        return out
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


# ----------------------------------------------------------------------------
# Metrics
# ----------------------------------------------------------------------------

def compute_scenario_metrics(scenario_id: str, plan: list[dict],
                             plan_unique: list[dict], plan_skipped: list[dict],
                             edits: list[dict], narrative: list[str]) -> dict:
    n_total = len(narrative)
    n_role_hit = sum(1 for p in plan if p["source"] == "role")
    n_archetype = sum(1 for p in plan if p["source"] == "archetype")
    n_reuse = sum(1 for p in plan if p["source"] == "reuse")
    n_none = sum(1 for p in plan if p["source"] == "none")

    n_built = len(edits)
    n_edit_ok = sum(1 for e in edits if e["edit_ok"])
    n_overflow = sum(1 for e in edits if e["overflow"])
    fill_ratios = [e["fill_ratio"] for e in edits if e["edit_ok"]]
    avg_fill_ratio = sum(fill_ratios) / len(fill_ratios) if fill_ratios else 0

    # role 풀 고갈 추적
    role_pool_exhaustion: dict[str, int] = {}
    role_request_count: Counter = Counter(narrative)
    for r, n_requested in role_request_count.items():
        n_satisfied = sum(
            1 for p in plan
            if p["role"] == r and p["source"] in ("role", "archetype")
        )
        if n_satisfied < n_requested:
            role_pool_exhaustion[r] = {
                "requested": n_requested,
                "satisfied_direct": n_satisfied,
                "shortage": n_requested - n_satisfied,
            }

    # 슬라이드당 paragraph 평균 (Mode A 미채움 추적)
    avg_n_paragraphs = (
        sum(e["n_paragraphs_total"] for e in edits) / len(edits) if edits else 0
    )

    # 점수 (가중)
    score_role = (n_role_hit / n_total) * 100
    score_fill = avg_fill_ratio * 100
    score_overflow = (1 - n_overflow / max(n_built, 1)) * 100
    # Hybrid composite (D는 별도 단계)
    composite_quant = round(
        score_role * 0.30 + score_fill * 0.40 + score_overflow * 0.30, 1
    )

    return {
        "scenario_id": scenario_id,
        "narrative_length": n_total,
        "metrics": {
            "A_role_match": {
                "role_hit": n_role_hit,
                "archetype_fallback": n_archetype,
                "reuse_fallback": n_reuse,
                "miss": n_none,
                "role_hit_pct": round(score_role, 1),
            },
            "B_slot_fill": {
                "avg_fill_ratio": round(avg_fill_ratio, 4),
                "avg_paragraphs_per_slide": round(avg_n_paragraphs, 1),
                "fill_pct_proxy": round(score_fill, 1),
                "comment": "Mode A는 슬라이드당 1개 paragraph만 채움. fill_ratio = 1/n_paragraphs.",
            },
            "C_overflow": {
                "n_overflow": n_overflow,
                "n_built": n_built,
                "overflow_rate_pct": round(n_overflow / max(n_built, 1) * 100, 1),
                "no_overflow_score": round(score_overflow, 1),
            },
            "F_role_pool_exhaustion": role_pool_exhaustion,
            "build_stats": {
                "planned": n_total,
                "built": n_built,
                "edit_ok": n_edit_ok,
                "skipped": len(plan_skipped),
            },
            "composite_quant_score": composite_quant,
        },
    }


# ----------------------------------------------------------------------------
# Main
# ----------------------------------------------------------------------------

def run_scenario(scenario_id: str, labels: list[dict], skeletons: dict,
                 slot_schemas: dict) -> dict:
    print()
    print("=" * 80)
    print(f"SCENARIO: {scenario_id}")
    print("=" * 80)

    sk = skeletons[scenario_id]
    narrative = sk["narrative_sequence"]
    sc = SCENARIO_CONTENT[scenario_id]

    print(f"  name: {sc['scenario_name']}")
    print(f"  narrative length: {len(narrative)}")
    print(f"  narrative: {narrative}")

    plan = select_deck(labels, narrative)
    n_role = sum(1 for p in plan if p["source"] == "role")
    n_arch = sum(1 for p in plan if p["source"] == "archetype")
    n_reuse = sum(1 for p in plan if p["source"] == "reuse")
    n_none = sum(1 for p in plan if p["source"] == "none")
    print(f"  retrieval: role={n_role} archetype={n_arch} reuse={n_reuse} none={n_none}")

    out_dir = OUTPUT_ROOT / scenario_id
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_out = out_dir / "deck.pptx"
    png_dir = out_dir / "pngs"

    t0 = time.time()
    build = build_pptx(plan, sc["content_by_role"], pptx_out, slot_schemas)
    dt_build = time.time() - t0
    print(f"  build: {dt_build:.1f}s")

    t0 = time.time()
    pngs = render_pngs(pptx_out, png_dir)
    dt_render = time.time() - t0
    print(f"  render: {dt_render:.1f}s, {len(pngs)} pngs")

    metrics = compute_scenario_metrics(
        scenario_id, plan, build["plan_unique"], build["plan_skipped"],
        build["edits"], narrative,
    )
    metrics["paths"] = {
        "pptx": str(pptx_out),
        "png_dir": str(png_dir),
        "pngs": [str(p) for p in pngs],
    }
    metrics["plan"] = plan
    metrics["edits"] = build["edits"]
    metrics["plan_skipped"] = build["plan_skipped"]
    metrics["timing"] = {"build_sec": round(dt_build, 1), "render_sec": round(dt_render, 1)}

    # 시나리오별 보고서
    report_path = out_dir / "report.json"
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(metrics, f, ensure_ascii=False, indent=2)
    print(f"  report: {report_path}")

    return metrics


def main():
    target = sys.argv[1] if len(sys.argv) > 1 else None

    print("=" * 80)
    print("Phase A3 Step 1 — 5 Benchmark Scenarios (Mode A only)")
    print("=" * 80)

    labels = load_catalog()
    skeletons = load_skeletons()
    slot_schemas = load_slot_schemas()
    print(f"loaded: {len(labels)} labels / {len(skeletons)} skeletons "
          f"/ {len(slot_schemas)} slides w/ slot schemas")

    # 전체 role 분포
    role_counter = Counter()
    for lab in labels:
        for r in lab.get("narrative_role", []):
            role_counter[r] += 1
    print("[overall role distribution]")
    for r, c in role_counter.most_common():
        print(f"  {r:>16}: {c:5d}")

    if target:
        if target not in SCENARIO_CONTENT:
            print(f"ERROR: unknown scenario {target}")
            sys.exit(1)
        scenarios = [target]
    else:
        scenarios = list(SCENARIO_CONTENT.keys())

    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)

    all_results = []
    for sid in scenarios:
        res = run_scenario(sid, labels, skeletons, slot_schemas)
        all_results.append(res)

    # 통합 scoreboard
    scoreboard = {
        "phase": "A3-Step1",
        "mode": "A_only",
        "n_scenarios": len(all_results),
        "results": all_results,
        "summary": {
            "avg_role_hit_pct": round(
                sum(r["metrics"]["A_role_match"]["role_hit_pct"] for r in all_results)
                / len(all_results), 1,
            ),
            "avg_fill_pct": round(
                sum(r["metrics"]["B_slot_fill"]["fill_pct_proxy"] for r in all_results)
                / len(all_results), 1,
            ),
            "avg_overflow_rate_pct": round(
                sum(r["metrics"]["C_overflow"]["overflow_rate_pct"] for r in all_results)
                / len(all_results), 1,
            ),
            "avg_composite_quant": round(
                sum(r["metrics"]["composite_quant_score"] for r in all_results)
                / len(all_results), 1,
            ),
        },
    }

    scoreboard_path = OUTPUT_ROOT / "scoreboard.json"
    with open(scoreboard_path, "w", encoding="utf-8") as f:
        json.dump(scoreboard, f, ensure_ascii=False, indent=2)

    print()
    print("=" * 80)
    print("SCOREBOARD")
    print("=" * 80)
    for r in all_results:
        m = r["metrics"]
        print(f"  {r['scenario_id']:>30} | role={m['A_role_match']['role_hit_pct']:5.1f}% "
              f"fill={m['B_slot_fill']['fill_pct_proxy']:5.1f}% "
              f"overflow={m['C_overflow']['overflow_rate_pct']:5.1f}% "
              f"comp={m['composite_quant_score']:5.1f}")
    print()
    print(f"AVG composite quant = {scoreboard['summary']['avg_composite_quant']}")
    print(f"  (Note: D PPTEval / E visual review은 별도 단계)")
    print(f"saved: {scoreboard_path}")


if __name__ == "__main__":
    main()
