"""Phase A3 Step 3 — paragraph-aware fill 5 벤치마크 재측정.

v1 (Mode A 단독): 슬라이드당 1개 paragraph만 채움 → fill 6.8%
v2 (paragraph fill): role별 슬롯 매핑 → 모든 fillable 슬롯 채움

Usage
-----
python scripts/benchmark_5_scenarios_v2.py [scenario_id]
"""
from __future__ import annotations

import json
import sys
import time
from collections import Counter
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation

from ppt_builder.template import edit_ops
from ppt_builder.template.editor import TemplateEditor
from ppt_builder.catalog.paragraph_query import ParagraphStore


def _shrink_font_if_overflow(para, new_text: str, max_chars: int | None,
                             min_pt: float = 8.0) -> bool:
    """텍스트가 max_chars 초과 시 폰트 자동 축소.

    축소 비율 = sqrt(max_chars / len(text)) (면적 기반 비례).
    min_pt 이하로는 안 내림.
    Returns: True if shrunk
    """
    if not max_chars or len(new_text) <= max_chars:
        return False
    import math
    scale = math.sqrt(max_chars / max(len(new_text), 1))
    runs = list(para.runs)
    if not runs:
        return False
    shrunk_any = False
    for r in runs:
        try:
            cur = r.font.size
            if cur is None:
                continue
            new_pt = max(min_pt, cur.pt * scale)
            if new_pt < cur.pt - 0.1:  # actually shrinking
                from pptx.util import Pt
                r.font.size = Pt(new_pt)
                shrunk_any = True
        except Exception:
            pass
    return shrunk_any


def _replace_paragraph_universal(slide, slot, new_text: str,
                                 shrink_overflow: bool = True) -> None:
    """edit_ops.replace_paragraph + table 셀 분기 + overflow 폰트 축소."""
    # 일반 shape: edit_ops 사용
    if slot.shape_kind != "TABLE":
        edit_ops.replace_paragraph(slide, slot.flat_idx, slot.paragraph_id, new_text)
        if shrink_overflow and slot.max_chars:
            # 편집된 paragraph 다시 찾아서 폰트 축소
            for fi, sh in edit_ops.iter_leaf_shapes(slide):
                if fi != slot.flat_idx:
                    continue
                if not getattr(sh, "has_text_frame", False):
                    return
                paras = sh.text_frame.paragraphs
                if slot.paragraph_id < len(paras):
                    _shrink_font_if_overflow(
                        paras[slot.paragraph_id], new_text, slot.max_chars
                    )
                return
        return
    # TABLE 셀
    target_shape = None
    for fi, sh in edit_ops.iter_leaf_shapes(slide):
        if fi == slot.flat_idx:
            target_shape = sh
            break
    if target_shape is None or not target_shape.has_table:
        raise edit_ops.SlideEditError(f"flat_idx={slot.flat_idx} not a table")
    if slot.table_row is None or slot.table_col is None:
        raise edit_ops.SlideEditError("table slot missing row/col")
    cell = target_shape.table.cell(slot.table_row, slot.table_col)
    paras = cell.text_frame.paragraphs
    if slot.paragraph_id >= len(paras):
        raise edit_ops.SlideEditError(
            f"paragraph_id={slot.paragraph_id} out of range in cell"
        )
    para = paras[slot.paragraph_id]
    runs = list(para.runs)
    if runs:
        first, *rest = runs
        for r in rest:
            r._r.getparent().remove(r._r)
        first.text = new_text
    else:
        run = para.add_run()
        run.text = new_text
    if shrink_overflow and slot.max_chars:
        _shrink_font_if_overflow(para, new_text, slot.max_chars)

# v1과 같은 자산 재사용
sys.path.insert(0, str(ROOT / "scripts"))
from benchmark_5_scenarios import (  # noqa: E402
    SCENARIO_CONTENT,
    ARCHETYPE_FALLBACK,
    load_skeletons,
    candidates_for_role,
    select_deck,
    _reorder_sldIdLst,
)


def load_catalog() -> list[dict]:
    """v5: vision-relabeled 카탈로그 로드."""
    with open(CATALOG_PATH, encoding="utf-8") as f:
        return json.load(f)["labels"]


# ----------------------------------------------------------------------------
# Track 1.2: Capacity-aware retrieval
# ----------------------------------------------------------------------------

def _capacity_fitness(slide_index: int, target_n: int,
                      store: ParagraphStore) -> float:
    """슬라이드 parallel 슬롯 capacity + max_chars vs target_n.

    핵심 안전장치: primary 슬롯의 평균 max_chars가 너무 작으면 (narrow 슬롯 함정)
    어떤 target_n이라도 페널티.
    """
    cap = store.slot_capacity(slide_index)
    primary_caps = [
        cap.get("chevron_label", 0),
        cap.get("card_header", 0),
        cap.get("callout_text", 0),
        cap.get("card_body", 0),
    ]
    primary = max(primary_caps) if primary_caps else 0

    # primary 슬롯의 평균 max_chars 측정 (narrow 슬롯 회피)
    fillable = store.fillable_slots(slide_index)
    primary_role = None
    primary_max = 0
    for r in ("chevron_label", "card_header", "callout_text", "card_body"):
        if cap.get(r, 0) == primary and primary > 0:
            primary_role = r
            primary_max = primary
            break

    if primary_role:
        slots = fillable.get(primary_role, [])
        max_chars_vals = [s.max_chars for s in slots if s.max_chars]
        avg_max_chars = (
            sum(max_chars_vals) / len(max_chars_vals) if max_chars_vals else 0
        )
        # narrow 슬롯 패널티 (max_chars < 20 = 약 한국어 5자)
        if avg_max_chars and avg_max_chars < 20:
            return 0.3
    else:
        avg_max_chars = 0

    if target_n <= 3:
        # 단일/소량: 풍부한 framework + 충분한 슬롯 폭
        if primary == 0:
            return 0.7  # parallel 없는 슬라이드 (cover/표) — fine for n<=3
        if avg_max_chars >= 30:
            return 1.0  # 양호한 폭의 framework
        return 0.6

    # n >= 4: capacity 매칭
    if primary == 0:
        return 0.3
    if primary >= target_n:
        excess = primary - target_n
        base = max(0.6, 1.0 - 0.02 * excess)
    else:
        shortage = target_n - primary
        base = max(0.2, 0.6 - 0.15 * shortage)
    # narrow 패널티 추가 적용
    if avg_max_chars and avg_max_chars < 30:
        base *= 0.7
    return base


def select_deck_diverse(labels: list[dict], narrative: list[str],
                        scenario_content: dict,
                        store: ParagraphStore) -> list[dict]:
    """Slide diversity-aware deck 선정.

    핵심:
    - role 풀 직접 매칭 시 top-K 중 narrow 슬롯 함정 회피 + 다양성 균형
    - 같은 group_signature 패밀리는 회피 (시각 단조로움 방지)
    - capacity가 컨텐츠 수에 가까운 것 약한 선호
    """
    used: set[int] = set()
    used_signatures: set[str] = set()
    plan = []
    role_use_count: dict[str, int] = {}
    by_role = scenario_content["content_by_role"]

    for role in narrative:
        target_n = len(by_role.get(role, [])) or 1

        cands, source = candidates_for_role(labels, role, used)

        if cands:
            # top-30 후보 평가
            top_cands = cands[:30]
            scored = []
            for c in top_cands:
                sidx = c["slide_index"]
                # narrow 슬롯 페널티
                cap = store.slot_capacity(sidx)
                fillable = store.fillable_slots(sidx)
                primary_role = None
                primary_max_chars = 0
                for r in ("chevron_label", "card_header", "callout_text"):
                    if cap.get(r, 0) > 0:
                        primary_role = r
                        slots = fillable.get(r, [])
                        max_chars_vals = [s.max_chars for s in slots if s.max_chars]
                        primary_max_chars = (
                            sum(max_chars_vals) / len(max_chars_vals)
                            if max_chars_vals else 0
                        )
                        break

                # narrow 페널티 (max_chars < 20 = 약 5자 한국어)
                narrow_penalty = 0.0
                if primary_max_chars and primary_max_chars < 20:
                    narrow_penalty = 0.4

                # capacity fitness (target_n vs primary capacity)
                primary_cap = max(
                    cap.get("chevron_label", 0),
                    cap.get("card_header", 0),
                    cap.get("callout_text", 0),
                    cap.get("card_body", 0),
                )
                fit = 1.0
                if target_n >= 2:
                    if primary_cap == 0:
                        fit = 0.6
                    elif primary_cap >= target_n:
                        fit = max(0.7, 1.0 - 0.02 * (primary_cap - target_n))
                    else:
                        fit = max(0.4, 0.7 - 0.1 * (target_n - primary_cap))

                # signature diversity 보너스
                sig_key = f"{c.get('macro')}:{tuple(sorted(c.get('archetype', [])))}"
                diversity_bonus = 0.0 if sig_key in used_signatures else 0.1

                score = (
                    0.4 * c.get("overall_confidence", 0)
                    + 0.4 * fit
                    + 0.2 * (1 - narrow_penalty)
                    + diversity_bonus
                )
                scored.append((score, c, sig_key, fit, narrow_penalty))

            scored.sort(key=lambda x: x[0], reverse=True)
            chosen = scored[0][1]
            chosen_sig = scored[0][2]
            used.add(chosen["slide_index"])
            used_signatures.add(chosen_sig)
            plan.append({
                "role": role,
                "slide_index": chosen["slide_index"],
                "source": source,
                "archetype": chosen.get("archetype", []),
                "macro": chosen.get("macro"),
                "confidence": chosen.get("overall_confidence", 0),
                "score": round(scored[0][0], 2),
                "capacity_fit": round(scored[0][3], 2),
                "target_n_items": target_n,
                "reuse_count": 0,
            })
        else:
            # reuse fallback
            reuse_pool = [
                l for l in labels if role in l.get("narrative_role", [])
            ]
            if not reuse_pool:
                reuse_pool = [
                    l for l in labels
                    if any(a in ARCHETYPE_FALLBACK.get(role, [])
                           for a in l.get("archetype", []))
                ]
            if reuse_pool:
                # 같은 reuse 풀에서 used_signatures 회피
                reuse_pool.sort(
                    key=lambda l: l.get("overall_confidence", 0), reverse=True
                )
                chosen = None
                for c in reuse_pool[:10]:
                    sig_key = f"{c.get('macro')}:{tuple(sorted(c.get('archetype', [])))}"
                    if sig_key not in used_signatures:
                        chosen = c
                        break
                if chosen is None:
                    chosen = reuse_pool[0]
                role_use_count[role] = role_use_count.get(role, 0) + 1
                plan.append({
                    "role": role,
                    "slide_index": chosen["slide_index"],
                    "source": "reuse",
                    "archetype": chosen.get("archetype", []),
                    "macro": chosen.get("macro"),
                    "confidence": chosen.get("overall_confidence", 0),
                    "score": chosen.get("overall_confidence", 0),
                    "capacity_fit": 0.5,
                    "target_n_items": target_n,
                    "reuse_count": role_use_count[role],
                })
            else:
                plan.append({
                    "role": role, "slide_index": None, "source": "none",
                    "archetype": [], "macro": None, "confidence": 0,
                    "score": 0, "capacity_fit": 0,
                    "target_n_items": target_n, "reuse_count": 0,
                })
    return plan


def select_deck_capacity_aware(labels: list[dict], narrative: list[str],
                               scenario_content: dict,
                               store: ParagraphStore) -> list[dict]:
    """capacity-aware 슬라이드 선정. Step 4 Track 1.2."""
    used: set[int] = set()
    plan = []
    role_use_count: dict[str, int] = {}

    by_role = scenario_content["content_by_role"]

    for role in narrative:
        # 이 role의 content 항목 수 파악
        role_items = by_role.get(role, [])
        target_n = len(role_items) if role_items else 1

        cands, source = candidates_for_role(labels, role, used)

        if cands:
            # capacity fitness로 재정렬 (top 30개 중에서)
            top_cands = cands[:30]
            scored = []
            for c in top_cands:
                fitness = _capacity_fitness(c["slide_index"], target_n, store)
                conf = c.get("overall_confidence", 0)
                # 합성 점수: capacity 50% + role confidence 50%
                score = 0.5 * fitness + 0.5 * conf
                scored.append((score, c, fitness))
            scored.sort(key=lambda x: x[0], reverse=True)
            chosen = scored[0][1]
            chosen_fitness = scored[0][2]
            used.add(chosen["slide_index"])
            plan.append({
                "role": role,
                "slide_index": chosen["slide_index"],
                "source": source,
                "archetype": chosen.get("archetype", []),
                "macro": chosen.get("macro"),
                "confidence": chosen.get("overall_confidence", 0),
                "capacity_fitness": round(chosen_fitness, 2),
                "target_n_items": target_n,
                "reuse_count": 0,
            })
        else:
            # fallback to reuse
            reuse_pool = [
                l for l in labels if role in l.get("narrative_role", [])
            ]
            if not reuse_pool:
                reuse_pool = [
                    l for l in labels
                    if any(a in ARCHETYPE_FALLBACK.get(role, [])
                           for a in l.get("archetype", []))
                ]
            if reuse_pool:
                # capacity-aware 선택
                top = reuse_pool[:30]
                scored = []
                for c in top:
                    fitness = _capacity_fitness(c["slide_index"], target_n, store)
                    score = 0.5 * fitness + 0.5 * c.get("overall_confidence", 0)
                    scored.append((score, c, fitness))
                scored.sort(key=lambda x: x[0], reverse=True)
                chosen = scored[0][1]
                role_use_count[role] = role_use_count.get(role, 0) + 1
                plan.append({
                    "role": role,
                    "slide_index": chosen["slide_index"],
                    "source": "reuse",
                    "archetype": chosen.get("archetype", []),
                    "macro": chosen.get("macro"),
                    "confidence": chosen.get("overall_confidence", 0),
                    "capacity_fitness": round(scored[0][2], 2),
                    "target_n_items": target_n,
                    "reuse_count": role_use_count[role],
                })
            else:
                plan.append({
                    "role": role, "slide_index": None, "source": "none",
                    "archetype": [], "macro": None, "confidence": 0,
                    "capacity_fitness": 0, "target_n_items": target_n,
                    "reuse_count": 0,
                })
    return plan


CATALOG_PATH = ROOT / "output" / "catalog" / "final_labels_v2.json"  # vision-relabeled
TEMPLATE_PATH = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
OUTPUT_ROOT = ROOT / "output" / "benchmark_v2"


# ----------------------------------------------------------------------------
# Content auto-split + expand
# ----------------------------------------------------------------------------

import re


def _split_compound_text(text: str, max_segments: int = 8) -> list[str]:
    """단일 long text를 segment로 분리.

    분리 우선순위: ' / ' > ' + ' > ', ' (공백 동반).
    각 segment는 최소 4자 이상 유지.
    """
    if not text:
        return []
    # ' / ' 우선
    if " / " in text and text.count(" / ") >= 1:
        parts = [p.strip() for p in text.split(" / ")]
    elif " + " in text and text.count(" + ") >= 1:
        parts = [p.strip() for p in text.split(" + ")]
    elif ", " in text and text.count(", ") >= 2:
        parts = [p.strip() for p in text.split(", ")]
    else:
        return [text]

    parts = [p for p in parts if len(p) >= 3]
    if len(parts) > max_segments:
        # 너무 많으면 앞 N개만 + 나머지 합침
        parts = parts[: max_segments - 1] + [" ".join(parts[max_segments - 1 :])]
    return parts if len(parts) > 1 else [text]


def _expand_to_capacity(items: list[str], target_n: int) -> list[str]:
    """items를 target_n 개로 expand. 부족 시 split 시도, 여전히 부족 시 그대로.

    절대 augmented 가짜 텍스트 추가 안 함 (사용자 신뢰 우선).
    """
    if not items or target_n <= len(items):
        return items[:target_n] if target_n > 0 else items

    # 단일 long text 분리 시도
    if len(items) == 1:
        split = _split_compound_text(items[0], max_segments=target_n)
        if len(split) > 1:
            return split[:target_n]

    # 여러 항목인데 일부에 compound가 있으면 추가 분리
    expanded = []
    for it in items:
        sub = _split_compound_text(it)
        expanded.extend(sub)
    return expanded[:target_n] if expanded else items


# ----------------------------------------------------------------------------
# Content-by-role expansion: scenario 컨텐츠를 role별 list로 확장
# ----------------------------------------------------------------------------

def expand_content_for_slide(role: str, scenario_content: dict,
                             role_use_counter: dict[str, int],
                             slot_capacity: dict[str, int]) -> dict[str, list[str]]:
    """선택된 슬라이드 role + capacity에 맞춰 fill용 content_by_role 생성.

    Title은 항상 채우고, parallel items는 단일 best-fit 슬롯 타입에만 할당
    (slot capacity와 가장 잘 맞는 role 선택). 중복 할당 방지.
    """
    by_role = scenario_content["content_by_role"]
    use_idx = role_use_counter.get(role, 0)
    role_use_counter[role] = use_idx + 1

    role_items = by_role.get(role, [])
    if not role_items:
        role_items = [role.upper()]

    # title: 시나리오명 또는 role별 첫 line
    title = scenario_content["scenario_name"]
    if role in ("opening", "closing", "divider"):
        title = role_items[min(use_idx, len(role_items) - 1)]

    out: dict[str, list[str] | str] = {"title": title}

    # parallel content: 시각 우선순위 슬롯 1개에만 할당 + capacity에 맞춰 expand
    if role not in ("opening", "closing", "divider", "agenda"):
        candidates = ["chevron_label", "card_header", "callout_text", "card_body"]
        chosen = None
        for c in candidates:
            if slot_capacity.get(c, 0) > 0:
                chosen = c
                break
        if chosen:
            # 슬롯 수에 맞춰 컨텐츠 expand (부족하면 split, 그래도 부족하면 그대로)
            target_n = slot_capacity[chosen]
            expanded = _expand_to_capacity(role_items, target_n)
            out[chosen] = expanded

    return out


# ----------------------------------------------------------------------------
# Build with paragraph-level fill
# ----------------------------------------------------------------------------

def build_pptx_v2(plan: list[dict], scenario_content: dict, store: ParagraphStore,
                  pptx_out: Path) -> dict:
    pptx_out.parent.mkdir(parents=True, exist_ok=True)

    seen = set()
    plan_unique = []
    for p in plan:
        if p["slide_index"] is None:
            continue
        if p["slide_index"] in seen:
            continue
        seen.add(p["slide_index"])
        plan_unique.append(p)

    desired_order = [p["slide_index"] for p in plan_unique]
    keep_unique_sorted = sorted(set(desired_order))

    editor = TemplateEditor(TEMPLATE_PATH)
    editor.keep_slides(keep_unique_sorted)
    _reorder_sldIdLst(editor.prs, keep_unique_sorted, desired_order)

    edit_results = []
    role_use_counter: dict[str, int] = {}

    for step_idx, item in enumerate(plan_unique):
        role = item["role"]
        sidx = item["slide_index"]
        step = step_idx + 1

        # 슬라이드의 슬롯 capacity (먼저)
        capacity = store.slot_capacity(sidx)
        # role별 컨텐츠 확장 (capacity-aware)
        content = expand_content_for_slide(role, scenario_content,
                                           role_use_counter, capacity)
        # None 값 제거
        content = {k: v for k, v in content.items() if v}

        # 매핑 — truncate + 폰트 축소 병행 (둘 다 안전장치)
        edits = store.match_content(sidx, content, truncate_overflow=True)

        # 실제 슬라이드 편집
        target_slide = editor.prs.slides[step_idx]
        n_ok = 0
        n_fail = 0
        per_edit_log = []
        edited_keys: set[tuple[int, int]] = set()
        # role -> 슬롯 인덱스 lookup (universal replace용)
        all_fillable = store.fillable_slots(sidx)
        flat_to_slot: dict[tuple[int, int], object] = {}
        for slist in all_fillable.values():
            for s_obj in slist:
                flat_to_slot[(s_obj.flat_idx, s_obj.paragraph_id)] = s_obj

        for e in edits:
            try:
                slot_obj = flat_to_slot.get((e["flat_idx"], e["paragraph_id"]))
                if slot_obj is None:
                    raise edit_ops.SlideEditError("slot lookup failed")
                _replace_paragraph_universal(target_slide, slot_obj, e["text"])
                n_ok += 1
                edited_keys.add((e["flat_idx"], e["paragraph_id"]))
                per_edit_log.append({
                    **e, "edit_ok": True, "edit_reason": "OK",
                })
            except Exception as ex:
                n_fail += 1
                per_edit_log.append({
                    **e, "edit_ok": False,
                    "edit_reason": f"{type(ex).__name__}: {ex}",
                })

        # 미매칭 fillable 슬롯의 ~~ placeholder 비우기 (시각적 깔끔함)
        n_blanked = 0
        fillable = store.fillable_slots(sidx)
        for role_name, slots in fillable.items():
            for slot in slots:
                key = (slot.flat_idx, slot.paragraph_id)
                if key in edited_keys:
                    continue
                # 원래 텍스트가 ~~ 류 placeholder일 때만 비움
                orig = (slot.text_original or "").strip()
                if orig in {"~~", "~~\x0b", "\x0b~~", "~"} or "~~" in orig:
                    try:
                        _replace_paragraph_universal(target_slide, slot, " ")
                        n_blanked += 1
                        edited_keys.add(key)
                    except Exception:
                        pass

        # Track 1.1: 잔존 ~~ 추가 청소 — fillable 외 슬롯도 ~~만 있는 placeholder 비움
        # (단, page_number는 보존 — 숫자가 의미)
        n_residual_cleaned = 0
        all_slots = store.slots(sidx)
        for slot in all_slots:
            key = (slot.flat_idx, slot.paragraph_id)
            if key in edited_keys:
                continue
            if slot.role in ("page_number", "footer", "date_text"):
                continue
            orig_raw = slot.text_original or ""
            orig = orig_raw.strip()
            # ~~만 있는 (또는 거의 ~~만 있는) placeholder
            is_pure_placeholder = (
                orig == "~~"
                or orig.replace("~", "").replace("\x0b", "").strip() == ""
                or (orig.startswith("~~") and len(orig) <= 4)
                or (orig.endswith("~~") and len(orig) <= 4)
            )
            # number-prefixed placeholder ("1. ~~", "2. ~~")
            is_numbered_placeholder = (
                len(orig) <= 8 and "~~" in orig
                and any(c.isdigit() for c in orig)
            )
            if is_pure_placeholder or is_numbered_placeholder:
                try:
                    _replace_paragraph_universal(target_slide, slot, " ")
                    n_residual_cleaned += 1
                    edited_keys.add(key)
                except Exception:
                    pass

        # 통계: visual_clean_ratio = (filled + blanked + ~~ 아닌 슬롯) / fillable
        # ~~ visible 만 visual debt로 카운트
        n_initially_clean = 0
        for role_name, slots in fillable.items():
            for slot in slots:
                key = (slot.flat_idx, slot.paragraph_id)
                if key in edited_keys:
                    continue  # 이미 채웠거나 비웠음
                orig = (slot.text_original or "").strip()
                # 빈 paragraph 또는 ~~ 없는 텍스트 → 시각적 clean
                if not orig:
                    n_initially_clean += 1
                elif "~~" not in orig:
                    n_initially_clean += 1

        fillable_total = sum(capacity.values())
        visual_clean = n_ok + n_blanked + n_initially_clean
        edit_results.append({
            "step": step,
            "role": role,
            "slide_index": sidx,
            "source": item["source"],
            "slot_capacity": capacity,
            "n_fillable": fillable_total,
            "n_edit_attempted": len(edits),
            "n_edit_ok": n_ok,
            "n_edit_fail": n_fail,
            "n_blanked": n_blanked,
            "n_initially_clean": n_initially_clean,
            "n_residual_cleaned": n_residual_cleaned,
            "fill_ratio": (n_ok / fillable_total) if fillable_total else 0,
            "visual_resolution_ratio": (
                visual_clean / fillable_total if fillable_total else 0
            ),
            "n_overflow_truncated": sum(
                1 for e in per_edit_log if e.get("original_overflow")
            ),
            "edits": per_edit_log,
        })

    editor.save(pptx_out)
    editor.cleanup()

    return {"plan_unique": plan_unique, "edits": edit_results, "pptx": str(pptx_out)}


# ----------------------------------------------------------------------------
# PNG render (v1과 동일)
# ----------------------------------------------------------------------------

def render_pngs(pptx_path: Path, png_dir: Path) -> list[Path]:
    import pythoncom
    import win32com.client

    png_dir.mkdir(parents=True, exist_ok=True)
    for old in png_dir.glob("*.png"):
        old.unlink()

    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            str(pptx_path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False,
        )
        total = presentation.Slides.Count
        out = []
        for i in range(1, total + 1):
            p = png_dir / f"step_{i:02d}.png"
            try:
                presentation.Slides(i).Export(str(p), "PNG", 1568, 1176)
                out.append(p)
            except Exception as e:
                print(f"  [err] slide {i}: {e}", flush=True)
        return out
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


# ----------------------------------------------------------------------------
# Metrics
# ----------------------------------------------------------------------------

def compute_v2_metrics(scenario_id: str, plan: list[dict],
                       edits: list[dict], narrative: list[str]) -> dict:
    n_total = len(narrative)
    n_role_hit = sum(1 for p in plan if p["source"] == "role")
    n_archetype = sum(1 for p in plan if p["source"] == "archetype")
    n_reuse = sum(1 for p in plan if p["source"] == "reuse")
    n_none = sum(1 for p in plan if p["source"] == "none")

    fillable_total = sum(e["n_fillable"] for e in edits)
    edited_total = sum(e["n_edit_ok"] for e in edits)
    blanked_total = sum(e["n_blanked"] for e in edits)
    initial_clean_total = sum(e.get("n_initially_clean", 0) for e in edits)
    fail_total = sum(e["n_edit_fail"] for e in edits)
    overflow_total = sum(e["n_overflow_truncated"] for e in edits)

    fill_ratio = edited_total / fillable_total if fillable_total else 0
    visual_resolution = (
        (edited_total + blanked_total + initial_clean_total) / fillable_total
        if fillable_total else 0
    )
    avg_per_slide_fillable = fillable_total / max(len(edits), 1)
    avg_per_slide_filled = edited_total / max(len(edits), 1)

    # 점수
    score_role = (n_role_hit / n_total) * 100
    score_fill = fill_ratio * 100
    # overflow 안 함 (truncate 처리해서 잘림 없으나 짧아짐 - 부분 페널티)
    score_overflow_penalty = (1 - overflow_total / max(edited_total, 1)) * 100

    score_visual = visual_resolution * 100
    composite_quant = round(
        score_role * 0.25 + score_fill * 0.30 + score_visual * 0.30
        + score_overflow_penalty * 0.15, 1
    )

    return {
        "scenario_id": scenario_id,
        "narrative_length": n_total,
        "metrics": {
            "A_role_match": {
                "role_hit": n_role_hit, "archetype_fallback": n_archetype,
                "reuse_fallback": n_reuse, "miss": n_none,
                "role_hit_pct": round(score_role, 1),
            },
            "B_slot_fill_v2": {
                "fillable_total": fillable_total,
                "filled_total": edited_total,
                "blanked_total": blanked_total,
                "edit_failed": fail_total,
                "fill_ratio": round(fill_ratio, 4),
                "fill_pct": round(score_fill, 1),
                "visual_resolution_ratio": round(visual_resolution, 4),
                "visual_resolution_pct": round(score_visual, 1),
                "avg_fillable_per_slide": round(avg_per_slide_fillable, 1),
                "avg_filled_per_slide": round(avg_per_slide_filled, 1),
            },
            "C_overflow": {
                "n_overflow_truncated": overflow_total,
                "n_edited": edited_total,
                "overflow_rate_pct": round(
                    overflow_total / max(edited_total, 1) * 100, 1
                ),
                "no_overflow_score": round(score_overflow_penalty, 1),
            },
            "composite_quant_score": composite_quant,
        },
    }


# ----------------------------------------------------------------------------
# Main
# ----------------------------------------------------------------------------

def run_scenario(scenario_id: str, labels: list[dict], skeletons: dict,
                 store: ParagraphStore) -> dict:
    print()
    print("=" * 80)
    print(f"SCENARIO v2: {scenario_id}")
    print("=" * 80)

    sk = skeletons[scenario_id]
    narrative = sk["narrative_sequence"]
    sc = SCENARIO_CONTENT[scenario_id]

    # diversity-aware select: top-K 후보 중 capacity 기반 선호 + role 별 풀 안 겹침
    plan = select_deck_diverse(labels, narrative, sc, store)
    n_role = sum(1 for p in plan if p["source"] == "role")
    n_arch = sum(1 for p in plan if p["source"] == "archetype")
    print(f"  retrieval: role={n_role} archetype={n_arch}")

    out_dir = OUTPUT_ROOT / scenario_id
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_out = out_dir / "deck.pptx"
    png_dir = out_dir / "pngs"

    t0 = time.time()
    build = build_pptx_v2(plan, sc, store, pptx_out)
    dt_build = time.time() - t0
    print(f"  build: {dt_build:.1f}s")

    t0 = time.time()
    pngs = render_pngs(pptx_out, png_dir)
    dt_render = time.time() - t0
    print(f"  render: {dt_render:.1f}s, {len(pngs)} pngs")

    metrics = compute_v2_metrics(scenario_id, plan, build["edits"], narrative)
    metrics["paths"] = {
        "pptx": str(pptx_out),
        "png_dir": str(png_dir),
        "pngs": [str(p) for p in pngs],
    }
    metrics["plan"] = plan
    metrics["edits_summary"] = [
        {k: v for k, v in e.items() if k != "edits"} for e in build["edits"]
    ]
    metrics["timing"] = {"build_sec": round(dt_build, 1), "render_sec": round(dt_render, 1)}

    report_path = out_dir / "report.json"
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(metrics, f, ensure_ascii=False, indent=2)
    print(f"  report: {report_path}")

    return metrics


def main():
    target = sys.argv[1] if len(sys.argv) > 1 else None

    print("=" * 80)
    print("Phase A3 Step 3 — 5 Benchmark v2 (paragraph-aware fill)")
    print("=" * 80)

    labels = load_catalog()
    skeletons = load_skeletons()
    store = ParagraphStore.load()
    print(f"loaded: {len(labels)} labels / {len(skeletons)} skeletons "
          f"/ {store.n_slides} slides w/ paragraph store")

    if target:
        scenarios = [target]
    else:
        scenarios = list(SCENARIO_CONTENT.keys())

    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)

    all_results = []
    for sid in scenarios:
        res = run_scenario(sid, labels, skeletons, store)
        all_results.append(res)

    scoreboard = {
        "phase": "A3-Step3-v2",
        "mode": "Mode A + paragraph-fill",
        "n_scenarios": len(all_results),
        "results": all_results,
        "summary": {
            "avg_role_hit_pct": round(
                sum(r["metrics"]["A_role_match"]["role_hit_pct"] for r in all_results)
                / len(all_results), 1,
            ),
            "avg_fill_pct": round(
                sum(r["metrics"]["B_slot_fill_v2"]["fill_pct"] for r in all_results)
                / len(all_results), 1,
            ),
            "avg_visual_resolution_pct": round(
                sum(r["metrics"]["B_slot_fill_v2"]["visual_resolution_pct"]
                    for r in all_results) / len(all_results), 1,
            ),
            "avg_overflow_rate_pct": round(
                sum(r["metrics"]["C_overflow"]["overflow_rate_pct"] for r in all_results)
                / len(all_results), 1,
            ),
            "avg_composite_quant": round(
                sum(r["metrics"]["composite_quant_score"] for r in all_results)
                / len(all_results), 1,
            ),
        },
    }

    scoreboard_path = OUTPUT_ROOT / "scoreboard.json"
    with open(scoreboard_path, "w", encoding="utf-8") as f:
        json.dump(scoreboard, f, ensure_ascii=False, indent=2)

    print()
    print("=" * 80)
    print("SCOREBOARD v2")
    print("=" * 80)
    for r in all_results:
        m = r["metrics"]
        print(
            f"  {r['scenario_id']:>30} | role={m['A_role_match']['role_hit_pct']:5.1f}% "
            f"fill={m['B_slot_fill_v2']['fill_pct']:5.1f}% "
            f"visual={m['B_slot_fill_v2']['visual_resolution_pct']:5.1f}% "
            f"overflow={m['C_overflow']['overflow_rate_pct']:5.1f}% "
            f"comp={m['composite_quant_score']:5.1f}"
        )
    print()
    print(f"AVG composite quant = {scoreboard['summary']['avg_composite_quant']}")
    print(f"saved: {scoreboard_path}")


if __name__ == "__main__":
    main()
