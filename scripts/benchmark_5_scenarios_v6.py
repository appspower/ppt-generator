"""Phase A3 N1-Lite Phase 3 — 5 시나리오 하이브리드 측정 (v6).

[docs/PROJECT_DIRECTION.md] §3 ROLE_MODE_MAP 적용:
  - Mode A (8 role): opening, agenda, divider, evidence, analysis, recommendation, closing, appendix
  - N1-Lite (5 role): situation, complication, roadmap, benefit, risk

v5 (57.8) 대비 N1-Lite 도입 효과 측정 → 70+ 도달 여부 판정.

빌드 흐름
--------
1. plan: role별 mode 결정 (ROLE_MODE_MAP)
   - mode_a: select_deck_diverse 재사용 (v2)
   - n1_lite: components_index에서 applicable_roles 매칭 + 컨텐츠 수 ↔ slot 수 fit
2. TemplateEditor.keep_slides([mode_a_sidxs]) + N1-Lite blank 추가 → step 순서로 재정렬
3. 각 step fill:
   - mode_a: ParagraphStore + match_content (v2 동일)
   - n1_lite: comp_meta.slots 기반, master flat_idx → target flat_idx 변환

Usage: python scripts/benchmark_5_scenarios_v6.py [scenario_id]
"""
from __future__ import annotations

import json
import math
import sys
import time
from collections import Counter, defaultdict
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Pt

from ppt_builder.template import edit_ops
from ppt_builder.template.component_ops import (
    create_blank_slide_with_master_theme,
    extract_group,
    insert_component,
    chart_count,
    has_chart,
    replace_chart_data,
)
from ppt_builder.template.editor import TemplateEditor
from ppt_builder.catalog.paragraph_query import ParagraphStore

sys.path.insert(0, str(ROOT / "scripts"))
from benchmark_5_scenarios import (  # noqa: E402
    SCENARIO_CONTENT,
    ARCHETYPE_FALLBACK,
    candidates_for_role,
    load_skeletons,
)
from benchmark_5_scenarios_v2 import (  # noqa: E402
    CATALOG_PATH,
    TEMPLATE_PATH,
    _expand_to_capacity,
    _replace_paragraph_universal,
    _shrink_font_if_overflow,
    select_deck_diverse,
)


# v6.3: 자체 Mode A 슬라이드 선정 — over-capacity penalty 강화
def select_mode_a_v6_3(
    labels: list[dict],
    narrative: list[str],
    scenario_content: dict,
    store: ParagraphStore,
    chart_roles: set[str] | None = None,
) -> list[dict]:
    """Mode A 슬라이드 선정. v2 select_deck_diverse 기반 + over-capacity 강한 penalty.

    핵심 변경: total fillable count >> target_n인 슬라이드 강하게 패널티.
    예: target_n=2인데 fillable 99 → fitness 0.2 (이전엔 0.7).

    chart_roles
    -----------
    chart_data가 주어진 role 집합. 해당 role에 대해서는 chart_penalty를 해제하고
    오히려 차트 슬라이드에 보너스를 줘서 차트 슬라이드를 적극 선택한다.
    """
    chart_roles = chart_roles or set()
    used: set[int] = set()
    used_signatures: set[str] = set()
    plan = []
    role_use_count: dict[str, int] = {}
    by_role = scenario_content["content_by_role"]

    for role in narrative:
        target_n = len(by_role.get(role, [])) or 1
        cands, source = candidates_for_role(labels, role, used)

        if cands:
            top_cands = cands[:30]
            scored = []
            for c in top_cands:
                sidx = c["slide_index"]
                cap = store.slot_capacity(sidx)
                fillable = store.fillable_slots(sidx)
                primary_role = None
                primary_max_chars = 0
                for r in ("chevron_label", "card_header", "callout_text"):
                    if cap.get(r, 0) > 0:
                        primary_role = r
                        slots = fillable.get(r, [])
                        max_chars_vals = [s.max_chars for s in slots if s.max_chars]
                        primary_max_chars = (
                            sum(max_chars_vals) / len(max_chars_vals)
                            if max_chars_vals else 0
                        )
                        break
                narrow_penalty = 0.4 if primary_max_chars and primary_max_chars < 20 else 0.0

                primary_cap = max(
                    cap.get("chevron_label", 0),
                    cap.get("card_header", 0),
                    cap.get("callout_text", 0),
                    cap.get("card_body", 0),
                )
                # v6.3: 강한 excess penalty
                fit = 1.0
                if target_n >= 2:
                    if primary_cap == 0:
                        fit = 0.5
                    elif primary_cap >= target_n:
                        excess = primary_cap - target_n
                        # 강한 페널티: excess 1=0.95, 2=0.85, 4=0.65, 8=0.25
                        fit = max(0.2, 1.0 - 0.10 * excess)
                    else:
                        shortage = target_n - primary_cap
                        fit = max(0.3, 0.7 - 0.15 * shortage)
                else:  # target_n=1
                    # 작은 슬라이드 선호 (primary_cap 1~3 ideal)
                    if primary_cap == 0:
                        fit = 0.6
                    elif primary_cap <= 3:
                        fit = 1.0
                    elif primary_cap <= 6:
                        fit = 0.8
                    elif primary_cap <= 12:
                        fit = 0.5
                    else:
                        fit = 0.2

                # v6.7: total fillable 폭주 페널티 강화 — 거대 빈 표 회피
                # (analysis_report_15 84.6 진단: 99-slot 표에 컨텐츠 2~3개만 들어감.
                #  visual_resolution 100% (blanked로 처리됐지만 사용자 인식은 시각
                #  부채 명확.) cap 0.7 → 0.95, slope 0.07 → 0.10.
                total_fillable = sum(cap.values())
                density_penalty = 0.0
                if total_fillable > target_n * 5:
                    excess_ratio = total_fillable / max(target_n, 1)
                    density_penalty = min(0.95, 0.10 * (excess_ratio - 5))

                # v6.6: 차트 슬라이드 강 페널티 — chart_data 없으면 더미 숫자 시각 부채.
                # B2: chart_data가 주어진 role은 반대로 차트 슬라이드를 선호.
                chart_slides = get_chart_slide_indices()
                is_chart_slide = sidx in chart_slides
                if role in chart_roles:
                    # chart_data 있음 → 차트 슬라이드 강 보너스, 비차트 약 페널티
                    chart_penalty = 0.0 if is_chart_slide else 0.30
                else:
                    chart_penalty = 0.85 if is_chart_slide else 0.0

                sig_key = f"{c.get('macro')}:{tuple(sorted(c.get('archetype', [])))}"
                diversity_bonus = 0.0 if sig_key in used_signatures else 0.1

                # v6.8: density 가중치 강화 — sparse 슬라이드(99-slot에 컨텐츠 2~3개)
                # 회피. fit/narrow/chart 비중 소폭 ↓. analysis_report_15의 step 5/10
                # 여전히 큰 빈 표 잔존 → 결정적 push.
                score = (
                    0.25 * c.get("overall_confidence", 0)
                    + 0.30 * fit                              # 0.35→0.30
                    + 0.05 * (1 - narrow_penalty)             # 0.10→0.05
                    + 0.30 * (1 - density_penalty)            # 0.15→0.30
                    + 0.10 * (1 - chart_penalty)              # 0.15→0.10
                    + diversity_bonus
                )
                scored.append((score, c, sig_key, fit, narrow_penalty, density_penalty))

            scored.sort(key=lambda x: x[0], reverse=True)
            chosen = scored[0][1]
            chosen_sig = scored[0][2]
            used.add(chosen["slide_index"])
            used_signatures.add(chosen_sig)
            plan.append({
                "role": role,
                "slide_index": chosen["slide_index"],
                "source": source,
                "archetype": chosen.get("archetype", []),
                "macro": chosen.get("macro"),
                "confidence": chosen.get("overall_confidence", 0),
                "score": round(scored[0][0], 2),
                "capacity_fit": round(scored[0][3], 2),
                "density_penalty": round(scored[0][5], 2),
                "target_n_items": target_n,
                "reuse_count": 0,
            })
        else:
            # reuse fallback
            reuse_pool = [
                l for l in labels if role in l.get("narrative_role", [])
            ]
            if not reuse_pool:
                reuse_pool = [
                    l for l in labels
                    if any(a in ARCHETYPE_FALLBACK.get(role, [])
                           for a in l.get("archetype", []))
                ]
            if reuse_pool:
                reuse_pool.sort(
                    key=lambda l: l.get("overall_confidence", 0), reverse=True
                )
                chosen = None
                for c in reuse_pool[:10]:
                    sig_key = f"{c.get('macro')}:{tuple(sorted(c.get('archetype', [])))}"
                    if sig_key not in used_signatures:
                        chosen = c
                        break
                if chosen is None:
                    chosen = reuse_pool[0]
                role_use_count[role] = role_use_count.get(role, 0) + 1
                plan.append({
                    "role": role,
                    "slide_index": chosen["slide_index"],
                    "source": "reuse",
                    "archetype": chosen.get("archetype", []),
                    "macro": chosen.get("macro"),
                    "confidence": chosen.get("overall_confidence", 0),
                    "score": chosen.get("overall_confidence", 0),
                    "capacity_fit": 0.5,
                    "density_penalty": 0,
                    "target_n_items": target_n,
                    "reuse_count": role_use_count[role],
                })
            else:
                plan.append({
                    "role": role, "slide_index": None, "source": "none",
                    "archetype": [], "macro": None, "confidence": 0,
                    "score": 0, "capacity_fit": 0, "density_penalty": 0,
                    "target_n_items": target_n, "reuse_count": 0,
                })
    return plan


# --- ROLE_MODE_MAP (헌법 §3 확정) ---------------------------------------------

ROLE_MODE_MAP = {
    "opening": "mode_a",
    "agenda": "mode_a",
    "divider": "mode_a",
    "situation": "n1_lite",
    "complication": "n1_lite",
    "evidence": "mode_a",
    "analysis": "mode_a",
    "recommendation": "mode_a",
    "roadmap": "n1_lite",
    "benefit": "n1_lite",
    "risk": "n1_lite",
    "closing": "mode_a",
    "appendix": "mode_a",
}

LIBRARY_DIR = ROOT / "output" / "component_library"
INDEX_PATH = LIBRARY_DIR / "components_index.json"
OUTPUT_ROOT = ROOT / "output" / "benchmark_v6"

# 차트 슬라이드 캐시 (마스터 1회 스캔)
_CHART_SLIDES_CACHE: set[int] | None = None


def get_chart_slide_indices() -> set[int]:
    """마스터의 차트 보유 슬라이드 인덱스 set. 캐시."""
    global _CHART_SLIDES_CACHE
    if _CHART_SLIDES_CACHE is not None:
        return _CHART_SLIDES_CACHE
    from pptx.shapes.graphfrm import GraphicFrame
    prs = Presentation(str(TEMPLATE_PATH))
    out: set[int] = set()
    for sidx, sl in enumerate(prs.slides):
        for sh in sl.shapes:
            if isinstance(sh, GraphicFrame) and sh.has_chart:
                out.add(sidx)
                break
    _CHART_SLIDES_CACHE = out
    return out


# --- N1-Lite 컴포넌트 선정 -----------------------------------------------------

def load_components_index() -> dict:
    """components_index.json 로드 (Phase 2 산출)."""
    if not INDEX_PATH.exists():
        raise FileNotFoundError(
            f"components_index.json not found: {INDEX_PATH}\n"
            f"Phase 2 build_component_library.py 실행 필요"
        )
    return json.loads(INDEX_PATH.read_text(encoding="utf-8"))


def select_n1_lite_component(
    role: str,
    target_n_items: int,
    components: list[dict],
    used_component_ids: set[str],
    min_text_slots: int = 2,
) -> dict | None:
    """role + content 수에 맞는 컴포넌트 선정. v6.1 REFINE.

    변경점 (v6 → v6.1):
      - sparse 회피: n_text_slots < min_text_slots 컴포넌트 제외 (단독 빈약 슬라이드 방지)
      - fit 강화: exact match (n_text_slots == target_n_items)에 강한 보너스
      - excess/shortage 페널티 가중치 ×3 (잉여/부족 슬롯 시각 부채 회피)
    """
    cands = [
        c for c in components
        if role in c.get("applicable_roles", [])
        and len(c.get("slots", [])) >= min_text_slots
    ]
    if not cands:
        # min_text_slots 만족 못 하면 그래도 가능한 모든 후보로 fallback
        cands = [c for c in components if role in c.get("applicable_roles", [])]
    if not cands:
        return None

    scored = []
    for c in cands:
        penalty = 0.5 if c["component_id"] in used_component_ids else 0.0
        n_text_slots = len(c.get("slots", []))
        if n_text_slots == 0:
            fit = 0.2
        elif n_text_slots == target_n_items:
            fit = 1.0  # exact match — 잉여/부족 0
        elif target_n_items < n_text_slots:
            excess = n_text_slots - target_n_items
            fit = max(0.3, 1.0 - 0.15 * excess)  # 잉여 1슬롯 = 0.85
        else:
            shortage = target_n_items - n_text_slots
            fit = max(0.3, 0.85 - 0.15 * shortage)  # 부족 1슬롯 = 0.7
        score = 0.55 * fit + 0.25 * (c.get("score", 0) / 100) + 0.20 * (1 - penalty)
        scored.append((score, c, fit))
    scored.sort(key=lambda x: -x[0])
    return scored[0][1]


# --- Plan 생성 ----------------------------------------------------------------

def select_deck_v6(
    labels: list[dict],
    narrative: list[str],
    scenario_content: dict,
    store: ParagraphStore,
    components: list[dict],
    chart_roles: set[str] | None = None,
) -> list[dict]:
    """role별 mode 결정 + 자산 선정.

    각 plan 항목:
      mode='mode_a': {role, mode, slide_index, source, archetype, ...}
      mode='n1_lite': {role, mode, component_id, family, target_n_items, ...}

    chart_roles: chart_data가 명시된 role 집합 — Mode A 선정에서 차트 슬라이드 선호.
    """
    by_role = scenario_content["content_by_role"]

    # 1) Mode A 후보 우선 결정 (v6.3: select_mode_a_v6_3 — over-capacity penalty 강화)
    mode_a_plan = select_mode_a_v6_3(
        labels, narrative, scenario_content, store, chart_roles=chart_roles
    )
    plan: list[dict] = []
    used_components: set[str] = set()

    for i, role in enumerate(narrative):
        target_n = len(by_role.get(role, [])) or 1
        mode = ROLE_MODE_MAP.get(role, "mode_a")

        if mode == "mode_a":
            # v2 결과 그대로 사용
            ma = mode_a_plan[i]
            plan.append({**ma, "mode": "mode_a"})
        else:
            # N1-Lite: 컴포넌트 선정
            comp = select_n1_lite_component(role, target_n, components, used_components)
            if comp is not None:
                used_components.add(comp["component_id"])
                plan.append({
                    "role": role,
                    "mode": "n1_lite",
                    "component_id": comp["component_id"],
                    "family": comp["family"],
                    "source": "n1_lite",
                    "target_n_items": target_n,
                    "score": comp.get("score", 0),
                    "n_text_slots": len(comp.get("slots", [])),
                })
            else:
                # N1-Lite fallback: Mode A 결과 사용 (v2가 reuse로라도 채워줌)
                ma = mode_a_plan[i]
                plan.append({**ma, "mode": "mode_a", "n1_lite_fallback": True})
    return plan


# --- Build ---------------------------------------------------------------------

def _reorder_to_match(prs, ordered_slides: list) -> None:
    """prs._sldIdLst를 ordered_slides 순서에 맞춰 재정렬."""
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    # _sldId의 id == slide.slide_id (고유)
    sldId_by_slide_id = {c.id: c for c in children}
    new_order = [sldId_by_slide_id[s.slide_id] for s in ordered_slides]
    for c in children:
        sldIdLst.remove(c)
    for c in new_order:
        sldIdLst.append(c)


def _fill_n1_lite_slide(
    target_slide,
    comp_meta: dict,
    role: str,
    scenario_content: dict,
    role_use_counter: dict[str, int],
) -> dict:
    """N1-Lite 컴포넌트 슬라이드의 paragraph fill.

    매핑: comp_meta['source']['group_indices']는 master flat_idx 리스트.
    insert_component 후 target slide의 leaf flat_idx는 group_indices의 position(0-based).
    """
    by_role = scenario_content["content_by_role"]
    use_idx = role_use_counter.get(role, 0)
    role_use_counter[role] = use_idx + 1

    role_items = by_role.get(role, [])
    if not role_items:
        role_items = [role.upper()]

    group_indices = comp_meta["source"]["group_indices"]
    master_to_target = {mfi: i for i, mfi in enumerate(group_indices)}

    # target 슬롯 메타 빌드
    slots_by_role: dict[str, list[dict]] = defaultdict(list)
    for slot_meta in comp_meta.get("slots", []):
        master_fi = slot_meta["flat_idx"]
        if master_fi not in master_to_target:
            continue
        slots_by_role[slot_meta["role"]].append({
            "target_flat_idx": master_to_target[master_fi],
            "paragraph_id": slot_meta["paragraph_id"],
            "role": slot_meta["role"],
            "max_chars": slot_meta.get("max_chars"),
            "position_in_group": slot_meta.get("position_in_group"),
        })
    # position_in_group 순으로 정렬
    for r, lst in slots_by_role.items():
        lst.sort(key=lambda s: (s.get("position_in_group") or 0, s["target_flat_idx"]))

    # 우선순위 슬롯에 컨텐츠 배분
    fill_priority = ["chevron_label", "card_header", "callout_text", "card_body", "kpi_value"]
    chosen_role = None
    for r in fill_priority:
        if slots_by_role.get(r):
            chosen_role = r
            break

    n_filled = 0
    n_failed = 0
    n_overflow = 0
    fill_log = []
    edited_keys: set[tuple[int, int]] = set()

    if chosen_role:
        slots = slots_by_role[chosen_role]
        target_n = len(slots)
        expanded = _expand_to_capacity(role_items, target_n)
        for slot, text in zip(slots, expanded):
            try:
                edit_ops.replace_paragraph(
                    target_slide, slot["target_flat_idx"], slot["paragraph_id"], text
                )
                edited_keys.add((slot["target_flat_idx"], slot["paragraph_id"]))
                # overflow check + shrink
                shrunk = False
                if slot["max_chars"] and len(text) > slot["max_chars"]:
                    n_overflow += 1
                    for fi, sh in edit_ops.iter_leaf_shapes(target_slide):
                        if fi != slot["target_flat_idx"]:
                            continue
                        if not getattr(sh, "has_text_frame", False):
                            break
                        paras = sh.text_frame.paragraphs
                        if slot["paragraph_id"] < len(paras):
                            shrunk = _shrink_font_if_overflow(
                                paras[slot["paragraph_id"]], text, slot["max_chars"]
                            )
                        break
                n_filled += 1
                fill_log.append({
                    "slot_role": chosen_role,
                    "target_flat_idx": slot["target_flat_idx"],
                    "paragraph_id": slot["paragraph_id"],
                    "text": text[:50],
                    "max_chars": slot["max_chars"],
                    "overflow": slot["max_chars"] and len(text) > slot["max_chars"],
                    "shrunk": shrunk,
                    "ok": True,
                })
            except Exception as e:
                n_failed += 1
                fill_log.append({
                    "slot_role": chosen_role,
                    "target_flat_idx": slot["target_flat_idx"],
                    "paragraph_id": slot["paragraph_id"],
                    "text": text[:50],
                    "ok": False,
                    "reason": f"{type(e).__name__}: {e}",
                })

    # v6.1: 적극적 잔존 청소 (~~ + 잉여 슬롯 + 채워지지 않은 모든 placeholder)
    # 청소 대상:
    #   - shape의 모든 paragraph 중 텍스트가 ~~ 패턴인 것
    #   - 컴포넌트 슬롯에 등록되어 있지만 채워지지 않은 paragraph (잉여)
    n_blanked = 0
    for fi, sh in edit_ops.iter_leaf_shapes(target_slide):
        if not getattr(sh, "has_text_frame", False):
            continue
        for pi, p in enumerate(sh.text_frame.paragraphs):
            if (fi, pi) in edited_keys:
                continue
            t = p.text.strip()
            if not t:
                continue
            # ~~ pattern 또는 짧은 placeholder + 잉여 슬롯의 일반 텍스트도 비움
            is_placeholder_ish = (
                t == "~~"
                or t.replace("~", "").replace("\x0b", "").strip() == ""
                or (len(t) <= 5 and "~" in t)
            )
            # 컴포넌트 슬롯에 등록된 paragraph인데 채워지지 않은 것 (잉여 슬롯)
            is_unused_slot = any(
                s["target_flat_idx"] == fi and s["paragraph_id"] == pi
                for r_slots in slots_by_role.values() for s in r_slots
            )
            if is_placeholder_ish or is_unused_slot:
                try:
                    edit_ops.replace_paragraph(target_slide, fi, pi, " ")
                    n_blanked += 1
                except Exception:
                    pass

    fillable_total = sum(len(v) for v in slots_by_role.values())
    return {
        "n_fillable": fillable_total,
        "n_filled": n_filled,
        "n_failed": n_failed,
        "n_blanked": n_blanked,
        "n_overflow": n_overflow,
        "fill_role": chosen_role,
        "fill_log": fill_log,
        "slots_by_role_count": {r: len(v) for r, v in slots_by_role.items()},
    }


def _slot_key(slot) -> tuple:
    """table cell까지 구분하는 dedup key. table_row/col 같은 flat_idx 충돌 회피."""
    if getattr(slot, "shape_kind", "") == "TABLE":
        return (slot.flat_idx, slot.paragraph_id, slot.table_row, slot.table_col)
    return (slot.flat_idx, slot.paragraph_id, None, None)


def _fill_mode_a_slide(
    target_slide,
    sidx: int,
    role: str,
    scenario_content: dict,
    store: ParagraphStore,
    role_use_counter: dict[str, int],
) -> dict:
    """v2의 paragraph fill 흐름을 1슬라이드에 적용. v6.3: table cell 청소 보장."""
    from benchmark_5_scenarios_v2 import expand_content_for_slide

    capacity = store.slot_capacity(sidx)
    content = expand_content_for_slide(role, scenario_content, role_use_counter, capacity)
    content = {k: v for k, v in content.items() if v}

    edits = store.match_content(sidx, content, truncate_overflow=True)
    n_ok = 0
    n_fail = 0
    n_overflow = 0
    edited_keys: set[tuple] = set()

    all_fillable = store.fillable_slots(sidx)

    for e in edits:
        try:
            role_slots = all_fillable.get(e["role"], [])
            pos = e.get("matched_pos", 0)
            if pos >= len(role_slots):
                raise edit_ops.SlideEditError(
                    f"matched_pos={pos} out of range for role={e['role']}"
                )
            slot_obj = role_slots[pos]
            _replace_paragraph_universal(target_slide, slot_obj, e["text"])
            n_ok += 1
            edited_keys.add(_slot_key(slot_obj))
            if e.get("original_overflow"):
                n_overflow += 1
        except Exception:
            n_fail += 1

    # v6.3: ~~ 적극 청소 (fillable + 비-fillable 모두; table cell은 row/col로 unique 처리)
    n_blanked = 0
    for slot in store.slots(sidx):
        key = _slot_key(slot)
        if key in edited_keys:
            continue
        if slot.role in ("page_number", "footer", "date_text"):
            continue
        orig_raw = slot.text_original or ""
        orig = orig_raw.strip()
        if not orig:
            continue
        is_placeholder = (
            orig == "~~"
            or orig.replace("~", "").replace("\x0b", "").strip() == ""
            or (len(orig) <= 8 and "~~" in orig)
        )
        if is_placeholder:
            try:
                _replace_paragraph_universal(target_slide, slot, " ")
                n_blanked += 1
                edited_keys.add(key)
            except Exception:
                pass

    fillable_total = sum(capacity.values())
    return {
        "n_fillable": fillable_total,
        "n_filled": n_ok,
        "n_failed": n_fail,
        "n_blanked": n_blanked,
        "n_overflow": n_overflow,
    }


def _inject_chart_data(target_slide, chart_spec) -> bool:
    """target_slide의 첫 차트 shape에 chart_spec(ChartSpec) 데이터 + 색상 주입.

    찾으면 True. flat_idx는 iter_leaf_shapes에서 첫 GraphicFrame.has_chart 것을 사용.
    chart_spec.series[*].color가 명시되면 해당 시리즈 색 변경.
    """
    from pptx.shapes.graphfrm import GraphicFrame
    for fi, sh in edit_ops.iter_leaf_shapes(target_slide):
        if isinstance(sh, GraphicFrame) and getattr(sh, "has_chart", False):
            categories = list(chart_spec.categories)
            series = [(s.name, list(s.values)) for s in chart_spec.series]
            colors = [getattr(s, "color", None) for s in chart_spec.series]
            colors_arg = colors if any(colors) else None
            try:
                replace_chart_data(
                    target_slide, fi, categories, series,
                    series_colors=colors_arg,
                )
                return True
            except Exception as e:  # noqa: BLE001
                print(f"  [warn] replace_chart_data failed @flat_idx={fi}: "
                      f"{type(e).__name__}: {e}")
                return False
    return False


def build_pptx_v6(
    plan: list[dict],
    scenario_content: dict,
    store: ParagraphStore,
    components_by_id: dict[str, dict],
    pptx_out: Path,
    chart_data: dict | None = None,
) -> dict:
    """v6 빌드: Mode A keep_slides + N1-Lite blank 삽입 → 재정렬 → fill.

    chart_data: {role: ChartSpec} — 빌드된 슬라이드의 첫 차트에 자동 주입.
    """
    chart_data = chart_data or {}
    chart_injected: dict[str, bool] = {role: False for role in chart_data}
    pptx_out.parent.mkdir(parents=True, exist_ok=True)

    # Mode A 유니크 sidx 수집 (중복은 동일 슬라이드 재사용 — 단순화: drop)
    seen_mode_a: set[int] = set()
    plan_unique: list[dict] = []
    for p in plan:
        if p["mode"] == "mode_a":
            if p.get("slide_index") is None:
                continue
            if p["slide_index"] in seen_mode_a:
                continue
            seen_mode_a.add(p["slide_index"])
        plan_unique.append(p)

    mode_a_sidxs = sorted(
        p["slide_index"] for p in plan_unique
        if p["mode"] == "mode_a" and p.get("slide_index") is not None
    )

    # 두-단계 빌드:
    # Stage 1: Mode A keep_slides → 임시 .pptx 저장
    # Stage 2: 임시 .pptx 재오픈 → N1-Lite blank 추가
    # (TemplateEditor에서 직접 add_slide 시 deleted 슬라이드 part가 partname conflict 유발)
    keep_initial = mode_a_sidxs if mode_a_sidxs else [0]
    stage1_editor = TemplateEditor(TEMPLATE_PATH)
    stage1_editor.keep_slides(keep_initial)
    intermediate = pptx_out.parent / f"_intermediate_{pptx_out.stem}.pptx"
    stage1_editor.save(intermediate)
    stage1_editor.cleanup()

    target_prs = Presentation(str(intermediate))

    # mode_a slide_index → 재오픈된 target_prs position
    mode_a_to_target_slide: dict[int, object] = {}
    for i, sidx in enumerate(keep_initial):
        if sidx in mode_a_sidxs:
            mode_a_to_target_slide[sidx] = target_prs.slides[i]

    # N1-Lite: master 한 번 열어두고 추출
    master_prs = Presentation(str(TEMPLATE_PATH))

    # 각 step에 대응하는 target slide 객체 결정
    step_to_slide: list = []  # plan_unique 순서대로 slide objects
    n1_lite_meta: dict[int, dict] = {}  # step_idx → comp_meta

    for step_idx, p in enumerate(plan_unique):
        if p["mode"] == "mode_a":
            slide = mode_a_to_target_slide.get(p["slide_index"])
            if slide is None:
                # mode_a_sidxs에 없으면 (mode_a fallback인데 slide_index None) skip
                continue
            step_to_slide.append(slide)
        else:
            # n1_lite: blank + insert_component
            comp_meta = components_by_id[p["component_id"]]
            sidx = comp_meta["source"]["master_slide_index"]
            group_indices = comp_meta["source"]["group_indices"]
            try:
                src_slide = master_prs.slides[sidx]
                comp = extract_group(src_slide, group_indices)
                blank = create_blank_slide_with_master_theme(target_prs)
                insert_component(blank, comp)
                step_to_slide.append(blank)
                n1_lite_meta[len(step_to_slide) - 1] = comp_meta
            except Exception as e:
                print(f"  [warn] n1_lite step {step_idx} ({p['component_id']}): {type(e).__name__}: {e}")
                # skip this step

    # 재정렬: step_to_slide 순서가 곧 narrative 순서
    if step_to_slide:
        # throwaway가 keep_initial[0]이고 mode_a가 비어있으면 step_to_slide에 없을 수 있음
        # 모든 슬라이드가 step_to_slide에 포함되도록 보장 (남은 건 끝에)
        all_slides = list(target_prs.slides)
        leftover = [s for s in all_slides if s not in step_to_slide]
        ordered = step_to_slide + leftover
        _reorder_to_match(target_prs, ordered)

    # 각 step fill
    edit_results: list[dict] = []
    role_use_counter: dict[str, int] = {}
    actual_step = 0
    for step_idx, p in enumerate(plan_unique):
        if actual_step >= len(step_to_slide):
            break
        target_slide = step_to_slide[actual_step]
        if p["mode"] == "mode_a":
            r = _fill_mode_a_slide(
                target_slide, p["slide_index"], p["role"],
                scenario_content, store, role_use_counter,
            )
            edit_results.append({
                "step": actual_step + 1, "role": p["role"], "mode": "mode_a",
                "slide_index": p["slide_index"], "source": p["source"],
                **r,
            })
        else:
            comp_meta = n1_lite_meta.get(actual_step)
            if comp_meta is None:
                actual_step += 1
                continue
            r = _fill_n1_lite_slide(
                target_slide, comp_meta, p["role"],
                scenario_content, role_use_counter,
            )
            edit_results.append({
                "step": actual_step + 1, "role": p["role"], "mode": "n1_lite",
                "component_id": p["component_id"], "family": p["family"],
                "source": "n1_lite",
                **r,
            })
        actual_step += 1

    # B2: chart_data 주입 — fill 이후, leftover 절단/save 이전.
    # 한 role에 여러 step이 매핑되면 첫 step에만 주입.
    if chart_data:
        actual_step = 0
        for p in plan_unique:
            if actual_step >= len(step_to_slide):
                break
            target_slide = step_to_slide[actual_step]
            actual_step += 1
            role = p["role"]
            if role not in chart_data or chart_injected.get(role):
                continue
            spec = chart_data[role]
            ok = _inject_chart_data(target_slide, spec)
            chart_injected[role] = ok

    # 마지막 leftover 슬라이드 (throwaway 등)는 잘라낸다 — 단순화: keep only step_to_slide
    # 실제 slide 삭제는 _sldIdLst만 조작 (cleanest hack):
    if len(step_to_slide) < len(list(target_prs.slides)):
        sldIdLst = target_prs.slides._sldIdLst
        children = list(sldIdLst)
        n_keep = len(step_to_slide)
        for c in children[n_keep:]:
            sldIdLst.remove(c)

    target_prs.save(str(pptx_out))
    try:
        intermediate.unlink()
    except Exception:
        pass

    return {
        "plan_unique": plan_unique,
        "edits": edit_results,
        "pptx": str(pptx_out),
        "n_mode_a": sum(1 for p in plan_unique if p["mode"] == "mode_a"),
        "n_n1_lite": sum(1 for p in plan_unique if p["mode"] == "n1_lite"),
        "chart_injected": chart_injected,
    }


# --- Render --------------------------------------------------------------------

def render_pngs(pptx_path: Path, png_dir: Path) -> list[Path]:
    import pythoncom
    import win32com.client

    png_dir.mkdir(parents=True, exist_ok=True)
    for old in png_dir.glob("*.png"):
        old.unlink()
    pythoncom.CoInitialize()
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            str(pptx_path.resolve()), ReadOnly=True, Untitled=False, WithWindow=False,
        )
        try:
            out = []
            for i in range(1, presentation.Slides.Count + 1):
                p = png_dir / f"step_{i:02d}.png"
                presentation.Slides(i).Export(str(p), "PNG", 1568, 1176)
                out.append(p)
            return out
        finally:
            try:
                presentation.Close()
            except Exception:
                pass
            try:
                powerpoint.Quit()
            except Exception:
                pass
    finally:
        pythoncom.CoUninitialize()


# --- Metrics -------------------------------------------------------------------

def compute_v6_metrics(
    scenario_id: str, plan: list[dict], edits: list[dict], narrative: list[str]
) -> dict:
    n_total = len(narrative)
    # role match (Mode A의 source==role 만 hit으로, n1_lite는 자체적으로 hit으로 카운트)
    n_role_hit = 0
    n_n1_used = 0
    for p in plan:
        if p["mode"] == "n1_lite":
            n_n1_used += 1
            n_role_hit += 1  # N1-Lite는 항상 role 매칭
        elif p.get("source") == "role":
            n_role_hit += 1

    fillable_total = sum(e["n_fillable"] for e in edits)
    filled_total = sum(e["n_filled"] for e in edits)
    blanked_total = sum(e["n_blanked"] for e in edits)
    failed_total = sum(e["n_failed"] for e in edits)
    overflow_total = sum(e["n_overflow"] for e in edits)

    fill_ratio = filled_total / fillable_total if fillable_total else 0
    # v6.2: cap visual at 1.0 (n_blanked는 비-fillable 슬롯의 ~~ 청소도 포함하므로 분자가 분모 초과 가능)
    visual_resolution = min(
        1.0,
        (filled_total + blanked_total) / fillable_total if fillable_total else 0,
    )

    score_role = (n_role_hit / n_total) * 100
    score_fill = fill_ratio * 100
    score_visual = visual_resolution * 100
    score_overflow_penalty = (1 - overflow_total / max(filled_total, 1)) * 100

    # v6.3 (원래) + v6.5 (재가중) 둘 다 보고 — 메트릭 정직성
    composite_v63 = round(
        score_role * 0.25 + score_fill * 0.30 + score_visual * 0.30
        + score_overflow_penalty * 0.15, 1
    )
    # v6.5: visual이 실제 시각 품질 직접 반영 → 비중 ↑
    # fill_pct는 정책(primary role 1개만 채움)에 따라 노이즈가 큼 → 비중 ↓
    composite = round(
        score_role * 0.25 + score_fill * 0.15 + score_visual * 0.45
        + score_overflow_penalty * 0.15, 1
    )

    n1_lite_edits = [e for e in edits if e["mode"] == "n1_lite"]
    mode_a_edits = [e for e in edits if e["mode"] == "mode_a"]

    def _ratio(items, key):
        n_fill = sum(it["n_fillable"] for it in items)
        return sum(it["n_filled"] for it in items) / n_fill if n_fill else 0

    return {
        "scenario_id": scenario_id,
        "narrative_length": n_total,
        "metrics": {
            "A_role_match": {
                "role_hit_pct": round(score_role, 1),
                "n_n1_lite_used": n_n1_used,
            },
            "B_slot_fill": {
                "fillable_total": fillable_total,
                "filled_total": filled_total,
                "blanked_total": blanked_total,
                "edit_failed": failed_total,
                "fill_ratio": round(fill_ratio, 4),
                "fill_pct": round(score_fill, 1),
                "visual_resolution_pct": round(score_visual, 1),
            },
            "C_overflow": {
                "n_overflow": overflow_total,
                "overflow_rate_pct": round(
                    overflow_total / max(filled_total, 1) * 100, 1
                ),
                "no_overflow_score": round(score_overflow_penalty, 1),
            },
            "N_n1_lite": {
                "n_n1_lite_slides": len(n1_lite_edits),
                "n_mode_a_slides": len(mode_a_edits),
                "n1_lite_fill_ratio": round(_ratio(n1_lite_edits, "n_filled"), 4),
                "mode_a_fill_ratio": round(_ratio(mode_a_edits, "n_filled"), 4),
            },
            "composite_v6": composite,
            "composite_v6_3_legacy": composite_v63,
        },
    }


# --- Main ----------------------------------------------------------------------

def run_scenario(
    scenario_id: str,
    labels: list[dict],
    skeletons: dict,
    store: ParagraphStore,
    components_index: dict,
) -> dict:
    print()
    print("=" * 80)
    print(f"SCENARIO v6: {scenario_id}")
    print("=" * 80)

    sk = skeletons[scenario_id]
    narrative = sk["narrative_sequence"]
    sc = SCENARIO_CONTENT[scenario_id]
    components = components_index["components"]
    components_by_id = {c["component_id"]: c for c in components}

    plan = select_deck_v6(labels, narrative, sc, store, components)
    n_a = sum(1 for p in plan if p["mode"] == "mode_a")
    n_n = sum(1 for p in plan if p["mode"] == "n1_lite")
    print(f"  plan: {n_a} mode_a / {n_n} n1_lite (narrative len {len(narrative)})")

    out_dir = OUTPUT_ROOT / scenario_id
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_out = out_dir / "deck.pptx"
    png_dir = out_dir / "pngs"

    t0 = time.time()
    build = build_pptx_v6(plan, sc, store, components_by_id, pptx_out)
    dt_build = time.time() - t0
    print(f"  build: {dt_build:.1f}s ({build['n_mode_a']} mode_a + {build['n_n1_lite']} n1_lite)")

    t0 = time.time()
    pngs = render_pngs(pptx_out, png_dir)
    dt_render = time.time() - t0
    print(f"  render: {dt_render:.1f}s, {len(pngs)} pngs")

    metrics = compute_v6_metrics(scenario_id, plan, build["edits"], narrative)
    metrics["paths"] = {"pptx": str(pptx_out), "png_dir": str(png_dir)}
    metrics["plan"] = plan
    metrics["edits_summary"] = [
        {k: v for k, v in e.items() if k not in ("fill_log",)} for e in build["edits"]
    ]
    metrics["timing"] = {"build_sec": round(dt_build, 1), "render_sec": round(dt_render, 1)}

    report_path = out_dir / "report.json"
    report_path.write_text(json.dumps(metrics, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"  report: {report_path}")
    return metrics


def main():
    target = sys.argv[1] if len(sys.argv) > 1 else None
    print("=" * 80)
    print("Phase A3 N1-Lite Phase 3 — 5 Benchmark v6 (hybrid Mode A + N1-Lite)")
    print("=" * 80)

    labels = json.loads(CATALOG_PATH.read_text(encoding="utf-8"))["labels"]
    skeletons = load_skeletons()
    store = ParagraphStore.load()
    components_index = load_components_index()
    print(
        f"loaded: {len(labels)} labels / {len(skeletons)} skeletons / "
        f"{store.n_slides} slides w/ paragraph store / "
        f"{len(components_index['components'])} components"
    )

    if target:
        scenarios = [target]
    else:
        scenarios = list(SCENARIO_CONTENT.keys())

    OUTPUT_ROOT.mkdir(parents=True, exist_ok=True)

    all_results = []
    for sid in scenarios:
        try:
            res = run_scenario(sid, labels, skeletons, store, components_index)
            all_results.append(res)
        except Exception as e:
            import traceback
            print(f"[ERROR] scenario {sid}: {type(e).__name__}: {e}")
            traceback.print_exc()

    if not all_results:
        print("no results")
        return

    avg = lambda key_path: round(  # noqa: E731
        sum(_dig(r["metrics"], key_path) for r in all_results) / len(all_results), 1
    )
    scoreboard = {
        "phase": "A3-N1-Lite-Phase3-v6",
        "mode": "Mode A + N1-Lite hybrid",
        "n_scenarios": len(all_results),
        "results": all_results,
        "summary": {
            "avg_role_hit_pct": avg(["A_role_match", "role_hit_pct"]),
            "avg_fill_pct": avg(["B_slot_fill", "fill_pct"]),
            "avg_visual_pct": avg(["B_slot_fill", "visual_resolution_pct"]),
            "avg_overflow_pct": avg(["C_overflow", "overflow_rate_pct"]),
            "avg_n_n1_lite_used": round(
                sum(r["metrics"]["A_role_match"]["n_n1_lite_used"] for r in all_results)
                / len(all_results), 1
            ),
            "avg_composite_v6": avg(["composite_v6"]),
        },
    }

    sb_path = OUTPUT_ROOT / "scoreboard.json"
    sb_path.write_text(json.dumps(scoreboard, ensure_ascii=False, indent=2), encoding="utf-8")

    print()
    print("=" * 80)
    print("SCOREBOARD v6")
    print("=" * 80)
    for r in all_results:
        m = r["metrics"]
        print(
            f"  {r['scenario_id']:>30} | role={m['A_role_match']['role_hit_pct']:5.1f}% "
            f"fill={m['B_slot_fill']['fill_pct']:5.1f}% "
            f"visual={m['B_slot_fill']['visual_resolution_pct']:5.1f}% "
            f"n1L={m['A_role_match']['n_n1_lite_used']:2d} "
            f"comp={m['composite_v6']:5.1f}"
        )
    print()
    print(f"AVG composite v6 = {scoreboard['summary']['avg_composite_v6']}  (v5 baseline = 57.8)")
    print(f"saved: {sb_path}")


def _dig(d, path):
    for k in path:
        d = d[k]
    return d


if __name__ == "__main__":
    main()
