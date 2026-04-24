"""Phase A2 심층 layout clustering 분석.

1251장 전수:
- shape type / name / geometry 분포
- A01~ 태그 (placeholder type, title_text 패턴) 검증
- layout 한글 이름 복원
- 10x10 grid occupancy 기반 unique layout 수 측정
- clustering feature 추천
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from collections import Counter, defaultdict
import json
import re
import os
import hashlib

PATH = 'c:/Users/y2kbo/Apps/PPT/docs/references/_master_templates/PPT 템플릿.pptx'
OUT_DIR = 'c:/Users/y2kbo/Apps/PPT/output/catalog'
os.makedirs(OUT_DIR, exist_ok=True)

# 정규화 기준 (1920x1080)
# 원본 슬라이드 크기를 읽어 scale
CODE_RE = re.compile(r'^[A-Z]\d{1,2}$')  # A01, B12 등


def decode_layout_name(name: str) -> str:
    """cp949/UTF-8 mojibake 복원 시도."""
    if not name:
        return ''
    # 이미 한글이면 그대로
    try:
        name.encode('ascii')
    except UnicodeEncodeError:
        return name
    # mojibake 복원 시도
    for src_enc, tgt_enc in [('latin-1', 'cp949'), ('latin-1', 'utf-8'), ('cp1252', 'cp949')]:
        try:
            decoded = name.encode(src_enc).decode(tgt_enc)
            if any('가' <= c <= '힣' for c in decoded):
                return decoded
        except Exception:
            pass
    return name


def normalize_xy(left, top, slide_w, slide_h):
    """EMU 좌표를 [0,1]로 정규화."""
    if left is None or slide_w == 0:
        return 0.0, 0.0
    return max(0.0, min(1.0, left / slide_w)), max(0.0, min(1.0, top / slide_h))


def grid_cell(cx, cy, grid=10):
    """중심점을 10x10 grid에 매핑."""
    gx = min(grid - 1, max(0, int(cx * grid)))
    gy = min(grid - 1, max(0, int(cy * grid)))
    return gy * grid + gx


def iter_shapes(shapes, depth=0):
    for shp in shapes:
        yield shp, depth
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from iter_shapes(shp.shapes, depth + 1)
            except Exception:
                pass


def analyze():
    prs = Presentation(PATH)
    slide_w = prs.slide_width  # EMU
    slide_h = prs.slide_height
    print(f"Slide dim (EMU): {slide_w} x {slide_h}")
    print(f"Slide dim (inch): {slide_w/914400:.2f} x {slide_h/914400:.2f}")

    shape_type_ctr = Counter()
    auto_shape_ctr = Counter()
    placeholder_type_ctr = Counter()
    shape_name_ctr = Counter()
    shape_name_prefix_ctr = Counter()  # 앞 3글자
    layout_name_ctr = Counter()
    layout_decoded_ctr = Counter()
    shape_count_per_slide = []
    first3_names_samples = []  # (slide_idx, [names])
    first3_texts_samples = []
    code_tag_positions = []  # A01 같은 태그가 발견된 위치

    # grid occupancy feature
    grid_sigs_plain = []  # binary 100-dim
    grid_sigs_typed = []  # (cell, shape_type) 튜플
    structure_sigs = []  # shape type histogram

    geometry_samples = defaultdict(list)  # shape_type → (left, top, w, h) 샘플

    layout_first_slides = defaultdict(list)  # layout_name → slide index

    for i, slide in enumerate(prs.slides):
        layout_raw = slide.slide_layout.name if slide.slide_layout else ''
        layout_decoded = decode_layout_name(layout_raw)
        layout_name_ctr[layout_raw] += 1
        layout_decoded_ctr[layout_decoded] += 1
        layout_first_slides[layout_decoded].append(i + 1)

        shapes_flat = list(iter_shapes(slide.shapes))
        shape_count_per_slide.append(len(shapes_flat))

        grid_plain = [0] * 100
        grid_typed_set = set()
        type_hist = Counter()

        first3_names = []
        first3_texts = []
        for idx, (shp, depth) in enumerate(shapes_flat):
            st = shp.shape_type
            st_str = str(st).split('.')[-1] if st else 'NONE'
            shape_type_ctr[st_str] += 1
            type_hist[st_str] += 1

            # AutoShape type
            if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    ast = shp.auto_shape_type
                    auto_shape_ctr[str(ast).split('.')[-1] if ast else 'NONE'] += 1
                except Exception:
                    auto_shape_ctr['ERR'] += 1

            # Placeholder
            if shp.is_placeholder:
                try:
                    ph_type = shp.placeholder_format.type
                    placeholder_type_ctr[str(ph_type).split('.')[-1] if ph_type else 'NONE'] += 1
                except Exception:
                    placeholder_type_ctr['ERR'] += 1

            # shape name
            nm = (shp.name or '').strip()
            shape_name_ctr[nm] += 1
            if len(nm) >= 3:
                shape_name_prefix_ctr[nm[:3]] += 1

            if idx < 3:
                first3_names.append(nm)
                try:
                    txt = shp.text_frame.text.strip() if shp.has_text_frame else ''
                except Exception:
                    txt = ''
                first3_texts.append(txt[:40])

            # geometry
            try:
                left = shp.left or 0
                top = shp.top or 0
                w = shp.width or 0
                h = shp.height or 0
                if len(geometry_samples[st_str]) < 5:
                    geometry_samples[st_str].append((left, top, w, h))
                # grid mapping (center point)
                cx = (left + w / 2) / slide_w if slide_w else 0
                cy = (top + h / 2) / slide_h if slide_h else 0
                cell = grid_cell(cx, cy)
                grid_plain[cell] = 1
                grid_typed_set.add((cell, st_str))
            except Exception:
                pass

            # Code tag detection (A01 in text)
            try:
                if shp.has_text_frame:
                    txt = shp.text_frame.text.strip()
                    if CODE_RE.match(txt) and len(code_tag_positions) < 40:
                        left = shp.left or 0
                        top = shp.top or 0
                        code_tag_positions.append({
                            'slide': i + 1,
                            'text': txt,
                            'shape_name': nm,
                            'left_norm': round(left / slide_w, 3) if slide_w else 0,
                            'top_norm': round(top / slide_h, 3) if slide_h else 0,
                        })
            except Exception:
                pass

        first3_names_samples.append((i + 1, first3_names))
        first3_texts_samples.append((i + 1, first3_texts))
        grid_sigs_plain.append(tuple(grid_plain))
        grid_sigs_typed.append(tuple(sorted(grid_typed_set)))
        structure_sigs.append(tuple(sorted(type_hist.items())))

    total = len(prs.slides)
    # ============================
    print("\n" + "=" * 70)
    print(f"TOTAL SLIDES: {total}")
    print("=" * 70)

    print("\n[1] Shape type 분포 (leaf 기준, 전체 shape):")
    for st, cnt in shape_type_ctr.most_common(20):
        print(f"  {st:40s} {cnt:7d}")

    print("\n[1b] AutoShape subtype 분포 (상위 15):")
    for ast, cnt in auto_shape_ctr.most_common(15):
        print(f"  {ast:40s} {cnt:7d}")

    print("\n[1c] Placeholder type 분포:")
    for pt, cnt in placeholder_type_ctr.most_common(20):
        print(f"  {pt:40s} {cnt:7d}")

    print("\n[1d] Shape 개수 분포 (슬라이드당):")
    sc = sorted(shape_count_per_slide)
    print(f"  min={sc[0]}, max={sc[-1]}, avg={sum(sc)/total:.1f}")
    print(f"  median={sc[total//2]}, p25={sc[total//4]}, p75={sc[3*total//4]}, p90={sc[int(0.9*total)]}")

    print("\n[1e] Shape name 패턴 - 상위 프리픽스 (앞 3글자):")
    for p, cnt in shape_name_prefix_ctr.most_common(20):
        print(f"  {p!r:15s} {cnt:7d}")

    print("\n[1f] Shape name TOP 20 (전체):")
    for nm, cnt in shape_name_ctr.most_common(20):
        print(f"  {nm[:50]!r:55s} {cnt:6d}")

    # 섹션 코드 (A01~Z99) 직접 검색
    code_like_names = [(nm, c) for nm, c in shape_name_ctr.items() if CODE_RE.match(nm)]
    print(f"\n[1g] Shape name이 섹션 코드 형태 (A01~Z99): {len(code_like_names)} unique, 총 {sum(c for _, c in code_like_names)}회")
    for nm, c in sorted(code_like_names, key=lambda x: -x[1])[:20]:
        print(f"  {nm:10s} {c}")

    # Layout
    print(f"\n[2] Layout (raw) unique: {len(layout_name_ctr)}")
    print(f"[2] Layout (decoded) unique: {len(layout_decoded_ctr)}")
    print("\n상위 Layout (decoded):")
    for nm, cnt in layout_decoded_ctr.most_common(20):
        first_slides = layout_first_slides[nm][:3]
        print(f"  {nm!r:60s} {cnt:5d}  (예: {first_slides})")

    # A03 tag 검증
    print(f"\n[4] 섹션 코드 태그 ({len(code_tag_positions)}개 발견, 텍스트 내):")
    for row in code_tag_positions[:15]:
        print(f"  slide #{row['slide']}: text={row['text']!r}, shape_name={row['shape_name']!r},"
              f" pos=({row['left_norm']}, {row['top_norm']})")

    # title_text 기반 code
    # slide_meta.json 열어 확인
    meta = json.load(open('c:/Users/y2kbo/Apps/PPT/output/catalog/slide_meta.json', encoding='utf-8'))
    title_code_ctr = Counter()
    title_code_positions = []
    for row in meta:
        t = (row.get('title_text') or '').strip()
        if CODE_RE.match(t):
            title_code_ctr[t[0]] += 1  # 첫 글자 (A/B/C...)
            if len(title_code_positions) < 20:
                title_code_positions.append((row['slide_index'], t))
    print(f"\n[4b] slide_meta.json title_text 중 섹션 코드 패턴: {sum(title_code_ctr.values())}개")
    print(f"  섹션(첫 글자) 분포: {dict(title_code_ctr.most_common())}")
    print(f"  예시:")
    for idx, t in title_code_positions[:12]:
        print(f"    slide_index={idx}: {t!r}")

    # first3 text (code 여부)
    code_in_first3 = 0
    for slide_idx, texts in first3_texts_samples:
        for t in texts:
            if CODE_RE.match(t.strip()):
                code_in_first3 += 1
                break
    print(f"\n[4c] 슬라이드 첫 3개 shape 중 섹션 코드 텍스트 포함: {code_in_first3}/{total} ({code_in_first3/total*100:.1f}%)")

    # ============================
    # Unique layout counting
    print("\n" + "=" * 70)
    print("[3] UNIQUE LAYOUT 수 측정")
    print("=" * 70)

    unique_plain = len(set(grid_sigs_plain))
    unique_typed = len(set(grid_sigs_typed))
    unique_struct = len(set(structure_sigs))
    print(f"  grid occupancy (10x10 binary):        unique = {unique_plain:4d} / {total}  ({unique_plain/total*100:.1f}%)")
    print(f"  grid + shape_type (cell, type):       unique = {unique_typed:4d} / {total}  ({unique_typed/total*100:.1f}%)")
    print(f"  shape type 히스토그램 (structure_sig):  unique = {unique_struct:4d} / {total}  ({unique_struct/total*100:.1f}%)")

    # distribution: 가장 많은 slide가 같은 sig 몇 개?
    plain_ctr = Counter(grid_sigs_plain)
    typed_ctr = Counter(grid_sigs_typed)
    struct_ctr = Counter(structure_sigs)
    print(f"\n  [grid plain] 최대 동일 sig: {plain_ctr.most_common(1)[0][1]} slides")
    print(f"  [grid plain] top10 sig 빈도: {[c for _, c in plain_ctr.most_common(10)]}")
    print(f"\n  [grid typed] 최대 동일 sig: {typed_ctr.most_common(1)[0][1]} slides")
    print(f"  [grid typed] top10 sig 빈도: {[c for _, c in typed_ctr.most_common(10)]}")
    print(f"\n  [struct sig] 최대 동일 sig: {struct_ctr.most_common(1)[0][1]} slides")
    print(f"  [struct sig] top10 sig 빈도: {[c for _, c in struct_ctr.most_common(10)]}")

    # grid 크기 별 실험 (5x5, 8x8, 12x12)
    print("\n[3b] Grid 크기별 unique 수:")
    for g in (4, 5, 6, 8, 10, 12, 16):
        sigs = []
        for i, slide in enumerate(prs.slides):
            shapes_flat = list(iter_shapes(slide.shapes))
            occ = [0] * (g * g)
            for shp, _ in shapes_flat:
                try:
                    left = shp.left or 0
                    top = shp.top or 0
                    w = shp.width or 0
                    h = shp.height or 0
                    cx = (left + w / 2) / slide_w
                    cy = (top + h / 2) / slide_h
                    gx = min(g - 1, max(0, int(cx * g)))
                    gy = min(g - 1, max(0, int(cy * g)))
                    occ[gy * g + gx] = 1
                except Exception:
                    pass
            sigs.append(tuple(occ))
        u = len(set(sigs))
        top_cluster = Counter(sigs).most_common(1)[0][1]
        print(f"  {g}x{g}: unique = {u:4d} / {total}  ({u/total*100:5.1f}%), 최대 클러스터 = {top_cluster}")

    # geometry 샘플
    print("\n[1h] Geometry 샘플 (shape_type별 최초 5개, EMU 원본):")
    for st, samples in list(geometry_samples.items())[:10]:
        print(f"  {st}:")
        for left, top, w, h in samples:
            print(f"    left={left}, top={top}, w={w}, h={h}")

    # 결과 JSON 저장
    out_summary = {
        'total_slides': total,
        'slide_dim_emu': [slide_w, slide_h],
        'shape_type_dist': dict(shape_type_ctr.most_common(20)),
        'auto_shape_dist': dict(auto_shape_ctr.most_common(20)),
        'placeholder_type_dist': dict(placeholder_type_ctr.most_common(20)),
        'shape_count_stats': {
            'min': sc[0], 'max': sc[-1], 'avg': sum(sc)/total,
            'median': sc[total//2], 'p25': sc[total//4],
            'p75': sc[3*total//4], 'p90': sc[int(0.9*total)],
        },
        'shape_name_prefix_top': dict(shape_name_prefix_ctr.most_common(20)),
        'shape_name_code_like_count': len(code_like_names),
        'shape_name_code_like_total_occurrences': sum(c for _, c in code_like_names),
        'layout_unique_raw': len(layout_name_ctr),
        'layout_unique_decoded': len(layout_decoded_ctr),
        'layout_top_decoded': dict(layout_decoded_ctr.most_common(20)),
        'code_tags_in_text': len(code_tag_positions),
        'code_tags_samples': code_tag_positions[:20],
        'title_code_section_dist': dict(title_code_ctr.most_common()),
        'title_code_total': sum(title_code_ctr.values()),
        'first3_code_coverage': f"{code_in_first3}/{total}",
        'unique_grid_10_plain': unique_plain,
        'unique_grid_10_typed': unique_typed,
        'unique_struct_sig': unique_struct,
        'grid_plain_top10_cluster_sizes': [c for _, c in plain_ctr.most_common(10)],
        'grid_typed_top10_cluster_sizes': [c for _, c in typed_ctr.most_common(10)],
        'struct_top10_cluster_sizes': [c for _, c in struct_ctr.most_common(10)],
    }
    out_path = os.path.join(OUT_DIR, 'layout_analysis_summary.json')
    with open(out_path, 'w', encoding='utf-8') as f:
        json.dump(out_summary, f, indent=2, ensure_ascii=False, default=str)
    print(f"\n저장: {out_path}")


if __name__ == '__main__':
    analyze()
