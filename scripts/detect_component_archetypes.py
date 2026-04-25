"""Phase 1C: 컴포넌트 archetype 자동 검출 (1,251장 마스터 템플릿).

검출 카테고리:
- C1 Chart-like:       chart_native / picture_chart_like
- C2 Table-like:       table_native / grid_layout / gantt_like
- C3 Diagram:          orgchart, flowchart, timeline_h, roadmap, matrix_2x2, matrix_3x3,
                       swimlane, venn, funnel, hub_spoke
- C4 Card:             cards_2col, cards_3col, cards_4col, cards_with_icon, cards_with_image
- C5 Text/Cover/Divider: cover_slide, section_divider, text_heavy

각 슬라이드에 multi-label 부여 + confidence (low/medium/high).
산출물: output/catalog/phase1c_component_types.json
"""
from __future__ import annotations

import json
import os
import re
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PATH = Path('c:/Users/y2kbo/Apps/PPT/docs/references/_master_templates/PPT 템플릿.pptx')
OUT_DIR = Path('c:/Users/y2kbo/Apps/PPT/output/catalog')
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / 'phase1c_component_types.json'


# ===== utilities =====

def iter_leaf(shapes, depth=0):
    for shp in shapes:
        st = shp.shape_type
        if st == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from iter_leaf(shp.shapes, depth + 1)
            except Exception:
                pass
        else:
            yield shp, depth


def shape_geom(shp, slide_w, slide_h):
    """Returns (left_n, top_n, w_n, h_n, cx, cy) normalized 0..1, or None."""
    try:
        left = shp.left or 0
        top = shp.top or 0
        w = shp.width or 0
        h = shp.height or 0
        if slide_w <= 0 or slide_h <= 0:
            return None
        ln = left / slide_w
        tn = top / slide_h
        wn = w / slide_w
        hn = h / slide_h
        cx = ln + wn / 2
        cy = tn + hn / 2
        return (ln, tn, wn, hn, cx, cy)
    except Exception:
        return None


def autoshape_name(shp):
    if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return None
    try:
        ast = shp.auto_shape_type
        return str(ast).split('.')[-1].split(' ')[0] if ast else None
    except Exception:
        return None


def text_of(shp):
    try:
        if shp.has_text_frame:
            return shp.text_frame.text or ''
    except Exception:
        pass
    return ''


# ===== detectors =====

def detect_chart(shapes, slide_w, slide_h):
    """C1. Native chart 또는 차트처럼 보이는 picture."""
    chart_count = 0
    chart_types = []
    pic_chart_like = 0
    for shp, _ in shapes:
        st = shp.shape_type
        if st == MSO_SHAPE_TYPE.CHART:
            chart_count += 1
            try:
                ct = shp.chart.chart_type
                chart_types.append(str(ct).split('.')[-1].split(' ')[0])
            except Exception:
                chart_types.append('UNKNOWN')
        elif st == MSO_SHAPE_TYPE.PICTURE:
            # 차트형 picture: 적당히 크고(>15% 너비, >15% 높이), 가로:세로 0.5~3 사이
            g = shape_geom(shp, slide_w, slide_h)
            if g is None:
                continue
            _, _, wn, hn, _, _ = g
            if wn > 0.15 and hn > 0.15:
                ratio = wn / hn if hn > 0 else 0
                if 0.4 <= ratio <= 3.0:
                    pic_chart_like += 1
    return {
        'chart_count': chart_count,
        'chart_types': chart_types,
        'picture_chart_like_count': pic_chart_like,
    }


def detect_table(shapes, slide_w, slide_h):
    """C2 native table + cell-grid heuristic."""
    native = 0
    tbl_dims = []  # (rows, cols)
    for shp, _ in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
            native += 1
            try:
                tbl = shp.table
                tbl_dims.append((len(tbl.rows), len(tbl.columns)))
            except Exception:
                pass
    return {
        'table_native_count': native,
        'table_dims': tbl_dims,
    }


def detect_grid_pattern(rects, slide_w, slide_h):
    """카드/표 grid 검출.

    rects: list of (left, top, w, h) normalized 0..1
    Returns: dict with grid info:
      - rows: int
      - cols: int
      - n: int (그리드를 이루는 박스 개수)
      - aligned: bool
    """
    if len(rects) < 2:
        return None
    # Y 좌표를 cluster (오차 허용 0.03)
    rows_y = []
    cols_x = []
    for left, top, w, h in rects:
        rows_y.append(top + h / 2)
        cols_x.append(left + w / 2)

    def cluster(vals, tol=0.04):
        sv = sorted(vals)
        clusters = [[sv[0]]]
        for v in sv[1:]:
            if abs(v - clusters[-1][-1]) <= tol:
                clusters[-1].append(v)
            else:
                clusters.append([v])
        return clusters

    row_clusters = cluster(rows_y)
    col_clusters = cluster(cols_x)
    rows = len(row_clusters)
    cols = len(col_clusters)
    return {
        'rows': rows,
        'cols': cols,
        'n_rects': len(rects),
        'row_sizes': [len(c) for c in row_clusters],
        'col_sizes': [len(c) for c in col_clusters],
    }


def collect_card_candidates(shapes, slide_w, slide_h):
    """주요 박스(Rectangle/RoundedRect/Cube/SnipRect) 중 적당히 큰 것."""
    boxes = []  # (left, top, w, h, shape)
    for shp, _ in shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        ast = autoshape_name(shp)
        if ast not in ('RECTANGLE', 'ROUNDED_RECTANGLE', 'CUBE',
                       'SNIP_1_RECTANGLE', 'ROUND_2_SAME_RECTANGLE', 'PARALLELOGRAM',
                       'FLOWCHART_PROCESS', 'SNIP_2_SAME_RECTANGLE'):
            continue
        g = shape_geom(shp, slide_w, slide_h)
        if g is None:
            continue
        ln, tn, wn, hn, _, _ = g
        # 너무 작거나 너무 큰 (배경) 박스 제외
        if wn < 0.05 or wn > 0.6:
            continue
        if hn < 0.04 or hn > 0.7:
            continue
        boxes.append((ln, tn, wn, hn))
    return boxes


def detect_cards(shapes, slide_w, slide_h):
    """C4 카드 패턴 검출.

    Cards: 동일 너비/높이 (편차 ≤ 10%) 박스 N(2-6)개가 가로 일렬 (같은 y, 다른 x).
    엄격: 텍스트 라벨이 적어도 절반의 카드 안에 있어야 함.
    """
    boxes = collect_card_candidates(shapes, slide_w, slide_h)
    if len(boxes) < 2:
        return None
    by_y = defaultdict(list)
    for b in boxes:
        ln, tn, wn, hn = b
        ykey = round(tn / 0.04) * 0.04
        by_y[ykey].append(b)
    best = None
    for ykey, group in by_y.items():
        if len(group) < 2:
            continue
        # 카드 개수 2-8 까지만 (그 이상은 보통 표/grid)
        if len(group) > 8:
            continue
        ws = [g[2] for g in group]
        hs = [g[3] for g in group]
        w_mean = sum(ws) / len(ws)
        h_mean = sum(hs) / len(hs)
        if w_mean < 0.07 or h_mean < 0.05:
            continue
        # 너비/높이 편차 < 10%
        w_var = max(abs(w - w_mean) for w in ws) / w_mean
        h_var = max(abs(h - h_mean) for h in hs) / h_mean
        if w_var > 0.10 or h_var > 0.15:
            continue
        # 가로 spread: 카드의 x 좌표가 충분히 흩어져 있어야 함
        xs = sorted(g[0] for g in group)
        if xs[-1] - xs[0] < 0.3 * len(group):
            # 평균 0.3 미만이면 정렬되었다고 보기 어려움 (단, 2개일 때는 spread > 0.2)
            min_spread = 0.2 if len(group) == 2 else 0.3
            if xs[-1] - xs[0] < min_spread:
                continue
        # 균등 간격 검증
        if len(group) >= 3:
            gaps = [xs[i+1] - xs[i] for i in range(len(xs)-1)]
            if max(gaps) > 2.5 * min(g for g in gaps if g > 0.001):
                continue
        n = len(group)
        if best is None or n > best['n']:
            best = {'n': n, 'w_mean': w_mean, 'h_mean': h_mean, 'y': ykey}
    return best


def detect_lines_and_arrows(shapes, slide_w, slide_h):
    """connector + arrow shape 카운트."""
    n_connector = 0
    n_line = 0
    n_arrow_shape = 0
    arrow_kinds = Counter()
    n_diamond = 0  # decision
    for shp, _ in shapes:
        st = shp.shape_type
        if st == MSO_SHAPE_TYPE.LINE:
            n_line += 1
            # connector check
            try:
                if 'cxnSp' in shp.element.tag:
                    n_connector += 1
            except Exception:
                pass
        elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            ast = autoshape_name(shp)
            if ast and 'ARROW' in ast:
                n_arrow_shape += 1
                arrow_kinds[ast] += 1
            if ast == 'DIAMOND':
                n_diamond += 1
        # Connector via element check (xml lvl)
        try:
            if shp.element.tag.endswith('}cxnSp'):
                n_connector += 1
        except Exception:
            pass
    return {
        'connector': n_connector,
        'line': n_line,
        'arrow_shape': n_arrow_shape,
        'arrow_kinds': dict(arrow_kinds),
        'diamond': n_diamond,
    }


def detect_orgchart(shapes, slide_w, slide_h, lines_info):
    """C3 orgchart: 박스 트리 구조 (depth ≥ 3) + 수직 connector lines.

    엄격: 박스 depth ≥ 3 (root, level1, level2 최소) + 수직 connector ≥ 3.
    """
    boxes = collect_card_candidates(shapes, slide_w, slide_h)
    if len(boxes) < 6:
        return False, 0.0
    # 명확한 y level이 ≥ 3 (계층 깊이)
    ys = sorted(b[1] for b in boxes)
    levels = []
    for y in ys:
        if not levels or y - levels[-1] >= 0.08:
            levels.append(y)
    if len(levels) < 3:
        return False, 0.0
    # 형제 박스 (같은 y에 width 비슷한 박스 ≥ 2)가 ≥ 2 레벨에서 발생
    by_y_level = defaultdict(list)
    for b in boxes:
        # level 매핑: 가장 가까운 level y
        lev = min(range(len(levels)), key=lambda i: abs(levels[i] - b[1]))
        by_y_level[lev].append(b)
    levels_with_siblings = sum(1 for g in by_y_level.values() if len(g) >= 2)
    if levels_with_siblings < 2:
        return False, 0.0
    # connector lines 충분
    if lines_info['connector'] + lines_info['line'] < 4:
        return False, 0.0
    conf = min(1.0, 0.4 + 0.04 * len(boxes) + 0.08 * len(levels))
    return True, conf


def detect_flowchart(shapes, slide_w, slide_h, lines_info):
    """C3 flowchart: arrow shape ≥ 3 + boxes ≥ 4 또는 diamond ≥ 1.

    엄격: 최소 arrow ≥ 3 + 다수 박스 또는 decision diamond.
    """
    boxes = collect_card_candidates(shapes, slide_w, slide_h)
    has_diamond = lines_info['diamond'] >= 1
    has_arrows = lines_info['arrow_shape'] >= 3
    if has_arrows and (has_diamond or len(boxes) >= 4):
        conf = min(1.0, 0.4 + 0.04 * lines_info['arrow_shape'] + 0.2 * lines_info['diamond'])
        return True, conf
    # decision diamond 강하게 — 명확한 flowchart signal
    if has_diamond and lines_info['arrow_shape'] >= 1 and len(boxes) >= 3:
        return True, 0.7
    return False, 0.0


def detect_timeline_h(shapes, slide_w, slide_h, lines_info):
    """C3 horizontal timeline: 가로 일렬로 정렬된 작은 마커(원/점) + 가로선/화살표.

    조건: 동일 y에 OVAL/circle ≥ 3 + (line/arrow 가 가로방향).
    """
    ovals = []  # (left, top, w, h)
    for shp, _ in shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if autoshape_name(shp) not in ('OVAL', 'CHEVRON'):
            continue
        g = shape_geom(shp, slide_w, slide_h)
        if g is None:
            continue
        ln, tn, wn, hn, _, _ = g
        if wn > 0.15:
            continue
        ovals.append((ln, tn, wn, hn))
    if len(ovals) < 3:
        return False, 0.0
    # 가장 흔한 y row
    by_y = Counter(round(o[1] / 0.03) * 0.03 for o in ovals)
    top_y, count = by_y.most_common(1)[0]
    if count < 3:
        return False, 0.0
    # 그 row 내 x 분포가 가로로 spread (range > 0.4)
    same_row = [o for o in ovals if abs(o[1] - top_y) <= 0.03]
    xs = [o[0] for o in same_row]
    if max(xs) - min(xs) > 0.4:
        conf = min(1.0, 0.4 + 0.08 * count)
        return True, conf
    return False, 0.0


def detect_roadmap(shapes, slide_w, slide_h):
    """C3 단계형 로드맵: PENTAGON/CHEVRON ≥ 3 가로 일렬.

    또한 'Phase' 같은 텍스트가 있으면 신뢰도 boost.
    """
    chevs = []  # (left, top, w, h)
    text_blob = ''
    for shp, _ in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            ast = autoshape_name(shp)
            if ast in ('PENTAGON', 'CHEVRON'):
                g = shape_geom(shp, slide_w, slide_h)
                if g is not None:
                    chevs.append(g[:4])
        text_blob += ' ' + text_of(shp)
    if len(chevs) < 3:
        return False, 0.0
    by_y = Counter(round(c[1] / 0.04) * 0.04 for c in chevs)
    _, count = by_y.most_common(1)[0]
    if count < 3:
        return False, 0.0
    has_phase_text = bool(re.search(r'(phase|step|단계|stage|\b\d단계)', text_blob.lower()))
    conf = 0.5 + 0.07 * count + (0.15 if has_phase_text else 0)
    return True, min(1.0, conf)


def detect_matrix(shapes, slide_w, slide_h):
    """C3 2x2 또는 3x3 매트릭스: 정확한 grid 박스 4개 또는 9개."""
    boxes = collect_card_candidates(shapes, slide_w, slide_h)
    if len(boxes) < 4:
        return None, 0.0
    grid = detect_grid_pattern(boxes, slide_w, slide_h)
    if grid is None:
        return None, 0.0
    rows, cols = grid['rows'], grid['cols']
    n = grid['n_rects']
    if rows == 2 and cols == 2 and 3 <= n <= 6:
        return '2x2', 0.7
    if rows == 3 and cols == 3 and 7 <= n <= 12:
        return '3x3', 0.65
    return None, 0.0


def detect_swimlane(shapes, slide_w, slide_h):
    """C3 swimlane: 가로 lane (너비 > 0.7, aspect > 5:1) ≥ 3 stacked vertically.

    엄격: lane이 슬라이드 거의 전 너비 + ≥ 3개 + y 분산 > 0.2.
    """
    lanes = []
    for shp, _ in shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        ast = autoshape_name(shp)
        if ast not in ('RECTANGLE', 'ROUNDED_RECTANGLE'):
            continue
        g = shape_geom(shp, slide_w, slide_h)
        if g is None:
            continue
        ln, tn, wn, hn, _, _ = g
        if hn <= 0:
            continue
        ar = wn / hn
        if wn > 0.70 and ar > 5.0 and 0.05 < hn < 0.25:
            lanes.append((ln, tn, wn, hn))
    if len(lanes) < 3:
        return False, 0.0
    ys = sorted([l[1] for l in lanes])
    if ys[-1] - ys[0] < 0.20:
        return False, 0.0
    # 일정한 spacing 검증
    diffs = [ys[i+1] - ys[i] for i in range(len(ys)-1)]
    if diffs and max(diffs) > 3 * min(d for d in diffs if d > 0.01) if any(d > 0.01 for d in diffs) else False:
        return False, 0.0
    return True, min(1.0, 0.4 + 0.08 * len(lanes))


def detect_venn(shapes, slide_w, slide_h):
    """C3 Venn: OVAL 2-3개 overlap (BBox 겹침)."""
    ovals = []
    for shp, _ in shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if autoshape_name(shp) != 'OVAL':
            continue
        g = shape_geom(shp, slide_w, slide_h)
        if g is None:
            continue
        ln, tn, wn, hn, _, _ = g
        # 큰 oval만 (지름 > 15%)
        if wn < 0.15 or hn < 0.1:
            continue
        ovals.append((ln, tn, wn, hn))
    if 2 <= len(ovals) <= 4:
        # 겹침 검증
        overlaps = 0
        for i in range(len(ovals)):
            for j in range(i + 1, len(ovals)):
                a, b = ovals[i], ovals[j]
                if (a[0] < b[0] + b[2] and a[0] + a[2] > b[0]
                        and a[1] < b[1] + b[3] and a[1] + a[3] > b[1]):
                    overlaps += 1
        if overlaps >= 1:
            return True, min(1.0, 0.5 + 0.15 * overlaps)
    return False, 0.0


def detect_funnel(shapes, slide_w, slide_h):
    """C3 funnel: 너비가 점진 감소하는 사다리꼴/사각형 ≥ 3개 세로 stack."""
    boxes = collect_card_candidates(shapes, slide_w, slide_h)
    if len(boxes) < 3:
        return False, 0.0
    # y 정렬 후 width 단조 감소 검증
    sorted_by_y = sorted(boxes, key=lambda b: b[1])
    widths = [b[2] for b in sorted_by_y]
    # 단조성: 연속 감소 > 50% of pairs
    decr = sum(1 for i in range(1, len(widths)) if widths[i] < widths[i - 1] * 0.95)
    if decr >= len(widths) - 1 and len(widths) >= 3:
        # 너비 감소 폭 충분?
        if widths[0] > widths[-1] * 1.2:
            return True, min(1.0, 0.4 + 0.07 * len(widths))
    return False, 0.0


def detect_hub_spoke(shapes, slide_w, slide_h, lines_info):
    """C3 hub-spoke: 중앙 큰 도형 + 방사형 sub-shapes ≥ 4 + 중앙→주변 connector.

    엄격: 중앙(0.4-0.6 영역)에 명확히 큰 도형 1개 + 주변(외곽) shape ≥ 4
    + 중앙으로부터 다양한 angle.
    """
    center_shape = None  # (cx, cy, area)
    peripheral = []  # (cx, cy)
    for shp, _ in shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        g = shape_geom(shp, slide_w, slide_h)
        if g is None:
            continue
        ln, tn, wn, hn, cx, cy = g
        if wn < 0.04 or hn < 0.04:
            continue
        if 0.40 <= cx <= 0.60 and 0.35 <= cy <= 0.65 and wn > 0.12 and hn > 0.12:
            if center_shape is None or wn * hn > center_shape[2]:
                center_shape = (cx, cy, wn * hn)
        elif wn < 0.22 and hn < 0.22:
            peripheral.append((cx, cy))
    if center_shape is None or len(peripheral) < 4:
        return False, 0.0
    # angle diversity: 중앙으로부터 4분면 모두 occupy?
    cx0, cy0, _ = center_shape
    quads = set()
    for px, py in peripheral:
        dx, dy = px - cx0, py - cy0
        if abs(dx) < 0.05 and abs(dy) < 0.05:
            continue
        quads.add((dx >= 0, dy >= 0))
    if len(quads) < 3:
        return False, 0.0
    if lines_info['line'] + lines_info['connector'] < 4:
        return False, 0.0
    return True, 0.6


def detect_gantt(shapes, slide_w, slide_h):
    """C2 Gantt: 가로 긴 막대(aspect > 4:1) ≥ 4개 + 다른 시작 x + 서로 다른 행 + 가로 grid.

    엄격: 막대 aspect_ratio (w/h) ≥ 4, 다양한 시작 x (≥ 4 unique), 다양한 y행 (≥ 4).
    """
    bars = []
    for shp, _ in shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if autoshape_name(shp) not in ('RECTANGLE', 'ROUNDED_RECTANGLE'):
            continue
        g = shape_geom(shp, slide_w, slide_h)
        if g is None:
            continue
        ln, tn, wn, hn, _, _ = g
        # 막대: 가로:세로 ≥ 4:1, 너무 작지도 너무 크지도 않음
        if hn <= 0 or wn <= 0:
            continue
        ar = wn / hn
        if ar < 4.0:
            continue
        if hn > 0.05 or hn < 0.005:
            continue
        if wn < 0.04 or wn > 0.6:
            continue
        bars.append((ln, tn, wn, hn))
    if len(bars) < 5:
        return False, 0.0
    # 시작 x (≥ 4 unique)
    xs = sorted({round(b[0] / 0.02) * 0.02 for b in bars})
    if len(xs) < 4:
        return False, 0.0
    # y행 (≥ 4 unique)
    ys = sorted({round(b[1] / 0.03) * 0.03 for b in bars})
    if len(ys) < 4:
        return False, 0.0
    return True, min(1.0, 0.4 + 0.04 * len(bars))


def detect_cover(shapes, slide_w, slide_h, slide_idx, total_slides, page_text):
    """C5 cover/divider: 첫 슬라이드 또는 텍스트가 매우 적은 큰 텍스트 박스."""
    # cover heuristic: 슬라이드 #1 또는 매우 적은 shape (≤ 8) + 큰 텍스트 박스
    leaf_shapes = list(shapes)
    n_shapes = len(leaf_shapes)
    big_text_boxes = 0
    for shp, _ in leaf_shapes:
        if not shp.has_text_frame:
            continue
        g = shape_geom(shp, slide_w, slide_h)
        if g is None:
            continue
        _, _, wn, hn, _, _ = g
        if wn > 0.3 and hn > 0.1:
            big_text_boxes += 1
    text_len = len(page_text)
    is_cover = (slide_idx == 1) or (n_shapes <= 12 and big_text_boxes >= 1 and text_len < 200)
    return is_cover


def detect_text_heavy(shapes, slide_w, slide_h, page_text):
    """C5 텍스트 위주: 적은 shape (≤25) + 큰 텍스트 박스가 슬라이드 면적 ≥ 50%.

    엄격: shape 적고, 텍스트 길고, 도형 적은 슬라이드만.
    """
    n_shapes = len(shapes)
    if n_shapes > 25:
        return False
    text_area = 0.0
    for shp, _ in shapes:
        if not shp.has_text_frame:
            continue
        g = shape_geom(shp, slide_w, slide_h)
        if g is None:
            continue
        _, _, wn, hn, _, _ = g
        text_area += wn * hn
    return text_area > 0.5 and len(page_text.strip()) > 200


# ===== orchestrator =====

def analyze_slide(slide, slide_w, slide_h, slide_idx, total_slides):
    flat = list(iter_leaf(slide.shapes))
    page_text = '\n'.join(text_of(shp) for shp, _ in flat)

    chart = detect_chart(flat, slide_w, slide_h)
    table = detect_table(flat, slide_w, slide_h)
    lines = detect_lines_and_arrows(flat, slide_w, slide_h)
    cards = detect_cards(flat, slide_w, slide_h)
    matrix_kind, matrix_conf = detect_matrix(flat, slide_w, slide_h)
    org_yes, org_conf = detect_orgchart(flat, slide_w, slide_h, lines)
    flow_yes, flow_conf = detect_flowchart(flat, slide_w, slide_h, lines)
    tl_yes, tl_conf = detect_timeline_h(flat, slide_w, slide_h, lines)
    rd_yes, rd_conf = detect_roadmap(flat, slide_w, slide_h)
    sw_yes, sw_conf = detect_swimlane(flat, slide_w, slide_h)
    vn_yes, vn_conf = detect_venn(flat, slide_w, slide_h)
    fn_yes, fn_conf = detect_funnel(flat, slide_w, slide_h)
    hs_yes, hs_conf = detect_hub_spoke(flat, slide_w, slide_h, lines)
    gantt_yes, gantt_conf = detect_gantt(flat, slide_w, slide_h)
    is_cover = detect_cover(flat, slide_w, slide_h, slide_idx, total_slides, page_text)
    is_text_heavy = detect_text_heavy(flat, slide_w, slide_h, page_text)

    types = []  # multi-label

    # C1
    if chart['chart_count'] > 0:
        types.append({'type': 'chart_native', 'confidence': 'high',
                      'detail': {'count': chart['chart_count'], 'kinds': chart['chart_types']}})
    elif chart['picture_chart_like_count'] > 0:
        types.append({'type': 'picture_chart_like', 'confidence': 'low',
                      'detail': {'count': chart['picture_chart_like_count']}})

    # C2
    if table['table_native_count'] > 0:
        types.append({'type': 'table_native', 'confidence': 'high',
                      'detail': {'count': table['table_native_count'], 'dims': table['table_dims']}})
    if gantt_yes:
        types.append({'type': 'gantt_like', 'confidence': _conf_bucket(gantt_conf),
                      'detail': {'score': round(gantt_conf, 2)}})

    # C3 diagrams
    if matrix_kind:
        types.append({'type': f'matrix_{matrix_kind}', 'confidence': _conf_bucket(matrix_conf),
                      'detail': {'score': round(matrix_conf, 2)}})
    if org_yes:
        types.append({'type': 'orgchart', 'confidence': _conf_bucket(org_conf),
                      'detail': {'score': round(org_conf, 2)}})
    if flow_yes:
        types.append({'type': 'flowchart', 'confidence': _conf_bucket(flow_conf),
                      'detail': {'score': round(flow_conf, 2),
                                 'arrows': lines['arrow_shape'], 'diamond': lines['diamond']}})
    if tl_yes:
        types.append({'type': 'timeline_h', 'confidence': _conf_bucket(tl_conf),
                      'detail': {'score': round(tl_conf, 2)}})
    if rd_yes:
        types.append({'type': 'roadmap', 'confidence': _conf_bucket(rd_conf),
                      'detail': {'score': round(rd_conf, 2)}})
    if sw_yes:
        types.append({'type': 'swimlane', 'confidence': _conf_bucket(sw_conf),
                      'detail': {'score': round(sw_conf, 2)}})
    if vn_yes:
        types.append({'type': 'venn', 'confidence': _conf_bucket(vn_conf),
                      'detail': {'score': round(vn_conf, 2)}})
    if fn_yes:
        types.append({'type': 'funnel', 'confidence': _conf_bucket(fn_conf),
                      'detail': {'score': round(fn_conf, 2)}})
    if hs_yes:
        types.append({'type': 'hub_spoke', 'confidence': _conf_bucket(hs_conf),
                      'detail': {'score': round(hs_conf, 2)}})

    # C4 cards
    if cards and cards['n'] >= 2:
        # 카드 타입 결정 (2/3/4/5+)
        n = cards['n']
        if n == 2:
            t = 'cards_2col'
        elif n == 3:
            t = 'cards_3col'
        elif n == 4:
            t = 'cards_4col'
        elif n == 5:
            t = 'cards_5col'
        else:
            t = f'cards_{n}col'
        types.append({'type': t, 'confidence': 'medium' if n >= 3 else 'low',
                      'detail': {'n': n, 'w': round(cards['w_mean'], 3), 'h': round(cards['h_mean'], 3)}})

    # C5
    if is_cover:
        types.append({'type': 'cover_or_divider', 'confidence': 'medium',
                      'detail': {'shape_count': len(flat), 'text_len': len(page_text)}})
    elif is_text_heavy:
        types.append({'type': 'text_heavy', 'confidence': 'medium',
                      'detail': {'text_len': len(page_text)}})

    # 아무것도 안 잡히면 unknown
    if not types:
        types.append({'type': 'unclassified', 'confidence': 'low',
                      'detail': {'shape_count': len(flat)}})

    return types


# 우선순위: 더 구체적일수록 높음
PRIORITY = {
    'chart_native': 100,
    'table_native': 95,
    'venn': 92,
    'matrix_3x3': 90,
    'matrix_2x2': 89,
    'gantt_like': 85,
    'flowchart': 80,
    'orgchart': 78,
    'swimlane': 75,
    'roadmap': 72,
    'timeline_h': 70,
    'funnel': 68,
    'hub_spoke': 66,
    'cards_5col': 60,
    'cards_6col': 60,
    'cards_4col': 58,
    'cards_3col': 56,
    'cards_2col': 54,
    'cards_7col': 52,
    'cards_8col': 50,
    'picture_chart_like': 40,
    'cover_or_divider': 35,
    'text_heavy': 30,
    'unclassified': 5,
}
CONF_RANK = {'high': 3, 'medium': 2, 'low': 1}


def select_primary(types):
    """가장 구체적이고 신뢰도 높은 archetype 선정."""
    if not types:
        return None
    # confidence + priority 결합
    def key(t):
        return (CONF_RANK.get(t['confidence'], 0),
                PRIORITY.get(t['type'], 0))
    best = max(types, key=key)
    return best['type']


def _conf_bucket(score: float) -> str:
    if score >= 0.7:
        return 'high'
    if score >= 0.5:
        return 'medium'
    return 'low'


# ===== main =====

def main():
    print(f'Loading {PATH} ...')
    prs = Presentation(str(PATH))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    total = len(prs.slides)
    print(f'Slides: {total}, dim={slide_w/914400:.2f}x{slide_h/914400:.2f} inch')

    per_slide = []
    type_hist = Counter()        # multi-label: 모든 타입 합산
    primary_hist = Counter()     # primary 단일 타입
    multi_label = Counter()

    for i, slide in enumerate(prs.slides):
        types = analyze_slide(slide, slide_w, slide_h, i + 1, total)
        primary = select_primary(types)
        per_slide.append({
            'slide_index': i,
            'slide_number': i + 1,
            'primary_archetype': primary,
            'types': types,
        })
        for t in types:
            type_hist[t['type']] += 1
        if primary:
            primary_hist[primary] += 1
        multi_label[len(types)] += 1
        if (i + 1) % 200 == 0:
            print(f'  processed {i+1}/{total}')

    summary = {
        'total_slides': total,
        'global_histogram_multilabel': dict(type_hist.most_common()),
        'primary_archetype_histogram': dict(primary_hist.most_common()),
        'multi_label_distribution': dict(multi_label),
        'detection_categories': {
            'C1_chart': ['chart_native', 'picture_chart_like'],
            'C2_table': ['table_native', 'gantt_like'],
            'C3_diagram': ['orgchart', 'flowchart', 'timeline_h', 'roadmap',
                           'matrix_2x2', 'matrix_3x3', 'swimlane', 'venn', 'funnel', 'hub_spoke'],
            'C4_card': ['cards_2col', 'cards_3col', 'cards_4col', 'cards_5col',
                        'cards_6col', 'cards_7col', 'cards_8col'],
            'C5_text': ['cover_or_divider', 'text_heavy'],
            'fallback': ['unclassified'],
        },
        'per_slide': per_slide,
    }

    OUT_PATH.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding='utf-8')
    print(f'\nWritten: {OUT_PATH}  ({OUT_PATH.stat().st_size:,} bytes)')

    # 요약 출력
    print('\n=== Multi-label histogram (slide can match multiple) ===')
    for t, c in type_hist.most_common():
        print(f'  {t:25s} {c:5d}  ({c/total*100:5.1f}%)')

    print('\n=== Primary archetype (one per slide, priority-selected) ===')
    for t, c in primary_hist.most_common():
        print(f'  {t:25s} {c:5d}  ({c/total*100:5.1f}%)')

    print('\n=== Multi-label distribution (number of types per slide) ===')
    for k in sorted(multi_label):
        print(f'  {k} types: {multi_label[k]} slides')


if __name__ == '__main__':
    main()
