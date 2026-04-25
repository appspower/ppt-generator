"""Phase A3 Step 2 — paragraph 단위 메타데이터 추출.

마스터 PPT 1,251장의 모든 paragraph에 대해
(slide_index, flat_idx, paragraph_id, text, font, bbox, shape_type, ...)
를 추출하여 output/catalog/paragraphs.json에 저장한다.

Step 2의 라벨러(label_paragraphs.py)가 이 데이터를 입력으로 사용.

CLI: python scripts/extract_paragraphs.py
"""
from __future__ import annotations

import json
import sys
import time
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.group import GroupShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.util import Emu

from ppt_builder.template import edit_ops


TEMPLATE_PATH = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
SLOT_SCHEMAS_PATH = ROOT / "output" / "catalog" / "slot_schemas.json"
OUTPUT_PATH = ROOT / "output" / "catalog" / "paragraphs.json"


def _shape_kind(shape) -> str:
    """단순화한 shape 종류 라벨."""
    try:
        if isinstance(shape, GroupShape):
            return "GROUP"
        if isinstance(shape, Picture):
            return "PICTURE"
        if isinstance(shape, GraphicFrame):
            if shape.has_table:
                return "TABLE"
            if shape.has_chart:
                return "CHART"
            return "GRAPHICFRAME"
        if shape.is_placeholder:
            return "PLACEHOLDER"
        if isinstance(shape, Shape):
            try:
                ast = shape.auto_shape_type
                if ast is not None:
                    return f"AUTOSHAPE:{ast.name}"
            except Exception:
                pass
            return "SHAPE"
        return type(shape).__name__.upper()
    except Exception:
        return "UNKNOWN"


def _placeholder_type(shape) -> str | None:
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            return ph.type.name if ph and ph.type else None
    except Exception:
        return None
    return None


def _bbox_of(shape, slide_w: int, slide_h: int) -> dict:
    """Return absolute (EMU) and percentage bbox."""
    try:
        left = shape.left or 0
        top = shape.top or 0
        width = shape.width or 0
        height = shape.height or 0
    except Exception:
        return {"left_emu": 0, "top_emu": 0, "width_emu": 0, "height_emu": 0,
                "left_pct": 0, "top_pct": 0, "width_pct": 0, "height_pct": 0}
    return {
        "left_emu": int(left),
        "top_emu": int(top),
        "width_emu": int(width),
        "height_emu": int(height),
        "left_pct": round(left / slide_w, 4) if slide_w else 0,
        "top_pct": round(top / slide_h, 4) if slide_h else 0,
        "width_pct": round(width / slide_w, 4) if slide_w else 0,
        "height_pct": round(height / slide_h, 4) if slide_h else 0,
    }


def _para_font(para) -> dict:
    """대표 폰트 정보 추출 (첫 run 기준 + max font_size)."""
    runs = list(para.runs)
    if not runs:
        # paragraph-level font
        try:
            sz = para.font.size
            return {
                "font_size_pt": sz.pt if sz else None,
                "bold": para.font.bold,
                "italic": para.font.italic,
                "n_runs": 0,
            }
        except Exception:
            return {"font_size_pt": None, "bold": None, "italic": None, "n_runs": 0}

    sizes = []
    for r in runs:
        try:
            if r.font.size:
                sizes.append(r.font.size.pt)
        except Exception:
            pass
    first = runs[0]
    try:
        bold = first.font.bold
    except Exception:
        bold = None
    try:
        italic = first.font.italic
    except Exception:
        italic = None
    try:
        color = None
        if first.font.color and first.font.color.type is not None:
            try:
                rgb = first.font.color.rgb
                if rgb is not None:
                    color = f"#{str(rgb)}"
            except Exception:
                color = None
    except Exception:
        color = None

    return {
        "font_size_pt": max(sizes) if sizes else None,
        "bold": bold,
        "italic": italic,
        "color": color,
        "n_runs": len(runs),
    }


def _alignment(para) -> str | None:
    try:
        if para.alignment is not None:
            return para.alignment.name
    except Exception:
        return None
    return None


def _level(para) -> int:
    try:
        return int(para.level or 0)
    except Exception:
        return 0


def _ancestor_group_signature(stack: list[int]) -> str | None:
    """Group 중첩 경로를 문자열로. None이면 직접 슬라이드 자식."""
    if not stack:
        return None
    return ">".join(str(i) for i in stack)


def _walk_paragraphs(shapes, group_stack: list[int], idx_state: list[int],
                     slide_w: int, slide_h: int) -> list[dict]:
    out = []
    for sh in shapes:
        if isinstance(sh, GroupShape):
            group_stack.append(idx_state[1])  # use a separate group counter
            idx_state[1] += 1
            out.extend(_walk_paragraphs(
                sh.shapes, group_stack, idx_state, slide_w, slide_h
            ))
            group_stack.pop()
            continue

        flat_idx = idx_state[0]
        idx_state[0] += 1

        kind = _shape_kind(sh)
        bbox = _bbox_of(sh, slide_w, slide_h)
        ph_type = _placeholder_type(sh)

        # TABLE: cells iteration
        if kind == "TABLE":
            try:
                tbl = sh.table
                n_rows = len(tbl.rows)
                n_cols = len(tbl.columns)
                for r_i, row in enumerate(tbl.rows):
                    for c_i, cell in enumerate(row.cells):
                        if cell.text_frame is None:
                            continue
                        for p_i, para in enumerate(cell.text_frame.paragraphs):
                            text = para.text or ""
                            font = _para_font(para)
                            out.append({
                                "flat_idx": flat_idx,
                                "paragraph_id": p_i,
                                "table_row": r_i,
                                "table_col": c_i,
                                "table_n_rows": n_rows,
                                "table_n_cols": n_cols,
                                "text": text,
                                "text_len": len(text),
                                "shape_kind": kind,
                                "placeholder_type": ph_type,
                                "bbox": bbox,
                                "group_path": _ancestor_group_signature(group_stack),
                                "font": font,
                                "alignment": _alignment(para),
                                "level": _level(para),
                            })
            except Exception:
                pass
            continue

        # 일반 text frame
        if not getattr(sh, "has_text_frame", False):
            continue
        try:
            tf = sh.text_frame
        except Exception:
            continue

        n_para = len(tf.paragraphs)
        for p_i, para in enumerate(tf.paragraphs):
            text = para.text or ""
            font = _para_font(para)
            out.append({
                "flat_idx": flat_idx,
                "paragraph_id": p_i,
                "n_paragraphs_in_shape": n_para,
                "text": text,
                "text_len": len(text),
                "shape_kind": kind,
                "placeholder_type": ph_type,
                "bbox": bbox,
                "group_path": _ancestor_group_signature(group_stack),
                "font": font,
                "alignment": _alignment(para),
                "level": _level(para),
            })

    return out


def main():
    print(f"[open] {TEMPLATE_PATH.name}", flush=True)
    prs = Presentation(str(TEMPLATE_PATH))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    print(f"  slide size: {slide_w} x {slide_h} EMU "
          f"({slide_w / 914400:.1f}\" x {slide_h / 914400:.1f}\")")

    # 기존 slot_schemas의 max_chars를 paragraph에 join할 수 있게 lookup
    with open(SLOT_SCHEMAS_PATH, encoding="utf-8") as f:
        schemas = json.load(f)
    max_chars_lookup: dict[tuple[int, int], int] = {}
    for s in schemas:
        max_chars_lookup[(s["slide_index"], s["flat_idx"])] = s.get("max_chars")

    n_slides = len(prs.slides)
    print(f"  slides: {n_slides}")

    all_paragraphs = []
    t0 = time.time()
    for s_i, slide in enumerate(prs.slides):
        # idx_state[0]: flat_idx counter, [1]: group counter
        idx_state = [0, 0]
        group_stack: list[int] = []
        paras = _walk_paragraphs(
            slide.shapes, group_stack, idx_state, slide_w, slide_h
        )
        # slide_index 부여 + max_chars join
        for p in paras:
            p["slide_index"] = s_i
            p["max_chars"] = max_chars_lookup.get((s_i, p["flat_idx"]))
        all_paragraphs.extend(paras)
        if (s_i + 1) % 100 == 0:
            print(f"  ... {s_i + 1}/{n_slides} ({time.time() - t0:.1f}s, "
                  f"{len(all_paragraphs)} paras)", flush=True)

    print(f"[done] extracted {len(all_paragraphs)} paragraphs in "
          f"{time.time() - t0:.1f}s")

    # 통계
    n_text = sum(1 for p in all_paragraphs if p["text"])
    n_empty = sum(1 for p in all_paragraphs if not p["text"])
    avg_per_slide = len(all_paragraphs) / n_slides
    print(f"  with text: {n_text}  empty: {n_empty}  avg/slide: {avg_per_slide:.1f}")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
        json.dump({
            "summary": {
                "n_slides": n_slides,
                "n_paragraphs": len(all_paragraphs),
                "n_with_text": n_text,
                "n_empty": n_empty,
                "avg_per_slide": round(avg_per_slide, 1),
                "slide_w_emu": slide_w,
                "slide_h_emu": slide_h,
            },
            "paragraphs": all_paragraphs,
        }, f, ensure_ascii=False, indent=2)
    print(f"[saved] {OUTPUT_PATH}")
    print(f"  size: {OUTPUT_PATH.stat().st_size / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
