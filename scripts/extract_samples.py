"""샘플 슬라이드 20장을 별도 .pptx로 추출하여 시각 확인 용도로 저장.

전략: 다양한 장에서 샘플링 (초반/중반/후반 고르게) + 특이 케이스(최대 shape).
"""
from pptx import Presentation
import copy
from lxml import etree
from pptx.oxml.ns import qn
import os

PATH = 'docs/references/_master_templates/PPT 템플릿.pptx'
OUT_DIR = 'docs/references/_master_templates/_samples'
os.makedirs(OUT_DIR, exist_ok=True)

# 고르게 샘플 20장: 1,64,127,...,1251 (1251/20 = 62.55 간격)
indices = [1, 63, 125, 188, 250, 313, 375, 438, 500, 563,
           625, 688, 750, 813, 875, 938, 1000, 1063, 1125, 1250]


def save_slide_subset(src_path, indices, out_path):
    """원본 PPT를 복제하고 지정 인덱스 외 슬라이드를 삭제."""
    prs = Presentation(src_path)
    keep_indices = set(i - 1 for i in indices)  # 0-based

    # xml_slides의 sldIdLst 순회하며 제거 (뒤에서부터)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for idx in reversed(range(len(slides_list))):
        if idx not in keep_indices:
            slide = prs.slides[idx]
            rId = slides_list[idx].get(qn('r:id'))
            prs.part.drop_rel(rId)
            xml_slides.remove(slides_list[idx])

    prs.save(out_path)
    return len(keep_indices)


out_path = os.path.join(OUT_DIR, 'samples_20.pptx')
n = save_slide_subset(PATH, indices, out_path)
print(f"Saved {n} slides to {out_path}")
print(f"Size: {os.path.getsize(out_path)/1024/1024:.1f} MB")
