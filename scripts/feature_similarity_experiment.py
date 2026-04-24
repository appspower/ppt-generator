"""Clustering feature 추천을 위한 pairwise similarity 실험.

- Feature A: grid 6x6 occupancy (binary 36-dim)
- Feature B: shape type histogram (normalized)
- Feature C: structure_sig (count exact match)
- Feature D: grid 8x8 + shape type mix (cell, type)

각 feature로 1251장 샘플링 페어의 similarity 분포와 추정 클러스터 수 체크.
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import Counter
import random
import json
import os

PATH = 'c:/Users/y2kbo/Apps/PPT/docs/references/_master_templates/PPT 템플릿.pptx'
OUT_DIR = 'c:/Users/y2kbo/Apps/PPT/output/catalog'
random.seed(42)


def iter_shapes(shapes):
    for shp in shapes:
        yield shp
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from iter_shapes(shp.shapes)
            except Exception:
                pass


def jaccard(a, b):
    if not a and not b:
        return 1.0
    inter = len(a & b)
    union = len(a | b)
    return inter / union if union else 1.0


def cosine_hist(a: Counter, b: Counter):
    keys = set(a) | set(b)
    if not keys:
        return 1.0
    dot = sum(a[k] * b[k] for k in keys)
    na = sum(v * v for v in a.values()) ** 0.5
    nb = sum(v * v for v in b.values()) ** 0.5
    if na == 0 or nb == 0:
        return 0.0
    return dot / (na * nb)


def main():
    prs = Presentation(PATH)
    slide_w, slide_h = prs.slide_width, prs.slide_height

    # Extract features per slide
    feat_grid6 = []         # set of cells
    feat_grid8 = []         # set of cells
    feat_type_hist = []     # Counter
    feat_typed_grid8 = []   # set of (cell, type)
    feat_geom_buckets = []  # set of (quadrant, type) coarse

    for slide in prs.slides:
        shapes = list(iter_shapes(slide.shapes))
        g6 = set()
        g8 = set()
        typed8 = set()
        hist = Counter()
        coarse = set()
        for shp in shapes:
            try:
                left, top = shp.left or 0, shp.top or 0
                w, h = shp.width or 0, shp.height or 0
                cx = (left + w/2) / slide_w
                cy = (top + h/2) / slide_h
                st = str(shp.shape_type).split('.')[-1] if shp.shape_type else 'NONE'
                hist[st] += 1
                # 6x6 and 8x8
                gx6, gy6 = min(5, max(0, int(cx*6))), min(5, max(0, int(cy*6)))
                gx8, gy8 = min(7, max(0, int(cx*8))), min(7, max(0, int(cy*8)))
                g6.add(gy6*6 + gx6)
                g8.add(gy8*8 + gx8)
                typed8.add((gy8*8 + gx8, st))
                # coarse: 3x3 + type
                gx3, gy3 = min(2, max(0, int(cx*3))), min(2, max(0, int(cy*3)))
                coarse.add((gy3*3 + gx3, st))
            except Exception:
                pass
        feat_grid6.append(g6)
        feat_grid8.append(g8)
        feat_typed_grid8.append(typed8)
        feat_type_hist.append(hist)
        feat_geom_buckets.append(coarse)

    N = len(prs.slides)
    # 5000 random pairs
    pairs = [(random.randint(0, N-1), random.randint(0, N-1)) for _ in range(5000)]
    pairs = [(a, b) for a, b in pairs if a != b]

    def describe(name, sims):
        sims = sorted(sims)
        n = len(sims)
        print(f"\n{name}:")
        print(f"  mean={sum(sims)/n:.3f}, median={sims[n//2]:.3f}")
        print(f"  p10={sims[n//10]:.3f}, p25={sims[n//4]:.3f}")
        print(f"  p75={sims[3*n//4]:.3f}, p90={sims[int(n*0.9)]:.3f}")
        # threshold별 "같은 클러스터" 비율
        for th in (0.5, 0.7, 0.8, 0.9, 0.95):
            hi = sum(1 for s in sims if s >= th)
            print(f"  sim>={th}: {hi} ({hi/n*100:.1f}%)")

    # Compute
    sim_g6 = [jaccard(feat_grid6[a], feat_grid6[b]) for a, b in pairs]
    sim_g8 = [jaccard(feat_grid8[a], feat_grid8[b]) for a, b in pairs]
    sim_typed8 = [jaccard(feat_typed_grid8[a], feat_typed_grid8[b]) for a, b in pairs]
    sim_coarse = [jaccard(feat_geom_buckets[a], feat_geom_buckets[b]) for a, b in pairs]
    sim_type = [cosine_hist(feat_type_hist[a], feat_type_hist[b]) for a, b in pairs]

    describe("[A] grid 6x6 Jaccard", sim_g6)
    describe("[B] grid 8x8 Jaccard", sim_g8)
    describe("[C] grid 8x8 + type (Jaccard)", sim_typed8)
    describe("[D] shape_type histogram cosine", sim_type)
    describe("[E] coarse 3x3 + type (Jaccard)", sim_coarse)

    # Estimated cluster count (threshold-based connected components)
    from itertools import combinations
    # 효율 위해 random subset 500장 사용
    sub = random.sample(range(N), 500)

    def cluster_count(feature_fn, metric_fn, threshold):
        # union-find
        parent = list(range(len(sub)))
        def find(x):
            while parent[x] != x:
                parent[x] = parent[parent[x]]
                x = parent[x]
            return x
        def union(a, b):
            ra, rb = find(a), find(b)
            if ra != rb:
                parent[ra] = rb
        # N^2 (500^2 = 250k)
        feats = [feature_fn(i) for i in sub]
        for i in range(len(sub)):
            for j in range(i+1, len(sub)):
                if metric_fn(feats[i], feats[j]) >= threshold:
                    union(i, j)
        roots = set(find(i) for i in range(len(sub)))
        return len(roots)

    print("\n\n[클러스터 수 추정 (500장 서브샘플, threshold 기반 union-find)]")
    for th in (0.5, 0.6, 0.7, 0.8):
        c_g6 = cluster_count(lambda i: feat_grid6[i], jaccard, th)
        c_g8 = cluster_count(lambda i: feat_grid8[i], jaccard, th)
        c_typed = cluster_count(lambda i: feat_typed_grid8[i], jaccard, th)
        c_coarse = cluster_count(lambda i: feat_geom_buckets[i], jaccard, th)
        c_type = cluster_count(lambda i: feat_type_hist[i], cosine_hist, th)
        print(f"  th={th}:")
        print(f"    grid6x6 Jaccard:        {c_g6}")
        print(f"    grid8x8 Jaccard:        {c_g8}")
        print(f"    grid8x8+type Jaccard:   {c_typed}")
        print(f"    coarse 3x3+type:        {c_coarse}")
        print(f"    shape_type cos-hist:    {c_type}")

    # 저장
    out = {
        'pair_count': len(pairs),
        'sample_subset_size': len(sub),
        'sim_summaries': {
            'grid6x6': {'mean': sum(sim_g6)/len(sim_g6), 'median': sorted(sim_g6)[len(sim_g6)//2]},
            'grid8x8': {'mean': sum(sim_g8)/len(sim_g8), 'median': sorted(sim_g8)[len(sim_g8)//2]},
            'grid8x8_typed': {'mean': sum(sim_typed8)/len(sim_typed8), 'median': sorted(sim_typed8)[len(sim_typed8)//2]},
            'shape_type_cos': {'mean': sum(sim_type)/len(sim_type), 'median': sorted(sim_type)[len(sim_type)//2]},
            'coarse3x3_typed': {'mean': sum(sim_coarse)/len(sim_coarse), 'median': sorted(sim_coarse)[len(sim_coarse)//2]},
        },
    }
    with open(os.path.join(OUT_DIR, 'feature_similarity.json'), 'w', encoding='utf-8') as f:
        json.dump(out, f, indent=2, ensure_ascii=False)


if __name__ == '__main__':
    main()
