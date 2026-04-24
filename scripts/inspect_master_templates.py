"""Phase A1: 1251장 마스터 템플릿 심층 검수."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import Counter
import re

PATH = 'docs/references/_master_templates/PPT 템플릿.pptx'


def inspect():
    prs = Presentation(PATH)
    shape_type_counter = Counter()
    shape_count_per_slide = []
    placeholder_count_per_slide = []
    text_len_per_slide = []
    has_group = 0
    has_image = 0
    has_table = 0
    has_chart = 0
    has_smartart = 0
    tilde_pattern = re.compile(r'~+')
    placeholder_samples = []

    def walk(shapes, state):
        for shp in shapes:
            state['shape_count'] += 1
            st = shp.shape_type
            shape_type_counter[str(st)] += 1
            if st == MSO_SHAPE_TYPE.GROUP:
                state['has_group'] = True
                walk(shp.shapes, state)
                continue
            if st == MSO_SHAPE_TYPE.PICTURE:
                state['has_image'] = True
            if st == MSO_SHAPE_TYPE.TABLE:
                state['has_table'] = True
            if st == MSO_SHAPE_TYPE.CHART:
                state['has_chart'] = True
            try:
                if hasattr(shp, 'element') and shp.element.tag.endswith('}graphicFrame'):
                    for el in shp.element.iter():
                        if 'diagram' in el.tag:
                            state['has_smartart'] = True
                            break
            except Exception:
                pass
            if shp.has_text_frame:
                tf_text = shp.text_frame.text
                state['text_total'] += len(tf_text)
                matches = tilde_pattern.findall(tf_text)
                if matches:
                    state['placeholder_count'] += len(matches)
                    if len(placeholder_samples) < 30 and tf_text.strip():
                        placeholder_samples.append((state['idx'], tf_text[:200]))

    for i, slide in enumerate(prs.slides):
        state = {
            'idx': i + 1, 'shape_count': 0, 'placeholder_count': 0, 'text_total': 0,
            'has_group': False, 'has_image': False, 'has_table': False,
            'has_chart': False, 'has_smartart': False,
        }
        walk(slide.shapes, state)
        shape_count_per_slide.append(state['shape_count'])
        placeholder_count_per_slide.append(state['placeholder_count'])
        text_len_per_slide.append(state['text_total'])
        if state['has_group']: has_group += 1
        if state['has_image']: has_image += 1
        if state['has_table']: has_table += 1
        if state['has_chart']: has_chart += 1
        if state['has_smartart']: has_smartart += 1

    total = len(prs.slides)
    print("=" * 60)
    print(f"SLIDE COUNT: {total}")
    print("=" * 60)
    print("\n[Shape 타입 분포 (전체)]")
    for st, cnt in shape_type_counter.most_common(20):
        print(f"  {st:50s} {cnt:6d}")
    print(f"\n[Shape 수 분포]")
    print(f"  평균: {sum(shape_count_per_slide)/total:.1f}")
    print(f"  최대: {max(shape_count_per_slide)}")
    print(f"  최소: {min(shape_count_per_slide)}")
    print(f"  중앙값: {sorted(shape_count_per_slide)[total//2]}")
    print(f"\n[복잡 요소 보유 슬라이드 비율]")
    print(f"  GROUP 포함:    {has_group:5d} ({has_group/total*100:.1f}%)")
    print(f"  PICTURE 포함:  {has_image:5d} ({has_image/total*100:.1f}%)")
    print(f"  TABLE 포함:    {has_table:5d} ({has_table/total*100:.1f}%)")
    print(f"  CHART 포함:    {has_chart:5d} ({has_chart/total*100:.1f}%)")
    print(f"  SMARTART 포함: {has_smartart:5d} ({has_smartart/total*100:.1f}%)")
    print(f"\n[텍스트 길이 분포]")
    print(f"  평균: {sum(text_len_per_slide)/total:.0f}자")
    print(f"  최대: {max(text_len_per_slide)}자")
    print(f"  최소: {min(text_len_per_slide)}자")
    print(f"\n[`~` placeholder 분포]")
    with_ph = sum(1 for c in placeholder_count_per_slide if c > 0)
    print(f"  placeholder 있는 슬라이드: {with_ph} ({with_ph/total*100:.1f}%)")
    print(f"  placeholder 총 개수: {sum(placeholder_count_per_slide)}")
    print(f"  슬라이드당 평균: {sum(placeholder_count_per_slide)/total:.1f}")
    print(f"  슬라이드당 최대: {max(placeholder_count_per_slide)}")
    print(f"\n[placeholder 샘플 (최대 30개)]")
    for slide_idx, sample in placeholder_samples[:30]:
        print(f"  슬라이드 {slide_idx}: {sample[:120]!r}")


if __name__ == '__main__':
    inspect()
