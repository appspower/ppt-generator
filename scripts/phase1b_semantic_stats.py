"""Phase 1B - OOXML 의미 구조 통계 수집 (1,251장).

슬라이드를 단순 shape list가 아니라 "사람이 의도한 묶음(group)"으로 본다.

측정 항목:
  B1. Group shape nesting / leaf 수
  B2. 텍스트 영역 위치 패턴 (제일 큰 텍스트, banner/sidebar/footer)
  B3. Shape 정렬 격자 자동 추론 (좌측 클러스터링 → N열 × M행)
  B4. 시그니처 요소 검출 (corner marker, accent strip, divider)
  B5. ~~ 분포 (8x8 grid)

산출물:
  output/catalog/phase1b_semantic_stats.json   (per-slide + global summary)
"""
from __future__ import annotations

import json
import os
import re
import statistics
from collections import Counter, defaultdict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PATH = 'c:/Users/y2kbo/Apps/PPT/docs/references/_master_templates/PPT 템플릿.pptx'
META_PATH = 'c:/Users/y2kbo/Apps/PPT/output/catalog/slide_meta.json'
OUT_DIR = 'c:/Users/y2kbo/Apps/PPT/output/catalog'
OUT_PER_SLIDE = os.path.join(OUT_DIR, 'phase1b_semantic_stats.json')

PLACEHOLDER_TOKEN_RE = re.compile(r'~~+')
CODE_RE = re.compile(r'^[A-Z]\d{1,2}$')

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def iter_leaf_shapes(shapes, depth=0, parent_group_id=None):
    """leaf shape만 yield. (shape, depth, parent_group_id)."""
    for shp in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                yield from iter_leaf_shapes(shp.shapes, depth + 1, id(shp))
            except Exception:
                pass
        else:
            yield shp, depth, parent_group_id


def iter_groups(shapes, depth=0):
    """group shape만 (group_shape, depth)."""
    for shp in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield shp, depth
            try:
                yield from iter_groups(shp.shapes, depth + 1)
            except Exception:
                pass


def safe_emu(v):
    return v if v is not None else 0


def text_of(shp):
    try:
        if shp.has_text_frame:
            return shp.text_frame.text or ''
    except Exception:
        return ''
    return ''


def max_font_size_pt(shp):
    """텍스트 프레임 내 모든 run을 훑어 최대 font size(pt) 반환. None이면 -1."""
    best = -1.0
    try:
        if not shp.has_text_frame:
            return best
        for para in shp.text_frame.paragraphs:
            for run in para.runs:
                try:
                    sz = run.font.size
                    if sz is not None:
                        pt = sz.pt
                        if pt > best:
                            best = pt
                except Exception:
                    pass
    except Exception:
        pass
    return best


def char_count(shp):
    try:
        return len(text_of(shp))
    except Exception:
        return 0


# ---------------------------------------------------------------------------
# B2 area classification
# ---------------------------------------------------------------------------
def classify_text_area(left_n, top_n, w_n, h_n):
    """가장 큰 텍스트의 위치를 9-zone으로 라벨."""
    cx = left_n + w_n / 2
    cy = top_n + h_n / 2
    # banner: 가로 폭이 슬라이드 70% 이상
    if w_n >= 0.70 and top_n < 0.20:
        return 'top_banner'
    if w_n >= 0.70 and top_n >= 0.80:
        return 'bottom_banner'
    if h_n >= 0.60 and left_n < 0.20:
        return 'left_sidebar'
    if h_n >= 0.60 and left_n >= 0.65:
        return 'right_sidebar'
    # 9-zone
    col = 'left' if cx < 1/3 else ('center' if cx < 2/3 else 'right')
    row = 'top' if cy < 1/3 else ('mid' if cy < 2/3 else 'bot')
    return f'{row}_{col}'


# ---------------------------------------------------------------------------
# B3 grid inference (1-D clustering on left/top)
# ---------------------------------------------------------------------------
def cluster_1d(values, tol):
    """단순 1D clustering: 정렬 후 인접 차이가 tol 이하면 같은 그룹.

    Returns: list of (centroid, members[]).
    """
    if not values:
        return []
    vs = sorted(values)
    clusters = [[vs[0]]]
    for v in vs[1:]:
        if v - clusters[-1][-1] <= tol:
            clusters[-1].append(v)
        else:
            clusters.append([v])
    return [(sum(c) / len(c), c) for c in clusters]


def infer_grid(left_norms, top_norms, tol=0.06, min_members=3):
    """left/top 좌표 클러스터링으로 (cols, rows) 추론.

    min_members 이상이 들어간 클러스터만 카운트.
    tol=0.06 (~6% of slide width/height): 작은 떨림은 동일 컬럼/행으로 묶음.
    min_members=3: 우연 정렬 2개는 그리드로 보지 않음.
    """
    cols = [c for c in cluster_1d(left_norms, tol) if len(c[1]) >= min_members]
    rows = [c for c in cluster_1d(top_norms, tol) if len(c[1]) >= min_members]
    return len(cols), len(rows)


# ---------------------------------------------------------------------------
# B4 signature elements
# ---------------------------------------------------------------------------
def detect_signatures(shapes_meta, slide_w_emu, slide_h_emu):
    """시그니처 요소 검출.

    shapes_meta: list of dict with left_n, top_n, w_n, h_n, w_emu, h_emu, has_text, text, font_pt
    Returns: dict.
    """
    sig = {
        'corner_top_left_marker': False,    # 좌상단 작은 박스 (w<0.10, h<0.06, top<0.10, left<0.10)
        'corner_top_right_marker': False,
        'left_accent_strip': False,         # 좌측 얇은 strip (w<0.04, h>0.50)
        'right_accent_strip': False,
        'top_divider': False,               # 가로 얇은 라인 (w>0.50, h<0.02, top<0.30)
        'bottom_divider': False,
        'page_number_loc': None,            # 'br', 'bl', 'bc', 'tr', 'tl', 'tc', None
    }
    for s in shapes_meta:
        l, t, w, h = s['left_n'], s['top_n'], s['w_n'], s['h_n']
        # corner marker (작은 영역, 텍스트가 코드 또는 짧은 라벨)
        if w < 0.10 and h < 0.07 and l < 0.10 and t < 0.10:
            sig['corner_top_left_marker'] = True
        if w < 0.10 and h < 0.07 and l > 0.85 and t < 0.10:
            sig['corner_top_right_marker'] = True
        # accent strip (얇고 길쭉)
        if w < 0.04 and h > 0.50 and l < 0.10:
            sig['left_accent_strip'] = True
        if w < 0.04 and h > 0.50 and l > 0.92:
            sig['right_accent_strip'] = True
        # divider
        if w > 0.50 and h < 0.02:
            if t < 0.30:
                sig['top_divider'] = True
            elif t > 0.70:
                sig['bottom_divider'] = True
        # 페이지 번호: 짧은 숫자 텍스트
        txt = (s.get('text') or '').strip()
        if txt and txt.isdigit() and 1 <= len(txt) <= 3 and w < 0.10 and h < 0.06:
            cx = l + w / 2
            cy = t + h / 2
            row = 't' if cy < 0.20 else ('b' if cy > 0.80 else 'm')
            col = 'l' if cx < 0.20 else ('r' if cx > 0.80 else 'c')
            if row in ('t', 'b'):
                sig['page_number_loc'] = row + col
    return sig


# ---------------------------------------------------------------------------
# B5 placeholder distribution
# ---------------------------------------------------------------------------
def grid_cell_idx(cx, cy, g=8):
    gx = min(g - 1, max(0, int(cx * g)))
    gy = min(g - 1, max(0, int(cy * g)))
    return gy * g + gx


def placeholder_grid_signature(shapes_meta, g=8):
    """~~가 들어간 텍스트 박스의 (cx, cy) → 8x8 grid cell 멤버십.

    Returns: (placeholder_count, set_of_cells, cluster_score)
      cluster_score = 1 - (unique_cells / total_placeholders)  if total>0 else 0
    """
    cells = []
    pct = 0
    for s in shapes_meta:
        txt = s.get('text') or ''
        n = len(PLACEHOLDER_TOKEN_RE.findall(txt))
        if n == 0:
            continue
        pct += n
        cx = s['left_n'] + s['w_n'] / 2
        cy = s['top_n'] + s['h_n'] / 2
        cells.append(grid_cell_idx(cx, cy, g))
    if not cells:
        return 0, [], 0.0
    unique_cells = set(cells)
    cluster_score = 1.0 - (len(unique_cells) / len(cells))
    return pct, sorted(unique_cells), round(cluster_score, 3)


# ---------------------------------------------------------------------------
# Archetype classification
# ---------------------------------------------------------------------------
def _classify_archetype(*, cols, rows, largest_zone, sig, ph_count, leaf_count, group_count_top):
    """슬라이드의 의미 archetype 추정. 우선순위 룰 기반."""
    # cover / divider류 — 텍스트 위주 + 본문 셀 거의 없음
    if leaf_count <= 5 and ph_count <= 8 and largest_zone in ('top_banner', 'mid_center', 'top_center'):
        return 'cover_or_divider'
    # 좌측 큰 텍스트 + 우측 본문 (좌측 banner/sidebar 형)
    if largest_zone in ('left_sidebar', 'mid_left', 'top_left') and cols >= 2:
        return 'left_title_right_body'
    # 본문 행렬
    if cols == 2 and rows == 2:
        return '2x2_matrix'
    if cols == 3 and rows == 3:
        return '3x3_matrix'
    if cols >= 4 and rows >= 4:
        return 'dense_grid'
    # 다열 비교
    if cols >= 3 and rows <= 2:
        return f'{cols}col_compare'
    if cols == 2 and rows <= 2:
        return '2col_compare'
    # 단일/소수 열, 다행 (목록)
    if cols == 1 and rows >= 3:
        return 'vertical_list'
    if cols <= 1 and rows <= 1 and ph_count <= 5:
        return 'single_block'
    # 기타
    if sig.get('left_accent_strip') or sig.get('right_accent_strip'):
        return 'accent_strip_layout'
    return 'mixed'


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def run():
    prs = Presentation(PATH)
    sw, sh = prs.slide_width, prs.slide_height
    print(f'slide dim emu: {sw}x{sh} ({sw/914400:.2f}x{sh/914400:.2f} in)')

    per_slide = []

    # global aggregates
    grid_shape_dist = Counter()             # (cols, rows) → count
    largest_text_zone_ctr = Counter()
    sig_ctr = Counter()
    page_num_loc_ctr = Counter()
    placeholder_global_count = []
    placeholder_unique_cells = []
    placeholder_clustered = []  # cluster_score >= 0.5
    group_leaf_counts_global = []
    group_depths_global = []
    group_count_per_slide = []

    for i, slide in enumerate(prs.slides):
        # leaf shapes ----------------------------------------------------------
        leaves = list(iter_leaf_shapes(slide.shapes))
        groups = list(iter_groups(slide.shapes))

        # B1: group nesting
        max_depth = 0
        leaves_per_group = defaultdict(int)
        for shp, depth, parent_gid in leaves:
            if parent_gid is not None:
                leaves_per_group[parent_gid] += 1
                if depth > max_depth:
                    max_depth = depth
        group_leaf_sizes = list(leaves_per_group.values())
        group_count_top = sum(1 for g, d in groups if d == 0)  # top-level group
        group_count_total = len(groups)
        group_count_per_slide.append(group_count_top)
        group_depths_global.append(max_depth)
        group_leaf_counts_global.extend(group_leaf_sizes)

        # build shapes_meta for B2-B5
        shapes_meta = []
        largest_text = None  # (font_pt or h_n proxy, dict)
        for shp, depth, _gid in leaves:
            try:
                left = safe_emu(shp.left)
                top = safe_emu(shp.top)
                w = safe_emu(shp.width)
                h = safe_emu(shp.height)
            except Exception:
                continue
            ln = left / sw if sw else 0
            tn = top / sh if sh else 0
            wn = w / sw if sw else 0
            hn = h / sh if sh else 0
            ln = max(0.0, min(1.5, ln))
            tn = max(0.0, min(1.5, tn))
            txt = text_of(shp)
            font_pt = max_font_size_pt(shp)
            has_text = bool(txt.strip())
            is_ph = bool(getattr(shp, 'is_placeholder', False))
            try:
                st = str(shp.shape_type).split('.')[-1] if shp.shape_type else 'NONE'
            except Exception:
                st = 'NONE'
            sm = {
                'left_n': ln, 'top_n': tn, 'w_n': wn, 'h_n': hn,
                'w_emu': w, 'h_emu': h,
                'has_text': has_text, 'text': txt[:200],
                'font_pt': font_pt,
                'is_placeholder': is_ph,
                'shape_type': st,
                'depth': depth,
            }
            shapes_meta.append(sm)

            if has_text:
                # 우선순위: 명시 font size > 텍스트 박스 높이
                key = font_pt if font_pt > 0 else hn * 100  # fallback proxy
                if largest_text is None or key > largest_text[0]:
                    largest_text = (key, sm)

        # B2: largest text zone
        largest_zone = None
        if largest_text is not None:
            sm = largest_text[1]
            largest_zone = classify_text_area(sm['left_n'], sm['top_n'], sm['w_n'], sm['h_n'])
            largest_text_zone_ctr[largest_zone] += 1

        # B3: grid inference - "본문 단위" 셀 추정
        #  - 너무 작은 marker(코너/accent), 너무 큰 background(>=80% area) 제외
        #  - 텍스트가 있는 박스 또는 8% 이상 크기의 도형
        grid_targets = []
        for s in shapes_meta:
            area = s['w_n'] * s['h_n']
            if area >= 0.80:               # background frame 제외
                continue
            if s['w_n'] < 0.04 and s['h_n'] < 0.04:   # tiny marker 제외
                continue
            if s['has_text'] or (s['w_n'] > 0.08 and s['h_n'] > 0.05):
                grid_targets.append(s)
        lefts = [s['left_n'] for s in grid_targets]
        tops = [s['top_n'] for s in grid_targets]
        cols, rows = infer_grid(lefts, tops)
        grid_shape_dist[(cols, rows)] += 1

        # B4: signatures
        sig = detect_signatures(shapes_meta, sw, sh)
        for k, v in sig.items():
            if k == 'page_number_loc':
                if v:
                    page_num_loc_ctr[v] += 1
            elif v:
                sig_ctr[k] += 1

        # B5: placeholder grid
        ph_count, ph_cells, ph_cluster = placeholder_grid_signature(shapes_meta, g=8)
        placeholder_global_count.append(ph_count)
        placeholder_unique_cells.extend(ph_cells)
        if ph_count > 0 and ph_cluster >= 0.5:
            placeholder_clustered.append(i)

        # archetype tagging (semantic shorthand)
        archetype = _classify_archetype(
            cols=cols, rows=rows,
            largest_zone=largest_zone,
            sig=sig,
            ph_count=ph_count,
            leaf_count=len(shapes_meta),
            group_count_top=group_count_top,
        )

        per_slide.append({
            'slide_index': i,
            'layout_name': slide.slide_layout.name if slide.slide_layout else '',
            'leaf_shape_count': len(shapes_meta),
            'group_count_top': group_count_top,
            'group_count_total': group_count_total,
            'max_group_depth': max_depth,
            'group_leaf_sizes': group_leaf_sizes,
            'largest_text_zone': largest_zone,
            'largest_text_font_pt': round(largest_text[0], 1) if largest_text else None,
            'grid_cols': cols,
            'grid_rows': rows,
            'sig': sig,
            'placeholder_count': ph_count,
            'placeholder_cells': ph_cells,
            'placeholder_cluster_score': ph_cluster,
            'archetype': archetype,
        })

    total = len(per_slide)

    # global summary
    glc = group_leaf_counts_global
    glc_stats = {
        'count': len(glc),
        'mean': round(statistics.mean(glc), 2) if glc else 0,
        'median': statistics.median(glc) if glc else 0,
        'p25': sorted(glc)[len(glc)//4] if glc else 0,
        'p75': sorted(glc)[3*len(glc)//4] if glc else 0,
        'p95': sorted(glc)[int(0.95*len(glc))] if glc else 0,
        'max': max(glc) if glc else 0,
    }

    gcps = group_count_per_slide
    pgc = placeholder_global_count
    placeholder_stats = {
        'mean': round(statistics.mean(pgc), 2),
        'median': statistics.median(pgc),
        'p25': sorted(pgc)[total//4],
        'p75': sorted(pgc)[3*total//4],
        'p95': sorted(pgc)[int(0.95*total)],
        'max': max(pgc),
        'slides_with_zero': sum(1 for v in pgc if v == 0),
        'slides_with_ge10': sum(1 for v in pgc if v >= 10),
    }

    # grid shape buckets
    grid_shape_named = Counter()
    for (c, r), cnt in grid_shape_dist.items():
        if c == 0 and r == 0:
            label = 'free_form'
        elif c <= 1 and r <= 1:
            label = 'single_block'
        elif c == 1:
            label = f'1col_x_{r}row'
        elif r == 1:
            label = f'{c}col_x_1row'
        elif c == 2 and r == 2:
            label = '2x2_matrix'
        elif c == 3 and r == 3:
            label = '3x3_matrix'
        elif c >= 2 and r >= 2:
            label = f'{c}col_x_{r}row_grid'
        else:
            label = f'{c}col_x_{r}row'
        grid_shape_named[label] += cnt

    archetype_ctr = Counter(p['archetype'] for p in per_slide)

    summary = {
        'total_slides': total,
        'slide_dim_emu': [sw, sh],
        'archetype_dist': dict(archetype_ctr.most_common()),
        # B1
        'group': {
            'slides_with_group': sum(1 for v in gcps if v > 0),
            'top_level_group_count_per_slide': {
                'mean': round(statistics.mean(gcps), 2),
                'median': statistics.median(gcps),
                'max': max(gcps),
                'p95': sorted(gcps)[int(0.95*total)],
            },
            'leaves_per_group': glc_stats,
            'max_depth_dist': dict(Counter(group_depths_global).most_common()),
        },
        # B2
        'largest_text_zone_dist': dict(largest_text_zone_ctr.most_common()),
        # B3
        'grid_shape_dist_top20': dict(Counter(
            {f'{c}x{r}': cnt for (c, r), cnt in grid_shape_dist.items()}
        ).most_common(20)),
        'grid_shape_named_top20': dict(grid_shape_named.most_common(20)),
        'free_form_count': grid_shape_dist.get((0, 0), 0),
        # B4
        'signature_counts': dict(sig_ctr.most_common()),
        'page_number_loc_dist': dict(page_num_loc_ctr.most_common()),
        # B5
        'placeholder_per_slide': placeholder_stats,
        'placeholder_clustered_slides_count': len(placeholder_clustered),
        'placeholder_cell_freq_top20': dict(Counter(placeholder_unique_cells).most_common(20)),
    }

    out = {
        'summary': summary,
        'per_slide': per_slide,
    }

    with open(OUT_PER_SLIDE, 'w', encoding='utf-8') as f:
        json.dump(out, f, ensure_ascii=False, indent=2, default=str)
    print(f'\nsaved -> {OUT_PER_SLIDE}')

    # console snapshot
    print('\n==== B1: GROUP NESTING ====')
    print(f"slides_with_group:       {summary['group']['slides_with_group']}/{total}")
    print(f"top group count median:  {summary['group']['top_level_group_count_per_slide']['median']}, "
          f"mean={summary['group']['top_level_group_count_per_slide']['mean']}, "
          f"max={summary['group']['top_level_group_count_per_slide']['max']}")
    print(f"leaves_per_group:        median={glc_stats['median']}, mean={glc_stats['mean']}, "
          f"p75={glc_stats['p75']}, p95={glc_stats['p95']}, max={glc_stats['max']}")
    print(f"max group depth dist:    {summary['group']['max_depth_dist']}")

    print('\n==== B2: LARGEST TEXT ZONE ====')
    for z, c in largest_text_zone_ctr.most_common():
        print(f'  {z:18s} {c:5d}  ({c/total*100:5.1f}%)')

    print('\n==== B3: GRID SHAPES (top 15) ====')
    for label, c in grid_shape_named.most_common(15):
        print(f'  {label:25s} {c:5d}  ({c/total*100:5.1f}%)')
    print(f"  free_form (0x0): {summary['free_form_count']}")

    print('\n==== B4: SIGNATURE ELEMENTS ====')
    for k, c in sig_ctr.most_common():
        print(f'  {k:25s} {c:5d}  ({c/total*100:5.1f}%)')
    print(f'  page_number_loc dist: {dict(page_num_loc_ctr.most_common())}')

    print('\n==== B5: PLACEHOLDER (~~) ====')
    print(f"  per slide: median={placeholder_stats['median']}, mean={placeholder_stats['mean']}, "
          f"p75={placeholder_stats['p75']}, p95={placeholder_stats['p95']}, max={placeholder_stats['max']}")
    print(f"  zero-placeholder slides: {placeholder_stats['slides_with_zero']}")
    print(f"  >=10 placeholders:       {placeholder_stats['slides_with_ge10']}")
    print(f"  clustered (score>=0.5):  {len(placeholder_clustered)}")

    print('\n==== ARCHETYPE DIST ====')
    for k, c in archetype_ctr.most_common():
        print(f'  {k:25s} {c:5d}  ({c/total*100:5.1f}%)')

    return summary


if __name__ == '__main__':
    run()
