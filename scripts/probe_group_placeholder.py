"""GROUP 내부 `~~` placeholder 구조 및 치환 가능성 확인.

핵심 질문:
1. 그룹 안의 텍스트도 python-pptx로 접근 가능한가?
2. 그룹 내부 Run 단위의 폰트/색상 속성 보존은?
3. `~~`를 실제 텍스트로 치환 후 PPT가 정상 열리는가?
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import copy, os

PATH = 'docs/references/_master_templates/PPT 템플릿.pptx'
OUT = 'docs/references/_master_templates/_samples/test_replace.pptx'

os.makedirs(os.path.dirname(OUT), exist_ok=True)

prs = Presentation(PATH)

# 첫 20장만 대상으로 실험
test_slides = list(prs.slides)[:20]

total_shapes = 0
group_shapes = 0
text_shapes = 0
placeholder_hits = 0
replaced = 0


def walk_and_report(shapes, depth=0):
    global total_shapes, group_shapes, text_shapes, placeholder_hits, replaced
    for shp in shapes:
        total_shapes += 1
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_shapes += 1
            walk_and_report(shp.shapes, depth + 1)
            continue
        if shp.has_text_frame:
            text_shapes += 1
            tf = shp.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if '~~' in run.text:
                        placeholder_hits += 1
                        # Run 단위 속성 확인
                        font = run.font
                        sample = {
                            'text': run.text[:50],
                            'size': font.size.pt if font.size else None,
                            'bold': font.bold,
                            'name': font.name,
                            'depth': depth,
                        }
                        if placeholder_hits <= 5:
                            print(f"    Run 샘플: {sample}")
                        # 실제 치환 테스트 (첫 5개만)
                        if replaced < 5:
                            original = run.text
                            run.text = run.text.replace('~~', '[TEST]')
                            replaced += 1
                            print(f"    치환: {original!r} → {run.text!r}")


for i, slide in enumerate(test_slides):
    print(f"\n[슬라이드 {i+1}]")
    walk_and_report(slide.shapes)

print(f"\n=== 요약 (첫 20장) ===")
print(f"  총 shape: {total_shapes}")
print(f"  group: {group_shapes}")
print(f"  text 포함: {text_shapes}")
print(f"  `~~` Run 수: {placeholder_hits}")
print(f"  치환 테스트: {replaced}")

prs.save(OUT)
print(f"\n저장: {OUT}")
print(f"Size: {os.path.getsize(OUT)/1024/1024:.1f} MB")
