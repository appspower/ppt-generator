"""Mode A (Whole-slide clone) 실측 시뮬레이션.

목표
----
'D철강 2030 넷제로 전환 로드맵' 10장 PPT를 transformation_roadmap_10
스켈레톤(opening → situation → ... → closing)으로 생성하여 실제 작동성을
측정한다.

파이프라인
----------
1. final_labels.json 로드 → narrative_role별 후보 풀 구성
2. 각 role에 대해 slide_index 1개 선택 (fallback 규칙 포함)
3. 마스터 템플릿에서 10장을 SlideCloner로 복제하여 새 .pptx 생성
4. edit_ops.replace_paragraph로 가상 컨텐츠 주입
   (가장 긴 텍스트 paragraph를 'title-ish'로 간주해 1차 교체;
    너무 공격적으로 다 바꾸면 시각 일관성 망가짐 → 보수적으로 1~2개만)
5. PowerPoint COM으로 PNG 렌더
6. 결과 출력: 매칭/실패 표 + PNG 경로 리스트

CLI: `python scripts/simulate_mode_a.py`
"""
from __future__ import annotations

import json
import sys
import time
from collections import Counter
from pathlib import Path

# 프로젝트 루트
ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation

from ppt_builder.template import edit_ops
from ppt_builder.template.editor import TemplateEditor


# ----------------------------------------------------------------------------
# 상수
# ----------------------------------------------------------------------------

CATALOG_PATH = ROOT / "output" / "catalog" / "final_labels.json"
TEMPLATE_PATH = (
    ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
)
OUTPUT_DIR = ROOT / "output" / "simulation"
PPTX_OUT = OUTPUT_DIR / "netzero_modeA.pptx"
PNG_DIR = OUTPUT_DIR / "netzero_modeA_pngs"


# transformation_roadmap_10 스켈레톤 — 10장 narrative
NARRATIVE_SEQUENCE = [
    "opening",
    "situation",
    "complication",
    "analysis",
    "recommendation",
    "roadmap",
    "roadmap",
    "benefit",
    "risk",
    "closing",
]


# 가상 컨텐츠 (1슬라이드당 1개의 hero 문장만 — Mode A의 보수적 교체 대상)
VIRTUAL_CONTENT = {
    "opening": "D철강 2030 넷제로 전환 로드맵",
    "situation": "2024년 배출량 1,200만톤 / 산업 평균 대비 +15%",
    "complication": "EU CBAM 2026 시행 / 비용 영향 연 800억원",
    "analysis": "현 배출 1,200만톤 vs 2030 목표 600만톤 / 갭 50%",
    "recommendation": "전기로 전환 + 수소환원제철 + 그린전력 PPA",
    "roadmap_phase1": "Phase 1 (2026-2027) 전기로 도입, 폐열회수 250만톤 감축",
    "roadmap_phase2": "Phase 2 (2028-2030) 수소환원 파일럿, 350만톤 추가 감축",
    "benefit": "탄소비용 절감 4,200억원 / ESG 등급 A 진입",
    "risk": "수소 단가 변동 / 정책 후퇴 / 설비 투자 회수기간 8년",
    "closing": "2030 600만톤 / 2040 Net-Zero / 다음 단계 의사결정",
}


# Fallback 규칙 — narrative_role 매칭이 부족할 때 archetype 기반으로 보강
ARCHETYPE_FALLBACK = {
    "opening": ["cover_divider"],
    "closing": ["cover_divider"],
    "situation": ["cards_3col", "cards_2col", "left_title_right_body"],
    "complication": ["cards_3col", "left_title_right_body", "vertical_list"],
    "analysis": ["matrix_2x2", "table_native", "dense_grid"],
    "recommendation": ["cards_3col", "vertical_list", "flowchart"],
    "roadmap": ["roadmap", "timeline_h"],
    "benefit": ["cards_3col", "cards_2col", "vertical_list"],
    "risk": ["matrix_2x2", "table_native", "cards_3col"],
}


# ----------------------------------------------------------------------------
# Retrieval
# ----------------------------------------------------------------------------

def load_catalog() -> list[dict]:
    with open(CATALOG_PATH, encoding="utf-8") as f:
        return json.load(f)["labels"]


def candidates_for_role(
    labels: list[dict], role: str, used: set[int]
) -> tuple[list[dict], str]:
    """role에 맞는 후보 리스트와 매칭 소스('role'/'archetype'/'none') 반환."""
    direct = [
        l for l in labels
        if role in l.get("narrative_role", []) and l["slide_index"] not in used
    ]
    if direct:
        # 신뢰도 높은 것 우선
        direct.sort(key=lambda l: l.get("overall_confidence", 0), reverse=True)
        return direct, "role"

    # fallback: archetype
    fallback_archs = ARCHETYPE_FALLBACK.get(role, [])
    fb = [
        l for l in labels
        if any(a in fallback_archs for a in l.get("archetype", []))
        and l["slide_index"] not in used
    ]
    if fb:
        fb.sort(key=lambda l: l.get("overall_confidence", 0), reverse=True)
        return fb, "archetype"

    return [], "none"


def select_deck(labels: list[dict]) -> list[dict]:
    """10장 narrative_sequence를 결정. 각 항목은 {role, slide_index, source}."""
    used: set[int] = set()
    plan = []
    for role in NARRATIVE_SEQUENCE:
        cands, source = candidates_for_role(labels, role, used)
        if not cands:
            plan.append({"role": role, "slide_index": None, "source": "none"})
            continue
        chosen = cands[0]
        used.add(chosen["slide_index"])
        plan.append({
            "role": role,
            "slide_index": chosen["slide_index"],
            "source": source,
            "archetype": chosen.get("archetype", []),
            "macro": chosen.get("macro"),
            "confidence": chosen.get("overall_confidence", 0),
        })
    return plan


# ----------------------------------------------------------------------------
# 텍스트 capacity 측정 (max_chars 프록시)
# ----------------------------------------------------------------------------

def slide_text_summary(src_prs, slide_index: int) -> dict:
    slide = src_prs.slides[slide_index]
    para_records = []
    for fi, sh in edit_ops.iter_leaf_shapes(slide):
        if not getattr(sh, "has_text_frame", False):
            continue
        try:
            tf = sh.text_frame
        except Exception:
            continue
        for pi, p in enumerate(tf.paragraphs):
            text = p.text
            if text:
                para_records.append({
                    "div_id": fi,
                    "para_id": pi,
                    "text": text,
                    "len": len(text),
                })
    para_records.sort(key=lambda r: r["len"], reverse=True)
    return {
        "slide_index": slide_index,
        "n_paragraphs": len(para_records),
        "max_text_len": max((r["len"] for r in para_records), default=0),
        "total_chars": sum(r["len"] for r in para_records),
        "top3": para_records[:3],
        "all": para_records,
    }


# ----------------------------------------------------------------------------
# 슬라이드 빌드
# ----------------------------------------------------------------------------

def _reorder_sldIdLst(prs, original_indices_in_keep_order: list[int],
                      desired_order: list[int]) -> None:
    """keep_slides 후의 슬라이드 순서를 desired_order(원본 슬라이드 인덱스 기준)로 재정렬.

    keep_slides는 원본 인덱스 오름차순으로 보존하므로, sldIdLst의 r:id 순서를
    desired_order에 맞게 재배열한다.

    original_indices_in_keep_order: keep_slides 직후의 0..N-1 위치 → 원본 인덱스 매핑
    desired_order: 원하는 원본 인덱스 시퀀스 (중복 X 가정)
    """
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    assert len(children) == len(original_indices_in_keep_order), \
        f"len mismatch: {len(children)} vs {len(original_indices_in_keep_order)}"
    pos_by_original = {orig: i for i, orig in enumerate(original_indices_in_keep_order)}
    new_order = [children[pos_by_original[orig]] for orig in desired_order]
    for c in children:
        sldIdLst.remove(c)
    for c in new_order:
        sldIdLst.append(c)


def build_pptx(plan: list[dict]) -> dict:
    """plan에 따라 마스터에서 10장 추출 + edit_ops로 컨텐츠 교체.

    TemplateEditor.keep_slides 사용 — 관계(rId, 이미지, 차트, 폰트) 100% 보존.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # 0. 우선 src_prs에서 텍스트 capacity 측정 (편집 전 원본 기준)
    print(f"[open] template = {TEMPLATE_PATH.name}", flush=True)
    src_prs = Presentation(str(TEMPLATE_PATH))

    summaries: dict[int, dict] = {}
    plan_with_idx = [p for p in plan if p["slide_index"] is not None]
    for p in plan_with_idx:
        summaries[p["slide_index"]] = slide_text_summary(src_prs, p["slide_index"])
    del src_prs

    # 1. TemplateEditor로 keep_slides + 재정렬
    desired_order = [p["slide_index"] for p in plan if p["slide_index"] is not None]
    keep_unique_sorted = sorted(set(desired_order))

    print(f"[editor] keep_slides {keep_unique_sorted}", flush=True)
    editor = TemplateEditor(TEMPLATE_PATH)
    editor.keep_slides(keep_unique_sorted)
    print(f"  kept = {editor.slide_count} slides", flush=True)

    # 재정렬: keep_unique_sorted 순서 → desired_order 순서로
    _reorder_sldIdLst(editor.prs, keep_unique_sorted, desired_order)
    print(f"  reordered to narrative sequence", flush=True)

    # 2. 각 슬라이드에 edit_ops 적용
    edit_results = []
    for step_idx, item in enumerate(plan):
        role = item["role"]
        sidx = item["slide_index"]
        step = step_idx + 1
        if sidx is None:
            edit_results.append({
                "step": step,
                "role": role,
                "slide_index": None,
                "edit_ok": False,
                "edit_reason": "no_candidate",
                "text_summary": None,
            })
            continue

        summary = summaries[sidx]

        # 가상 컨텐츠 결정
        if role == "roadmap":
            already = sum(1 for r in edit_results if r["role"] == "roadmap")
            content_key = "roadmap_phase1" if already == 0 else "roadmap_phase2"
        else:
            content_key = role
        new_text = VIRTUAL_CONTENT.get(content_key, role.upper())

        # 재정렬 후 슬라이드 위치 = step_idx
        # (한 원본 인덱스가 두 번 등장할 수 없는 가정 = 우리 plan에서는 모두 unique)
        target_slide = editor.prs.slides[step_idx]

        edit_ok = False
        edit_reason = ""
        target_div_id = None
        target_para_id = None
        if summary["top3"]:
            cand = summary["top3"][0]
            target_div_id = cand["div_id"]
            target_para_id = cand["para_id"]
            try:
                edit_ops.replace_paragraph(
                    target_slide, cand["div_id"], cand["para_id"], new_text
                )
                edit_ok = True
                edit_reason = (
                    f"replaced div={cand['div_id']} para={cand['para_id']} "
                    f"(orig='{cand['text'][:20]}', new_len={len(new_text)})"
                )
            except Exception as e:
                edit_reason = f"replace_paragraph FAILED: {type(e).__name__}: {e}"
        else:
            edit_reason = "no text paragraph found"

        print(
            f"  [edit] step {step:2d} {role:>14}  src#{sidx:4d}  "
            f"div={target_div_id} para={target_para_id} "
            f"{'OK' if edit_ok else 'FAIL'}",
            flush=True,
        )

        edit_results.append({
            "step": step,
            "role": role,
            "slide_index": sidx,
            "content_key": content_key,
            "new_text": new_text,
            "edit_ok": edit_ok,
            "edit_reason": edit_reason,
            "target_div_id": target_div_id,
            "target_para_id": target_para_id,
            "text_summary": {
                "n_paragraphs": summary["n_paragraphs"],
                "max_text_len": summary["max_text_len"],
                "total_chars": summary["total_chars"],
                "top3": [
                    {"div": p["div_id"], "para": p["para_id"], "len": p["len"],
                     "preview": p["text"][:40]}
                    for p in summary["top3"]
                ],
            },
        })

    # 3. 저장
    editor.save(PPTX_OUT)
    editor.cleanup()
    print(f"[save] {PPTX_OUT}", flush=True)
    return {"plan": plan, "edits": edit_results, "pptx": str(PPTX_OUT)}


# ----------------------------------------------------------------------------
# PNG 렌더
# ----------------------------------------------------------------------------

def render_pngs(pptx_path: Path, png_dir: Path) -> list[Path]:
    import pythoncom
    import win32com.client

    png_dir.mkdir(parents=True, exist_ok=True)
    # 기존 PNG 삭제 (재실행 시 stale 방지)
    for old in png_dir.glob("*.png"):
        old.unlink()

    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(
            str(pptx_path.resolve()),
            ReadOnly=True,
            Untitled=False,
            WithWindow=False,
        )
        total = presentation.Slides.Count
        print(f"[render] {total} slides", flush=True)
        out = []
        for i in range(1, total + 1):
            p = png_dir / f"step_{i:02d}.png"
            try:
                presentation.Slides(i).Export(str(p), "PNG", 1568, 1176)
                out.append(p)
            except Exception as e:
                print(f"  [err] slide {i}: {e}", flush=True)
        return out
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                powerpoint.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


# ----------------------------------------------------------------------------
# Main
# ----------------------------------------------------------------------------

def main() -> None:
    print("=" * 70)
    print("Mode A 실측 시뮬레이션: D철강 2030 넷제로 전환 로드맵 10장")
    print("=" * 70)

    labels = load_catalog()
    print(f"[catalog] {len(labels)} slides loaded")

    # role 분포
    role_counter = Counter()
    for lab in labels:
        for r in lab.get("narrative_role", []):
            role_counter[r] += 1
    print("[role distribution]")
    for role in NARRATIVE_SEQUENCE:
        print(f"  {role:>14}: {role_counter.get(role, 0):4d}")

    # 1. 검색
    plan = select_deck(labels)
    print()
    print("=" * 70)
    print("[plan] 10장 선택 결과")
    print("=" * 70)
    n_role = sum(1 for p in plan if p["source"] == "role")
    n_arch = sum(1 for p in plan if p["source"] == "archetype")
    n_none = sum(1 for p in plan if p["source"] == "none")
    print(f"role hit: {n_role}/10  archetype fallback: {n_arch}/10  miss: {n_none}/10")
    print()
    for i, p in enumerate(plan, 1):
        idx_str = f"{p['slide_index']:4d}" if p["slide_index"] is not None else "----"
        print(
            f"  {i:2d}. {p['role']:>14}  src#{idx_str}  "
            f"src={p['source']:<10}  archetype={p.get('archetype', [])}  "
            f"conf={p.get('confidence', 0):.2f}"
        )

    # 2. 빌드
    t0 = time.time()
    result = build_pptx(plan)
    dt = time.time() - t0
    print(f"[build done] {dt:.1f}s")

    # 3. PNG 렌더
    t0 = time.time()
    pngs = render_pngs(PPTX_OUT, PNG_DIR)
    dt = time.time() - t0
    print(f"[render done] {dt:.1f}s, {len(pngs)} PNGs")

    # 4. JSON 보고서 저장
    report = {
        "plan": plan,
        "edits": result["edits"],
        "pngs": [str(p) for p in pngs],
        "metrics": {
            "role_hit": n_role,
            "archetype_fallback": n_arch,
            "miss": n_none,
            "edit_ok": sum(1 for e in result["edits"] if e["edit_ok"]),
        },
    }
    report_path = OUTPUT_DIR / "netzero_modeA_report.json"
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)
    print(f"[report] {report_path}")


if __name__ == "__main__":
    main()
