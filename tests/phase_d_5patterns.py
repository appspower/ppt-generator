"""Phase D — 같은 주제(Palantir SAP 전환)를 5가지 패턴으로 자동 생성.

목적: Layer 2 + Phase A/B/C가 진짜로 "매번 다른 화면 구성"을 만들 수
있는지 검증한다. 동일한 주제의 정보를 5개 패턴으로 다르게 표현.

각 슬라이드는:
1. 패턴 함수 호출 → 빌드
2. validate_visual()로 시각 검증 (Layer 1+3)
3. design_check.inspect_design()으로 디자인 점검 (Phase B)
4. PDF 변환

5장 모두 통과해야 Phase D 성공.
"""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.patterns import (
    SlideFooter,
    SlideHeader,
    ComparisonSpec,
    ExecutiveSpec,
    ProcessSpec,
    QuadrantSpec,
    TimelineSpec,
    comparison_matrix,
    executive_summary,
    process_flow,
    quadrant_story,
    timeline_phases,
)
from ppt_builder.design_check import inspect_design
from ppt_builder.visual_validate import validate_visual


# ============================================================
# 공통 footer
# ============================================================

FOOTER = SlideFooter(
    source="출처: Palantir AIP ERP Migration Suite, Unit8 Case Study, SAPPHIRE 2025",
    right="HD현대",
)


# ============================================================
# 1. Executive Summary
# ============================================================

def build_executive():
    spec = ExecutiveSpec(
        header=SlideHeader(
            title="Palantir 단일 플랫폼으로 SAP 전환 일정 14% · 테스트 70% · DT 50% 단축 효과를 확보",
            category="1. Palantir 활용안 — Executive Summary",
            nav_path=["1. 제안 개요", "2. 활용 전략"],
        ),
        hero_label="WHY NOW",
        hero_headline="SAP 전환의 3대 병목을\n단일 Ontology로 해소",
        hero_subtitle="테스트·Cutover·거버넌스를 단일 플랫폼으로 통합 관리",
        bottlenecks=[
            {
                "num": "01", "title": "테스트 자동화",
                "kpi": "공수 70%↓ · 정확도 99.8%",
                "bullets": [
                    "AIP가 Blueprint→비즈니스 규칙 자동 추출",
                    "LLM-as-Judge 0~10점 연속 검증, Jira 자동 등록",
                ],
            },
            {
                "num": "02", "title": "Cutover 오케스트레이션",
                "kpi": "초과율 10%↓ · DT 50%↓",
                "bullets": [
                    "Ontology Task DAG (MRP/잔액이관/마스터동결)",
                    "Workshop Critical Path 라이브 + 리허설 자동 비교",
                ],
            },
            {
                "num": "03", "title": "프로젝트 거버넌스",
                "kpi": "보고 실시간 · 일정 14%↓ (2.5개월)",
                "bullets": [
                    "Health 대시보드 + AIP Agent 리스크 질의",
                    "Readiness Scorecard AI 자동 산출",
                ],
            },
        ],
        kpis=[
            {"value": "14%", "label": "전체 일정 단축", "detail": "18 → 15.5개월"},
            {"value": "70%", "label": "테스트 공수 절감", "detail": "수작업 → AIP"},
            {"value": "50%", "label": "Cutover DT 단축", "detail": "초과율 → 10%↓"},
            {"value": "2~3주", "label": "Quick Win 입증", "detail": "L1 즉시 ROI"},
        ],
        roadmap_phases=[
            {"tag": "L1", "name": "가시화", "duration": "2~3주",
             "deliverables": ["Health 대시보드", "Config Register", "보고 자동화"]},
            {"tag": "L2", "name": "자동화", "duration": "4~6주",
             "deliverables": ["AIP 테스트 생성", "FDD 초안", "결함 Triage"]},
            {"tag": "L3", "name": "최적화", "duration": "8~12주",
             "deliverables": ["Cutover 앱", "Process Mining", "Readiness Scorecard"]},
            {"tag": "L4", "name": "지능화", "duration": "9개월+",
             "deliverables": ["AI Go/No-Go", "Hypercare 탐지", "운영 안정화"]},
        ],
        takeaway="Quick Win L1(2~3주)으로 가치 입증 → L2 Build → L3 Test → L4 Go-Live 점진 확대 — 단일 플랫폼으로 3대 리스크 통합 관리",
        footer=FOOTER,
    )
    return spec, executive_summary


# ============================================================
# 2. Timeline Phases
# ============================================================

def build_timeline():
    spec = TimelineSpec(
        header=SlideHeader(
            title="Quick Win에서 AI 지능화까지 4단계 의존성 체인으로 점진 확대를 추진",
            category="2. 도입 로드맵 — PMO View",
            nav_path=["2. 활용 전략", "3. 단계별 추진 계획"],
        ),
        intro="각 단계 산출물이 다음 단계의 입력 — Level 1 미완료 시 Level 2 이후 착수 불가 (의존성 체인)",
        phases=[
            {
                "tag": "L1", "name": "가시화", "duration": "착수 즉시 (2~3주)",
                "objective": "프로젝트 현황을 실시간으로 가시화하여 PMO 판단 속도 향상",
                "deliverables": [
                    "Foundry Health Dashboard",
                    "Config Decision Register",
                    "Weekly Report Auto-Gen",
                ],
                "metrics": "PMO 보고 3일 → 실시간",
                "prerequisites": "Jira/ADO API 접근 권한 발급",
                "gate": "대시보드 갱신 ≤5분 + 경영진 승인",
                "team": "Palantir FE 1명 + PwC PMO 1명",
            },
            {
                "tag": "L2", "name": "자동화", "duration": "Build (4~6주)",
                "objective": "테스트케이스/FDD 작성을 AIP로 자동화하여 공수 대폭 절감",
                "deliverables": [
                    "AIP Test Case Generator",
                    "FDD Auto-Draft (5분/건)",
                    "Defect Triage Cluster",
                ],
                "metrics": "테스트 작성 60~70%↓",
                "prerequisites": "L1 Ontology 완료 + Blueprint 확보",
                "gate": "커버리지 ≥80% + 정확도 ≥95%",
                "team": "Palantir BE 2명 + PwC 리드 1명",
            },
            {
                "tag": "L3", "name": "최적화", "duration": "Test (8~12주)",
                "objective": "Cutover 리허설 + 프로세스 마이닝으로 Go-Live 리스크 정량화",
                "deliverables": [
                    "Cutover Orchestration App",
                    "Process Mining Gap Report",
                    "Readiness Scorecard",
                ],
                "metrics": "초과율 10%↓ DT 50%↓",
                "prerequisites": "L2 테스트 앱 안정화 + Cutover 시트",
                "gate": "Mock Cutover 3회 + 초과율 ≤15%",
                "team": "Palantir BE+DE 2명 + PwC Basis",
            },
            {
                "tag": "L4", "name": "지능화", "duration": "Go-Live+ (9개월~)",
                "objective": "운영 단계에서 AI Agent로 이상 탐지 및 의사결정 지원",
                "deliverables": [
                    "AI Go/No-Go Engine",
                    "Hypercare Anomaly Detector",
                    "Ops Stability Dashboard",
                ],
                "metrics": "Hypercare KPI 7일 SLA",
                "prerequisites": "L1~3 전체 데이터 축적 완료",
                "gate": "KPI 7일 SLA + Critical 0건",
                "team": "Palantir 1명 + HD현대 IT Ops 2명",
            },
        ],
        takeaway="Level 1이 Quick Win이자 전체 기반 — 2~3주 내 가치 입증 실패 시 전략 재검토, 성공 시 자연스러운 확대 동력 확보",
        footer=FOOTER,
    )
    return spec, timeline_phases


# ============================================================
# 3. Comparison Matrix
# ============================================================

def build_comparison():
    spec = ComparisonSpec(
        header=SlideHeader(
            title="Palantir는 Ontology+AIP+Workshop 단일 플랫폼으로 3사 도구를 통합 — 차별적 우위 확보",
            category="3. Tooling 비교 — IT View",
            nav_path=["2. 활용 전략", "4. 기술 스택 비교"],
        ),
        intro="SAP 전환 핵심 영역에서 시장 대표 도구와 Palantir Foundry+AIP의 역할을 정량 비교",
        criteria_labels=[
            "테스트 자동화",
            "프로세스 마이닝",
            "Cutover 관리",
            "거버넌스",
            "AI/LLM 통합",
            "단일 플랫폼",
        ],
        options=[
            {
                "name": "Tricentis", "summary": "테스트 실행 전문",
                "criteria": ["✓ 강함", "✗ 없음", "✗ 없음", "△ 부분", "△ 부분", "✗"],
                "highlight": False,
            },
            {
                "name": "Celonis", "summary": "프로세스 마이닝 전문",
                "criteria": ["✗ 없음", "✓ 강함", "✗ 없음", "△ 부분", "△ 부분", "✗"],
                "highlight": False,
            },
            {
                "name": "SAP Signavio", "summary": "프로세스 모델링",
                "criteria": ["✗ 없음", "△ 모델 기반", "✗ 없음", "△ 부분", "✗ 없음", "✗"],
                "highlight": False,
            },
            {
                "name": "Palantir", "summary": "Ontology+AIP+Workshop",
                "criteria": ["✓ AIP", "✓ Foundry", "✓ Ontology", "✓ Agent", "✓ 네이티브", "✓"],
                "highlight": True,
            },
        ],
        takeaway="Palantir = 3사 허브 포지셔닝 — Tricentis/Celonis는 보완 도구로 활용, Palantir가 통합 오케스트레이션",
        footer=FOOTER,
    )
    return spec, comparison_matrix


# ============================================================
# 4. Process Flow
# ============================================================

def build_process():
    spec = ProcessSpec(
        header=SlideHeader(
            title="Blueprint→결함 등록 5단계 자동 파이프라인으로 테스트 작성 공수 70% 절감을 달성",
            category="4. 테스트 자동화 흐름 — Engineer View",
            nav_path=["2. 활용 전략", "5. 모듈별 상세 설계"],
        ),
        intro="기존 수작업 3~6개월을 5단계 자동 파이프라인으로 압축 — 각 단계별 도구와 산출물",
        steps=[
            {
                "name": "Blueprint 수집",
                "actor": "PwC 컨설턴트",
                "tools": "PDF/Word\nFoundry Pipeline",
                "output": "구조화된 Blueprint Corpus",
                "duration": "3~5일",
                "prerequisites": "Blueprint 문서 공급사에서 수령 완료",
                "risks": "문서 품질 낮으면 규칙 추출 정확도↓",
            },
            {
                "name": "규칙 추출",
                "actor": "AIP Logic",
                "tools": "Interpretation AI\nLLM Chain",
                "output": "Rule Registry (Ontology 객체)",
                "duration": "2~3일",
                "prerequisites": "Blueprint Corpus 구조화 완료",
                "metrics": "Precision ≥95% 달성 목표",
            },
            {
                "name": "TC 자동 생성",
                "actor": "AIP Logic",
                "tools": "T-Code Mapping\nLLM Generator",
                "output": "TestCase 객체 (입력/기대출력 쌍)",
                "duration": "1~2일",
                "prerequisites": "Rule Registry + T-Code 매핑 테이블",
                "metrics": "TC 커버리지 ≥80% 목표",
            },
            {
                "name": "연속 검증",
                "actor": "LLM-as-Judge",
                "tools": "Rubric 0~10\nκ≥0.85",
                "output": "Pass/Fail + 신뢰도 점수",
                "duration": "수 시간",
                "prerequisites": "TC 객체 100건↑ 생성 완료",
                "example": "에너지社: 20K+ 레코드 2주/$4K",
            },
            {
                "name": "결함 등록",
                "actor": "Foundry Connector",
                "tools": "Jira REST API\nWebhook",
                "output": "Jira 이슈 + 객체 링크",
                "duration": "실시간",
                "prerequisites": "Jira 프로젝트 + API 토큰 발급",
                "risks": "Jira 권한 미확보 시 수동 등록 필요",
            },
        ],
        takeaway="총 1주 이내 전체 파이프라인 가동 — 수작업 3~6개월 대비 60~70%↓, 정확도 99.8% 유지",
        footer=FOOTER,
    )
    return spec, process_flow


# ============================================================
# 5. Quadrant Story
# ============================================================

def build_quadrant():
    spec = QuadrantSpec(
        header=SlideHeader(
            title="8개 Palantir 활용 모듈을 ROI×난이도로 분류, Quick Win 영역에서 시작을 권장",
            category="5. 모듈 우선순위 — Investment View",
            nav_path=["2. 활용 전략", "6. 투자 우선순위"],
        ),
        intro="단기 ROI와 구현 난이도 두 축으로 8개 활용 모듈을 분류 — Quick Win 영역에서 시작 권장",
        x_axis_label="구현 난이도",
        y_axis_label="단기 ROI",
        x_low="LOW", x_high="HIGH",
        y_low="LOW", y_high="HIGH",
        quadrants=[
            {  # TL: 고ROI / 저난이도 = QUICK WIN
                "title": "QUICK WIN",
                "highlight": True,
                "items": [
                    "Health 대시보드 (L1)",
                    "Config Register (L1)",
                    "주간 보고 자동화 (L1)",
                ],
                "description": "2~3주 내 즉시 가치 입증 — Jira REST 연동으로 PMO 보고 실시간화",
                "action": "킥오프 직후 Palantir 1명 + PwC PMO 1명 투입하여 착수",
                "metrics": "PMO 주 3~4h 절감, 보고 3일→실시간",
            },
            {  # TR: 고ROI / 고난이도 = STRATEGIC BET
                "title": "STRATEGIC BET",
                "highlight": False,
                "items": [
                    "AIP 테스트 자동화 (L2)",
                    "Cutover 오케스트레이션 (L3)",
                ],
                "description": "L1 성공 후 확대 — Blueprint 품질과 리허설 계획이 전제",
                "action": "L1 Gate 통과 후 Palantir 2명 추가 투입",
                "metrics": "테스트 60~70%↓, Cutover 초과율 10%↓",
            },
            {  # BL: 저ROI / 저난이도 = NICE TO HAVE
                "title": "NICE TO HAVE",
                "highlight": False,
                "items": [
                    "FDD 자동 초안 (L2)",
                    "결함 Triage 클러스터 (L2)",
                ],
                "description": "L2 진행 시 함께 구축 가능 — 독립 투자 대비 효과 낮음",
                "action": "L2 테스트 자동화와 번들로 진행 권장",
            },
            {  # BR: 저ROI / 고난이도 = DEFER
                "title": "DEFER",
                "highlight": False,
                "items": [
                    "AI Go/No-Go (L4)",
                    "Hypercare 이상 탐지 (L4)",
                ],
                "description": "L1~3 전 데이터 축적 후에야 의미 — 현 단계 착수 불가",
                "action": "Go-Live 3개월 전 재평가하여 착수 여부 결정",
            },
        ],
        insight="L1 Quick Win 3개 모듈을 즉시 착수 → 가치 입증 후 L2 Strategic Bet 단계적 확대 — DEFER는 L1~3 데이터 축적 후 재평가",
        footer=FOOTER,
    )
    return spec, quadrant_story


# ============================================================
# 메인 — 5장 빌드 + 검증
# ============================================================

def build_one(builder_fn, output: Path):
    """builder_fn() = (spec, pattern_func) 튜플 반환."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    spec, pattern_func = builder_fn()
    pattern_func(slide, spec)

    prs.save(output)
    return output


def main():
    out_dir = Path("output/phase_d")
    out_dir.mkdir(parents=True, exist_ok=True)

    cases = [
        ("01_executive", build_executive, "executive"),
        ("02_timeline", build_timeline, "timeline"),
        ("03_comparison", build_comparison, "comparison"),
        ("04_process", build_process, "process"),
        ("05_quadrant", build_quadrant, "quadrant"),
    ]

    print("=" * 70)
    print("Phase D — 같은 주제(Palantir SAP 전환) × 5 패턴")
    print("=" * 70)

    all_passed = True
    for name, builder_fn, kind in cases:
        out = out_dir / f"{name}.pptx"
        try:
            build_one(builder_fn, out)
        except Exception as e:
            print(f"\n[{name}] BUILD FAILED: {e}")
            all_passed = False
            continue

        # Layer 3: 시각 검증 (정적만 — PDF는 마지막에 일괄)
        visual = validate_visual(out, convert_pdf=False)
        # Phase B: 디자인 점검 (패턴별 임계값)
        design = inspect_design(str(out), pattern_kind=kind)

        ok = not visual.issues and design.passed
        status = "PASS" if ok else "FAIL"
        if not ok:
            all_passed = False

        print(f"\n[{name}] {status}")
        print(f"  visual: {len(visual.issues)} issues")
        for i in visual.issues[:3]:
            print(f"    - {i}")
        print(f"  design: passed={design.passed}, issues={len(design.issues)}")
        for i in design.issues[:3]:
            print(f"    - {i}")
        print(f"  metrics: {dict((k, v) for k, v in design.metrics.items() if 'density' in k or 'top_color' in k)}")

    # PDF 변환 (마지막에 일괄)
    print("\n" + "=" * 70)
    print("PDF 변환 중...")
    print("=" * 70)
    for name, _, _ in cases:
        out = out_dir / f"{name}.pptx"
        if not out.exists():
            continue
        try:
            visual = validate_visual(out, convert_pdf=True)
            print(f"  {name}.pdf  — {'OK' if visual.pdf_available else 'SKIP'}")
        except Exception as e:
            print(f"  {name}.pdf  — FAIL ({e})")

    print("\n" + "=" * 70)
    print(f"Phase D 결과: {'ALL PASSED' if all_passed else 'SOME FAILED'}")
    print("=" * 70)
    return 0 if all_passed else 1


if __name__ == "__main__":
    sys.exit(main())
