"""Layer 2 프로토타입 v1 — 절대좌표 디자인 자유도 검증.

같은 주제(Palantir SAP 전환)를 R최종(main_r7)과 완전히 다른 레이아웃으로 만든다.
컴포넌트 프리셋 없이, 백지 위에 절대좌표로 직접 배치한다.

R최종 main_r7과의 차이:
- 3개 카드 균등 배치 → 비대칭 hero(좌측 큰 KPI) + 우측 누적 stack
- 둥근 카드 → 굵은 각진 사각형 + heavy border
- 색상 절제 → strong contrast (검정 hero + 흰 본문)
- 중앙정렬 takeaway → 좌측 stripe + 큰 숫자
"""

import sys
from pathlib import Path

# 프로젝트 루트를 sys.path에 추가
_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas
from ppt_builder.visual_validate import validate_visual


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    c = Canvas(slide)

    # ----------- 헤더 영역 -----------
    c.title(
        "Palantir 투입으로 SAP 전환 일정 14% · 테스트 70% · DT 50% 단축",
        x=0.3, y=0.2, w=9.4, h=0.45, size=15,
        underline_accent=True,
    )
    # breadcrumb
    c.text(
        "① Palantir 활용안 — Executive Summary",
        x=0.3, y=0.75, w=9.4, h=0.25,
        size=9, color="grey", align="left",
    )

    # ----------- 본문: 비대칭 hero + 누적 stack -----------
    #
    # ┌──────────────────┬───────────────────────┐
    # │                  │  ▎ 14%   일정 단축    │
    # │   HERO BOX       │  ─────────────────    │
    # │   "WHY NOW"      │  ▎ 70%   테스트 공수↓│
    # │   (검정 배경)     │  ─────────────────    │
    # │                  │  ▎ 50%   다운타임↓   │
    # └──────────────────┴───────────────────────┘

    HERO_X, HERO_Y = 0.3, 1.15
    HERO_W, HERO_H = 4.3, 5.0

    # Hero 박스 — 검정 배경, 흰 글씨, heavy border
    c.box(
        x=HERO_X, y=HERO_Y, w=HERO_W, h=HERO_H,
        fill="dark", border=None,
    )
    # 좌측 굵은 오렌지 stripe (시그니처)
    c.box(
        x=HERO_X, y=HERO_Y, w=0.18, h=HERO_H,
        fill="accent", border=None,
    )
    # Hero 라벨 칩
    c.label_chip(
        "WHY NOW",
        x=HERO_X + 0.4, y=HERO_Y + 0.35, w=1.5, h=0.32,
        fill="accent", text_color="white",
    )
    # Hero 큰 제목
    c.text(
        "SAP 전환의\n3대 병목을\nPalantir 단일\n플랫폼으로 해소",
        x=HERO_X + 0.4, y=HERO_Y + 0.85, w=HERO_W - 0.55, h=2.5,
        size=22, bold=True, color="white", anchor="top",
    )
    # Hero 부연
    c.text(
        "테스트 자동화 · Cutover 오케스트레이션 · 거버넌스를\n"
        "단일 Ontology로 통합 — Quick Win(2~3주)부터 점진 확대",
        x=HERO_X + 0.4, y=HERO_Y + 3.6, w=HERO_W - 0.55, h=1.2,
        size=10, color="grey_light", anchor="top",
    )

    # ----------- 우측 누적 KPI stack (3개) -----------
    STACK_X = 4.85
    STACK_W = 4.85
    KPI_H = 1.55
    GAP = 0.18
    KPI_Y_START = HERO_Y

    kpis = [
        ("14%", "전체 일정 단축", "18개월 → 15.5개월 (약 2.5개월)"),
        ("70%", "테스트 작성 공수 절감", "수작업 3~6개월 → 수 일, 정확도 99.8%"),
        ("50%", "Cutover 다운타임 감소", "초과율 30~40% → 10% 이하, DT 50%↓"),
    ]
    for i, (val, label, detail) in enumerate(kpis):
        ky = KPI_Y_START + i * (KPI_H + GAP)
        # 박스 (흰색 + 굵은 검정 테두리 + 좌측 오렌지 stripe)
        c.box(
            x=STACK_X, y=ky, w=STACK_W, h=KPI_H,
            fill="white", border=2.0, border_color="black",
        )
        c.box(
            x=STACK_X, y=ky, w=0.12, h=KPI_H,
            fill="accent", border=None,
        )
        # 큰 숫자
        c.text(
            val,
            x=STACK_X + 0.25, y=ky + 0.18, w=1.6, h=1.1,
            size=36, bold=True, color="accent",
            font="Georgia", anchor="middle",
        )
        # 라벨
        c.text(
            label,
            x=STACK_X + 1.95, y=ky + 0.22, w=STACK_W - 2.1, h=0.45,
            size=12, bold=True, color="black", anchor="top",
        )
        # 디테일
        c.text(
            detail,
            x=STACK_X + 1.95, y=ky + 0.7, w=STACK_W - 2.1, h=0.7,
            size=9, color="grey", anchor="top",
        )
        # 우측 끝 작은 인덱스
        c.text(
            f"0{i+1}",
            x=STACK_X + STACK_W - 0.5, y=ky + 0.1, w=0.4, h=0.3,
            size=10, bold=True, color="grey_mid", align="right",
        )

    # ----------- 하단 takeaway 바 -----------
    BAR_Y = 6.55
    c.box(
        x=0.3, y=BAR_Y, w=9.4, h=0.45,
        fill="dark", border=None,
    )
    c.box(
        x=0.3, y=BAR_Y, w=0.12, h=0.45,
        fill="accent", border=None,
    )
    c.text(
        "Quick Win 거버넌스(2~3주) → 테스트 자동화(4~6주) → Cutover 최적화(8~12주) — "
        "점진 투입으로 리스크 최소화",
        x=0.5, y=BAR_Y, w=9.1, h=0.45,
        size=10, bold=True, color="white", anchor="middle",
    )

    # ----------- 푸터 -----------
    c.divider_h(x=0.3, y=7.1, w=9.4, color="border", width=0.75)
    c.text(
        "Strictly Private and Confidential",
        x=0.3, y=7.18, w=3.5, h=0.18,
        size=7, color="grey",
    )
    c.text(
        "출처: Palantir AIP ERP Migration Suite, Unit8 Case Study, SAPPHIRE 2025",
        x=3.0, y=7.18, w=5.5, h=0.18,
        size=7, color="grey",
    )
    c.text(
        "pwc",
        x=0.3, y=7.32, w=0.8, h=0.15,
        size=7, bold=True, color="accent",
    )
    c.text(
        "HD현대",
        x=8.8, y=7.32, w=1.0, h=0.15,
        size=7, bold=True, color="black", align="right",
    )

    return prs


def main():
    out_dir = Path("output/proto")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "layer2_proto_v1.pptx"

    prs = build()
    prs.save(out)
    print(f"Saved: {out}")

    # 자동 시각 검증
    report = validate_visual(out, convert_pdf=True)
    print(f"\nIssues: {len(report.issues)}")
    for iss in report.issues:
        print(f"  - {iss}")
    print(f"\nPDF: {report.pdf_path}")
    print(f"Severity: {report.severity_count()}")


if __name__ == "__main__":
    main()
