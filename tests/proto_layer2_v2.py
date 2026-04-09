"""Layer 2 프로토타입 v2 — 사용자 피드백 반영판.

피드백 → 적용:
1. 두께: 0.75pt 보통, 굵은 라인 제거
2. 색상: 진한 오렌지 제외, 회색 위계 (Dark/Black/Grey/Grey-mid/Grey-light)
3. 밀도: Hero 빈공간 → 3대 병목 mini-list로 채움
4. 우측: KPI 3 stack → KPI 2×2 콤팩트 + 그 아래 mini-timeline + Quick Win 콜아웃
"""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas
from ppt_builder.visual_validate import validate_visual


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    c = Canvas(slide)

    # ====================================================
    # 헤더
    # ====================================================
    c.title(
        "Palantir 투입으로 SAP 전환 일정 14% · 테스트 70% · DT 50% 단축",
        x=0.3, y=0.2, w=9.4, h=0.45, size=15,
        underline=True, underline_color="dark",
    )
    c.text(
        "① Palantir 활용안 — Executive Summary",
        x=0.3, y=0.75, w=9.4, h=0.25,
        size=9, color="grey", align="left",
    )

    # ====================================================
    # 본문 좌측 — Hero (밀도 채움)
    # ====================================================
    HERO_X, HERO_Y = 0.3, 1.15
    HERO_W, HERO_H = 4.3, 5.15

    # 다크 배경
    c.box(
        x=HERO_X, y=HERO_Y, w=HERO_W, h=HERO_H,
        fill="dark", border=None,
    )

    # WHY NOW 칩 (좌상단)
    c.label_chip(
        "WHY NOW",
        x=HERO_X + 0.3, y=HERO_Y + 0.3, w=1.3, h=0.28,
        fill="grey", text_color="white",
    )

    # 큰 헤드라인 (3줄로 압축)
    c.text(
        "SAP 전환의 3대 병목을\n단일 Ontology로 해소",
        x=HERO_X + 0.3, y=HERO_Y + 0.72, w=HERO_W - 0.55, h=1.2,
        size=20, bold=True, color="white", anchor="top",
    )

    # ─── 가는 구분선 ───
    c.box(
        x=HERO_X + 0.3, y=HERO_Y + 1.95, w=HERO_W - 0.6, h=0.012,
        fill="grey", border=None,
    )

    # 3대 병목 mini-list (정보로 빈공간 채움)
    bottlenecks = [
        ("01", "테스트 자동화", "Blueprint→AIP Rule→LLM 테스트케이스 생성·검증\n수작업 3~6개월 → 수 일, 정확도 99.8%"),
        ("02", "Cutover 오케스트레이션", "Ontology Task DAG + Critical Path 라이브\n초과율 30~40% → 10%↓, 다운타임 50%↓"),
        ("03", "프로젝트 거버넌스", "Health 대시보드 + AIP Agent 리스크 질의\nPMO 보고 실시간, 일정 14%↓ (2.5개월)"),
    ]
    item_h = 0.92
    item_start_y = HERO_Y + 2.1
    for i, (num, title, detail) in enumerate(bottlenecks):
        iy = item_start_y + i * item_h
        # 번호
        c.text(
            num,
            x=HERO_X + 0.3, y=iy, w=0.4, h=0.25,
            size=10, bold=True, color="grey_light", anchor="top",
        )
        # 제목
        c.text(
            title,
            x=HERO_X + 0.75, y=iy, w=HERO_W - 1.0, h=0.27,
            size=11, bold=True, color="white", anchor="top",
        )
        # 디테일
        c.text(
            detail,
            x=HERO_X + 0.75, y=iy + 0.28, w=HERO_W - 1.0, h=0.6,
            size=8, color="grey_light", anchor="top",
        )

    # ====================================================
    # 본문 우측 — KPI 2×2 + mini-timeline
    # ====================================================
    R_X = 4.85
    R_W = 4.85
    R_Y = 1.15

    # ─── 우측 영역 라벨 ───
    c.text(
        "정량 효과 (4대 KPI)",
        x=R_X, y=R_Y, w=R_W, h=0.22,
        size=9, bold=True, color="grey", anchor="top",
    )

    # ─── KPI 2×2 ───
    KPI_AREA_Y = R_Y + 0.27
    KPI_W = (R_W - 0.15) / 2  # 가로 2개 + 갭 0.15
    KPI_H = 1.4
    KPI_GAP_X = 0.15
    KPI_GAP_Y = 0.15

    kpis = [
        ("14%", "전체 일정 단축", "18 → 15.5개월\n2.5개월 절감"),
        ("70%", "테스트 공수 절감", "수작업 → AIP\n정확도 99.8%"),
        ("50%", "Cutover DT 단축", "초과율 30~40%\n→ 10% 이하"),
        ("2~3주", "Quick Win 가치 입증", "L1 대시보드\n구현 즉시 ROI"),
    ]
    for i, (val, label, detail) in enumerate(kpis):
        col = i % 2
        row = i // 2
        kx = R_X + col * (KPI_W + KPI_GAP_X)
        ky = KPI_AREA_Y + row * (KPI_H + KPI_GAP_Y)

        # 흰 박스 + 회색 보통 테두리
        c.box(
            x=kx, y=ky, w=KPI_W, h=KPI_H,
            fill="white", border=0.75, border_color="grey_mid",
        )
        # 좌측 미세 dark stripe (시그니처, 가늘게)
        c.box(
            x=kx, y=ky, w=0.06, h=KPI_H,
            fill="dark", border=None,
        )
        # 큰 숫자 (검정, 오렌지 X)
        c.text(
            val,
            x=kx + 0.18, y=ky + 0.12, w=KPI_W - 0.3, h=0.55,
            size=24 if len(val) <= 3 else 20,
            bold=True, color="black",
            font="Georgia", anchor="top",
        )
        # 라벨
        c.text(
            label,
            x=kx + 0.18, y=ky + 0.7, w=KPI_W - 0.3, h=0.28,
            size=9, bold=True, color="black", anchor="top",
        )
        # 디테일
        c.text(
            detail,
            x=kx + 0.18, y=ky + 0.97, w=KPI_W - 0.3, h=0.4,
            size=7, color="grey", anchor="top",
        )

    # ─── mini-timeline (4단계 chevron 대신 가는 단계 바) ───
    TL_Y = KPI_AREA_Y + 2 * KPI_H + KPI_GAP_Y + 0.18

    c.text(
        "도입 로드맵 (점진 확대)",
        x=R_X, y=TL_Y, w=R_W, h=0.22,
        size=9, bold=True, color="grey", anchor="top",
    )

    TL_BAR_Y = TL_Y + 0.3
    TL_BAR_H = 0.5
    n_phases = 4
    phase_w = (R_W - 0.06 * (n_phases - 1)) / n_phases
    phase_gap = 0.06
    phases = [
        ("L1", "가시화", "2~3주"),
        ("L2", "자동화", "4~6주"),
        ("L3", "최적화", "8~12주"),
        ("L4", "지능화", "9개월+"),
    ]
    fill_intensity = ["dark", "grey", "grey_mid", "grey_light"]
    text_color = ["white", "white", "black", "black"]
    for i, ((tag, name, dur), fill_c, tc) in enumerate(
        zip(phases, fill_intensity, text_color)
    ):
        px = R_X + i * (phase_w + phase_gap)
        c.box(
            x=px, y=TL_BAR_Y, w=phase_w, h=TL_BAR_H,
            fill=fill_c, border=None,
        )
        # 단계 태그
        c.text(
            tag,
            x=px + 0.05, y=TL_BAR_Y + 0.04, w=phase_w - 0.1, h=0.18,
            size=8, bold=True, color=tc, anchor="top",
        )
        # 단계명
        c.text(
            name,
            x=px + 0.05, y=TL_BAR_Y + 0.21, w=phase_w - 0.1, h=0.18,
            size=9, bold=True, color=tc, anchor="top",
        )
        # 기간
        c.text(
            dur,
            x=px + 0.05, y=TL_BAR_Y + 0.36, w=phase_w - 0.1, h=0.14,
            size=7, color=tc, anchor="top",
        )

    # ====================================================
    # 하단 takeaway 바
    # ====================================================
    BAR_Y = 6.55
    c.box(
        x=0.3, y=BAR_Y, w=9.4, h=0.42,
        fill="dark", border=None,
    )
    c.text(
        "Quick Win L1(2~3주)으로 가치 입증 → L2 Build → L3 Test → L4 Go-Live 점진 확대 — "
        "단일 플랫폼으로 3대 리스크 통합 관리",
        x=0.5, y=BAR_Y, w=9.1, h=0.42,
        size=10, bold=True, color="white", anchor="middle",
    )

    # ====================================================
    # 푸터
    # ====================================================
    c.divider_h(x=0.3, y=7.1, w=9.4, color="border", width=0.5)
    c.text(
        "Strictly Private and Confidential",
        x=0.3, y=7.18, w=3.5, h=0.18,
        size=7, color="grey",
    )
    c.text(
        "출처: Palantir AIP ERP Migration Suite, Unit8 Case Study, SAPPHIRE 2025",
        x=3.0, y=7.18, w=5.5, h=0.18,
        size=7, color="grey",
    )
    c.text(
        "pwc",
        x=0.3, y=7.32, w=0.8, h=0.15,
        size=7, bold=True, color="dark",
    )
    c.text(
        "HD현대",
        x=8.8, y=7.32, w=1.0, h=0.15,
        size=7, bold=True, color="black", align="right",
    )

    return prs


def main():
    out_dir = Path("output/proto")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "layer2_proto_v2.pptx"

    prs = build()
    prs.save(out)
    print(f"Saved: {out}")

    report = validate_visual(out, convert_pdf=True)
    print(f"\nIssues: {len(report.issues)}")
    for iss in report.issues:
        print(f"  - {iss}")
    print(f"\nPDF: {report.pdf_path}")
    print(f"Severity: {report.severity_count()}")


if __name__ == "__main__":
    main()
