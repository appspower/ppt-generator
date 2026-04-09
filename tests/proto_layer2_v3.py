"""Layer 2 프로토타입 v3 — 사용자 디테일 피드백 반영판.

피드백 → 적용:
1. 완전 블랙 → 미디엄 다크 그레이 (grey_800, grey_700) 위계
2. 콘텐츠 양 부족 → 각 항목에 sub-bullet 2개씩 (▪)
3. 번호 01/02/03 → 원형 OVAL 안에
4. 박스 안 박스 조합 → heading 박스 + sub-list 박스 분리
5. L1/L2 박스 → CHEVRON 화살표로 (시퀀스 시각화)
6. 도입 로드맵 아래 빈 공간 → 각 단계 산출물 mini-list 추가
7. WHY NOW → 회색 톤 (grey_400)
"""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas
from ppt_builder.visual_validate import validate_visual


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    c = Canvas(slide)

    # ====================================================
    # 헤더
    # ====================================================
    c.title(
        "Palantir 투입으로 SAP 전환 일정 14% · 테스트 70% · DT 50% 단축",
        x=0.3, y=0.2, w=9.4, h=0.45, size=15,
        underline=True, underline_color="grey_700",
    )
    c.text(
        "① Palantir 활용안 — Executive Summary",
        x=0.3, y=0.75, w=9.4, h=0.25,
        size=9, color="grey_700", align="left",
    )

    # ====================================================
    # 본문 좌측 — Hero (회색 톤 + 박스 조합 + 원형 번호)
    # ====================================================
    HERO_X, HERO_Y = 0.3, 1.15
    HERO_W, HERO_H = 4.45, 5.15

    # 미디엄 다크 그레이 배경 (완전 블랙 회피)
    c.box(
        x=HERO_X, y=HERO_Y, w=HERO_W, h=HERO_H,
        fill="grey_800", border=None,
    )

    # WHY NOW 칩 (회색 톤)
    c.label_chip(
        "WHY NOW",
        x=HERO_X + 0.3, y=HERO_Y + 0.3, w=1.3, h=0.28,
        fill="grey_400", text_color="white",
    )

    # 헤드라인 (3줄)
    c.text(
        "SAP 전환의 3대 병목을\n단일 Ontology로 해소",
        x=HERO_X + 0.3, y=HERO_Y + 0.7, w=HERO_W - 0.55, h=1.0,
        size=18, bold=True, color="white", anchor="top",
    )
    # 부연 (한 줄)
    c.text(
        "테스트·Cutover·거버넌스를 단일 플랫폼으로 통합 관리",
        x=HERO_X + 0.3, y=HERO_Y + 1.62, w=HERO_W - 0.55, h=0.25,
        size=9, color="grey_200", anchor="top",
    )

    # ─── 가는 구분선 ───
    c.box(
        x=HERO_X + 0.3, y=HERO_Y + 1.95, w=HERO_W - 0.6, h=0.012,
        fill="grey_400", border=None,
    )

    # 3대 병목 — 원형 번호 + heading 박스 + sub-bullet
    bottlenecks = [
        (
            "01", "테스트 자동화",
            ["AIP가 Blueprint→비즈니스 규칙 자동 추출",
             "LLM-as-Judge 0~10점 연속 검증, 결함 Jira 자동 등록"],
            "공수 70%↓ · 정확도 99.8%",
        ),
        (
            "02", "Cutover 오케스트레이션",
            ["Ontology Task DAG (MRP/잔액이관/마스터동결)",
             "Workshop Critical Path 라이브 + 리허설 자동 비교"],
            "초과율 10%↓ · DT 50%↓",
        ),
        (
            "03", "프로젝트 거버넌스",
            ["Health 대시보드 + AIP Agent 리스크 질의",
             "Readiness Scorecard AI 자동 산출, 출처 추적"],
            "보고 실시간 · 일정 14%↓ (2.5개월)",
        ),
    ]
    item_h = 1.05
    item_start_y = HERO_Y + 2.12
    for i, (num, title, bullets, kpi) in enumerate(bottlenecks):
        iy = item_start_y + i * item_h

        # 원형 번호 (좌측)
        c.circle(
            x=HERO_X + 0.3, y=iy + 0.04, d=0.38,
            fill="grey_900",
            border=1.0, border_color="grey_400",
            text=num, text_color="white", text_size=10,
        )
        # 제목 (오른쪽)
        c.text(
            title,
            x=HERO_X + 0.78, y=iy, w=HERO_W - 1.1, h=0.27,
            size=11, bold=True, color="white", anchor="top",
        )
        # KPI 우측 정렬
        c.text(
            kpi,
            x=HERO_X + 0.78, y=iy + 0.27, w=HERO_W - 1.1, h=0.2,
            size=8, bold=True, color="grey_200", anchor="top",
        )
        # Sub-bullets (▪)
        for bi, bul in enumerate(bullets):
            c.text(
                f"▪  {bul}",
                x=HERO_X + 0.78, y=iy + 0.5 + bi * 0.25, w=HERO_W - 1.0, h=0.22,
                size=8, color="grey_200", anchor="top",
            )

    # ====================================================
    # 본문 우측 — KPI 2×2 + Roadmap chevron + 산출물
    # ====================================================
    R_X = 4.95
    R_W = 4.75
    R_Y = 1.15

    # ─── 우측 영역 라벨 ───
    c.text(
        "정량 효과 (4대 KPI)",
        x=R_X, y=R_Y, w=R_W, h=0.22,
        size=10, bold=True, color="grey_800", anchor="top",
    )

    # ─── KPI 2×2 ───
    KPI_AREA_Y = R_Y + 0.3
    KPI_W = (R_W - 0.15) / 2
    KPI_H = 1.25
    KPI_GAP_X = 0.15
    KPI_GAP_Y = 0.15

    kpis = [
        ("14%", "전체 일정 단축", "18 → 15.5개월"),
        ("70%", "테스트 공수 절감", "수작업 → AIP 자동"),
        ("50%", "Cutover DT 단축", "초과율 → 10% 이하"),
        ("2~3주", "Quick Win 가치 입증", "L1 즉시 ROI 확보"),
    ]
    for i, (val, label, detail) in enumerate(kpis):
        col = i % 2
        row = i // 2
        kx = R_X + col * (KPI_W + KPI_GAP_X)
        ky = KPI_AREA_Y + row * (KPI_H + KPI_GAP_Y)

        # 흰 박스 + 회색 보통 테두리
        c.box(
            x=kx, y=ky, w=KPI_W, h=KPI_H,
            fill="white", border=0.75, border_color="grey_mid",
        )
        # 좌측 미세 grey stripe
        c.box(
            x=kx, y=ky, w=0.08, h=KPI_H,
            fill="grey_700", border=None,
        )
        # 큰 숫자 (검정 → grey_900으로 약간 부드럽게)
        c.text(
            val,
            x=kx + 0.2, y=ky + 0.18, w=KPI_W - 0.3, h=0.55,
            size=24 if len(val) <= 3 else 19,
            bold=True, color="grey_900",
            font="Georgia", anchor="middle",
        )
        # 라벨
        c.text(
            label,
            x=kx + 0.2, y=ky + 0.72, w=KPI_W - 0.3, h=0.25,
            size=9, bold=True, color="grey_900", anchor="top",
        )
        # 디테일
        c.text(
            detail,
            x=kx + 0.2, y=ky + 0.96, w=KPI_W - 0.3, h=0.22,
            size=7, color="grey_700", anchor="top",
        )

    # ─── 도입 로드맵 — Chevron 화살표 + 산출물 mini-list ───
    RM_Y = KPI_AREA_Y + 2 * KPI_H + KPI_GAP_Y + 0.22

    c.text(
        "도입 로드맵 (점진 확대)",
        x=R_X, y=RM_Y, w=R_W, h=0.22,
        size=10, bold=True, color="grey_800", anchor="top",
    )

    # Chevron 4개 — 각이 진 화살표 시퀀스
    CHEV_Y = RM_Y + 0.3
    CHEV_H = 0.42
    n = 4
    chev_overlap = 0.08  # 화살표 겹침으로 연결감
    chev_w = (R_W + chev_overlap * (n - 1)) / n
    phases = [
        ("L1", "가시화", "2~3주", "grey_800"),
        ("L2", "자동화", "4~6주", "grey_700"),
        ("L3", "최적화", "8~12주", "grey_400"),
        ("L4", "지능화", "9개월+", "grey_200"),
    ]
    text_colors = ["white", "white", "white", "grey_900"]
    for i, ((tag, name, dur, fill_c), tc) in enumerate(zip(phases, text_colors)):
        cx = R_X + i * (chev_w - chev_overlap)
        c.chevron(
            x=cx, y=CHEV_Y, w=chev_w, h=CHEV_H,
            fill=fill_c, text=f"{tag}  {name}", text_color=tc, text_size=9,
        )

    # 각 단계의 산출물 mini-list (chevron 아래)
    DELIV_Y = CHEV_Y + CHEV_H + 0.12
    DELIV_H = 0.85
    deliverables = [
        ["Health 대시보드", "Config Register", "주간 보고 자동화"],
        ["AIP 테스트 생성", "FDD 자동 초안", "결함 Triage"],
        ["Cutover 앱", "Process Mining", "Readiness Scorecard"],
        ["AI Go/No-Go", "Hypercare 탐지", "운영 안정화"],
    ]
    deliv_w = (R_W - 0.15 * (n - 1)) / n
    for i, items in enumerate(deliverables):
        dx = R_X + i * (deliv_w + 0.15)
        # 박스 (아주 옅은 그레이)
        c.box(
            x=dx, y=DELIV_Y, w=deliv_w, h=DELIV_H,
            fill="grey_100", border=0.5, border_color="grey_mid",
        )
        # 산출물 리스트
        for di, item in enumerate(items):
            c.text(
                f"▪ {item}",
                x=dx + 0.07, y=DELIV_Y + 0.08 + di * 0.23, w=deliv_w - 0.14, h=0.22,
                size=7, color="grey_900", anchor="top",
            )

    # ====================================================
    # 하단 takeaway 바 (회색 톤)
    # ====================================================
    BAR_Y = 6.55
    c.box(
        x=0.3, y=BAR_Y, w=9.4, h=0.42,
        fill="grey_800", border=None,
    )
    c.text(
        "Quick Win L1(2~3주)으로 가치 입증 → L2 Build → L3 Test → L4 Go-Live 점진 확대 — "
        "단일 플랫폼으로 3대 리스크 통합 관리",
        x=0.5, y=BAR_Y, w=9.1, h=0.42,
        size=10, bold=True, color="white", anchor="middle",
    )

    # ====================================================
    # 푸터
    # ====================================================
    c.divider_h(x=0.3, y=7.1, w=9.4, color="border", width=0.5)
    c.text(
        "Strictly Private and Confidential",
        x=0.3, y=7.18, w=3.5, h=0.18,
        size=7, color="grey_700",
    )
    c.text(
        "출처: Palantir AIP ERP Migration Suite, Unit8 Case Study, SAPPHIRE 2025",
        x=3.0, y=7.18, w=5.5, h=0.18,
        size=7, color="grey_700",
    )
    c.text(
        "pwc",
        x=0.3, y=7.32, w=0.8, h=0.15,
        size=7, bold=True, color="grey_900",
    )
    c.text(
        "HD현대",
        x=8.8, y=7.32, w=1.0, h=0.15,
        size=7, bold=True, color="grey_900", align="right",
    )

    return prs


def main():
    out_dir = Path("output/proto")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "layer2_proto_v3.pptx"

    prs = build()
    prs.save(out)
    print(f"Saved: {out}")

    report = validate_visual(out, convert_pdf=True)
    print(f"\nIssues: {len(report.issues)}")
    for iss in report.issues:
        print(f"  - {iss}")
    print(f"\nPDF: {report.pdf_path}")
    print(f"Severity: {report.severity_count()}")


if __name__ == "__main__":
    main()
