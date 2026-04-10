"""품질 상한선 테스트 — 3개 새 주제 × 3개 패턴.

현재 워크플로우의 평준화된 상한선을 확인하기 위한 테스트.
리서치 데이터를 기반으로 최대한 짜임새 있게 구성.
"""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.patterns import (
    SlideFooter,
    SlideHeader,
    ComparisonSpec,
    ProcessSpec,
    QuadrantSpec,
    comparison_matrix,
    process_flow,
    quadrant_story,
)
from ppt_builder.design_check import inspect_design
from ppt_builder.visual_validate import validate_visual


# ============================================================
# 1. AI 도입 성숙도 비교 — comparison_matrix
# ============================================================

def build_ai_maturity():
    spec = ComparisonSpec(
        header=SlideHeader(
            title="반도체가 AI 도입 최선두(82%), 조선은 31%로 최하위 — 격차의 본질은 데이터 구조",
            category="글로벌 제조업 AI 성숙도 — Industry Benchmark",
            nav_path=["1. AI 전략", "3. 산업별 성숙도 비교"],
        ),
        intro="4대 제조 산업의 AI 운영 배포 수준을 6개 차원으로 정량 비교 — 반도체의 압도적 선도와 조선의 구조적 한계",
        criteria_labels=[
            "AI 운영 배포율",
            "주요 유즈케이스",
            "AI 투자 ROI (3년)",
            "AI 인력 비율",
            "데이터 인프라",
            "핵심 병목",
        ],
        options=[
            {
                "name": "자동차",
                "summary": "예지보전 중심 ★★★★",
                "criteria": [
                    "68% (OEM 기준)",
                    "예지보전 (투자의 38%)",
                    "220~310%",
                    "4.2% (1인당 $1,850/년)",
                    "클라우드 71%, 데이터레이크 63%",
                    "OT-IT 통합 (실패 43%)",
                ],
                "highlight": False,
            },
            {
                "name": "반도체",
                "summary": "수율 최적화 ★★★★★",
                "criteria": [
                    "82% (Top 20 팹)",
                    "수율 최적화 (APC 표준화)",
                    "450~600%",
                    "11.3% (TSMC 2,200명+)",
                    "온프레미스 우세, 데이터레이크 89%",
                    "XAI 설명가능성 (도입 28%)",
                ],
                "highlight": True,
            },
            {
                "name": "화학",
                "summary": "공정 최적화 ★★★",
                "criteria": [
                    "51% (Top 50 기준)",
                    "에너지 절감 (BASF 18%↓)",
                    "180~240%",
                    "2.8% (1인당 $2,100/년)",
                    "클라우드 58%, DCS 연동 44%",
                    "안전 규제 (재인증 14개월)",
                ],
                "highlight": False,
            },
            {
                "name": "조선",
                "summary": "자동화 초기 ★★",
                "criteria": [
                    "31% (한국 Big 3 기준)",
                    "용접 로봇 + AI 설계",
                    "280~350% (운항 기준)",
                    "1.4% (연 500명 채용 목표)",
                    "클라우드 38%, 디지털트윈 26%",
                    "ETO 구조 (데이터 단절)",
                ],
                "highlight": False,
            },
        ],
        takeaway="반도체(82%)와 조선(31%)의 격차는 기술이 아닌 '반복 데이터 구조' 차이 — 조선은 ETO 특성상 데이터 혁신 없이 AI ROI 실현 불가",
        footer=SlideFooter(
            source="출처: McKinsey 2024, PwC Industrial AI Survey 2024, SEMI 2024, DNV Maritime 2024",
            right="PwC",
        ),
    )
    return spec, comparison_matrix


# ============================================================
# 2. 클라우드 마이그레이션 5단계 — process_flow
# ============================================================

def build_cloud_migration():
    spec = ProcessSpec(
        header=SlideHeader(
            title="Assessment→Optimize 5단계 체계적 전환으로 TCO 30% 절감과 배포 주기 10배 단축을 달성",
            category="엔터프라이즈 클라우드 마이그레이션 — Methodology",
            nav_path=["2. 클라우드 전략", "4. 전환 방법론"],
        ),
        intro="평균 14~20개월 소요되는 엔터프라이즈 마이그레이션을 5단계 파이프라인으로 체계화 — 단계별 산출물과 게이트 기준",
        steps=[
            {
                "name": "Assessment",
                "actor": "EA + Cloud Strategy Lead",
                "tools": "AWS Migration Evaluator\nCAST Highlight\nServiceNow CMDB",
                "output": "앱 인벤토리 + TCO 분석 + 비즈니스 케이스",
                "duration": "8~12주",
                "prerequisites": "경영진 스폰서십 + 예산 승인",
                "risks": "Shadow IT 누락 (CMDB 대비 15~30%↑)",
                "metrics": "인벤토리 커버리지 ≥95%",
            },
            {
                "name": "Planning",
                "actor": "Cloud Architect + CISO",
                "tools": "Well-Architected Tool\nTerraform\nJira/Confluence",
                "output": "6R 로드맵 + 타깃 아키텍처 + 거버넌스 모델",
                "duration": "8~12주",
                "prerequisites": "앱 인벤토리 + TCO 분석 완료",
                "risks": "과도한 Refactor 계획 → 예산 초과",
                "metrics": "100% 워크로드 6R 분류",
            },
            {
                "name": "Build",
                "actor": "Infra Architect + DevOps",
                "tools": "Control Tower\nTerraform/CDK\nCheckov/tfsec",
                "output": "랜딩 존 + IaC 코드베이스 + 보안 기준선",
                "duration": "8~16주",
                "prerequisites": "타깃 아키텍처 + IP 주소 계획",
                "risks": "CIDR 충돌, ClickOps 잔존",
                "metrics": "CIS Benchmark ≥85%",
            },
            {
                "name": "Migration",
                "actor": "Migration Engineer + DBA",
                "tools": "AWS MGN\nDMS (CDC)\nSnowball Edge",
                "output": "Wave별 완료 보고서 + 컷오버 런북",
                "duration": "16~36주",
                "prerequisites": "랜딩 존 완료 + 네트워크 안정화",
                "risks": "Hidden Dependencies, DB CDC 지연",
                "metrics": "컷오버 성공률 100%",
            },
            {
                "name": "Optimize",
                "actor": "FinOps + Cloud Native Arch",
                "tools": "Cost Explorer\nCompute Optimizer\nDatadog/Grafana",
                "output": "FinOps 대시보드 + 성능 기준선 + 현대화 로드맵",
                "duration": "6~12개월+",
                "prerequisites": "마이그레이션 완료 + 모니터링 운영",
                "risks": "CoE 공백 → 비용 통제 실패",
                "metrics": "클라우드 낭비율 <15%",
            },
        ],
        takeaway="마이그레이션 실패율 50%의 주원인은 Assessment 부실 — Shadow IT 파악 + 6R 분류 없이 착수하면 Wave 3에서 일정 붕괴",
        footer=SlideFooter(
            source="출처: Gartner 2023, McKinsey Cloud 2023, AWS MAP Guide, FinOps Foundation 2023",
            right="PwC",
        ),
    )
    return spec, process_flow


# ============================================================
# 3. DX 투자 포트폴리오 — quadrant_story
# ============================================================

def build_dx_portfolio():
    spec = QuadrantSpec(
        header=SlideHeader(
            title="DX 투자의 74%가 목표 미달성 — Quick Win 선행 후 Strategic Bet 순차 전환이 성패를 결정",
            category="제조업 디지털 전환 — Investment Portfolio",
            nav_path=["3. DX 전략", "5. 투자 우선순위"],
        ),
        intro="10개 DX 이니셔티브를 ROI×구현 복잡도로 분류 — 초기 2년간 Quick Win에 예산 40~50% 집중이 핵심",
        x_axis_label="구현 복잡도",
        y_axis_label="기대 ROI",
        x_low="LOW", x_high="HIGH",
        y_low="LOW", y_high="HIGH",
        quadrants=[
            {  # TL: 고ROI / 저복잡도 = QUICK WIN
                "title": "QUICK WIN",
                "highlight": True,
                "items": [
                    "예측적 설비 유지보수 (ROI 350%, 14개월)",
                    "실시간 생산 대시보드 (ROI 250%, 8개월)",
                    "백오피스 RPA (ROI 200%, 6개월)",
                ],
                "description": "기존 인프라 활용, 6~12개월 내 ROI 입증 — 변화관리 부담 최소",
                "action": "즉시 착수: 파일럿 3개월 → ROI 검증 → 전사 확대",
            },
            {  # TR: 고ROI / 고복잡도 = STRATEGIC BET
                "title": "STRATEGIC BET",
                "highlight": False,
                "items": [
                    "디지털 트윈 (ROI 500%, 24~36개월)",
                    "AI 품질 검사 (ROI 400%, 18~30개월)",
                    "스마트 팩토리 통합 (ROI 300%, 24~48개월)",
                ],
                "description": "IoT+데이터+도메인 인력 동시 필요 — Quick Win 역량 축적 후 착수",
                "action": "단계적 투자: 핵심 라인 1개 파일럿 → 비즈니스 케이스 → 전사 로드맵",
            },
            {  # BL: 저ROI / 저복잡도 = NICE TO HAVE
                "title": "NICE TO HAVE",
                "highlight": False,
                "items": [
                    "문서 디지털화 (ROI 100%, 12개월)",
                    "임직원 셀프서비스 포털 (ROI 80%)",
                    "기초 BI 리포팅 (ROI 120%)",
                ],
                "description": "핵심 프로세스 영향 제한적 — 대형 이니셔티브에 번들링 권장",
                "action": "여력 시 추진: Quick Win 또는 Smart Factory에 포함하여 별도 예산 지양",
            },
            {  # BR: 저ROI / 고복잡도 = DEFER
                "title": "DEFER",
                "highlight": False,
                "items": [
                    "자율 제조 무인화 (페이백 5~10년)",
                    "블록체인 공급망 (PoC→전환 8%)",
                ],
                "description": "기술 성숙도 미달 또는 네트워크 효과 전제 — 현 시점 ROI 논거 미약",
                "action": "보류: Smart Factory 완성 후 중장기 비전으로 재평가",
            },
        ],
        insight="DX 성공 기업은 초기 2년간 Quick Win에 예산 40~50% 집중하여 조직 역량+경영진 신뢰를 확보한 후 Strategic Bet으로 무게중심 이동",
        footer=SlideFooter(
            source="출처: McKinsey Digital 2023, Gartner 2024, IDC Smart Manufacturing 2024, WEF Lighthouse 2024",
            right="PwC",
        ),
    )
    return spec, quadrant_story


# ============================================================
# 메인 — 3장 빌드 + 검증
# ============================================================

def build_one(builder_fn, output: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec, pattern_func = builder_fn()
    pattern_func(slide, spec)
    prs.save(output)
    return output


def main():
    out_dir = Path("output/quality_ceiling")
    out_dir.mkdir(parents=True, exist_ok=True)

    cases = [
        ("01_ai_maturity", build_ai_maturity, "comparison"),
        ("02_cloud_migration", build_cloud_migration, "process"),
        ("03_dx_portfolio", build_dx_portfolio, "quadrant"),
    ]

    print("=" * 70)
    print("Quality Ceiling Test — 3개 새 주제 × 3개 패턴")
    print("=" * 70)

    all_passed = True
    for name, builder_fn, kind in cases:
        out = out_dir / f"{name}.pptx"
        try:
            build_one(builder_fn, out)
        except Exception as e:
            print(f"\n[{name}] BUILD FAILED: {e}")
            import traceback
            traceback.print_exc()
            all_passed = False
            continue

        visual = validate_visual(out, convert_pdf=False)
        design = inspect_design(str(out), pattern_kind=kind)

        ok = not visual.issues and design.passed
        status = "PASS" if ok else "FAIL"
        if not ok:
            all_passed = False

        print(f"\n[{name}] {status}")
        print(f"  visual: {len(visual.issues)} issues")
        for i in visual.issues[:5]:
            print(f"    - {i}")
        print(f"  design: passed={design.passed}, issues={len(design.issues)}")
        for i in design.issues[:5]:
            print(f"    - {i}")
        print(f"  metrics: density={design.metrics.get('slide_1_density', '?')}")

    # PDF 변환
    print("\n" + "=" * 70)
    print("PDF 변환 중...")
    print("=" * 70)
    for name, _, _ in cases:
        out = out_dir / f"{name}.pptx"
        if not out.exists():
            continue
        try:
            visual = validate_visual(out, convert_pdf=True)
            print(f"  {name}.pdf  — {'OK' if visual.pdf_available else 'SKIP'}")
        except Exception as e:
            print(f"  {name}.pdf  — FAIL ({e})")

    print("\n" + "=" * 70)
    print(f"결과: {'ALL PASSED' if all_passed else 'SOME FAILED'}")
    print("=" * 70)
    return 0 if all_passed else 1


if __name__ == "__main__":
    sys.exit(main())
