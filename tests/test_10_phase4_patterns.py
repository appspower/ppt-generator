"""10개 Phase 4 신규 패턴 통합 테스트."""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.patterns import (
    SlideFooter, SlideHeader,
    MaturityModelSpec, maturity_model,
    MilestoneTimelineSpec, milestone_timeline,
    RagStatusSpec, rag_status_table,
    ArchStackSpec, architecture_stack,
    GanttSpec, gantt_roadmap,
    CycleSpec, cycle_diagram,
    ValueChainSpec, value_chain,
    BubbleChartSpec, bubble_chart,
    TreeSpec, tree_diagram,
    HarveyBallSpec, harvey_ball_matrix,
)
from ppt_builder.visual_validate import validate_visual

FOOTER = SlideFooter(source="출처: PwC Analysis 2024", right="PwC")


def build_maturity_model():
    spec = MaturityModelSpec(
        header=SlideHeader(
            title="HD현대 디지털 성숙도 현재 L2(Managed) — SAP 전환과 병행하여 L4(Predictive)까지 18개월 내 도달 목표",
            category="디지털 전환 — 성숙도 진단",
            nav_path=["1. 전략 방향", "2. 성숙도 평가"],
        ),
        intro="디지털 성숙도 5단계 모델 기반 현재 수준 진단과 목표 경로 — SAP S/4HANA 전환이 L3→L4 도약의 핵심 레버",
        stages=[
            {"title": "Ad-hoc", "level": "L1",
             "detail": "Excel 기반 수작업 보고, 부서별 사일로 데이터, 표준화 미비. 의사결정이 경험과 직감에 의존."},
            {"title": "Managed", "level": "L2",
             "detail": "MDM 초기 구축, ERP 데이터 표준화 착수. 부서 간 공유 시작되었으나 실시간 연동 부재."},
            {"title": "Analytical", "level": "L3",
             "detail": "BI 대시보드 운영, KPI 기반 경영 정착. Self-Service 분석 환경 구축, 데이터 거버넌스 확립."},
            {"title": "Predictive", "level": "L4",
             "detail": "ML 파이프라인 운영, 수요 예측·설비 예지보전 모델 적용. Proactive Alert 기반 의사결정."},
            {"title": "AI-Driven", "level": "L5",
             "detail": "자율 의사결정 시스템, 실시간 최적화 엔진. Prescriptive Analytics로 최적 행동 자동 추천."},
        ],
        current=1,
        target=3,
        takeaway="L2→L4 전환의 핵심은 SAP S/4HANA 실시간 데이터 + Foundry 분석 플랫폼 — 18개월·$2.5M 투자로 예측 경영 체제 진입",
        footer=FOOTER,
    )
    return spec, maturity_model


def build_milestone_timeline():
    spec = MilestoneTimelineSpec(
        header=SlideHeader(
            title="SAP S/4HANA 전환 18개월 로드맵 — 6대 마일스톤 기준으로 단계별 Go/No-Go 게이트 운영",
            category="프로젝트 일정 — 마일스톤 타임라인",
            nav_path=["2. 추진 계획", "1. 마스터 일정"],
        ),
        intro="2024 Q1 착수부터 2025 Q2 Go-Live까지 6대 마일스톤 — 각 게이트 통과 기준과 핵심 산출물 정의",
        milestones=[
            {"date": "2024 Q1", "title": "킥오프 완료",
             "detail": "프로젝트 헌장 승인, PMO 체계 수립, Foundry 환경 프로비저닝 완료",
             "highlight": True},
            {"date": "2024 Q2", "title": "Blueprint 확정",
             "detail": "Fit-Gap 분석 180건 완료, To-Be 프로세스 135개 확정, 커스터마이징 범위 동결",
             "highlight": False},
            {"date": "2024 Q3", "title": "Build 완료",
             "detail": "ABAP 개발 92건, 인터페이스 47건, 데이터 마이그레이션 설계 완료",
             "highlight": False},
            {"date": "2024 Q4", "title": "통합 테스트 완료",
             "detail": "AIP 자동생성 TC 20K건 수행, 결함 해소율 97%, UAT 현업 서명 확보",
             "highlight": True},
            {"date": "2025 Q1", "title": "Cutover 리허설",
             "detail": "3회 리허설 완료, DT 12시간→6시간 단축, 롤백 시나리오 검증 완료",
             "highlight": False},
            {"date": "2025 Q2", "title": "Go-Live",
             "detail": "본번 전환 완료, Hypercare 4주 운영, KPI 모니터링 대시보드 가동",
             "highlight": True},
        ],
        takeaway="6대 마일스톤 중 킥오프·통합테스트·Go-Live가 Critical Gate — Palantir 모듈이 테스트 자동화와 Cutover 최적화를 직접 지원",
        footer=FOOTER,
    )
    return spec, milestone_timeline


def build_rag_status():
    spec = RagStatusSpec(
        header=SlideHeader(
            title="SAP 전환 8대 워크스트림 중 데이터 마이그레이션과 인터페이스가 Amber — 집중 관리 필요",
            category="PMO 대시보드 — RAG 상태 현황",
            nav_path=["3. 거버넌스", "1. 워크스트림 현황"],
        ),
        intro="2024년 Q3 기준 8대 워크스트림의 일정·예산·품질·리스크 4개 차원 RAG 현황 — 주간 스티어링 보고용",
        columns=["워크스트림", "일정", "예산", "품질", "리스크"],
        rows=[
            {"name": "FI/CO 재무회계", "values": ["G", "G", "G", "G"]},
            {"name": "MM 구매관리", "values": ["G", "G", "A", "G"]},
            {"name": "SD 영업물류", "values": ["G", "A", "G", "G"]},
            {"name": "PP 생산계획", "values": ["A", "G", "G", "A"]},
            {"name": "데이터 마이그레이션", "values": ["A", "A", "R", "A"]},
            {"name": "인터페이스 개발", "values": ["A", "G", "A", "R"]},
            {"name": "권한·보안 설계", "values": ["G", "G", "G", "G"]},
            {"name": "변화관리·교육", "values": ["G", "G", "G", "A"]},
        ],
        takeaway="데이터 마이그레이션 품질(R)과 인터페이스 리스크(R) 즉시 대응 필요 — Foundry Pipeline 자동 검증으로 품질 이슈 선제 해소 권장",
        footer=FOOTER,
    )
    return spec, rag_status_table


def build_architecture_stack():
    spec = ArchStackSpec(
        header=SlideHeader(
            title="SAP S/4HANA 5계층 기술 아키텍처 — Presentation부터 Infrastructure까지 End-to-End 설계",
            category="기술 아키텍처 — 스택 다이어그램",
            nav_path=["1. 기술 전략", "2. 아키텍처 설계"],
        ),
        intro="SAP S/4HANA 전환 후 목표 아키텍처 5계층 — 각 레이어별 핵심 컴포넌트와 기술 스택 정의",
        layers=[
            {"title": "Presentation", "badge": "L1",
             "items": ["SAP Fiori", "Launchpad", "SAPUI5", "Mobile Apps"]},
            {"title": "Application", "badge": "L2",
             "items": ["S/4HANA Core", "SuccessFactors", "Ariba", "BTP Extension"]},
            {"title": "Integration", "badge": "L3",
             "items": ["SAP BTP Integration Suite", "API Management", "Event Mesh"]},
            {"title": "Data & Analytics", "badge": "L4",
             "items": ["SAP BW/4HANA", "SAC", "Foundry Ontology", "Data Intelligence"]},
            {"title": "Infrastructure", "badge": "L5",
             "items": ["Azure (IaaS)", "HANA Cloud", "Kubernetes", "Disaster Recovery"]},
        ],
        side_label="보안·거버넌스",
        takeaway="5계층 스택에서 L3(Integration)과 L4(Data)가 Foundry 연동의 핵심 접점 — BTP Integration Suite를 통한 실시간 데이터 흐름 보장",
        footer=FOOTER,
    )
    return spec, architecture_stack


def build_gantt_roadmap():
    spec = GanttSpec(
        header=SlideHeader(
            title="SAP 전환 4개 스트림이 6개 분기에 걸쳐 병렬 수행 — Q5 Go-Live를 향한 통합 로드맵",
            category="프로그램 로드맵 — Gantt 차트",
            nav_path=["2. 추진 계획", "2. 통합 로드맵"],
        ),
        intro="Infra·Application·Data·Change Management 4개 스트림의 Q1~Q6 일정 — Go-Live(Q5) 기준 역산 설계",
        phases=["Q1", "Q2", "Q3", "Q4", "Q5", "Q6"],
        streams=[
            {"name": "Infra 구축",
             "bars": [
                 {"start": 0, "end": 2, "label": "환경 구축", "highlight": False},
                 {"start": 2, "end": 3, "label": "성능 튜닝", "highlight": False},
                 {"start": 4, "end": 6, "label": "운영 전환", "highlight": False},
             ]},
            {"name": "Application 개발",
             "bars": [
                 {"start": 0, "end": 1, "label": "Fit-Gap", "highlight": False},
                 {"start": 1, "end": 3, "label": "Build", "highlight": True},
                 {"start": 3, "end": 4.5, "label": "테스트", "highlight": True},
             ]},
            {"name": "Data 마이그레이션",
             "bars": [
                 {"start": 0, "end": 1.5, "label": "분석·설계", "highlight": False},
                 {"start": 1.5, "end": 3, "label": "ETL 개발", "highlight": False},
                 {"start": 3, "end": 4.5, "label": "검증·전환", "highlight": True},
             ]},
            {"name": "Change Mgmt",
             "bars": [
                 {"start": 0, "end": 1, "label": "영향 분석", "highlight": False},
                 {"start": 1, "end": 4, "label": "교육·소통", "highlight": False},
                 {"start": 4, "end": 6, "label": "Hypercare", "highlight": False},
             ]},
        ],
        milestones=[
            {"phase": 1, "label": "Blueprint 확정"},
            {"phase": 3, "label": "Build 완료"},
            {"phase": 5, "label": "Go-Live"},
        ],
        takeaway="Application Build(Q2~Q3)와 Data 검증(Q4)이 Critical Path — 두 스트림의 병렬 완료가 Q5 Go-Live의 전제 조건",
        footer=FOOTER,
    )
    return spec, gantt_roadmap


def build_cycle_diagram():
    spec = CycleSpec(
        header=SlideHeader(
            title="SAP 운영 안정화를 위한 PDCA 지속적 개선 사이클 — Foundry가 Check 단계를 자동화",
            category="운영 모델 — PDCA 순환",
            nav_path=["4. 운영 전략", "1. 지속적 개선"],
        ),
        intro="Go-Live 이후 운영 안정화와 지속적 개선을 위한 PDCA 사이클 — 매 분기 1회전 목표",
        center={"title": "PDCA", "subtitle": "지속적 개선"},
        stages=[
            {"title": "Plan (계획)",
             "detail": "KPI 목표 설정, 개선 과제 도출, Foundry 대시보드에서 이상 패턴 식별 후 과제화",
             "badge": "P"},
            {"title": "Do (실행)",
             "detail": "SAP 설정 변경, 프로세스 개선 시행, BTP Extension 개발, 현업 파일럿 적용",
             "badge": "D"},
            {"title": "Check (확인)",
             "detail": "Foundry 자동 모니터링으로 KPI 달성률 측정, 편차 분석, 원인 진단 리포트 자동 생성",
             "badge": "C"},
            {"title": "Act (조치)",
             "detail": "표준화·수평전개 결정, 미달 과제 재계획, 거버넌스 위원회 보고 및 차기 사이클 반영",
             "badge": "A"},
        ],
        takeaway="Check 단계에서 Foundry 자동화가 수작업 분석 대비 80% 시간 절감 — 분기 사이클을 월간으로 단축 가능",
        footer=FOOTER,
    )
    return spec, cycle_diagram


def build_value_chain():
    spec = ValueChainSpec(
        header=SlideHeader(
            title="HD현대 조선 가치사슬에서 조달·건조가 원가의 70% — SAP S/4HANA 전환의 최우선 대상 영역",
            category="전략 분석 — Porter 가치사슬",
            nav_path=["1. 전략 방향", "3. 가치사슬 분석"],
        ),
        intro="Porter 가치사슬 프레임워크로 HD현대 조선의 주요활동 5단계와 지원활동 3개 영역을 매핑 — SAP 전환 우선순위 도출",
        primary=[
            {"title": "수주·설계",
             "detail": "선박 수주 영업, 기본설계·상세설계, 선주 요구사항 관리. SAP SD·PS 모듈 적용 대상.",
             "highlight": False},
            {"title": "조달·구매",
             "detail": "강재·기자재 구매, 협력사 관리, 글로벌 소싱. 원가 35% 차지. SAP MM·SRM 핵심 영역.",
             "highlight": True},
            {"title": "건조·생산",
             "detail": "블록 제작, 탑재, 의장, 도장. 원가 35% 차지. SAP PP·PM 모듈로 공정 최적화.",
             "highlight": True},
            {"title": "시운전·인도",
             "detail": "해상 시운전, 선급 검사, 인도. SAP QM·PS 모듈로 품질 추적 및 프로젝트 정산.",
             "highlight": False},
            {"title": "A/S·보증",
             "detail": "보증 기간 관리, 부품 공급, 기술 지원. SAP CS·PM 모듈로 서비스 이력 관리.",
             "highlight": False},
        ],
        support=[
            {"title": "기술 인프라",
             "detail": "PLM·CAD 시스템, SAP S/4HANA, BTP 통합 플랫폼, Foundry 분석 환경"},
            {"title": "인적자원 관리",
             "detail": "기능 인력 양성, 디지털 역량 교육, SuccessFactors 기반 성과 관리"},
            {"title": "재무·경영지원",
             "detail": "원가 관리, 투자 의사결정, SAP FI/CO 기반 실시간 경영 정보 제공"},
        ],
        margin_label="마진",
        takeaway="조달(35%)+건조(35%)=원가 70% — 이 두 영역의 SAP MM·PP 고도화가 마진 개선의 최대 레버",
        footer=FOOTER,
    )
    return spec, value_chain


def build_bubble_chart():
    spec = BubbleChartSpec(
        header=SlideHeader(
            title="SAP 8대 모듈 중 MM·PP가 높은 ROI와 높은 복잡도를 동시 보유 — 전문 인력 집중 배치 필요",
            category="포트폴리오 분석 — ROI×복잡도 버블 차트",
            nav_path=["1. 전략 방향", "4. 모듈 우선순위"],
        ),
        intro="SAP 8대 모듈의 기대 ROI(x축)와 구현 복잡도(y축)를 매핑 — 버블 크기는 영향 범위(사용자 수) 반영",
        x_label="기대 ROI (배)",
        y_label="구현 복잡도 (1~10)",
        bubbles=[
            {"label": "MM", "x": 3.8, "y": 7.5, "size": 9, "highlight": True},
            {"label": "SD", "x": 3.2, "y": 6.0, "size": 7, "highlight": True},
            {"label": "PP", "x": 4.2, "y": 8.5, "size": 8, "highlight": True},
            {"label": "FI", "x": 2.8, "y": 4.0, "size": 6, "highlight": False},
            {"label": "CO", "x": 2.5, "y": 3.5, "size": 5, "highlight": False},
            {"label": "HR", "x": 1.8, "y": 5.0, "size": 4, "highlight": False},
            {"label": "PM", "x": 3.0, "y": 6.5, "size": 6, "highlight": False},
            {"label": "QM", "x": 2.2, "y": 4.5, "size": 3, "highlight": False},
        ],
        narratives=[
            {"title": "High ROI + High Complexity (MM, PP)",
             "detail": "ROI 3.8~4.2배로 최고 수준이나 복잡도 7.5~8.5로 리스크 동반. 전문 컨설턴트 집중 배치 및 Foundry 테스트 자동화 필수."},
            {"title": "High ROI + Mid Complexity (SD, PM)",
             "detail": "ROI 3.0~3.2배, 복잡도 6.0~6.5 중간 수준. 표준 프로세스 적용으로 빠른 가치 실현 가능."},
            {"title": "Stable Modules (FI, CO, QM, HR)",
             "detail": "ROI 1.8~2.8배, 복잡도 3.5~5.0. 글로벌 표준 프로세스 적용 가능하여 리스크 낮음. 병렬 구현 권장."},
        ],
        takeaway="MM·PP 모듈이 ROI 최고(4.2배)이나 복잡도도 최고(8.5) — 이 두 모듈에 테스트 자동화·전문 인력을 집중 투입해야 전환 성공 보장",
        footer=FOOTER,
    )
    return spec, bubble_chart


def build_tree_diagram():
    spec = TreeSpec(
        header=SlideHeader(
            title="SAP 전환 리스크를 기술·조직·프로세스·외부 4개 축으로 MECE 분해 — 12개 세부 리스크 식별",
            category="리스크 관리 — Issue Tree",
            nav_path=["3. 거버넌스", "2. 리스크 분해"],
        ),
        intro="SAP S/4HANA 전환 프로젝트의 리스크를 MECE 원칙으로 4대 카테고리·12개 세부 항목으로 구조화",
        root={"title": "전환 리스크"},
        branches=[
            {"title": "기술 리스크",
             "children": [
                 "데이터 마이그레이션 정합성 실패",
                 "인터페이스 호환성 오류 (Legacy↔S/4)",
                 "성능 저하 (HANA 사이징 부족)",
             ],
             "highlight": True},
            {"title": "조직 리스크",
             "children": [
                 "핵심 인력 이탈·리소스 부족",
                 "현업 저항·변화관리 실패",
                 "경영진 지원 약화",
             ],
             "highlight": False},
            {"title": "프로세스 리스크",
             "children": [
                 "요구사항 Scope Creep",
                 "테스트 커버리지 부족",
                 "Cutover 리허설 미흡",
             ],
             "highlight": True},
            {"title": "외부 리스크",
             "children": [
                 "SAP 라이선스 정책 변경",
                 "협력사 시스템 연동 지연",
                 "규제·컴플라이언스 변경",
             ],
             "highlight": False},
        ],
        takeaway="기술 리스크(데이터·인터페이스)와 프로세스 리스크(테스트·Cutover)가 발생 확률 Top 4 — Foundry 자동화로 직접 억제 가능",
        footer=FOOTER,
    )
    return spec, tree_diagram


def build_harvey_ball_matrix():
    spec = HarveyBallSpec(
        header=SlideHeader(
            title="ERP 4사 비교에서 SAP S/4HANA가 기능·확장성·생태계 3개 축 최고점 — HD현대 최적 선택지",
            category="벤더 비교 — Harvey Ball 매트릭스",
            nav_path=["1. 전략 방향", "5. 솔루션 선정"],
        ),
        intro="SAP S/4HANA·Oracle Cloud ERP·MS Dynamics 365·Workday 4개 솔루션의 5대 평가 기준 정량 비교",
        row_labels=["SAP S/4HANA", "Oracle Cloud ERP", "MS Dynamics 365", "Workday"],
        col_labels=["기능 완성도", "확장성", "TCO 효율", "생태계", "UX"],
        scores=[
            [4, 4, 2, 4, 3],
            [4, 3, 3, 3, 3],
            [3, 3, 4, 2, 4],
            [2, 2, 3, 2, 4],
        ],
        highlight_row=0,
        takeaway="SAP S/4HANA가 기능(4)·확장성(4)·생태계(4)에서 최고 — TCO(2)는 열위이나 조선·중공업 특화 기능과 글로벌 파트너 생태계가 보상",
        footer=FOOTER,
    )
    return spec, harvey_ball_matrix


# ============================================================
# 메인
# ============================================================

def build_one(builder_fn, output: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec, pattern_func = builder_fn()
    pattern_func(slide, spec)
    prs.save(output)
    return output


def main():
    out_dir = Path("output/phase4_patterns")
    out_dir.mkdir(parents=True, exist_ok=True)

    cases = [
        ("13_maturity_model", build_maturity_model),
        ("14_milestone_timeline", build_milestone_timeline),
        ("15_rag_status", build_rag_status),
        ("16_arch_stack", build_architecture_stack),
        ("17_gantt_roadmap", build_gantt_roadmap),
        ("18_cycle_diagram", build_cycle_diagram),
        ("19_value_chain", build_value_chain),
        ("20_bubble_chart", build_bubble_chart),
        ("21_tree_diagram", build_tree_diagram),
        ("22_harvey_ball", build_harvey_ball_matrix),
    ]

    print("=" * 70)
    print("10 Phase 4 Patterns Test")
    print("=" * 70)

    all_passed = True
    for name, builder_fn in cases:
        out = out_dir / f"{name}.pptx"
        try:
            build_one(builder_fn, out)
        except Exception as e:
            print(f"\n[{name}] BUILD FAILED: {e}")
            import traceback
            traceback.print_exc()
            all_passed = False
            continue

        visual = validate_visual(out, convert_pdf=False)
        ok = not visual.issues
        status = "PASS" if ok else "FAIL"
        if not ok:
            all_passed = False

        print(f"\n[{name}] {status}")
        print(f"  visual: {len(visual.issues)} issues")
        for i in visual.issues[:3]:
            print(f"    - {i}")

    # PDF
    print("\n" + "=" * 70)
    print("PDF...")
    print("=" * 70)
    for name, _ in cases:
        out = out_dir / f"{name}.pptx"
        if not out.exists():
            continue
        try:
            visual = validate_visual(out, convert_pdf=True)
            print(f"  {name}.pdf  -- {'OK' if visual.pdf_available else 'SKIP'}")
        except Exception as e:
            print(f"  {name}.pdf  -- FAIL ({e})")

    print("\n" + "=" * 70)
    print(f"Result: {'ALL PASSED' if all_passed else 'SOME FAILED'}")
    print("=" * 70)
    return 0 if all_passed else 1


if __name__ == "__main__":
    sys.exit(main())
