"""7개 신규 패턴 통합 테스트."""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.patterns import (
    SlideFooter, SlideHeader,
    HubSpokeSpec, hub_spoke,
    BeforeAfterSpec, before_after,
    KpiDashboardSpec, kpi_dashboard,
    WaterfallSpec, waterfall_bridge,
    SwimlaneSpec, swimlane,
    PyramidSpec, pyramid_layers,
    DataNarrativeSpec, data_narrative,
)
from ppt_builder.design_check import inspect_design
from ppt_builder.visual_validate import validate_visual

FOOTER = SlideFooter(source="출처: PwC Analysis 2024", right="PwC")


def build_hub_spoke():
    spec = HubSpokeSpec(
        header=SlideHeader(
            title="Foundry Ontology를 중심으로 6개 모듈이 단일 데이터 허브에 연결되어 통합 운영",
            category="시스템 아키텍처 — Hub & Spoke",
            nav_path=["1. 기술 전략", "3. 통합 아키텍처"],
        ),
        intro="Palantir Foundry Ontology가 중심 허브 역할 — 6개 기능 모듈이 공통 데이터 모델을 공유하여 사일로 제거",
        hub={"title": "Foundry\nOntology", "subtitle": "통합 데이터 허브"},
        spokes=[
            {"badge": "L1", "title": "Health Dashboard", "detail": "프로젝트 현황 실시간 가시화, PMO 보고 자동화"},
            {"badge": "L1", "title": "Config Register", "detail": "설정 결정 이력 추적, 변경 영향 분석"},
            {"badge": "L2", "title": "AIP Test Generator", "detail": "Blueprint→TC 자동 생성, LLM 검증"},
            {"badge": "L2", "title": "Defect Triage", "detail": "결함 클러스터링, 우선순위 자동 분류"},
            {"badge": "L3", "title": "Cutover App", "detail": "Task DAG 오케스트레이션, 리허설 비교"},
            {"badge": "L4", "title": "AI Go/No-Go", "detail": "Readiness Scorecard, 이상 탐지"},
        ],
        takeaway="Ontology 중심 아키텍처 = 모듈 간 데이터 사일로 제로 — 신규 모듈 추가 시 연동 비용이 O(1)로 수렴",
        footer=FOOTER,
    )
    return spec, hub_spoke


def build_before_after():
    spec = BeforeAfterSpec(
        header=SlideHeader(
            title="수작업 3~6개월 테스트 프로세스를 AIP 자동화로 1~2주 파이프라인으로 전환",
            category="테스트 자동화 — AS-IS vs TO-BE",
            nav_path=["2. 활용 전략", "5. 프로세스 혁신"],
        ),
        intro="기존 수작업 중심 테스트 프로세스와 AIP 자동화 이후 프로세스를 5개 차원에서 정량 비교",
        before_title="AS-IS (수작업)",
        after_title="TO-BE (AIP 자동화)",
        arrow_label="전환",
        before_items=[
            {"label": "TC 작성", "detail": "컨설턴트가 Blueprint 수동 분석 후 Excel 작성", "kpi": "3~6개월 소요"},
            {"label": "검증 방식", "detail": "육안 검토 + 동료 리뷰, 누락 빈번", "kpi": "정확도 85%"},
            {"label": "결함 등록", "detail": "Jira 수동 입력, 객체 링크 없음", "kpi": "등록 지연 2~3일"},
            {"label": "커버리지", "detail": "시간 제약으로 핵심 시나리오만 작성", "kpi": "커버리지 40~60%"},
            {"label": "비용 구조", "detail": "컨설턴트 인건비 중심, 규모 비례 증가", "kpi": "건당 $50~100"},
        ],
        after_items=[
            {"label": "TC 작성", "detail": "AIP가 Blueprint→규칙→TC 자동 생성", "kpi": "1~2주 완료"},
            {"label": "검증 방식", "detail": "LLM-as-Judge 0~10점 연속 검증", "kpi": "정확도 99.8%"},
            {"label": "결함 등록", "detail": "Foundry→Jira REST API 자동 연동", "kpi": "실시간 등록"},
            {"label": "커버리지", "detail": "전수 시나리오 자동 생성, 엣지케이스 포함", "kpi": "커버리지 95%+"},
            {"label": "비용 구조", "detail": "플랫폼 비용 고정, 규모 무관", "kpi": "건당 $0.2"},
        ],
        takeaway="AIP 전환 시 TC 작성 공수 70%↓, 정확도 85%→99.8%, 커버리지 60%→95% — 비용은 건당 $50→$0.2로 250배 효율화",
        footer=FOOTER,
    )
    return spec, before_after


def build_kpi_dashboard():
    spec = KpiDashboardSpec(
        header=SlideHeader(
            title="Palantir 도입 6개월 후 6대 핵심 KPI 전원 목표 달성 — Quick Win 전략의 정량 성과",
            category="성과 대시보드 — KPI Scorecard",
            nav_path=["3. 성과 보고", "1. KPI 종합"],
        ),
        intro="L1 Quick Win 착수 후 6개월 시점 기준, 6대 핵심 성과 지표의 달성 현황과 트렌드",
        kpis=[
            {"value": "14%", "label": "전체 일정 단축", "subtitle": "18개월 → 15.5개월",
             "detail": "목표 10% 초과 달성", "trend": "up"},
            {"value": "70%", "label": "테스트 공수 절감", "subtitle": "수작업 → AIP 자동화",
             "detail": "L2 완료 기준", "trend": "up"},
            {"value": "99.8%", "label": "TC 정확도", "subtitle": "LLM-as-Judge 검증",
             "detail": "기존 85% 대비", "trend": "up"},
            {"value": "50%", "label": "Cutover DT 단축", "subtitle": "초과율 10%↓",
             "detail": "3회 리허설 완료", "trend": "up"},
            {"value": "실시간", "label": "PMO 보고 주기", "subtitle": "기존 3일 → 실시간",
             "detail": "Health Dashboard 가동", "trend": "up"},
            {"value": "$4K", "label": "TC 생성 비용", "subtitle": "20K+ 레코드 2주",
             "detail": "기존 대비 250배↓", "trend": "down"},
        ],
        bottom_note="* 모든 수치는 L1~L2 완료 기준 (2024년 6월). L3 Cutover는 예상치 기반.",
        takeaway="6대 KPI 전원 Green — Quick Win(L1) 2주 만에 가치 입증 후 L2 확대 결정이 자연스럽게 승인",
        footer=FOOTER,
    )
    return spec, kpi_dashboard


def build_waterfall():
    spec = WaterfallSpec(
        header=SlideHeader(
            title="SAP 전환 일정 18개월에서 Palantir 도입 효과로 15.5개월까지 2.5개월(14%) 단축을 달성",
            category="일정 효과 분석 — Waterfall Bridge",
            nav_path=["3. 성과 보고", "2. 일정 분석"],
        ),
        intro="기존 18개월 SAP 전환 일정에서 각 Palantir 모듈이 기여한 일정 단축 효과를 분해",
        start={"label": "기존 계획", "value": 18},
        steps=[
            {"label": "테스트\n자동화", "value": -1.2, "detail": "TC 생성 70%↓"},
            {"label": "Cutover\n최적화", "value": -0.8, "detail": "DT 50%↓"},
            {"label": "거버넌스\n자동화", "value": -0.3, "detail": "보고 실시간"},
            {"label": "리스크\n조기 감지", "value": -0.2, "detail": "Agent 알림"},
        ],
        end={"label": "최종 일정", "value": 15.5},
        unit="개월",
        takeaway="테스트 자동화(1.2개월)가 최대 기여 → Cutover(0.8개월) → 거버넌스(0.3개월) 순 — 총 2.5개월(14%) 단축",
        footer=FOOTER,
    )
    return spec, waterfall_bridge


def build_swimlane():
    spec = SwimlaneSpec(
        header=SlideHeader(
            title="4개 조직이 5단계에 걸쳐 협업하는 SAP 전환 거버넌스 — 각 단계별 역할과 산출물을 정의",
            category="거버넌스 모델 — RACI Swimlane",
            nav_path=["1. 추진 체계", "2. 역할 분담"],
        ),
        intro="PwC, HD현대 IT, Palantir, 현업 4개 조직의 단계별 역할을 시각화 — 협업 공백 제로 설계",
        lanes=["PwC 컨설팅", "HD현대 IT", "Palantir", "현업 부서"],
        phases=["Assessment", "Build", "Test", "Cutover", "Hypercare"],
        activities=[
            {"lane": 0, "phase": 0, "text": "현황 진단\nGap 분석", "highlight": True},
            {"lane": 0, "phase": 1, "text": "Blueprint\n검증"},
            {"lane": 0, "phase": 2, "text": "TC 리뷰\n품질 관리", "highlight": True},
            {"lane": 0, "phase": 3, "text": "Cutover\n총괄 PMO"},
            {"lane": 0, "phase": 4, "text": "안정화\n지원"},
            {"lane": 1, "phase": 0, "text": "인프라\n준비"},
            {"lane": 1, "phase": 1, "text": "시스템\n구축"},
            {"lane": 1, "phase": 2, "text": "환경\n관리"},
            {"lane": 1, "phase": 3, "text": "Go-Live\n실행", "highlight": True},
            {"lane": 1, "phase": 4, "text": "운영\n전환"},
            {"lane": 2, "phase": 0, "text": "Ontology\n설계"},
            {"lane": 2, "phase": 1, "text": "AIP 모듈\n개발", "highlight": True},
            {"lane": 2, "phase": 2, "text": "TC 자동\n생성"},
            {"lane": 2, "phase": 3, "text": "Cutover\nApp"},
            {"lane": 2, "phase": 4, "text": "AI Agent\n운영"},
            {"lane": 3, "phase": 0, "text": "요구사항\n정의"},
            {"lane": 3, "phase": 1, "text": "프로세스\n확인"},
            {"lane": 3, "phase": 2, "text": "UAT\n수행", "highlight": True},
            {"lane": 3, "phase": 3, "text": "데이터\n검증"},
            {"lane": 3, "phase": 4, "text": "현업\n안정화"},
        ],
        takeaway="PwC가 PMO 총괄, Palantir가 기술 플랫폼, HD현대 IT가 실행, 현업이 검증 — 4자 협업으로 공백 제로",
        footer=FOOTER,
    )
    return spec, swimlane


def build_pyramid():
    spec = PyramidSpec(
        header=SlideHeader(
            title="5단계 데이터 성숙도 모델에서 HD현대는 현재 L2 — L4 도달 시 AI 기반 의사결정 가능",
            category="데이터 성숙도 — Maturity Model",
            nav_path=["1. 전략", "4. 성숙도 진단"],
        ),
        intro="데이터 기반 조직으로의 전환을 위한 5단계 성숙도 모델 — 현재 위치 진단과 목표 경로 제시",
        layers=[
            {"badge": "L5", "title": "AI-Driven", "detail": "자율 의사결정, 실시간 최적화, Prescriptive Analytics"},
            {"badge": "L4", "title": "Predictive", "detail": "예측 모델 운영, ML Pipeline, Proactive Alert"},
            {"badge": "L3", "title": "Analytical", "detail": "BI 대시보드, KPI 기반 경영, Self-Service 분석"},
            {"badge": "L2", "title": "Managed", "detail": "데이터 표준화, MDM 구축, 부서간 공유 시작"},
            {"badge": "L1", "title": "Ad-hoc", "detail": "Excel 중심, 사일로 데이터, 수작업 보고"},
        ],
        side_notes=[
            {"label": "현재 수준", "value": "L2"},
            {"label": "목표 수준", "value": "L4"},
            {"label": "도달 기간", "value": "18개월"},
            {"label": "필요 투자", "value": "$2.5M"},
        ],
        takeaway="L2→L4 전환에 18개월·$2.5M 투자 필요 — Palantir Foundry가 L3(분석)→L4(예측)를 플랫폼 레벨에서 가속",
        footer=FOOTER,
    )
    return spec, pyramid_layers


def build_data_narrative():
    spec = DataNarrativeSpec(
        header=SlideHeader(
            title="SAP 전환 프로젝트 리스크 Top 7 중 테스트 지연이 발생 확률 85%로 최고 — 집중 대응 필요",
            category="리스크 분석 — Data-Driven View",
            nav_path=["2. 리스크 관리", "3. 정량 분석"],
        ),
        intro="과거 30개 SAP 전환 프로젝트 데이터 기반으로 리스크 발생 확률을 정량화 — Top 7 리스크와 대응 인사이트",
        chart_title="리스크 발생 확률 (%)",
        chart_data=[
            {"label": "테스트 지연", "value": 85, "highlight": True},
            {"label": "데이터 마이그레이션", "value": 72, "highlight": True},
            {"label": "Cutover 초과", "value": 65, "highlight": False},
            {"label": "인력 부족", "value": 58, "highlight": False},
            {"label": "요구사항 변경", "value": 52, "highlight": False},
            {"label": "인터페이스 오류", "value": 45, "highlight": False},
            {"label": "성능 이슈", "value": 38, "highlight": False},
        ],
        chart_unit="%",
        narratives=[
            {"title": "테스트 지연 (85%)", "detail": "수작업 TC 작성이 병목 — AIP 자동화로 70% 공수 절감하여 근본 해소. 30개 프로젝트 중 26개에서 발생."},
            {"title": "데이터 마이그레이션 (72%)", "detail": "레거시 데이터 정합성 문제가 주원인. Foundry Pipeline으로 사전 검증 자동화 권장."},
            {"title": "Cutover 초과 (65%)", "detail": "리허설 미흡이 직접 원인. Workshop Cutover App으로 3회 이상 리허설 의무화."},
        ],
        takeaway="Top 3 리스크(테스트·데이터·Cutover) 모두 Palantir 모듈로 직접 대응 가능 — 발생 확률 50% 이하로 억제 목표",
        footer=FOOTER,
    )
    return spec, data_narrative


# ============================================================
# 메인
# ============================================================

def build_one(builder_fn, output: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spec, pattern_func = builder_fn()
    pattern_func(slide, spec)
    prs.save(output)
    return output


def main():
    out_dir = Path("output/new_patterns")
    out_dir.mkdir(parents=True, exist_ok=True)

    cases = [
        ("06_hub_spoke", build_hub_spoke),
        ("07_before_after", build_before_after),
        ("08_kpi_dashboard", build_kpi_dashboard),
        ("09_waterfall", build_waterfall),
        ("10_swimlane", build_swimlane),
        ("11_pyramid", build_pyramid),
        ("12_data_narrative", build_data_narrative),
    ]

    print("=" * 70)
    print("7 New Patterns Test")
    print("=" * 70)

    all_passed = True
    for name, builder_fn in cases:
        out = out_dir / f"{name}.pptx"
        try:
            build_one(builder_fn, out)
        except Exception as e:
            print(f"\n[{name}] BUILD FAILED: {e}")
            import traceback
            traceback.print_exc()
            all_passed = False
            continue

        visual = validate_visual(out, convert_pdf=False)
        ok = not visual.issues
        status = "PASS" if ok else "FAIL"
        if not ok:
            all_passed = False

        print(f"\n[{name}] {status}")
        print(f"  visual: {len(visual.issues)} issues")
        for i in visual.issues[:3]:
            print(f"    - {i}")

    # PDF
    print("\n" + "=" * 70)
    print("PDF...")
    print("=" * 70)
    for name, _ in cases:
        out = out_dir / f"{name}.pptx"
        if not out.exists():
            continue
        try:
            visual = validate_visual(out, convert_pdf=True)
            print(f"  {name}.pdf  -- {'OK' if visual.pdf_available else 'SKIP'}")
        except Exception as e:
            print(f"  {name}.pdf  -- FAIL ({e})")

    print("\n" + "=" * 70)
    print(f"Result: {'ALL PASSED' if all_passed else 'SOME FAILED'}")
    print("=" * 70)
    return 0 if all_passed else 1


if __name__ == "__main__":
    sys.exit(main())
