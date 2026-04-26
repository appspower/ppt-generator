"""End-to-end 통합 테스트 — build_scenario가 실제 .pptx를 산출하는지 검증.

데이터 파일 의존성 (catalog + paragraph store + components_index)이 있으므로
누락 시 skip. 운영 환경에서는 항상 통과해야 한다.

검증 항목
--------
1. minimal narrative → .pptx 생성, 슬라이드 수 == narrative 길이
2. skeleton_id 경로 (narrative_sequence 없음)
3. chart_data 주입 시 실제 .pptx 차트의 categories가 입력값으로 교체됨
4. run.py CLI 의 --list-presets 정상 동작 (subprocess)
5. run.py CLI 의 --preset 빌드 정상 동작 (subprocess, PNG 없이)
"""
from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

CATALOG_PATH = ROOT / "output" / "catalog" / "final_labels.json"
INDEX_PATH = ROOT / "output" / "component_library" / "components_index.json"
TEMPLATE_PATH = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"

requires_data = pytest.mark.skipif(
    not (CATALOG_PATH.exists() and INDEX_PATH.exists() and TEMPLATE_PATH.exists()),
    reason="catalog / components_index / master template 누락 — 데이터 빌드 필요",
)


@pytest.fixture(scope="module")
def out_dir(tmp_path_factory) -> Path:
    return tmp_path_factory.mktemp("e2e_decks")


@requires_data
def test_e2e_minimal_narrative_build(out_dir: Path):
    """narrative_sequence + content_by_role만 주고 .pptx 생성 → 슬라이드 수 확인."""
    from pptx import Presentation

    from ppt_builder.api import build_scenario
    from ppt_builder.models import ScenarioInput

    scenario = ScenarioInput(
        scenario_name="e2e_minimal",
        narrative_sequence=["opening", "evidence", "recommendation", "closing"],
        content_by_role={
            "opening": ["e2e 테스트 발표"],
            "evidence": ["근거 1", "근거 2"],
            "recommendation": ["권고 사항 1"],
            "closing": ["마무리 멘트"],
        },
    )
    out = out_dir / "minimal.pptx"
    result = build_scenario(scenario, output=out)

    assert result.pptx.exists(), f"pptx not created: {result.pptx}"
    assert result.pptx == out
    assert len(result.plan) == 4, f"plan length != 4: {result.plan}"

    prs = Presentation(str(out))
    assert len(prs.slides) == 4, (
        f"slide count != 4: actual {len(prs.slides)}"
    )


@requires_data
def test_e2e_skeleton_id_path(out_dir: Path):
    """skeleton_id만 주고 narrative_sequence는 자동 로드."""
    from pptx import Presentation

    from ppt_builder.api import build_scenario
    from ppt_builder.models import ScenarioInput

    scenario = ScenarioInput(
        scenario_name="e2e_skel",
        skeleton_id="analysis_report_15",
        content_by_role={
            "opening": ["분기 검토"],
            "evidence": ["근거"],
        },
    )
    out = out_dir / "skel.pptx"
    result = build_scenario(scenario, output=out)

    assert out.exists()
    assert len(result.narrative) > 0
    # analysis_report_15는 보통 15 step 시퀀스 — skeleton 그대로 사용
    assert len(result.plan) >= 3
    prs = Presentation(str(out))
    assert len(prs.slides) == len(result.plan)


@requires_data
def test_e2e_chart_data_injection_visible_in_pptx(out_dir: Path):
    """chart_data 주입 + color 적용이 실제 .pptx에 반영됐는지 확인."""
    from pptx import Presentation
    from pptx.shapes.graphfrm import GraphicFrame

    from ppt_builder.api import build_scenario
    from ppt_builder.models import ChartSeriesSpec, ChartSpec, ScenarioInput

    scenario = ScenarioInput(
        scenario_name="e2e_chart",
        narrative_sequence=["opening", "evidence", "closing"],
        content_by_role={
            "opening": ["차트 테스트"],
            "evidence": ["분기별 매출 추이"],
            "closing": ["요약"],
        },
        chart_data={
            "evidence": ChartSpec(
                categories=["Q1_E2E", "Q2_E2E", "Q3_E2E", "Q4_E2E"],
                series=[
                    ChartSeriesSpec(
                        name="e2e_매출",
                        values=[111.0, 222.0, 333.0, 444.0],
                        color="#D04A02",  # PwC accent — color 주입 검증
                    )
                ],
            )
        },
    )
    out = out_dir / "chart.pptx"
    result = build_scenario(scenario, output=out)

    if not result.chart_injected.get("evidence"):
        pytest.skip(
            "evidence role에 차트 슬라이드가 매칭되지 않음 — "
            f"plan: {[(p['mode'], p.get('source')) for p in result.plan]}"
        )

    # .pptx를 다시 열어서 차트 categories + values + color 검증
    prs = Presentation(str(out))
    found = False
    for slide in prs.slides:
        for sh in slide.shapes:
            if isinstance(sh, GraphicFrame) and sh.has_chart:
                chart = sh.chart
                cats = list(chart.plots[0].categories)
                if any("Q1_E2E" in str(c) for c in cats):
                    found = True
                    series0 = chart.series[0]
                    assert list(series0.values)[0] == 111.0, (
                        f"series[0] != 111.0: {list(series0.values)}"
                    )
                    # color 적용 확인 (fill 또는 line 둘 중 하나는 #D04A02)
                    rgb_str = None
                    try:
                        rgb_str = str(series0.format.fill.fore_color.rgb).upper()
                    except Exception:
                        pass
                    if rgb_str != "D04A02":
                        try:
                            rgb_str = str(series0.format.line.color.rgb).upper()
                        except Exception:
                            pass
                    assert rgb_str == "D04A02", (
                        f"expected D04A02, got fill/line rgb={rgb_str}"
                    )
                    break
        if found:
            break

    assert found, (
        "chart_injected['evidence']=True지만 .pptx에서 Q1_E2E 카테고리 발견 못함"
    )


@requires_data
def test_e2e_cli_list_presets():
    """python run.py --list-presets 가 0 exit + preset 목록 출력."""
    res = subprocess.run(
        [sys.executable, str(ROOT / "run.py"), "--list-presets"],
        capture_output=True, text=True, encoding="utf-8",
        cwd=str(ROOT),
        env={"PYTHONIOENCODING": "utf-8", **__import__("os").environ},
        timeout=30,
    )
    assert res.returncode == 0, f"non-zero exit: {res.returncode}\n{res.stderr}"
    assert "analysis_report_15" in res.stdout, (
        f"preset 목록에 analysis_report_15 없음:\n{res.stdout}"
    )
    assert "consulting_proposal_30" in res.stdout


@requires_data
def test_e2e_cli_preset_build(out_dir: Path):
    """python run.py --preset analysis_report_15 (작은 시나리오, PNG 없이)."""
    out = out_dir / "cli_preset.pptx"
    res = subprocess.run(
        [
            sys.executable, str(ROOT / "run.py"),
            "--preset", "analysis_report_15",
            "--output", str(out),
            "--quiet",
        ],
        capture_output=True, text=True, encoding="utf-8",
        cwd=str(ROOT),
        env={"PYTHONIOENCODING": "utf-8", **__import__("os").environ},
        timeout=180,
    )
    assert res.returncode == 0, (
        f"build 실패 exit={res.returncode}\nstdout={res.stdout}\nstderr={res.stderr}"
    )
    assert out.exists(), f"output 미생성: {out}"
    # 파일 크기는 최소 50KB는 나와야 (Master 의존성으로 5MB+ 가까이 보통)
    assert out.stat().st_size > 50_000, f"output too small: {out.stat().st_size} bytes"


@requires_data
def test_e2e_cli_unknown_preset_returns_error_code():
    """존재하지 않는 preset → 0이 아닌 exit + 안내 메시지."""
    res = subprocess.run(
        [sys.executable, str(ROOT / "run.py"), "--preset", "__no_such_preset__"],
        capture_output=True, text=True, encoding="utf-8",
        cwd=str(ROOT),
        env={"PYTHONIOENCODING": "utf-8", **__import__("os").environ},
        timeout=30,
    )
    assert res.returncode != 0
    combined = res.stdout + res.stderr
    assert "preset" in combined and "없음" in combined
