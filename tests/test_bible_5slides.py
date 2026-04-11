"""바이블 워크플로우 검증 — 5개 독립 주제, 각 1장씩.

각 슬라이드는 Step 0~6 전체 프로세스를 거쳐 생성된다.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.composer import SlideComposer
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_kpi_card, comp_kpi_row, comp_native_chart, comp_bullet_list,
    comp_comparison_grid, comp_pyramid, comp_waterfall, comp_heatmap_grid,
    comp_hub_spoke_diagram, comp_chevron_flow, comp_hero_block,
    comp_styled_card,
)

OUTPUT = Path(__file__).parent.parent / "output" / "bible_5slides.pptx"


def make(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ============================================================
# Slide 1: 반도체 공급망 재편 — 미중한 3국 투자 경쟁
# Step 0: 주제=반도체 공급망, 청중=C-Level, 목적=투자 의사결정
# Step 1: Deloitte 2026, SIA, InvestKOREA, PwC 기반 리서치
# Step 2: comparison_3 + 수치 → l_layout
# Step 3: comp_comparison_grid(좌) + comp_kpi_row(우상) + comp_bullet_list(우하)
# ============================================================
def slide_01_semiconductor(prs):
    """미국·중국·한국 반도체 투자 경쟁 비교."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="미국·중국·한국이 반도체 주도권을 위해 $1조+ 규모의 투자 경쟁에 돌입했다",
        category="Semiconductor"))
    zones = comp.layout("l_layout", left_ratio=0.55, top_ratio=0.55)

    comp_comparison_grid(comp.canvas,
                         columns=[
                             {"name": "미국", "summary": "CHIPS Act 주도",
                              "criteria": [
                                  "연방 $50B + 민간 $630B (140개 프로젝트)",
                                  "TSMC $65B·Intel $7.9B·Samsung $4.7B 유치",
                                  "2030년까지 숙련 인력 11.5만명 충원 필요",
                              ]},
                             {"name": "한국", "summary": "메모리 패권 수성", "highlight": True,
                              "criteria": [
                                  "글로벌 메모리 60.5% (DRAM 70.5%, NAND 52.6%)",
                                  "월 수출 $150억 돌파 — YoY +33% (2025.8)",
                                  "파운드리 점유율 17.3% — TSMC 대비 열위",
                              ]},
                             {"name": "중국", "summary": "자국 역량 확대",
                              "criteria": [
                                  "텅스텐 글로벌 생산 79% 독점, 수출 40% 제한",
                                  "DUV 다중패터닝으로 성숙 공정 자급 확대",
                                  "한국 기초역량 추월 (KISTEP 보고서)",
                              ]},
                         ],
                         row_labels=["투자/역량", "핵심 전략", "리스크 요인"],
                         region=zones["left_full"])

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "$630B", "label": "미국 민간 투자 총액", "trend": "up"},
        {"value": "60.5%", "label": "한국 메모리 점유율", "trend": "flat"},
    ], region=zones["right_top"])

    comp_bullet_list(comp.canvas, title="한국 반도체의 전략적 시사점",
                     items=[
                         "메모리 패권은 견고하나, 파운드리(17.3%)에서 TSMC(62%)와 격차 확대",
                         "미국 CHIPS Act 보조금 수혜 중이나, 지분 취득 조건 리스크 부상",
                         "중국 소재 무기화(텅스텐·갈륨)에 대한 공급망 다변화 시급",
                         "AI 데이터센터 투자 $500B(2026) → HBM 수요 폭증이 한국에 기회",
                     ], region=zones["right_bottom"])

    comp.takeaway("메모리 패권 수성 + 파운드리 확대 투자 + 소재 공급망 다변화가 한국의 3대 과제")
    comp.footer(SlideFooter(source="Deloitte 2026 Outlook, SIA, InvestKOREA, PwC Semiconductor 2026"))


# ============================================================
# Slide 2: AI 데이터센터 에너지 위기 — 빅테크 원자력 전환
# Step 0: 주제=AI DC 에너지, 청중=인프라 임원, 목적=에너지 전략 수립
# Step 1: IEA, Gartner, Deloitte, Bloomberg 기반 리서치
# Step 2: data_kpi + comparison_3 → t_layout
# Step 3: comp_kpi_row(상) + comp_comparison_grid(하좌) + comp_bullet_list(하우)
# ============================================================
def slide_02_ai_energy(prs):
    """AI 데이터센터 에너지 소비 급증과 빅테크 원자력 전환."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="AI 데이터센터가 2026년 1,100TWh를 소비하며, 빅테크가 원자력으로 전환하고 있다",
        category="Energy & AI"))
    zones = comp.layout("t_layout", top_ratio=0.22, right_ratio=0.40)

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "1,100TWh", "label": "2026 DC 전력 소비 (IEA)", "detail": "일본 전체 전력 규모", "trend": "up"},
        {"value": "5x", "label": "AI 서버 전력 증가 (2025→2030)", "detail": "93→432 TWh", "trend": "up"},
        {"value": "$500B", "label": "2026 AI DC 투자 (글로벌)", "detail": "칩 비중 50~60%", "trend": "up"},
    ], region=zones["top"])

    comp_comparison_grid(comp.canvas,
                         columns=[
                             {"name": "Microsoft", "summary": "원자력 PPA",
                              "criteria": [
                                  "Constellation Energy 2GW PPA (2040)",
                                  "역대 최대 기업 원자력 계약",
                                  "Three Mile Island 재가동",
                              ]},
                             {"name": "Google", "summary": "차세대 SMR", "highlight": True,
                              "criteria": [
                                  "Kairos Power 500MW 용융염 원자로",
                                  "2035년 상업 운전 목표",
                                  "24/7 탄소프리 에너지 전략",
                              ]},
                             {"name": "Amazon", "summary": "대규모 SMR",
                              "criteria": [
                                  "X-energy SMR 5GW 투자",
                                  "버지니아·워싱턴 DC 근접 배치",
                                  "AWS 전력 자급 전략",
                              ]},
                         ],
                         row_labels=["핵심 계약", "규모/목표", "전략 특징"],
                         region=zones["bottom_left"])

    comp_bullet_list(comp.canvas, title="시사점 및 리스크",
                     items=[
                         "SMR 상용화는 2030년 이후 — 현재 NRC 인증은 NuScale 1건뿐",
                         "그리드 연결 대기 최대 10년 — 전력 공급이 DC 확장의 병목",
                         "한국 DC 시장도 전력 제약 직면 — 수도권 전력 할당 한계",
                         "'All of the above' 전략: 원자력+재생E+가스 병행이 현실적",
                     ], region=zones["bottom_right"])

    comp.takeaway("AI 전력 수요가 2년 만에 2.7배 증가 — 원자력이 유일한 베이스로드 해법으로 부상")
    comp.footer(SlideFooter(source="IEA Energy and AI 2026, Gartner Nov 2025, Deloitte Nuclear+DC"))


# ============================================================
# Slide 3: 글로벌 탄소시장 가격 전망 — EU ETS 가격 급등
# Step 0: 주제=탄소시장, 청중=ESG/재무 임원, 목적=탄소 비용 리스크 평가
# Step 1: Enerdata, Homaio, CarbonCredits, ABN AMRO 리서치
# Step 2: data_trend → two_column
# Step 3: comp_waterfall(좌) + comp_bullet_list(우)
# ============================================================
def slide_03_carbon_price(prs):
    """EU 탄소 가격 급등과 한국 제조업 비용 영향."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="EU 탄소 가격이 2026년 €91로 상승하며, 한국 수출 제조업의 비용 부담이 가중된다",
        category="Carbon Market"))
    zones = comp.layout("two_column", split=0.55)

    comp_waterfall(comp.canvas,
                   start={"label": "2024\n€68/t", "value": 68},
                   steps=[
                       {"label": "배출권\n공급 -8%", "value": 9},
                       {"label": "CBAM\n완전시행", "value": 7},
                       {"label": "ETS2\n확대", "value": 7},
                       {"label": "2030\n목표강화", "value": 54},
                   ],
                   end={"label": "2030E\n€145/t", "value": 145},
                   unit="€",
                   region=zones["left"])

    comp_bullet_list(comp.canvas, title="한국 제조업 시사점",
                     items=[
                         "EU ETS €91(2026) vs 자발적시장 $20 — 4.5배 격차로 이중 부담",
                         "2026년 배출권 공급 8% 감소 → 기계적 가격 상승 불가피",
                         "CBAM 완전시행(2026.1) — 한국 철강·시멘트 수출에 톤당 €80+ 부과",
                         "K-ETS ₩8,684/t는 EU €91 대비 1/8 수준 — 양 시장 가격 수렴 압력",
                         "2035년 €200 전망 — 지금 내부 탄소 가격제(ICP) 도입이 급선무",
                     ], region=zones["right"])

    comp.takeaway("탄소 가격 €145(2030)·€200(2035) 경로에서 내부 탄소 가격제 미도입 기업은 경쟁력 상실")
    comp.footer(SlideFooter(source="Enerdata Carbon Price Forecast 2026, Homaio EU ETS Guide, ABN AMRO ESG"))


# ============================================================
# Slide 4: AI 신약 개발 — R&D 비용 40% 절감과 임상 성공률 2배
# Step 0: 주제=AI 신약개발, 청중=제약 R&D 임원, 목적=AI 투자 판단
# Step 1: McKinsey, Axis Intelligence, WEF, NaturalAntibody 리서치
# Step 2: process_linear → full
# Step 3: comp_chevron_flow(show_details=True, full)
# ============================================================
def slide_04_ai_pharma(prs):
    """AI 신약 개발의 비용·시간·성공률 혁신."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="AI가 신약 개발 비용을 40% 절감하고 임상 성공률을 2배로 높이고 있다",
        category="Pharma & AI"))

    # full 레이아웃에서 main 영역을 KPI + chevron으로 분할
    zones = comp.layout("full")
    m = zones["main"]
    kpi_r = Region(m.x, m.y, m.w, 0.5)
    main_r = Region(m.x, m.y + 0.6, m.w, m.h - 0.6)

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "-40%", "label": "R&D 비용 절감", "detail": "$2B → $1.2B", "trend": "down"},
        {"value": "3~6yr", "label": "개발 기간 단축", "detail": "기존 10~15년", "trend": "down"},
        {"value": "80%", "label": "Phase I 성공률", "detail": "기존 40~65%", "trend": "up"},
        {"value": "$60~110B", "label": "연간 절감 잠재력", "detail": "McKinsey GenAI 추정", "trend": "up"},
    ], region=kpi_r)

    comp_chevron_flow(comp.canvas, phases=[
        {"tag": "표적발굴", "label": "Target ID",
         "details": ["AI 단백질 구조 예측", "AlphaFold 200M+ 구조", "발굴 기간 4년→6개월"]},
        {"tag": "후보물질", "label": "Lead Opt",
         "details": ["생성형 AI 분자 설계", "Exscientia <12개월", "전임상 비용 50% 절감"]},
        {"tag": "임상시험", "label": "Clinical",
         "details": ["AI 환자 매칭", "Phase I 성공률 80~90%", "전통 40~65% 대비 2배"]},
        {"tag": "허가/출시", "label": "Approval",
         "details": ["2026~27 FDA 첫 승인 60%", "173개 AI 파이프라인", "시장 $16.5B(2034)"]},
    ], show_details=True, style="gradient", region=main_r)

    comp.takeaway("AI 신약 173개 파이프라인 중 2026~27 첫 FDA 승인 예상 — 제약 R&D 패러다임 전환의 변곡점")
    comp.footer(SlideFooter(source="McKinsey GenAI 2025, Axis Intelligence 2026, WEF Drug Discovery, AllAboutAI"))


# ============================================================
# Slide 5: 우주산업 상업화 — Space Economy $1조 시대
# Step 0: 주제=우주산업, 청중=신사업/투자 임원, 목적=투자 기회 평가
# Step 1: Space Foundation, PwC Space, SpaceNexus, SIA 리서치
# Step 2: data_hero + structure → grid_2x2
# Step 3: comp_styled_card × 4 (시장규모/발사/위성통신/한국기회)
# ============================================================
def slide_05_space_economy(prs):
    """우주산업 $626B → $1조 성장과 핵심 투자 영역."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="우주산업이 $626B에서 2032년 $1조로 성장하며, 위성통신이 핵심 동력이다",
        category="Space Economy"))
    zones = comp.layout("grid_2x2")

    comp_styled_card(comp.canvas, title="시장 규모",
                     kpi_value="$626B → $1T",
                     bullets=[
                         "2025년 $626B — 전년 대비 +2.1%",
                         "상업 부문 78%, 정부 22%",
                         "Morgan Stanley·BoA 모두 2032 $1T+ 전망",
                     ], style="dark", region=zones["tl"])

    comp_styled_card(comp.canvas, title="발사 혁명",
                     kpi_value="28hr/1회",
                     bullets=[
                         "2025 H1: 329회 궤도 발사 (321 성공)",
                         "4,517개 위성 배치 — 역대 최다",
                         "SpaceX 81회(전체 54%) 독주",
                     ], style="subtle", region=zones["tr"])

    comp_styled_card(comp.canvas, title="위성통신 = 핵심 동력",
                     kpi_value="$10.4B",
                     bullets=[
                         "Starlink 9,500기 — 전체 위성의 65%",
                         "Amazon Kuiper·OneWeb 경쟁 본격화",
                         "저궤도 광대역이 매출의 핵심 엔진",
                     ], style="light", region=zones["bl"])

    comp_styled_card(comp.canvas, title="한국 기회 영역",
                     number="KR",
                     bullets=[
                         "누리호 4차 성공 — 소형위성 시장 진입",
                         "한화에어로·KAI 방산 수출 연계",
                         "6G 위성통신 R&D — NTN 선점",
                     ], style="numbered", region=zones["br"])

    comp.takeaway("위성통신(Starlink $10.4B)이 우주산업 성장 엔진 — 한국은 발사체·NTN·방산 3축 전략 필요")
    comp.footer(SlideFooter(source="Space Foundation Q2 2025, PwC Space Trends, SIA 28th Report, SpaceNexus"))


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_01_semiconductor(prs)
    slide_02_ai_energy(prs)
    slide_03_carbon_price(prs)
    slide_04_ai_pharma(prs)
    slide_05_space_economy(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")

    # Step 5-A: EVALUATE
    from ppt_builder.evaluate import evaluate_pptx, print_report
    report = evaluate_pptx(str(OUTPUT))
    print_report(report)


if __name__ == "__main__":
    main()
