"""바이블 워크플로우 2차 실증 — 5개 독립 주제, 각 1장.

컴포넌트 밀도 개선(자동 피팅+볼드 마크업) 반영.
산출물: PPT + PDF + PNG 3종 필수.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.composer import SlideComposer
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_kpi_card, comp_kpi_row, comp_bullet_list,
    comp_comparison_grid, comp_waterfall, comp_heatmap_grid,
    comp_chevron_flow, comp_styled_card,
)

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "bible_round2"


def make(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ============================================================
# S1: 희토류 공급 위기 — 중국 60% 생산, 90% 정제 독점
# Layout: two_column | heatmap_grid(좌) + bullet_list(우)
# ============================================================
def s1_rare_earth(prs):
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="중국이 희토류 생산 60%·정제 90%를 독점하며, 서방의 탈중국 공급망 구축이 시급하다",
        category="Critical Minerals"))
    zones = comp.layout("two_column", split=0.50)

    comp_heatmap_grid(comp.canvas,
                      row_labels=["네오디뮴(Nd)", "디스프로슘(Dy)", "테르븀(Tb)", "스칸듐(Sc)", "갈륨(Ga)"],
                      col_labels=["중국 점유율", "대체 공급 가능성", "서방 정제 역량"],
                      values=[
                          [0.7, 0.4, 0.3],
                          [0.9, 0.8, 0.9],
                          [0.9, 0.9, 0.9],
                          [0.6, 0.5, 0.7],
                          [0.8, 0.6, 0.6],
                      ],
                      cell_texts=[
                          ["70%", "호주·미국 확대중", "MP Materials"],
                          ["90%", "매우 제한적", "거의 전무"],
                          ["95%", "거의 불가", "전무"],
                          ["60%", "필리핀·호주", "초기 단계"],
                          ["80%", "일본·독일 개발중", "제한적"],
                      ],
                      region=zones["left"])

    comp_bullet_list(comp.canvas, title="서방 탈중국 전략 현황 (2026)",
                     items=[
                         "**미국 DoD $4억 투자** — MP Materials에 Mountain Pass 정제 확장 + 자석 공장 건설. 2026년 국내 완전 통합 자석 공급망 목표",
                         "**호주 Lynas USA** — 텍사스 정제 시설 건설 중. 미국 내 유일한 중국 외 중·희토류 정제 후보",
                         "**EU 핵심원자재법(CRMA)** — 2030년까지 역내 정제 40%, 재활용 25% 달성 목표. 단, 2026년 기준 서방 중희토류 상용 정제 시설 0개",
                         "**가격 충격** — 2025년 수출 통제 시 유럽 희토류 가격 중국 대비 6배 급등. 자동차·방산 업계 가동률 감소 불가피",
                         "**한국 영향** — 전기차 모터·방산 부품의 네오디뮴 자석 100% 수입 의존. 재활용 기술 + 동맹국 다변화가 유일한 중기 해법",
                     ], region=zones["right"])

    comp.takeaway("2026년 기준 서방의 중희토류 상용 정제 시설은 0개 — 탈중국은 5~10년 프로젝트이며 단기 충격 불가피")
    comp.footer(SlideFooter(source="IEA Critical Minerals 2025, CSIS Rare Earth Analysis, EU CRMA, Clark Hill"))


# ============================================================
# S2: 글로벌 물 위기 — $7조 투자 갭
# Layout: l_layout | waterfall(좌) + kpi_row(우상) + bullet_list(우하)
# ============================================================
def s2_water_crisis(prs):
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="글로벌 물 인프라에 연간 $4,350억이 부족하며, 투자 갭이 $7조에 달한다",
        category="Water Security"))
    zones = comp.layout("l_layout", left_ratio=0.55, top_ratio=0.45)

    comp_waterfall(comp.canvas,
                   start={"label": "현재 투자\n(연간)", "value": 480},
                   steps=[
                       {"label": "노후 교체\n필요", "value": 150},
                       {"label": "인구 증가\n대응", "value": 100},
                       {"label": "기후 적응\n투자", "value": 85},
                       {"label": "위생·수질\n개선", "value": 100},
                   ],
                   end={"label": "필요 투자\n(연간)", "value": 915},
                   unit="B",
                   region=zones["left_full"])

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "$7T", "label": "글로벌 투자 갭 (World Bank)", "trend": "up"},
        {"value": "2.2B", "label": "안전한 물 접근 부족 인구", "trend": "up"},
    ], region=zones["right_top"])

    comp_bullet_list(comp.canvas, title="투자 확대 전략",
                     items=[
                         "**민간 자본 참여 확대** — 현재 물 투자의 91%가 공공 부문, 민간 2% 미만. 블렌디드 파이낸스로 민간 참여 유도 시급",
                         "**ROI 근거 확보** — 물 인프라 €1 투자당 €1.30 GVA 창출, 투자 갭 해소 시 2억 600만 일자리 + €8.4T GDP 성장",
                         "**디지털 물 관리** — AI 기반 누수 탐지(NRW 30%→10%)로 기존 인프라 효율화. 신규 투자 30% 절감 효과",
                         "**한국 기회** — K-water 기술 수출(스마트 물 관리), 해수담수화 기술 중동·아프리카 시장 진출",
                     ], region=zones["right_bottom"])

    comp.takeaway("물 인프라 €1 투자당 €1.30 경제 효과 — 투자 갭 해소가 2억 일자리 창출의 기회")
    comp.footer(SlideFooter(source="World Bank Water 2025, WEF Water Investment, Cambridge-Acea Report"))


# ============================================================
# S3: 자율주행 상용화 — Waymo 2,500대, Level 4 시대 개막
# Layout: full | kpi_row(상) + chevron_flow(하, show_details)
# ============================================================
def s3_autonomous_vehicle(prs):
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="Waymo가 주 25만회 유료 운행을 달성하며, Level 4 자율주행의 상용화가 시작됐다",
        category="Autonomous Vehicles"))
    zones = comp.layout("full")
    m = zones["main"]
    kpi_r = Region(m.x, m.y, m.w, 0.5)
    main_r = Region(m.x, m.y + 0.6, m.w, m.h - 0.6)

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "2,500대", "label": "Waymo 미국 내 로보택시", "trend": "up"},
        {"value": "250K/주", "label": "주간 유료 운행 횟수", "trend": "up"},
        {"value": "$719B", "label": "2033년 시장 규모", "trend": "up"},
        {"value": "38.2%", "label": "Level 4 CAGR (2026-30)", "trend": "up"},
    ], region=kpi_r)

    comp_chevron_flow(comp.canvas, phases=[
        {"tag": "L2+", "label": "ADAS 보급",
         "details": ["Tesla Autopilot·현대 HDA2 등 양산 보급", "운전자 상시 감독 필수, 부분 자동화", "2025년 신차 70%+ 탑재"]},
        {"tag": "L4", "label": "로보택시", "highlight": True,
         "details": ["Waymo 7개 도시, 2026년 20개 도시 목표", "Tesla 오스틴 로보택시 2025.6 런칭", "Aurora 무인 화물트럭 텍사스 상용 운행"]},
        {"tag": "규제", "label": "법제도 정비",
         "details": ["NHTSA AV STEP: 핸들·페달 없는 차량 허용 규정", "주(州)별 규제 불일치가 확장 병목", "한국 RideFlux L4 파일럿 2026.3 확인"]},
        {"tag": "생태계", "label": "시장 재편",
         "details": ["GM Cruise 철수 → 시장 통합 가속", "Waymo-Toyota 전략 제휴 (플랫폼 공유)", "중국 Baidu Apollo Go·WeRide 해외 진출"]},
    ], show_details=True, style="gradient", region=main_r)

    comp.takeaway("L4 로보택시가 2026년 20개 도시로 확장 — 한국은 규제 샌드박스 속도가 시장 선점을 결정")
    comp.footer(SlideFooter(source="Waymo Blog 2025, Technavio L4 Market, GreenCars AV Report, NHTSA AV STEP"))


# ============================================================
# S4: ESG 공시 규제 — EU CSRD vs 미국 vs 한국 비교
# Layout: t_layout | kpi_row(상) + comparison_grid(하좌) + bullet_list(하우)
# ============================================================
def s4_esg_regulation(prs):
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="EU CSRD 2년 유예에도 ESG 공시 의무화는 불가역적이며, 한국 기업의 준비가 시급하다",
        category="ESG Regulation"))
    zones = comp.layout("t_layout", top_ratio=0.12, right_ratio=0.38)

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "FY2027", "label": "EU CSRD 적용 (대기업)", "trend": "flat"},
        {"value": "FY2026", "label": "미국 CA SB253 Scope1·2", "trend": "flat"},
        {"value": "FY2025", "label": "한국 KSSB 자산 2조+", "trend": "flat"},
    ], region=zones["top"])

    comp_comparison_grid(comp.canvas,
                         columns=[
                             {"name": "EU (CSRD)", "summary": "FY2027 대기업",
                              "criteria": [
                                  "**Omnibus 간소화로 대상 90% 축소** — 직원 1,000명+ 대기업만 적용. 상장 SME 면제. 제3국 기업은 FY2028부터 (매출 €1.5억+)",
                                  "**ESRS 이중 중요성 평가 의무** — 재무 영향 + 사회·환경 영향 모두 보고. 제한적 인증(Limited Assurance) 2026년 기준 도입",
                                  "**2년 유예(Stop-the-Clock)** — FY2025→FY2027 지연. 선행 기업은 이미 준비 완료 — 유예를 기회로 활용 vs 지연 리스크",
                              ]},
                             {"name": "미국", "summary": "주(州)별 분절", "highlight": True,
                              "criteria": [
                                  "**SEC 연방 규정 계류 중** — 캘리포니아 SB253(매출 $10억+)이 사실상 기준. Scope 1·2는 FY2025, Scope 3는 2027",
                                  "**SB261 기후 재무 리스크 공시** — 매출 $5억+ 의무. 제9순회항소법원 2025.11 시행 유예 — 법적 불확실성 지속",
                                  "**뉴욕주 2027.6 GHG 의무화 확정** — 주별 규제 파편화로 다주 운영 기업은 복수 프레임워크 대응 필요",
                              ]},
                             {"name": "한국 (KSSB)", "summary": "자산 2조+ 의무",
                              "criteria": [
                                  "**2025년 자산 2조+ 상장사 의무** — KSSB 기준은 ISSB S1·S2 기반. 글로벌 호환성 확보",
                                  "**2027년 5,000억+ 확대** — 2030년 전체 코스피 적용 예정. 중견기업은 3~5년 준비 기간 보유",
                                  "**단일 중요성 기반** — EU CSRD와 달리 Financial Materiality만 적용. 이중 중요성 전환 여부가 향후 쟁점",
                              ]},
                         ],
                         row_labels=["적용 범위·일정", "보고 기준·인증", "핵심 쟁점"],
                         region=zones["bottom_left"])

    comp_bullet_list(comp.canvas, title="한국 기업 대응 체크리스트",
                     items=[
                         "**KSSB 의무 보고 착수** — 자산 2조+ 기업은 FY2025 데이터 수집 이미 시작. 자산 5,000억+ 기업은 2027 대비 시스템 구축 필요",
                         "**Scope 3 데이터 확보** — 공급망 배출 데이터가 가장 큰 난관. 1차 협력사부터 단계적 수집 체계 구축 권장",
                         "**EU CSRD 연동 준비** — 유럽 수출 기업은 FY2028 제3국 기준 대비 필수. ESRS 이중 중요성 평가 역량 내재화",
                         "**ESG 데이터 인프라** — 수기 수집→자동화 플랫폼 전환. AI 기반 ESG 데이터 분석으로 보고 품질+효율 확보",
                     ], region=zones["bottom_right"])

    comp.takeaway("EU 2년 유예는 '면제'가 아닌 '준비 기간' — 선행 기업은 이미 FY2027 보고 체계 완료, 지금 시작하지 않으면 규제 리스크 직면")
    comp.footer(SlideFooter(source="EU CSRD Omnibus 2025, CA SB253/261, KSSB, Harvard Law ESG Review"))


# ============================================================
# S5: 양자컴퓨팅 시장 — $1B 돌파, 산업별 적용 확산
# Layout: grid_2x2 | styled_card × 4
# ============================================================
def s5_quantum_computing(prs):
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="양자컴퓨팅 시장이 $1B를 돌파하며, 금융·제약·소재 분야에서 상용 적용이 시작됐다",
        category="Quantum Computing"))
    zones = comp.layout("grid_2x2")

    comp_styled_card(comp.canvas, title="시장 현황",
                     kpi_value="$1B+ (2025)",
                     bullets=[
                         "2024년 $650~750M → 2025년 $1B 돌파",
                         "2030년 $20.2B 전망 (CAGR 41.8%)",
                         "QaaS(서비스형 양자) 모델로 진입장벽 하락",
                         "IBM·Google·Microsoft 3강 + 스타트업 생태계",
                     ], style="dark", region=zones["tl"])

    comp_styled_card(comp.canvas, title="기술 이정표",
                     kpi_value="4,158 큐빗",
                     bullets=[
                         "IBM Kookaburra: 3칩 연결 4,158큐빗 (2025)",
                         "Google Willow: 최초 오류 임계점 돌파 (2024.12)",
                         "큐빗 수 증가 → 오류율 감소 입증 = 확장성 증명",
                         "양자 우위(Quantum Advantage) 실용 영역 진입 시작",
                     ], style="subtle", region=zones["tr"])

    comp_styled_card(comp.canvas, title="산업 적용 사례",
                     number="USE",
                     bullets=[
                         "**금융(25% 점유)** — JPMorgan+IBM: 옵션 가격 산정·리스크 분석에서 몬테카를로 능가",
                         "**제약** — 분자 시뮬레이션으로 신약 후보물질 탐색 10배 가속",
                         "**소재** — 배터리·촉매 신소재 설계, 양자 화학 시뮬레이션",
                         "**물류** — 경로 최적화·공급망 스케줄링 NP-hard 문제 접근",
                     ], style="light", region=zones["bl"])

    comp_styled_card(comp.canvas, title="한국 전략 시사점",
                     number="KR",
                     bullets=[
                         "**KIST·KAIST 양자연구** — 국내 양자 컴퓨터 독자 개발 추진 중이나 글로벌 대비 3~5년 격차",
                         "**IBM Quantum Network 참여** — 삼성·SK 등 QaaS 활용 PoC 단계. 금융·반도체 시뮬레이션 우선",
                         "**양자 인력 양성** — 글로벌 양자 인력 수요 대비 공급 부족. 석박사급 연구 인력 확보 경쟁 심화",
                         "**양자 내성 암호(PQC)** — NIST PQC 표준 2024 확정. 금융·공공 시스템 전환 로드맵 수립 필요",
                     ], style="numbered", region=zones["br"])

    comp.takeaway("양자컴퓨팅 $1B 시장 진입 — 한국은 QaaS 기반 PoC와 양자내성암호 전환을 동시에 추진해야")
    comp.footer(SlideFooter(source="MarketsandMarkets QC 2025, InsightAce Analytics, SpinQ Industry Trends, PatentPC"))


# ============================================================
# 메인 실행 — PPT + PDF + PNG 3종 생성
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s1_rare_earth(prs)
    s2_water_crisis(prs)
    s3_autonomous_vehicle(prs)
    s4_esg_regulation(prs)
    s5_quantum_computing(prs)

    # Step 4: GENERATE — PPT
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path}")

    # Step 5-A: EVALUATE
    from ppt_builder.evaluate import evaluate_pptx, print_report
    report = evaluate_pptx(str(pptx_path))
    print_report(report)

    # Step 4 산출물: PDF
    import pythoncom, win32com.client
    pythoncom.CoInitialize()
    pdf_path = pptx_path.with_suffix(".pdf")
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    p = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
    p.SaveAs(str(pdf_path.resolve()), 32)
    p.Close()
    print(f"PDF:  {pdf_path}")

    # Step 4 산출물: PNG
    from ppt_builder.track_c.png_export import pptx_to_pngs
    png_dir = OUTPUT_DIR / f"{NAME}_pngs"
    paths = pptx_to_pngs(pptx_path, png_dir)
    print(f"PNGs: {png_dir} ({len(paths)}장)")


if __name__ == "__main__":
    main()
