"""바이블 워크플로우 단일 슬라이드 테스트.

산출물: .pptx + .pdf + 슬라이드별 .png
바이블 Step 4(GENERATE) → Step 5(EVALUATE + 시각검증) 전체 수행.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.composer import SlideComposer
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_kpi_row, comp_comparison_grid, comp_bullet_list,
    comp_heatmap_grid, comp_chevron_flow,
)

OUTPUT_DIR = Path(__file__).parent.parent / "output"


def generate_pptx(prs, name: str) -> Path:
    """PPTX 저장."""
    pptx_path = OUTPUT_DIR / f"{name}.pptx"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path}")
    return pptx_path


def generate_pdf(pptx_path: Path) -> Path:
    """PowerPoint COM으로 PDF 변환."""
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    pdf_path = pptx_path.with_suffix(".pdf")
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    pres = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
    pres.SaveAs(str(pdf_path.resolve()), 32)  # ppSaveAsPDF
    pres.Close()
    print(f"PDF:  {pdf_path}")
    return pdf_path


def generate_pngs(pptx_path: Path) -> list[Path]:
    """PowerPoint COM으로 슬라이드별 PNG 추출."""
    from ppt_builder.track_c.png_export import pptx_to_pngs

    png_dir = pptx_path.parent / f"{pptx_path.stem}_pngs"
    paths = pptx_to_pngs(pptx_path, png_dir)
    print(f"PNGs: {png_dir} ({len(paths)}장)")
    return paths


def evaluate_and_report(pptx_path: Path) -> dict:
    """Step 5-A: evaluate.py 자동 평가."""
    from ppt_builder.evaluate import evaluate_pptx, print_report

    report = evaluate_pptx(str(pptx_path))
    print_report(report)
    return report


# ============================================================
# 슬라이드 정의
# ============================================================

def make(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_global_supply_chain_risk(prs):
    """글로벌 공급망 리스크 — 지정학·기후·AI 3중 위협.

    Step 0: 주제=공급망 리스크, 청중=COO/SCM 임원, 목적=공급망 전략 수립
    Step 1: WEF Global Risks 2026, McKinsey Supply Chain, Gartner SCM
    Step 2: comparison + data → two_column (heatmap + bullet)
    Step 3: comp_heatmap_grid(좌) + comp_bullet_list(우)
    """
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="지정학·기후·AI 3중 위협이 글로벌 공급망의 연간 $1.6조 손실을 야기하고 있다",
        category="Supply Chain Risk"))
    zones = comp.layout("two_column", split=0.55)

    comp_heatmap_grid(comp.canvas,
                      row_labels=["지정학 리스크", "기후 재해", "AI/사이버", "원자재 집중", "인력 부족"],
                      col_labels=["영향도", "발생확률", "대비수준"],
                      values=[
                          [0.9, 0.8, 0.3],   # 지정학: 영향 극대, 확률 높음, 대비 미흡
                          [0.8, 0.9, 0.2],   # 기후: 빈도 급증, 대비 최저
                          [0.7, 0.7, 0.4],   # AI/사이버: 신규 위협
                          [0.6, 0.5, 0.5],   # 원자재: 중간
                          [0.5, 0.8, 0.3],   # 인력: 확률 높지만 영향 중간
                      ],
                      cell_texts=[
                          ["Critical", "High", "Low"],
                          ["Critical", "Critical", "Very Low"],
                          ["High", "High", "Medium"],
                          ["High", "Medium", "Medium"],
                          ["Medium", "High", "Low"],
                      ],
                      region=zones["left"])

    comp_bullet_list(comp.canvas, title="COO를 위한 5대 공급망 전략 (2026)",
                     items=[
                         "**다중 소싱 의무화** — 단일 국가 의존도 40% 상한 설정. 미중 디커플링으로 중국 의존 반도체·배터리·희토류 공급망 재편 필수. Friend-shoring 전략 병행",
                         "**기후 복원력 내장** — 주요 거점별 기후 리스크 시나리오(홍수·가뭄·폭염) 연 2회 모의. 2025년 유럽 홍수로 자동차 부품 6주 중단 사례 — 안전재고 2주→4주 확대",
                         "**AI Control Tower 구축** — 실시간 Tier 1~3 공급업체 모니터링. McKinsey: AI 도입 기업의 공급 중단 대응 속도 65% 향상. 예측형 리스크 알림으로 선제 대응",
                         "**사이버 공급망 보안** — 공급업체 1,200개 중 보안 평가 완료 34%에 불과(WEF). SBOM 의무화 + 3자 보안 SLA 계약 조건 삽입. 분기별 자동 스캔 전환",
                         "**인력 리스킬링 + RPA** — 물류·조달의 디지털 역량 전환. RPA로 발주·검수·정산 자동화, 인력 30% 효율화. 잔여 인력은 전략적 소싱·리스크 분석에 집중",
                     ], region=zones["right"])

    comp.takeaway("연간 $1.6조 공급망 손실의 80%는 예측 가능한 리스크 — AI Control Tower + Multi-sourcing이 핵심 방어선")
    comp.footer(SlideFooter(source="WEF Global Risks 2026, McKinsey Supply Chain Pulse 2025, Gartner SCM Survey"))


# ============================================================
# 메인 실행
# ============================================================

def main():
    name = "bible_supply_chain"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_global_supply_chain_risk(prs)

    # Step 4: GENERATE
    pptx_path = generate_pptx(prs, name)

    # Step 5-A: EVALUATE
    report = evaluate_and_report(pptx_path)

    # Step 5-B: 시각 검증용 PDF + PNG 생성
    generate_pdf(pptx_path)
    generate_pngs(pptx_path)

    print(f"\n산출물:")
    print(f"  PPT: output/{name}.pptx")
    print(f"  PDF: output/{name}.pdf")
    print(f"  PNG: output/{name}_pngs/")


if __name__ == "__main__":
    main()
