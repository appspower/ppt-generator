"""Blueprint Phase 1 테스트 — 새 레이아웃 + 컴포넌트 검증.

3가지 복합 조합 패턴을 생성하여 시각 검증:
  슬라이드 1: CP2 — 번호 그리드 (grid_nxm + numbered_cell)
  슬라이드 2: CP3 — 타임라인 밴드 (timeline_band + timeline_marker + bullet_list)
  슬라이드 3: CP4 — 이종 패널 (l_layout + native_chart + kpi_card + styled_card)
  슬라이드 4: CP1 변형 — 중앙+4방향 (center_peripheral_4 + bullet_list)
  슬라이드 5: T자 레이아웃 (t_layout + kpi_row + bullet_list)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.composer import SlideComposer
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_numbered_cell,
    comp_timeline_marker,
    comp_icon_header_card,
    comp_kpi_card,
    comp_kpi_row,
    comp_bullet_list,
    comp_styled_card,
)

OUTPUT = Path(__file__).parent.parent / "output" / "blueprint_phase1_test.pptx"


def make_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def slide_1_numbered_grid(prs: Presentation):
    """CP2: 3×2 번호 그리드."""
    slide = make_slide(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="6-Step Process Framework",
        category="Methodology",
    ))

    zones = comp.layout("grid_nxm", rows=2, cols=3, gap=0.0)

    steps = [
        ("00", "Discover", "현행 프로세스 분석 및 개선 기회 발굴"),
        ("01", "Define", "목표 프로세스 정의 및 KPI 설정"),
        ("02", "Design", "솔루션 아키텍처 및 상세 설계"),
        ("03", "Develop", "시스템 구현 및 단위 테스트"),
        ("04", "Deploy", "통합 테스트 및 Go-Live 실행"),
        ("05", "Deliver", "안정화 및 지속적 개선"),
    ]
    colors = ["grey_200", "grey_400", "grey_700", "accent_mid", "accent", "accent"]

    for i, (num, hdr, body) in enumerate(steps):
        ri, ci = divmod(i, 3)
        comp_numbered_cell(
            comp.canvas,
            number=num, header=hdr, body=body,
            bg_color=colors[i],
            region=zones[f"r{ri}c{ci}"],
        )

    comp.footer(SlideFooter())


def slide_2_timeline(prs: Presentation):
    """CP3: 타임라인 밴드 + 상하 교차 콘텐츠."""
    slide = make_slide(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="디지털 전환 로드맵 2024-2028",
        category="Roadmap",
    ))

    zones = comp.layout("timeline_band", steps=5, band_ratio=0.07)

    # 밴드 마커
    comp_timeline_marker(
        comp.canvas,
        labels=["2024", "2025", "2026", "2027", "2028"],
        style="bar",
        highlight_idx=2,
        region=zones["band"],
    )

    # 각 스텝 콘텐츠
    step_data = [
        ("AS-IS 분석", ["현행 시스템 진단", "Gap 분석 완료"]),
        ("Foundation", ["ERP Core 구축", "데이터 마이그레이션"]),
        ("Optimization", ["프로세스 자동화", "AI/ML 파일럿"]),
        ("Scale-Up", ["전사 확산", "CoE 설립"]),
        ("Innovation", ["자율 운영", "Predictive Analytics"]),
    ]
    for i, (hdr, items) in enumerate(step_data):
        comp_bullet_list(
            comp.canvas,
            title=hdr,
            items=items,
            region=zones[f"step_{i}"],
        )

    comp.takeaway("2026년 Optimization 단계가 전환의 핵심 변곡점")
    comp.footer(SlideFooter())


def slide_3_heterogeneous(prs: Presentation):
    """CP4: 이종 패널 (L자 레이아웃)."""
    slide = make_slide(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="사업 성과 종합 분석",
        category="Performance",
    ))

    zones = comp.layout("l_layout", left_ratio=0.45, top_ratio=0.55)

    # 좌측: 차트
    from ppt_builder.components import comp_native_chart
    comp_native_chart(
        comp.canvas,
        chart_type="vertical_bar",
        chart_kwargs={
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "values": [120, 145, 160, 180],
            "highlight_idx": 3,
            "series_name": "2025 매출",
        },
        region=zones["left_full"],
    )

    # 우측 상단: KPI
    comp_kpi_card(
        comp.canvas,
        value="20%",
        label="YoY 성장률",
        detail="전년 대비 +180억원",
        trend="up",
        region=zones["right_top"],
    )

    # 우측 하단: 스타일 카드
    comp_styled_card(
        comp.canvas,
        title="핵심 성과",
        body="디지털 전환 투자 ROI 340% 달성.\n신규 고객 획득 2.3배 증가.\n운영 비용 15% 절감.",
        style="accent",
        region=zones["right_bottom"],
    )

    comp.footer(SlideFooter())


def slide_4_center_peripheral(prs: Presentation):
    """CP1 변형: center_peripheral_4 + bullet_list."""
    slide = make_slide(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="4대 전략 방향",
        category="Strategy",
    ))

    zones = comp.layout("center_peripheral_4", center_ratio=0.35)

    # 중앙에 큰 메시지
    c = comp.canvas
    cr = zones["center"]
    c.box(x=0, y=0, w=cr.w, h=cr.h,
          fill="accent", border=None, region=cr)
    c.text("Digital\nTransformation\n2026", x=0.1, y=cr.h * 0.2,
           w=cr.w - 0.2, h=cr.h * 0.6,
           size=14, bold=True, color="white", anchor="middle", region=cr)

    # 4방향 텍스트
    directions = {
        "left": ("Process Innovation", [
            "RPA 기반 업무 자동화",
            "E2E 프로세스 재설계",
            "Lean Six Sigma 적용",
        ]),
        "right": ("Data & Analytics", [
            "통합 데이터 플랫폼",
            "AI/ML 분석 내재화",
            "실시간 의사결정 지원",
        ]),
        "top": ("Technology Platform", [
            "Cloud-Native 전환",
            "마이크로서비스 아키텍처",
        ]),
        "bottom": ("People & Culture", [
            "디지털 역량 교육",
            "Change Management",
        ]),
    }
    for pos, (hdr, items) in directions.items():
        comp_bullet_list(
            comp.canvas,
            title=hdr,
            items=items,
            region=zones[pos],
        )

    comp.footer(SlideFooter())


def slide_5_t_layout(prs: Presentation):
    """T자 레이아웃: 상단 KPI + 하단 좌우."""
    slide = make_slide(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="프로젝트 현황 대시보드",
        category="PMO",
    ))

    zones = comp.layout("t_layout", top_ratio=0.30, right_ratio=0.45)

    # 상단: KPI 행
    comp_kpi_row(
        comp.canvas,
        kpis=[
            {"value": "87%", "label": "일정 준수율", "trend": "up"},
            {"value": "94%", "label": "품질 달성률", "trend": "up"},
            {"value": "₩12.3B", "label": "집행 예산", "trend": "flat"},
            {"value": "3건", "label": "미해결 이슈", "trend": "down"},
        ],
        region=zones["top"],
    )

    # 하단 좌: 진행 사항
    comp_bullet_list(
        comp.canvas,
        title="이번 주 주요 진행사항",
        items=[
            "SAP S/4HANA 개발 서버 구축 완료",
            "MM/PP 모듈 Fit-Gap 분석 1차 완료",
            "데이터 마이그레이션 매핑 규칙 확정",
            "통합 테스트 시나리오 45건 작성",
        ],
        region=zones["bottom_left"],
    )

    # 하단 우: 리스크
    comp_bullet_list(
        comp.canvas,
        title="주요 리스크 및 대응",
        items=[
            "[HIGH] 레거시 인터페이스 복잡도 → 전담 팀 배정",
            "[MED] 현업 참여도 저조 → 주간 워크숍 강화",
            "[LOW] 테스트 환경 지연 → 클라우드 대안 검토",
        ],
        region=zones["bottom_right"],
    )

    comp.takeaway("일정 87% 준수 중이나 레거시 인터페이스 리스크 집중 관리 필요")
    comp.footer(SlideFooter())


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_1_numbered_grid(prs)
    slide_2_timeline(prs)
    slide_3_heterogeneous(prs)
    slide_4_center_peripheral(prs)
    slide_5_t_layout(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
