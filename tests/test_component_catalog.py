"""컴포넌트 카탈로그 PPT — 42개 실물 렌더링. PPT+PDF+PNG."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from ppt_builder.primitives import Canvas, Region
from ppt_builder.components import *

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "component_catalog"
R = Region(0.3, 1.2, 9.4, 5.8)


def make(prs, title, cat):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    c.text(cat, x=0.3, y=0.15, w=9.4, h=0.25, size=9, color="grey_700", anchor="top")
    c.text(title, x=0.3, y=0.38, w=9.4, h=0.40, size=16, bold=True, color="grey_900", anchor="top")
    c.line(x1=0.3, y1=0.85, x2=9.7, y2=0.85, color="accent", width=2.0)
    return c


def safe(prs, title, cat, fn):
    c = make(prs, title, cat)
    try:
        fn(c, R)
    except Exception as e:
        c.text(f"ERROR: {type(e).__name__}: {str(e)[:150]}",
               x=0.3, y=2.0, w=9.4, h=1.5, size=9, color="negative", anchor="top")


def build(prs):
    r = R

    safe(prs, "comp_kpi_card (3개 변형)", "Atomic", lambda c, r: [
        comp_kpi_card(c, value="25%", label="탄소 배출 비중", detail="~9 Gt CO₂", trend="down", region=Region(r.x, r.y, 3.0, 1.5)),
        comp_kpi_card(c, value="$4.5T", label="필요 투자", detail="연간 IEA", trend="up", region=Region(r.x+3.3, r.y, 3.0, 1.5)),
        comp_kpi_card(c, value="41%", label="감축 목표", trend="flat", region=Region(r.x+6.6, r.y, 3.0, 1.5)),
    ])

    safe(prs, "comp_kpi_row (4개 가로)", "Atomic", lambda c, r:
        comp_kpi_row(c, kpis=[
            {"value": "$10.5T", "label": "사이버범죄 비용", "trend": "up"},
            {"value": "+47%", "label": "AI 공격 증가율", "trend": "up"},
            {"value": "4M", "label": "보안 인력 부족", "trend": "up"},
            {"value": "$520B", "label": "보안 지출 2026E", "trend": "up"},
        ], region=Region(r.x, r.y, r.w, 1.5)))

    safe(prs, "comp_mini_table", "Atomic", lambda c, r:
        comp_mini_table(c, headers=["구분", "2024", "2025", "2026E"],
                        rows=[["매출", "₩500B", "₩620B", "₩780B"], ["영업이익", "₩45B", "₩68B", "₩95B"], ["이익률", "9.0%", "11.0%", "12.2%"]],
                        region=Region(r.x, r.y, r.w, 2.0)))

    safe(prs, "comp_bullet_list (**볼드** 지원)", "Atomic", lambda c, r:
        comp_bullet_list(c, title="CISO 우선순위", items=[
            "**사이버 복원력** — 공격 후 복구 속도가 핵심",
            "**제로트러스트** — Never Trust, Always Verify",
            "**AI 보안 자동화** — 인력 400만 부족 보완",
            "**클라우드 거버넌스** — 설정 오류가 침해 1위",
            "**공급망 보안** — SBOM 의무화 + 3자 SLA",
        ], region=Region(r.x, r.y, r.w*0.55, r.h)))

    safe(prs, "comp_bar_chart_h", "Atomic", lambda c, r:
        comp_bar_chart_h(c, data=[
            {"label": "철강", "value": 3.3, "highlight": True}, {"label": "시멘트", "value": 2.5},
            {"label": "화학", "value": 1.0}, {"label": "자동차", "value": 0.5},
        ], unit="Gt CO₂", region=Region(r.x, r.y, r.w*0.5, 3.0)))

    safe(prs, "comp_stat_row", "Atomic", lambda c, r:
        comp_stat_row(c, stats=[
            {"label": "매출", "value": "+20%"}, {"label": "고객", "value": "1,200"},
            {"label": "NPS", "value": "72"}, {"label": "이탈", "value": "3.2%"},
        ], region=Region(r.x, r.y, r.w, 0.6)))

    safe(prs, "comp_callout", "Atomic", lambda c, r:
        comp_callout(c, title="핵심 인사이트", body="AI 공격 47% 급증 — 기존 방어 체계 전환 필요",
                     region=Region(r.x, r.y, r.w, 1.2)))

    safe(prs, "comp_rag_row", "Atomic", lambda c, r:
        comp_rag_row(c, label="프로젝트 상태", values=["G", "A", "G", "R"],
                     region=Region(r.x, r.y, r.w, 0.5)))

    safe(prs, "comp_numbered_items", "Atomic", lambda c, r:
        comp_numbered_items(c, items=[
            {"title": "현황 진단", "body": "Scope 1·2·3 인벤토리"},
            {"title": "목표 설정", "body": "SBTi 감축 경로"},
            {"title": "실행 계획", "body": "EAF·CCUS 파일럿"},
        ], region=Region(r.x, r.y, r.w, 2.5)))

    safe(prs, "comp_section_header", "Atomic", lambda c, r:
        comp_section_header(c, title="02. 시장 분석", region=Region(r.x, r.y, r.w, 1.0)))

    safe(prs, "comp_progress_bar", "Atomic", lambda c, r:
        comp_progress_bar(c, label="진행률", value=0.65, target=0.65,
                          region=Region(r.x, r.y, r.w*0.6, 0.8)))

    safe(prs, "comp_vertical_bars", "Atomic", lambda c, r:
        comp_vertical_bars(c, data=[
            {"label": "Q1", "value": 45}, {"label": "Q2", "value": 62},
            {"label": "Q3", "value": 78, "highlight": True}, {"label": "Q4", "value": 55},
        ], unit="건", region=Region(r.x, r.y, r.w*0.5, 3.0)))

    safe(prs, "comp_heat_row", "Atomic", lambda c, r:
        comp_heat_row(c, label="리스크", values=[0.2, 0.5, 0.8, 0.3, 0.9],
                      col_labels=["기술", "시장", "규제", "인력", "보안"],
                      region=Region(r.x, r.y, r.w, 0.6)))

    safe(prs, "comp_gauge", "Atomic", lambda c, r:
        comp_gauge(c, value=0.73, label="달성률", target=0.73,
                   region=Region(r.x, r.y, 3.0, 2.5)))

    safe(prs, "comp_tag_group", "Atomic", lambda c, r:
        comp_tag_group(c, tags=["AI/ML", "Cloud", "DevOps", "Cybersecurity", "IoT"],
                       region=Region(r.x, r.y, r.w*0.6, 1.0)))

    safe(prs, "comp_comparison_row", "Atomic", lambda c, r:
        comp_comparison_row(c, label="배포", value_a="On-Prem", value_b="Cloud",
                            region=Region(r.x, r.y, r.w, 0.5)))

    safe(prs, "comp_metric_delta", "Atomic", lambda c, r:
        comp_metric_delta(c, label="매출", current=620, previous=500, unit="B",
                          region=Region(r.x, r.y, 3.0, 1.0)))

    safe(prs, "comp_timeline_mini", "Atomic", lambda c, r:
        comp_timeline_mini(c, phases=["Q1 킥오프", "Q2 파일럿", "Q3 확산", "Q4 최적화"], current=2,
                           region=Region(r.x, r.y, r.w, 0.8)))

    safe(prs, "comp_icon_list", "Atomic", lambda c, r:
        comp_icon_list(c, items=[
            {"icon": "▶", "title": "속도", "body": "3배 빠른 처리"},
            {"icon": "◆", "title": "정확도", "body": "오류율 0.1%"},
            {"icon": "●", "title": "비용", "body": "40% 절감"},
        ], region=Region(r.x, r.y, r.w*0.5, 2.5)))

    safe(prs, "comp_data_card", "Atomic", lambda c, r:
        comp_data_card(c, value=620, label="매출", previous=500, target=700, unit="B",
                       detail="영업이익 ₩68B", region=Region(r.x, r.y, 4.0, 2.5)))

    safe(prs, "comp_icon_card", "Atomic", lambda c, r:
        comp_icon_card(c, icon="CL", title="클라우드 전환",
                       body="온프레미스→하이브리드 전환으로 비용 40% 절감",
                       region=Region(r.x, r.y, 3.5, 2.0)))

    safe(prs, "comp_icon_row", "Atomic", lambda c, r:
        comp_icon_row(c, items=[
            {"icon": "▶", "label": "속도"}, {"icon": "◆", "label": "품질"},
            {"icon": "●", "label": "비용"}, {"icon": "■", "label": "리스크"},
        ], region=Region(r.x, r.y, r.w, 0.8)))

    safe(prs, "comp_styled_card (4가지 스타일)", "Atomic", lambda c, r: [
        comp_styled_card(c, title="Dark", kpi_value="$1B+", bullets=["시장 돌파", "CAGR 41.8%"], style="dark", region=Region(r.x, r.y, 2.2, 2.5)),
        comp_styled_card(c, title="Light", kpi_value="90%", bullets=["회수율 목표", "EU 규정"], style="light", region=Region(r.x+2.4, r.y, 2.2, 2.5)),
        comp_styled_card(c, title="Subtle", bullets=["연회색 배경", "보조 정보"], style="subtle", region=Region(r.x+4.8, r.y, 2.2, 2.5)),
        comp_styled_card(c, title="Numbered", number="01", bullets=["번호 강조", "순서 중요"], style="numbered", region=Region(r.x+7.2, r.y, 2.2, 2.5)),
    ])

    safe(prs, "comp_native_chart (세로 바)", "Atomic", lambda c, r:
        comp_native_chart(c, chart_type="vertical_bar", chart_kwargs={
            "categories": ["철강", "시멘트", "화학", "자동차"], "values": [3.3, 2.5, 1.0, 0.5],
            "highlight_idx": 0, "series_name": "배출량 (Gt CO₂)",
        }, region=Region(r.x, r.y, r.w*0.5, 3.5)))

    safe(prs, "comp_numbered_cell (2개)", "Atomic", lambda c, r: [
        comp_numbered_cell(c, number="01", header="진단", body="Scope 1·2·3\n기준연도\n갭 분석", region=Region(r.x, r.y, 3.0, 2.5)),
        comp_numbered_cell(c, number="02", header="목표", body="SBTi\n탄소가격\nKPI", bg_color="accent", region=Region(r.x+3.3, r.y, 3.0, 2.5)),
    ])

    safe(prs, "comp_timeline_marker", "Atomic", lambda c, r:
        comp_timeline_marker(c, labels=["2024 파일럿", "2025 확산", "2026 최적화", "2027 글로벌"],
                             highlight_idx=1, region=Region(r.x, r.y, r.w, 1.5)))

    safe(prs, "comp_icon_header_card", "Atomic", lambda c, r:
        comp_icon_header_card(c, icon="DT", header="디지털 트윈",
                              body="에너지 15~30% 절감\n예측 정비\nFoxconn 30%+",
                              region=Region(r.x, r.y, 4.0, 2.5)))

    # ── Compound ──
    safe(prs, "comp_chevron_flow (기본)", "Compound", lambda c, r:
        comp_chevron_flow(c, phases=[{"tag": "D", "label": "진단"}, {"tag": "P", "label": "계획"},
                                    {"tag": "E", "label": "실행"}, {"tag": "O", "label": "최적화"}],
                          style="gradient", region=Region(r.x, r.y, r.w, 0.55)))

    safe(prs, "comp_chevron_flow (show_details)", "Compound", lambda c, r:
        comp_chevron_flow(c, phases=[
            {"tag": "P1", "label": "진단", "details": ["현행 분석", "Gap 도출", "벤치마크"]},
            {"tag": "P2", "label": "설계", "details": ["솔루션 설계", "아키텍처"]},
            {"tag": "P3", "label": "구현", "details": ["개발", "테스트", "마이그레이션"]},
            {"tag": "P4", "label": "안정화", "details": ["Go-Live", "모니터링"]},
        ], show_details=True, style="gradient", region=Region(r.x, r.y, r.w, 3.5)))

    safe(prs, "comp_hero_block", "Compound", lambda c, r:
        comp_hero_block(c, label="ACTION REQUIRED", headline="2030년까지 41% 감축,\n연간 $4.5조 투자 필요",
                        sub_points=["전기화·순환경제 = Quick Win", "수소·CCUS 파일럿 임계점", "내부 탄소 가격제 = 기초"],
                        bg_color="grey_800", region=Region(r.x, r.y, r.w, 3.5)))

    safe(prs, "comp_hub_spoke_diagram", "Compound", lambda c, r:
        comp_hub_spoke_diagram(c, center="통합\n플랫폼", center_sub="Core",
                               spokes=[{"title": "ERP", "detail": "재무"}, {"title": "CRM", "detail": "고객"},
                                       {"title": "SCM", "detail": "공급망"}, {"title": "AI", "detail": "분석"}, {"title": "IoT", "detail": "설비"}],
                               region=Region(r.x, r.y, r.w*0.6, 4.5)))

    safe(prs, "comp_comparison_grid (**볼드**)", "Compound", lambda c, r:
        comp_comparison_grid(c, columns=[
            {"name": "On-Prem", "summary": "자체", "criteria": ["**자체 운영** — 보안 강점", "**높은 비용** — 서버 구매", "**느린 확장** — 조달 2~3월"]},
            {"name": "Hybrid", "summary": "혼합", "highlight": True, "criteria": ["**혼합** — 핵심 온프렘+부하 클라우드", "**중간 비용** — 점진 투자", "**유연 확장** — 버스트"]},
            {"name": "Cloud", "summary": "위탁", "criteria": ["**완전 위탁** — 운영 최소", "**낮은 비용** — 종량제", "**즉시 확장** — 분 단위"]},
        ], row_labels=["운영", "비용", "확장성"], region=Region(r.x, r.y, r.w, 3.5)))

    safe(prs, "comp_architecture_stack", "Compound", lambda c, r:
        comp_architecture_stack(c, layers=[
            {"name": "Presentation", "items": ["React", "Next.js"]}, {"name": "API", "items": ["Kong", "Auth"]},
            {"name": "Logic", "items": ["Node.js", "Python"]}, {"name": "Data", "items": ["PostgreSQL", "Redis"]},
            {"name": "Infra", "items": ["K8s", "AWS"]},
        ], region=Region(r.x, r.y, r.w*0.5, 4.5)))

    safe(prs, "comp_pyramid", "Compound", lambda c, r:
        comp_pyramid(c, layers=[
            {"title": "전기화", "detail": "CO₂ 70~80% 저감"}, {"title": "순환경제", "detail": "에너지 95% 절감"},
            {"title": "그린수소", "detail": "$2.50/kg"}, {"title": "CCUS", "detail": "50~90% 포집"},
            {"title": "디지털 트윈", "detail": "15~30% 절감"},
        ], region=Region(r.x, r.y, r.w*0.5, 4.5)))

    safe(prs, "comp_cycle_arrows", "Compound", lambda c, r:
        comp_cycle_arrows(c, center="PDCA", center_sub="Cycle",
                          stages=[{"label": "Plan"}, {"label": "Do"}, {"label": "Check"}, {"label": "Act"}],
                          region=Region(r.x, r.y, r.w*0.5, 4.5)))

    safe(prs, "comp_waterfall", "Compound", lambda c, r:
        comp_waterfall(c, start={"label": "현재\n(₩100B)", "value": 100},
                       steps=[{"label": "자동화", "value": -15}, {"label": "클라우드", "value": -10}, {"label": "인력", "value": -8}],
                       end={"label": "목표\n(₩67B)", "value": 67}, unit="B",
                       region=Region(r.x, r.y, r.w*0.7, 4.0)))

    safe(prs, "comp_before_after", "Compound", lambda c, r:
        comp_before_after(c, before_title="AS-IS", after_title="TO-BE",
                          before_items=[{"label": "수동", "detail": "엑셀 기반", "kpi": "5일"}, {"label": "사일로", "detail": "단절"}, {"label": "사후", "detail": "문제 후 조치"}],
                          after_items=[{"label": "자동화", "detail": "RPA+AI", "kpi": "2시간"}, {"label": "통합", "detail": "실시간"}, {"label": "예측", "detail": "AI 탐지"}],
                          arrow_label="전환", region=Region(r.x, r.y, r.w, 4.0)))

    safe(prs, "comp_gantt_bars", "Compound", lambda c, r:
        comp_gantt_bars(c, phases=["Q1", "Q2", "Q3", "Q4"],
                        streams=[
                            {"name": "인프라", "bars": [{"start": 0, "end": 2, "label": "셋업"}]},
                            {"name": "개발", "bars": [{"start": 1, "end": 3, "label": "MVP", "highlight": True}]},
                            {"name": "테스트", "bars": [{"start": 2, "end": 3.5, "label": "UAT"}]},
                            {"name": "배포", "bars": [{"start": 3, "end": 4, "label": "Go-Live"}]},
                        ], milestones=[{"phase": 3, "label": "Go-Live"}],
                        region=Region(r.x, r.y, r.w, 3.5)))

    safe(prs, "comp_value_chain", "Compound", lambda c, r:
        comp_value_chain(c, primary=[
            {"name": "Inbound", "items": ["조달"]}, {"name": "Ops", "items": ["제조"]},
            {"name": "Outbound", "items": ["물류"]}, {"name": "Mktg", "items": ["브랜딩"]},
            {"name": "Service", "items": ["A/S"]},
        ], support=[{"name": "인프라", "detail": "경영"}, {"name": "HR", "detail": "채용"}, {"name": "기술", "detail": "R&D"}],
        region=Region(r.x, r.y, r.w, 4.5)))

    safe(prs, "comp_logic_tree", "Compound", lambda c, r:
        comp_logic_tree(c, root="매출 증대", branches=[
            {"label": "신규 고객", "leaves": ["마케팅", "채널", "다변화"]},
            {"label": "기존 고객", "leaves": ["업셀링", "크로스셀링", "이탈 방지"]},
            {"label": "가격", "leaves": ["프리미엄화", "번들링"]},
        ], region=Region(r.x, r.y, r.w, 4.0)))

    safe(prs, "comp_quadrant_matrix", "Compound", lambda c, r:
        comp_quadrant_matrix(c, x_axis="실행 용이성", y_axis="영향도", quadrants=[
            {"label": "Quick Win", "items": ["자동화", "비용절감"]},
            {"label": "전략 투자", "items": ["AI", "클라우드"]},
            {"label": "필수 유지", "items": ["보안", "컴플라이언스"]},
            {"label": "재검토", "items": ["레거시"]},
        ], region=Region(r.x, r.y, r.w*0.6, 4.5)))

    safe(prs, "comp_funnel", "Compound", lambda c, r:
        comp_funnel(c, stages=[
            {"label": "TAM", "value": "₩500B", "detail": "전체 시장"},
            {"label": "SAM", "value": "₩200B", "detail": "접근 가능"},
            {"label": "SOM", "value": "₩50B", "detail": "목표"},
            {"label": "Revenue", "value": "₩15B", "detail": "현재"},
        ], region=Region(r.x, r.y, r.w*0.5, 4.0)))

    safe(prs, "comp_callout_annotation", "Compound", lambda c, r:
        comp_callout_annotation(c, annotations=[
            {"x": 0.5, "y": 0.5, "label": "핵심", "detail": "주목"},
            {"x": 3.5, "y": 1.5, "label": "출처", "detail": "IEA 2025"},
        ], region=Region(r.x, r.y, r.w*0.5, 3.0)))

    safe(prs, "comp_heatmap_grid", "Compound", lambda c, r:
        comp_heatmap_grid(c, row_labels=["기술", "시장", "규제", "인력"],
                          col_labels=["영향도", "확률", "대비"],
                          values=[[0.9,0.7,0.3],[0.5,0.6,0.5],[0.6,0.8,0.4],[0.4,0.7,0.3]],
                          cell_texts=[["Critical","High","Low"],["Medium","High","Medium"],["High","Critical","Low"],["Medium","High","Low"]],
                          region=Region(r.x, r.y, r.w*0.6, 3.5)))


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    build(prs)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path} ({len(prs.slides)}장)")

    import pythoncom, win32com.client
    pythoncom.CoInitialize()
    pdf_path = pptx_path.with_suffix(".pdf")
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    p = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
    p.SaveAs(str(pdf_path.resolve()), 32)
    p.Close()
    print(f"PDF:  {pdf_path}")

    from ppt_builder.track_c.png_export import pptx_to_pngs
    png_dir = OUTPUT_DIR / f"{NAME}_pngs"
    paths = pptx_to_pngs(pptx_path, png_dir)
    print(f"PNGs: {png_dir} ({len(paths)}장)")


if __name__ == "__main__":
    main()
