"""component_ops 통합 테스트 — 실제 마스터 템플릿 1,251장에서 그룹 추출/삽입.

시나리오:
  마스터 슬라이드 #12 (5-멤버 chevron-like 그룹, flat_idx 6~10)에서
  그룹 통째 추출 → 같은 마스터 기반 빈 슬라이드에 삽입 → 저장+재오픈 검증.

검증 포인트:
  - extract_group이 공통 부모(<p:grpSp>) 인식 → 합성 wrap 대신 통째 deepcopy
  - 한글 폰트 (<a:ea>) 보존
  - placeholder 텍스트 ('~~') 보존
  - a16:creationId 재생성 (원본과 다름)
  - insert 후 dst의 leaf 5개 늘어남
  - 저장+재오픈 후 PPT 손상 없음

실행: python tests/test_component_ops_master_integration.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from ppt_builder.template import TemplateEditor  # noqa: E402
from ppt_builder.template.component_ops import (  # noqa: E402
    create_blank_slide_with_master_theme,
    extract_group,
    insert_component,
)
from ppt_builder.template.edit_ops import iter_leaf_shapes  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
MASTER = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
OUT = ROOT / "output" / "component_ops_tests"
OUT.mkdir(parents=True, exist_ok=True)

CREATIONID_TAG = "{http://schemas.microsoft.com/office/drawing/2014/main}creationId"
SOURCE_SLIDE = 12  # 5-member chevron-like group at flat_idx 6~10
SOURCE_GROUP_INDICES = [6, 7, 8, 9, 10]


def has_korean(s: str) -> bool:
    return any("\uac00" <= c <= "\ud7a3" for c in s)


def main() -> int:
    if not MASTER.exists():
        print(f"[SKIP] master not found: {MASTER}")
        return 0

    print(f"Opening master: {MASTER}")
    src_prs = Presentation(str(MASTER))
    src_slide = src_prs.slides[SOURCE_SLIDE]
    src_leaves = list(iter_leaf_shapes(src_slide))
    print(f"  source slide #{SOURCE_SLIDE}: {len(src_leaves)} leaves")

    # 원본 그룹 멤버 텍스트 / 폰트 snapshot
    orig_texts = []
    orig_ea_fonts: set[str] = set()
    for fi, sh in src_leaves:
        if fi in SOURCE_GROUP_INDICES:
            if sh.has_text_frame:
                orig_texts.append(sh.text_frame.text)
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        ea = r._r.find(".//" + qn("a:rPr") + "/" + qn("a:ea"))
                        if ea is not None and ea.get("typeface"):
                            orig_ea_fonts.add(ea.get("typeface"))
    print(f"  selected member texts: {orig_texts}")
    print(f"  selected member EA fonts: {orig_ea_fonts}")

    # 원본 creationId snapshot (있으면)
    orig_ids: set[str] = set()
    for fi, sh in src_leaves:
        if fi in SOURCE_GROUP_INDICES:
            for el in sh._element.iter(CREATIONID_TAG):
                v = el.get("id")
                if v:
                    orig_ids.add(v)
    print(f"  selected member creationIds (originals): {len(orig_ids)}")

    # 1) 추출
    comp = extract_group(src_slide, SOURCE_GROUP_INDICES)
    assert comp.element.tag == qn("p:grpSp"), (
        f"expected <p:grpSp> (whole-group path), got {comp.element.tag}"
    )
    print(f"  extracted: tag={comp.element.tag.split('}')[-1]}, "
          f"bbox={comp.bbox_emu}, rids={len(comp.source_rids)}")

    # creationId 모두 새 GUID인지
    new_ids = {el.get("id") for el in comp.element.iter(CREATIONID_TAG)}
    if orig_ids:
        assert new_ids.isdisjoint(orig_ids), (
            f"creationId collision after extract: {new_ids & orig_ids}"
        )
        print(f"  creationIds regenerated: {len(new_ids)} unique, no collision")

    # 2) 같은 마스터에서 source 슬라이드 1장만 보존 + blank 추가
    #    (TemplateEditor.keep_slides는 slide-rels를 정리해 ZipFile duplicate 경고를 방지)
    editor = TemplateEditor(str(MASTER))
    editor.keep_slides([SOURCE_SLIDE])
    target_prs = editor.prs
    blank = create_blank_slide_with_master_theme(target_prs)
    pre_leaves = list(iter_leaf_shapes(blank))
    print(f"  target blank slide: {len(pre_leaves)} pre-existing leaves")

    # 3) 삽입
    new_idx = insert_component(blank, comp)
    post_leaves = list(iter_leaf_shapes(blank))
    print(f"  inserted: new_idx={new_idx}, total leaves now={len(post_leaves)}")
    assert len(post_leaves) - len(pre_leaves) == len(SOURCE_GROUP_INDICES), (
        f"expected {len(SOURCE_GROUP_INDICES)} new leaves, "
        f"got {len(post_leaves) - len(pre_leaves)}"
    )

    # 4) 텍스트 / EA 폰트 검증
    inserted_texts: list[str] = []
    inserted_ea: set[str] = set()
    for _, sh in post_leaves:
        if sh.has_text_frame:
            inserted_texts.append(sh.text_frame.text)
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    ea = r._r.find(".//" + qn("a:rPr") + "/" + qn("a:ea"))
                    if ea is not None and ea.get("typeface"):
                        inserted_ea.add(ea.get("typeface"))
    for ot in orig_texts:
        assert ot in inserted_texts, f"original text {ot!r} missing after insert"
    if orig_ea_fonts:
        assert orig_ea_fonts.issubset(inserted_ea), (
            f"EA font lost: orig={orig_ea_fonts}, after={inserted_ea}"
        )
        print(f"  EA fonts preserved: {orig_ea_fonts}")

    # 5) 두 번째 insert (creationId 중복 회피 검증)
    second_idx = insert_component(blank, comp)
    print(f"  second insert: new_idx={second_idx}")
    final_leaves = list(iter_leaf_shapes(blank))
    assert len(final_leaves) - len(post_leaves) == len(SOURCE_GROUP_INDICES)

    all_ids = []
    for _, sh in final_leaves:
        for el in sh._element.iter(CREATIONID_TAG):
            v = el.get("id")
            if v:
                all_ids.append(v)
    if all_ids:
        assert len(set(all_ids)) == len(all_ids), (
            f"duplicate creationIds across two inserts: "
            f"{len(all_ids) - len(set(all_ids))} collisions"
        )
        print(f"  total creationIds: {len(all_ids)}, all unique [OK]")

    # 6) 저장 + 재오픈
    out_path = OUT / "master_integration.pptx"
    editor.save(out_path)
    print(f"  saved: {out_path}")

    reopened = Presentation(str(out_path))
    # 슬라이드 2장: [0]=원본 source slide, [1]=blank+컴포넌트 2회 삽입
    assert len(reopened.slides) == 2, len(reopened.slides)
    rs = reopened.slides[1]
    re_leaves = list(iter_leaf_shapes(rs))
    print(f"  reopened slide#1: {len(re_leaves)} leaves")
    assert len(re_leaves) == 2 * len(SOURCE_GROUP_INDICES), len(re_leaves)

    # 재오픈 후 텍스트 다시 확인
    re_texts: list[str] = []
    for _, sh in re_leaves:
        if sh.has_text_frame:
            re_texts.append(sh.text_frame.text)
    for ot in orig_texts:
        assert re_texts.count(ot) >= 2, (
            f"text {ot!r} expected >=2 occurrences after 2x insert, got {re_texts.count(ot)}"
        )

    editor.cleanup()
    print()
    print("[OK] master integration test passed")
    return 0


if __name__ == "__main__":
    sys.exit(main())
