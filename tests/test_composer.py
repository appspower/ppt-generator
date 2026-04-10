"""SlideComposer 조합 테스트 — 5가지 복합 슬라이드.

기존 22개 패턴과 다른 구조의 슬라이드를 Composer로 조합하여 생성.
패턴 + 컴포넌트의 자유 조합을 검증.
"""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.composer import SlideComposer
from ppt_builder.primitives import Region
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_kpi_row, comp_bar_chart_h, comp_bullet_list,
    comp_callout, comp_stat_row, comp_section_header,
)
from ppt_builder.visual_validate import validate_visual

FOOTER = SlideFooter(source="출처: PwC Analysis 2024", right="PwC")


# ============================================================
# 1. top_bottom: 상단 KPI Row + 하단 바 차트
# ============================================================

def build_kpi_plus_chart():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="SAP 전환 4대 KPI 전원 Green — 테스트 자동화가 가장 큰 기여",
        category="성과 종합 — KPI + 상세 분석",
        nav_path=["3. 성과 보고", "1. 종합 대시보드"],
    ))
    comp.intro("4대 핵심 KPI 달성 현황(상단)과 모듈별 ROI 상세 분석(하단)을 단일 뷰로 통합")

    zones = comp.layout("top_bottom", split=0.3)

    # 상단: KPI 4개
    comp_kpi_row(comp.canvas, kpis=[
        {"value": "14%", "label": "일정 단축", "detail": "18→15.5개월", "trend": "up"},
        {"value": "70%", "label": "TC 공수↓", "detail": "수작업→AIP", "trend": "up"},
        {"value": "50%", "label": "DT 단축", "detail": "초과율 10%↓", "trend": "up"},
        {"value": "$4K", "label": "TC 비용", "detail": "20K건 2주", "trend": "down"},
    ], region=zones["top"])

    # 하단: 바 차트
    comp_bar_chart_h(comp.canvas, title="모듈별 ROI 기여도 (%)", data=[
        {"label": "테스트 자동화", "value": 85, "highlight": True},
        {"label": "Cutover 최적화", "value": 72, "highlight": True},
        {"label": "Health Dashboard", "value": 58},
        {"label": "Config Register", "value": 45},
        {"label": "Defect Triage", "value": 38},
        {"label": "주간 보고 자동화", "value": 32},
    ], unit="%", region=zones["bottom"])

    comp.takeaway("테스트 자동화(85%)와 Cutover(72%)가 ROI Top 2 — 두 모듈이 전체 가치의 70%+ 기여")
    comp.footer(FOOTER)
    return prs


# ============================================================
# 2. two_column: 좌 불릿 리스트 + 우 Callout 강조
# ============================================================

def build_analysis_plus_insight():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="SAP 전환 성공의 3대 조건 — 데이터 품질, 변화관리, 테스트 자동화를 동시에 확보",
        category="전략 분석 — Critical Success Factors",
        nav_path=["1. 전략", "2. 성공 조건"],
    ))
    comp.intro("30개 SAP 전환 프로젝트 분석 결과, 성공과 실패를 가르는 3대 핵심 조건을 도출")

    zones = comp.layout("two_column", split=0.5)

    # 좌: 현황 분석 불릿
    comp.canvas.push_region(zones["left"])
    comp.canvas.section_label("현황 분석", x=0, y=0, w=zones["left"].w, size=10)
    comp.canvas.pop_region()

    comp_bullet_list(comp.canvas, title="프로젝트 실패 패턴 (Top 5)",
                     items=[
                         "테스트 커버리지 부족 → Go-Live 후 결함 폭증 (85%)",
                         "데이터 마이그레이션 정합성 미검증 (72%)",
                         "Cutover 리허설 3회 미만 → DT 초과 (65%)",
                         "핵심 인력 이탈로 지식 단절 (58%)",
                         "요구사항 변경 통제 실패 → Scope Creep (52%)",
                     ],
                     region=zones["left"].sub(0, 0.4, zones["left"].w, zones["left"].h - 0.4))

    # 우: 3대 성공 조건 Callout
    comp.canvas.push_region(zones["right"])
    comp.canvas.section_label("3대 성공 조건", x=0, y=0, w=zones["right"].w, size=10)
    comp.canvas.pop_region()

    callout_h = (zones["right"].h - 0.5) / 3
    comp_callout(comp.canvas, title="1. 데이터 품질 선행 확보",
                 body="마이그레이션 전 Foundry Pipeline으로 정합성 사전 검증 자동화. 오류율 0.1% 이하 달성 후 전환 착수.",
                 bar_color="grey_900",
                 region=zones["right"].sub(0, 0.4, zones["right"].w, callout_h - 0.08))

    comp_callout(comp.canvas, title="2. 변화관리 조기 착수",
                 body="킥오프 시점부터 현업 Champion 지정. 교육 프로그램 3단계 운영. 저항 지표 월간 모니터링.",
                 bar_color="grey_700",
                 region=zones["right"].sub(0, 0.4 + callout_h, zones["right"].w, callout_h - 0.08))

    comp_callout(comp.canvas, title="3. 테스트 자동화 필수",
                 body="AIP 기반 TC 자동 생성으로 커버리지 95%+ 확보. LLM-as-Judge 연속 검증으로 정확도 99.8%.",
                 bar_color="grey_400",
                 region=zones["right"].sub(0, 0.4 + callout_h * 2, zones["right"].w, callout_h - 0.08))

    comp.takeaway("3대 조건 중 하나라도 미충족 시 프로젝트 실패 확률 3배↑ — Palantir가 1번(데이터)과 3번(테스트)을 직접 해소")
    comp.footer(FOOTER)
    return prs


# ============================================================
# 3. sidebar_left: 좌 통계 + 우 불릿 상세
# ============================================================

def build_stats_plus_detail():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="Palantir 도입 효과 4대 지표 — 정량 성과와 정성 인사이트를 통합 보고",
        category="성과 보고 — Quantitative + Qualitative",
        nav_path=["3. 성과", "2. 통합 뷰"],
    ))

    zones = comp.layout("sidebar_left", sidebar_w=2.8)

    # 좌측: 통계
    comp.canvas.push_region(zones["sidebar"])
    comp.canvas.section_label("핵심 지표", x=0, y=0, w=zones["sidebar"].w)
    stats = [
        {"value": "14%", "label": "일정 단축"},
        {"value": "70%", "label": "공수 절감"},
        {"value": "99.8%", "label": "TC 정확도"},
        {"value": "50%", "label": "DT 단축"},
    ]
    for i, st in enumerate(stats):
        comp.canvas.stat_block(
            value=st["value"], label=st["label"],
            x=0.1, y=0.45 + i * 1.1, w=zones["sidebar"].w - 0.2, h=0.9,
        )
    comp.canvas.pop_region()

    # 우측: 상세 해설
    comp.canvas.push_region(zones["main"])
    comp.canvas.section_label("상세 인사이트", x=0, y=0, w=zones["main"].w)
    comp.canvas.pop_region()

    insights = [
        ("일정 14% 단축의 의미", "18개월→15.5개월. 테스트 자동화(1.2개월) + Cutover 최적화(0.8개월) + 거버넌스(0.3개월) + 리스크 감지(0.2개월) 합산. 전체 프로젝트 비용 약 $2M 절감 효과."),
        ("테스트 공수 70% 절감", "AIP가 Blueprint→규칙→TC 자동 생성. 기존 컨설턴트 3~6개월 수작업을 1~2주로 압축. 건당 비용 $50→$0.2로 250배 효율화."),
        ("TC 정확도 99.8%", "LLM-as-Judge 0~10점 연속 검증. 기존 육안 리뷰(85%) 대비 14.8%p 향상. 결함 누락률 15%→0.2%로 75배 개선."),
        ("Cutover DT 50% 단축", "Workshop App 기반 3회 리허설로 12시간→6시간. 초과율 25%→10%로 감소. Go-Live 리스크 정량화 달성."),
    ]
    for i, (title, body) in enumerate(insights):
        h = (zones["main"].h - 0.4) / len(insights)
        comp_callout(comp.canvas, title=title, body=body,
                     bar_color="grey_900" if i == 0 else "grey_400",
                     region=zones["main"].sub(0, 0.4 + i * h, zones["main"].w, h - 0.08))

    comp.takeaway("4대 지표 모두 목표 초과 달성 — Quick Win(L1) 성공이 L2~L3 확대의 자연스러운 동력")
    comp.footer(FOOTER)
    return prs


# ============================================================
# 4. three_column: 3열 Callout 비교
# ============================================================

def build_three_options():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="3가지 테스트 자동화 접근법 중 AIP 기반 전수 자동화가 ROI 최고 — Option B 권장",
        category="옵션 분석 — Test Automation Strategy",
        nav_path=["2. 전략", "3. 옵션 비교"],
    ))
    comp.intro("수작업 유지(A), AIP 전수 자동화(B), 하이브리드(C) 3가지 접근법의 비용·효과·리스크 비교")

    zones = comp.layout("three_column")

    options = [
        {
            "title": "Option A\n수작업 유지",
            "bullets": ["기존 프로세스 유지", "컨설턴트 3~6개월", "커버리지 40~60%", "건당 $50~100", "리스크: 일정 지연 85%"],
            "bar": "grey_400",
        },
        {
            "title": "Option B (권장)\nAIP 전수 자동화",
            "bullets": ["Blueprint→TC 자동생성", "1~2주 완료", "커버리지 95%+", "건당 $0.2", "리스크: 초기 설정 2주"],
            "bar": "grey_900",
        },
        {
            "title": "Option C\n하이브리드",
            "bullets": ["핵심 TC는 AIP 자동", "엣지케이스 수동 보완", "커버리지 80~90%", "건당 $15~25", "리스크: 경계 불명확"],
            "bar": "grey_700",
        },
    ]

    for i, (key, zone) in enumerate(zones.items()):
        opt = options[i]
        # 헤더 박스
        comp.canvas.box(x=0, y=0, w=zone.w, h=0.8,
                        fill="grey_900" if i == 1 else "grey_200",
                        border=None, region=zone)
        comp.canvas.text(opt["title"], x=0.15, y=0, w=zone.w - 0.3, h=0.8,
                         size=12 if i == 1 else 11, bold=True,
                         color="white" if i == 1 else "grey_900",
                         anchor="middle", region=zone)
        # 불릿
        comp_bullet_list(comp.canvas, items=opt["bullets"],
                         region=zone.sub(0, 0.9, zone.w, zone.h - 1.0),
                         item_size=9)

    comp.takeaway("Option B(AIP 전수 자동화)가 ROI 350%, 커버리지 95%, 건당 $0.2로 3가지 기준 모두 최우위")
    comp.footer(FOOTER)
    return prs


# ============================================================
# 5. grid_2x2: 4사분면 컴포넌트 혼합
# ============================================================

def build_mixed_grid():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="SAP 전환 현황 4분면 대시보드 — KPI, 리스크, 일정, 품질을 단일 뷰로 통합",
        category="PMO 대시보드 — 4-Quadrant View",
        nav_path=["3. 거버넌스", "1. 종합 현황"],
    ))

    zones = comp.layout("grid_2x2")

    # TL: KPI 2개
    comp.canvas.push_region(zones["tl"])
    comp.canvas.section_label("핵심 KPI", x=0, y=0, w=zones["tl"].w)
    comp.canvas.pop_region()
    comp_kpi_row(comp.canvas, kpis=[
        {"value": "92%", "label": "마일스톤 달성률", "trend": "up"},
        {"value": "97%", "label": "결함 해소율", "trend": "up"},
    ], region=zones["tl"].sub(0, 0.35, zones["tl"].w, zones["tl"].h - 0.4))

    # TR: 리스크 불릿
    comp.canvas.push_region(zones["tr"])
    comp.canvas.section_label("Top 리스크", x=0, y=0, w=zones["tr"].w)
    comp.canvas.pop_region()
    comp_bullet_list(comp.canvas, items=[
        "데이터 마이그레이션 정합성 (R)",
        "인터페이스 호환성 테스트 지연 (A)",
        "핵심 인력 2명 이탈 예고 (A)",
        "Cutover 리허설 1회차 미완 (A)",
    ], region=zones["tr"].sub(0, 0.35, zones["tr"].w, zones["tr"].h - 0.4))

    # BL: 바 차트
    comp_bar_chart_h(comp.canvas, title="워크스트림 진척률 (%)", data=[
        {"label": "FI/CO", "value": 95, "highlight": True},
        {"label": "MM", "value": 82},
        {"label": "SD", "value": 78},
        {"label": "PP", "value": 65},
        {"label": "Data", "value": 55, "highlight": True},
    ], unit="%", region=zones["bl"])

    # BR: 일정 요약
    comp.canvas.push_region(zones["br"])
    comp.canvas.section_label("일정 현황", x=0, y=0, w=zones["br"].w)
    comp.canvas.pop_region()
    comp_bullet_list(comp.canvas, items=[
        "Blueprint: 완료 (2024 Q2)",
        "Build: 진행중 82% (2024 Q3)",
        "테스트: 착수 예정 (2024 Q4)",
        "Go-Live: 2025 Q2 예정",
    ], region=zones["br"].sub(0, 0.35, zones["br"].w, zones["br"].h - 0.4))

    comp.takeaway("전체 진척률 82%, KPI Green — 단, 데이터 마이그레이션(R)과 인력 리스크(A) 즉시 대응 필요")
    comp.footer(FOOTER)
    return prs


# ============================================================
# 메인
# ============================================================

def main():
    out_dir = Path("output/composer_test")
    out_dir.mkdir(parents=True, exist_ok=True)

    cases = [
        ("01_kpi_plus_chart", build_kpi_plus_chart),
        ("02_analysis_insight", build_analysis_plus_insight),
        ("03_stats_detail", build_stats_plus_detail),
        ("04_three_options", build_three_options),
        ("05_mixed_grid", build_mixed_grid),
    ]

    print("=" * 70)
    print("SlideComposer Test — 5 Composite Slides")
    print("=" * 70)

    all_passed = True
    for name, builder_fn in cases:
        out = out_dir / f"{name}.pptx"
        try:
            prs = builder_fn()
            prs.save(out)
        except Exception as e:
            print(f"\n[{name}] BUILD FAILED: {e}")
            import traceback
            traceback.print_exc()
            all_passed = False
            continue

        visual = validate_visual(out, convert_pdf=False)
        ok = not visual.issues
        status = "PASS" if ok else "FAIL"
        if not ok:
            all_passed = False

        print(f"\n[{name}] {status}")
        print(f"  visual: {len(visual.issues)} issues")
        for i in visual.issues[:3]:
            print(f"    - {i}")

    # PDF
    print("\n" + "=" * 70)
    print("PDF...")
    for name, _ in cases:
        out = out_dir / f"{name}.pptx"
        if not out.exists():
            continue
        try:
            visual = validate_visual(out, convert_pdf=True)
            print(f"  {name}.pdf  -- {'OK' if visual.pdf_available else 'SKIP'}")
        except Exception as e:
            print(f"  {name}.pdf  -- FAIL ({e})")

    print("\n" + "=" * 70)
    print(f"Result: {'ALL PASSED' if all_passed else 'SOME FAILED'}")
    print("=" * 70)
    return 0 if all_passed else 1


if __name__ == "__main__":
    sys.exit(main())
