"""Compound Component 테스트 — 5개 컴포넌트 단독 + 조합 검증.

  슬라이드 1: comp_chevron_flow (단독, gradient)
  슬라이드 2: comp_hero_block (단독)
  슬라이드 3: comp_hub_spoke_diagram (단독, 5 spoke)
  슬라이드 4: comp_comparison_grid (단독, 3열)
  슬라이드 5: comp_architecture_stack (단독, 5층)
  슬라이드 6: 조합 — t_layout: chevron(상단) + hero(하단좌) + kpi(하단우)
  슬라이드 7: 조합 — center_peripheral_4: hub_spoke(중앙) + bullets(4방향)
  슬라이드 8: 조합 — top_bottom: comparison_grid(상단) + architecture_stack(하단)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.composer import SlideComposer
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_chevron_flow,
    comp_hero_block,
    comp_hub_spoke_diagram,
    comp_comparison_grid,
    comp_architecture_stack,
    comp_kpi_row,
    comp_bullet_list,
)

OUTPUT = Path(__file__).parent.parent / "output" / "compound_components_test.pptx"


def make(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def s1_chevron(prs):
    slide = make(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(title="디지털 전환 4단계 프로세스", category="Process"))
    zones = comp.layout("full")
    comp_chevron_flow(comp.canvas, phases=[
        {"tag": "01", "label": "Discover", "details": ["현행 분석", "Gap 도출"]},
        {"tag": "02", "label": "Design", "details": ["솔루션 설계", "아키텍처"]},
        {"tag": "03", "label": "Deliver", "details": ["구현", "테스트"]},
        {"tag": "04", "label": "Deploy", "details": ["Go-Live", "안정화"]},
    ], show_details=True, region=zones["main"])
    comp.footer(SlideFooter())


def s2_hero(prs):
    slide = make(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(title="핵심 전략 메시지", category="Strategy"))
    zones = comp.layout("full")
    comp_hero_block(comp.canvas,
        label="TRANSFORMATION",
        headline="클라우드 네이티브 전환으로\n운영 비용 40% 절감 실현",
        sub_points=[
            "마이크로서비스 아키텍처 기반 시스템 재설계",
            "컨테이너 오케스트레이션(K8s) 도입",
            "DevOps 파이프라인 자동화 (CI/CD)",
            "모니터링 & 옵저버빌리티 체계 구축",
        ],
        region=zones["main"])
    comp.footer(SlideFooter())


def s3_hub_spoke(prs):
    slide = make(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(title="통합 플랫폼 생태계 구조", category="Architecture"))
    zones = comp.layout("full")
    comp_hub_spoke_diagram(comp.canvas,
        center="통합\n플랫폼",
        center_sub="Core System",
        spokes=[
            {"title": "ERP", "detail": "SAP S/4HANA\n재무, 구매, 생산"},
            {"title": "CRM", "detail": "Salesforce\n고객 관리"},
            {"title": "SCM", "detail": "공급망 최적화\n실시간 추적"},
            {"title": "Analytics", "detail": "BI + AI/ML\n예측 분석"},
            {"title": "IoT", "detail": "설비 모니터링\n실시간 데이터"},
        ],
        region=zones["main"])
    comp.footer(SlideFooter())


def s4_comparison(prs):
    slide = make(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(title="도입 방식 3안 비교 분석", category="Comparison"))
    zones = comp.layout("full")
    comp_comparison_grid(comp.canvas,
        columns=[
            {"name": "Big Bang", "summary": "일괄 전환", "criteria": [
                "12개월", "₩150억", "높음", "빠른 전환"]},
            {"name": "Phased", "summary": "단계 전환", "highlight": True, "criteria": [
                "18개월", "₩120억", "중간", "안정적"]},
            {"name": "Parallel", "summary": "병행 운영", "criteria": [
                "24개월", "₩180억", "낮음", "안전"]},
        ],
        row_labels=["기간", "비용", "리스크", "장점"],
        region=zones["main"])
    comp.footer(SlideFooter())


def s5_stack(prs):
    slide = make(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(title="시스템 아키텍처 레이어", category="Technology"))
    zones = comp.layout("full")
    comp_architecture_stack(comp.canvas,
        layers=[
            {"title": "Presentation", "items": ["React", "Next.js", "Tailwind"]},
            {"title": "API Gateway", "items": ["Kong", "Rate Limit", "Auth"]},
            {"title": "Business Logic", "items": ["Spring Boot", "gRPC", "Kafka"]},
            {"title": "Data Layer", "items": ["PostgreSQL", "Redis", "S3"]},
            {"title": "Infrastructure", "items": ["Kubernetes", "Terraform", "ArgoCD"]},
        ],
        style="gradient",
        region=zones["main"])
    comp.footer(SlideFooter())


def s6_combo_t(prs):
    """조합: t_layout — chevron(상단) + hero(하단좌) + kpi(하단우)."""
    slide = make(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(title="전환 로드맵과 기대 효과", category="Executive"))
    zones = comp.layout("t_layout", top_ratio=0.22, right_ratio=0.40)

    comp_chevron_flow(comp.canvas, phases=[
        {"tag": "P1", "label": "진단"},
        {"tag": "P2", "label": "설계"},
        {"tag": "P3", "label": "구현"},
        {"tag": "P4", "label": "안정화"},
    ], style="gradient", region=zones["top"])

    comp_hero_block(comp.canvas,
        label="WHY",
        headline="레거시 시스템이\n성장의 병목",
        sub_points=[
            "연간 유지보수비 ₩45억",
            "신기능 배포 주기 6개월",
            "시스템 장애 월 평균 3회",
        ],
        region=zones["bottom_left"])

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "40%", "label": "비용 절감", "trend": "up"},
        {"value": "3x", "label": "배포 속도", "trend": "up"},
    ], region=zones["bottom_right"])

    comp.takeaway("P1 진단 단계에서 레거시 병목을 정량화하여 경영진 의사결정 지원")
    comp.footer(SlideFooter())


def s7_combo_hub(prs):
    """조합: center_peripheral_4 — hub_spoke(중앙) + bullets(4방향)."""
    slide = make(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(title="디지털 CoE 거버넌스 구조", category="Governance"))
    zones = comp.layout("center_peripheral_4", center_ratio=0.42)

    comp_hub_spoke_diagram(comp.canvas,
        center="Digital\nCoE",
        spokes=[
            {"title": "전략", "detail": "DX 로드맵"},
            {"title": "기술", "detail": "아키텍처"},
            {"title": "인재", "detail": "역량 개발"},
            {"title": "프로세스", "detail": "자동화"},
        ],
        region=zones["center"])

    sides = {
        "left": ("전략 기획", ["DX 비전 수립", "투자 우선순위", "KPI 체계"]),
        "right": ("기술 관리", ["표준 아키텍처", "기술 스택 관리", "보안 정책"]),
        "top": ("인재 육성", ["디지털 교육", "CoP 운영"]),
        "bottom": ("프로세스", ["자동화 파이프라인", "변경 관리"]),
    }
    for pos, (title, items) in sides.items():
        comp_bullet_list(comp.canvas, title=title, items=items,
                         region=zones[pos])

    comp.footer(SlideFooter())


def s8_combo_compare_stack(prs):
    """조합: top_bottom — comparison(상단) + architecture_stack(하단)."""
    slide = make(prs)
    comp = SlideComposer(slide)
    comp.header(SlideHeader(title="클라우드 전환: 옵션 비교 + 타겟 아키텍처", category="Analysis"))
    zones = comp.layout("top_bottom", split=0.50)

    comp_comparison_grid(comp.canvas,
        columns=[
            {"name": "On-Prem", "criteria": ["자체 운영", "높음", "느림"]},
            {"name": "Hybrid", "highlight": True, "criteria": ["혼합", "중간", "유연"]},
            {"name": "Full Cloud", "criteria": ["완전 위탁", "낮음", "빠름"]},
        ],
        row_labels=["운영 방식", "초기 비용", "확장성"],
        region=zones["top"])

    comp_architecture_stack(comp.canvas,
        layers=[
            {"title": "Frontend", "items": ["Web", "Mobile", "API"]},
            {"title": "Backend", "items": ["Microservices", "Serverless"]},
            {"title": "Cloud", "items": ["AWS", "Azure", "GCP"]},
        ],
        style="gradient",
        region=zones["bottom"])

    comp.takeaway("Hybrid 방식이 비용-확장성 균형 최적, 3계층 클라우드 아키텍처 적용")
    comp.footer(SlideFooter())


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s1_chevron(prs)
    s2_hero(prs)
    s3_hub_spoke(prs)
    s4_comparison(prs)
    s5_stack(prs)
    s6_combo_t(prs)
    s7_combo_hub(prs)
    s8_combo_compare_stack(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
