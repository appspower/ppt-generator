"""밀도 극대화 테스트 — 사이버보안 1장.

바이블 워크플로우 + 컴포넌트 밀도 개선 반영.
핵심: 각 셀/불릿에 2~3줄 텍스트를 넣어 공간을 채운다.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.composer import SlideComposer
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_kpi_row, comp_comparison_grid, comp_bullet_list,
)

OUTPUT = Path(__file__).parent.parent / "output" / "density_test.pptx"


def make(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_cybersecurity(prs):
    """사이버보안: AI 위협 급증 + 기업 대응. 밀도 극대화."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="AI 기반 사이버 공격이 47% 급증하며, 400만 인력 부족 속 기업 보안 투자가 $520B로 확대되고 있다",
        category="Cybersecurity"))
    zones = comp.layout("t_layout", top_ratio=0.12, right_ratio=0.38)

    # 상단: KPI 4개 — 콤팩트하게 핵심 수치
    comp_kpi_row(comp.canvas, kpis=[
        {"value": "$10.5T", "label": "사이버범죄 연간 비용 (2025)", "trend": "up"},
        {"value": "+47%", "label": "AI 기반 공격 증가율 (YoY)", "trend": "up"},
        {"value": "4M", "label": "글로벌 보안 인력 부족", "trend": "up"},
        {"value": "$520B", "label": "기업 보안 지출 (2026E)", "trend": "up"},
    ], region=zones["top"])

    # 하좌: 비교 그리드 — 3대 위협 × 3차원 분석
    # ★ 핵심: 각 criteria를 40~60자(2~3줄)로 작성하여 셀을 채운다
    comp_comparison_grid(comp.canvas,
                         columns=[
                             {"name": "AI 피싱/딥페이크", "summary": "공격 자동화",
                              "criteria": [
                                  "GenAI로 피싱 메일 1,265% 폭증 — 기존 필터 우회율 78%. CEO 음성 딥페이크로 $25M 송금 사기 발생 (홍콩 다국적기업)",
                                  "공격 준비 시간 45일→5일로 단축. 자동화 스크립트가 취약점 스캔~익스플로잇까지 무인 수행. 중소기업 특히 취약",
                                  "AI 피싱 탐지 솔루션 도입 시급 — 기존 rule-based 방어 무력화. Behavioral AI 기반 이메일 분석으로 탐지율 92%+ 가능",
                              ]},
                             {"name": "랜섬웨어", "summary": "피해 급증", "highlight": True,
                              "criteria": [
                                  "전체 공격의 35% 차지, YoY +84% 증가. 건당 평균 피해 $1.18M — SME의 경우 51%가 랜섬웨어로 인한 비용. 복구 평균 23일 소요",
                                  "이중 갈취(Double Extortion) 전략 보편화 — 데이터 암호화 + 유출 협박 병행. 지불 후에도 40%는 데이터 미복구. RaaS 모델로 진입장벽 하락",
                                  "백업 격리(Air-gapped) + EDR + 제로트러스트 3중 방어 필수. 사고 대응 플레이북 분기별 모의훈련 — MTTR 72hr→24hr 목표",
                              ]},
                             {"name": "공급망 공격", "summary": "3자 리스크",
                              "criteria": [
                                  "사이버 복원력 강화의 2위 과제(46%, WEF). SolarWinds·MOVEit 사태 이후 SBOM(소프트웨어 부품 명세) 의무화 확산",
                                  "평균 기업의 3자 벤더 수 1,200개 이상. 이 중 보안 평가 완료 비율 34%에 불과. 공급망 1개 침해 → 평균 11개 기업 연쇄 피해",
                                  "3자 보안 평가 자동화(TPRM) + 계약서 SLA에 보안 요건 명시. 연 1회 감사 → 실시간 모니터링 전환 필요",
                              ]},
                         ],
                         row_labels=["현황·규모", "공격 기법·진화", "대응 전략"],
                         region=zones["bottom_left"])

    # 하우: CISO 6대 우선순위 — 각 항목 2줄 이상
    comp_bullet_list(comp.canvas, title="CISO 6대 우선순위 (Gartner 2025)",
                     items=[
                         "사이버 복원력(Cyber Resilience) — 공격을 막는 것보다 '공격 후 복구 속도'가 핵심. BCP/DR 자동화로 MTTR 단축 (Gartner 1순위)",
                         "제로트러스트 아키텍처 — 네트워크 경계 보안에서 'Never Trust, Always Verify'로 전환. Mayo Clinic·JPMorgan 등 AI 기반 ZTA 도입 완료",
                         "AI 보안 자동화(SOAR/XDR) — 인력 400만 부족을 AI로 보완. 탐지~대응 자동화로 분석관 1인당 처리량 5배 향상",
                         "클라우드 보안 거버넌스 — 멀티클라우드 환경의 설정 오류가 침해 원인 1위. CSPM·CNAPP 도입으로 가시성 확보",
                         "공급망 3자 리스크 관리 — SBOM 의무화 대비 + 실시간 벤더 보안 모니터링. 계약 단계부터 보안 SLA 명문화",
                         "보안 인력 리스킬링 — 기존 IT 인력의 보안 역량 전환 교육. 59% 이직 고려 중(ISC²) — 번아웃 방지와 경력 경로 설계 병행",
                     ], region=zones["bottom_right"])

    comp.takeaway("AI 공격 47% 급증 + 인력 400만 부족 = '사이버 복원력' 중심의 자동화·제로트러스트 전환이 생존 조건")
    comp.footer(SlideFooter(source="Cybersecurity Ventures 2026, WEF Global Cyber Outlook 2026, ISC² Workforce 2025, Gartner CISO Survey"))


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cybersecurity(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")

    from ppt_builder.evaluate import evaluate_pptx, print_report
    report = evaluate_pptx(str(OUTPUT))
    print_report(report)


if __name__ == "__main__":
    main()
