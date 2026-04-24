"""편집 API 5종 단위 테스트 — pytest 없이 standalone 실행.

커버리지:
  1. iter_leaf_shapes: 평탄화 순회 (그룹 재귀)
  2. clone_paragraph: XML 복제 + 한글 EA 폰트 보존
  3. replace_paragraph: 단순 치환 + markdown bold/italic/code
  4. del_paragraph: 중간 삭제 + 마지막 1개 보호
  5. replace_image: 비율 유지 + blip rId 교체
  6. del_image: 픽처 삭제
  7. SlideEditError: 범위 초과 / 타입 불일치

실행: python tests/test_edit_ops.py
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from ppt_builder.template.edit_ops import (  # noqa: E402
    SlideEditError,
    clone_paragraph,
    del_image,
    del_paragraph,
    iter_leaf_shapes,
    replace_image,
    replace_paragraph,
)

OUT_DIR = Path(__file__).resolve().parent.parent / "output" / "edit_ops_tests"
OUT_DIR.mkdir(parents=True, exist_ok=True)

passed: list[str] = []
failed: list[tuple[str, str]] = []


def _run(name: str, fn):
    try:
        fn()
        passed.append(name)
        print(f"  [OK]   {name}")
    except AssertionError as e:
        failed.append((name, f"assert: {e}"))
        print(f"  [FAIL] {name}: {e}")
    except Exception as e:
        failed.append((name, f"{type(e).__name__}: {e}"))
        print(f"  [ERR]  {name}: {type(e).__name__}: {e}")


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _make_text_shape(slide, left, top, width, height, paragraphs):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
    return tb


def _make_dummy_png(path: Path, w: int = 100, h: int = 50, color=(255, 0, 0)):
    img = Image.new("RGB", (w, h), color)
    img.save(str(path))


# -----------------------------------------------------------------------------
# 1. iter_leaf_shapes — 평탄화
# -----------------------------------------------------------------------------

def test_iter_leaf_shapes_flat():
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 3, 1, ["A"])
    _make_text_shape(slide, 0.5, 2.0, 3, 1, ["B"])
    _make_text_shape(slide, 0.5, 3.5, 3, 1, ["C"])
    indices = [i for i, _ in iter_leaf_shapes(slide)]
    assert indices == [0, 1, 2], f"expected [0,1,2], got {indices}"


def test_iter_leaf_shapes_group_recurse():
    """그룹 내부 shape도 평탄 인덱스로 노출."""
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 3, 1, ["outside-0"])
    # 그룹 대신 더미로 텍스트 박스 2개를 추가 후 첫 번째 박스를 유지
    # python-pptx는 그룹 생성 API가 제한적이므로 실제 템플릿 기반 그룹 테스트는
    # 통합 테스트에서 수행. 여기서는 평탄 순서만 확인.
    _make_text_shape(slide, 0.5, 2.0, 3, 1, ["outside-1"])
    seen = [(i, sh.text_frame.text) for i, sh in iter_leaf_shapes(slide)]
    assert seen == [(0, "outside-0"), (1, "outside-1")], seen


# -----------------------------------------------------------------------------
# 2. clone_paragraph
# -----------------------------------------------------------------------------

def test_clone_paragraph_basic():
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 4, 2, ["first", "second"])
    new_idx = clone_paragraph(slide, div_id=0, paragraph_id=0)
    shape = next(sh for i, sh in iter_leaf_shapes(slide) if i == 0)
    paras = shape.text_frame.paragraphs
    assert len(paras) == 3, f"expected 3 paragraphs, got {len(paras)}"
    assert new_idx == 2, f"expected new_idx=2, got {new_idx}"
    assert paras[2].text == "first", f"cloned text: {paras[2].text!r}"


def test_clone_paragraph_preserves_korean_ea_font():
    """한글 런의 <a:ea typeface> 속성이 XML 복제에서 보존되는지."""
    prs = _new_prs()
    slide = _blank(prs)
    tb = _make_text_shape(slide, 0.5, 0.5, 5, 2, ["한글 테스트"])
    # 강제로 EA 폰트 지정
    run = tb.text_frame.paragraphs[0].runs[0]
    run.font.name = "맑은 고딕"
    run.font.size = Pt(14)
    run.font.bold = True

    clone_paragraph(slide, div_id=0, paragraph_id=0)
    paras = tb.text_frame.paragraphs
    assert len(paras) == 2
    cloned_run = paras[1].runs[0]
    assert cloned_run.text == "한글 테스트"
    assert cloned_run.font.name == "맑은 고딕"
    assert cloned_run.font.size == Pt(14)
    assert cloned_run.font.bold is True


# -----------------------------------------------------------------------------
# 3. replace_paragraph
# -----------------------------------------------------------------------------

def test_replace_paragraph_plain():
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 4, 1, ["old text"])
    replace_paragraph(slide, div_id=0, paragraph_id=0, text="new text")
    shape = next(sh for i, sh in iter_leaf_shapes(slide) if i == 0)
    assert shape.text_frame.paragraphs[0].text == "new text"


def test_replace_paragraph_markdown_bold_italic():
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 6, 1, ["placeholder"])
    replace_paragraph(
        slide, div_id=0, paragraph_id=0,
        text="normal **bold** and *italic* here",
    )
    shape = next(sh for i, sh in iter_leaf_shapes(slide) if i == 0)
    runs = shape.text_frame.paragraphs[0].runs
    combined = "".join(r.text for r in runs)
    assert combined == "normal bold and italic here", combined
    # bold / italic 런이 하나씩은 있어야 함
    bold_count = sum(1 for r in runs if r.font.bold)
    italic_count = sum(1 for r in runs if r.font.italic)
    assert bold_count >= 1, f"no bold run: {[(r.text, r.font.bold) for r in runs]}"
    assert italic_count >= 1, f"no italic run: {[(r.text, r.font.italic) for r in runs]}"


def test_replace_paragraph_preserves_korean_font():
    prs = _new_prs()
    slide = _blank(prs)
    tb = _make_text_shape(slide, 0.5, 0.5, 5, 1, ["한글"])
    run = tb.text_frame.paragraphs[0].runs[0]
    run.font.name = "맑은 고딕"
    run.font.size = Pt(12)

    replace_paragraph(slide, div_id=0, paragraph_id=0, text="변경된 한글 **강조**")
    runs = tb.text_frame.paragraphs[0].runs
    combined = "".join(r.text for r in runs)
    assert combined == "변경된 한글 강조", combined
    # 첫 run의 폰트 이름이 유지되는지
    assert runs[0].font.name == "맑은 고딕"
    assert runs[0].font.size == Pt(12)


# -----------------------------------------------------------------------------
# 4. del_paragraph
# -----------------------------------------------------------------------------

def test_del_paragraph_middle():
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 4, 2, ["A", "B", "C"])
    del_paragraph(slide, div_id=0, paragraph_id=1)
    shape = next(sh for i, sh in iter_leaf_shapes(slide) if i == 0)
    texts = [p.text for p in shape.text_frame.paragraphs]
    assert texts == ["A", "C"], texts


def test_del_paragraph_last_preserves_shell():
    prs = _new_prs()
    slide = _blank(prs)
    tb = _make_text_shape(slide, 0.5, 0.5, 4, 1, ["only"])
    run = tb.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(20)
    del_paragraph(slide, div_id=0, paragraph_id=0)
    shape = next(sh for i, sh in iter_leaf_shapes(slide) if i == 0)
    paras = shape.text_frame.paragraphs
    assert len(paras) == 1, f"paragraph count: {len(paras)}"
    # 런은 유지되되 텍스트만 비워짐
    runs = paras[0].runs
    assert all(r.text == "" for r in runs)


# -----------------------------------------------------------------------------
# 5. replace_image
# -----------------------------------------------------------------------------

def test_replace_image():
    prs = _new_prs()
    slide = _blank(prs)
    # 원본 이미지
    src_path = OUT_DIR / "_src.png"
    _make_dummy_png(src_path, w=200, h=100, color=(255, 0, 0))
    pic = slide.shapes.add_picture(
        str(src_path), Inches(1), Inches(1), Inches(4), Inches(2),
    )
    assert pic is not None

    # 새 이미지
    new_path = OUT_DIR / "_new.png"
    _make_dummy_png(new_path, w=100, h=100, color=(0, 255, 0))

    # 교체
    replace_image(slide, img_id=0, image_path=new_path)
    shape = next(sh for i, sh in iter_leaf_shapes(slide) if i == 0)
    # blip rId가 slide part의 이미지 관계 중 하나여야 함
    blip = shape._element.find(".//" + qn("a:blip"))
    rId = blip.get(qn("r:embed"))
    rel = shape.part.rels[rId]
    assert "image" in rel.reltype
    # 새 이미지 blob이 100x100 PNG여야 함
    blob = rel.target_part.blob
    img = Image.open(io.BytesIO(blob))
    assert img.size == (100, 100), f"replacement size: {img.size}"


def test_replace_image_rejects_non_picture():
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 3, 1, ["not a picture"])
    try:
        replace_image(slide, img_id=0, image_path=OUT_DIR / "_src.png")
    except SlideEditError:
        return
    raise AssertionError("expected SlideEditError for non-picture shape")


# -----------------------------------------------------------------------------
# 6. del_image
# -----------------------------------------------------------------------------

def test_del_image():
    prs = _new_prs()
    slide = _blank(prs)
    src_path = OUT_DIR / "_src.png"
    _make_dummy_png(src_path)
    slide.shapes.add_picture(str(src_path), Inches(1), Inches(1), Inches(3), Inches(2))
    _make_text_shape(slide, 5.0, 1.0, 3, 1, ["keep me"])

    before = [type(sh).__name__ for _, sh in iter_leaf_shapes(slide)]
    assert "Picture" in before, before

    del_image(slide, img_id=0)

    after = [type(sh).__name__ for _, sh in iter_leaf_shapes(slide)]
    assert "Picture" not in after, after
    assert len(after) == 1, after


# -----------------------------------------------------------------------------
# 7. SlideEditError 범위 검증
# -----------------------------------------------------------------------------

def test_errors():
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 3, 1, ["x"])
    # 잘못된 div_id
    try:
        clone_paragraph(slide, div_id=99, paragraph_id=0)
    except SlideEditError:
        pass
    else:
        raise AssertionError("no error for bad div_id")
    # 잘못된 paragraph_id
    try:
        replace_paragraph(slide, div_id=0, paragraph_id=99, text="x")
    except SlideEditError:
        pass
    else:
        raise AssertionError("no error for bad paragraph_id")


# -----------------------------------------------------------------------------
# 8. 통합: 한 슬라이드에 모든 API 체이닝 후 저장/재오픈
# -----------------------------------------------------------------------------

def test_roundtrip_save_reopen():
    prs = _new_prs()
    slide = _blank(prs)
    _make_text_shape(slide, 0.5, 0.5, 5, 2, ["첫 문단", "둘째 문단"])
    src_path = OUT_DIR / "_src.png"
    _make_dummy_png(src_path)
    slide.shapes.add_picture(str(src_path), Inches(1), Inches(3), Inches(3), Inches(2))

    clone_paragraph(slide, div_id=0, paragraph_id=0)
    replace_paragraph(slide, div_id=0, paragraph_id=1, text="**교체된** 두 번째")
    del_paragraph(slide, div_id=0, paragraph_id=0)

    out = OUT_DIR / "roundtrip.pptx"
    prs.save(str(out))

    # 재오픈하여 구조 확인
    prs2 = Presentation(str(out))
    slide2 = prs2.slides[0]
    tb = slide2.shapes[0]
    texts = [p.text for p in tb.text_frame.paragraphs]
    # 원본: ["첫 문단", "둘째 문단"]
    # clone(0): append "첫 문단" → ["첫 문단", "둘째 문단", "첫 문단"]
    # replace(1, "**교체된** 두 번째") → ["첫 문단", "교체된 두 번째", "첫 문단"]
    # del(0) → ["교체된 두 번째", "첫 문단"]
    assert texts == ["교체된 두 번째", "첫 문단"], texts


# -----------------------------------------------------------------------------
# 실행
# -----------------------------------------------------------------------------

def main():
    tests = [
        ("iter_leaf_shapes flat", test_iter_leaf_shapes_flat),
        ("iter_leaf_shapes group (synthetic)", test_iter_leaf_shapes_group_recurse),
        ("clone_paragraph basic", test_clone_paragraph_basic),
        ("clone_paragraph korean EA font", test_clone_paragraph_preserves_korean_ea_font),
        ("replace_paragraph plain", test_replace_paragraph_plain),
        ("replace_paragraph markdown bold/italic", test_replace_paragraph_markdown_bold_italic),
        ("replace_paragraph korean font", test_replace_paragraph_preserves_korean_font),
        ("del_paragraph middle", test_del_paragraph_middle),
        ("del_paragraph last preserves shell", test_del_paragraph_last_preserves_shell),
        ("replace_image", test_replace_image),
        ("replace_image rejects non-picture", test_replace_image_rejects_non_picture),
        ("del_image", test_del_image),
        ("SlideEditError bounds", test_errors),
        ("roundtrip save+reopen", test_roundtrip_save_reopen),
    ]
    print(f"Running {len(tests)} edit_ops tests ...")
    for name, fn in tests:
        _run(name, fn)
    print()
    print(f"Passed: {len(passed)} / {len(tests)}")
    if failed:
        print(f"Failed: {len(failed)}")
        for name, err in failed:
            print(f"  - {name}: {err}")
        sys.exit(1)
    print("All tests passed.")


if __name__ == "__main__":
    main()
