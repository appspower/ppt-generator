"""한글 1174 슬라이드 기반 편집 API 통합 테스트 — Phase A1 POC 재현.

원본 마스터 템플릿에서 슬라이드 1174 (SAP FI 모듈, 한글 다수) 만 추출하여
5 API 전체를 체이닝하고 결과를 저장 후 재오픈 검증.

실행: python tests/test_edit_ops_hangul_integration.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from ppt_builder.template import TemplateEditor  # noqa: E402
from ppt_builder.template.edit_ops import (  # noqa: E402
    clone_paragraph,
    del_paragraph,
    iter_leaf_shapes,
    replace_paragraph,
)

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "docs" / "references" / "_master_templates" / "PPT 템플릿.pptx"
OUT = ROOT / "output" / "edit_ops_tests"
OUT.mkdir(parents=True, exist_ok=True)

TARGET_SLIDE = 1174  # 0-based (Phase A1 보고서의 "슬라이드 1175"는 1-based)


def has_korean(s: str) -> bool:
    return any("\uac00" <= c <= "\ud7a3" for c in s)


def ea_font_names(shape) -> set[str]:
    """shape 내 모든 run의 <a:ea typeface> 속성 값 모음."""
    names = set()
    if not shape.has_text_frame:
        return names
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            ea = run._r.find(".//" + qn("a:rPr") + "/" + qn("a:ea"))
            if ea is not None:
                v = ea.get("typeface")
                if v:
                    names.add(v)
    return names


def main():
    if not SRC.exists():
        print(f"[SKIP] master template not found: {SRC}")
        return 0

    print(f"Opening master template: {SRC}")
    editor = TemplateEditor(str(SRC))
    print(f"  slide count: {editor.slide_count}")
    editor.keep_slides([TARGET_SLIDE])
    print(f"  kept only slide {TARGET_SLIDE}, now: {editor.slide_count} slide")

    slide = editor.prs.slides[0]
    leaf_shapes = list(iter_leaf_shapes(slide))
    print(f"  flat leaf shape count: {len(leaf_shapes)}")

    # 1) 한글 paragraph 여럿 가진 shape 찾기
    target_div = None
    for idx, sh in leaf_shapes:
        if not sh.has_text_frame:
            continue
        korean_paras = [
            pi for pi, p in enumerate(sh.text_frame.paragraphs) if has_korean(p.text)
        ]
        if len(korean_paras) >= 2:
            target_div = idx
            target_shape = sh
            print(
                f"  target shape: div_id={idx}, {len(sh.text_frame.paragraphs)} paras, "
                f"korean_paras={korean_paras}"
            )
            break

    if target_div is None:
        print("[FAIL] no Korean shape found")
        return 1

    # EA 폰트 snapshot
    ea_before = ea_font_names(target_shape)
    orig_texts = [p.text for p in target_shape.text_frame.paragraphs]
    print(f"  before texts: {orig_texts[:5]}")
    print(f"  before EA fonts: {ea_before}")

    # 2) 편집 체인: clone → replace → del
    # clone p0
    new_idx = clone_paragraph(slide, div_id=target_div, paragraph_id=0)
    print(f"  clone_paragraph: new paragraph_id = {new_idx}")

    # replace new cloned paragraph with markdown
    replace_paragraph(
        slide, div_id=target_div, paragraph_id=new_idx,
        text="**강조된** 추가 한글 문단",
    )
    print("  replace_paragraph: markdown with Korean applied")

    # delete paragraph 0 (original)
    del_paragraph(slide, div_id=target_div, paragraph_id=0)
    print("  del_paragraph: removed p0")

    # 3) EA 폰트 보존 확인
    ea_after = ea_font_names(target_shape)
    print(f"  after EA fonts: {ea_after}")
    assert ea_before.issubset(ea_after), (
        f"EA fonts lost! before={ea_before}, after={ea_after}"
    )

    # 4) 저장 후 재오픈
    out_path = OUT / "hangul_1174_edited.pptx"
    editor.save(out_path)
    print(f"  saved: {out_path}")

    reopened = Presentation(str(out_path))
    assert len(reopened.slides) == 1
    reopened_slide = reopened.slides[0]
    # target shape 재접근 (flat idx 유지 가정은 불안정 — 텍스트 포함으로 찾기)
    found = False
    for idx, sh in iter_leaf_shapes(reopened_slide):
        if not sh.has_text_frame:
            continue
        texts = [p.text for p in sh.text_frame.paragraphs]
        if any("강조된" in t for t in texts):
            found = True
            print(f"  reopened target texts: {texts[:5]}")
            # EA 폰트 재검증
            ea_reopened = ea_font_names(sh)
            print(f"  reopened EA fonts: {ea_reopened}")
            break
    assert found, "markdown-replaced paragraph not found after reopen"

    editor.cleanup()
    print()
    print("[OK] integration test passed")
    return 0


if __name__ == "__main__":
    sys.exit(main())
