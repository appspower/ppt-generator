"""Exhibit-First 프로세스 검증 — 3장.

이전 12장에서 사용하지 않은 컴포넌트가 자연스럽게 선택되는지 검증.
- Slide A: comp_gantt_bars (타임라인) — 0회 → 최초 사용
- Slide B: comp_before_after (전/후) — 0회 → 최초 사용
- Slide C: comp_funnel (퍼널) — 0회 → 최초 사용
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.composer import SlideComposer
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_kpi_card, comp_kpi_row, comp_bullet_list,
    comp_gantt_bars, comp_before_after, comp_funnel,
    comp_styled_card,
)

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "exhibit_first_test"


def make(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ============================================================
# Slide A: 제약 특허 절벽 — comp_gantt_bars (최초 사용!)
# Phase A: "2026~28 특허 절벽으로 $200B+ 매출이 위험"
# Phase B: 시간축 위 이벤트 → 타임라인 → comp_gantt_bars
# Phase C: gantt(주인공) + kpi_row(보강) → top_bottom
# ============================================================
def slide_a_patent_cliff(prs):
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="빅파마 $200B 매출의 특허 절벽이 2026~28년에 집중되며, M&A $240B가 급증하고 있다",
        category="Pharma Patent Cliff"))
    zones = comp.layout("top_bottom", split=0.18)

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "$200B+", "label": "특허 만료 위험 매출", "trend": "up"},
        {"value": "$240B", "label": "2025 M&A 거래 (YoY+81%)", "trend": "up"},
        {"value": "$29.5B", "label": "Keytruda 단일 최대 리스크", "trend": "up"},
    ], region=zones["top"])

    # phases = 연도 열, streams = 약물별 행, bars = 특허 만료 기간
    comp_gantt_bars(comp.canvas,
                    phases=["2025", "2026", "2027", "2028", "2029"],
                    streams=[
                        {"name": "Keytruda\n(Merck $29.5B)", "bars": [
                            {"start": 1, "end": 3, "label": "특허 만료 → SC 전환 시도", "highlight": True}]},
                        {"name": "Eliquis\n(BMS $13B)", "bars": [
                            {"start": 1, "end": 2, "label": "제네릭 진입"}]},
                        {"name": "Opdivo\n(BMS $9B)", "bars": [
                            {"start": 2, "end": 4, "label": "바이오시밀러 경쟁"}]},
                        {"name": "Januvia\n(Merck $2.3B)", "bars": [
                            {"start": 1, "end": 1.5, "label": "제네릭 런칭"}]},
                        {"name": "Trulicity\n(Lilly $5.2B)", "bars": [
                            {"start": 2, "end": 3, "label": "GLP-1 전환"}]},
                    ],
                    milestones=[
                        {"phase": 3, "label": "Keytruda IRA 가격 적용 (2028.1)"},
                    ],
                    region=zones["bottom"])

    comp.takeaway("Keytruda SC 전환(2042년 특허)이 성공 여부가 Merck의 향후 14년을 결정 — M&A 외 자체 파이프라인 확보 시급")
    comp.footer(SlideFooter(source="GEN Top 20 Patent Cliff, Labiotech, DeepCeutix $300B Report, BioSpace"))


# ============================================================
# Slide B: AI 에이전트 vs SaaS — comp_before_after (최초 사용!)
# Phase A: "AI 에이전트가 SaaS를 대체하며 구조를 재편한다"
# Phase B: 전/후 대비 → comp_before_after
# Phase C: before_after(주인공) + kpi_row(보강) → l_layout
# ============================================================
def slide_b_ai_agent_saas(prs):
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="AI 에이전트가 SaaS 시장을 재편하며, 기업 소프트웨어 예산의 50%가 재배분되고 있다",
        category="AI Agents vs SaaS"))
    zones = comp.layout("l_layout", left_ratio=0.60, top_ratio=0.50)

    comp_before_after(comp.canvas,
                      before_title="기존 SaaS 모델",
                      before_items=[
                          {"label": "도구별 라이센스 구매", "detail": "CRM, ERP, HR 각각 별도 계약"},
                          {"label": "사람이 도구에 접속", "detail": "데이터 입력·조회·보고서 수동 생성"},
                          {"label": "도구 간 연동 비용", "detail": "통합 프로젝트 별도 발주, 유지보수 지속"},
                          {"label": "비용 선형 증가", "detail": "라이센스 수 × 사용자 수", "kpi": "평균 IT 예산의 32%"},
                      ],
                      after_title="AI 에이전트 모델",
                      after_items=[
                          {"label": "에이전트가 도구 대신 조작", "detail": "CRM 업데이트, 보고서, 일정 자동 처리"},
                          {"label": "의도만 전달", "detail": "자연어로 지시 → 에이전트가 실행·검증"},
                          {"label": "오케스트레이션 자동화", "detail": "도구 간 연결을 에이전트가 자동 조율"},
                          {"label": "비용 1/10", "detail": "에이전트 수 기반 과금", "kpi": "SaaS 라이센스 50% 감축"},
                      ],
                      arrow_label="전환",
                      region=zones["left_full"])

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "$2T", "label": "SW 시가총액 증발 (2026.2)", "trend": "down"},
        {"value": "46.3%", "label": "AI 에이전트 시장 CAGR", "trend": "up"},
    ], region=zones["right_top"])

    comp_bullet_list(comp.canvas, title="기업 대응 전략",
                     items=[
                         "**결정론적 시스템(ERP·HR)은 유지** — LLM은 의도 해석, 실제 실행은 기존 시스템이 담당하는 하이브리드 아키텍처",
                         "**확률적 시스템(보고서·분석)은 대체** — Publicis Sapient: Adobe 등 SaaS 라이센스 50% 감축, GenAI로 전환",
                         "**오케스트레이션 레이어 구축** — 기존 시스템 위에 AI 에이전트가 동작하는 중간층이 핵심 투자 영역",
                     ], region=zones["right_bottom"])

    comp.takeaway("SaaS '도구 접속' 모델에서 AI 에이전트 '의도 전달' 모델로 전환 — 결정론적 시스템은 유지, 확률적 시스템은 대체")
    comp.footer(SlideFooter(source="Bain AI+SaaS 2025, Deloitte TMT Predictions 2026, Zylos Research"))


# ============================================================
# Slide C: EV 배터리 리사이클링 — comp_funnel (최초 사용!)
# Phase A: "$35B 시장으로 성장, 리튬 회수율 50%→90% 목표"
# Phase B: 전체→수거→분해→회수 = 단계별 축소 → funnel
# Phase C: funnel(주인공) + kpi_row(보강) → t_layout
# ============================================================
def slide_c_battery_recycling(prs):
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="EV 배터리 리사이클링이 $43B 시장으로 성장하며, EU가 리튬 회수 50%를 의무화했다",
        category="Battery Recycling"))
    zones = comp.layout("t_layout", top_ratio=0.15, right_ratio=0.42)

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "$43B", "label": "2030 시장 규모", "detail": "CAGR 15.3%", "trend": "up"},
        {"value": "50%", "label": "EU 리튬 회수 목표 (2027)", "trend": "up"},
        {"value": "90%", "label": "Co·Ni·Cu 회수율 목표", "trend": "up"},
    ], region=zones["top"])

    comp_funnel(comp.canvas, stages=[
        {"label": "폐배터리 발생", "value": "1,200만톤/yr (2030)", "detail": "EV 보급 확대로 폐배터리 급증 — 연 15%+ 성장"},
        {"label": "수거·운송", "value": "수거율 45%", "detail": "수거 인프라 구축 중. 운송 안전 규제가 비용 상승 요인"},
        {"label": "방전·해체", "value": "비용의 30%", "detail": "모듈→셀 분리. 안전 비용이 전체 리사이클링 원가의 30%"},
        {"label": "습식/건식 공정", "value": "블랙매스 추출", "detail": "Redwood Materials·Li-Cycle이 기술 선도"},
        {"label": "원자재 회수", "value": "Li 50%·Co 90%", "detail": "EU 의무: 리튬 50%, 코발트·니켈 90% 회수. 배터리급 순도"},
    ], region=zones["bottom_left"])

    comp_bullet_list(comp.canvas, title="한국 산업 시사점",
                     items=[
                         "**성일하이텍·에코프로** — 국내 배터리 리사이클링 선도 기업. 유럽·미국 현지 공장 투자 확대 중",
                         "**EU 배터리 규정(2025~)** — 재활용 효율·원자재 회수율 의무화. 한국 수출 배터리에도 적용",
                         "**IRA 핵심 광물 요건** — 미국 EV 보조금 수령 조건에 재활용 원자재 포함. 리사이클링이 보조금 확보 전략",
                         "**도시 광산(Urban Mining)** — 폐배터리에서 추출한 리튬이 채굴 대비 CO₂ 70% 절감. ESG+경제성 동시 달성",
                     ], region=zones["bottom_right"])

    comp.takeaway("폐배터리 리사이클링은 환경 의무이자 $43B 사업 기회 — EU·IRA 규제가 시장을 견인")
    comp.footer(SlideFooter(source="Knowledge Sourcing EV Recycling 2030, EU Battery Regulation, Precedence Research"))


# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_a_patent_cliff(prs)
    slide_b_ai_agent_saas(prs)
    slide_c_battery_recycling(prs)

    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path}")

    from ppt_builder.evaluate import evaluate_pptx, print_report
    report = evaluate_pptx(str(pptx_path))
    print_report(report)

    import pythoncom, win32com.client
    pythoncom.CoInitialize()
    pdf_path = pptx_path.with_suffix(".pdf")
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    p = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
    p.SaveAs(str(pdf_path.resolve()), 32)
    p.Close()
    print(f"PDF:  {pdf_path}")

    from ppt_builder.track_c.png_export import pptx_to_pngs
    png_dir = OUTPUT_DIR / f"{NAME}_pngs"
    paths = pptx_to_pngs(pptx_path, png_dir)
    print(f"PNGs: {png_dir} ({len(paths)}장)")


if __name__ == "__main__":
    main()
