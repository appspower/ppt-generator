"""HD현대 SAP S/4HANA 전환 제안서 — 10장 풀 덱 생성.

워크플로우 전체를 태우는 통합 테스트.
22개 패턴 + Composer 조합을 골고루 사용.
"""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.patterns import (
    SlideFooter, SlideHeader,
    ExecutiveSpec, executive_summary,
    BeforeAfterSpec, before_after,
    TimelineSpec, timeline_phases,
    HubSpokeSpec, hub_spoke,
    ProcessSpec, process_flow,
    ValueChainSpec, value_chain,
    QuadrantSpec, quadrant_story,
    GanttSpec, gantt_roadmap,
    TreeSpec, tree_diagram,
)
from ppt_builder.composer import SlideComposer, apply_zone_tone
from ppt_builder.components import (
    comp_data_card, comp_progress_bar, comp_icon_list,
    comp_metric_delta,
)
from ppt_builder.visual_validate import validate_visual

FOOTER = SlideFooter(
    source="출처: Palantir AIP ERP Migration Suite, PwC Analysis 2024",
    right="HD현대",
)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ================================================================
    # Slide 1: Executive Summary
    # ================================================================
    executive_summary(add_slide(prs), ExecutiveSpec(
        header=SlideHeader(
            title="Palantir 단일 플랫폼으로 SAP 전환 일정 14%·테스트 70%·DT 50% 단축 효과를 확보",
            category="1. 제안 개요 — Executive Summary",
            nav_path=["1. 제안 개요", "2. 활용 전략"],
        ),
        hero_label="WHY NOW",
        hero_headline="SAP 전환의 3대 병목을\n단일 Ontology로 해소",
        hero_subtitle="테스트·Cutover·거버넌스를 단일 플랫폼으로 통합 관리",
        bottlenecks=[
            {"num": "01", "title": "테스트 자동화",
             "kpi": "공수 70%↓ · 정확도 99.8%",
             "bullets": ["AIP가 Blueprint→규칙 자동 추출", "LLM-as-Judge 연속 검증"]},
            {"num": "02", "title": "Cutover 오케스트레이션",
             "kpi": "초과율 10%↓ · DT 50%↓",
             "bullets": ["Ontology Task DAG", "리허설 자동 비교"]},
            {"num": "03", "title": "프로젝트 거버넌스",
             "kpi": "보고 실시간 · 일정 14%↓",
             "bullets": ["Health 대시보드", "Readiness Scorecard"]},
        ],
        kpis=[
            {"value": "14%", "label": "전체 일정 단축", "detail": "18→15.5개월"},
            {"value": "70%", "label": "테스트 공수 절감", "detail": "수작업→AIP"},
            {"value": "50%", "label": "Cutover DT 단축", "detail": "초과율→10%↓"},
            {"value": "2~3주", "label": "Quick Win 입증", "detail": "L1 즉시 ROI"},
        ],
        roadmap_phases=[
            {"tag": "L1", "name": "가시화", "duration": "2~3주",
             "deliverables": ["Health 대시보드", "Config Register", "보고 자동화"]},
            {"tag": "L2", "name": "자동화", "duration": "4~6주",
             "deliverables": ["AIP 테스트 생성", "FDD 초안", "결함 Triage"]},
            {"tag": "L3", "name": "최적화", "duration": "8~12주",
             "deliverables": ["Cutover 앱", "Process Mining"]},
            {"tag": "L4", "name": "지능화", "duration": "9개월+",
             "deliverables": ["AI Go/No-Go", "Hypercare 탐지"]},
        ],
        takeaway="Quick Win L1(2~3주) 가치 입증 → L2~L4 점진 확대 — 단일 플랫폼 3대 리스크 통합 관리",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 2: Before / After
    # ================================================================
    before_after(add_slide(prs), BeforeAfterSpec(
        header=SlideHeader(
            title="수작업 3~6개월 테스트를 AIP 자동화로 1~2주 파이프라인으로 전환",
            category="2. 핵심 변화 — AS-IS vs TO-BE",
            nav_path=["1. 제안 개요", "3. 프로세스 혁신"],
        ),
        intro="기존 수작업 프로세스와 AIP 자동화 이후를 5개 차원에서 정량 비교",
        before_title="AS-IS (수작업)",
        after_title="TO-BE (AIP 자동화)",
        arrow_label="전환",
        before_items=[
            {"label": "TC 작성", "detail": "Blueprint 수동 분석 → Excel", "kpi": "3~6개월"},
            {"label": "검증 방식", "detail": "육안 리뷰, 누락 빈번", "kpi": "정확도 85%"},
            {"label": "결함 등록", "detail": "Jira 수동 입력", "kpi": "지연 2~3일"},
            {"label": "커버리지", "detail": "핵심 시나리오만", "kpi": "40~60%"},
            {"label": "비용", "detail": "인건비 중심", "kpi": "건당 $50~100"},
        ],
        after_items=[
            {"label": "TC 작성", "detail": "AIP Blueprint→규칙→TC 자동", "kpi": "1~2주"},
            {"label": "검증 방식", "detail": "LLM-as-Judge 0~10점", "kpi": "정확도 99.8%"},
            {"label": "결함 등록", "detail": "Jira REST API 자동", "kpi": "실시간"},
            {"label": "커버리지", "detail": "전수+엣지케이스", "kpi": "95%+"},
            {"label": "비용", "detail": "플랫폼 고정비", "kpi": "건당 $0.2"},
        ],
        takeaway="AIP 전환 시 공수 70%↓, 정확도 85→99.8%, 커버리지 60→95% — 건당 $50→$0.2 (250배↓)",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 3: Timeline Phases (4단계 로드맵)
    # ================================================================
    timeline_phases(add_slide(prs), TimelineSpec(
        header=SlideHeader(
            title="Quick Win에서 AI 지능화까지 4단계 의존성 체인으로 점진 확대를 추진",
            category="3. 도입 로드맵 — PMO View",
            nav_path=["2. 활용 전략", "1. 단계별 계획"],
        ),
        intro="각 단계 산출물이 다음 단계의 입력 — Level 1 미완료 시 Level 2 이후 착수 불가",
        phases=[
            {"tag": "L1", "name": "가시화", "duration": "2~3주",
             "objective": "PMO 판단 속도 향상",
             "deliverables": ["Health Dashboard", "Config Register", "보고 자동화"],
             "metrics": "보고 3일→실시간",
             "prerequisites": "Jira API 권한", "gate": "갱신 ≤5분", "team": "Palantir 1+PwC 1"},
            {"tag": "L2", "name": "자동화", "duration": "4~6주",
             "objective": "테스트 공수 대폭 절감",
             "deliverables": ["AIP TC Generator", "FDD Auto-Draft", "Defect Triage"],
             "metrics": "TC 작성 70%↓",
             "prerequisites": "L1 Ontology+Blueprint", "gate": "커버리지≥80%", "team": "Palantir 2+PwC 1"},
            {"tag": "L3", "name": "최적화", "duration": "8~12주",
             "objective": "Go-Live 리스크 정량화",
             "deliverables": ["Cutover App", "Process Mining", "Readiness Scorecard"],
             "metrics": "초과율 10%↓",
             "prerequisites": "L2 테스트 안정화", "gate": "Mock 3회", "team": "Palantir 2+PwC Basis"},
            {"tag": "L4", "name": "지능화", "duration": "9개월+",
             "objective": "AI Agent 이상 탐지",
             "deliverables": ["AI Go/No-Go", "Hypercare Detector", "Ops Dashboard"],
             "metrics": "KPI 7일 SLA",
             "prerequisites": "L1~3 데이터 축적", "gate": "Critical 0건", "team": "Palantir 1+IT Ops 2"},
        ],
        takeaway="Level 1이 Quick Win이자 전체 기반 — 2~3주 내 가치 입증이 전략 확대의 전제 조건",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 4: Hub & Spoke (시스템 아키텍처)
    # ================================================================
    hub_spoke(add_slide(prs), HubSpokeSpec(
        header=SlideHeader(
            title="Foundry Ontology 중심으로 6개 모듈이 단일 데이터 허브에 연결",
            category="4. 시스템 아키텍처 — Hub & Spoke",
            nav_path=["2. 활용 전략", "2. 통합 구조"],
        ),
        intro="Foundry Ontology가 중심 허브 — 6개 모듈이 공통 데이터 모델을 공유하여 사일로 제거",
        hub={"title": "Foundry\nOntology", "subtitle": "통합 데이터 허브"},
        spokes=[
            {"badge": "L1", "title": "Health Dashboard", "detail": "실시간 가시화, PMO 자동화"},
            {"badge": "L1", "title": "Config Register", "detail": "설정 이력, 변경 영향 분석"},
            {"badge": "L2", "title": "AIP Test Gen", "detail": "TC 자동 생성, LLM 검증"},
            {"badge": "L2", "title": "Defect Triage", "detail": "결함 클러스터, 자동 분류"},
            {"badge": "L3", "title": "Cutover App", "detail": "Task DAG, 리허설 비교"},
            {"badge": "L4", "title": "AI Go/No-Go", "detail": "Scorecard, 이상 탐지"},
        ],
        takeaway="Ontology 중심 = 사일로 제로 — 신규 모듈 추가 시 연동 비용 O(1)",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 5: Process Flow (테스트 자동화 파이프라인)
    # ================================================================
    process_flow(add_slide(prs), ProcessSpec(
        header=SlideHeader(
            title="Blueprint→결함 등록 5단계 파이프라인으로 테스트 공수 70% 절감을 달성",
            category="5. 테스트 자동화 — Engineer View",
            nav_path=["2. 활용 전략", "3. 모듈 상세"],
        ),
        intro="수작업 3~6개월을 5단계 자동 파이프라인으로 압축",
        steps=[
            {"name": "Blueprint 수집", "actor": "PwC 컨설턴트",
             "tools": "PDF/Word\nFoundry Pipeline", "output": "구조화 Blueprint Corpus",
             "duration": "3~5일", "prerequisites": "문서 수령 완료"},
            {"name": "규칙 추출", "actor": "AIP Logic",
             "tools": "Interpretation AI\nLLM Chain", "output": "Rule Registry",
             "duration": "2~3일", "metrics": "Precision ≥95%"},
            {"name": "TC 자동 생성", "actor": "AIP Logic",
             "tools": "T-Code Mapping\nLLM Generator", "output": "TestCase 객체",
             "duration": "1~2일", "metrics": "커버리지 ≥80%"},
            {"name": "연속 검증", "actor": "LLM-as-Judge",
             "tools": "Rubric 0~10\nκ≥0.85", "output": "Pass/Fail+신뢰도",
             "duration": "수 시간", "example": "에너지社 20K+건 2주/$4K"},
            {"name": "결함 등록", "actor": "Foundry Connector",
             "tools": "Jira REST API\nWebhook", "output": "Jira 이슈+링크",
             "duration": "실시간", "risks": "Jira 권한 미확보 시 수동"},
        ],
        takeaway="총 1주 내 전체 파이프라인 가동 — 수작업 대비 60~70%↓, 정확도 99.8%",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 6: Value Chain (조선 가치사슬)
    # ================================================================
    value_chain(add_slide(prs), ValueChainSpec(
        header=SlideHeader(
            title="HD현대 조선 가치사슬에서 조달·건조가 원가 70% — SAP 전환 최우선 대상",
            category="6. 전략 분석 — Porter 가치사슬",
            nav_path=["1. 제안 개요", "4. 가치사슬"],
        ),
        intro="Porter 가치사슬로 5단계 주요활동과 3개 지원활동을 매핑 — SAP 전환 우선순위 도출",
        primary=[
            {"title": "수주·설계", "detail": "선박 수주, 기본·상세설계\nSAP SD·PS 적용"},
            {"title": "조달·구매", "detail": "강재·기자재 구매\n원가 35%. SAP MM 핵심", "highlight": True},
            {"title": "건조·생산", "detail": "블록 제작, 탑재, 의장\n원가 35%. SAP PP·PM", "highlight": True},
            {"title": "시운전·인도", "detail": "해상 시운전, 선급 검사\nSAP QM·PS"},
            {"title": "A/S·보증", "detail": "보증 관리, 부품 공급\nSAP CS·PM"},
        ],
        support=[
            {"title": "기술 인프라", "detail": "PLM·CAD, SAP S/4HANA, BTP, Foundry"},
            {"title": "인적자원 관리", "detail": "기능 인력 양성, SuccessFactors 성과 관리"},
            {"title": "재무·경영지원", "detail": "원가 관리, SAP FI/CO 실시간 경영 정보"},
        ],
        margin_label="마진",
        takeaway="조달(35%)+건조(35%)=원가 70% — SAP MM·PP 고도화가 마진 개선의 최대 레버",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 7: Quadrant Story (투자 우선순위)
    # ================================================================
    quadrant_story(add_slide(prs), QuadrantSpec(
        header=SlideHeader(
            title="8개 모듈을 ROI×난이도로 분류 — Quick Win 영역에서 시작을 권장",
            category="7. 투자 우선순위 — Investment View",
            nav_path=["2. 활용 전략", "4. 포트폴리오"],
        ),
        intro="단기 ROI와 구현 난이도 두 축으로 8개 모듈 분류",
        x_axis_label="구현 난이도", y_axis_label="단기 ROI",
        x_low="LOW", x_high="HIGH", y_low="LOW", y_high="HIGH",
        quadrants=[
            {"title": "QUICK WIN", "highlight": True,
             "items": ["Health 대시보드 (L1)", "Config Register (L1)", "주간 보고 자동화 (L1)"],
             "description": "2~3주 내 즉시 가치 입증", "action": "킥오프 직후 착수"},
            {"title": "STRATEGIC BET",
             "items": ["AIP 테스트 자동화 (L2)", "Cutover 오케스트레이션 (L3)"],
             "description": "L1 성공 후 확대", "action": "L1 Gate 통과 후 Palantir 2명 추가"},
            {"title": "NICE TO HAVE",
             "items": ["FDD 자동 초안 (L2)", "결함 Triage 클러스터 (L2)"],
             "description": "L2와 번들 추진 권장"},
            {"title": "DEFER",
             "items": ["AI Go/No-Go (L4)", "Hypercare 이상 탐지 (L4)"],
             "description": "L1~3 데이터 축적 후 재평가"},
        ],
        insight="L1 Quick Win 즉시 착수 → 가치 입증 후 L2 Strategic Bet 단계적 확대",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 8: Gantt Roadmap (실행 로드맵)
    # ================================================================
    gantt_roadmap(add_slide(prs), GanttSpec(
        header=SlideHeader(
            title="4개 스트림이 6분기에 걸쳐 병렬 수행 — Q5 Go-Live 목표",
            category="8. 실행 로드맵 — Gantt Chart",
            nav_path=["2. 활용 전략", "5. 통합 일정"],
        ),
        intro="Infra·Application·Data·Change 4개 스트림의 Q1~Q6 통합 일정",
        phases=["Q1", "Q2", "Q3", "Q4", "Q5", "Q6"],
        streams=[
            {"name": "Infra 구축", "bars": [
                {"start": 0, "end": 2, "label": "환경 구축"},
                {"start": 2, "end": 3, "label": "성능 튜닝"},
                {"start": 4, "end": 6, "label": "운영 전환"},
            ]},
            {"name": "Application", "bars": [
                {"start": 0, "end": 1, "label": "Fit-Gap"},
                {"start": 1, "end": 3, "label": "Build", "highlight": True},
                {"start": 3, "end": 5, "label": "테스트", "highlight": True},
            ]},
            {"name": "Data 마이그레이션", "bars": [
                {"start": 0, "end": 2, "label": "분석·설계"},
                {"start": 2, "end": 3, "label": "ETL 개발"},
                {"start": 3, "end": 5, "label": "검증·전환", "highlight": True},
            ]},
            {"name": "Change Mgmt", "bars": [
                {"start": 0, "end": 2, "label": "영향 분석"},
                {"start": 2, "end": 4, "label": "교육·소통"},
                {"start": 4, "end": 6, "label": "Hypercare"},
            ]},
        ],
        milestones=[
            {"phase": 1, "label": "Blueprint 확정"},
            {"phase": 3, "label": "Build 완료"},
            {"phase": 5, "label": "Go-Live"},
        ],
        takeaway="Application Build(Q2~Q3)와 Data 검증(Q4)이 Critical Path — 병렬 완료가 Go-Live 전제",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 9: Tree Diagram (리스크 분해)
    # ================================================================
    tree_diagram(add_slide(prs), TreeSpec(
        header=SlideHeader(
            title="전환 리스크를 기술·조직·프로세스·외부 4축으로 MECE 분해",
            category="9. 리스크 관리 — Issue Tree",
            nav_path=["3. 거버넌스", "1. 리스크 분해"],
        ),
        intro="SAP S/4HANA 전환 리스크를 MECE 원칙으로 4대 카테고리·12개 항목으로 구조화",
        root={"title": "전환 리스크"},
        branches=[
            {"title": "기술 리스크", "highlight": True,
             "children": ["데이터 정합성 실패", "인터페이스 호환성", "성능 저하 (사이징)"]},
            {"title": "조직 리스크",
             "children": ["핵심 인력 이탈", "현업 저항", "경영진 지원 약화"]},
            {"title": "프로세스 리스크", "highlight": True,
             "children": ["Scope Creep", "테스트 커버리지 부족", "Cutover 리허설 미흡"]},
            {"title": "외부 리스크",
             "children": ["SAP 라이선스 변경", "협력사 연동 지연", "규제 변경"]},
        ],
        takeaway="기술(데이터·IF)과 프로세스(테스트·Cutover)가 Top 4 — Foundry 자동화로 직접 억제",
        footer=FOOTER,
    ))

    # ================================================================
    # Slide 10: 기대 효과 종합 (Composer — data_card + progress_bar)
    # ================================================================
    slide10 = add_slide(prs)
    comp = SlideComposer(slide10)
    comp.header(SlideHeader(
        title="Palantir 도입 후 4대 KPI 전원 목표 초과 달성 — 워크스트림별 진행률 실시간 관리",
        category="10. 기대 효과 — Performance Summary",
        nav_path=["3. 성과 보고", "1. 종합 현황"],
    ))

    zones = comp.layout("top_bottom", split=0.32)

    # 상단: 스마트 데이터 카드 4개
    top = zones["top"]
    gap = 0.15
    card_w = (top.w - gap * 3) / 4
    cards = [
        {"value": 15.5, "label": "전환 일정 (개월)", "previous": 18, "target": 16,
         "unit": "개월", "higher_is_better": False, "detail": "18→15.5 (목표 16)"},
        {"value": 70, "label": "테스트 공수 절감", "previous": 0, "target": 60,
         "unit": "%", "detail": "목표 60→실제 70%"},
        {"value": 99.8, "label": "TC 정확도", "previous": 85, "target": 95,
         "unit": "%", "detail": "기존 85→현재 99.8%"},
        {"value": 50, "label": "Cutover DT 단축", "previous": 0, "target": 40,
         "unit": "%", "detail": "목표 40→실제 50%"},
    ]
    for i, card in enumerate(cards):
        comp_data_card(comp.canvas, region=top.sub(i * (card_w + gap), 0, card_w, top.h), **card)

    # 하단: 워크스트림 진행률
    bottom = zones["bottom"]
    apply_zone_tone(comp.canvas, bottom, "subtle", border=False)
    comp.canvas.push_region(bottom)
    comp.canvas.section_label("워크스트림 진행률", x=0.1, y=0.08, w=bottom.w - 0.2)
    comp.canvas.pop_region()

    streams = [
        ("FI/CO 재무회계", 95, 80), ("MM 구매관리", 82, 80),
        ("SD 영업물류", 78, 80), ("PP 생산계획", 65, 80),
        ("데이터 마이그레이션", 55, 80), ("인터페이스 개발", 72, 80),
    ]
    bar_h = (bottom.h - 0.5) / len(streams)
    for i, (name, val, tgt) in enumerate(streams):
        comp_progress_bar(comp.canvas, label=name, value=val, target=tgt,
                          region=bottom.sub(0.1, 0.42 + i * bar_h, bottom.w - 0.2, bar_h))

    comp.takeaway("4대 KPI Green + 워크스트림 6개 중 4개 Green — PP(65%)·데이터(55%) 집중 대응 필요")
    comp.footer(FOOTER)

    return prs


def main():
    out_dir = Path("output/full_deck")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "HD현대_SAP_전환_제안서_10장.pptx"

    print("=" * 70)
    print("HD현대 SAP S/4HANA 전환 제안서 — 10장 풀 덱 생성")
    print("=" * 70)

    try:
        prs = build_deck()
        prs.save(out)
        print(f"\n생성 완료: {out}")
    except Exception as e:
        print(f"\nBUILD FAILED: {e}")
        import traceback
        traceback.print_exc()
        return 1

    # 검증
    visual = validate_visual(out, convert_pdf=False)
    print(f"\nvisual: {len(visual.issues)} issues")
    for i in visual.issues[:10]:
        print(f"  - {i}")

    # PDF 변환
    print("\nPDF 변환 중...")
    try:
        visual = validate_visual(out, convert_pdf=True)
        print(f"PDF: {'OK' if visual.pdf_available else 'SKIP'}")
    except Exception as e:
        print(f"PDF: FAIL ({e})")

    print("\n" + "=" * 70)
    print(f"결과: {'PASS' if not visual.issues else 'ISSUES FOUND'}")
    print("=" * 70)
    return 0


if __name__ == "__main__":
    sys.exit(main())
