"""Layer 2 (primitives + design_check + patterns) 테스트.

기존 test_workflow.py와 같은 직접 실행 모드를 따른다 (pytest 선택적).
"""

from __future__ import annotations

import sys
from pathlib import Path

try:
    import pytest  # type: ignore
    HAS_PYTEST = True
except ImportError:
    HAS_PYTEST = False

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import COLORS, Canvas, color
from ppt_builder.design_check import (
    DesignDecision,
    DesignReport,
    decide_density,
    decide_emphasis_color,
    decide_layout_archetype,
    decide_number_marker,
    inspect_design,
)
from ppt_builder.patterns import (
    SlideFooter,
    SlideHeader,
    ComparisonSpec,
    ExecutiveSpec,
    ProcessSpec,
    QuadrantSpec,
    TimelineSpec,
    comparison_matrix,
    executive_summary,
    process_flow,
    quadrant_story,
    timeline_phases,
)


# ============================================================
# Helpers
# ============================================================


def _new_slide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _save(prs, tmp_path: Path, name: str) -> Path:
    out = tmp_path / name
    prs.save(out)
    return out


# ============================================================
# primitives.py
# ============================================================


def test_color_helper_string_alias():
    assert color("white") == COLORS["white"]
    assert color("grey_800") == COLORS["grey_800"]


def test_color_helper_hex():
    rgb = color("#FD5108")
    assert rgb == COLORS["accent"]


def test_color_helper_unknown_raises():
    try:
        color("nope")
    except ValueError:
        return
    raise AssertionError("expected ValueError")


def test_canvas_box_text_creates_shapes(tmp_path: Path):
    prs, slide = _new_slide()
    c = Canvas(slide)
    c.box(x=0.3, y=0.3, w=2, h=1)
    c.text("hello", x=0.3, y=0.3, w=2, h=1)
    out = _save(prs, tmp_path, "canvas_basic.pptx")
    assert out.exists()
    re = Presentation(out)
    assert len(re.slides[0].shapes) >= 2


def test_canvas_circle_chevron(tmp_path: Path):
    prs, slide = _new_slide()
    c = Canvas(slide)
    c.circle(x=1, y=1, d=0.4, text="01")
    c.chevron(x=2, y=1, w=2, h=0.5, text="L1")
    out = _save(prs, tmp_path, "circle_chevron.pptx")
    assert out.exists()


def test_canvas_phase_a_composites(tmp_path: Path):
    prs, slide = _new_slide()
    c = Canvas(slide)
    c.badge("NEW", x=0.3, y=0.3)
    c.callout_box(x=0.3, y=1, w=4, h=1.5, title="t", body="body")
    c.arrow_chain(["A", "B", "C"], x=5, y=1, w=4.5, h=0.5)
    c.dot_grid(x=5, y=2, filled=2)
    c.mini_table(
        x=0.3, y=3, w=4, h=1.5,
        headers=["c1", "c2"], rows=[["1", "2"], ["3", "4"]],
    )
    c.stat_block(value="42%", label="x", x=5, y=3, w=2)
    c.numbered_list([("a", "d1"), ("b", "d2")], x=0.3, y=5, w=4)
    c.section_label("LBL", x=5, y=5, w=4)
    out = _save(prs, tmp_path, "composites.pptx")
    assert out.exists()


# ============================================================
# design_check.py
# ============================================================


def test_decide_number_marker_chevron_for_sequence():
    d = decide_number_marker(item_count=4, has_sequence=True, space_h=0.5)
    assert d.recommendation["shape"] == "chevron"


def test_decide_number_marker_circle_for_list():
    d = decide_number_marker(item_count=3, has_sequence=False, space_h=1.0)
    assert d.recommendation["shape"] == "circle"


def test_decide_emphasis_color_grey_default():
    d = decide_emphasis_color(palette="grey")
    assert "grey" in d.recommendation["primary"]


def test_decide_density_action_words():
    enrich = decide_density(available_area=10, intended_chars=100)
    assert enrich.recommendation["action"] == "enrich"
    over = decide_density(available_area=10, intended_chars=4000)
    assert over.recommendation["action"] == "reduce_or_expand"
    good = decide_density(available_area=10, intended_chars=1500)
    assert good.recommendation["action"] == "ok"


def test_decide_layout_archetype_executive():
    d = decide_layout_archetype(intent="executive", item_count=3)
    assert "hero" in d.recommendation["pattern"].lower()


def test_inspect_design_returns_report(tmp_path: Path):
    prs, slide = _new_slide()
    c = Canvas(slide)
    c.box(x=0.3, y=0.3, w=9, h=6, fill="white", border=0.75, border_color="grey_mid")
    c.text("test content", x=0.5, y=0.5, w=8, h=0.5, size=14)
    out = _save(prs, tmp_path, "design_test.pptx")
    report = inspect_design(str(out))
    assert isinstance(report, DesignReport)
    assert "slide_1_density" in report.metrics


def test_inspect_design_catches_forbidden_color(tmp_path: Path):
    prs, slide = _new_slide()
    c = Canvas(slide)
    c.box(x=1, y=1, w=2, h=2, fill="accent", border=None)  # FD5108 (금지)
    out = _save(prs, tmp_path, "forbidden.pptx")
    report = inspect_design(str(out))
    assert any(i.category == "contrast" for i in report.issues)


def test_inspect_design_pattern_kind_floor(tmp_path: Path):
    """comparison 패턴은 짧은 셀이 정상 — density floor가 낮아야 함."""
    prs, slide = _new_slide()
    c = Canvas(slide)
    c.box(x=0.3, y=0.3, w=9, h=6, fill="white", border=None)
    # 매우 짧은 텍스트만
    c.text("✓", x=2, y=2, w=1, h=0.3, size=10)
    out = _save(prs, tmp_path, "low_density.pptx")
    # default: high severity
    r1 = inspect_design(str(out))
    assert any(i.category == "density" for i in r1.issues)
    # comparison: 통과 (floor=6)
    r2 = inspect_design(str(out), pattern_kind="comparison")
    # density는 둘 다 같지만, comparison은 floor가 낮음
    high_density_issues = [i for i in r2.issues if i.category == "density" and i.severity == "high"]
    assert len(high_density_issues) <= len([i for i in r1.issues if i.category == "density" and i.severity == "high"])


# ============================================================
# patterns.py — 5개 패턴 모두 빌드 가능 확인
# ============================================================


def _minimal_footer():
    return SlideFooter(source="src", right="HD")


def _minimal_header(title="Test"):
    return SlideHeader(title=title, breadcrumb="bc")


def test_pattern_executive_builds(tmp_path: Path):
    prs, slide = _new_slide()
    spec = ExecutiveSpec(
        header=_minimal_header("Exec"),
        hero_label="WHY",
        hero_headline="head",
        hero_subtitle="sub",
        bottlenecks=[
            {"num": "01", "title": "A", "kpi": "k", "bullets": ["b1", "b2"]},
            {"num": "02", "title": "B", "kpi": "k", "bullets": ["b1", "b2"]},
            {"num": "03", "title": "C", "kpi": "k", "bullets": ["b1", "b2"]},
        ],
        kpis=[
            {"value": "1%", "label": "L"},
            {"value": "2%", "label": "L"},
            {"value": "3%", "label": "L"},
            {"value": "4%", "label": "L"},
        ],
        roadmap_phases=[
            {"tag": "L1", "name": "X", "duration": "1w", "deliverables": ["d"]},
            {"tag": "L2", "name": "Y", "duration": "2w", "deliverables": ["d"]},
            {"tag": "L3", "name": "Z", "duration": "3w", "deliverables": ["d"]},
            {"tag": "L4", "name": "W", "duration": "4w", "deliverables": ["d"]},
        ],
        takeaway="take",
        footer=_minimal_footer(),
    )
    executive_summary(slide, spec)
    out = _save(prs, tmp_path, "exec.pptx")
    assert out.exists()


def test_pattern_timeline_builds(tmp_path: Path):
    prs, slide = _new_slide()
    spec = TimelineSpec(
        header=_minimal_header("TL"),
        intro="intro",
        phases=[
            {"tag": "L1", "name": "X", "duration": "1w", "objective": "o",
             "deliverables": ["d1"], "metrics": "m"},
            {"tag": "L2", "name": "Y", "duration": "2w", "objective": "o",
             "deliverables": ["d2"], "metrics": "m"},
            {"tag": "L3", "name": "Z", "duration": "3w", "objective": "o",
             "deliverables": ["d3"], "metrics": "m"},
            {"tag": "L4", "name": "W", "duration": "4w", "objective": "o",
             "deliverables": ["d4"], "metrics": "m"},
        ],
        takeaway="t",
        footer=_minimal_footer(),
    )
    timeline_phases(slide, spec)
    out = _save(prs, tmp_path, "tl.pptx")
    assert out.exists()


def test_pattern_comparison_builds(tmp_path: Path):
    prs, slide = _new_slide()
    spec = ComparisonSpec(
        header=_minimal_header("CMP"),
        intro="i",
        criteria_labels=["c1", "c2", "c3"],
        options=[
            {"name": "A", "summary": "s", "criteria": ["1", "2", "3"]},
            {"name": "B", "summary": "s", "criteria": ["1", "2", "3"]},
            {"name": "C", "summary": "s", "criteria": ["1", "2", "3"], "highlight": True},
        ],
        takeaway="t",
        footer=_minimal_footer(),
    )
    comparison_matrix(slide, spec)
    out = _save(prs, tmp_path, "cmp.pptx")
    assert out.exists()


def test_pattern_process_builds(tmp_path: Path):
    prs, slide = _new_slide()
    spec = ProcessSpec(
        header=_minimal_header("PROC"),
        intro="i",
        steps=[
            {"name": "S1", "actor": "a", "tools": "t", "output": "o", "duration": "1d"},
            {"name": "S2", "actor": "a", "tools": "t", "output": "o", "duration": "2d"},
            {"name": "S3", "actor": "a", "tools": "t", "output": "o", "duration": "3d"},
        ],
        takeaway="t",
        footer=_minimal_footer(),
    )
    process_flow(slide, spec)
    out = _save(prs, tmp_path, "proc.pptx")
    assert out.exists()


def test_pattern_quadrant_builds(tmp_path: Path):
    prs, slide = _new_slide()
    spec = QuadrantSpec(
        header=_minimal_header("Q"),
        intro="i",
        x_axis_label="X", y_axis_label="Y",
        x_low="lo", x_high="hi", y_low="lo", y_high="hi",
        quadrants=[
            {"title": "TL", "items": ["a"], "highlight": True},
            {"title": "TR", "items": ["a"]},
            {"title": "BL", "items": ["a"]},
            {"title": "BR", "items": ["a"]},
        ],
        insight="ins",
        footer=_minimal_footer(),
    )
    quadrant_story(slide, spec)
    out = _save(prs, tmp_path, "q.pptx")
    assert out.exists()


# ============================================================
# 직접 실행 모드
# ============================================================


def _run_all_tests_directly() -> int:
    import inspect
    import tempfile
    import traceback

    test_fns = [
        (name, fn)
        for name, fn in globals().items()
        if name.startswith("test_") and callable(fn)
    ]

    passed = failed = 0
    for name, fn in test_fns:
        sig = inspect.signature(fn)
        try:
            if "tmp_path" in sig.parameters:
                with tempfile.TemporaryDirectory() as td:
                    fn(Path(td))
            else:
                fn()
            print(f"PASS  {name}")
            passed += 1
        except Exception as e:
            print(f"FAIL  {name}: {type(e).__name__}: {e}")
            traceback.print_exc()
            failed += 1

    print(f"\n{passed} passed, {failed} failed")
    return 0 if failed == 0 else 1


if __name__ == "__main__":
    sys.exit(_run_all_tests_directly())
