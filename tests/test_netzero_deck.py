"""글로벌 제조업 넷제로 전환 전략 — 워크플로우 6단계 준수.

완전히 새로운 주제, 새로운 리서치 기반.
Step 1~3 명시적 수행 후 Step 4 GENERATE.
"""

import sys
from pathlib import Path
_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches
from ppt_builder.primitives import Canvas, Region
from ppt_builder.patterns import (
    SlideFooter, SlideHeader,
    ExecutiveSpec, executive_summary,
    ComparisonSpec, comparison_matrix,
    HarveyBallSpec, harvey_ball_matrix,
    QuadrantSpec, quadrant_story,
    WaterfallSpec, waterfall_bridge,
    DataNarrativeSpec, data_narrative,
    ChevronTimelineSpec, chevron_timeline,
    GridProcessSpec, grid_process,
)
from ppt_builder.composer import SlideComposer, apply_zone_tone
from ppt_builder.components import (
    comp_data_card, comp_icon_card, comp_icon_row,
    comp_styled_card, comp_styled_card_row,
    comp_bar_chart_h, comp_bullet_list, comp_icon_list,
)
from ppt_builder.charts.native import chart_vertical_bar
from ppt_builder.visual_validate import validate_visual
from ppt_builder.design_check import inspect_design

FOOTER = SlideFooter(
    source="출처: IEA 2024, BloombergNEF 2025, McKinsey Net Zero 2024, POSCO IR 2024",
    right="PwC",
)

def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ================================================================
    # 1. Executive Summary — 제조업 탄소 위기 + 해법 + 로드맵
    # ================================================================
    executive_summary(add_slide(prs), ExecutiveSpec(
        header=SlideHeader(
            title="제조업 탄소 배출 9Gt(글로벌 25%)을 2030년까지 36% 감축하지 않으면 1.5°C 목표를 달성할 수 없음",
            category="1. 현황 진단 — Executive Summary",
            nav_path=["1. 진단", "2. 전략"],
        ),
        hero_label="NET ZERO 2030",
        hero_headline="제조업의 3대 탈탄소\n병목을 동시에 해소",
        hero_subtitle="기술·규제·투자 갭을 단일 전략으로 통합 관리",
        bottlenecks=[
            {"num": "01", "title": "기술 갭",
             "kpi": "그린수소 $4~9/kg → 목표 $2.50/kg",
             "bullets": ["5대 기술 중 상업화된 것은 전기화만", "CCUS 포집 비용 $40~120/톤으로 높음"]},
            {"num": "02", "title": "규제 충격",
             "kpi": "EU CBAM 2026 + K-ETS Phase 4",
             "bullets": ["미대응 시 철강 마진 10~18%p 하락", "Scope 3 공시 의무화 확대 중"]},
            {"num": "03", "title": "투자 갭",
             "kpi": "$2.1조 vs 필요 $4.5조 (갭 $2.4조)",
             "bullets": ["제조업 투자 비중 7%에 불과", "광업·공급망 추가 $3~4조 필요"]},
        ],
        kpis=[
            {"value": "9Gt", "label": "제조업 연간 CO₂", "detail": "글로벌 배출의 25%"},
            {"value": "25%", "label": "글로벌 배출 비중", "detail": "철강+시멘트+화학=20%"},
            {"value": "$2.4조", "label": "연간 투자 갭", "detail": "현재 $2.1조 vs 필요 $4.5조"},
            {"value": "2026", "label": "EU CBAM 시행", "detail": "€100/톤 벌금 부과"},
        ],
        roadmap_phases=[
            {"tag": "2025", "name": "측정·파일럿", "duration": "즉시",
             "deliverables": ["내부 탄소 가격 도입", "Scope 1·2 측정 체계", "수소 파일럿 FID"]},
            {"tag": "2027", "name": "기술 검증", "duration": "2년",
             "deliverables": ["수소 $2.50/kg 달성 여부", "CCUS 상업화 FID"]},
            {"tag": "2030", "name": "대규모 전환", "duration": "3년",
             "deliverables": ["재생에너지 3배", "산업 CO₂ 35%↓"]},
        ],
        takeaway="2025~2027년이 결정적 분기점 — 기술 파일럿 성공 여부가 2030 달성 가능성을 좌우함",
        footer=FOOTER,
    ))

    # ================================================================
    # 2. 규제 비교 — EU ETS vs K-ETS vs China ETS
    # ================================================================
    comparison_matrix(add_slide(prs), ComparisonSpec(
        header=SlideHeader(
            title="3대 탄소시장이 동시에 강화 — EU CBAM 2026 시행으로 수출 제조업에 직접 타격을 줌",
            category="2. 규제 환경 — 글로벌 ETS 비교",
            nav_path=["1. 진단", "3. 규제 분석"],
        ),
        intro="EU ETS·K-ETS·China ETS 3대 탄소시장의 가격·커버리지·특이점을 정량 비교 — 한국 수출기업의 이중 부담 분석",
        criteria_labels=[
            "탄소 가격 (2024)",
            "커버리지",
            "Phase/단계",
            "CBAM 영향",
            "Scope 3 요구",
            "벌금/제재",
        ],
        options=[
            {"name": "EU ETS", "summary": "가장 성숙한 체계",
             "criteria": ["~$80/톤 (최고)", "EU 배출 40%", "Phase 4 (2021~2030)",
                          "2026 완전 시행", "CSRD 의무화", "€100/톤 미제출 벌금"]},
            {"name": "K-ETS", "summary": "한국 배출 80% 커버",
             "criteria": ["~$6~11/톤", "국가 배출 80% (최광범)", "Phase 3→4 전환",
                          "수출기업 이중 부담", "자발적 단계", "할당 초과 시 과징금"],
             "highlight": True},
            {"name": "China ETS", "summary": "2024 확대 중",
             "criteria": ["$14.6/톤 (사상최고)", "국가 배출 40→60%", "2024 철강·시멘트 추가",
                          "해당 없음", "미도입", "강도 낮음"]},
        ],
        takeaway="K-ETS 탄소 가격($6~11)은 EU($80)의 1/10 수준이나, CBAM으로 EU 수출 시 EU 가격 적용 — 한국 철강·시멘트 수출기업 직접 타격",
        footer=FOOTER,
    ))

    # ================================================================
    # 3. 기술 비교 — Harvey Ball (5기술 × 4기준)
    # ================================================================
    harvey_ball_matrix(add_slide(prs), HarveyBallSpec(
        header=SlideHeader(
            title="5대 탈탄소 기술 중 전기화만 상업 성숙 — 수소·CCUS는 2027년 검증이 분기점임",
            category="3. 기술 지형 — 5대 기술 성숙도",
            nav_path=["2. 전략", "1. 기술 평가"],
        ),
        intro="그린수소·CCUS·전기화·순환경제·디지털트윈 5대 기술의 비용경쟁력·기술성숙도·감축잠재력·투자규모를 평가",
        row_labels=["그린수소 (DRI)", "CCUS (포집·저장)", "산업 전기화 (EAF·열펌프)",
                     "순환경제 (재활용)", "디지털 트윈 (에너지 최적화)"],
        col_labels=["비용 경쟁력", "기술 성숙도", "감축 잠재력", "투자 규모"],
        scores=[
            [1, 2, 4, 3],  # 수소: 비용 높음, 기술 중간, 감축 높음, 투자 큼
            [1, 2, 3, 3],  # CCUS: 비용 높음, 기술 중간, 감축 중상, 투자 큼
            [3, 4, 3, 2],  # 전기화: 비용 좋음, 기술 성숙, 감축 중상, 투자 중간
            [4, 3, 2, 1],  # 순환: 비용 최고, 기술 상, 감축 중간, 투자 작음
            [3, 3, 2, 1],  # DT: 비용 좋음, 기술 상, 감축 중간, 투자 작음
        ],
        highlight_row=2,  # 전기화 강조
        takeaway="전기화(EAF·열펌프)가 유일한 상업 성숙 기술 — 수소는 2027년 $2.50/kg 달성 여부가 대규모 전환의 전제 조건",
        footer=FOOTER,
    ))

    # ================================================================
    # 4. 투자 현황 — Composer(sidebar_left): 통계 + 바 차트
    # ================================================================
    slide4 = add_slide(prs)
    comp4 = SlideComposer(slide4)
    comp4.header(SlideHeader(
        title="글로벌 에너지 전환 투자 $2.1조 돌파했으나 제조업 비중은 7%에 불과 — 갭 $2.4조를 메워야 함",
        category="4. 투자 분석 — Investment Gap",
        nav_path=["2. 전략", "2. 투자 현황"],
    ))
    comp4.intro("2024년 사상 최초 $2조 돌파, 그러나 IEA NZE 경로 대비 연간 $2.4조 부족 — 제조업 집중 투자 시급")

    zones4 = comp4.layout("sidebar_left", sidebar_w=3.0)

    # 좌: 핵심 통계
    apply_zone_tone(comp4.canvas, zones4["sidebar"], "subtle", border=False)
    comp4.canvas.push_region(zones4["sidebar"])
    comp4.canvas.section_label("핵심 수치", x=0.1, y=0.08, w=zones4["sidebar"].w - 0.2)
    comp4.canvas.pop_region()
    stats = [
        {"value": 2.1, "label": "2024 투자 ($조)", "previous": 1.8, "target": 4.5, "unit": "조", "higher_is_better": True},
        {"value": 7, "label": "제조업 비중 (%)", "previous": 5, "target": 20, "unit": "%"},
        {"value": 2.4, "label": "연간 갭 ($조)", "unit": "조", "detail": "목표 $4.5조 대비"},
        {"value": 818, "label": "중국 투자 ($B)", "previous": 680, "unit": "B", "detail": "글로벌 39%"},
    ]
    card_h = (zones4["sidebar"].h - 0.5) / 4
    for i, st in enumerate(stats):
        comp_data_card(comp4.canvas,
                       region=zones4["sidebar"].sub(0.05, 0.42 + i * card_h, zones4["sidebar"].w - 0.1, card_h - 0.08),
                       **st)

    # 우: 네이티브 바 차트 (섹터별 투자)
    chart_vertical_bar(slide4,
                       categories=["전동화\n운송", "재생\n에너지", "전력망", "에너지\n저장", "수소", "CCUS"],
                       values=[757, 728, 390, 230, 42, 18],
                       highlight_idx=0,
                       region=zones4["main"])

    comp4.takeaway("전동화·재생에너지가 투자의 70%+ 집중 — 수소($42B)·CCUS($18B)는 전체의 3%로 심각한 과소 투자")
    comp4.footer(FOOTER)

    # ================================================================
    # 5. 섹터 비교 — Grid Process (4섹터 진도)
    # ================================================================
    grid_process(add_slide(prs), GridProcessSpec(
        header=SlideHeader(
            title="자동차가 R&D 4.4%로 탈탄소 선두, 시멘트는 0.6%로 최하위 — 섹터별 격차가 심각함",
            category="5. 섹터 비교 — 탈탄소 진도",
            nav_path=["2. 전략", "3. 산업별 분석"],
        ),
        intro="철강·시멘트·화학·자동차 4대 섹터의 R&D 투자·기술 성숙도·감축 실적·대표 기업을 비교",
        items=[
            {"number": "철강", "title": "ArcelorMittal",
             "detail": "R&D 1.3%. 감축 5.4% 달성.\nEAF 전환 착수. HyREX 파일럿.\n2030 목표 '점점 불가능' 인정"},
            {"number": "시멘트", "title": "Heidelberg Materials",
             "detail": "R&D 0.6% (최저). 넷제로\nevoZero 출시. Brevik CCUS\n가동. 감축 잠재력 50~90%"},
            {"number": "화학", "title": "BASF",
             "detail": "Scope 1·2 25% 감축 목표.\n전해조 €5,900만 투자.\n2030년까지 €20~30억 추가"},
            {"number": "자동차", "title": "BMW / 현대차",
             "detail": "R&D 4.4% (최고). 전 수명\n주기 전략. EV 전환 가속.\n현대 EV 30만대 (2024)"},
        ],
        takeaway="R&D 투자 격차 7배 (자동차 4.4% vs 시멘트 0.6%) — 시멘트·철강은 CCUS 없이 목표 달성 불가",
        footer=FOOTER,
    ))

    # ================================================================
    # 6. 한국 기업 — Composer(grid_2x2): 아이콘 카드 4개
    # ================================================================
    slide6 = add_slide(prs)
    comp6 = SlideComposer(slide6)
    comp6.header(SlideHeader(
        title="한국 4대 제조기업이 넷제로 레이스에서 각기 다른 기술 경로를 선택하여 추진 중임",
        category="6. 한국 현황 — K-제조업 넷제로",
        nav_path=["1. 진단", "4. 한국 분석"],
    ))
    comp6.intro("POSCO(수소 DRI)·현대차(EV)·삼성SDI(배터리 재활용)·SK(그린 투자) — 4사 전략 비교")

    zones6 = comp6.layout("grid_2x2")
    companies = [
        ("database", "POSCO — 수소 철강", "HyREX 수소환원 기술 세계 최초 도전\n2026 파일럿 30만톤, 2030 상업 100만톤\n그룹 투자 KRW 121조원 ($930억)", "dark"),
        ("rocket", "현대차 — EV 전환", "2024 순수 EV 30만대 (+20% YoY)\nSK온 합작 35GWh ($50억)\n전 수명주기 탄소관리 전략", "accent"),
        ("refresh", "삼성SDI — 배터리 순환", "사용 종료 배터리 ESS 재활용\n리튬·니켈·코발트 재추출\n현대차 협력 순환경제 구축", "subtle"),
        ("growth", "SK — 그린 투자 포트폴리오", "배터리 재활용 30GWh (2025)\nKRW 3,000억 매출 목표\n특허 50건, 그린수소 R&D", "light"),
    ]
    for i, (zone_key, zone) in enumerate(zones6.items()):
        icon, title, body, style = companies[i]
        comp_icon_card(comp6.canvas, icon=icon, title=title, body=body, style=style, region=zone)

    comp6.takeaway("POSCO 수소 DRI이 가장 대담한 도전 — 2026 파일럿 성공 시 글로벌 철강업 판도 변화의 시작점")
    comp6.footer(FOOTER)

    # ================================================================
    # 7. 탄소 가격 시나리오 — Waterfall (마진 충격 분해)
    # ================================================================
    waterfall_bridge(add_slide(prs), WaterfallSpec(
        header=SlideHeader(
            title="탄소 가격 $100/톤 시 철강 마진 10%p 하락 — 선제 투자 없이는 경쟁력 상실이 불가피함",
            category="7. 재무 영향 — 탄소 가격 시나리오",
            nav_path=["2. 전략", "4. 재무 분석"],
        ),
        intro="현재 마진 기준에서 탄소 가격 단계별 상승 시 철강 기업의 EBITDA 마진 변화를 분해",
        start={"label": "현재 마진", "value": 18},
        steps=[
            {"label": "탄소비용\n$50/톤", "value": -3, "detail": "직접 비용↑"},
            {"label": "에너지\n가격 전가", "value": -2, "detail": "전기료 상승"},
            {"label": "CBAM\n인증서", "value": -3, "detail": "EU 수출분"},
            {"label": "원자재\n가격 전가", "value": -2, "detail": "스크랩가↑"},
            {"label": "그린 프리미엄\n회수", "value": 2, "detail": "5~30% 프리미엄"},
        ],
        end={"label": "조정 마진", "value": 10},
        unit="%p",
        takeaway="마진 18%p→10%p (44% 하락) — 그린 프리미엄(+2%p)으로 일부 회복 가능하나 선제 감축 투자가 유일한 해법",
        footer=FOOTER,
    ))

    # ================================================================
    # 8. Scope 3 공급망 — Data Narrative
    # ================================================================
    data_narrative(add_slide(prs), DataNarrativeSpec(
        header=SlideHeader(
            title="공급망 Scope 3이 자체 배출의 26배 — 공급업체 67%+ 참여 없이 넷제로를 달성할 수 없음",
            category="8. 공급망 도전 — Scope 3 분석",
            nav_path=["2. 전략", "5. 공급망"],
        ),
        intro="CDP+BCG 2024 분석: 기업 총 GHG의 75%가 Scope 3 — 공급업체 협력이 넷제로의 핵심 전제 조건",
        chart_title="산업별 Scope 3 비중 (%)",
        chart_data=[
            {"label": "소매·소비재", "value": 95, "highlight": True},
            {"label": "자동차", "value": 85, "highlight": True},
            {"label": "전자·IT", "value": 80},
            {"label": "화학", "value": 75},
            {"label": "철강", "value": 65},
            {"label": "시멘트", "value": 50},
        ],
        chart_unit="%",
        narratives=[
            {"title": "Scope 3 = 직접 배출의 26배", "detail": "대부분 기업의 총 GHG 중 75%+가 공급망에서 발생. 자체 감축만으로는 목표 달성 불가. 공급업체 40%만 기후 협력 중."},
            {"title": "SBTi 67% 참여 의무", "detail": "Scope 3이 40% 초과 시 공급업체 67%+ 참여 목표 필수. 2026년까지 참여, 2028년까지 자체 목표 설정·공시 요구."},
            {"title": "협력 기업의 7배 효과", "detail": "공급업체와 기후 협력하는 기업이 Scope 3 목표 설정 확률 7배 높음. 그린 조달 기준 선제 도입이 경쟁 우위."},
        ],
        takeaway="공급업체 67%+ 참여 의무화(SBTi) — 2026년까지 참여 체계 구축, 2028년까지 공급업체 자체 목표 설정이 필수",
        footer=FOOTER,
    ))

    # ================================================================
    # 9. 로드맵 — Chevron Timeline (2025→2030)
    # ================================================================
    chevron_timeline(add_slide(prs), ChevronTimelineSpec(
        header=SlideHeader(
            title="2025~2027년이 결정적 분기점 — 기술 파일럿 성공이 2030 달성 가능성을 좌우함",
            category="9. 로드맵 — 2030 넷제로 마일스톤",
            nav_path=["3. 실행", "1. 로드맵"],
        ),
        intro="2025(즉시 착수) → 2027(기술 검증) → 2030(대규모 전환) 3단계 로드맵 — 각 시점의 Go/No-Go 기준 명시",
        phases=[
            {"year": "2025", "title": "측정·파일럿 착수",
             "detail": "내부 탄소 가격 $50~80/톤 도입\nScope 1·2 측정·공시 체계 구축\n수소·CCUS 파일럿 FID\n공급업체 30%+ 참여 착수",
             "position": "top"},
            {"year": "2026", "title": "EU CBAM 시행",
             "detail": "CBAM 인증서 구매 의무\n€100/톤 벌금 적용 시작\nK-ETS Phase 4 전환\nScope 3 공시 준비",
             "position": "bottom"},
            {"year": "2027", "title": "기술 검증 분기점",
             "detail": "그린수소 $2.50/kg 달성 여부\nPOSCO HyREX 파일럿 성공\nCCUS 상업 FID 결정\n대규모 투자 Go/No-Go",
             "position": "top"},
            {"year": "2028", "title": "상업화 확대",
             "detail": "EAF 대규모 전환 착수\nCCUS 연 100만톤 포집\n공급업체 67% 목표 달성\nSBTi 인증 확보",
             "position": "bottom"},
            {"year": "2030", "title": "넷제로 중간 목표",
             "detail": "산업 CO₂ 35% 감축\n재생에너지 3배 확대\n수소 철강 100만톤/년\n2050 넷제로 경로 확정",
             "position": "top"},
        ],
        takeaway="2027년이 최대 분기점 — 수소 $2.50/kg + POSCO HyREX 성공 시 대규모 전환 착수, 실패 시 경로 재설정 필요",
        footer=FOOTER,
    ))

    # ================================================================
    # 10. 3대 전략 축 — Composer(three_column): styled_card
    # ================================================================
    slide10 = add_slide(prs)
    comp10 = SlideComposer(slide10)
    comp10.header(SlideHeader(
        title="포트폴리오 전환·기술 리더십·규제 선제 대응 3축 동시 추진으로 넷제로 경쟁력을 확보함",
        category="10. 전략 제언 — Three Pillars",
        nav_path=["3. 실행", "2. 전략 축"],
    ))
    comp10.intro("한국 대형 제조 복합기업이 2030 넷제로를 달성하기 위한 3대 전략 축과 구체적 실행 과제")

    zones10 = comp10.layout("three_column")
    pillars = [
        {"title": "포트폴리오 전환", "body": "탄소 집약 사업 축소 + 저탄소 신사업 확대\n\n▪ 고로→EAF 전환 로드맵\n▪ 그린수소 DRI 투자 확대\n▪ 배터리 재활용 사업 신설\n▪ 탄소 크레딧 거래 역량 구축",
         "style": "dark", "number": "01"},
        {"title": "기술 리더십", "body": "선제 R&D 투자로 경쟁 우위 확보\n\n▪ 수소·CCUS·전기화 파일럿 FID\n▪ 디지털 트윈 에너지 최적화\n▪ 순환경제 재활용 기술 개발\n▪ 그린 프리미엄 제품 출시",
         "style": "accent", "number": "02"},
        {"title": "규제 선제 대응", "body": "K-ETS·CBAM·Scope 3 선제 준수\n\n▪ 내부 탄소 가격 $80/톤 도입\n▪ Scope 1·2·3 통합 공시 체계\n▪ 공급업체 67%+ 참여 프로그램\n▪ SBTi 인증 취득 (2026 목표)",
         "style": "dark", "number": "03"},
    ]
    for i, (zone_key, zone) in enumerate(zones10.items()):
        p = pillars[i]
        comp_styled_card(comp10.canvas, title=p["title"], body=p["body"],
                         style=p["style"], number=p["number"], region=zone)

    comp10.takeaway("3축 동시 추진 — 포트폴리오 전환(장기) + 기술 리더십(중기) + 규제 선제 대응(단기)의 시간 축 조합이 핵심")
    comp10.footer(FOOTER)

    return prs


def main():
    out_dir = Path("output/netzero_deck")
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / "제조업_넷제로_전환_10장.pptx"

    print("=" * 70)
    print("Step 4: GENERATE — 제조업 넷제로 전환 전략 10장")
    print("=" * 70)

    try:
        prs = build_deck()
        prs.save(out)
        print(f"\n생성 완료: {out}")
    except Exception as e:
        print(f"\nBUILD FAILED: {e}")
        import traceback
        traceback.print_exc()
        return 1

    # Step 5: EVALUATE
    print("\n" + "=" * 70)
    print("Step 5: EVALUATE")
    print("=" * 70)

    visual = validate_visual(out, convert_pdf=False)
    print(f"\nVisual issues: {len(visual.issues)}")
    for i in visual.issues[:10]:
        print(f"  - {i}")

    design = inspect_design(str(out))
    print(f"\nDesign: passed={design.passed}, issues={len(design.issues)}")
    for iss in design.issues:
        print(f"  - {iss}")
    for k, v in sorted(design.metrics.items()):
        if 'density' in k:
            print(f"  {k}: {v}")

    # PDF
    print("\nPDF 변환 중...")
    try:
        visual = validate_visual(out, convert_pdf=True)
        print(f"PDF: {'OK' if visual.pdf_available else 'SKIP'}")
    except Exception as e:
        print(f"PDF: FAIL ({e})")

    # Step 6 판단
    if design.passed and not visual.issues:
        print("\n✅ Step 5 PASS — Step 6 REFINE 불필요")
    else:
        print(f"\n❌ Step 5 FAIL ({len(design.issues)} issues) — Step 6 REFINE 필요")
        high_issues = [i for i in design.issues if '[HIGH' in i]
        if high_issues:
            print(f"  HIGH 이슈 {len(high_issues)}건:")
            for h in high_issues:
                print(f"    {h}")

    return 0


if __name__ == "__main__":
    sys.exit(main())
