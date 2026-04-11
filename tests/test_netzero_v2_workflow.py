"""넷제로 제조업 전환 10장 덱 — 워크플로우 6단계 준수.

Step 1: ANALYZE — docs/research/netzero_manufacturing_2030_research.md 기반
Step 2: PLAN — 10장 슬라이드 계획 + 덱 리듬 규칙 적용
Step 3: SELECT — slide_designer.md 매칭 테이블 기반 레이아웃+컴포넌트 선택
Step 4: GENERATE — SlideComposer + comp_xxx 코드 생성 (이 파일)
Step 5: EVALUATE — evaluate_pptx() 실행
Step 6: REFINE — 이슈 확인 후 수정

Composition Selection:
  S1:  표지 (커버)
  S2:  l_layout — comp_kpi_card(Hero 25%) + comp_native_chart(바) + comp_bullet_list
  S3:  full — comp_comparison_grid (EU/K-ETS/China 3열 비교)
  S4:  full — comp_pyramid (5대 기술 성숙도 피라미드)
  S5:  full — comp_waterfall (투자 갭 Bridge)
  S6:  full — comp_heatmap_grid (섹터별 진도 히트맵)
  S7:  center_peripheral_4 — comp_hub_spoke_diagram(한국 제조) + comp_bullet_list(4사)
  S8:  t_layout — comp_chevron_flow(상) + comp_kpi_row(하좌) + comp_bullet_list(하우)
  S9:  two_column — comp_heatmap_grid(좌) + comp_bullet_list(우)
  S10: full — comp_hero_block (결론)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region
from ppt_builder.composer import SlideComposer
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_kpi_card, comp_kpi_row, comp_native_chart, comp_bullet_list,
    comp_comparison_grid, comp_pyramid, comp_waterfall, comp_heatmap_grid,
    comp_hub_spoke_diagram, comp_chevron_flow, comp_hero_block,
    comp_styled_card,
)

OUTPUT = Path(__file__).parent.parent / "output" / "netzero_v2_workflow.pptx"


def make(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def s01_cover(prs):
    """S1: 표지."""
    s = make(prs)
    c = Canvas(s)
    c.box(x=0, y=0, w=10, h=7.5, fill="grey_900", border=None)
    c.text("글로벌 제조업\n넷제로 전환 전략 2030",
           x=0.5, y=2.0, w=9.0, h=2.5,
           size=28, bold=True, color="white", anchor="middle")
    c.text("PwC Consulting Korea", x=0.5, y=4.8, w=9.0, h=0.4,
           size=12, color="grey_400", anchor="top")
    c.text("2026년 4월", x=0.5, y=5.3, w=9.0, h=0.3,
           size=10, color="grey_400", anchor="top")


def s02_status(prs):
    """S2: 현황 — 제조업은 글로벌 탄소 배출의 25%를 차지한다."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="제조업은 글로벌 탄소 배출의 25%를 차지하며, 전환이 시급하다",
        category="Current State"))
    zones = comp.layout("l_layout", left_ratio=0.30, top_ratio=0.55)

    comp_kpi_card(comp.canvas, value="25%", label="글로벌 탄소 배출 중 제조업 비중",
                  detail="~9 Gt CO₂ (2022, IEA)", trend="down", region=zones["left_full"])

    comp_native_chart(comp.canvas, chart_type="vertical_bar",
                      chart_kwargs={
                          "categories": ["철강", "시멘트", "화학", "자동차"],
                          "values": [3.3, 2.5, 1.0, 0.5],
                          "highlight_idx": 0,
                          "series_name": "연간 배출량 (Gt CO₂)",
                      }, region=zones["right_top"])

    comp_bullet_list(comp.canvas, title="핵심 현황",
                     items=[
                         "에너지집약 3대 섹터(철강·시멘트·화학)가 산업 배출의 66%",
                         "1.5°C 달성을 위해 2030년까지 연평균 4% 감축 필요",
                         "프로세스 배출(화학반응) ~2 Gt — 기술 혁신 없이 감축 불가",
                     ], region=zones["right_bottom"])

    comp.takeaway("2030년까지 직접 산업 배출 41% 감축이 필요하나, 현재 감축 속도는 연 0.7% 증가세")
    comp.footer(SlideFooter(source="IEA Net Zero Roadmap 2023, UNECE 2022"))


def s03_regulation(prs):
    """S3: 규제 압력 — CBAM/K-ETS/China ETS 동시 강화."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="탄소 국경세·배출권거래제·공시 의무가 동시에 강화되고 있다",
        category="Regulation"))
    zones = comp.layout("full")

    comp_comparison_grid(comp.canvas,
                         columns=[
                             {"name": "EU CBAM", "summary": "2026.1 시행",
                              "criteria": ["€80/톤", "EU 배출 40%", "시멘트·철강·알루미늄", "톤당 €100 벌금"]},
                             {"name": "K-ETS", "summary": "Phase 4 진행", "highlight": True,
                              "criteria": ["₩8,684/톤", "국가 배출 80%", "6종 온실가스", "2026~2030 신규 할당"]},
                             {"name": "China ETS", "summary": "2024 확대",
                              "criteria": ["¥104.5/톤", "국가 배출 60%", "철강·시멘트 추가", "2027년 절대상한"]},
                         ],
                         row_labels=["탄소 가격", "커버리지", "대상 확대", "주요 변화"],
                         region=zones["main"])

    comp.takeaway("K-ETS Phase 4 + EU CBAM 완전 시행 = 한국 수출 제조업에 이중 부담")
    comp.footer(SlideFooter(source="ICAP 2024, EU Commission, SEC 2024"))


def s04_technology(prs):
    """S4: 기술 지형 — 5대 탈탄소 기술 성숙도."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="5대 탈탄소 기술이 상업화 임계점에 접근하고 있다",
        category="Technology"))
    zones = comp.layout("two_column", split=0.45)

    comp_pyramid(comp.canvas, layers=[
        {"title": "전기화 (EAF·열펌프)", "detail": "가장 성숙, CO₂ 70~80% 저감"},
        {"title": "순환경제 (재활용)", "detail": "알루미늄 95% 에너지 절감"},
        {"title": "그린수소 (DRI)", "detail": "2030년 $2.50/kg 목표"},
        {"title": "CCUS (탄소포집)", "detail": "시멘트 50~90% 포집 가능"},
        {"title": "디지털 트윈", "detail": "에너지 15~30% 절감"},
    ], region=zones["left"])

    comp_bullet_list(comp.canvas, title="기술별 핵심 데이터",
                     items=[
                         "전기아크로(EAF): BF-BOF 대비 CO₂ 70~80% 저감",
                         "그린수소: 2030년 8개 시장에서 그레이수소 대비 경쟁력 달성",
                         "CCUS: 고순도 공정 $15~25/톤, 시멘트 $38~86/톤",
                         "산업용 열펌프: 보일러 대비 3~5배 에너지 효율",
                         "디지털 트윈: Foxconn 연 30%+ 에너지 절감",
                     ], region=zones["right"])

    comp.takeaway("전기화와 순환경제는 즉시 적용 가능, 수소·CCUS는 2030년 상업화 임계점")
    comp.footer(SlideFooter(source="BNEF 2024, IEA CCUS 2024, McKinsey 2024"))


def s05_investment(prs):
    """S5: 투자 현황 — 2.1조 달러 vs 필요 4.5조 달러."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="연간 $2.1조 투자 중이나 넷제로 달성에는 $4.5조가 필요하다",
        category="Investment"))
    zones = comp.layout("full")

    comp_waterfall(comp.canvas,
                   start={"label": "현재 투자\n($2.1T)", "value": 2100},
                   steps=[
                       {"label": "재생에너지\n확대", "value": 800},
                       {"label": "전력망\n투자", "value": 400},
                       {"label": "산업 전환\n투자", "value": 600},
                       {"label": "운송 전환\n투자", "value": 500},
                   ],
                   end={"label": "필요 투자\n($4.5T)", "value": 4500},
                   unit="B",
                   region=zones["main"])

    comp.takeaway("연간 $2.4조의 투자 갭 — 제조업 전환 투자가 전체의 7%에 불과")
    comp.footer(SlideFooter(source="BNEF 2025, IEA Net Zero 2023"))


def s06_sector(prs):
    """S6: 섹터 비교 — 철강/시멘트/화학/자동차 진도 히트맵."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="철강은 투자를 선도하나, 시멘트의 탈탄소화가 가장 뒤처져 있다",
        category="Sector Analysis"))
    zones = comp.layout("full")

    comp_heatmap_grid(comp.canvas,
                      row_labels=["철강", "시멘트", "화학", "자동차"],
                      col_labels=["R&D 투자", "기술 성숙", "2030 목표", "현재 진도"],
                      values=[
                          [0.3, 0.4, 0.5, 0.5],   # 철강: 중간
                          [0.9, 0.7, 0.8, 0.9],   # 시멘트: 매우 뒤처짐
                          [0.4, 0.5, 0.4, 0.6],   # 화학: 중간
                          [0.1, 0.2, 0.2, 0.1],   # 자동차: 선도
                      ],
                      cell_texts=[
                          ["1.3%", "중간", "가변적", "중간"],
                          ["0.6%", "초기", "-25%", "최저"],
                          ["높음", "중간", "-35%", "중하"],
                          ["4.4%", "성숙", "EV가속", "최상"],
                      ],
                      region=zones["main"])

    comp.takeaway("시멘트 R&D 투자(0.6%)는 자동차(4.4%)의 1/7 — 정책 지원 없이는 전환 불가")
    comp.footer(SlideFooter(source="WEF Net-Zero Industry Tracker 2024"))


def s07_korea(prs):
    """S7: 한국 현황 — 4대 기업 넷제로."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="한국 제조 4사가 각자의 경로로 넷제로 레이스에 진입했다",
        category="Korea"))
    zones = comp.layout("center_peripheral_4", center_ratio=0.38)

    comp_hub_spoke_diagram(comp.canvas,
                           center="K-제조\n넷제로",
                           center_sub="2050 목표",
                           spokes=[
                               {"title": "POSCO", "detail": "HyREX\n수소DRI"},
                               {"title": "현대차", "detail": "EV 30만대\n배터리 JV"},
                               {"title": "삼성SDI", "detail": "배터리\n재활용"},
                               {"title": "HD현대", "detail": "친환경\n선박"},
                           ],
                           region=zones["center"])

    sides = {
        "left": ("POSCO", ["HyREX 2026 파일럿 30만톤", "2030 상업화 100만톤", "투자 ₩121조"]),
        "right": ("현대차", ["EV 30만대 (YoY +20%)", "SK온 JV $50억", "V2X 전략"]),
        "top": ("삼성SDI", ["배터리 재활용", "ESS 재활용"]),
        "bottom": ("HD현대", ["친환경 선박", "암모니아 추진"]),
    }
    for pos, (title, items) in sides.items():
        comp_bullet_list(comp.canvas, title=title, items=items, region=zones[pos])

    comp.footer(SlideFooter(source="GMK Center, Hyundai Newsroom 2024"))


def s08_roadmap(prs):
    """S8: 전환 로드맵 — 4단계."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="넷제로 전환은 진단→계획→실행→최적화 4단계로 추진해야 한다",
        category="Roadmap"))
    zones = comp.layout("t_layout", top_ratio=0.22, right_ratio=0.45)

    comp_chevron_flow(comp.canvas, phases=[
        {"tag": "D", "label": "진단"},
        {"tag": "P", "label": "계획"},
        {"tag": "E", "label": "실행"},
        {"tag": "O", "label": "최적화"},
    ], style="gradient", region=zones["top"])

    comp_kpi_row(comp.canvas, kpis=[
        {"value": "41%", "label": "2030 감축 목표", "trend": "down"},
        {"value": "$2.4T", "label": "연간 투자 갭", "trend": "flat"},
    ], region=zones["bottom_left"])

    comp_bullet_list(comp.canvas, title="단계별 핵심 활동",
                     items=[
                         "[진단] Scope 1·2·3 배출 인벤토리 구축",
                         "[계획] SBTi 기반 과학적 감축 경로 설정",
                         "[실행] EAF·열펌프·CCUS 파일럿 투자",
                         "[최적화] 디지털 트윈 기반 실시간 모니터링",
                     ], region=zones["bottom_right"])

    comp.takeaway("진단 단계의 배출 인벤토리 정확도가 전체 전환 성공의 기초")
    comp.footer(SlideFooter(source="SBTi Framework, PwC Analysis"))


def s09_risk(prs):
    """S9: 리스크 관리."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="기술 성숙도와 탄소 가격 불확실성이 가장 큰 전환 리스크다",
        category="Risk"))
    zones = comp.layout("two_column", split=0.55)

    comp_heatmap_grid(comp.canvas,
                      row_labels=["기술 성숙도", "탄소 가격", "규제 변동", "공급망", "인력 부족"],
                      col_labels=["영향도", "발생확률"],
                      values=[
                          [0.9, 0.7],
                          [0.8, 0.8],
                          [0.5, 0.6],
                          [0.6, 0.5],
                          [0.4, 0.7],
                      ],
                      cell_texts=[
                          ["Critical", "High"],
                          ["Critical", "Critical"],
                          ["Medium", "High"],
                          ["High", "Medium"],
                          ["Medium", "High"],
                      ],
                      region=zones["left"])

    comp_bullet_list(comp.canvas, title="리스크 대응 전략",
                     items=[
                         "[기술] 다중 기술 포트폴리오 — 단일 기술 의존 회피",
                         "[가격] 내부 탄소 가격제 도입 — 투자 의사결정 가격 신호",
                         "[규제] 규제 시나리오 분석 — 3단계(낙관/기본/비관)",
                         "[공급망] 저탄소 소재 사전 계약 — 장기 공급 확보",
                         "[인력] 그린 스킬 교육 프로그램 — 기존 인력 전환",
                     ], region=zones["right"])

    comp.takeaway("내부 탄소 가격제(ICP) 도입이 리스크 관리의 첫 번째 조치")
    comp.footer(SlideFooter(source="PwC Climate Risk Assessment Framework"))


def s10_conclusion(prs):
    """S10: 결론."""
    s = make(prs)
    comp = SlideComposer(s)
    comp.header(SlideHeader(
        title="제조업 넷제로 전환은 선택이 아닌 생존의 문제다",
        category="Conclusion"))
    zones = comp.layout("full")

    comp_hero_block(comp.canvas,
                    label="ACTION REQUIRED",
                    headline="2030년까지 41% 감축,\n연간 $4.5조 투자가 필요하다",
                    sub_points=[
                        "전기화·순환경제는 지금 즉시 시작 가능한 Quick Win",
                        "수소·CCUS는 2026~2028 파일럿 투자 결정이 임계점",
                        "내부 탄소 가격제 + Scope 3 인벤토리 = 전환의 기초",
                        "K-ETS Phase 4 + EU CBAM = 한국 수출 제조업 이중 부담",
                        "선도 기업(ArcelorMittal, Heidelberg, POSCO)은 이미 레이스 중",
                    ],
                    bg_color="grey_800",
                    region=zones["main"])

    comp.takeaway("지금 시작하지 않으면, 2030년 규제 장벽에서 시장 접근이 차단된다")
    comp.footer(SlideFooter(source="PwC Net Zero Manufacturing Practice"))


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    s01_cover(prs)
    s02_status(prs)
    s03_regulation(prs)
    s04_technology(prs)
    s05_investment(prs)
    s06_sector(prs)
    s07_korea(prs)
    s08_roadmap(prs)
    s09_risk(prs)
    s10_conclusion(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")

    # Step 5: EVALUATE
    from ppt_builder.evaluate import evaluate_pptx, print_report
    report = evaluate_pptx(str(OUTPUT))
    print_report(report)


if __name__ == "__main__":
    main()
