"""PwC Tier-1A 컴포넌트 5종 — 독립 함수로 구현 후 PPT 렌더링 + 평가."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from ppt_builder.primitives import Canvas, Region

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "pwc_tier1a"


# ============================================================
# Color tokens (PwC palette shortcuts)
# ============================================================
ORANGE_GRADIENT = [
    "white",        # cell 0 — 밝음
    "#FFE0C2",      # cell 1 — 연한 주황
    "accent_light",  # cell 2 — 중간 밝은 주황
    "accent_mid",   # cell 3 — 중간 주황
    "accent",       # cell 4 — 진한 주황
    "grey_900",     # cell 5 — 가장 어두운
]

# 텍스트 색상: 밝은 배경→어두운 글씨, 어두운 배경→흰 글씨
ORANGE_TEXT = [
    "grey_900",  # cell 0
    "grey_900",  # cell 1
    "grey_900",  # cell 2
    "white",     # cell 3
    "white",     # cell 4
    "white",     # cell 5
]


# ============================================================
# 슬라이드 헤더 유틸
# ============================================================
def _slide_header(prs, title: str) -> Canvas:
    """빈 슬라이드 + 상단 제목 렌더, Canvas 반환."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    c.text("PwC Tier-1A Component", x=0.3, y=0.15, w=9.4, h=0.25,
           size=9, color="grey_700", anchor="top")
    c.text(title, x=0.3, y=0.38, w=9.4, h=0.40,
           size=16, bold=True, color="grey_900", anchor="top")
    c.line(x1=0.3, y1=0.85, x2=9.7, y2=0.85, color="accent", width=2.0)
    return c


# ============================================================
# 1. pwc_process_grid_6cell — 3x2 그리드, 오렌지 그라데이션
# ============================================================
def pwc_process_grid_6cell(c: Canvas, region: Region) -> float:
    """3x2 프로세스 그리드. 셀별 오렌지 그라데이션 + 번호/헤더/본문."""
    cells = [
        {"num": "01", "header": "현황 진단",
         "body": "AS-IS 업무 프로세스 분석\n핵심 Pain Point 도출\n데이터 성숙도 평가"},
        {"num": "02", "header": "전략 수립",
         "body": "디지털 전환 로드맵 설계\n우선순위 과제 선정\nROI 시뮬레이션"},
        {"num": "03", "header": "아키텍처 설계",
         "body": "TO-BE 시스템 구조 정의\n클라우드 마이그레이션 계획\n데이터 통합 방안"},
        {"num": "04", "header": "파일럿 실행",
         "body": "MVP 범위 확정 및 개발\nPoC 검증 (4~6주)\n사용자 피드백 반영"},
        {"num": "05", "header": "전사 확산",
         "body": "단계적 Roll-out 계획\n변화관리 프로그램 운영\nKPI 모니터링 체계"},
        {"num": "06", "header": "운영 최적화",
         "body": "성과 측정 및 리포팅\n지속 개선 사이클 구축\nAI/자동화 고도화"},
    ]

    cols, rows = 3, 2
    gap = 0.15
    cell_w = (region.w - gap * (cols - 1)) / cols
    cell_h = (region.h - gap * (rows - 1)) / rows

    for idx, cell in enumerate(cells):
        col = idx % cols
        row = idx // cols
        cx = col * (cell_w + gap)
        cy = row * (cell_h + gap)

        fill = ORANGE_GRADIENT[idx]
        txt_color = ORANGE_TEXT[idx]

        # 셀 박스
        c.box(x=cx, y=cy, w=cell_w, h=cell_h,
              fill=fill, border=None, region=region)

        # 큰 번호
        c.text(cell["num"], x=cx + 0.15, y=cy + 0.12, w=1.0, h=0.45,
               size=28, bold=True, color=txt_color, font="Georgia",
               anchor="top", region=region)

        # 헤더
        c.text(cell["header"], x=cx + 0.15, y=cy + 0.60, w=cell_w - 0.30, h=0.25,
               size=10, bold=True, color=txt_color, anchor="top", region=region)

        # 본문 (줄바꿈)
        c.text(cell["body"], x=cx + 0.15, y=cy + 0.88, w=cell_w - 0.30, h=cell_h - 1.05,
               size=8, color=txt_color, anchor="top", region=region)

    return region.h


# ============================================================
# 2. pwc_timeline_bar — 수평 바 타임라인 (5개년)
# ============================================================
def pwc_timeline_bar(c: Canvas, region: Region) -> float:
    """수평 바 타임라인. 5년 세그먼트 + 상/하 텍스트 블록."""
    years = ["2024", "2025", "2026", "2027", "2028"]
    bar_colors = ["grey_400", "grey_700", "accent_mid", "accent", "grey_900"]

    # 바 위/아래 텍스트
    above = [
        {"header": "인프라 현대화", "body": "레거시 시스템 분석\n클라우드 전환 계획 수립"},
        {"header": "플랫폼 구축", "body": "데이터 레이크하우스 구축\nAPI Gateway 표준화"},
    ]
    below = [
        {"header": "AI/ML 도입", "body": "수요예측 모델 파일럿\n품질검사 자동화 PoC"},
        {"header": "전사 확산", "body": "전 사업부 Roll-out\n변화관리 프로그램"},
        {"header": "자율 운영", "body": "AIOps 기반 운영\n지속적 최적화 사이클"},
    ]

    r = region
    bar_y = 2.0  # 바의 y 위치 (region 내 상대)
    bar_h = 0.45
    seg_w = r.w / len(years)

    # --- 바 세그먼트 + 연도 라벨 ---
    for i, (year, bc) in enumerate(zip(years, bar_colors)):
        sx = i * seg_w
        c.box(x=sx, y=bar_y, w=seg_w + 0.01, h=bar_h,
              fill=bc, border=None, region=r)
        c.text(year, x=sx, y=bar_y, w=seg_w, h=bar_h,
               size=10, bold=True, color="white", align="center",
               anchor="middle", region=r)

    # --- 위 텍스트 블록 (2개) ---
    above_positions = [0.3, 4.2]
    block_w = 3.5
    for i, ab in enumerate(above):
        bx = above_positions[i]
        c.text(ab["header"], x=bx, y=0.0, w=block_w, h=0.28,
               size=10, bold=True, color="accent", anchor="top", region=r)
        c.text(ab["body"], x=bx, y=0.30, w=block_w, h=0.80,
               size=8, color="grey_900", anchor="top", region=r)
        # 연결선
        c.line(x1=bx + 0.3, y1=1.15, x2=bx + 0.3, y2=bar_y,
               color="grey_400", width=0.75, region=r)

    # --- 아래 텍스트 블록 (3개) ---
    below_positions = [0.3, 3.2, 6.3]
    for i, bl in enumerate(below):
        bx = below_positions[i]
        # 연결선
        c.line(x1=bx + 0.3, y1=bar_y + bar_h, x2=bx + 0.3, y2=bar_y + bar_h + 0.35,
               color="grey_400", width=0.75, region=r)
        c.text(bl["header"], x=bx, y=bar_y + bar_h + 0.40, w=2.8, h=0.28,
               size=10, bold=True, color="accent", anchor="top", region=r)
        c.text(bl["body"], x=bx, y=bar_y + bar_h + 0.70, w=2.8, h=0.80,
               size=8, color="grey_900", anchor="top", region=r)

    return region.h


# ============================================================
# 3. pwc_timeline_zigzag — 지그재그 5단계
# ============================================================
def pwc_timeline_zigzag(c: Canvas, region: Region) -> float:
    """5단계 지그재그 타임라인. 홀수=아래, 짝수=위."""
    steps = [
        {"num": "01", "header": "문제 정의",
         "body": "비즈니스 이슈 구조화\n가설 수립\n분석 프레임 설계"},
        {"num": "02", "header": "데이터 수집",
         "body": "내부 데이터 추출\n외부 벤치마크 확보\n인터뷰 실시"},
        {"num": "03", "header": "분석 및 인사이트",
         "body": "정량/정성 분석 병행\n패턴 도출 및 해석\nSo-What 정리"},
        {"num": "04", "header": "전략 도출",
         "body": "옵션 평가 (3안)\n실행 계획 수립\n이해관계자 합의"},
        {"num": "05", "header": "실행 지원",
         "body": "PMO 구성\n마일스톤 관리\n성과 KPI 트래킹"},
    ]

    r = region
    n = len(steps)
    step_w = 1.65
    total_w = step_w * n + 0.15 * (n - 1)
    x_start = (r.w - total_w) / 2  # 중앙 정렬

    y_top = 0.0      # 짝수(2,4) — 위
    y_bottom = 2.8    # 홀수(1,3,5) — 아래
    connector_y = 2.2  # 커넥터 수평 라인 y

    for i, step in enumerate(steps):
        sx = x_start + i * (step_w + 0.15)
        is_top = (i % 2 == 1)  # 0-indexed: 1,3 = 짝수 스텝(2,4)
        sy = y_top if is_top else y_bottom

        # 큰 오렌지 번호
        c.text(step["num"], x=sx, y=sy, w=step_w, h=0.45,
               size=24, bold=True, color="accent", font="Georgia",
               anchor="top", region=r)

        # 헤더
        c.text(step["header"], x=sx, y=sy + 0.48, w=step_w, h=0.28,
               size=10, bold=True, color="grey_900", anchor="top", region=r)

        # 본문
        c.text(step["body"], x=sx, y=sy + 0.78, w=step_w, h=1.0,
               size=8, color="grey_700", anchor="top", region=r)

        # 수직 연결선 (번호→커넥터 라인)
        if is_top:
            c.line(x1=sx + 0.3, y1=sy + 1.8, x2=sx + 0.3, y2=connector_y,
                   color="accent_mid", width=1.0, region=r)
        else:
            c.line(x1=sx + 0.3, y1=connector_y, x2=sx + 0.3, y2=sy,
                   color="accent_mid", width=1.0, region=r)

        # 커넥터 사각형 (오렌지 작은 사각)
        sq_size = 0.18
        c.box(x=sx + 0.3 - sq_size / 2, y=connector_y - sq_size / 2,
              w=sq_size, h=sq_size, fill="accent", border=None, region=r)

    # 수평 커넥터 라인 (전체)
    first_x = x_start + 0.3
    last_x = x_start + (n - 1) * (step_w + 0.15) + 0.3
    c.line(x1=first_x, y1=connector_y, x2=last_x, y2=connector_y,
           color="accent_mid", width=1.5, region=r)

    return region.h


# ============================================================
# 4. pwc_kpi_callout_box — 오렌지 배경 KPI 콜아웃
# ============================================================
def pwc_kpi_callout_box(
    c: Canvas,
    *,
    value: str,
    label: str,
    detail: str = "",
    fill: str = "accent",
    region: Region,
) -> float:
    """솔리드 오렌지 배경 + 큰 흰색 숫자 KPI 콜아웃 박스."""
    r = region

    # 배경 박스
    c.box(x=0, y=0, w=r.w, h=r.h, fill=fill, border=None, region=r)

    # 큰 숫자
    v_size = 28 if len(value) <= 5 else 22
    c.text(value, x=0.15, y=0.15, w=r.w - 0.30, h=0.55,
           size=v_size, bold=True, color="white", font="Georgia",
           align="center", anchor="top", region=r)

    # 라벨
    c.text(label, x=0.15, y=0.75, w=r.w - 0.30, h=0.30,
           size=9, bold=True, color="white", align="center",
           anchor="top", region=r)

    # 디테일
    if detail:
        c.text(detail, x=0.15, y=1.08, w=r.w - 0.30, h=0.25,
               size=7, color="white", align="center",
               anchor="top", region=r)

    return r.h


# ============================================================
# 5. pwc_donut_radial — 중심 원 + 3방향 텍스트 블록
# ============================================================
def pwc_donut_radial(c: Canvas, region: Region) -> float:
    """중심 원 + 3방향(120도 간격) 텍스트 블록."""
    r = region

    # 중심 원
    center_d = 1.6
    cx = (r.w - center_d) / 2
    cy = (r.h - center_d) / 2
    c.circle(x=cx, y=cy, d=center_d,
             fill="accent", border=None,
             text="Core\nStrategy", text_color="white", text_size=12,
             region=r)

    # 3개 텍스트 블록 (120도 간격 위치 근사)
    blocks = [
        {  # 상단 중앙
            "x": r.w / 2 - 1.5, "y": 0.0,
            "title": "디지털 혁신",
            "items": "AI/ML 기반 의사결정\n클라우드 네이티브 전환\n데이터 드리븐 경영",
        },
        {  # 좌하단
            "x": 0.0, "y": 3.2,
            "title": "운영 효율화",
            "items": "프로세스 자동화 (RPA)\n공급망 최적화\n비용 구조 혁신",
        },
        {  # 우하단
            "x": r.w - 3.2, "y": 3.2,
            "title": "고객 경험",
            "items": "옴니채널 통합\n개인화 서비스\n실시간 VOC 분석",
        },
    ]

    block_w = 3.0
    block_h = 1.6

    for blk in blocks:
        bx, by = blk["x"], blk["y"]

        # 제목
        c.text(blk["title"], x=bx, y=by, w=block_w, h=0.30,
               size=10, bold=True, color="accent", anchor="top", region=r)

        # 항목 (불릿)
        lines = blk["items"].split("\n")
        for li, line in enumerate(lines):
            c.text(f"  {line}", x=bx, y=by + 0.32 + li * 0.22, w=block_w, h=0.22,
                   size=8, color="grey_900", anchor="top", region=r)

        # 연결선 (블록 중심 → 원 중심 방향)
        bcx = bx + block_w / 2
        bcy = by + block_h / 2
        ocx = r.w / 2
        ocy = r.h / 2
        # 라인은 블록 가장자리에서 원 가장자리까지 (근사)
        dx = ocx - bcx
        dy = ocy - bcy
        dist = max((dx**2 + dy**2) ** 0.5, 0.01)
        # 원 반지름 = center_d / 2
        cr = center_d / 2 + 0.05
        # 블록에서 약간 안쪽
        line_sx = bcx + dx * 0.15
        line_sy = bcy + dy * 0.15
        # 원 가장자리
        line_ex = ocx - dx / dist * cr
        line_ey = ocy - dy / dist * cr

        c.line(x1=line_sx, y1=line_sy, x2=line_ex, y2=line_ey,
               color="accent_mid", width=1.0, region=r)

    return r.h


# ============================================================
# 빌드 + 평가
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Process Grid 6cell ---
    c1 = _slide_header(prs, "디지털 전환 6단계 프로세스 — 진단부터 최적화까지")
    r1 = Region(0.3, 1.1, 9.4, 5.8)
    pwc_process_grid_6cell(c1, r1)

    # --- Slide 2: Timeline Bar ---
    c2 = _slide_header(prs, "5개년 디지털 전환 로드맵 — 인프라에서 자율 운영까지")
    r2 = Region(0.3, 1.1, 9.4, 5.8)
    pwc_timeline_bar(c2, r2)

    # --- Slide 3: Timeline Zigzag ---
    c3 = _slide_header(prs, "문제 해결 5단계 방법론 — 구조화된 접근")
    r3 = Region(0.3, 1.1, 9.4, 5.8)
    pwc_timeline_zigzag(c3, r3)

    # --- Slide 4: KPI Callout Box (4개 가로 배열) ---
    c4 = _slide_header(prs, "핵심 성과 지표 — 전환 효과 정량화")
    r4 = Region(0.3, 1.5, 9.4, 5.0)
    kpis = [
        {"value": "47%", "label": "운영 비용 절감", "detail": "3년 누적 기준", "fill": "accent"},
        {"value": "3.2x", "label": "ROI 달성", "detail": "투자 대비 수익률", "fill": "accent_mid"},
        {"value": "92%", "label": "프로세스 자동화", "detail": "수작업 대비 감소", "fill": "grey_700"},
        {"value": "$180M", "label": "연간 절감액", "detail": "2028년 목표", "fill": "grey_900"},
    ]
    kpi_w = (r4.w - 0.15 * 3) / 4
    kpi_h = 1.5
    for i, kpi in enumerate(kpis):
        kr = Region(r4.x + i * (kpi_w + 0.15), r4.y, kpi_w, kpi_h)
        pwc_kpi_callout_box(c4, value=kpi["value"], label=kpi["label"],
                            detail=kpi["detail"], fill=kpi["fill"], region=kr)

    # 하단에 설명 텍스트
    c4.text("출처: PwC Digital Transformation Benchmark 2025  |  대상: 글로벌 500대 기업",
            x=0.3, y=6.5, w=9.4, h=0.3, size=7, color="grey_400", anchor="top")

    # --- Slide 5: Donut Radial ---
    c5 = _slide_header(prs, "전략 프레임워크 — 3대 핵심 축 통합 추진")
    r5 = Region(0.3, 1.1, 9.4, 5.8)
    pwc_donut_radial(c5, r5)

    # --- 저장 ---
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path} ({len(prs.slides)} slides)")

    # --- 평가 ---
    from ppt_builder.evaluate import evaluate_pptx
    report = evaluate_pptx(str(pptx_path))
    print(f"\n=== Evaluate Report ===")
    print(f"Score: {report['score']}/100  (pass={report['pass']})")
    print(f"Slides: {report['slide_count']}")
    if report["issues"]:
        print("Issues:")
        for iss in report["issues"]:
            print(f"  - {iss}")
    else:
        print("No issues found.")


if __name__ == "__main__":
    main()
