"""PwC-style Tier 1B 컴포넌트 4개 — waterfall, icon_text_3col, cover, data_table."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from ppt_builder.primitives import Canvas, Region, color, COLORS
from ppt_builder.evaluate import evaluate_pptx, print_report

OUTPUT = Path(__file__).parent.parent / "output" / "pwc_tier1b.pptx"

# PwC palette overrides
PWC_RED = RGBColor(0xD9, 0x3A, 0x2B)      # PwC red (logo, accents)
PWC_ORANGE = COLORS["accent"]              # FD5108 — existing accent (orange)
PWC_ORANGE_MID = COLORS["accent_mid"]      # FE7C39
PWC_GREY_900 = COLORS["grey_900"]          # near-black charcoal
PWC_GREY_800 = COLORS["grey_800"]          # dark grey
PWC_GREY_100 = COLORS["grey_100"]          # very light grey


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return slide, Canvas(slide)


# ============================================================
# 6. pwc_waterfall
# ============================================================

def pwc_waterfall(
    c: Canvas,
    *,
    start: dict,
    steps: list[dict],
    end: dict,
    unit: str = "",
    region: Region,
) -> float:
    """Waterfall/Bridge chart with PwC colors.

    Increases: accent (orange), Decreases: negative (red-ish),
    Start/End totals: grey_900 (dark charcoal). NOT green.
    Same API as comp_waterfall.
    """
    r = region
    all_items = [start] + steps + [end]
    n = len(all_items)
    if n == 0:
        return 0.0

    bar_gap = 0.06
    bar_w = (r.w - bar_gap * (n - 1)) / n

    # Cumulative values
    cumulative = [start["value"]]
    for s in steps:
        cumulative.append(cumulative[-1] + s["value"])
    all_values = cumulative + [end["value"]]
    max_val = max(abs(v) for v in all_values) if all_values else 1

    label_h = 0.30
    chart_h = r.h - label_h - 0.05
    baseline_y = chart_h * 0.72

    def val_to_y(v):
        if max_val == 0:
            return baseline_y
        return baseline_y - v * (chart_h * 0.60) / max_val

    # Baseline
    c.line(x1=0, y1=baseline_y, x2=r.w, y2=baseline_y,
           color="grey_mid", width=0.5, region=r)

    running = start["value"]
    for i, item in enumerate(all_items):
        bx = i * (bar_w + bar_gap)
        is_start = (i == 0)
        is_end = (i == n - 1)

        if is_start or is_end:
            # Totals: dark charcoal (PwC style)
            val = item["value"]
            bar_top = val_to_y(val)
            bh = abs(baseline_y - bar_top)
            c.box(x=bx, y=min(bar_top, baseline_y), w=bar_w, h=max(bh, 0.04),
                  fill="grey_900", border=None, region=r)
        else:
            step_val = item["value"]
            prev_cum = running
            running += step_val
            bar_top = val_to_y(max(prev_cum, running))
            bar_bot = val_to_y(min(prev_cum, running))
            # PwC: orange for increase, negative (red) for decrease
            bar_fill = "accent" if step_val > 0 else "negative"
            c.box(x=bx, y=bar_top, w=bar_w, h=max(abs(bar_bot - bar_top), 0.04),
                  fill=bar_fill, border=None, region=r)
            # Connector line
            if i > 0:
                prev_x = (i - 1) * (bar_w + bar_gap) + bar_w
                c.line(x1=prev_x, y1=val_to_y(prev_cum), x2=bx, y2=val_to_y(prev_cum),
                       color="grey_400", width=0.5, region=r)

        # Label
        c.text(item["label"], x=bx, y=chart_h + 0.03, w=bar_w, h=label_h,
               size=7, bold=True, color="grey_900", align="center", anchor="top", region=r)

        # Value annotation
        val = item["value"]
        prefix = "+" if (not is_start and not is_end) and val > 0 else ""
        val_str = f"{prefix}{val:,.0f}{unit}"
        vy = val_to_y(abs(val) if (is_start or is_end) else max(running, running - val)) - 0.22
        c.text(val_str, x=bx, y=vy, w=bar_w, h=0.20,
               size=8, bold=True, color="grey_900", align="center", anchor="top", region=r)

    return r.h


# ============================================================
# 7. pwc_icon_text_3col
# ============================================================

def pwc_icon_text_3col(
    c: Canvas,
    *,
    columns: list[dict],
    region: Region,
) -> float:
    """3 equal columns: orange circle with 2-letter abbreviation, header, body.

    columns: [{"abbr": "AI", "header": "...", "body": "..."}, ...]
    """
    r = region
    n = min(len(columns), 3)
    gap = 0.3
    col_w = (r.w - gap * (n - 1)) / n

    for i, col in enumerate(columns[:3]):
        cx = i * (col_w + gap)

        # Orange circle with abbreviation
        circle_d = 0.7
        circle_x = cx + (col_w - circle_d) / 2
        c.circle(
            x=circle_x, y=0.0, d=circle_d,
            fill="accent", border=None,
            text=col["abbr"], text_color="white", text_size=16, text_bold=True,
            region=r,
        )

        # Header (bold 10pt)
        c.text(
            col["header"],
            x=cx, y=circle_d + 0.25, w=col_w, h=0.35,
            size=10, bold=True, color="grey_900", align="center", anchor="top",
            region=r,
        )

        # Body text (8pt, grey)
        c.text(
            col["body"],
            x=cx + 0.1, y=circle_d + 0.65, w=col_w - 0.2, h=r.h - circle_d - 0.75,
            size=8, color="grey_700", align="center", anchor="top",
            region=r,
        )

    return r.h


# ============================================================
# 8. pwc_cover_slide
# ============================================================

def pwc_cover_slide(
    c: Canvas,
    *,
    title: str,
    author: str = "",
    date: str = "",
) -> None:
    """PwC-style cover slide — white bg, red 'pwc' logo, parallelogram accents.

    Renders directly on the slide (no region needed, full-slide layout).
    """
    # White background (blank layout already white)

    # Title — left-aligned, bold 24pt
    c.text(
        title,
        x=0.6, y=2.2, w=5.5, h=1.2,
        size=24, bold=True, color="grey_900", align="left", anchor="top",
    )

    # Author + date subtitle
    subtitle_parts = []
    if author:
        subtitle_parts.append(author)
    if date:
        subtitle_parts.append(date)
    if subtitle_parts:
        c.text(
            " | ".join(subtitle_parts),
            x=0.6, y=3.5, w=5.5, h=0.4,
            size=10, color="grey_700", align="left", anchor="top",
        )

    # Two overlapping orange parallelogram shapes — center-right decorative
    slide = c.slide
    para1 = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(5.8), Inches(1.5), Inches(4.0), Inches(3.5),
    )
    para1.fill.solid()
    para1.fill.fore_color.rgb = PWC_ORANGE
    para1.line.fill.background()
    para1.rotation = 0.0

    para2 = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        Inches(6.5), Inches(2.0), Inches(3.5), Inches(3.0),
    )
    para2.fill.solid()
    para2.fill.fore_color.rgb = PWC_ORANGE_MID
    para2.line.fill.background()
    para2.rotation = 0.0

    # "pwc" text — bottom-left, red, size 12
    c.text(
        "pwc",
        x=0.4, y=6.7, w=1.0, h=0.4,
        size=12, bold=True, color=PWC_RED, align="left", anchor="bottom",
    )

    # Thin accent line below title area
    c.line(x1=0.6, y1=3.95, x2=6.0, y2=3.95, color="grey_mid", width=0.5)


# ============================================================
# 9. pwc_data_table
# ============================================================

def pwc_data_table(
    c: Canvas,
    *,
    headers: list[str],
    rows: list[list[str]],
    region: Region,
    col_ratios: list[float] | None = None,
) -> float:
    """Full-width data table with PwC styling.

    Dark charcoal header (grey_900), white 10pt bold text.
    Alternating rows (grey_100 / white). Border 0.5pt grey_mid. Body 8pt.
    """
    r = region
    n_cols = len(headers)
    n_rows = len(rows)

    if col_ratios and len(col_ratios) == n_cols:
        total = sum(col_ratios)
        widths = [r.w * (cr / total) for cr in col_ratios]
    else:
        widths = [r.w / n_cols] * n_cols

    header_h = 0.35
    body_h = (r.h - header_h) / max(n_rows, 1)

    # Header row — dark charcoal fill, white bold text
    cx = 0.0
    for i, hdr in enumerate(headers):
        c.box(x=cx, y=0, w=widths[i], h=header_h,
              fill="grey_900", border=0.5, border_color="grey_mid", region=r)
        c.text(hdr, x=cx + 0.06, y=0, w=widths[i] - 0.12, h=header_h,
               size=10, bold=True, color="white", anchor="middle", region=r)
        cx += widths[i]

    # Body rows — alternating grey_100 / white
    for ri, row in enumerate(rows):
        ry = header_h + ri * body_h
        row_fill = "grey_100" if ri % 2 == 0 else "white"
        cx = 0.0
        for ci, cell in enumerate(row):
            c.box(x=cx, y=ry, w=widths[ci], h=body_h,
                  fill=row_fill, border=0.5, border_color="grey_mid", region=r)
            c.text(str(cell), x=cx + 0.06, y=ry, w=widths[ci] - 0.12, h=body_h,
                   size=8, color="grey_900", anchor="middle", region=r)
            cx += widths[ci]

    return r.h


# ============================================================
# Slide builders
# ============================================================

def slide_1_waterfall(prs):
    """PwC Waterfall — 비용 브릿지 분석."""
    slide, c = _new_slide(prs)
    c.title("디지털 전환 비용은 클라우드 마이그레이션이 주도한다", size=14)
    c.text("FY2025 IT 투자 브릿지 분석 (단위: 억원)", x=0.3, y=0.75, w=9.4, h=0.25,
           size=9, color="grey_700")

    pwc_waterfall(
        c,
        start={"label": "FY24\n실적", "value": 850},
        steps=[
            {"label": "클라우드\n마이그레이션", "value": 220},
            {"label": "ERP\n고도화", "value": 150},
            {"label": "AI/ML\n투자", "value": 95},
            {"label": "레거시\n절감", "value": -180},
            {"label": "인건비\n조정", "value": -55},
        ],
        end={"label": "FY25\n계획", "value": 1080},
        unit="",
        region=Region(0.3, 1.2, 9.4, 4.5),
    )

    # Bottom insight callout
    c.box(x=0.3, y=6.0, w=9.4, h=1.0, fill="grey_100", border=0.5, border_color="grey_mid")
    c.text(
        "핵심 시사점: 클라우드 마이그레이션(+220억)이 전체 증가분의 51%를 차지하며, "
        "레거시 절감(-180억)으로 순증가를 억제하는 투트랙 전략이 유효함. "
        "FY26부터 클라우드 OPEX 전환에 따른 CAPEX 감소 효과 기대.",
        x=0.5, y=6.1, w=9.0, h=0.8, size=8, color="grey_900",
    )


def slide_2_icon_text(prs):
    """PwC 3-Column Icon + Text — 전략 축."""
    slide, c = _new_slide(prs)
    c.title("디지털 전환 3대 전략 축: AI, 클라우드, 공급망 혁신", size=14)
    c.text("HD현대 그룹 2025-2027 중기 전략 프레임워크", x=0.3, y=0.75, w=9.4, h=0.25,
           size=9, color="grey_700")

    pwc_icon_text_3col(
        c,
        columns=[
            {
                "abbr": "AI",
                "header": "AI-Driven Operations",
                "body": (
                    "생산 공정 예측 모델 도입으로 불량률 40% 감소 목표. "
                    "자연어 기반 ERP 인터페이스로 현장 작업자 생산성 향상. "
                    "2026년까지 전 사업장 AI 플랫폼 표준화 완료."
                ),
            },
            {
                "abbr": "CL",
                "header": "Cloud Migration",
                "body": (
                    "온프레미스 SAP ECC에서 S/4HANA Cloud로 전환. "
                    "인프라 비용 35% 절감 + 글로벌 통합 데이터 레이크 구축. "
                    "Phase 1(국내) 완료, Phase 2(해외법인) 2026 착수."
                ),
            },
            {
                "abbr": "SC",
                "header": "Supply Chain Resilience",
                "body": (
                    "디지털 트윈 기반 공급망 시뮬레이션으로 리드타임 25% 단축. "
                    "Tier-2 협력사까지 실시간 가시성 확보. "
                    "지정학적 리스크 대응 다변화 시나리오 수립."
                ),
            },
        ],
        region=Region(0.3, 1.2, 9.4, 5.5),
    )

    # Bottom source line
    c.text("출처: PwC Advisory 내부 분석, 2025.04", x=0.3, y=7.0, w=9.4, h=0.25,
           size=7, color="grey_400", align="right")


def slide_3_cover(prs):
    """PwC Cover Slide."""
    slide, c = _new_slide(prs)
    pwc_cover_slide(
        c,
        title="HD현대 그룹\nDigital Transformation\nRoadmap 2025-2027",
        author="PwC Advisory | Digital Consulting",
        date="April 2026",
    )


def slide_4_data_table(prs):
    """PwC Data Table — 사업부별 디지털 성숙도."""
    slide, c = _new_slide(prs)
    c.title("사업부별 디지털 성숙도 평가: 조선 부문이 가장 앞서있다", size=14)
    c.text("2025년 1분기 디지털 성숙도 진단 결과 (5점 척도)", x=0.3, y=0.75, w=9.4, h=0.25,
           size=9, color="grey_700")

    pwc_data_table(
        c,
        headers=["사업부", "데이터 인프라", "프로세스 자동화", "AI 활용도"],
        rows=[
            ["HD현대중공업 (조선)", "4.2", "3.8", "3.5"],
            ["HD현대인프라코어 (건기)", "3.5", "3.2", "2.8"],
            ["HD현대일렉트릭 (전력)", "3.8", "3.5", "3.1"],
            ["HD한국조선해양 (지주)", "3.0", "2.8", "2.5"],
            ["HD현대오일뱅크 (에너지)", "3.6", "3.0", "2.7"],
        ],
        col_ratios=[2.5, 1.5, 1.5, 1.5],
        region=Region(0.3, 1.2, 9.4, 3.0),
    )

    # Insight box below the table
    c.box(x=0.3, y=4.5, w=9.4, h=2.5, fill="grey_100", border=0.5, border_color="grey_mid")
    c.text(
        "주요 발견사항",
        x=0.5, y=4.6, w=9.0, h=0.3, size=10, bold=True, color="grey_900",
    )
    c.text(
        "1. 조선 부문(4.2점)은 IoT 센서 기반 실시간 모니터링 도입으로 데이터 인프라 선도\n"
        "2. AI 활용도는 전 사업부 평균 2.9점으로 가장 큰 개선 여지 — 2026년 집중 투자 필요\n"
        "3. 건기 부문은 현장 디지털화 격차가 크며, 스마트 팩토리 파일럿 확대 권고\n"
        "4. 지주사의 데이터 거버넌스(3.0점) 강화가 그룹 시너지의 선결 과제\n"
        "5. 에너지 부문은 탄소 배출 모니터링 자동화로 ESG 대응 가속 가능",
        x=0.5, y=5.0, w=9.0, h=1.8, size=8, color="grey_900",
    )


# ============================================================
# Main
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_3_cover(prs)       # Slide 1: Cover
    slide_1_waterfall(prs)   # Slide 2: Waterfall
    slide_2_icon_text(prs)   # Slide 3: Icon + Text
    slide_4_data_table(prs)  # Slide 4: Data Table

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")

    # Evaluate
    report = evaluate_pptx(str(OUTPUT))
    print_report(report)
    return report


if __name__ == "__main__":
    main()
