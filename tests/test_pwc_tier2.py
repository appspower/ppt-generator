"""PwC Tier 2 Components — 8 slides (10~17) with PwC-style patterns.

Components:
  10. pwc_process_strip_4col — 4 equal grey monochrome columns
  11. pwc_diamond_quadrant — 4 corner cards + center diamond
  12. pwc_tab_variants — 6 tab/button styles on one slide
  13. pwc_hexagon_callout — 3 orange box callouts with numbers
  14. pwc_checklist — Checklist with sections and checkmarks
  15. pwc_nested_list — Accent border container with nested rows
  16. pwc_hierarchy_tree — 3-level hierarchy (1-3-6 boxes)
  17. pwc_connectors — Connector patterns using line() and arrow()
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from ppt_builder.primitives import Canvas, Region
from ppt_builder.evaluate import evaluate_pptx

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "pwc_tier2"


def make_slide(prs, title_text):
    """Create a blank slide with PwC-style header."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    c.text("PwC Tier 2 Components", x=0.3, y=0.15, w=9.4, h=0.25,
           size=9, color="grey_700", anchor="top")
    c.text(title_text, x=0.3, y=0.38, w=9.4, h=0.40,
           size=16, bold=True, color="grey_900", anchor="top")
    c.line(x1=0.3, y1=0.85, x2=9.7, y2=0.85, color="accent", width=2.0)
    return c


# ────────────────────────────────────────────────────────
# 10. pwc_process_strip_4col
# ────────────────────────────────────────────────────────

def slide_process_strip_4col(prs):
    c = make_slide(prs, "운영 프로세스 4단계 — 단계별 핵심 활동")

    cols = [
        {"fill": "grey_700", "num_color": "white", "txt_color": "white",
         "num": "01", "header": "현행 분석",
         "body": "업무 프로세스 매핑\n데이터 품질 진단\nGap 도출 및 정량화"},
        {"fill": "grey_400", "num_color": "white", "txt_color": "white",
         "num": "02", "header": "솔루션 설계",
         "body": "To-Be 프로세스 정의\n시스템 아키텍처 확정\n마이그레이션 전략 수립"},
        {"fill": "grey_200", "num_color": "grey_900", "txt_color": "grey_900",
         "num": "03", "header": "구현 및 테스트",
         "body": "모듈별 개발 실행\n통합 테스트 3단계\n사용자 수용 검증"},
        {"fill": "grey_100", "num_color": "grey_700", "txt_color": "grey_700",
         "num": "04", "header": "안정화 운영",
         "body": "Go-Live 지원 체계\n성과 모니터링 KPI\n지속적 개선 로드맵"},
    ]

    col_w = 2.2
    gap = 0.15
    start_x = 0.3
    top_y = 1.2

    for i, col in enumerate(cols):
        cx = start_x + i * (col_w + gap)
        # Column background
        c.box(x=cx, y=top_y, w=col_w, h=5.5, fill=col["fill"], border=None)
        # Large number
        c.text(col["num"], x=cx, y=top_y + 0.3, w=col_w, h=0.6,
               size=24, bold=True, color=col["num_color"], align="center", anchor="top")
        # Header
        c.text(col["header"], x=cx + 0.15, y=top_y + 1.1, w=col_w - 0.3, h=0.35,
               size=10, bold=True, color=col["txt_color"], anchor="top")
        # Body (3 lines)
        c.text(col["body"], x=cx + 0.15, y=top_y + 1.6, w=col_w - 0.3, h=2.0,
               size=8, color=col["txt_color"], anchor="top")


# ────────────────────────────────────────────────────────
# 11. pwc_diamond_quadrant
# ────────────────────────────────────────────────────────

def slide_diamond_quadrant(prs):
    c = make_slide(prs, "전략 목표 프레임워크 — 4대 핵심 축")

    # Center diamond
    center_x = 5.0
    center_y = 3.8
    c.diamond(cx=center_x, cy=center_y, size=1.0,
              fill="accent", text="Goal", text_color="white", text_size=12)

    # Corner cards: (x, y, title, body)
    cards = [
        (0.5, 1.2, "성장 전략", "시장 확대 및 신규\n사업 포트폴리오 구축"),
        (6.2, 1.2, "운영 혁신", "디지털 전환 기반\n프로세스 효율화 달성"),
        (0.5, 5.0, "인재 확보", "핵심 인력 유치 및\n조직 역량 강화"),
        (6.2, 5.0, "리스크 관리", "컴플라이언스 준수\n사이버 보안 강화"),
    ]

    for cx, cy, title, body in cards:
        c.box(x=cx, y=cy, w=3.2, h=1.8, fill="grey_100", border=0.75, border_color="grey_200")
        c.text(title, x=cx + 0.15, y=cy + 0.15, w=2.9, h=0.35,
               size=11, bold=True, color="grey_900", anchor="top")
        c.text(body, x=cx + 0.15, y=cy + 0.55, w=2.9, h=1.0,
               size=9, color="grey_700", anchor="top")

    # Connector lines from diamond to each card (implied by proximity)
    c.line(x1=center_x, y1=center_y - 0.55, x2=center_x, y2=center_y - 1.2, color="grey_400", width=1.0)
    c.line(x1=center_x, y1=center_y + 0.55, x2=center_x, y2=center_y + 1.2, color="grey_400", width=1.0)
    c.line(x1=center_x - 0.55, y1=center_y, x2=center_x - 2.5, y2=center_y, color="grey_400", width=1.0)
    c.line(x1=center_x + 0.55, y1=center_y, x2=center_x + 2.5, y2=center_y, color="grey_400", width=1.0)


# ────────────────────────────────────────────────────────
# 12. pwc_tab_variants
# ────────────────────────────────────────────────────────

def slide_tab_variants(prs):
    c = make_slide(prs, "UI 탭/버튼 스타일 가이드 — 6가지 변형")

    tab_w = 1.5
    tab_h = 0.4
    start_x = 0.5
    gap = 0.25
    row1_y = 1.5
    row2_y = 3.5
    label_offset = 0.55

    # Row 1: Plain Grey, Outlined, Orange Text
    # 1. Plain grey
    x1 = start_x
    c.box(x=x1, y=row1_y, w=tab_w, h=tab_h, fill="grey_200", border=None)
    c.text("Plain Grey", x=x1, y=row1_y, w=tab_w, h=tab_h,
           size=9, bold=True, color="grey_700", align="center", anchor="middle")
    c.text("fill=grey_200, no border", x=x1, y=row1_y + label_offset, w=tab_w + 0.5, h=0.3,
           size=7, color="grey_400", anchor="top")

    # 2. Outlined
    x2 = x1 + tab_w + gap + 1.0
    c.box(x=x2, y=row1_y, w=tab_w, h=tab_h, fill="white", border=1.0, border_color="grey_400")
    c.text("Outlined", x=x2, y=row1_y, w=tab_w, h=tab_h,
           size=9, bold=True, color="grey_700", align="center", anchor="middle")
    c.text("white fill, grey border", x=x2, y=row1_y + label_offset, w=tab_w + 0.5, h=0.3,
           size=7, color="grey_400", anchor="top")

    # 3. Orange text
    x3 = x2 + tab_w + gap + 1.0
    c.box(x=x3, y=row1_y, w=tab_w, h=tab_h, fill="white", border=0.75, border_color="grey_200")
    c.text("Orange Text", x=x3, y=row1_y, w=tab_w, h=tab_h,
           size=9, bold=True, color="accent", align="center", anchor="middle")
    c.text("white fill, accent text", x=x3, y=row1_y + label_offset, w=tab_w + 0.5, h=0.3,
           size=7, color="grey_400", anchor="top")

    # Row 2: Orange Fill, Orange+Grey Split, Grey with Accent Border
    # 4. Orange fill
    x4 = start_x
    c.box(x=x4, y=row2_y, w=tab_w, h=tab_h, fill="accent", border=None)
    c.text("Orange Fill", x=x4, y=row2_y, w=tab_w, h=tab_h,
           size=9, bold=True, color="white", align="center", anchor="middle")
    c.text("fill=accent, white text", x=x4, y=row2_y + label_offset, w=tab_w + 0.5, h=0.3,
           size=7, color="grey_400", anchor="top")

    # 5. Orange+Grey split
    x5 = x4 + tab_w + gap + 1.0
    half_w = tab_w / 2
    c.box(x=x5, y=row2_y, w=half_w, h=tab_h, fill="accent", border=None)
    c.text("Active", x=x5, y=row2_y, w=half_w, h=tab_h,
           size=8, bold=True, color="white", align="center", anchor="middle")
    c.box(x=x5 + half_w, y=row2_y, w=half_w, h=tab_h, fill="grey_200", border=None)
    c.text("Idle", x=x5 + half_w, y=row2_y, w=half_w, h=tab_h,
           size=8, bold=True, color="grey_700", align="center", anchor="middle")
    c.text("accent + grey split", x=x5, y=row2_y + label_offset, w=tab_w + 0.5, h=0.3,
           size=7, color="grey_400", anchor="top")

    # 6. Grey with accent border
    x6 = x5 + tab_w + gap + 1.0
    c.box(x=x6, y=row2_y, w=tab_w, h=tab_h, fill="grey_100", border=1.5, border_color="accent")
    c.text("Accent Bdr", x=x6, y=row2_y, w=tab_w, h=tab_h,
           size=9, bold=True, color="grey_900", align="center", anchor="middle")
    c.text("grey fill, accent border", x=x6, y=row2_y + label_offset, w=tab_w + 0.5, h=0.3,
           size=7, color="grey_400", anchor="top")


# ────────────────────────────────────────────────────────
# 13. pwc_hexagon_callout
# ────────────────────────────────────────────────────────

def slide_hexagon_callout(prs):
    c = make_slide(prs, "핵심 성과 지표 — 3대 전략 KPI")

    items = [
        {"num": "01", "label": "탄소 배출 41% 감축", "detail": "Scope 1+2 기준, 2030년 목표 대비 현재 달성률 기반 추정"},
        {"num": "02", "label": "연간 투자 $4.5T 필요", "detail": "IEA Net Zero 시나리오, 개도국 포함 글로벌 에너지 전환 총액"},
        {"num": "03", "label": "재생에너지 비중 60%", "detail": "발전 믹스 내 태양광+풍력 점유율, IRENA 2030 전망 기준"},
    ]

    for i, item in enumerate(items):
        top_y = 1.3 + i * 2.0
        # Orange box (hexagon approximation)
        c.box(x=0.5, y=top_y, w=1.2, h=1.2, fill="accent", border=None)
        c.text(item["num"], x=0.5, y=top_y, w=1.2, h=1.2,
               size=28, bold=True, color="white", align="center", anchor="middle")
        # Label to the right
        c.text(item["label"], x=2.0, y=top_y + 0.1, w=7.0, h=0.4,
               size=10, bold=True, color="grey_900", anchor="top")
        # Detail below label
        c.text(item["detail"], x=2.0, y=top_y + 0.55, w=7.0, h=0.5,
               size=8, color="grey_700", anchor="top")


# ────────────────────────────────────────────────────────
# 14. pwc_checklist
# ────────────────────────────────────────────────────────

def slide_checklist(prs):
    c = make_slide(prs, "Go-Live 준비 체크리스트 — 3대 영역 점검")

    sections = [
        {
            "subtitle": "데이터 마이그레이션",
            "items": ["✓  마스터 데이터 정합성 검증 완료", "✓  트랜잭션 데이터 이관 검증", "✓  데이터 롤백 절차 수립"],
            "bg": True,
        },
        {
            "subtitle": "시스템 통합 테스트",
            "items": ["✓  End-to-End 시나리오 테스트 통과", "✓  인터페이스 연동 검증 (ERP↔MES↔WMS)", "✓  성능 부하 테스트 기준 충족"],
            "bg": False,
        },
        {
            "subtitle": "조직 변화 관리",
            "items": ["✓  핵심 사용자 교육 완료 (200명)", "✓  업무 매뉴얼 배포 및 확인", "✓  헬프데스크 운영 체계 가동"],
            "bg": True,
        },
    ]

    cur_y = 1.2
    for sec in sections:
        section_h = 1.6
        # Alternating background band
        if sec["bg"]:
            c.box(x=0.3, y=cur_y, w=9.4, h=section_h, fill="grey_100", border=None)

        # Bold subtitle
        c.text(sec["subtitle"], x=0.5, y=cur_y + 0.1, w=8.0, h=0.35,
               size=11, bold=True, color="grey_900", anchor="top")

        # Checkmark items (indented)
        for j, item in enumerate(sec["items"]):
            c.text(item, x=0.8, y=cur_y + 0.5 + j * 0.35, w=8.5, h=0.3,
                   size=9, color="grey_700", anchor="top")

        cur_y += section_h + 0.15


# ────────────────────────────────────────────────────────
# 15. pwc_nested_list
# ────────────────────────────────────────────────────────

def slide_nested_list(prs):
    c = make_slide(prs, "프로젝트 위험 요소 상세 — 5대 리스크 영역")

    # Outer container with accent border
    container_x = 0.5
    container_y = 1.2
    container_w = 9.0
    container_h = 5.8
    c.box(x=container_x, y=container_y, w=container_w, h=container_h,
          fill="white", border=2.0, border_color="accent")

    rows = [
        ("기술 리스크", "레거시 시스템 호환성 문제로 인터페이스 오류 발생 가능성 높음"),
        ("일정 리스크", "핵심 모듈 커스터마이징 지연 시 Go-Live 일정 3개월 이상 지연 우려"),
        ("인력 리스크", "SAP 전문 컨설턴트 수급 부족, 내부 핵심 사용자 업무 병행 부담"),
        ("데이터 리스크", "마스터 데이터 표준화 미흡 — 10개 법인 코드 체계 불일치"),
        ("변화관리 리스크", "현업 저항 및 교육 부족으로 시스템 활용도 저하 가능"),
    ]

    row_h = 1.05
    inner_pad = 0.15
    for i, (subtitle, body) in enumerate(rows):
        ry = container_y + inner_pad + i * (row_h + 0.08)
        # Grey band for subtitle
        c.box(x=container_x + inner_pad, y=ry, w=container_w - 2 * inner_pad, h=0.35,
              fill="grey_200", border=None)
        c.text(subtitle, x=container_x + inner_pad + 0.1, y=ry, w=container_w - 0.6, h=0.35,
               size=10, bold=True, color="grey_900", anchor="middle")
        # Body text below
        c.text(body, x=container_x + inner_pad + 0.1, y=ry + 0.4, w=container_w - 0.6, h=0.55,
               size=8, color="grey_700", anchor="top")


# ────────────────────────────────────────────────────────
# 16. pwc_hierarchy_tree
# ────────────────────────────────────────────────────────

def slide_hierarchy_tree(prs):
    c = make_slide(prs, "조직 거버넌스 구조 — 3계층 의사결정 체계")

    # Level 1: 1 box (accent)
    l1_w = 2.5
    l1_h = 0.8
    l1_x = (10 - l1_w) / 2
    l1_y = 1.3
    c.box(x=l1_x, y=l1_y, w=l1_w, h=l1_h, fill="accent", border=None)
    c.text("Steering\nCommittee", x=l1_x, y=l1_y, w=l1_w, h=l1_h,
           size=11, bold=True, color="white", align="center", anchor="middle")

    # Level 2: 3 boxes (grey_700)
    l2_labels = ["프로그램 관리", "기술 아키텍처", "변화 관리"]
    l2_w = 2.5
    l2_h = 0.7
    l2_y = 2.7
    l2_gap = 0.35
    l2_total = 3 * l2_w + 2 * l2_gap
    l2_start = (10 - l2_total) / 2

    for i, label in enumerate(l2_labels):
        lx = l2_start + i * (l2_w + l2_gap)
        c.box(x=lx, y=l2_y, w=l2_w, h=l2_h, fill="grey_700", border=None)
        c.text(label, x=lx, y=l2_y, w=l2_w, h=l2_h,
               size=10, bold=True, color="white", align="center", anchor="middle")

    # Connectors L1 -> L2
    l1_cx = l1_x + l1_w / 2
    l1_bot = l1_y + l1_h
    for i in range(3):
        lx = l2_start + i * (l2_w + l2_gap) + l2_w / 2
        c.line(x1=l1_cx, y1=l1_bot, x2=lx, y2=l2_y, color="grey_400", width=1.0)

    # Level 3: 6 boxes (grey_200), 2 under each L2
    l3_labels = ["일정 관리", "예산 관리", "인프라", "보안", "교육 훈련", "커뮤니케이션"]
    l3_w = 1.15
    l3_h = 0.6
    l3_y = 4.0
    l3_gap = 0.15

    for i in range(3):
        parent_cx = l2_start + i * (l2_w + l2_gap) + l2_w / 2
        parent_bot = l2_y + l2_h
        child_start_x = parent_cx - l3_w - l3_gap / 2

        for j in range(2):
            idx = i * 2 + j
            cx = child_start_x + j * (l3_w + l3_gap)
            c.box(x=cx, y=l3_y, w=l3_w, h=l3_h, fill="grey_200", border=None)
            c.text(l3_labels[idx], x=cx, y=l3_y, w=l3_w, h=l3_h,
                   size=9, bold=True, color="grey_900", align="center", anchor="middle")
            # Connector L2 -> L3
            child_cx = cx + l3_w / 2
            c.line(x1=parent_cx, y1=parent_bot, x2=child_cx, y2=l3_y, color="grey_400", width=0.75)

    # Footer note
    c.text("* 각 워크스트림은 주간 보고 체계로 Steering Committee에 에스컬레이션",
           x=0.5, y=6.5, w=9.0, h=0.4, size=7, color="grey_400", anchor="top")


# ────────────────────────────────────────────────────────
# 17. pwc_connectors
# ────────────────────────────────────────────────────────

def slide_connectors(prs):
    c = make_slide(prs, "연결 패턴 — 시스템 간 인터페이스 표현")

    # Pattern (a): two grey boxes connected by red dashed line with arrow
    c.text("(a) Dashed connector with arrow", x=0.5, y=1.2, w=4.0, h=0.3,
           size=9, bold=True, color="grey_700", anchor="top")

    c.box(x=0.5, y=1.8, w=2.0, h=1.2, fill="grey_200", border=0.75, border_color="grey_400")
    c.text("ERP\nSystem", x=0.5, y=1.8, w=2.0, h=1.2,
           size=10, bold=True, color="grey_900", align="center", anchor="middle")

    c.box(x=5.0, y=1.8, w=2.0, h=1.2, fill="grey_200", border=0.75, border_color="grey_400")
    c.text("MES\nSystem", x=5.0, y=1.8, w=2.0, h=1.2,
           size=10, bold=True, color="grey_900", align="center", anchor="middle")

    # Red dashed line (use line + overlay approach since dashed is not natively in Canvas)
    # Draw multiple short segments to simulate dashing
    dash_y = 2.4
    seg_len = 0.25
    gap_len = 0.15
    cx = 2.6
    while cx + seg_len < 4.9:
        c.line(x1=cx, y1=dash_y, x2=cx + seg_len, y2=dash_y, color="accent", width=1.5)
        cx += seg_len + gap_len
    # Arrow at end
    c.arrow(x1=cx, y1=dash_y, x2=4.95, y2=dash_y, color="accent", width=1.5)

    # Pattern (b): grey box -> orange box connected by straight red line
    c.text("(b) Solid connector with arrow", x=0.5, y=4.0, w=4.0, h=0.3,
           size=9, bold=True, color="grey_700", anchor="top")

    c.box(x=0.5, y=4.6, w=2.0, h=1.2, fill="grey_200", border=0.75, border_color="grey_400")
    c.text("Legacy\nDB", x=0.5, y=4.6, w=2.0, h=1.2,
           size=10, bold=True, color="grey_900", align="center", anchor="middle")

    c.box(x=5.0, y=4.6, w=2.0, h=1.2, fill="accent", border=None)
    c.text("Cloud\nPlatform", x=5.0, y=4.6, w=2.0, h=1.2,
           size=10, bold=True, color="white", align="center", anchor="middle")

    # Straight red line with arrow
    c.arrow(x1=2.55, y1=5.2, x2=4.95, y2=5.2, color="accent", width=2.0)

    # Label on the connector
    c.text("API Gateway", x=3.0, y=4.85, w=1.5, h=0.3,
           size=8, bold=True, color="accent", align="center", anchor="top")


# ────────────────────────────────────────────────────────
# Main
# ────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_process_strip_4col(prs)
    slide_diamond_quadrant(prs)
    slide_tab_variants(prs)
    slide_hexagon_callout(prs)
    slide_checklist(prs)
    slide_nested_list(prs)
    slide_hierarchy_tree(prs)
    slide_connectors(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path} ({len(prs.slides)} slides)")

    # Evaluate
    report = evaluate_pptx(str(pptx_path))
    print(f"\nEvaluate score: {report['score']}/100  pass={report['pass']}")
    if report["issues"]:
        for iss in report["issues"]:
            print(f"  - {iss}")


if __name__ == "__main__":
    main()
