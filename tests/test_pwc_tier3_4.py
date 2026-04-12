"""PwC-Style Tier 3-4 Components — 6 slides (18~23).

18. pwc_icon_library     — 5×4 grid of 2-letter icon abbreviations in orange circles
19. pwc_pictogram_library — 4×3 grid of pictograms in accent_mid circles
20. pwc_waffle_chart     — 5×4 unit chart with "73%" large text
21. pwc_donut_simple     — Central percentage circle + 4 compass-position segments
22. pwc_hexagonal_cycle  — 6 boxes in hexagonal arrangement around a center
23. pwc_diamond_filled   — Diamond shape from 4 triangles + center + 4 corner texts

Output: output/pwc_tier3_4.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "pwc_tier3_4"


def make(prs, title_text: str, subtitle: str = ""):
    """Create blank slide with PwC-style header."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    # Dark top bar
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    # Title
    c.text(title_text, x=0.4, y=0.25, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.68, w=9.2, h=0.25,
               size=9, color="grey_700")
    # Separator
    c.line(x1=0.4, y1=0.98, x2=9.6, y2=0.98, color="grey_200", width=1.0)
    return s, c


# ============================================================
# Slide 18: pwc_icon_library — 5×4 grid, orange circles + labels
# ============================================================
def slide_18_icon_library(prs):
    s, c = make(prs, "Icon Library — 20 Standard Icons",
                "2-letter abbreviations for consistent visual language")

    icons = [
        ("AL", "Alert"), ("CL", "Cloud"), ("DB", "Database"), ("LK", "Lock"), ("CH", "Chart"),
        ("US", "User"), ("ML", "Mail"), ("SE", "Search"), ("ST", "Settings"), ("FI", "File"),
        ("GL", "Globe"), ("SH", "Shield"), ("WI", "Wifi"), ("CA", "Camera"), ("PR", "Print"),
        ("HO", "Home"), ("DL", "Download"), ("UP", "Upload"), ("PH", "Phone"), ("ED", "Edit"),
    ]

    cols, rows = 5, 4
    start_x, start_y = 0.6, 1.3
    col_w, row_h = 1.8, 1.4
    circle_d = 0.55

    for idx, (abbr, label) in enumerate(icons):
        row = idx // cols
        col = idx % cols
        cx = start_x + col * col_w
        cy = start_y + row * row_h

        # Orange circle with white abbreviation
        c.circle(x=cx + (col_w - circle_d) / 2, y=cy,
                 d=circle_d, fill="accent", border=None,
                 text=abbr, text_color="white", text_size=12, text_bold=True)
        # Label below
        c.text(label, x=cx, y=cy + circle_d + 0.08,
               w=col_w, h=0.25, size=7, color="grey_700",
               align="center", anchor="top")

    # Footer note
    c.text("※ 실제 프로젝트에서는 Segoe MDL2 Assets 벡터 아이콘으로 대체",
           x=0.4, y=7.0, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# Slide 19: pwc_pictogram_library — 4×3 grid, accent_mid circles
# ============================================================
def slide_19_pictogram_library(prs):
    s, c = make(prs, "Pictogram Library — Industry & Technology Symbols",
                "12 pictograms for sector-specific decks")

    items = [
        ("AI", "Artificial\nIntelligence"), ("CO2", "Carbon\nEmissions"),
        ("Au", "Automation"), ("IoT", "Internet of\nThings"),
        ("VR", "Virtual\nReality"), ("Tax", "Tax &\nRegulatory"),
        ("Sh", "Shipping &\nLogistics"), ("Bk", "Mobility"),
        ("Cl", "Cloud\nPlatform"), ("Ch", "Chemicals"),
        ("At", "Nuclear &\nEnergy"), ("Tm", "Team &\nPeople"),
    ]

    cols, rows = 4, 3
    start_x, start_y = 0.7, 1.3
    col_w, row_h = 2.2, 1.85
    circle_d = 0.75

    for idx, (abbr, label) in enumerate(items):
        row = idx // cols
        col = idx % cols
        cx = start_x + col * col_w
        cy = start_y + row * row_h

        # Larger accent_mid circle
        c.circle(x=cx + (col_w - circle_d) / 2, y=cy,
                 d=circle_d, fill="accent_mid", border=None,
                 text=abbr, text_color="white", text_size=14, text_bold=True)
        # Label below
        c.text(label, x=cx, y=cy + circle_d + 0.1,
               w=col_w, h=0.55, size=7, color="grey_700",
               align="center", anchor="top")


# ============================================================
# Slide 20: pwc_waffle_chart — 5×4 unit chart + 73% large text
# ============================================================
def slide_20_waffle_chart(prs):
    s, c = make(prs, "목표 달성률 73% — 2026 상반기 기준 순조로운 진행",
                "Waffle Chart: 각 셀 = 5%, 20셀 = 100%")

    pct = 73
    filled = round(pct / 5)  # 14.6 -> 15 cells filled out of 20

    # Large percentage on the left
    c.text("73%", x=0.5, y=2.0, w=3.0, h=2.0,
           size=60, bold=True, color="accent", align="center", anchor="middle")
    c.text("목표 달성률", x=0.5, y=4.0, w=3.0, h=0.5,
           size=14, bold=True, color="grey_900", align="center")
    c.text("2026 H1 기준 | 연말 목표 100%", x=0.5, y=4.5, w=3.0, h=0.35,
           size=9, color="grey_700", align="center")

    # 5×4 waffle grid on the right
    cols, rows = 5, 4
    sq_size = 0.65
    gap = 0.1
    grid_x = 4.5
    grid_y = 1.6

    cell_idx = 0
    for row in range(rows):
        for col in range(cols):
            x = grid_x + col * (sq_size + gap)
            y = grid_y + row * (sq_size + gap)
            is_filled = cell_idx < filled
            fill_color = "accent" if is_filled else "grey_200"
            c.box(x=x, y=y, w=sq_size, h=sq_size,
                  fill=fill_color, border=None)
            # Show percentage text in cell
            c.text(f"{cell_idx * 5 + 5}%", x=x, y=y, w=sq_size, h=sq_size,
                   size=7, bold=True,
                   color="white" if is_filled else "grey_400",
                   align="center", anchor="middle")
            cell_idx += 1

    # Legend
    c.box(x=4.5, y=5.0, w=0.25, h=0.25, fill="accent", border=None)
    c.text("달성 (5%/셀)", x=4.85, y=5.0, w=1.5, h=0.25,
           size=8, color="grey_700", anchor="middle")
    c.box(x=6.5, y=5.0, w=0.25, h=0.25, fill="grey_200", border=None)
    c.text("미달성", x=6.85, y=5.0, w=1.5, h=0.25,
           size=8, color="grey_700", anchor="middle")


# ============================================================
# Slide 21: pwc_donut_simple — Central circle + 4 compass segments
# ============================================================
def slide_21_donut_simple(prs):
    s, c = make(prs, "반도체 시장 점유율 — DRAM이 60.5%로 수익성 주도",
                "Donut Chart (simplified): 4 segments approximated by compass boxes")

    # Central circle
    center_d = 2.2
    center_x = (10 - center_d) / 2
    center_y = 2.3
    c.circle(x=center_x, y=center_y, d=center_d,
             fill="grey_900", border=None,
             text="60.5%", text_color="white", text_size=28, text_bold=True)
    c.text("전체 영업이익률", x=center_x, y=center_y + center_d + 0.1,
           w=center_d, h=0.3, size=9, color="grey_700", align="center")

    # 4 segment boxes in compass positions around center
    segments = [
        ("DRAM", "70%", "accent", 5.0 - 0.7, 1.4),       # top
        ("NAND", "53%", "grey_700", 7.2, 3.0),            # right
        ("Foundry", "17%", "grey_400", 5.0 - 0.7, 5.0),   # bottom
        ("Other", "12%", "grey_200", 1.4, 3.0),            # left
    ]

    seg_w, seg_h = 1.6, 1.0
    for label, value, fill_color, sx, sy in segments:
        c.box(x=sx, y=sy, w=seg_w, h=seg_h, fill=fill_color, border=None)
        text_color = "white" if fill_color in ("accent", "grey_700") else "grey_900"
        c.text(value, x=sx, y=sy + 0.08, w=seg_w, h=0.45,
               size=18, bold=True, color=text_color, align="center")
        c.text(label, x=sx, y=sy + 0.55, w=seg_w, h=0.3,
               size=9, bold=True, color=text_color, align="center")

    # Connecting lines from center to segments
    cx_mid = center_x + center_d / 2
    cy_mid = center_y + center_d / 2
    c.line(x1=cx_mid, y1=center_y, x2=cx_mid, y2=1.4 + seg_h,
           color="grey_mid", width=0.5)         # top
    c.line(x1=center_x + center_d, y1=cy_mid, x2=7.2, y2=cy_mid,
           color="grey_mid", width=0.5)          # right
    c.line(x1=cx_mid, y1=center_y + center_d, x2=cx_mid, y2=5.0,
           color="grey_mid", width=0.5)          # bottom
    c.line(x1=center_x, y1=cy_mid, x2=1.4 + seg_w, y2=cy_mid,
           color="grey_mid", width=0.5)          # left

    # Source
    c.text("Source: IDC Semiconductor Tracker 2025, Company Filings",
           x=0.4, y=6.8, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# Slide 22: pwc_hexagonal_cycle — 6 nodes + center in hex layout
# ============================================================
def slide_22_hexagonal_cycle(prs):
    s, c = make(prs, "디지털 전환 6대 핵심 역량 — Core Platform 중심의 유기적 연결",
                "Hexagonal Cycle: 6 capabilities around a core")

    # Center
    cx, cy = 4.25, 3.3
    cw, ch = 1.5, 1.0
    c.box(x=cx, y=cy, w=cw, h=ch, fill="accent", border=None)
    c.text("Core\nPlatform", x=cx, y=cy, w=cw, h=ch,
           size=11, bold=True, color="white", align="center", anchor="middle")

    # 6 surrounding nodes in hexagonal positions
    # Approximate hex: top, upper-right, lower-right, bottom, lower-left, upper-left
    nodes = [
        ("01", "Data &\nAnalytics", 4.25, 1.3),          # top
        ("02", "Cloud\nInfra", 6.6, 2.1),                 # upper-right
        ("03", "Cyber\nSecurity", 6.6, 4.5),              # lower-right
        ("04", "AI/ML\nOps", 4.25, 5.3),                  # bottom
        ("05", "IoT &\nEdge", 1.9, 4.5),                  # lower-left
        ("06", "UX &\nDesign", 1.9, 2.1),                 # upper-left
    ]

    nw, nh = 1.5, 1.0
    for num, label, nx, ny in nodes:
        c.box(x=nx, y=ny, w=nw, h=nh, fill="grey_700", border=None)
        c.text(num, x=nx + 0.05, y=ny + 0.05, w=0.35, h=0.25,
               size=8, bold=True, color="accent")
        c.text(label, x=nx, y=ny + 0.1, w=nw, h=nh - 0.15,
               size=9, bold=True, color="white", align="center", anchor="middle")

    # Connecting lines from center to each node
    center_mx = cx + cw / 2
    center_my = cy + ch / 2
    for _, _, nx, ny in nodes:
        nmx = nx + nw / 2
        nmy = ny + nh / 2
        c.line(x1=center_mx, y1=center_my, x2=nmx, y2=nmy,
               color="grey_400", width=0.75)

    # Outer ring description
    c.text("각 역량은 Core Platform과 양방향 데이터 교환 — Loosely Coupled, API-First 원칙",
           x=0.4, y=6.8, w=9.2, h=0.3, size=8, color="grey_700", align="center")


# ============================================================
# Slide 23: pwc_diamond_filled — Diamond shape + 4 corner texts
# ============================================================
def slide_23_diamond_filled(prs):
    s, c = make(prs, "전략 목표 달성을 위한 4대 축 — Goal 중심 Diamond Framework",
                "Diamond: 4 directional pillars converging on a central goal")

    # Center diamond using Canvas.diamond()
    c.diamond(cx=5.0, cy=3.8, size=1.5, fill="accent_mid",
              text="Goal", text_color="white", text_size=14)

    # 4 directional boxes (top, right, bottom, left)
    dir_boxes = [
        ("Growth", "매출 성장\n신사업 확대", 4.15, 1.5, "accent"),     # top
        ("Efficiency", "원가 절감\n프로세스 최적화", 7.0, 3.15, "grey_700"),  # right
        ("Innovation", "R&D 투자\n기술 리더십", 4.15, 5.4, "grey_700"),   # bottom
        ("People", "인재 확보\n조직 역량", 1.3, 3.15, "grey_700"),       # left
    ]

    bw, bh = 1.7, 1.2
    for title_txt, detail, bx, by, fill_col in dir_boxes:
        c.box(x=bx, y=by, w=bw, h=bh, fill=fill_col, border=None)
        c.text(title_txt, x=bx, y=by + 0.1, w=bw, h=0.35,
               size=11, bold=True, color="white", align="center")
        c.text(detail, x=bx, y=by + 0.5, w=bw, h=0.6,
               size=8, color="white", align="center", anchor="top")

    # Arrows from direction boxes to center diamond
    # Top → center
    c.arrow(x1=5.0, y1=1.5 + bh, x2=5.0, y2=3.8 - 0.75 - 0.1)
    # Right → center
    c.arrow(x1=7.0, y1=3.15 + bh / 2, x2=5.0 + 0.75 + 0.1, y2=3.8)
    # Bottom → center
    c.arrow(x1=5.0, y1=5.4, x2=5.0, y2=3.8 + 0.75 + 0.1)
    # Left → center
    c.arrow(x1=1.3 + bw, y1=3.15 + bh / 2, x2=5.0 - 0.75 - 0.1, y2=3.8)

    # 4 corner detail texts
    corners = [
        ("시장 목표", "2028년까지\n매출 $50B 돌파", 0.5, 1.2),
        ("효율 목표", "OPEX 15%\n절감 (3년)", 7.8, 1.2),
        ("혁신 목표", "R&D 비중\n매출 대비 12%", 0.5, 5.8),
        ("인재 목표", "핵심 인재\n이탈률 < 5%", 7.8, 5.8),
    ]

    for title_txt, body, tx, ty in corners:
        c.text(title_txt, x=tx, y=ty, w=1.8, h=0.25,
               size=8, bold=True, color="accent")
        c.text(body, x=tx, y=ty + 0.28, w=1.8, h=0.5,
               size=7, color="grey_700")


# ============================================================
# Evaluate
# ============================================================
def evaluate_and_report(pptx_path: Path) -> dict:
    from ppt_builder.evaluate import evaluate_pptx, print_report
    report = evaluate_pptx(str(pptx_path))
    print_report(report)
    return report


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_18_icon_library(prs)
    slide_19_pictogram_library(prs)
    slide_20_waffle_chart(prs)
    slide_21_donut_simple(prs)
    slide_22_hexagonal_cycle(prs)
    slide_23_diamond_filled(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX saved: {pptx_path}")

    # Evaluate
    report = evaluate_and_report(pptx_path)
    print(f"\nOutput: {pptx_path}")


if __name__ == "__main__":
    main()
