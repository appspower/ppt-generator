"""스마트 컴포넌트 + Zone 톤 + 디자인 토큰 통합 테스트.

기존 composer_test 대비 시각 품질 개선 확인.
"""

import sys
from pathlib import Path

_repo_root = Path(__file__).resolve().parent.parent
if str(_repo_root) not in sys.path:
    sys.path.insert(0, str(_repo_root))

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.composer import SlideComposer, apply_zone_tone
from ppt_builder.primitives import Region
from ppt_builder.patterns import SlideHeader, SlideFooter
from ppt_builder.components import (
    comp_kpi_row, comp_bar_chart_h, comp_bullet_list,
    comp_callout, comp_stat_row, comp_section_header,
    comp_progress_bar, comp_vertical_bars, comp_gauge,
    comp_data_card, comp_metric_delta, comp_icon_list,
    comp_timeline_mini, comp_heat_row, comp_tag_group,
)
from ppt_builder.visual_validate import validate_visual

FOOTER = SlideFooter(source="출처: PwC Analysis 2024", right="PwC")


# ============================================================
# 1. 스마트 KPI 대시보드 (data_card + progress_bar + gauge)
# ============================================================

def build_smart_dashboard():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="SAP 전환 4대 KPI 전원 목표 초과 달성 — 자동 트렌드 분석 기반 실시간 모니터링",
        category="스마트 대시보드 — Auto-Analyzed KPI",
        nav_path=["3. 성과", "1. 실시간 KPI"],
    ))

    zones = comp.layout("top_bottom", split=0.35)

    # 상단: 스마트 데이터 카드 4개 (자동 색상/트렌드)
    top = zones["top"]
    gap = 0.15
    card_w = (top.w - gap * 3) / 4
    cards = [
        {"value": 15.5, "label": "전환 일정 (개월)", "previous": 18, "target": 16,
         "unit": "개월", "higher_is_better": False, "detail": "목표 16 → 실제 15.5"},
        {"value": 70, "label": "테스트 공수 절감", "previous": 0, "target": 60,
         "unit": "%", "higher_is_better": True, "detail": "목표 60% → 실제 70%"},
        {"value": 99.8, "label": "TC 정확도", "previous": 85, "target": 95,
         "unit": "%", "higher_is_better": True, "detail": "기존 85% → 현재 99.8%"},
        {"value": 50, "label": "Cutover DT 단축", "previous": 0, "target": 40,
         "unit": "%", "higher_is_better": True, "detail": "목표 40% → 실제 50%"},
    ]
    for i, card in enumerate(cards):
        cr = top.sub(i * (card_w + gap), 0, card_w, top.h)
        comp_data_card(comp.canvas, region=cr, **card)

    # 하단: 좌 progress bar 4개 + 우 gauge 2개
    bottom = zones["bottom"]
    bz = {
        "left": Region(bottom.x, bottom.y, bottom.w * 0.55, bottom.h),
        "right": Region(bottom.x + bottom.w * 0.58, bottom.y, bottom.w * 0.42, bottom.h),
    }

    # 좌: 워크스트림 진행률
    apply_zone_tone(comp.canvas, bz["left"], "subtle", border=False)
    comp.canvas.push_region(bz["left"])
    comp.canvas.section_label("워크스트림 진행률", x=0.1, y=0.08, w=bz["left"].w - 0.2)
    comp.canvas.pop_region()

    streams = [
        ("FI/CO 재무회계", 95, 80),
        ("MM 구매관리", 82, 80),
        ("SD 영업물류", 78, 80),
        ("PP 생산계획", 65, 80),
        ("데이터 마이그레이션", 55, 80),
    ]
    bar_h = (bz["left"].h - 0.5) / len(streams)
    for i, (name, val, tgt) in enumerate(streams):
        comp_progress_bar(comp.canvas, label=name, value=val, target=tgt,
                          region=bz["left"].sub(0.1, 0.42 + i * bar_h,
                                                bz["left"].w - 0.2, bar_h))

    # 우: 게이지 2개
    apply_zone_tone(comp.canvas, bz["right"], "light", border=False)
    comp.canvas.push_region(bz["right"])
    comp.canvas.section_label("달성률", x=0.1, y=0.08, w=bz["right"].w - 0.2)
    comp.canvas.pop_region()

    gauge_h = (bz["right"].h - 0.5) / 2
    comp_gauge(comp.canvas, value=92, label="마일스톤 달성률", target=85,
               region=bz["right"].sub(0.1, 0.42, bz["right"].w - 0.2, gauge_h))
    comp_gauge(comp.canvas, value=97, label="결함 해소율", target=90,
               region=bz["right"].sub(0.1, 0.42 + gauge_h, bz["right"].w - 0.2, gauge_h))

    comp.takeaway("5개 워크스트림 중 3개 Green(≥80%), PP(65%)와 데이터(55%)는 Amber — 집중 대응 필요")
    comp.footer(FOOTER)
    return prs


# ============================================================
# 2. 변화량 분석 + 아이콘 리스트 (metric_delta + icon_list)
# ============================================================

def build_delta_analysis():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="Palantir 도입 전후 6대 지표 변화 — 테스트·Cutover·거버넌스 전 영역 개선",
        category="성과 분석 — Before vs After Metrics",
        nav_path=["3. 성과", "3. 변화량 분석"],
    ))
    comp.intro("6대 핵심 지표의 도입 전(baseline) vs 도입 후(current) 정량 비교 — 자동 트렌드 분석")

    zones = comp.layout("two_column", split=0.55)

    # 좌: metric_delta 6개
    apply_zone_tone(comp.canvas, zones["left"], "subtle", border=False)
    comp.canvas.push_region(zones["left"])
    comp.canvas.section_label("지표별 변화량", x=0.1, y=0.08, w=zones["left"].w - 0.2)
    comp.canvas.pop_region()

    metrics = [
        {"label": "전환 일정", "current": 15.5, "previous": 18.0, "unit": "개월", "higher_is_better": False},
        {"label": "TC 작성 공수", "current": 2.0, "previous": 24.0, "unit": "주", "higher_is_better": False},
        {"label": "TC 정확도", "current": 99.8, "previous": 85.0, "unit": "%", "higher_is_better": True},
        {"label": "커버리지", "current": 95.0, "previous": 55.0, "unit": "%", "higher_is_better": True},
        {"label": "Cutover DT", "current": 6.0, "previous": 12.0, "unit": "시간", "higher_is_better": False},
        {"label": "보고 주기", "current": 0.1, "previous": 3.0, "unit": "일", "higher_is_better": False},
    ]
    delta_h = (zones["left"].h - 0.5) / len(metrics)
    for i, m in enumerate(metrics):
        comp_metric_delta(comp.canvas, region=zones["left"].sub(
            0.1, 0.42 + i * delta_h, zones["left"].w - 0.2, delta_h), **m)

    # 우: 핵심 인사이트 아이콘 리스트
    apply_zone_tone(comp.canvas, zones["right"], "light", border=False)
    comp.canvas.push_region(zones["right"])
    comp.canvas.section_label("핵심 인사이트", x=0.1, y=0.08, w=zones["right"].w - 0.2)
    comp.canvas.pop_region()

    comp_icon_list(comp.canvas, items=[
        {"text": "테스트 공수 24주→2주 (92%↓) — AIP 자동 생성이 최대 기여", "icon": "check"},
        {"text": "TC 정확도 85%→99.8% — LLM-as-Judge 연속 검증 효과", "icon": "check"},
        {"text": "Cutover DT 12시간→6시간 — 3회 리허설 + Workshop App", "icon": "check"},
        {"text": "보고 주기 3일→실시간 — Health Dashboard 가동", "icon": "check"},
        {"text": "PP 모듈 진척률 65% — 생산계획 복잡도 예상 초과", "icon": "warn"},
        {"text": "데이터 마이그레이션 55% — 정합성 검증 지연 중", "icon": "error"},
    ], region=zones["right"].sub(0.1, 0.42, zones["right"].w - 0.2, zones["right"].h - 0.5))

    comp.takeaway("6대 지표 중 4개 Green, 2개 Warning — PP·데이터 마이그레이션에 Foundry Pipeline 자동 검증 즉시 투입 권장")
    comp.footer(FOOTER)
    return prs


# ============================================================
# 3. 히트맵 + 세로 바 + 타임라인 (grid_2x2)
# ============================================================

def build_rich_dashboard():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    comp = SlideComposer(slide)
    comp.header(SlideHeader(
        title="SAP 모듈 8개의 성숙도×진척률 히트맵 — MM·PP가 핵심 투자 대상",
        category="종합 현황 — Multi-View Dashboard",
        nav_path=["3. 거버넌스", "2. 종합 현황"],
    ))

    zones = comp.layout("grid_2x2")

    # TL: 히트맵
    apply_zone_tone(comp.canvas, zones["tl"], "light")
    comp.canvas.push_region(zones["tl"])
    comp.canvas.section_label("모듈별 진척률 히트맵", x=0.08, y=0.08, w=zones["tl"].w - 0.16)
    comp.canvas.pop_region()

    modules = [
        ("FI/CO", [95, 90, 88]),
        ("MM", [82, 75, 70]),
        ("SD", [78, 80, 85]),
        ("PP", [65, 55, 50]),
    ]
    heat_h = (zones["tl"].h - 0.8) / len(modules)
    # 헤더
    comp.canvas.text("Q1", x=zones["tl"].w * 0.35, y=0.38, w=zones["tl"].w * 0.2, h=0.2,
                     size=7, bold=True, color="grey_700", align="center", anchor="top",
                     region=zones["tl"])
    comp.canvas.text("Q2", x=zones["tl"].w * 0.55, y=0.38, w=zones["tl"].w * 0.2, h=0.2,
                     size=7, bold=True, color="grey_700", align="center", anchor="top",
                     region=zones["tl"])
    comp.canvas.text("Q3", x=zones["tl"].w * 0.75, y=0.38, w=zones["tl"].w * 0.2, h=0.2,
                     size=7, bold=True, color="grey_700", align="center", anchor="top",
                     region=zones["tl"])
    for i, (name, vals) in enumerate(modules):
        comp_heat_row(comp.canvas, label=name, values=vals, max_val=100,
                      region=zones["tl"].sub(0.08, 0.6 + i * heat_h,
                                             zones["tl"].w - 0.16, heat_h))

    # TR: 세로 바 차트
    apply_zone_tone(comp.canvas, zones["tr"], "subtle")
    comp.canvas.push_region(zones["tr"])
    comp.canvas.section_label("모듈별 ROI (배)", x=0.08, y=0.08, w=zones["tr"].w - 0.16)
    comp.canvas.pop_region()
    comp_vertical_bars(comp.canvas, data=[
        {"label": "MM", "value": 4.2},
        {"label": "PP", "value": 3.8},
        {"label": "SD", "value": 3.2},
        {"label": "FI", "value": 2.8},
    ], unit="배", region=zones["tr"].sub(0.08, 0.4, zones["tr"].w - 0.16, zones["tr"].h - 0.5))

    # BL: 미니 타임라인
    apply_zone_tone(comp.canvas, zones["bl"], "light")
    comp.canvas.push_region(zones["bl"])
    comp.canvas.section_label("프로젝트 일정", x=0.08, y=0.08, w=zones["bl"].w - 0.16)
    comp.canvas.pop_region()
    comp_timeline_mini(comp.canvas, phases=["킥오프", "Blueprint", "Build", "Test", "Go-Live"],
                       current=2,
                       region=zones["bl"].sub(0.08, 0.4, zones["bl"].w - 0.16, zones["bl"].h - 0.5))

    # BR: 태그 그룹 + 불릿
    apply_zone_tone(comp.canvas, zones["br"], "subtle")
    comp.canvas.push_region(zones["br"])
    comp.canvas.section_label("핵심 리스크", x=0.08, y=0.08, w=zones["br"].w - 0.16)
    comp.canvas.pop_region()
    comp_icon_list(comp.canvas, items=[
        {"text": "데이터 정합성 미검증", "icon": "error"},
        {"text": "PP 복잡도 예상 초과", "icon": "warn"},
        {"text": "인력 2명 이탈 예고", "icon": "warn"},
    ], region=zones["br"].sub(0.1, 0.4, zones["br"].w - 0.2, zones["br"].h - 0.5))

    comp.takeaway("MM(ROI 4.2배)·PP(3.8배)가 투자 최우선이나 PP 진척률 65%로 리스크 — 테스트 자동화 우선 투입")
    comp.footer(FOOTER)
    return prs


# ============================================================
# 메인
# ============================================================

def main():
    out_dir = Path("output/smart_composer")
    out_dir.mkdir(parents=True, exist_ok=True)

    cases = [
        ("01_smart_dashboard", build_smart_dashboard),
        ("02_delta_analysis", build_delta_analysis),
        ("03_rich_dashboard", build_rich_dashboard),
    ]

    print("=" * 70)
    print("Smart Composer Test — 3 High-Quality Slides")
    print("=" * 70)

    all_passed = True
    for name, builder_fn in cases:
        out = out_dir / f"{name}.pptx"
        try:
            prs = builder_fn()
            prs.save(out)
        except Exception as e:
            print(f"\n[{name}] BUILD FAILED: {e}")
            import traceback
            traceback.print_exc()
            all_passed = False
            continue

        visual = validate_visual(out, convert_pdf=False)
        ok = not visual.issues
        if not ok:
            all_passed = False
        print(f"\n[{name}] {'PASS' if ok else 'FAIL'}")
        print(f"  visual: {len(visual.issues)} issues")
        for i in visual.issues[:3]:
            print(f"    - {i}")

    print("\n" + "=" * 70)
    print("PDF...")
    for name, _ in cases:
        out = out_dir / f"{name}.pptx"
        if not out.exists():
            continue
        try:
            visual = validate_visual(out, convert_pdf=True)
            print(f"  {name}.pdf  -- {'OK' if visual.pdf_available else 'SKIP'}")
        except Exception as e:
            print(f"  {name}.pdf  -- FAIL ({e})")

    print(f"\nResult: {'ALL PASSED' if all_passed else 'SOME FAILED'}")
    return 0 if all_passed else 1


if __name__ == "__main__":
    sys.exit(main())
