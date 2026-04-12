"""Tier A — 10 chart-based components.

A1. comp_mekko_chart          — Marimekko (5 variable-width columns × 3 segments)
A2. comp_stacked_bar_100      — 100% stacked horizontal bars (5 bars × 4 segments)
A3. comp_grouped_bar          — Grouped bars (4 years × 3 series)
A4. comp_line_bar_combo       — Bar + line overlay (dual-axis feel)
A5. comp_bubble_chart         — 2D scatter with bubble size (6 bubbles)
A6. comp_tornado_chart        — Centered sensitivity bars (5 variables)
A7. comp_doughnut_with_detail — Big center % + 4 external segment labels
A8. comp_area_chart           — Stacked area (simulated with rectangles), 5 pts × 3 series
A9. comp_bridge_chart         — EBITDA bridge (7 bars, start/end + gains/losses)
A10. comp_sparkline_inline    — KPI table rows with inline mini sparklines

Output: output/tier_a_charts.pptx (+ .pdf, + _pngs/)
Run:  PYTHONPATH=. python tests/test_tier_a_charts.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas, Region

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "tier_a_charts"


# ============================================================
# Slide helper
# ============================================================
def make(prs, title_text: str, subtitle: str = ""):
    """Create blank slide with header bar + title + optional subtitle.

    Returns (slide, canvas, content_region_for_chart).
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    c.text(title_text, x=0.4, y=0.22, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.66, w=9.2, h=0.25,
               size=9, color="grey_700")
    c.line(x1=0.4, y1=0.96, x2=9.6, y2=0.96, color="grey_200", width=1.0)
    # Main content region
    region = Region(0.5, 1.15, 9.0, 5.9)
    return s, c, region


# ============================================================
# A1. comp_mekko_chart
# ============================================================
def comp_mekko_chart(c: Canvas, region: Region,
                     columns: list[tuple[str, float, list[tuple[str, float]]]]):
    """Marimekko: column widths proportional to totals, segments stack vertically.

    columns: list of (col_label, col_total_weight, [(seg_label, seg_pct), ...])
             seg_pct values should sum to 100.
    """
    # Reserve header + footer space
    header_h = 0.32
    footer_h = 0.28
    chart_x = region.x
    chart_y = region.y + header_h
    chart_w = region.w
    chart_h = region.h - header_h - footer_h

    total_weight = sum(w for _, w, _ in columns)
    seg_colors = ["grey_900", "grey_700", "accent"]

    cursor_x = chart_x
    for col_label, col_w, segments in columns:
        col_width = chart_w * (col_w / total_weight)

        # Column label (top)
        c.text(f"{col_label}\n{col_w:.0f}B",
               x=cursor_x, y=region.y, w=col_width, h=header_h,
               size=8, bold=True, color="grey_900",
               align="center", anchor="middle")

        # Stacked segments (sum to 100%)
        seg_y = chart_y
        for i, (seg_label, pct) in enumerate(segments):
            seg_h = chart_h * (pct / 100.0)
            fill = seg_colors[i % len(seg_colors)]
            c.box(x=cursor_x, y=seg_y, w=col_width, h=seg_h,
                  fill=fill, border=0.5, border_color="white")
            # Segment label inside (if tall enough)
            if seg_h > 0.35 and col_width > 0.8:
                txt_color = "white"
                c.text(f"{seg_label}\n{pct:.0f}%",
                       x=cursor_x, y=seg_y, w=col_width, h=seg_h,
                       size=8, bold=True, color=txt_color,
                       align="center", anchor="middle")
            seg_y += seg_h
        cursor_x += col_width

    # Footer legend
    legend_y = region.y + region.h - footer_h + 0.04
    lx = chart_x
    for i, name in enumerate(["Premium", "Mid-tier", "Entry"]):
        c.box(x=lx, y=legend_y + 0.05, w=0.14, h=0.14,
              fill=seg_colors[i], border=None)
        c.text(name, x=lx + 0.18, y=legend_y, w=1.2, h=0.22,
               size=8, color="grey_700", anchor="middle")
        lx += 1.4


def slide_a1_mekko(prs):
    s, c, r = make(prs,
                   "세그먼트별 시장 점유율 — 프리미엄 비중이 가장 높음",
                   "Marimekko: 카테고리 폭 = 시장 규모, 세그먼트 높이 = 가격대 비중")
    columns = [
        ("자동차",    180, [("Premium", 45), ("Mid-tier", 40), ("Entry", 15)]),
        ("전자",      120, [("Premium", 30), ("Mid-tier", 50), ("Entry", 20)]),
        ("화학",       90, [("Premium", 25), ("Mid-tier", 45), ("Entry", 30)]),
        ("철강",       60, [("Premium", 15), ("Mid-tier", 50), ("Entry", 35)]),
        ("기타",       50, [("Premium", 10), ("Mid-tier", 40), ("Entry", 50)]),
    ]
    comp_mekko_chart(c, r, columns)
    c.text("Source: 산업별 매출 데이터 2025, 내부 분석",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A2. comp_stacked_bar_100
# ============================================================
def comp_stacked_bar_100(c: Canvas, region: Region,
                          bars: list[tuple[str, list[float]]],
                          series_labels: list[str]):
    """100% stacked horizontal bars.

    bars: [(bar_label, [pct1, pct2, pct3, pct4]), ...]  pcts sum to 100.
    """
    n = len(bars)
    label_w = 1.4
    chart_x = region.x + label_w
    chart_w = region.w - label_w - 0.3  # leave room for right axis/legend
    row_h = 0.5
    gap = 0.2
    top_pad = 0.25
    legend_h = 0.5

    colors = ["accent", "grey_700", "grey_400", "grey_200"]
    text_colors = ["white", "white", "white", "grey_900"]

    # Legend on top
    lx = chart_x
    for i, lab in enumerate(series_labels):
        c.box(x=lx, y=region.y, w=0.14, h=0.14,
              fill=colors[i], border=None)
        c.text(lab, x=lx + 0.18, y=region.y - 0.02, w=1.4, h=0.2,
               size=8, color="grey_700", anchor="middle")
        lx += 1.6

    y = region.y + top_pad + 0.1
    for bar_label, pcts in bars:
        # Left row label
        c.text(bar_label, x=region.x, y=y, w=label_w - 0.1, h=row_h,
               size=10, bold=True, color="grey_900", anchor="middle")
        # Stacked segments
        sx = chart_x
        for i, pct in enumerate(pcts):
            seg_w = chart_w * (pct / 100.0)
            c.box(x=sx, y=y, w=seg_w, h=row_h,
                  fill=colors[i % len(colors)], border=0.5, border_color="white")
            if seg_w > 0.35:
                c.text(f"{pct:.0f}%",
                       x=sx, y=y, w=seg_w, h=row_h,
                       size=9, bold=True,
                       color=text_colors[i % len(text_colors)],
                       align="center", anchor="middle")
            sx += seg_w
        y += row_h + gap


def slide_a2_stacked_bar(prs):
    s, c, r = make(prs,
                   "부서별 예산 배분 — R&D가 Engineering에서 과반, Operations는 인건비 중심",
                   "100% Stacked Bar: 5개 부서의 비용 구성비 비교")
    bars = [
        ("Engineering", [55, 25, 12, 8]),
        ("Sales",       [20, 45, 25, 10]),
        ("Marketing",   [15, 30, 40, 15]),
        ("Operations",  [10, 55, 20, 15]),
        ("G&A",         [12, 38, 30, 20]),
    ]
    series_labels = ["R&D/인건비", "운영비", "마케팅", "기타"]
    comp_stacked_bar_100(c, r, bars, series_labels)
    c.text("Source: FY25 부서별 예산 계획",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A3. comp_grouped_bar
# ============================================================
def comp_grouped_bar(c: Canvas, region: Region,
                      categories: list[str],
                      series: list[tuple[str, list[float]]],
                      y_max: float | None = None):
    """Grouped vertical bars.

    categories: x-axis labels (len=N_groups)
    series: [(series_label, [val per category]), ...]  len=N_series
    """
    legend_w = 1.7
    chart_x = region.x + 0.4           # y-axis label room
    chart_y = region.y + 0.2
    chart_w = region.w - 0.4 - legend_w
    chart_h = region.h - 0.7            # x-axis label room
    axis_y = chart_y + chart_h

    n_cats = len(categories)
    n_series = len(series)
    group_w = chart_w / n_cats
    bar_w = (group_w * 0.7) / n_series
    inner_pad = (group_w - bar_w * n_series) / 2

    if y_max is None:
        y_max = max(v for _, vs in series for v in vs) * 1.15

    colors = ["accent", "grey_700", "grey_400"]

    # y-axis + gridlines (4 steps)
    for step in range(0, 5):
        gy = axis_y - (chart_h * step / 4)
        c.line(x1=chart_x, y1=gy, x2=chart_x + chart_w, y2=gy,
               color="grey_200", width=0.5)
        val = y_max * step / 4
        c.text(f"{val:.0f}",
               x=region.x, y=gy - 0.1, w=0.35, h=0.2,
               size=7, color="grey_700", align="right", anchor="middle")
    # x-axis baseline
    c.line(x1=chart_x, y1=axis_y, x2=chart_x + chart_w, y2=axis_y,
           color="grey_700", width=1.0)

    # bars
    for gi, cat in enumerate(categories):
        gx = chart_x + gi * group_w + inner_pad
        for si, (_, vals) in enumerate(series):
            v = vals[gi]
            bh = chart_h * (v / y_max)
            bx = gx + si * bar_w
            by = axis_y - bh
            c.box(x=bx, y=by, w=bar_w, h=bh,
                  fill=colors[si % len(colors)], border=None)
            # value label above bar
            c.text(f"{v:.0f}",
                   x=bx - 0.1, y=by - 0.22, w=bar_w + 0.2, h=0.2,
                   size=7, bold=True, color="grey_900", align="center")
        # category label
        c.text(cat,
               x=chart_x + gi * group_w, y=axis_y + 0.08,
               w=group_w, h=0.3,
               size=9, bold=True, color="grey_900", align="center")

    # Legend (right side)
    lx = chart_x + chart_w + 0.2
    ly = chart_y + 0.2
    for si, (lab, _) in enumerate(series):
        c.box(x=lx, y=ly, w=0.18, h=0.18,
              fill=colors[si % len(colors)], border=None)
        c.text(lab, x=lx + 0.24, y=ly - 0.02, w=1.4, h=0.22,
               size=9, color="grey_700", anchor="middle")
        ly += 0.35


def slide_a3_grouped_bar(prs):
    s, c, r = make(prs,
                   "연도별·사업부별 매출 — Enterprise 부문이 4년 연속 두 자릿수 성장",
                   "Grouped Bar: 3개 사업부의 연간 매출 (단위: 억 원)")
    categories = ["2022", "2023", "2024", "2025E"]
    series = [
        ("Enterprise",  [120, 145, 178, 210]),
        ("Consumer",    [95,  102, 108, 115]),
        ("Public",      [60,  72,  85,  92]),
    ]
    comp_grouped_bar(c, r, categories, series, y_max=250)
    c.text("Source: 연간 실적 공시 및 FY25 가이던스",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A4. comp_line_bar_combo
# ============================================================
def comp_line_bar_combo(c: Canvas, region: Region,
                         categories: list[str],
                         bar_values: list[float],
                         line_values: list[float],
                         bar_label: str = "매출",
                         line_label: str = "성장률 %"):
    """Bar (primary) + line (secondary y)."""
    chart_x = region.x + 0.5
    chart_y = region.y + 0.3
    chart_w = region.w - 1.0
    chart_h = region.h - 0.9
    axis_y = chart_y + chart_h

    n = len(categories)
    col_w = chart_w / n
    bar_w = col_w * 0.55

    bar_max = max(bar_values) * 1.2
    line_max = max(line_values) * 1.3
    line_min = min(0, min(line_values) * 1.2)

    # y-axis gridlines
    for step in range(0, 5):
        gy = axis_y - (chart_h * step / 4)
        c.line(x1=chart_x, y1=gy, x2=chart_x + chart_w, y2=gy,
               color="grey_200", width=0.5)
        # left axis = bar
        lval = bar_max * step / 4
        c.text(f"{lval:.0f}",
               x=region.x, y=gy - 0.1, w=0.45, h=0.2,
               size=7, color="grey_700", align="right", anchor="middle")
        # right axis = line
        rval = line_min + (line_max - line_min) * step / 4
        c.text(f"{rval:.0f}%",
               x=chart_x + chart_w + 0.05, y=gy - 0.1, w=0.5, h=0.2,
               size=7, color="accent", align="left", anchor="middle")

    c.line(x1=chart_x, y1=axis_y, x2=chart_x + chart_w, y2=axis_y,
           color="grey_700", width=1.0)

    # bars
    for i, (cat, v) in enumerate(zip(categories, bar_values)):
        bh = chart_h * (v / bar_max)
        bx = chart_x + i * col_w + (col_w - bar_w) / 2
        by = axis_y - bh
        c.box(x=bx, y=by, w=bar_w, h=bh, fill="grey_400", border=None)
        c.text(f"{v:.0f}",
               x=bx - 0.15, y=by - 0.22, w=bar_w + 0.3, h=0.2,
               size=7, bold=True, color="grey_900", align="center")
        c.text(cat,
               x=chart_x + i * col_w, y=axis_y + 0.08, w=col_w, h=0.25,
               size=9, bold=True, color="grey_900", align="center")

    # line (on secondary axis)
    def line_y(v: float) -> float:
        frac = (v - line_min) / (line_max - line_min)
        return axis_y - chart_h * frac

    pts = []
    for i, v in enumerate(line_values):
        px = chart_x + i * col_w + col_w / 2
        py = line_y(v)
        pts.append((px, py))

    # draw line segments
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        c.line(x1=x1, y1=y1, x2=x2, y2=y2, color="accent", width=2.0)

    # markers + labels
    for i, (px, py) in enumerate(pts):
        c.circle(x=px - 0.08, y=py - 0.08, d=0.16,
                 fill="accent", border=None)
        c.text(f"{line_values[i]:.1f}%",
               x=px - 0.4, y=py - 0.35, w=0.8, h=0.22,
               size=8, bold=True, color="accent", align="center")

    # legend
    ly = region.y + region.h - 0.3
    c.box(x=region.x + 0.2, y=ly, w=0.18, h=0.18, fill="grey_400", border=None)
    c.text(bar_label, x=region.x + 0.45, y=ly - 0.02, w=1.2, h=0.22,
           size=8, color="grey_700", anchor="middle")
    c.box(x=region.x + 2.0, y=ly + 0.06, w=0.3, h=0.06, fill="accent", border=None)
    c.text(line_label, x=region.x + 2.4, y=ly - 0.02, w=1.4, h=0.22,
           size=8, color="grey_700", anchor="middle")


def slide_a4_line_bar(prs):
    s, c, r = make(prs,
                   "매출은 지속 확대되나 성장률은 둔화 추세",
                   "Combo Chart: 매출 규모(bar) + 성장률(line) 동시 표시")
    categories = ["2021", "2022", "2023", "2024", "2025E"]
    bar_values = [120, 145, 178, 210, 240]
    line_values = [12.5, 20.8, 22.8, 18.0, 14.3]
    comp_line_bar_combo(c, r, categories, bar_values, line_values)
    c.text("Source: 연간 재무 보고서",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A5. comp_bubble_chart
# ============================================================
def comp_bubble_chart(c: Canvas, region: Region,
                       bubbles: list[tuple[str, float, float, float]],
                       x_label: str = "X",
                       y_label: str = "Y",
                       x_max: float | None = None,
                       y_max: float | None = None):
    """2D scatter with bubble size.

    bubbles: [(label, x, y, size_value), ...]
    """
    chart_x = region.x + 0.6
    chart_y = region.y + 0.2
    chart_w = region.w - 0.8
    chart_h = region.h - 0.9
    axis_y = chart_y + chart_h

    if x_max is None:
        x_max = max(b[1] for b in bubbles) * 1.2
    if y_max is None:
        y_max = max(b[2] for b in bubbles) * 1.2
    size_max = max(b[3] for b in bubbles)

    # gridlines
    for step in range(0, 5):
        gy = axis_y - (chart_h * step / 4)
        c.line(x1=chart_x, y1=gy, x2=chart_x + chart_w, y2=gy,
               color="grey_200", width=0.5)
        c.text(f"{y_max * step / 4:.0f}",
               x=region.x, y=gy - 0.1, w=0.55, h=0.2,
               size=7, color="grey_700", align="right", anchor="middle")
    for step in range(0, 5):
        gx = chart_x + (chart_w * step / 4)
        c.line(x1=gx, y1=chart_y, x2=gx, y2=axis_y,
               color="grey_200", width=0.5)
        c.text(f"{x_max * step / 4:.0f}",
               x=gx - 0.3, y=axis_y + 0.05, w=0.6, h=0.2,
               size=7, color="grey_700", align="center")

    # axes
    c.line(x1=chart_x, y1=axis_y, x2=chart_x + chart_w, y2=axis_y,
           color="grey_700", width=1.0)
    c.line(x1=chart_x, y1=chart_y, x2=chart_x, y2=axis_y,
           color="grey_700", width=1.0)

    # axis labels
    c.text(x_label,
           x=chart_x, y=axis_y + 0.32, w=chart_w, h=0.25,
           size=9, bold=True, color="grey_900", align="center")
    c.text(y_label,
           x=region.x - 0.05, y=chart_y + chart_h / 2 - 0.15, w=0.7, h=0.25,
           size=9, bold=True, color="grey_900", align="left")

    # bubbles
    for label, bx_v, by_v, bs_v in bubbles:
        # diameter: 0.3 ~ 1.1 in
        d = 0.3 + (bs_v / size_max) * 0.8
        px = chart_x + chart_w * (bx_v / x_max) - d / 2
        py = axis_y - chart_h * (by_v / y_max) - d / 2
        c.circle(x=px, y=py, d=d,
                 fill="accent", border=0.5, border_color="white")
        # label near bubble
        c.text(label,
               x=px - 0.2, y=py + d + 0.02, w=d + 0.4, h=0.2,
               size=8, bold=True, color="grey_900", align="center")


def slide_a5_bubble(prs):
    s, c, r = make(prs,
                   "시장 매력도 vs 자사 경쟁력 — B/D는 집중 투자, E/F는 수확 국면",
                   "Bubble: X=시장성장률, Y=자사점유율, 버블 크기=매출 규모")
    bubbles = [
        ("A",  8,  15,  30),
        ("B", 18,  35, 120),
        ("C", 12,  22,  50),
        ("D", 25,  48,  90),
        ("E",  5,  55,  70),
        ("F",  3,  40,  45),
    ]
    comp_bubble_chart(c, r, bubbles,
                       x_label="시장 성장률 (%)",
                       y_label="자사 점유율 (%)",
                       x_max=30, y_max=60)
    c.text("Source: 내부 시장 분석 2025",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A6. comp_tornado_chart
# ============================================================
def comp_tornado_chart(c: Canvas, region: Region,
                        variables: list[tuple[str, float, float]]):
    """Sensitivity tornado.

    variables: [(label, neg_impact, pos_impact), ...] — neg is negative number.
    Sorted by total absolute impact (largest first/top).
    """
    # Sort by combined absolute impact
    sorted_vars = sorted(variables,
                         key=lambda v: abs(v[1]) + abs(v[2]),
                         reverse=True)
    n = len(sorted_vars)

    label_w = 1.8
    chart_x = region.x + label_w
    chart_y = region.y + 0.3
    chart_w = region.w - label_w - 0.3
    chart_h = region.h - 0.7
    center_x = chart_x + chart_w / 2

    max_abs = max(max(abs(v[1]), abs(v[2])) for v in sorted_vars)
    half_w = chart_w / 2

    row_h = chart_h / n * 0.7
    row_gap = chart_h / n

    # center axis
    c.line(x1=center_x, y1=chart_y, x2=center_x, y2=chart_y + chart_h,
           color="grey_700", width=1.2)

    # axis scale labels at top
    c.text(f"-{max_abs:.0f}",
           x=chart_x - 0.2, y=chart_y - 0.3, w=0.5, h=0.22,
           size=7, color="grey_700")
    c.text("0",
           x=center_x - 0.15, y=chart_y - 0.3, w=0.3, h=0.22,
           size=7, color="grey_700", align="center")
    c.text(f"+{max_abs:.0f}",
           x=chart_x + chart_w - 0.25, y=chart_y - 0.3, w=0.5, h=0.22,
           size=7, color="grey_700", align="right")

    for i, (label, neg, pos) in enumerate(sorted_vars):
        row_y = chart_y + i * row_gap + (row_gap - row_h) / 2

        # left label
        c.text(label,
               x=region.x, y=row_y, w=label_w - 0.1, h=row_h,
               size=10, bold=True, color="grey_900",
               align="right", anchor="middle")

        # negative bar (extends left from center)
        nw = half_w * (abs(neg) / max_abs)
        c.box(x=center_x - nw, y=row_y, w=nw, h=row_h,
              fill="negative", border=None)
        c.text(f"{neg:+.1f}",
               x=center_x - nw - 0.6, y=row_y, w=0.55, h=row_h,
               size=8, bold=True, color="negative",
               align="right", anchor="middle")

        # positive bar (right)
        pw = half_w * (pos / max_abs)
        c.box(x=center_x, y=row_y, w=pw, h=row_h,
              fill="accent", border=None)
        c.text(f"+{pos:.1f}",
               x=center_x + pw + 0.05, y=row_y, w=0.55, h=row_h,
               size=8, bold=True, color="accent",
               align="left", anchor="middle")


def slide_a6_tornado(prs):
    s, c, r = make(prs,
                   "EBITDA 민감도 — 원자재 가격과 환율이 최대 리스크 요인",
                   "Tornado: 각 변수 ±10% 변동 시 EBITDA 영향 (단위: 억 원)")
    variables = [
        ("원자재 가격",  -35.0, 28.0),
        ("환율 (USD)",   -25.0, 30.0),
        ("판매 물량",    -18.0, 22.0),
        ("인건비",       -12.0,  8.0),
        ("물류비",        -8.0,  6.0),
    ]
    comp_tornado_chart(c, r, variables)
    c.text("Source: FY25 경영계획 민감도 분석",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A7. comp_doughnut_with_detail
# ============================================================
def comp_doughnut_with_detail(c: Canvas, region: Region,
                                center_value: str,
                                center_label: str,
                                segments: list[tuple[str, float, str]]):
    """Big center number + 4 colored 'segments' around (simulated).

    We approximate a doughnut with 4 quadrant rectangles + white center circle.
    segments: [(label, pct, anchor), ...] anchor in {"N","E","S","W"}
    """
    # Center of chart
    cx = region.x + region.w / 2
    cy = region.y + region.h / 2
    outer_d = 3.2
    inner_d = 1.9
    outer_x = cx - outer_d / 2
    outer_y = cy - outer_d / 2

    # Outer ring — 4 quadrant "wedges" approximated with quarter-sized circles
    colors = ["accent", "accent_mid", "grey_700", "grey_400"]
    # Use 4 full circles offset to create a ring feel isn't clean.
    # Instead: draw outer accent circle, overlay 3 pie-wedge substitutes as
    # colored boxes clipped-visually via rectangles at cardinal positions.
    # Simpler: draw 4 colored rectangles behind + big white circle to mask center.
    # Draw a large outer circle per segment color at cardinal offset.

    # Approach: draw 4 semicircle-ish rectangles at N/E/S/W
    q_w = outer_d / 2
    quads = [
        # (x, y, w, h, color) — N, E, S, W
        (cx - q_w, cy - q_w, q_w, q_w, colors[0]),   # NW
        (cx,      cy - q_w, q_w, q_w, colors[1]),    # NE
        (cx,      cy,       q_w, q_w, colors[2]),    # SE
        (cx - q_w, cy,      q_w, q_w, colors[3]),    # SW
    ]
    for qx, qy, qw, qh, fc in quads:
        c.box(x=qx, y=qy, w=qw, h=qh, fill=fc, border=None)

    # White masking circle in middle to create doughnut hole
    c.circle(x=cx - inner_d / 2, y=cy - inner_d / 2, d=inner_d,
             fill="white", border=2.0, border_color="grey_200")

    # Center big number + label
    c.text(center_value,
           x=cx - 1.2, y=cy - 0.55, w=2.4, h=0.7,
           size=32, bold=True, color="grey_900",
           align="center", anchor="middle")
    c.text(center_label,
           x=cx - 1.2, y=cy + 0.15, w=2.4, h=0.3,
           size=10, color="grey_700",
           align="center", anchor="middle")

    # External labels (top-right, bottom-right, bottom-left, top-left)
    anchors = [
        ("NW", cx - outer_d / 2 - 1.6, cy - outer_d / 2 + 0.1),
        ("NE", cx + outer_d / 2 + 0.1, cy - outer_d / 2 + 0.1),
        ("SE", cx + outer_d / 2 + 0.1, cy + outer_d / 2 - 0.5),
        ("SW", cx - outer_d / 2 - 1.6, cy + outer_d / 2 - 0.5),
    ]
    for i, (label, pct, _anchor) in enumerate(segments[:4]):
        _, lx, ly = anchors[i]
        # color chip
        c.box(x=lx, y=ly + 0.04, w=0.16, h=0.16, fill=colors[i], border=None)
        # label
        c.text(label,
               x=lx + 0.22, y=ly, w=1.4, h=0.22,
               size=9, bold=True, color="grey_900", anchor="middle")
        # pct
        c.text(f"{pct:.1f}%",
               x=lx + 0.22, y=ly + 0.22, w=1.4, h=0.22,
               size=11, bold=True, color=colors[i], anchor="middle")


def slide_a7_doughnut(prs):
    s, c, r = make(prs,
                   "ESG 목표 달성률 60.5% — 환경 부문이 상대적으로 저조",
                   "Doughnut: 4개 영역별 목표 대비 실적 비중")
    segments = [
        ("Governance",   25.4, "NW"),
        ("Social",       18.3, "NE"),
        ("Environment",  10.1, "SE"),
        ("Economic",      6.7, "SW"),
    ]
    comp_doughnut_with_detail(c, r,
                                center_value="60.5%",
                                center_label="전체 달성률",
                                segments=segments)
    c.text("Source: ESG 진단 보고서 2025 Q2",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A8. comp_area_chart
# ============================================================
def comp_area_chart(c: Canvas, region: Region,
                     x_labels: list[str],
                     series: list[tuple[str, list[float]]]):
    """Stacked area (simulated with vertical rectangles per x-band).

    For each x interval between x_i and x_{i+1}, we draw a rectangle per series
    whose height = value at x_i (step-area). Good enough for a presentation feel.
    """
    chart_x = region.x + 0.5
    chart_y = region.y + 0.3
    chart_w = region.w - 0.8
    chart_h = region.h - 0.9
    axis_y = chart_y + chart_h
    n = len(x_labels)

    # totals per x (for stacked max)
    totals = [sum(s[1][i] for s in series) for i in range(n)]
    y_max = max(totals) * 1.15

    # gridlines + y-labels
    for step in range(0, 5):
        gy = axis_y - (chart_h * step / 4)
        c.line(x1=chart_x, y1=gy, x2=chart_x + chart_w, y2=gy,
               color="grey_200", width=0.5)
        c.text(f"{y_max * step / 4:.0f}",
               x=region.x, y=gy - 0.1, w=0.45, h=0.2,
               size=7, color="grey_700", align="right", anchor="middle")
    c.line(x1=chart_x, y1=axis_y, x2=chart_x + chart_w, y2=axis_y,
           color="grey_700", width=1.0)

    # band width per x interval
    band_w = chart_w / (n - 1) if n > 1 else chart_w
    # Use full-light accent for fills (simulating transparency via colors)
    colors = ["accent_light", "grey_400", "grey_200"]
    border_colors = ["accent", "grey_700", "grey_400"]

    # Draw stacked rectangles per band using average of endpoints (step)
    for i in range(n - 1):
        bx = chart_x + i * band_w
        # For each series, rectangle height avg of values at i, i+1
        cumulative = 0.0
        for si, (_, vals) in enumerate(series):
            v = (vals[i] + vals[i + 1]) / 2.0
            bh = chart_h * (v / y_max)
            by = axis_y - cumulative - bh
            c.box(x=bx, y=by, w=band_w, h=bh,
                  fill=colors[si % len(colors)],
                  border=0.5, border_color=border_colors[si % len(border_colors)])
            cumulative += bh

    # x-axis labels
    for i, lab in enumerate(x_labels):
        gx = chart_x + i * band_w
        c.text(lab,
               x=gx - 0.3, y=axis_y + 0.08, w=0.6, h=0.22,
               size=9, bold=True, color="grey_900", align="center")

    # legend
    ly = region.y + region.h - 0.3
    lx = region.x + 0.3
    for si, (lab, _) in enumerate(series):
        c.box(x=lx, y=ly, w=0.2, h=0.18,
              fill=colors[si], border=0.5,
              border_color=border_colors[si])
        c.text(lab, x=lx + 0.26, y=ly - 0.02, w=1.4, h=0.22,
               size=8, color="grey_700", anchor="middle")
        lx += 1.7


def slide_a8_area(prs):
    s, c, r = make(prs,
                   "클라우드 사용량은 가속 성장, 온프레미스는 완만한 감소",
                   "Stacked Area: 3개 인프라 유형별 월별 사용량 추이")
    x_labels = ["Jan", "Apr", "Jul", "Oct", "Dec"]
    series = [
        ("Public Cloud",  [80, 120, 180, 240, 300]),
        ("Private Cloud", [60,  75,  95, 110, 120]),
        ("On-Prem",       [90,  85,  70,  55,  45]),
    ]
    comp_area_chart(c, r, x_labels, series)
    c.text("Source: FinOps 월간 사용량 리포트 2025",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A9. comp_bridge_chart
# ============================================================
def comp_bridge_chart(c: Canvas, region: Region,
                       items: list[tuple[str, float, str]]):
    """EBITDA bridge / waterfall.

    items: [(label, value, kind), ...] where kind in
        "start" | "end" | "inc" | "dec".
    'start' and 'end' are absolute; 'inc'/'dec' are deltas (inc positive, dec
    usually negative).
    """
    chart_x = region.x + 0.4
    chart_y = region.y + 0.4
    chart_w = region.w - 0.6
    chart_h = region.h - 1.0
    axis_y = chart_y + chart_h
    n = len(items)

    # compute running totals
    running = 0.0
    tops = []  # top of each bar
    bottoms = []  # bottom of each bar
    for label, v, kind in items:
        if kind == "start":
            bottoms.append(0.0)
            tops.append(v)
            running = v
        elif kind == "end":
            bottoms.append(0.0)
            tops.append(v)
            running = v
        elif kind == "inc":
            bottoms.append(running)
            tops.append(running + v)
            running = running + v
        elif kind == "dec":
            # v is negative; bar runs from running+v up to running
            new_r = running + v
            bottoms.append(new_r)
            tops.append(running)
            running = new_r

    y_max = max(tops) * 1.15
    y_min = 0

    col_w = chart_w / n
    bar_w = col_w * 0.6

    # gridlines
    for step in range(0, 5):
        gy = axis_y - (chart_h * step / 4)
        c.line(x1=chart_x, y1=gy, x2=chart_x + chart_w, y2=gy,
               color="grey_200", width=0.5)
        c.text(f"{y_max * step / 4:.0f}",
               x=region.x - 0.05, y=gy - 0.1, w=0.4, h=0.2,
               size=7, color="grey_700", align="right", anchor="middle")
    c.line(x1=chart_x, y1=axis_y, x2=chart_x + chart_w, y2=axis_y,
           color="grey_700", width=1.0)

    color_map = {
        "start": "grey_900",
        "end":   "grey_900",
        "inc":   "accent",
        "dec":   "negative",
    }

    def y_of(v: float) -> float:
        frac = (v - y_min) / (y_max - y_min)
        return axis_y - chart_h * frac

    prev_top = None
    for i, ((label, v, kind), top, bot) in enumerate(zip(items, tops, bottoms)):
        bx = chart_x + i * col_w + (col_w - bar_w) / 2
        by_top = y_of(top)
        by_bot = y_of(bot)
        bh = by_bot - by_top
        c.box(x=bx, y=by_top, w=bar_w, h=bh,
              fill=color_map[kind], border=None)

        # value label above bar
        if kind == "inc":
            disp = f"+{v:.0f}"
        elif kind == "dec":
            disp = f"{v:.0f}"
        else:
            disp = f"{v:.0f}"
        c.text(disp,
               x=bx - 0.2, y=by_top - 0.22, w=bar_w + 0.4, h=0.2,
               size=8, bold=True, color=color_map[kind], align="center")

        # category label
        c.text(label,
               x=chart_x + i * col_w, y=axis_y + 0.08, w=col_w, h=0.3,
               size=8, bold=True, color="grey_900", align="center")

        # connector line to next bar base (except for last)
        if i < n - 1:
            next_kind = items[i + 1][2]
            if next_kind in ("inc", "dec"):
                # connect top of this bar to top of this running total
                cy = y_of(running_after(items[:i + 1]))
                c.line(x1=bx + bar_w, y1=cy,
                       x2=chart_x + (i + 1) * col_w + (col_w - bar_w) / 2,
                       y2=cy, color="grey_400", width=0.75)


def running_after(partial_items):
    """Helper: compute running total after a list of items."""
    r = 0.0
    for label, v, kind in partial_items:
        if kind == "start" or kind == "end":
            r = v
        elif kind == "inc":
            r = r + v
        elif kind == "dec":
            r = r + v
    return r


def slide_a9_bridge(prs):
    s, c, r = make(prs,
                   "EBITDA Bridge FY24 → FY25 — 볼륨 성장이 원가 상승을 상쇄",
                   "Waterfall: 증가 요인(accent) / 감소 요인(red) / 시작·종료(dark)")
    items = [
        ("FY24 EBITDA", 520, "start"),
        ("Volume",      +85, "inc"),
        ("Price",       +40, "inc"),
        ("Mix",         +15, "inc"),
        ("COGS",        -60, "dec"),
        ("SG&A",        -25, "dec"),
        ("FY25 EBITDA", 575, "end"),
    ]
    comp_bridge_chart(c, r, items)
    c.text("Source: FY25 경영계획 Bridge Analysis",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# A10. comp_sparkline_inline
# ============================================================
def comp_sparkline_inline(c: Canvas, region: Region,
                           rows: list[tuple[str, str, list[float], float]]):
    """Table of KPI rows: label | big value | inline sparkline | delta%.

    rows: [(kpi_label, big_value_str, trend_values, delta_pct), ...]
    """
    n = len(rows)
    row_h = region.h / n * 0.85
    row_gap = region.h / n * 0.15

    # column widths
    col_label_w = 2.6
    col_value_w = 1.6
    col_spark_w = 3.2
    col_delta_w = 1.2
    col_value_x = region.x + col_label_w
    col_spark_x = col_value_x + col_value_w + 0.2
    col_delta_x = col_spark_x + col_spark_w + 0.2

    # header
    header_y = region.y
    c.text("KPI", x=region.x, y=header_y, w=col_label_w, h=0.28,
           size=9, bold=True, color="grey_700", anchor="middle")
    c.text("현재", x=col_value_x, y=header_y, w=col_value_w, h=0.28,
           size=9, bold=True, color="grey_700", anchor="middle")
    c.text("12개월 추이", x=col_spark_x, y=header_y, w=col_spark_w, h=0.28,
           size=9, bold=True, color="grey_700", anchor="middle")
    c.text("YoY", x=col_delta_x, y=header_y, w=col_delta_w, h=0.28,
           size=9, bold=True, color="grey_700", anchor="middle",
           align="right")
    # header separator
    c.line(x1=region.x, y1=header_y + 0.32,
           x2=region.x + region.w, y2=header_y + 0.32,
           color="grey_700", width=1.0)

    start_y = header_y + 0.42

    for i, (label, value, trend, delta) in enumerate(rows):
        y = start_y + i * (row_h + row_gap)

        # label
        c.text(label,
               x=region.x, y=y, w=col_label_w, h=row_h,
               size=11, bold=True, color="grey_900", anchor="middle")

        # big value
        c.text(value,
               x=col_value_x, y=y, w=col_value_w, h=row_h,
               size=18, bold=True, color="grey_900", anchor="middle")

        # sparkline: draw mini line inside the spark area
        spark_pad_x = 0.15
        spark_pad_y = row_h * 0.2
        spark_x0 = col_spark_x + spark_pad_x
        spark_y0 = y + spark_pad_y
        spark_w = col_spark_w - spark_pad_x * 2
        spark_h = row_h - spark_pad_y * 2

        tmin = min(trend)
        tmax = max(trend)
        trange = tmax - tmin if tmax > tmin else 1.0
        m = len(trend)
        pts = []
        for k, val in enumerate(trend):
            px = spark_x0 + (spark_w * k / (m - 1))
            py = spark_y0 + spark_h - (spark_h * (val - tmin) / trange)
            pts.append((px, py))

        # baseline (very light)
        c.line(x1=spark_x0, y1=spark_y0 + spark_h,
               x2=spark_x0 + spark_w, y2=spark_y0 + spark_h,
               color="grey_200", width=0.5)

        # sparkline segments
        line_color = "accent" if delta >= 0 else "negative"
        for k in range(len(pts) - 1):
            x1, y1 = pts[k]
            x2, y2 = pts[k + 1]
            c.line(x1=x1, y1=y1, x2=x2, y2=y2,
                   color=line_color, width=1.5)

        # end marker
        ex, ey = pts[-1]
        c.circle(x=ex - 0.06, y=ey - 0.06, d=0.12,
                 fill=line_color, border=None)

        # delta
        delta_color = "accent" if delta >= 0 else "negative"
        delta_txt = f"{delta:+.1f}%"
        c.text(delta_txt,
               x=col_delta_x, y=y, w=col_delta_w, h=row_h,
               size=12, bold=True, color=delta_color,
               align="right", anchor="middle")

        # row separator (light)
        if i < n - 1:
            sep_y = y + row_h + row_gap / 2
            c.line(x1=region.x, y1=sep_y,
                   x2=region.x + region.w, y2=sep_y,
                   color="grey_200", width=0.5)


def slide_a10_sparkline(prs):
    s, c, r = make(prs,
                   "핵심 KPI 대시보드 — 매출·마진 상승 추세, 인당 생산성 개선",
                   "Inline Sparkline: 4개 KPI의 현재값·12개월 추이·YoY 변동")
    rows = [
        ("매출",
         "₩2.4T",
         [180, 185, 195, 200, 210, 220, 218, 225, 235, 240, 250, 260],
         +12.5),
        ("영업이익률",
         "14.2%",
         [11.8, 12.0, 12.5, 12.8, 13.0, 13.2, 13.5, 13.7, 13.8, 14.0, 14.1, 14.2],
         +2.4),
        ("이탈률",
         "3.1%",
         [5.2, 5.0, 4.8, 4.5, 4.2, 4.0, 3.9, 3.7, 3.5, 3.3, 3.2, 3.1],
         -2.1),
        ("인당 매출",
         "₩850M",
         [720, 730, 740, 755, 765, 780, 790, 805, 815, 830, 840, 850],
         +18.1),
    ]
    comp_sparkline_inline(c, r, rows)
    c.text("Source: 내부 KPI 대시보드 (2025-04 기준)",
           x=0.4, y=7.05, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_a1_mekko(prs)
    slide_a2_stacked_bar(prs)
    slide_a3_grouped_bar(prs)
    slide_a4_line_bar(prs)
    slide_a5_bubble(prs)
    slide_a6_tornado(prs)
    slide_a7_doughnut(prs)
    slide_a8_area(prs)
    slide_a9_bridge(prs)
    slide_a10_sparkline(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path} ({len(prs.slides)} slides)")

    # PDF via PowerPoint COM
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        pdf_path = pptx_path.with_suffix(".pdf")
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        p = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        p.SaveAs(str(pdf_path.resolve()), 32)  # ppSaveAsPDF = 32
        p.Close()
        print(f"PDF:  {pdf_path}")
    except Exception as e:
        print(f"PDF export skipped: {e}")

    # PNG export
    try:
        from ppt_builder.track_c.png_export import pptx_to_pngs
        png_dir = OUTPUT_DIR / f"{NAME}_pngs"
        paths = pptx_to_pngs(pptx_path, png_dir)
        print(f"PNGs: {png_dir} ({len(paths)} images)")
    except Exception as e:
        print(f"PNG export skipped: {e}")


if __name__ == "__main__":
    main()
