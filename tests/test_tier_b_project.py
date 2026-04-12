"""Tier B — 8 project-management components.

B1. comp_harvey_balls        — 5×4 evaluation matrix with Harvey Ball glyphs
B2. comp_raci_matrix         — 6 tasks × 5 roles RACI
B3. comp_org_chart           — 4-level hierarchy (1 + 3 + 6)
B4. comp_stakeholder_map     — 2×2 influence/interest grid with dots
B5. comp_raid_log            — Risks / Assumptions / Issues / Dependencies 2×2
B6. comp_workstream_progress — 5 progress bars + RAG indicators
B7. comp_issue_tree_scored   — Issue tree with Harvey Ball scores on leaves
B8. comp_kanban_3col         — To Do / Doing / Done columns with task cards

Output: output/tier_b_project.pptx + .pdf + _pngs/
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "tier_b_project"


# ============================================================
# Common header
# ============================================================
def make(prs, title_text: str, subtitle: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    # Top accent bar
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    # Title
    c.text(title_text, x=0.4, y=0.22, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.66, w=9.2, h=0.25,
               size=9, color="grey_700")
    # Separator
    c.line(x1=0.4, y1=0.98, x2=9.6, y2=0.98, color="grey_200", width=1.0)
    return s, c


# Harvey Ball glyph map: 0~4 -> ○ ◔ ◑ ◕ ●
HARVEY = ["○", "◔", "◑", "◕", "●"]


# ============================================================
# B1. comp_harvey_balls — 5 rows × 4 criteria
# ============================================================
def slide_b1_harvey_balls(prs):
    s, c = make(prs, "공급사 평가 매트릭스 — 기술력이 C사 강점, 가격은 B사가 우위",
                "Harvey Balls: ● full / ◕ 3/4 / ◑ half / ◔ 1/4 / ○ empty")

    rows = ["공급사 A", "공급사 B", "공급사 C", "공급사 D", "공급사 E"]
    cols = ["기술력", "가격 경쟁력", "납기 신뢰", "사후지원"]
    # Score grid (0~4)
    scores = [
        [3, 2, 4, 3],
        [2, 4, 3, 2],
        [4, 2, 3, 4],
        [1, 3, 2, 3],
        [3, 3, 4, 2],
    ]

    grid_x, grid_y = 1.2, 1.5
    col_w = 1.9
    row_h = 0.85
    row_label_w = 1.8
    header_h = 0.55

    # Header band
    c.box(x=grid_x, y=grid_y, w=col_w * len(cols), h=header_h,
          fill="grey_900", border=None)
    for j, col in enumerate(cols):
        c.text(col, x=grid_x + j * col_w, y=grid_y, w=col_w, h=header_h,
               size=10, bold=True, color="white", align="center", anchor="middle")

    # Corner label for row header column
    c.box(x=grid_x - row_label_w, y=grid_y, w=row_label_w, h=header_h,
          fill="grey_900", border=None)
    c.text("공급사", x=grid_x - row_label_w, y=grid_y, w=row_label_w, h=header_h,
           size=10, bold=True, color="white", align="center", anchor="middle")

    # Rows
    for i, row_label in enumerate(rows):
        ry = grid_y + header_h + i * row_h
        # Row label cell
        c.box(x=grid_x - row_label_w, y=ry, w=row_label_w, h=row_h,
              fill="grey_100", border=0.5, border_color="grey_200")
        c.text(row_label, x=grid_x - row_label_w + 0.15, y=ry, w=row_label_w - 0.2, h=row_h,
               size=10, bold=True, color="grey_900", anchor="middle")
        # Cells
        for j in range(len(cols)):
            cx = grid_x + j * col_w
            c.box(x=cx, y=ry, w=col_w, h=row_h,
                  fill="white", border=0.5, border_color="grey_200")
            glyph = HARVEY[scores[i][j]]
            c.text(glyph, x=cx, y=ry, w=col_w, h=row_h,
                   size=14, bold=True, color="grey_900",
                   align="center", anchor="middle")

    # Legend
    legend_y = 6.55
    c.text("평가 범례:", x=0.4, y=legend_y, w=1.2, h=0.3,
           size=9, bold=True, color="grey_900", anchor="middle")
    legend_items = [("●", "Excellent"), ("◕", "Good"), ("◑", "Average"),
                    ("◔", "Below"), ("○", "Poor")]
    for k, (g, lbl) in enumerate(legend_items):
        lx = 1.5 + k * 1.5
        c.text(g, x=lx, y=legend_y, w=0.3, h=0.3,
               size=12, bold=True, color="grey_900", anchor="middle")
        c.text(lbl, x=lx + 0.3, y=legend_y, w=1.1, h=0.3,
               size=9, color="grey_700", anchor="middle")


# ============================================================
# B2. comp_raci_matrix — 6 tasks × 5 roles
# ============================================================
def slide_b2_raci_matrix(prs):
    s, c = make(prs, "RACI 매트릭스 — 단일 Accountable 원칙 준수 확인",
                "R=Responsible / A=Accountable / C=Consulted / I=Informed")

    tasks = [
        "요구사항 정의",
        "아키텍처 설계",
        "개발 및 단위 테스트",
        "통합/QA 테스트",
        "배포 및 이관",
        "운영 인수인계",
    ]
    roles = ["PM", "BA", "Dev", "QA", "Ops"]
    cells = [
        ["A", "R", "C", "I", "I"],
        ["A", "C", "R", "C", "C"],
        ["I", "C", "R", "C", "I"],
        ["I", "I", "C", "R", "A"],
        ["A", "I", "R", "C", "R"],
        ["C", "I", "I", "I", "A"],
    ]

    grid_x, grid_y = 2.0, 1.5
    row_label_w = 2.6
    col_w = 0.85
    row_h = 0.6
    header_h = 0.5

    code_fill = {
        "R": "accent",
        "A": "grey_900",
        "C": "grey_400",
        "I": "grey_200",
    }
    code_text_color = {
        "R": "white",
        "A": "white",
        "C": "white",
        "I": "grey_900",
    }

    # Header band (roles)
    c.box(x=grid_x, y=grid_y, w=col_w * len(roles), h=header_h,
          fill="grey_900", border=None)
    for j, role in enumerate(roles):
        c.text(role, x=grid_x + j * col_w, y=grid_y, w=col_w, h=header_h,
               size=11, bold=True, color="white", align="center", anchor="middle")

    # Task header cell
    c.box(x=grid_x - row_label_w, y=grid_y, w=row_label_w, h=header_h,
          fill="grey_900", border=None)
    c.text("Task", x=grid_x - row_label_w + 0.15, y=grid_y, w=row_label_w - 0.2, h=header_h,
           size=11, bold=True, color="white", anchor="middle")

    # Rows
    for i, task in enumerate(tasks):
        ry = grid_y + header_h + i * row_h
        c.box(x=grid_x - row_label_w, y=ry, w=row_label_w, h=row_h,
              fill="grey_100", border=0.5, border_color="grey_200")
        c.text(task, x=grid_x - row_label_w + 0.15, y=ry, w=row_label_w - 0.2, h=row_h,
               size=9, color="grey_900", anchor="middle")
        for j, code in enumerate(cells[i]):
            cx = grid_x + j * col_w
            if code:
                c.box(x=cx, y=ry, w=col_w, h=row_h,
                      fill=code_fill[code], border=0.5, border_color="white")
                c.text(code, x=cx, y=ry, w=col_w, h=row_h,
                       size=12, bold=True, color=code_text_color[code],
                       align="center", anchor="middle")
            else:
                c.box(x=cx, y=ry, w=col_w, h=row_h,
                      fill="white", border=0.5, border_color="grey_200")

    # Legend
    legend_y = 6.55
    legend = [("R", "Responsible"), ("A", "Accountable"),
              ("C", "Consulted"), ("I", "Informed")]
    for k, (code, lbl) in enumerate(legend):
        lx = 1.0 + k * 2.1
        c.box(x=lx, y=legend_y, w=0.4, h=0.4, fill=code_fill[code], border=None)
        c.text(code, x=lx, y=legend_y, w=0.4, h=0.4,
               size=11, bold=True, color=code_text_color[code],
               align="center", anchor="middle")
        c.text(lbl, x=lx + 0.5, y=legend_y, w=1.6, h=0.4,
               size=9, color="grey_900", anchor="middle")


# ============================================================
# B3. comp_org_chart — 1 + 3 + 6 hierarchy
# ============================================================
def slide_b3_org_chart(prs):
    s, c = make(prs, "조직 구조 — CEO 직속 3본부, 6개 팀 운영",
                "4-level hierarchy with L-shaped connectors")

    bw, bh = 1.55, 0.65
    mid_gap_y = 1.45
    bot_gap_y = 2.5

    # Level 1 — CEO
    ceo_x = (10 - bw) / 2
    ceo_y = 1.3
    c.box(x=ceo_x, y=ceo_y, w=bw, h=bh,
          fill="grey_100", border=1.5, border_color="accent")
    c.text("CEO", x=ceo_x, y=ceo_y, w=bw, h=bh,
           size=12, bold=True, color="grey_900", align="center", anchor="middle")

    # Level 2 — 3 directors
    mids = ["전략본부장", "기술본부장", "운영본부장"]
    mid_y = ceo_y + mid_gap_y
    mid_xs = [1.6, (10 - bw) / 2, 10 - bw - 1.6]
    for mx, label in zip(mid_xs, mids):
        c.box(x=mx, y=mid_y, w=bw, h=bh,
              fill="grey_100", border=1.5, border_color="accent")
        c.text(label, x=mx, y=mid_y, w=bw, h=bh,
               size=10, bold=True, color="grey_900", align="center", anchor="middle")

    # CEO -> mids: trunk down then horizontal then up into each mid
    trunk_y = ceo_y + bh + 0.25
    c.line(x1=ceo_x + bw / 2, y1=ceo_y + bh, x2=ceo_x + bw / 2, y2=trunk_y,
           color="grey_700", width=1.2)
    leftmost_mid_cx = mid_xs[0] + bw / 2
    rightmost_mid_cx = mid_xs[-1] + bw / 2
    c.line(x1=leftmost_mid_cx, y1=trunk_y, x2=rightmost_mid_cx, y2=trunk_y,
           color="grey_700", width=1.2)
    for mx in mid_xs:
        c.line(x1=mx + bw / 2, y1=trunk_y, x2=mx + bw / 2, y2=mid_y,
               color="grey_700", width=1.2)

    # Level 3 — 6 teams (2 under each mid)
    teams = [
        ("전략기획팀", "사업개발팀"),
        ("R&D팀", "플랫폼팀"),
        ("영업팀", "지원팀"),
    ]
    bot_y = mid_y + bot_gap_y - 1.35
    # Positioning: under each mid, two boxes side by side
    for i, mx in enumerate(mid_xs):
        t1, t2 = teams[i]
        # two boxes centered under mid
        spacing = 0.15
        tw = 1.35
        t1x = mx + bw / 2 - tw - spacing / 2
        t2x = mx + bw / 2 + spacing / 2
        ty = bot_y
        for tx, tlabel in [(t1x, t1), (t2x, t2)]:
            c.box(x=tx, y=ty, w=tw, h=bh,
                  fill="grey_100", border=1.5, border_color="accent")
            c.text(tlabel, x=tx, y=ty, w=tw, h=bh,
                   size=9, bold=True, color="grey_900", align="center", anchor="middle")

        # Connectors from mid to two teams (L-shape)
        mid_bottom_x = mx + bw / 2
        mid_bottom_y = mid_y + bh
        connector_y = mid_bottom_y + 0.22
        c.line(x1=mid_bottom_x, y1=mid_bottom_y, x2=mid_bottom_x, y2=connector_y,
               color="grey_700", width=1.0)
        c.line(x1=t1x + tw / 2, y1=connector_y, x2=t2x + tw / 2, y2=connector_y,
               color="grey_700", width=1.0)
        for tx in (t1x, t2x):
            c.line(x1=tx + tw / 2, y1=connector_y, x2=tx + tw / 2, y2=ty,
                   color="grey_700", width=1.0)

    c.text("※ 각 팀은 본부장 직속 보고, 본부 간 협업은 주간 운영회의 통해 정렬",
           x=0.4, y=6.9, w=9.2, h=0.3, size=8, color="grey_700", align="center")


# ============================================================
# B4. comp_stakeholder_map — 2×2 influence/interest
# ============================================================
def slide_b4_stakeholder_map(prs):
    s, c = make(prs, "스테이크홀더 맵 — CIO/CFO 적극 관리 필요, 구매팀 만족 유지",
                "Power / Interest Grid — 2×2 stakeholder prioritization")

    # Axis frame
    grid_x, grid_y = 1.5, 1.4
    grid_w, grid_h = 7.5, 4.8
    mid_x = grid_x + grid_w / 2
    mid_y = grid_y + grid_h / 2

    # Quadrant fills
    quads = [
        # (x, y, w, h, fill, title, desc)
        (grid_x, grid_y, grid_w / 2, grid_h / 2, "grey_100",
         "Keep Satisfied", "High 영향력 / Low 관심"),
        (mid_x, grid_y, grid_w / 2, grid_h / 2, "zone_negative",
         "Manage Closely", "High 영향력 / High 관심"),
        (grid_x, mid_y, grid_w / 2, grid_h / 2, "grey_200",
         "Monitor", "Low 영향력 / Low 관심"),
        (mid_x, mid_y, grid_w / 2, grid_h / 2, "zone_positive",
         "Keep Informed", "Low 영향력 / High 관심"),
    ]
    for qx, qy, qw, qh, fill, title_q, desc in quads:
        c.box(x=qx, y=qy, w=qw, h=qh, fill=fill,
              border=0.75, border_color="grey_400")
        c.text(title_q, x=qx + 0.1, y=qy + 0.1, w=qw - 0.2, h=0.3,
               size=11, bold=True, color="grey_900")
        c.text(desc, x=qx + 0.1, y=qy + 0.4, w=qw - 0.2, h=0.25,
               size=8, color="grey_700")

    # Axis labels
    c.text("관심도 Low  →  High",
           x=grid_x, y=grid_y + grid_h + 0.1, w=grid_w, h=0.3,
           size=10, bold=True, color="grey_900", align="center")
    # Y-axis label — vertical text approximation
    c.text("영향력\nHigh\n\n\nLow",
           x=grid_x - 0.9, y=grid_y, w=0.8, h=grid_h,
           size=10, bold=True, color="grey_900", align="center", anchor="middle")

    # Stakeholder dots (x_frac, y_frac, name)
    # Coordinates within the grid
    stakeholders = [
        (0.78, 0.22, "CIO"),         # Manage Closely
        (0.62, 0.32, "CFO"),         # Manage Closely
        (0.20, 0.25, "이사회"),        # Keep Satisfied
        (0.85, 0.70, "현업 PM"),       # Keep Informed
        (0.45, 0.78, "구매팀"),        # Monitor (near boundary)
        (0.30, 0.65, "감사팀"),        # Monitor
    ]
    dot_d = 0.3
    for fx, fy, name in stakeholders:
        cx_abs = grid_x + fx * grid_w - dot_d / 2
        cy_abs = grid_y + fy * grid_h - dot_d / 2
        c.circle(x=cx_abs, y=cy_abs, d=dot_d, fill="accent", border=None)
        c.text(name, x=cx_abs + dot_d + 0.05, y=cy_abs - 0.02,
               w=1.2, h=0.3, size=9, bold=True, color="grey_900", anchor="middle")


# ============================================================
# B5. comp_raid_log — 2×2 R/A/I/D
# ============================================================
def slide_b5_raid_log(prs):
    s, c = make(prs, "RAID 로그 — 위험 관리 및 이슈 추적 현황",
                "Risks / Assumptions / Issues / Dependencies — Week 12")

    panels = [
        # (x, y, title, items)
        (0.4, 1.3, "Risks", [
            "공급사 납기 지연 (확률 中, 영향 高)",
            "핵심 개발자 이탈 (확률 低, 영향 高)",
            "법규 변경에 따른 재설계 (확률 中, 영향 中)",
        ]),
        (5.1, 1.3, "Assumptions", [
            "사내 인프라 IaaS 환경 2025Q2 준비 완료",
            "경영진 예산 승인 $2.5M 유지",
            "파트너사 API v3 릴리스 일정 준수",
        ]),
        (0.4, 4.3, "Issues", [
            "QA 테스트 환경 구성 지연 (Day -5)",
            "요구사항 변경 요청 12건 대기",
            "부서간 협업 프로세스 미정립",
        ]),
        (5.1, 4.3, "Dependencies", [
            "ERP 마스터 데이터 이관 (재무팀)",
            "네트워크 보안 정책 승인 (정보보안팀)",
            "UAT 참여자 확보 (현업 10명)",
        ]),
    ]

    pw, ph = 4.5, 2.9
    for px, py, title_p, items in panels:
        c.box(x=px, y=py, w=pw, h=ph,
              fill="white", border=1.0, border_color="grey_400")
        # Accent header bar
        c.box(x=px, y=py, w=pw, h=0.45, fill="accent", border=None)
        c.text(title_p, x=px + 0.2, y=py, w=pw - 0.4, h=0.45,
               size=12, bold=True, color="white", anchor="middle")
        # Items
        for k, item in enumerate(items):
            iy = py + 0.6 + k * 0.65
            c.circle(x=px + 0.25, y=iy + 0.08, d=0.12,
                     fill="accent", border=None)
            c.text(item, x=px + 0.5, y=iy, w=pw - 0.6, h=0.55,
                   size=9, color="grey_900", anchor="top")


# ============================================================
# B6. comp_workstream_progress — 5 rows with bars + RAG
# ============================================================
def slide_b6_workstream_progress(prs):
    s, c = make(prs, "워크스트림 진척도 — WS3 Red, 즉시 리소스 재배분 필요",
                "5 workstreams with progress bars and RAG status")

    workstreams = [
        ("WS1 — 요구사항 분석", 95, "green"),
        ("WS2 — 아키텍처 설계", 80, "green"),
        ("WS3 — 개발 (Core)", 42, "red"),
        ("WS4 — 테스트/QA", 30, "amber"),
        ("WS5 — 인프라/배포", 60, "amber"),
    ]
    rag_color = {"green": "positive", "amber": "warning", "red": "negative"}
    rag_label = {"green": "On Track", "amber": "At Risk", "red": "Off Track"}

    row_x = 0.4
    label_w = 2.7
    bar_x = row_x + label_w + 0.15
    bar_w = 5.3
    rag_x = bar_x + bar_w + 0.25
    rag_d = 0.35
    row_y0 = 1.5
    row_h = 0.85

    # Header row
    c.text("Workstream", x=row_x, y=row_y0 - 0.4, w=label_w, h=0.3,
           size=9, bold=True, color="grey_700")
    c.text("Progress", x=bar_x, y=row_y0 - 0.4, w=bar_w, h=0.3,
           size=9, bold=True, color="grey_700")
    c.text("Status", x=rag_x - 0.05, y=row_y0 - 0.4, w=1.2, h=0.3,
           size=9, bold=True, color="grey_700")

    for i, (name, pct, status) in enumerate(workstreams):
        ry = row_y0 + i * row_h
        # Label
        c.box(x=row_x, y=ry, w=label_w, h=row_h - 0.15,
              fill="grey_100", border=None)
        c.text(name, x=row_x + 0.15, y=ry, w=label_w - 0.2, h=row_h - 0.15,
               size=10, bold=True, color="grey_900", anchor="middle")
        # Bar background
        bar_h = 0.4
        bar_y = ry + (row_h - 0.15 - bar_h) / 2
        c.box(x=bar_x, y=bar_y, w=bar_w, h=bar_h,
              fill="grey_200", border=None)
        # Filled
        filled_w = bar_w * pct / 100.0
        fill_col = rag_color[status]
        if filled_w > 0.01:
            c.box(x=bar_x, y=bar_y, w=filled_w, h=bar_h,
                  fill=fill_col, border=None)
        # Percentage text at right of the bar
        c.text(f"{pct}%",
               x=bar_x + bar_w + 0.05, y=bar_y, w=0.6, h=bar_h,
               size=10, bold=True, color="grey_900", anchor="middle")
        # Wait — rag_x is past bar end. Move % inside bar area end:
        # (the above places it at rag_x, which overlaps the circle). Remove it.
        # We'll place % text inside bar instead via text box.
        # RAG indicator circle
        cx = rag_x + 0.6
        cy = ry + (row_h - 0.15 - rag_d) / 2
        c.circle(x=cx, y=cy, d=rag_d, fill=fill_col, border=None)
        c.text(rag_label[status], x=cx + rag_d + 0.1, y=cy,
               w=1.2, h=rag_d, size=8, bold=True,
               color="grey_900", anchor="middle")

    # Footer note
    c.text("Source: PMO Dashboard | As of Week 12 | 업데이트 주기: 매주 금요일",
           x=0.4, y=6.85, w=9.2, h=0.3, size=8, color="grey_400")


# ============================================================
# B7. comp_issue_tree_scored — root + 3 branches + leaves w/ Harvey
# ============================================================
def slide_b7_issue_tree_scored(prs):
    s, c = make(prs, "매출 감소 원인 트리 — 신규고객 감소가 최대 영향 요인",
                "Issue tree with Harvey Ball impact scores on leaves")

    # Root
    root_x, root_y = 0.6, 3.3
    rw, rh = 2.0, 0.9
    c.box(x=root_x, y=root_y, w=rw, h=rh,
          fill="grey_900", border=None)
    c.text("매출 감소", x=root_x, y=root_y, w=rw, h=rh,
           size=14, bold=True, color="white", align="center", anchor="middle")

    # 3 branches
    branches = [
        ("신규고객 감소", [("광고 ROI 하락", 4), ("영업 인력 축소", 3)]),
        ("이탈률 증가", [("경쟁사 진입", 3), ("서비스 품질", 2)]),
        ("ARPU 하락", [("할인 프로모션 확대", 2), ("상위 플랜 이탈", 1)]),
    ]

    # Branch boxes
    br_x = 4.0
    bw_b, bh_b = 2.2, 0.7
    leaf_x = 7.0
    lw, lh = 2.6, 0.55
    branch_ys = [1.4, 3.45, 5.5]

    # connector root -> trunk
    root_right_x = root_x + rw
    trunk_x = (root_right_x + br_x) / 2
    root_center_y = root_y + rh / 2
    c.line(x1=root_right_x, y1=root_center_y, x2=trunk_x, y2=root_center_y,
           color="grey_700", width=1.2)
    c.line(x1=trunk_x, y1=branch_ys[0] + bh_b / 2,
           x2=trunk_x, y2=branch_ys[-1] + bh_b / 2,
           color="grey_700", width=1.2)

    for idx, (by, (bname, leaves)) in enumerate(zip(branch_ys, branches)):
        # Branch line
        c.line(x1=trunk_x, y1=by + bh_b / 2, x2=br_x, y2=by + bh_b / 2,
               color="grey_700", width=1.2)
        # Branch box
        c.box(x=br_x, y=by, w=bw_b, h=bh_b,
              fill="accent", border=None)
        c.text(bname, x=br_x, y=by, w=bw_b, h=bh_b,
               size=11, bold=True, color="white", align="center", anchor="middle")

        # Leaf trunk (from branch to middle of leaves)
        branch_right_x = br_x + bw_b
        leaf_trunk_x = (branch_right_x + leaf_x) / 2
        # Compute leaf y positions (stacked close to branch y)
        n = len(leaves)
        leaf_spacing = 0.7
        first_leaf_y = by + bh_b / 2 - ((n - 1) * leaf_spacing) / 2 - lh / 2
        leaf_ys = [first_leaf_y + k * leaf_spacing for k in range(n)]

        c.line(x1=branch_right_x, y1=by + bh_b / 2,
               x2=leaf_trunk_x, y2=by + bh_b / 2,
               color="grey_400", width=1.0)
        c.line(x1=leaf_trunk_x, y1=leaf_ys[0] + lh / 2,
               x2=leaf_trunk_x, y2=leaf_ys[-1] + lh / 2,
               color="grey_400", width=1.0)

        for (lname, score), ly in zip(leaves, leaf_ys):
            c.line(x1=leaf_trunk_x, y1=ly + lh / 2,
                   x2=leaf_x, y2=ly + lh / 2,
                   color="grey_400", width=1.0)
            c.box(x=leaf_x, y=ly, w=lw, h=lh,
                  fill="grey_100", border=0.75, border_color="grey_400")
            c.text(lname, x=leaf_x + 0.15, y=ly, w=lw - 0.7, h=lh,
                   size=9, color="grey_900", anchor="middle")
            # Harvey ball score on the right of leaf
            c.text(HARVEY[score], x=leaf_x + lw - 0.55, y=ly, w=0.5, h=lh,
                   size=14, bold=True, color="accent",
                   align="center", anchor="middle")

    c.text("영향도: ● 최고 / ◕ 높음 / ◑ 중간 / ◔ 낮음 / ○ 미미",
           x=0.4, y=6.85, w=9.2, h=0.3, size=8, color="grey_700", align="center")


# ============================================================
# B8. comp_kanban_3col — To Do / Doing / Done
# ============================================================
def slide_b8_kanban_3col(prs):
    s, c = make(prs, "칸반 보드 — Doing 4건, WIP 한계(5) 이내 안정 운영",
                "Three-column Kanban with task cards")

    columns = [
        ("To Do", "grey_700", False, [
            ("API 스펙 확정", "BA · Due 4/18"),
            ("보안 정책 리뷰", "Ops · Due 4/20"),
            ("회귀 테스트 계획", "QA · Due 4/22"),
            ("UAT 대상자 섭외", "PM · Due 4/25"),
        ]),
        ("Doing", "accent", True, [
            ("로그인 모듈 구현", "Dev · 진행 70%"),
            ("성능 벤치마크", "QA · 진행 40%"),
            ("매뉴얼 v1 작성", "BA · 진행 55%"),
            ("인프라 설정", "Ops · 진행 30%"),
        ]),
        ("Done", "grey_700", False, [
            ("요구사항 명세서", "BA · 4/02"),
            ("아키텍처 리뷰", "Dev · 4/05"),
            ("계약 체결", "PM · 4/08"),
        ]),
    ]

    col_x0 = 0.4
    col_w = 3.05
    col_gap = 0.15
    col_y = 1.3
    col_h = 5.7
    header_h = 0.55

    for i, (title_c, color_c, highlight, cards) in enumerate(columns):
        cx = col_x0 + i * (col_w + col_gap)

        # Column background
        c.box(x=cx, y=col_y, w=col_w, h=col_h,
              fill="grey_100", border=0.75, border_color="grey_200")

        # Column header
        c.box(x=cx, y=col_y, w=col_w, h=header_h,
              fill=color_c, border=None)
        count_label = f"{title_c}  ({len(cards)})"
        c.text(count_label, x=cx + 0.15, y=col_y, w=col_w - 0.3, h=header_h,
               size=12, bold=True, color="white", anchor="middle")

        # Cards
        card_h = 0.85
        card_gap = 0.15
        card_y0 = col_y + header_h + 0.2
        for k, (ctitle, csub) in enumerate(cards):
            cy = card_y0 + k * (card_h + card_gap)
            # Card
            border_col = "accent" if highlight else "grey_400"
            border_w = 1.5 if highlight else 0.75
            c.box(x=cx + 0.15, y=cy, w=col_w - 0.3, h=card_h,
                  fill="white", border=border_w, border_color=border_col)
            # Left accent stripe for Doing
            if highlight:
                c.box(x=cx + 0.15, y=cy, w=0.08, h=card_h,
                      fill="accent", border=None)
            # Title + subtitle
            text_offset = 0.28 if highlight else 0.15
            c.text(ctitle, x=cx + 0.15 + text_offset, y=cy + 0.1,
                   w=col_w - 0.35 - text_offset, h=0.35,
                   size=10, bold=True, color="grey_900")
            c.text(csub, x=cx + 0.15 + text_offset, y=cy + 0.45,
                   w=col_w - 0.35 - text_offset, h=0.3,
                   size=8, color="grey_700")


# ============================================================
# Evaluate
# ============================================================
def evaluate_and_report(pptx_path: Path) -> dict:
    try:
        from ppt_builder.evaluate import evaluate_pptx, print_report
        report = evaluate_pptx(str(pptx_path))
        print_report(report)
        return report
    except Exception as e:
        print(f"(evaluate skipped: {e})")
        return {}


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_b1_harvey_balls(prs)
    slide_b2_raci_matrix(prs)
    slide_b3_org_chart(prs)
    slide_b4_stakeholder_map(prs)
    slide_b5_raid_log(prs)
    slide_b6_workstream_progress(prs)
    slide_b7_issue_tree_scored(prs)
    slide_b8_kanban_3col(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path} ({len(prs.slides)}장)")

    # Evaluate
    evaluate_and_report(pptx_path)

    # PDF via PowerPoint COM
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        pdf_path = pptx_path.with_suffix(".pdf")
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        p = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        p.SaveAs(str(pdf_path.resolve()), 32)
        p.Close()
        print(f"PDF:  {pdf_path}")
    except Exception as e:
        print(f"(PDF export skipped: {e})")

    # PNGs
    try:
        from ppt_builder.track_c.png_export import pptx_to_pngs
        png_dir = OUTPUT_DIR / f"{NAME}_pngs"
        paths = pptx_to_pngs(pptx_path, png_dir)
        print(f"PNGs: {png_dir} ({len(paths)}장)")
    except Exception as e:
        print(f"(PNG export skipped: {e})")


if __name__ == "__main__":
    main()
