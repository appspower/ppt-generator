"""Tier C — Strategy Framework Components (10 slides).

C1. comp_swot_matrix         — 2x2 SWOT with diagonal accent lines
C2. comp_porter_5_forces     — Central industry box + 4 directional forces
C3. comp_bcg_matrix          — Growth/Share 2x2 with product bubbles
C4. comp_ansoff_matrix       — Product/Market 2x2 with risk color
C5. comp_mckinsey_7s         — Central Shared Values + 6 surrounding
C6. comp_magic_quadrant      — Gartner Leaders/Challengers/Niche/Visionaries
C7. comp_value_chain_porter  — Primary + support activities + margin
C8. comp_horizon_3           — H1/H2/H3 overlapping curves
C9. comp_pestel              — 2x3 PESTEL factor grid
C10. comp_balanced_scorecard — Central vision + 4 perspectives

Output: output/tier_c_strategy.pptx (+ PDF + PNGs)
"""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "tier_c_strategy"


# ============================================================
# Slide header helper
# ============================================================
def make(prs, title_text: str, subtitle: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    c.text(title_text, x=0.4, y=0.22, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.66, w=9.2, h=0.25,
               size=9, color="grey_700")
    c.line(x1=0.4, y1=0.96, x2=9.6, y2=0.96, color="grey_200", width=1.0)
    return s, c


# ============================================================
# C1. comp_swot_matrix
# ============================================================
def comp_swot_matrix(c: Canvas, *, x: float, y: float, w: float, h: float,
                     strengths, weaknesses, opportunities, threats):
    cell_w = w / 2
    cell_h = h / 2
    cells = [
        ("STRENGTHS (강점)", strengths, "grey_700", 0, 0),
        ("WEAKNESSES (약점)", weaknesses, "negative", 1, 0),
        ("OPPORTUNITIES (기회)", opportunities, "accent", 0, 1),
        ("THREATS (위협)", threats, "grey_900", 1, 1),
    ]
    for title, bullets, band_color, col, row in cells:
        cx = x + col * cell_w
        cy = y + row * cell_h
        c.box(x=cx, y=cy, w=cell_w, h=cell_h,
              fill="white", border=1.0, border_color="grey_mid")
        # colored header band
        c.box(x=cx, y=cy, w=cell_w, h=0.38,
              fill=band_color, border=None)
        c.text(title, x=cx + 0.12, y=cy + 0.06, w=cell_w - 0.24, h=0.28,
               size=10, bold=True, color="white")
        # bullets
        body = "\n".join(f"▪  {b}" for b in bullets)
        c.text(body, x=cx + 0.18, y=cy + 0.48, w=cell_w - 0.32, h=cell_h - 0.55,
               size=8.5, color="grey_900")
    # Diagonal accent lines from center
    cx_center = x + cell_w
    cy_center = y + cell_h
    c.line(x1=x, y1=y, x2=x + w, y2=y + h,
           color="accent_light", width=0.5)
    c.line(x1=x + w, y1=y, x2=x, y2=y + h,
           color="accent_light", width=0.5)
    # Vertical + horizontal divider emphasized
    c.line(x1=cx_center, y1=y, x2=cx_center, y2=y + h,
           color="grey_900", width=1.5)
    c.line(x1=x, y1=cy_center, x2=x + w, y2=cy_center,
           color="grey_900", width=1.5)


def slide_c1_swot(prs):
    s, c = make(prs, "SWOT Matrix — HD현대 ERP 전환 전략 포지셔닝",
                "내부 역량(S/W) × 외부 환경(O/T) 기반 전략 도출")
    comp_swot_matrix(
        c, x=0.5, y=1.2, w=9.0, h=5.8,
        strengths=[
            "SAP 도입 20년 노하우 축적",
            "전사 표준 프로세스 정립",
            "현업 Key User 800명 확보",
            "IT 거버넌스 성숙도 Level 4",
        ],
        weaknesses=[
            "레거시 ECC 커스터마이징 과도 (4,200 Z오브젝트)",
            "부서간 마스터 데이터 불일치 심각",
            "Cloud 운영 역량 부족",
            "변화 관리 조직 미비",
        ],
        opportunities=[
            "S/4HANA Cloud 2028 의무 전환",
            "그룹사 통합 ERP 시너지 $120M",
            "생성형 AI Joule 무상 제공",
            "ESG 공시 자동화 요구 증가",
        ],
        threats=[
            "ECC 유지보수 2030 종료 리스크",
            "글로벌 경쟁사 3년 앞서 전환 완료",
            "Clean Core 원칙 현업 저항",
            "라이선스 비용 연 18% 상승",
        ],
    )
    c.text("시사점: 강점(O)×기회(O) 영역에 집중, 약점(W)을 위협(T)이 되기 전에 선제 해결",
           x=0.5, y=7.1, w=9.0, h=0.3, size=9, bold=True, color="accent")


# ============================================================
# C2. comp_porter_5_forces
# ============================================================
def comp_porter_5_forces(c: Canvas, *, x: float, y: float, w: float, h: float,
                         center_title, new_entrants, suppliers, buyers, substitutes):
    # Center box
    cw, ch = 2.6, 1.2
    cx = x + (w - cw) / 2
    cy = y + (h - ch) / 2
    c.box(x=cx, y=cy, w=cw, h=ch, fill="accent", border=None)
    c.text(center_title["title"], x=cx, y=cy + 0.15, w=cw, h=0.35,
           size=12, bold=True, color="white", align="center")
    c.text(center_title["detail"], x=cx, y=cy + 0.55, w=cw, h=0.55,
           size=8, color="white", align="center")

    # 4 peripheral boxes (N / S / E / W)
    box_w, box_h = 2.4, 1.1
    top_x = x + (w - box_w) / 2
    top_y = y
    bot_x = top_x
    bot_y = y + h - box_h
    left_x = x
    left_y = y + (h - box_h) / 2
    right_x = x + w - box_w
    right_y = left_y

    positions = [
        ("neighbor_of 신규 진입자", new_entrants, top_x, top_y, "N"),
        ("공급자 교섭력", suppliers, left_x, left_y, "W"),
        ("구매자 교섭력", buyers, right_x, right_y, "E"),
        ("대체재 위협", substitutes, bot_x, bot_y, "S"),
    ]
    for title, data, bx, by, _ in positions:
        c.box(x=bx, y=by, w=box_w, h=box_h, fill="white",
              border=1.5, border_color="grey_900")
        c.text(data["title"], x=bx + 0.1, y=by + 0.08, w=box_w - 0.2, h=0.28,
               size=10, bold=True, color="accent")
        body = "\n".join(f"•  {b}" for b in data["bullets"])
        c.text(body, x=bx + 0.12, y=by + 0.38, w=box_w - 0.24, h=box_h - 0.46,
               size=8, color="grey_900")

    # Arrows pointing INTO center (4 directions)
    # Top → center top
    c.arrow(x1=top_x + box_w / 2, y1=top_y + box_h + 0.05,
            x2=cx + cw / 2, y2=cy - 0.05, color="accent", width=2.0)
    # Bottom → center bottom
    c.arrow(x1=bot_x + box_w / 2, y1=bot_y - 0.05,
            x2=cx + cw / 2, y2=cy + ch + 0.05, color="accent", width=2.0)
    # Left → center left
    c.arrow(x1=left_x + box_w + 0.05, y1=left_y + box_h / 2,
            x2=cx - 0.05, y2=cy + ch / 2, color="accent", width=2.0)
    # Right → center right
    c.arrow(x1=right_x - 0.05, y1=right_y + box_h / 2,
            x2=cx + cw + 0.05, y2=cy + ch / 2, color="accent", width=2.0)


def slide_c2_porter(prs):
    s, c = make(prs, "Porter's 5 Forces — 국내 ERP SI 산업 구조 분석",
                "HD현대 기준 산업 내 경쟁 강도 High")
    comp_porter_5_forces(
        c, x=0.5, y=1.2, w=9.0, h=5.9,
        center_title={"title": "산업 내 경쟁", "detail": "Samsung SDS / LG CNS / SK C&C\n경쟁 강도: High"},
        new_entrants={"title": "신규 진입자 (Medium)", "bullets": [
            "글로벌 Accenture 한국 확장",
            "Cloud-native 스타트업 부상",
        ]},
        suppliers={"title": "공급자 교섭력 (High)", "bullets": [
            "SAP 라이선스 독점 구조",
            "Hyperscaler 3사 과점",
        ]},
        buyers={"title": "구매자 교섭력 (High)", "bullets": [
            "대기업 고객 집중도 상위 20社",
            "RFP 단가 경쟁 심화",
        ]},
        substitutes={"title": "대체재 위협 (Medium)", "bullets": [
            "Oracle Fusion / MS Dynamics",
            "산업특화 Best-of-breed SaaS",
        ]},
    )


# ============================================================
# C3. comp_bcg_matrix
# ============================================================
def comp_bcg_matrix(c: Canvas, *, x: float, y: float, w: float, h: float,
                    products):
    # Frame
    mx, my = 0.7, 0.35  # label margins
    inner_x = x + mx
    inner_y = y + my
    inner_w = w - mx - 0.3
    inner_h = h - my - 0.7

    cell_w = inner_w / 2
    cell_h = inner_h / 2

    # Quadrant fills
    quads = [
        ("Question Marks", "grey_400", 0, 0),   # top-left
        ("Stars", "accent", 1, 0),              # top-right
        ("Dogs", "grey_200", 0, 1),             # bottom-left
        ("Cash Cows", "grey_700", 1, 1),        # bottom-right
    ]
    for name, fill, col, row in quads:
        qx = inner_x + col * cell_w
        qy = inner_y + row * cell_h
        c.box(x=qx, y=qy, w=cell_w, h=cell_h,
              fill=fill, border=1.0, border_color="grey_900")
        text_color = "white" if fill in ("accent", "grey_700", "grey_900") else "grey_900"
        c.text(name, x=qx + 0.1, y=qy + 0.1, w=cell_w - 0.2, h=0.3,
               size=12, bold=True, color=text_color)

    # Axis labels
    c.text("시장 성장률 (높음)", x=x, y=inner_y - 0.05, w=mx - 0.05, h=0.3,
           size=8, bold=True, color="grey_900", align="right")
    c.text("시장 성장률 (낮음)", x=x, y=inner_y + inner_h - 0.3, w=mx - 0.05, h=0.3,
           size=8, bold=True, color="grey_900", align="right")
    c.text("← 상대적 시장 점유율 (낮음)                     (높음) →",
           x=inner_x, y=inner_y + inner_h + 0.15, w=inner_w, h=0.3,
           size=9, bold=True, color="grey_900", align="center")
    # Axis arrows
    c.arrow(x1=inner_x - 0.1, y1=inner_y + inner_h, x2=inner_x - 0.1, y2=inner_y,
            color="grey_900", width=1.5)
    c.arrow(x1=inner_x, y1=inner_y + inner_h + 0.1,
            x2=inner_x + inner_w, y2=inner_y + inner_h + 0.1,
            color="grey_900", width=1.5)

    # Place product bubbles: product = dict(name, share 0-1, growth 0-1, size_rel)
    for p in products:
        # share: 0 (left) ~ 1 (right)
        # growth: 0 (bottom) ~ 1 (top)
        px = inner_x + p["share"] * inner_w
        py = inner_y + (1 - p["growth"]) * inner_h
        d = 0.35 + p.get("size_rel", 0.5) * 0.55
        c.circle(x=px - d / 2, y=py - d / 2, d=d,
                 fill="white", border=1.5, border_color="grey_900",
                 text=p["abbr"], text_color="grey_900", text_size=9)
        c.text(p["name"], x=px - 0.8, y=py + d / 2 + 0.02, w=1.6, h=0.2,
               size=7, color="grey_700", align="center")


def slide_c3_bcg(prs):
    s, c = make(prs, "BCG Growth-Share Matrix — HD현대 사업 포트폴리오",
                "시장 성장률 × 상대 점유율 기준 자원 배분 의사결정")
    comp_bcg_matrix(
        c, x=0.5, y=1.2, w=9.0, h=5.7,
        products=[
            {"abbr": "SH", "name": "조선(상선)", "share": 0.82, "growth": 0.78, "size_rel": 1.0},    # Star
            {"abbr": "OF", "name": "해양플랜트", "share": 0.18, "growth": 0.72, "size_rel": 0.6},     # Question
            {"abbr": "EN", "name": "엔진기계", "share": 0.88, "growth": 0.28, "size_rel": 0.85},     # Cash cow
            {"abbr": "RB", "name": "로봇사업", "share": 0.25, "growth": 0.85, "size_rel": 0.35},     # Question
            {"abbr": "CS", "name": "건설기계(국내)", "share": 0.15, "growth": 0.15, "size_rel": 0.4}, # Dog
            {"abbr": "GR", "name": "친환경선", "share": 0.70, "growth": 0.90, "size_rel": 0.7},       # Star
        ],
    )
    c.text("전략: Stars(조선·친환경선) 추가 투자, Cash Cow(엔진) 수익 확보, Dog 사업 재검토",
           x=0.5, y=7.1, w=9.0, h=0.3, size=9, bold=True, color="accent")


# ============================================================
# C4. comp_ansoff_matrix
# ============================================================
def comp_ansoff_matrix(c: Canvas, *, x: float, y: float, w: float, h: float,
                       quadrants):
    mx, my = 1.1, 0.35
    inner_x = x + mx
    inner_y = y + my
    inner_w = w - mx - 0.3
    inner_h = h - my - 0.6
    cw = inner_w / 2
    ch = inner_h / 2

    # quadrants = list of 4: (title, risk, desc, col, row, risk_color)
    risk_map = {"Low": "positive", "Medium": "warning",
                "High": "negative", "Very High": "grey_900"}

    for q in quadrants:
        qx = inner_x + q["col"] * cw
        qy = inner_y + q["row"] * ch
        c.box(x=qx, y=qy, w=cw, h=ch,
              fill="white", border=1.0, border_color="grey_mid")
        # Risk chip
        rc = risk_map[q["risk"]]
        c.box(x=qx, y=qy, w=cw, h=0.42, fill=rc, border=None)
        c.text(q["title"], x=qx + 0.12, y=qy + 0.06, w=cw - 1.3, h=0.30,
               size=11, bold=True, color="white")
        c.text(f"Risk: {q['risk']}", x=qx + cw - 1.3, y=qy + 0.08, w=1.2, h=0.28,
               size=8.5, bold=True, color="white", align="right")
        c.text(q["desc"], x=qx + 0.15, y=qy + 0.52, w=cw - 0.3, h=ch - 0.6,
               size=9, color="grey_900")

    # Axis labels
    c.text("기존\n제품", x=x + 0.05, y=inner_y + 0.15, w=mx - 0.1, h=ch - 0.3,
           size=10, bold=True, color="grey_900", align="center", anchor="middle")
    c.text("신규\n제품", x=x + 0.05, y=inner_y + ch + 0.15, w=mx - 0.1, h=ch - 0.3,
           size=10, bold=True, color="grey_900", align="center", anchor="middle")
    c.text("기존 시장", x=inner_x, y=inner_y + inner_h + 0.1, w=cw, h=0.3,
           size=10, bold=True, color="grey_900", align="center")
    c.text("신규 시장", x=inner_x + cw, y=inner_y + inner_h + 0.1, w=cw, h=0.3,
           size=10, bold=True, color="grey_900", align="center")


def slide_c4_ansoff(prs):
    s, c = make(prs, "Ansoff Matrix — 성장 전략 옵션 (제품 × 시장)",
                "기존/신규 × 제품/시장 조합별 성장 경로와 리스크")
    comp_ansoff_matrix(
        c, x=0.5, y=1.2, w=9.0, h=5.7,
        quadrants=[
            {"title": "Market Penetration", "risk": "Low",
             "desc": "현 고객 대상 점유율 확대.\n가격 경쟁력·고객 충성 프로그램 강화.\nTarget: 국내 대형 조선 시장 점유율 +5%p",
             "col": 0, "row": 0},
            {"title": "Product Development", "risk": "Medium",
             "desc": "기존 고객에게 신제품 공급.\n친환경 LNG·메탄올 추진선 등 신규 라인업.\n3년내 신제품 매출 $2.1B 목표",
             "col": 0, "row": 1},
            {"title": "Market Development", "risk": "High",
             "desc": "현 제품으로 신규 지역/세그먼트 진입.\n중동·인도 해양 플랜트 신규 수주 추진.\n현지 JV 파트너 2곳 확보",
             "col": 1, "row": 0},
            {"title": "Diversification", "risk": "Very High",
             "desc": "신제품 + 신시장 동시 진출.\n로봇·수소 선박·원전 SMR 사업 진출.\nR&D 투자 연 $450M, 성공률 18%",
             "col": 1, "row": 1},
        ],
    )


# ============================================================
# C5. comp_mckinsey_7s
# ============================================================
def comp_mckinsey_7s(c: Canvas, *, x: float, y: float, w: float, h: float,
                     center, outer):
    cx = x + w / 2
    cy = y + h / 2
    # Connecting lines from center to each outer circle (draw first so beneath)
    radius = min(w, h) * 0.36
    outer_d = 1.35
    center_d = 1.55

    positions = []
    for i in range(6):
        angle = -math.pi / 2 + i * (2 * math.pi / 6)  # start top
        ox = cx + radius * math.cos(angle)
        oy = cy + radius * math.sin(angle)
        positions.append((ox, oy))

    # Lines (center to each outer)
    for ox, oy in positions:
        c.line(x1=cx, y1=cy, x2=ox, y2=oy, color="grey_400", width=1.0)
    # Interconnections between adjacent outer circles
    for i in range(6):
        ox1, oy1 = positions[i]
        ox2, oy2 = positions[(i + 1) % 6]
        c.line(x1=ox1, y1=oy1, x2=ox2, y2=oy2, color="grey_200", width=0.75)

    # Center circle
    c.circle(x=cx - center_d / 2, y=cy - center_d / 2, d=center_d,
             fill="accent", border=None,
             text=center["title"], text_color="white", text_size=11)
    c.text(center["detail"], x=cx - 1.3, y=cy + center_d / 2 - 0.4, w=2.6, h=0.4,
           size=8, color="white", align="center")

    # Outer circles
    for (ox, oy), item in zip(positions, outer):
        c.circle(x=ox - outer_d / 2, y=oy - outer_d / 2, d=outer_d,
                 fill="white", border=1.5, border_color="grey_900",
                 text=item["title"], text_color="accent", text_size=10)
        # detail below the circle name
        c.text(item["detail"], x=ox - outer_d / 2 - 0.2, y=oy + 0.05, w=outer_d + 0.4, h=0.5,
               size=7.5, color="grey_700", align="center")


def slide_c5_7s(prs):
    s, c = make(prs, "McKinsey 7S — ERP 전환 조직 정렬 진단",
                "Hard(전략·조직·시스템) + Soft(공유가치·스타일·인력·스킬) 정합성 평가")
    comp_mckinsey_7s(
        c, x=0.5, y=1.0, w=9.0, h=6.2,
        center={"title": "Shared\nValues",
                "detail": "고객가치 최우선\nClean Core 원칙"},
        outer=[
            {"title": "Strategy", "detail": "S/4HANA\n2027 Go-live"},
            {"title": "Structure", "detail": "PMO+CoE\nHybrid 운영"},
            {"title": "Systems", "detail": "Cloud\nIntegration Suite"},
            {"title": "Style", "detail": "Top-down\n+Agile Squad"},
            {"title": "Staff", "detail": "Key User\n800명 확보"},
            {"title": "Skills", "detail": "Fiori/ABAP\n역량 Re-skill"},
        ],
    )


# ============================================================
# C6. comp_magic_quadrant
# ============================================================
def comp_magic_quadrant(c: Canvas, *, x: float, y: float, w: float, h: float,
                        vendors):
    mx, my = 0.6, 0.35
    inner_x = x + mx
    inner_y = y + my
    inner_w = w - mx - 0.3
    inner_h = h - my - 0.6
    cw = inner_w / 2
    ch = inner_h / 2

    # Quadrants
    c.box(x=inner_x, y=inner_y, w=cw, h=ch,
          fill="grey_100", border=1.0, border_color="grey_mid")
    c.box(x=inner_x + cw, y=inner_y, w=cw, h=ch,
          fill="accent_light", border=1.0, border_color="grey_mid")
    c.box(x=inner_x, y=inner_y + ch, w=cw, h=ch,
          fill="grey_200", border=1.0, border_color="grey_mid")
    c.box(x=inner_x + cw, y=inner_y + ch, w=cw, h=ch,
          fill="grey_100", border=1.0, border_color="grey_mid")

    c.text("CHALLENGERS", x=inner_x + 0.15, y=inner_y + 0.1, w=cw - 0.3, h=0.3,
           size=10, bold=True, color="grey_700")
    c.text("LEADERS", x=inner_x + cw + 0.15, y=inner_y + 0.1, w=cw - 0.3, h=0.3,
           size=10, bold=True, color="accent")
    c.text("NICHE PLAYERS", x=inner_x + 0.15, y=inner_y + ch + 0.1, w=cw - 0.3, h=0.3,
           size=10, bold=True, color="grey_700")
    c.text("VISIONARIES", x=inner_x + cw + 0.15, y=inner_y + ch + 0.1, w=cw - 0.3, h=0.3,
           size=10, bold=True, color="grey_700")

    # Axis
    c.arrow(x1=inner_x - 0.1, y1=inner_y + inner_h, x2=inner_x - 0.1, y2=inner_y,
            color="grey_900", width=1.5)
    c.arrow(x1=inner_x, y1=inner_y + inner_h + 0.1,
            x2=inner_x + inner_w, y2=inner_y + inner_h + 0.1,
            color="grey_900", width=1.5)
    c.text("Ability to Execute →", x=x - 0.1, y=inner_y + 0.2, w=mx - 0.1, h=0.4,
           size=8, bold=True, color="grey_900", align="right")
    c.text("Completeness of Vision →",
           x=inner_x, y=inner_y + inner_h + 0.15, w=inner_w, h=0.3,
           size=9, bold=True, color="grey_900", align="center")

    # Vendor dots
    for v in vendors:
        vx = inner_x + v["vision"] * inner_w
        vy = inner_y + (1 - v["execution"]) * inner_h
        c.circle(x=vx - 0.12, y=vy - 0.12, d=0.24,
                 fill="accent", border=1.0, border_color="grey_900")
        c.text(v["name"], x=vx + 0.15, y=vy - 0.12, w=1.5, h=0.25,
               size=8, bold=True, color="grey_900")


def slide_c6_magic_quadrant(prs):
    s, c = make(prs, "Magic Quadrant — 국내 ERP SI 벤더 경쟁 포지션",
                "비전 완성도 × 실행 역량 2축 평가 (2026년 1Q)")
    comp_magic_quadrant(
        c, x=0.5, y=1.2, w=9.0, h=5.7,
        vendors=[
            {"name": "Samsung SDS", "vision": 0.78, "execution": 0.85},   # Leader
            {"name": "LG CNS", "vision": 0.68, "execution": 0.72},        # Leader/Challenger
            {"name": "SK C&C", "vision": 0.45, "execution": 0.70},        # Challenger
            {"name": "현대오토에버", "vision": 0.58, "execution": 0.55},   # mid
            {"name": "Accenture Korea", "vision": 0.80, "execution": 0.40}, # Visionary
            {"name": "로컬 부티크", "vision": 0.25, "execution": 0.30},    # Niche
        ],
    )


# ============================================================
# C7. comp_value_chain_porter
# ============================================================
def comp_value_chain_porter(c: Canvas, *, x: float, y: float, w: float, h: float,
                            support, primary):
    # Support strip (top)
    support_h = 0.55
    support_w = w - 1.0  # leave 1.0" for margin triangle
    each_sw = support_w / len(support)
    for i, item in enumerate(support):
        sx = x + i * each_sw
        c.box(x=sx, y=y, w=each_sw - 0.02, h=support_h,
              fill="grey_200", border=1.0, border_color="grey_900")
        c.text(item["title"], x=sx + 0.05, y=y + 0.06, w=each_sw - 0.12, h=0.2,
               size=9, bold=True, color="grey_900")
        c.text(item["desc"], x=sx + 0.05, y=y + 0.28, w=each_sw - 0.12, h=0.25,
               size=7, color="grey_700")

    # Primary activities chevrons
    gap_y = 0.25
    chev_y = y + support_h + gap_y
    chev_h = h - support_h - gap_y
    chev_w = support_w / len(primary)
    overlap = 0.15
    colors = ["accent", "accent_mid", "accent", "accent_mid", "accent"]

    for i, (item, fill) in enumerate(zip(primary, colors)):
        cx = x + i * (chev_w - overlap)
        c.chevron(x=cx, y=chev_y, w=chev_w, h=chev_h,
                  fill=fill, text="", text_color="white")
        c.text(item["title"], x=cx + 0.1, y=chev_y + 0.15, w=chev_w - 0.5, h=0.3,
               size=10, bold=True, color="white")
        c.text(item["desc"], x=cx + 0.1, y=chev_y + 0.5, w=chev_w - 0.6, h=chev_h - 0.6,
               size=8, color="white")

    # Margin triangle (right edge, spanning full height)
    tri_x = x + support_w + 0.02
    tri_w = w - support_w - 0.02
    # Use a simple box with accent fill + label
    c.box(x=tri_x, y=y, w=tri_w, h=h,
          fill="grey_900", border=None)
    c.text("M\nA\nR\nG\nI\nN",
           x=tri_x, y=y, w=tri_w, h=h,
           size=12, bold=True, color="white",
           align="center", anchor="middle")


def slide_c7_value_chain(prs):
    s, c = make(prs, "Value Chain Analysis — HD현대 조선사업 가치사슬",
                "주요/지원 활동별 원가 구조 및 차별화 포인트 매핑")
    comp_value_chain_porter(
        c, x=0.5, y=1.2, w=9.0, h=5.7,
        support=[
            {"title": "Infrastructure", "desc": "ERP·PLM·MES 통합 플랫폼"},
            {"title": "HR Management", "desc": "용접·설계 마이스터 양성"},
            {"title": "Technology", "desc": "디지털트윈·AI 검사"},
            {"title": "Procurement", "desc": "그룹 통합 구매, 후판 Hedging"},
        ],
        primary=[
            {"title": "Inbound", "desc": "후판·엔진 부품\n입고 물류"},
            {"title": "Operations", "desc": "설계→가공→조립\n도크 공정"},
            {"title": "Outbound", "desc": "진수·시운전\n인도"},
            {"title": "Marketing & Sales", "desc": "선주 영업\nLife-cycle 계약"},
            {"title": "Service", "desc": "A/S·수리\n친환경 개조"},
        ],
    )


# ============================================================
# C8. comp_horizon_3
# ============================================================
def comp_horizon_3(c: Canvas, *, x: float, y: float, w: float, h: float,
                   horizons):
    # Baseline timeline
    axis_y = y + h - 0.55
    c.line(x1=x + 0.5, y1=axis_y, x2=x + w - 0.3, y2=axis_y,
           color="grey_900", width=1.5)
    c.arrow(x1=x + w - 0.3, y1=axis_y, x2=x + w - 0.1, y2=axis_y,
            color="grey_900", width=1.5)
    # Timeline labels
    ticks = ["Now", "+1Y", "+2Y", "+3Y", "+5Y"]
    tick_w = (w - 0.8) / (len(ticks) - 1)
    for i, label in enumerate(ticks):
        tx = x + 0.5 + i * tick_w
        c.line(x1=tx, y1=axis_y - 0.05, x2=tx, y2=axis_y + 0.05,
               color="grey_900", width=1.0)
        c.text(label, x=tx - 0.4, y=axis_y + 0.1, w=0.8, h=0.25,
               size=8, color="grey_900", align="center")
    c.text("Time →", x=x + w - 0.9, y=axis_y + 0.35, w=0.8, h=0.25,
           size=8, bold=True, color="grey_900")

    # Value axis label
    c.text("Value\n↑", x=x - 0.1, y=y + 0.1, w=0.5, h=0.5,
           size=8, bold=True, color="grey_900")

    # Three horizon bands: approximated by stacked boxes giving "rising" feel
    # H1: tall left, tapers right (accent, most solid)
    # H2: low-mid, peaks mid-right
    # H3: low near term, rising late
    colors = ["accent", "accent_mid", "accent_light"]

    # H1: rectangular band left 50%, tapered via trapezoid-like shape (we fake with two boxes)
    h1_w = w * 0.55
    h1_h = 1.3
    c.box(x=x + 0.3, y=axis_y - h1_h, w=h1_w * 0.75, h=h1_h,
          fill=colors[0], border=None)
    # tapered tail
    c.box(x=x + 0.3 + h1_w * 0.75, y=axis_y - h1_h + 0.4, w=h1_w * 0.25, h=h1_h - 0.4,
          fill=colors[0], border=None)

    # H2: emerging — mid area rising
    h2_start_x = x + w * 0.28
    h2_w = w * 0.55
    h2_h = 1.7
    # stepped rising
    c.box(x=h2_start_x, y=axis_y - h2_h * 0.35, w=h2_w * 0.35, h=h2_h * 0.35,
          fill=colors[1], border=None)
    c.box(x=h2_start_x + h2_w * 0.35, y=axis_y - h2_h * 0.75, w=h2_w * 0.35, h=h2_h * 0.75,
          fill=colors[1], border=None)
    c.box(x=h2_start_x + h2_w * 0.70, y=axis_y - h2_h, w=h2_w * 0.30, h=h2_h,
          fill=colors[1], border=None)

    # H3: new — rising late
    h3_start_x = x + w * 0.55
    h3_w = w * 0.40
    h3_h = 2.2
    c.box(x=h3_start_x, y=axis_y - h3_h * 0.20, w=h3_w * 0.35, h=h3_h * 0.20,
          fill=colors[2], border=None)
    c.box(x=h3_start_x + h3_w * 0.35, y=axis_y - h3_h * 0.55, w=h3_w * 0.30, h=h3_h * 0.55,
          fill=colors[2], border=None)
    c.box(x=h3_start_x + h3_w * 0.65, y=axis_y - h3_h, w=h3_w * 0.35, h=h3_h,
          fill=colors[2], border=None)

    # Labels on each horizon
    # H1 label
    c.text("Horizon 1 — Extend & Defend", x=x + 0.5, y=axis_y - 1.25, w=3.5, h=0.28,
           size=11, bold=True, color="white")
    c.text(horizons["h1"], x=x + 0.5, y=axis_y - 0.95, w=3.5, h=0.6,
           size=8, color="white")
    # H2 label
    c.text("Horizon 2 — Build Emerging", x=x + w * 0.35, y=axis_y - 1.65, w=3.5, h=0.28,
           size=11, bold=True, color="white")
    c.text(horizons["h2"], x=x + w * 0.35, y=axis_y - 1.35, w=3.5, h=0.6,
           size=8, color="white")
    # H3 label
    c.text("Horizon 3 — Create New", x=x + w * 0.58, y=axis_y - 2.15, w=3.5, h=0.28,
           size=11, bold=True, color="grey_900")
    c.text(horizons["h3"], x=x + w * 0.58, y=axis_y - 1.85, w=3.5, h=0.6,
           size=8, color="grey_900")


def slide_c8_horizon(prs):
    s, c = make(prs, "Three Horizons of Growth — HD현대 성장 로드맵",
                "H1(Core) → H2(Adjacent) → H3(Transformational) 단계별 투자 배분")
    comp_horizon_3(
        c, x=0.5, y=1.2, w=9.0, h=5.6,
        horizons={
            "h1": "기존 조선·엔진 사업 수익성 유지.\n원가 혁신 연 -3%, OPEX 최적화.",
            "h2": "친환경 선박·스마트 야드 확장.\n2028년 매출 비중 35% 목표.",
            "h3": "수소 선박·SMR·로봇 신사업.\n10년 내 신규 매출 $15B 창출.",
        },
    )


# ============================================================
# C9. comp_pestel
# ============================================================
def comp_pestel(c: Canvas, *, x: float, y: float, w: float, h: float, factors):
    cols, rows = 3, 2
    cell_w = w / cols
    cell_h = h / rows
    for i, f in enumerate(factors):
        col = i % cols
        row = i // cols
        fx = x + col * cell_w
        fy = y + row * cell_h
        c.box(x=fx + 0.05, y=fy + 0.05, w=cell_w - 0.1, h=cell_h - 0.1,
              fill="white", border=1.0, border_color="grey_mid")
        # Letter badge + title
        c.box(x=fx + 0.05, y=fy + 0.05, w=0.5, h=0.5,
              fill="accent", border=None)
        c.text(f["letter"], x=fx + 0.05, y=fy + 0.05, w=0.5, h=0.5,
               size=18, bold=True, color="white", align="center", anchor="middle")
        c.text(f["title"], x=fx + 0.6, y=fy + 0.12, w=cell_w - 0.7, h=0.35,
               size=12, bold=True, color="accent")
        # Bullets
        body = "\n".join(f"•  {b}" for b in f["bullets"])
        c.text(body, x=fx + 0.2, y=fy + 0.68, w=cell_w - 0.3, h=cell_h - 0.8,
               size=8.5, color="grey_900")


def slide_c9_pestel(prs):
    s, c = make(prs, "PESTEL Analysis — HD현대 거시환경 분석",
                "정치·경제·사회·기술·환경·법률 6대 영역 Macro Risk Scan")
    comp_pestel(
        c, x=0.3, y=1.1, w=9.4, h=5.9,
        factors=[
            {"letter": "P", "title": "Political", "bullets": [
                "미·중 탈동조화 가속",
                "IRA / CRMA 공급망 재편",
                "중동 지정학 리스크 상존",
            ]},
            {"letter": "E", "title": "Economic", "bullets": [
                "한국 기준금리 3.50%",
                "원·달러 환율 1,380원",
                "후판 가격 +12% YoY",
            ]},
            {"letter": "S", "title": "Social", "bullets": [
                "조선 기능인력 고령화",
                "외국인 근로자 확대 정책",
                "ESG 가치 소비 확산",
            ]},
            {"letter": "T", "title": "Technological", "bullets": [
                "Gen-AI Joule 전사 확산",
                "디지털트윈 상용화",
                "자율운항 선박 상용 승인",
            ]},
            {"letter": "E", "title": "Environmental", "bullets": [
                "IMO 2030 온실가스 -40%",
                "EU ETS 해운 포함",
                "메탄올·암모니아 연료 의무화",
            ]},
            {"letter": "L", "title": "Legal", "bullets": [
                "중대재해처벌법 강화",
                "개인정보보호법(PIPA) 개정",
                "美 Jones Act 조선 보조 검토",
            ]},
        ],
    )


# ============================================================
# C10. comp_balanced_scorecard
# ============================================================
def comp_balanced_scorecard(c: Canvas, *, x: float, y: float, w: float, h: float,
                            vision, perspectives):
    cx = x + w / 2
    cy = y + h / 2

    cw, ch = 2.4, 1.2
    # Center
    c.box(x=cx - cw / 2, y=cy - ch / 2, w=cw, h=ch,
          fill="accent", border=None)
    c.text(vision["title"], x=cx - cw / 2, y=cy - ch / 2 + 0.12, w=cw, h=0.4,
           size=13, bold=True, color="white", align="center")
    c.text(vision["detail"], x=cx - cw / 2, y=cy - ch / 2 + 0.55, w=cw, h=0.6,
           size=8.5, color="white", align="center")

    # 4 perspective boxes at N/S/E/W
    box_w, box_h = 3.0, 2.0
    positions = {
        "top": (cx - box_w / 2, y),
        "bottom": (cx - box_w / 2, y + h - box_h),
        "left": (x, cy - box_h / 2),
        "right": (x + w - box_w, cy - box_h / 2),
    }

    for key, data in perspectives.items():
        bx, by = positions[key]
        c.box(x=bx, y=by, w=box_w, h=box_h,
              fill="white", border=1.5, border_color="grey_900")
        c.box(x=bx, y=by, w=box_w, h=0.38, fill="grey_900", border=None)
        c.text(data["title"], x=bx + 0.15, y=by + 0.06, w=box_w - 0.3, h=0.28,
               size=11, bold=True, color="white")

        # Objectives / Measures / Targets
        col_w = (box_w - 0.3) / 3
        col_y = by + 0.48
        headers = [("Objectives", data["objectives"]),
                   ("Measures", data["measures"]),
                   ("Targets", data["targets"])]
        for i, (hd, items) in enumerate(headers):
            cx2 = bx + 0.15 + i * col_w
            c.text(hd, x=cx2, y=col_y, w=col_w, h=0.22,
                   size=8, bold=True, color="accent")
            body = "\n".join(f"• {it}" for it in items)
            c.text(body, x=cx2, y=col_y + 0.24, w=col_w, h=box_h - 0.75,
                   size=7, color="grey_900")

    # Connect center to each perspective
    c.line(x1=cx, y1=y + box_h, x2=cx, y2=cy - ch / 2, color="grey_900", width=1.0)
    c.line(x1=cx, y1=cy + ch / 2, x2=cx, y2=y + h - box_h, color="grey_900", width=1.0)
    c.line(x1=x + box_w, y1=cy, x2=cx - cw / 2, y2=cy, color="grey_900", width=1.0)
    c.line(x1=cx + cw / 2, y1=cy, x2=x + w - box_w, y2=cy, color="grey_900", width=1.0)


def slide_c10_bsc(prs):
    s, c = make(prs, "Balanced Scorecard — HD현대 전략 실행 대시보드",
                "4대 관점별 목표·측정지표·타깃 (BSC 2026)")
    comp_balanced_scorecard(
        c, x=0.3, y=1.05, w=9.4, h=6.1,
        vision={"title": "Vision & Strategy",
                "detail": "Global Top-tier\nClean Shipbuilder\n2030"},
        perspectives={
            "top": {
                "title": "Financial",
                "objectives": ["매출 성장", "수익성 개선", "현금흐름"],
                "measures": ["매출액", "영업이익률", "FCF"],
                "targets": ["$42B", "9.5%", "+$1.8B"],
            },
            "bottom": {
                "title": "Learning & Growth",
                "objectives": ["핵심인재", "디지털 역량", "조직 몰입"],
                "measures": ["Retention", "AI 교육 이수", "eNPS"],
                "targets": [">95%", "100%", "+42"],
            },
            "left": {
                "title": "Customer",
                "objectives": ["고객 만족", "브랜드 가치", "신규 고객"],
                "measures": ["CSAT", "Brand Index", "신규 선주"],
                "targets": [">88", "Top 3", "+12社"],
            },
            "right": {
                "title": "Internal Process",
                "objectives": ["공정 혁신", "품질", "납기"],
                "measures": ["Dock 회전율", "불량률", "On-time"],
                "targets": ["+8%", "<0.5%", ">98%"],
            },
        },
    )


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_c1_swot(prs)
    slide_c2_porter(prs)
    slide_c3_bcg(prs)
    slide_c4_ansoff(prs)
    slide_c5_7s(prs)
    slide_c6_magic_quadrant(prs)
    slide_c7_value_chain(prs)
    slide_c8_horizon(prs)
    slide_c9_pestel(prs)
    slide_c10_bsc(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX: {pptx_path} ({len(prs.slides)} slides)")

    # PDF via PowerPoint COM
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        pdf_path = pptx_path.with_suffix(".pdf")
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        p = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        p.SaveAs(str(pdf_path.resolve()), 32)
        p.Close()
        print(f"PDF:  {pdf_path}")
    except Exception as e:
        print(f"PDF generation skipped: {e}")

    # PNGs
    try:
        from ppt_builder.track_c.png_export import pptx_to_pngs
        png_dir = OUTPUT_DIR / f"{NAME}_pngs"
        paths = pptx_to_pngs(pptx_path, png_dir)
        print(f"PNGs: {png_dir} ({len(paths)} images)")
    except Exception as e:
        print(f"PNG export skipped: {e}")


if __name__ == "__main__":
    main()
