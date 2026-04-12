"""Tier D — Change Management Components (7 slides).

D1. comp_maturity_model_5level  — 5-step staircase (Initial → Optimizing)
D2. comp_capability_heatmap     — 6 capabilities × 4 maturity levels
D3. comp_s_curve                — 5-phase bell + cumulative S-curve
D4. comp_change_curve           — Kubler-Ross change curve (5 stages)
D5. comp_rollout_wave           — 3 overlapping wave bands on timeline
D6. comp_governance_body        — 4-tier governance hierarchy
D7. comp_change_impact_assessment — 5 depts × 6 elements matrix

Output: output/tier_d_change.pptx + PDF + PNGs.
"""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "tier_d_change"


# ------------------------------------------------------------
# Slide scaffold
# ------------------------------------------------------------
def make(prs, title_text: str, subtitle: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    c.text(title_text, x=0.4, y=0.22, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.66, w=9.2, h=0.25,
               size=9, color="grey_700")
    c.line(x1=0.4, y1=0.98, x2=9.6, y2=0.98, color="grey_200", width=1.0)
    return s, c


# ============================================================
# D1. comp_maturity_model_5level
# ============================================================
def slide_d1_maturity_model(prs):
    s, c = make(
        prs,
        "조직 성숙도 모델 — 5단계 진화 경로",
        "Initial(임시적) → Optimizing(지속 혁신)까지 단계적 역량 성장",
    )

    levels = [
        ("Level 1", "Initial",      "임시적·영웅 의존",   ["비표준 프로세스", "개별 역량 의존", "재현성 낮음"]),
        ("Level 2", "Repeatable",   "프로젝트 단위 반복",  ["기본 관리 체계", "일정/비용 추적", "부서별 고립"]),
        ("Level 3", "Defined",      "조직 표준화",        ["전사 표준 정의", "교육/훈련 체계", "프로세스 자산"]),
        ("Level 4", "Managed",      "정량 관리",          ["KPI 기반 통제", "변동성 예측", "품질 목표 달성"]),
        ("Level 5", "Optimizing",   "지속 혁신",          ["데이터 기반 개선", "혁신 내재화", "업계 벤치마크"]),
    ]

    # Staircase geometry
    base_y = 6.3
    start_x = 0.6
    step_w = 1.78
    min_h = 0.7
    max_h = 2.8
    colors = ["grey_200", "grey_400", "accent_light", "accent_mid", "accent"]
    text_colors = ["grey_900", "grey_900", "grey_900", "white", "white"]

    for i, ((lv, name, desc, chars)) in enumerate(levels):
        h = min_h + (max_h - min_h) * i / (len(levels) - 1)
        x = start_x + i * step_w
        y = base_y - h

        c.box(x=x, y=y, w=step_w - 0.08, h=h,
              fill=colors[i], border=0.75, border_color="grey_mid")

        # Level tag at top of step
        c.text(lv, x=x + 0.1, y=y + 0.08, w=step_w - 0.28, h=0.22,
               size=8, bold=True, color=text_colors[i])
        # Stage name
        c.text(name, x=x + 0.1, y=y + 0.32, w=step_w - 0.28, h=0.28,
               size=12, bold=True, color=text_colors[i])
        # One-line desc
        c.text(desc, x=x + 0.1, y=y + 0.62, w=step_w - 0.28, h=0.22,
               size=8, color=text_colors[i])

        # Characteristics list (only where there's room)
        if h >= 1.3:
            char_text = "\n".join("• " + ch for ch in chars)
            c.text(char_text, x=x + 0.1, y=y + 0.92, w=step_w - 0.28, h=h - 1.0,
                   size=7, color=text_colors[i])

    # Baseline
    c.line(x1=start_x, y1=base_y + 0.02, x2=start_x + len(levels) * step_w,
           y2=base_y + 0.02, color="grey_700", width=1.5)
    # Upward arrow below
    c.arrow(x1=start_x, y1=base_y + 0.35,
            x2=start_x + len(levels) * step_w - 0.1, y2=base_y + 0.35,
            color="accent", width=2.0)
    c.text("성숙도 향상 방향", x=start_x, y=base_y + 0.5, w=step_w * len(levels),
           h=0.22, size=8, bold=True, color="accent", align="center")


# ============================================================
# D2. comp_capability_heatmap
# ============================================================
def slide_d2_capability_heatmap(prs):
    s, c = make(
        prs,
        "역량 히트맵 — 현재 vs 1년 후 vs 3년 후 vs 목표",
        "6개 핵심 역량의 시계열 성숙도 평가 (1~4, 색이 진할수록 높음)",
    )

    caps = [
        ("전략",      [2, 3, 3, 4]),
        ("운영",      [1, 2, 3, 4]),
        ("기술",      [2, 2, 3, 4]),
        ("사람",      [1, 2, 2, 3]),
        ("데이터",    [1, 2, 3, 4]),
        ("거버넌스",  [2, 2, 3, 4]),
    ]
    cols = ["Current", "1yr", "3yr", "Target"]

    # Grid geometry
    left_col_w = 1.6
    cell_w = 1.7
    cell_h = 0.7
    start_x = 1.0
    start_y = 1.5

    # Header row
    c.text("역량 \\ 시점", x=start_x, y=start_y, w=left_col_w, h=cell_h,
           size=9, bold=True, color="grey_900", anchor="middle", align="left")
    for j, col in enumerate(cols):
        cx = start_x + left_col_w + j * cell_w
        c.box(x=cx, y=start_y, w=cell_w - 0.05, h=cell_h,
              fill="grey_900", border=None)
        c.text(col, x=cx, y=start_y, w=cell_w - 0.05, h=cell_h,
               size=10, bold=True, color="white", anchor="middle", align="center")

    # Score → color (1~4)
    score_colors = {
        1: ("grey_100", "grey_900"),
        2: ("grey_200", "grey_900"),
        3: ("accent_mid", "white"),
        4: ("accent", "white"),
    }

    for i, (cap, scores) in enumerate(caps):
        ry = start_y + (i + 1) * cell_h
        # Row label
        c.box(x=start_x, y=ry, w=left_col_w, h=cell_h,
              fill="grey_100", border=0.5, border_color="grey_mid")
        c.text(cap, x=start_x + 0.12, y=ry, w=left_col_w - 0.2, h=cell_h,
               size=10, bold=True, color="grey_900", anchor="middle")
        # Score cells
        for j, sc in enumerate(scores):
            cx = start_x + left_col_w + j * cell_w
            fill, txt_col = score_colors[sc]
            c.box(x=cx, y=ry, w=cell_w - 0.05, h=cell_h,
                  fill=fill, border=0.5, border_color="grey_mid")
            c.text(str(sc), x=cx, y=ry, w=cell_w - 0.05, h=cell_h,
                   size=14, bold=True, color=txt_col,
                   anchor="middle", align="center")

    # Legend
    legend_y = start_y + (len(caps) + 1) * cell_h + 0.25
    c.text("점수 기준:", x=start_x, y=legend_y, w=1.1, h=0.25,
           size=9, bold=True, color="grey_900")
    legend_labels = [
        ("1 = 초기",     "grey_100"),
        ("2 = 정의됨",   "grey_200"),
        ("3 = 관리됨",   "accent_mid"),
        ("4 = 최적화",   "accent"),
    ]
    for k, (lbl, col) in enumerate(legend_labels):
        lx = start_x + 1.1 + k * 1.7
        c.box(x=lx, y=legend_y + 0.02, w=0.3, h=0.22, fill=col,
              border=0.5, border_color="grey_mid")
        c.text(lbl, x=lx + 0.36, y=legend_y, w=1.3, h=0.28,
               size=8, color="grey_700")


# ============================================================
# D3. comp_s_curve
# ============================================================
def slide_d3_s_curve(prs):
    s, c = make(
        prs,
        "혁신 수용 곡선 — S-Curve & Bell Curve",
        "Rogers의 혁신 확산 이론: 채택자 5단계 분포",
    )

    # Chart frame
    px0, py0 = 1.2, 1.6          # origin top-left of chart area
    pw, ph = 7.8, 4.4             # chart inner width/height
    baseline_y = py0 + ph         # bottom y

    # Axes
    c.line(x1=px0, y1=baseline_y, x2=px0 + pw, y2=baseline_y,
           color="grey_700", width=1.2)
    c.line(x1=px0, y1=py0, x2=px0, y2=baseline_y,
           color="grey_700", width=1.2)

    # 5 phase boundaries (cumulative %): 0, 2.5, 16, 50, 84, 100
    phases = [
        ("Innovators",      2.5,  0.0,   2.5),
        ("Early Adopters",  13.5, 2.5,   16.0),
        ("Early Majority",  34.0, 16.0,  50.0),
        ("Late Majority",   34.0, 50.0,  84.0),
        ("Laggards",        16.0, 84.0,  100.0),
    ]

    # Vertical dividers & phase labels + %
    shade = ["accent_light", "accent_mid", "accent", "accent_mid", "accent_light"]
    for i, (name, pct, c0, c1) in enumerate(phases):
        x_start = px0 + pw * c0 / 100.0
        x_end   = px0 + pw * c1 / 100.0
        band_w = x_end - x_start
        # Light band at very bottom (phase zone)
        c.box(x=x_start, y=baseline_y + 0.05, w=band_w - 0.02, h=0.25,
              fill=shade[i], border=None)
        c.text(name, x=x_start, y=baseline_y + 0.32, w=band_w, h=0.22,
               size=8, bold=True, color="grey_900", align="center")
        c.text(f"{pct}%", x=x_start, y=baseline_y + 0.52, w=band_w, h=0.22,
               size=9, bold=True, color="accent", align="center")
        # Divider verticals inside chart
        if i > 0:
            c.line(x1=x_start, y1=py0, x2=x_start, y2=baseline_y,
                   color="grey_200", width=0.75)

    # Bell curve (approximate) — sample normal pdf over [-3,3], map to chart
    import random
    n = 60
    bell_points = []
    for i in range(n + 1):
        t = -3.0 + 6.0 * i / n
        y = math.exp(-0.5 * t * t)              # 0..1
        # map t to chart x: normal cdf-ish linear stretch so peak at 50%
        cdf = 0.5 * (1 + math.erf(t / math.sqrt(2))) * 100
        x = px0 + pw * cdf / 100.0
        py = baseline_y - y * (ph * 0.55)
        bell_points.append((x, py))
    # draw bell as polyline
    for i in range(len(bell_points) - 1):
        x1, y1 = bell_points[i]
        x2, y2 = bell_points[i + 1]
        c.line(x1=x1, y1=y1, x2=x2, y2=y2, color="accent_mid", width=2.0)

    # Cumulative S-curve — normal cdf
    s_points = []
    for i in range(n + 1):
        t = -3.0 + 6.0 * i / n
        cdf = 0.5 * (1 + math.erf(t / math.sqrt(2)))
        cdf_pct = cdf * 100
        x = px0 + pw * cdf_pct / 100.0
        py = baseline_y - cdf * ph * 0.9
        s_points.append((x, py))
    for i in range(len(s_points) - 1):
        x1, y1 = s_points[i]
        x2, y2 = s_points[i + 1]
        c.line(x1=x1, y1=y1, x2=x2, y2=y2, color="accent", width=2.5)

    # Legend
    c.box(x=px0 + 0.2, y=py0 + 0.1, w=2.8, h=0.7, fill="white",
          border=0.75, border_color="grey_mid")
    c.line(x1=px0 + 0.35, y1=py0 + 0.28, x2=px0 + 0.75, y2=py0 + 0.28,
           color="accent_mid", width=2.0)
    c.text("Bell Curve (분포)", x=px0 + 0.85, y=py0 + 0.18, w=1.9, h=0.22,
           size=8, color="grey_900")
    c.line(x1=px0 + 0.35, y1=py0 + 0.55, x2=px0 + 0.75, y2=py0 + 0.55,
           color="accent", width=2.5)
    c.text("S-Curve (누적 채택)", x=px0 + 0.85, y=py0 + 0.45, w=1.9, h=0.22,
           size=8, color="grey_900")

    # Y-axis label
    c.text("채택자\n비율", x=px0 - 1.0, y=py0 + ph / 2 - 0.3, w=0.8, h=0.6,
           size=8, bold=True, color="grey_700", align="right")


# ============================================================
# D4. comp_change_curve (Kubler-Ross)
# ============================================================
def slide_d4_change_curve(prs):
    s, c = make(
        prs,
        "변화 곡선 — Kubler-Ross 5단계",
        "변화 수용까지 조직/개인이 거치는 정서·성과 궤적",
    )

    px0, py0 = 1.0, 1.7
    pw, ph = 8.3, 4.2
    base_y = py0 + ph * 0.5   # horizontal mid-line (baseline performance)

    # Axes
    c.line(x1=px0, y1=py0 + ph, x2=px0 + pw, y2=py0 + ph,
           color="grey_700", width=1.2)   # time axis bottom
    c.line(x1=px0, y1=py0, x2=px0, y2=py0 + ph,
           color="grey_700", width=1.2)   # perf axis
    # Baseline dashed-ish
    c.line(x1=px0, y1=base_y, x2=px0 + pw, y2=base_y,
           color="grey_400", width=0.75)

    c.text("시간 →", x=px0 + pw - 0.8, y=py0 + ph + 0.05, w=0.8, h=0.22,
           size=8, color="grey_700")
    c.text("성과 /\n정서", x=px0 - 0.9, y=py0 + ph / 2 - 0.2, w=0.8, h=0.5,
           size=8, bold=True, color="grey_700", align="right")

    # 5 stages — x position (0-1 normalized) and y offset from base (+ down = lower performance)
    stages = [
        ("Denial",      0.08, +0.25, "부정 / 충격\n현실 회피"),
        ("Resistance",  0.25, +1.20, "저항 / 분노\n저점 도달"),
        ("Exploration", 0.50, +0.40, "탐색 / 수용\n학습 시작"),
        ("Commitment",  0.75, -0.60, "몰입 / 실행\n성과 회복"),
        ("Integration", 0.95, -1.15, "통합 / 내재화\n이전 수준 초과"),
    ]

    # Build curve points through stages using cubic interpolation
    stage_pts = []
    for name, tx, dy, note in stages:
        x = px0 + tx * pw
        y = base_y + dy * (ph * 0.35)
        # clamp inside chart
        y = max(py0 + 0.1, min(py0 + ph - 0.1, y))
        stage_pts.append((x, y))

    # Interpolate smooth polyline: 15 segments between each consecutive pair using catmull-rom
    def catmull(p0, p1, p2, p3, t):
        return (
            0.5 * ((2 * p1[0]) + (-p0[0] + p2[0]) * t +
                   (2 * p0[0] - 5 * p1[0] + 4 * p2[0] - p3[0]) * t * t +
                   (-p0[0] + 3 * p1[0] - 3 * p2[0] + p3[0]) * t * t * t),
            0.5 * ((2 * p1[1]) + (-p0[1] + p2[1]) * t +
                   (2 * p0[1] - 5 * p1[1] + 4 * p2[1] - p3[1]) * t * t +
                   (-p0[1] + 3 * p1[1] - 3 * p2[1] + p3[1]) * t * t * t),
        )

    # Pad endpoints
    padded = [stage_pts[0]] + stage_pts + [stage_pts[-1]]
    curve = []
    segs = 18
    for i in range(len(padded) - 3):
        p0, p1, p2, p3 = padded[i], padded[i + 1], padded[i + 2], padded[i + 3]
        for k in range(segs):
            t = k / segs
            curve.append(catmull(p0, p1, p2, p3, t))
    curve.append(stage_pts[-1])

    # Draw curve
    for i in range(len(curve) - 1):
        x1, y1 = curve[i]
        x2, y2 = curve[i + 1]
        c.line(x1=x1, y1=y1, x2=x2, y2=y2, color="accent", width=2.5)

    # Stage dots + annotations
    for i, (name, tx, dy, note) in enumerate(stages):
        x, y = stage_pts[i]
        c.circle(x=x - 0.12, y=y - 0.12, d=0.24, fill="accent",
                 border=1.0, border_color="white", text="")
        # Stage label above/below depending on curve position
        ann_above = y > base_y
        ann_y = y - 0.75 if ann_above else y + 0.25
        c.box(x=x - 0.85, y=ann_y, w=1.7, h=0.58,
              fill="white", border=0.75, border_color="accent")
        c.text(f"{i + 1}. {name}", x=x - 0.8, y=ann_y + 0.03, w=1.6, h=0.22,
               size=8, bold=True, color="accent", align="center")
        c.text(note, x=x - 0.8, y=ann_y + 0.25, w=1.6, h=0.32,
               size=7, color="grey_900", align="center")


# ============================================================
# D5. comp_rollout_wave
# ============================================================
def slide_d5_rollout_wave(prs):
    s, c = make(
        prs,
        "롤아웃 웨이브 플랜 — 3 Waves 병행 배포",
        "Wave별 범위 확대 + 후행 Wave가 선행 학습을 흡수",
    )

    px0, py0 = 0.7, 1.6
    pw, ph = 8.7, 4.8

    # Timeline baseline
    c.line(x1=px0, y1=py0 + ph - 0.8, x2=px0 + pw, y2=py0 + ph - 0.8,
           color="grey_700", width=1.2)
    # Month markers
    months = ["Q1", "Q2", "Q3", "Q4", "Q1+1", "Q2+1"]
    for i, m in enumerate(months):
        mx = px0 + pw * i / (len(months) - 1)
        c.line(x1=mx, y1=py0 + ph - 0.85, x2=mx, y2=py0 + ph - 0.75,
               color="grey_700", width=1.0)
        c.text(m, x=mx - 0.35, y=py0 + ph - 0.7, w=0.7, h=0.22,
               size=8, bold=True, color="grey_700", align="center")

    waves = [
        {
            "name": "Wave 1",
            "period": "Q1–Q2",
            "t_start": 0.02, "t_end": 0.40,
            "y_top": py0 + 0.4,
            "color": "accent_light",
            "scope": ["본사 재무팀", "2개 Legal Entity", "핵심 프로세스 7개"],
        },
        {
            "name": "Wave 2",
            "period": "Q2–Q4",
            "t_start": 0.25, "t_end": 0.70,
            "y_top": py0 + 1.4,
            "color": "accent_mid",
            "scope": ["영업/구매 부서", "5개 법인 확장", "통합 리포팅"],
        },
        {
            "name": "Wave 3",
            "period": "Q4–Q2+1",
            "t_start": 0.55, "t_end": 0.98,
            "y_top": py0 + 2.4,
            "color": "accent",
            "scope": ["해외 법인 10개", "생산/물류 전체", "고도화 분석"],
        },
    ]

    bottom_y = py0 + ph - 0.85

    # Draw wave bands (arc-like: rectangle + rounded top using simple rectangle + top curve approximation via smaller rectangles)
    for w in waves:
        x1 = px0 + pw * w["t_start"]
        x2 = px0 + pw * w["t_end"]
        width = x2 - x1
        top_y = w["y_top"]
        band_h = bottom_y - top_y

        # Main band
        c.box(x=x1, y=top_y, w=width, h=band_h,
              fill=w["color"], border=0.75, border_color="grey_mid",
              shape="rounded")

        # Label box on the band
        label_txt = w["name"]
        txt_col = "white" if w["color"] == "accent" else "grey_900"
        c.text(label_txt, x=x1 + 0.1, y=top_y + 0.08, w=1.2, h=0.3,
               size=13, bold=True, color=txt_col)
        c.text(w["period"], x=x1 + 0.1, y=top_y + 0.4, w=1.5, h=0.22,
               size=8, color=txt_col)

        # Scope items
        scope_text = "\n".join("• " + sc for sc in w["scope"])
        c.text(scope_text, x=x1 + 0.1, y=top_y + 0.7, w=width - 0.2,
               h=band_h - 0.75,
               size=8, color=txt_col)

    # Title labels above waves
    c.text("범위 →", x=px0 - 0.55, y=py0 + ph / 2 - 0.15, w=0.6, h=0.3,
           size=8, bold=True, color="grey_700", align="right")


# ============================================================
# D6. comp_governance_body
# ============================================================
def slide_d6_governance_body(prs):
    s, c = make(
        prs,
        "프로그램 거버넌스 체계 — 4-Tier 의사결정 구조",
        "Steering → Program → Working Group → Task Force",
    )

    tiers = [
        {
            "name": "Steering Committee",
            "role": "전략 방향 / 투자 승인",
            "members": ["CEO", "CFO", "CIO", "사업부문장 3명"],
            "freq": "월 1회 / 필요 시",
            "fill": "accent",
            "text_color": "white",
            "w_ratio": 0.5,
        },
        {
            "name": "Program Management Office",
            "role": "프로그램 총괄 / 리스크·일정 관리",
            "members": ["PgM", "PMO Lead", "변화관리 리드", "재무 PM"],
            "freq": "주 1회 정례",
            "fill": "accent_mid",
            "text_color": "white",
            "w_ratio": 0.7,
        },
        {
            "name": "Working Groups (5개)",
            "role": "기능별 설계·실행",
            "members": ["재무", "공급망", "HR", "IT Infra", "데이터"],
            "freq": "주 2회 / 스프린트",
            "fill": "accent_light",
            "text_color": "grey_900",
            "w_ratio": 0.85,
        },
        {
            "name": "Task Forces (10+)",
            "role": "이슈 단위 즉시 해결",
            "members": ["현업 SME", "외부 컨설턴트", "개발팀"],
            "freq": "이슈 발생 시 / 데일리 스탠드업",
            "fill": "grey_200",
            "text_color": "grey_900",
            "w_ratio": 1.0,
        },
    ]

    center_x = 5.0
    y = 1.3
    tier_h = 1.3
    gap = 0.15
    max_w = 7.8

    for i, t in enumerate(tiers):
        w = max_w * t["w_ratio"]
        x = center_x - w / 2
        c.box(x=x, y=y, w=w, h=tier_h,
              fill=t["fill"], border=1.0, border_color="grey_mid")

        # Name
        c.text(t["name"], x=x + 0.2, y=y + 0.08, w=w - 0.4, h=0.3,
               size=13, bold=True, color=t["text_color"])
        # Role
        c.text(t["role"], x=x + 0.2, y=y + 0.38, w=w - 0.4, h=0.25,
               size=9, color=t["text_color"])

        # Members (left) + freq (right)
        members_text = "구성원: " + " / ".join(t["members"])
        c.text(members_text, x=x + 0.2, y=y + 0.66, w=w * 0.65, h=0.5,
               size=8, color=t["text_color"])
        c.text("개최: " + t["freq"], x=x + w * 0.66, y=y + 0.66,
               w=w * 0.32, h=0.5,
               size=8, bold=True, color=t["text_color"], align="right")

        # Connector arrow to next tier
        if i < len(tiers) - 1:
            c.arrow(x1=center_x, y1=y + tier_h,
                    x2=center_x, y2=y + tier_h + gap,
                    color="grey_700", width=1.5)
        y += tier_h + gap


# ============================================================
# D7. comp_change_impact_assessment
# ============================================================
def slide_d7_change_impact(prs):
    s, c = make(
        prs,
        "변화 영향도 평가 — 부서 × 변화 요소 매트릭스",
        "5개 부서에 대한 6개 변화 요소의 영향 강도 (None ~ Critical)",
    )

    depts = ["재무", "영업", "공급망", "HR", "IT"]
    elements = ["프로세스", "시스템", "조직구조", "역할정의", "KPI", "정책"]

    # Impact levels: 0~4
    # 샘플 매트릭스 (행=부서, 열=요소)
    matrix = [
        [4, 4, 2, 3, 3, 2],   # 재무
        [3, 3, 1, 2, 4, 1],   # 영업
        [4, 3, 3, 3, 2, 2],   # 공급망
        [2, 1, 4, 4, 3, 3],   # HR
        [3, 4, 2, 2, 1, 2],   # IT
    ]

    impact_spec = {
        0: ("None",     "white",       "grey_900"),
        1: ("Low",      "grey_200",    "grey_900"),
        2: ("Medium",   "accent_mid",  "white"),
        3: ("High",     "accent",      "white"),
        4: ("Critical", "negative",    "white"),
    }

    # Grid geometry
    grid_x = 0.9
    grid_y = 1.6
    cell_w = 1.05
    cell_h = 0.75
    row_label_w = 1.1

    # Header: elements
    for j, el in enumerate(elements):
        cx = grid_x + row_label_w + j * cell_w
        c.box(x=cx, y=grid_y, w=cell_w - 0.03, h=cell_h,
              fill="grey_900", border=None)
        c.text(el, x=cx, y=grid_y, w=cell_w - 0.03, h=cell_h,
               size=9, bold=True, color="white",
               align="center", anchor="middle")

    # Rows
    for i, dept in enumerate(depts):
        ry = grid_y + (i + 1) * cell_h
        # Row label
        c.box(x=grid_x, y=ry, w=row_label_w, h=cell_h,
              fill="grey_100", border=0.5, border_color="grey_mid")
        c.text(dept, x=grid_x + 0.1, y=ry, w=row_label_w - 0.15, h=cell_h,
               size=10, bold=True, color="grey_900", anchor="middle")

        for j, score in enumerate(matrix[i]):
            cx = grid_x + row_label_w + j * cell_w
            label, fill, txt_col = impact_spec[score]
            c.box(x=cx, y=ry, w=cell_w - 0.03, h=cell_h,
                  fill=fill, border=0.5, border_color="grey_mid")
            c.text(label, x=cx, y=ry, w=cell_w - 0.03, h=cell_h,
                   size=8, bold=True, color=txt_col,
                   align="center", anchor="middle")

    # Legend (right side)
    legend_x = grid_x + row_label_w + len(elements) * cell_w + 0.3
    legend_y = grid_y
    c.text("영향도 등급", x=legend_x, y=legend_y, w=1.7, h=0.25,
           size=10, bold=True, color="grey_900")

    for k in range(5):
        ly = legend_y + 0.35 + k * 0.55
        label, fill, txt_col = impact_spec[k]
        c.box(x=legend_x, y=ly, w=0.45, h=0.4,
              fill=fill, border=0.5, border_color="grey_mid")
        c.text(label, x=legend_x + 0.52, y=ly + 0.05, w=1.2, h=0.3,
               size=9, bold=True, color="grey_900")

    # Summary note
    c.text("※ 전략 / 운영 / 기술 / 사람 / 데이터 / 거버넌스 차원의 변화 영향을 부서별로 정량 평가",
           x=grid_x, y=grid_y + (len(depts) + 1) * cell_h + 0.2,
           w=9.0, h=0.3, size=8, color="grey_700")


# ============================================================
# Evaluate + Export
# ============================================================
def evaluate_and_report(pptx_path: Path):
    from ppt_builder.evaluate import evaluate_pptx, print_report
    report = evaluate_pptx(str(pptx_path))
    print_report(report)
    return report


def export_pdf_and_pngs(pptx_path: Path):
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        pdf_path = pptx_path.with_suffix(".pdf")
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        p = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        p.SaveAs(str(pdf_path.resolve()), 32)
        p.Close()
        try:
            ppt.Quit()
        except Exception:
            pass
        print(f"PDF:  {pdf_path}")
    except Exception as e:
        print(f"[WARN] PDF export skipped: {e}")

    try:
        from ppt_builder.track_c.png_export import pptx_to_pngs
        png_dir = OUTPUT_DIR / f"{NAME}_pngs"
        paths = pptx_to_pngs(pptx_path, png_dir)
        print(f"PNGs: {png_dir} ({len(paths)}장)")
    except Exception as e:
        print(f"[WARN] PNG export skipped: {e}")


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_d1_maturity_model(prs)
    slide_d2_capability_heatmap(prs)
    slide_d3_s_curve(prs)
    slide_d4_change_curve(prs)
    slide_d5_rollout_wave(prs)
    slide_d6_governance_body(prs)
    slide_d7_change_impact(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX saved: {pptx_path}")

    evaluate_and_report(pptx_path)
    export_pdf_and_pngs(pptx_path)

    print(f"\nOutput: {pptx_path}")


if __name__ == "__main__":
    main()
