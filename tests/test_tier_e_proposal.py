"""Tier E — Proposal / Sales Components (6 slides).

E1. comp_team_page          — 6-member team grid (2x3)
E2. comp_credentials_wall   — 12 client logo placeholders (3x4)
E3. comp_case_study_card    — Challenge / Solution / Results stacked card
E4. comp_pricing_tiers      — 3 pricing tiers side-by-side, middle highlighted
E5. comp_testimonial_quote  — Large quote + author block
E6. comp_roi_summary        — Investment / Benefit split + bottom ROI strip

Output: output/tier_e_proposal.pptx (+ .pdf + PNGs under output/tier_e_proposal/)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "tier_e_proposal"


# ============================================================
# Slide header helper
# ============================================================
def make(prs, title_text: str, subtitle: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    # accent top bar
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    # title
    c.text(title_text, x=0.4, y=0.22, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.66, w=9.2, h=0.25,
               size=9, color="grey_700")
    # separator
    c.line(x1=0.4, y1=0.98, x2=9.6, y2=0.98,
           color="grey_200", width=1.0)
    return s, c


# ============================================================
# E1. comp_team_page — 6 members (2 rows × 3 cols)
# ============================================================
def slide_e1_team_page(prs):
    _, c = make(prs, "Engagement Team — 6인의 Senior Consultant 편성",
                "SAP S/4HANA, Finance, SCM, HR 분야별 전문가 구성")

    members = [
        ("KS", "김준성", "Engagement Partner",
         "SAP 15년 · HD현대 ERP 전환 PM\nS/4HANA 대규모 이행 7건 완수"),
        ("LJ", "이재훈", "Finance Lead",
         "CO/FI 모듈 12년 · 삼성전자 출신\nIFRS 전환 프로젝트 다수 리딩"),
        ("PM", "박미래", "SCM Lead",
         "MM/PP/WM 10년 · 현대차 공급망\nS&OP 프로세스 재설계 전문"),
        ("CH", "최하영", "HCM Lead",
         "SuccessFactors 공인 · 8년\n글로벌 HR 표준화 프로젝트 4건"),
        ("YS", "윤서진", "Tech Architect",
         "BTP · CAP · Fiori 9년\nGreenfield 아키텍처 설계 전문"),
        ("JW", "정원석", "Change Manager",
         "Prosci ADKAR · 11년\n대기업 전환 Change Mgmt 6건"),
    ]

    cols, rows = 3, 2
    start_x, start_y = 0.5, 1.3
    col_w, row_h = 3.05, 2.85
    gap_x = 0.15

    circle_d = 0.9

    for idx, (initials, name, title, detail) in enumerate(members):
        r = idx // cols
        col = idx % cols
        cx = start_x + col * (col_w + gap_x)
        cy = start_y + r * row_h

        # card
        c.box(x=cx, y=cy, w=col_w, h=row_h - 0.2,
              fill="white", border=0.75, border_color="grey_mid")

        # initials circle
        c.circle(x=cx + (col_w - circle_d) / 2, y=cy + 0.25, d=circle_d,
                 fill="accent", border=None,
                 text=initials, text_color="white",
                 text_size=18, text_bold=True)

        # name
        c.text(name, x=cx + 0.1, y=cy + 1.25, w=col_w - 0.2, h=0.3,
               size=13, bold=True, color="grey_900",
               align="center", anchor="top")

        # title (accent)
        c.text(title, x=cx + 0.1, y=cy + 1.58, w=col_w - 0.2, h=0.25,
               size=9, bold=True, color="accent",
               align="center", anchor="top")

        # divider
        c.line(x1=cx + 0.4, y1=cy + 1.90, x2=cx + col_w - 0.4, y2=cy + 1.90,
               color="grey_200", width=0.75)

        # experience detail
        c.text(detail, x=cx + 0.15, y=cy + 1.98, w=col_w - 0.3, h=0.6,
               size=8, color="grey_700", align="center", anchor="top")

    # footer
    c.text("※ 전체 팀 평균 SAP 경력 10.8년 · S/4HANA 이행 프로젝트 누적 42건",
           x=0.4, y=7.1, w=9.2, h=0.25,
           size=8, color="grey_700", align="center")


# ============================================================
# E2. comp_credentials_wall — 12 client logos (4 cols × 3 rows)
# ============================================================
def slide_e2_credentials_wall(prs):
    _, c = make(prs, "주요 고객사 크리덴셜 — Industry Leader 12사 S/4HANA 이행 레퍼런스",
                "제조·금융·유통·에너지 등 산업별 Top-tier 고객 기반")

    clients = [
        ("Samsung", "Electronics"),
        ("Hyundai", "Motors"),
        ("LG", "Energy Solution"),
        ("SK", "Hynix"),
        ("POSCO", "Holdings"),
        ("Kia", "Corporation"),
        ("Lotte", "Chemical"),
        ("KB", "Financial"),
        ("Shinhan", "Bank"),
        ("CJ", "Logistics"),
        ("Doosan", "Enerbility"),
        ("Hanwha", "Aerospace"),
    ]

    cols, rows = 4, 3
    start_x, start_y = 0.5, 1.35
    col_w, row_h = 2.2, 1.65
    gap_x = 0.12
    gap_y = 0.2

    for idx, (name, sub) in enumerate(clients):
        r = idx // cols
        col = idx % cols
        bx = start_x + col * (col_w + gap_x)
        by = start_y + r * (row_h + gap_y)

        # logo placeholder rectangle
        c.box(x=bx, y=by, w=col_w, h=row_h,
              fill="grey_100", border=1.0, border_color="grey_mid")

        # brand name
        c.text(name, x=bx, y=by + 0.35, w=col_w, h=0.55,
               size=18, bold=True, color="grey_900",
               align="center", anchor="middle")
        # subtext
        c.text(sub, x=bx, y=by + 0.92, w=col_w, h=0.3,
               size=8, color="grey_700",
               align="center", anchor="top")
        # accent underline
        c.box(x=bx + col_w / 2 - 0.25, y=by + row_h - 0.2, w=0.5, h=0.04,
              fill="accent", border=None)

    c.text("※ NDA 서약된 레퍼런스는 별도 요청 시 공개 | 총 42개 S/4HANA 이행 완수",
           x=0.4, y=7.2, w=9.2, h=0.22,
           size=7, color="grey_400", align="center")


# ============================================================
# E3. comp_case_study_card — Challenge / Solution / Results
# ============================================================
def slide_e3_case_study_card(prs):
    _, c = make(prs, "Case Study — HD현대 S/4HANA Greenfield 18개월 성공 사례",
                "Challenge · Solution · Results 3단계 요약")

    # top meta bar: logo + industry + duration
    meta_y = 1.25
    c.box(x=0.5, y=meta_y, w=9.0, h=0.5,
          fill="white", border=0.75, border_color="grey_mid")
    # mini logo box
    c.box(x=0.62, y=meta_y + 0.08, w=1.3, h=0.34,
          fill="grey_900", border=None)
    c.text("HD HYUNDAI", x=0.62, y=meta_y + 0.08, w=1.3, h=0.34,
           size=9, bold=True, color="white",
           align="center", anchor="middle")
    # industry
    c.text("Industry", x=2.15, y=meta_y + 0.05, w=1.2, h=0.2,
           size=7, color="grey_700")
    c.text("Heavy Industry · Shipbuilding", x=2.15, y=meta_y + 0.24, w=2.4, h=0.22,
           size=9, bold=True, color="grey_900")
    # duration
    c.text("Duration", x=4.8, y=meta_y + 0.05, w=1.2, h=0.2,
           size=7, color="grey_700")
    c.text("2024.03 – 2025.09 (18 months)", x=4.8, y=meta_y + 0.24, w=3.0, h=0.22,
           size=9, bold=True, color="grey_900")
    # scope chip
    c.label_chip("Greenfield · 전사", x=7.95, y=meta_y + 0.11,
                 w=1.45, h=0.28, fill="accent")

    # === Challenge (top, grey_700 bg, white text) ===
    ch_y, ch_h = 1.85, 1.45
    c.box(x=0.5, y=ch_y, w=9.0, h=ch_h,
          fill="grey_700", border=None)
    c.text("CHALLENGE", x=0.65, y=ch_y + 0.08, w=2.0, h=0.25,
           size=9, bold=True, color="accent")
    c.text("20년 된 ECC 시스템의 한계 · 12개국 법인 데이터 분절",
           x=0.65, y=ch_y + 0.33, w=8.7, h=0.32,
           size=12, bold=True, color="white")
    c.text("▪  ECC 6.0 (2005년 도입) 커스터마이징 3,200건 누적\n"
           "▪  12개 해외법인 마스터데이터 불일치 · 월결산 14일 소요\n"
           "▪  제조·재무·구매 통합 리포팅 불가 · 경영진 의사결정 지연",
           x=0.65, y=ch_y + 0.68, w=8.7, h=0.75,
           size=9, color="white")

    # === Solution (middle, grey_100 bg) ===
    so_y, so_h = 3.40, 1.55
    c.box(x=0.5, y=so_y, w=9.0, h=so_h,
          fill="grey_100", border=None)
    c.text("SOLUTION", x=0.65, y=so_y + 0.08, w=2.0, h=0.25,
           size=9, bold=True, color="accent")
    c.text("S/4HANA Greenfield + Central Finance · 단일 Instance 전환",
           x=0.65, y=so_y + 0.33, w=8.7, h=0.32,
           size=12, bold=True, color="grey_900")
    # 3 solution pillars
    pillars = [
        ("01", "Single Instance", "전 법인 통합 1 Client"),
        ("02", "Fit-to-Standard", "커스터마이징 87% 삭감"),
        ("03", "Clean Core", "BTP Extension 기반"),
    ]
    pw = (9.0 - 0.3 * 2 - 0.4) / 3
    for i, (num, pt, pd) in enumerate(pillars):
        px = 0.7 + i * (pw + 0.15)
        py = so_y + 0.75
        c.box(x=px, y=py, w=0.35, h=0.35, fill="accent", border=None)
        c.text(num, x=px, y=py, w=0.35, h=0.35,
               size=10, bold=True, color="white",
               align="center", anchor="middle")
        c.text(pt, x=px + 0.45, y=py - 0.02, w=pw - 0.5, h=0.25,
               size=10, bold=True, color="grey_900")
        c.text(pd, x=px + 0.45, y=py + 0.22, w=pw - 0.5, h=0.25,
               size=8, color="grey_700")

    # === Results (bottom, accent bg, white text, big KPIs) ===
    re_y, re_h = 5.05, 1.90
    c.box(x=0.5, y=re_y, w=9.0, h=re_h,
          fill="accent", border=None)
    c.text("RESULTS", x=0.65, y=re_y + 0.08, w=2.0, h=0.25,
           size=9, bold=True, color="white")
    c.text("ROI 287% 달성 · 월결산 14일 → 3일 단축",
           x=0.65, y=re_y + 0.33, w=8.7, h=0.32,
           size=12, bold=True, color="white")

    kpis = [
        ("287%", "3년 누적 ROI"),
        ("14→3日", "월결산 소요일"),
        ("87%", "커스터마이징 삭감"),
        ("₩42B", "연간 운영비 절감"),
    ]
    kw = (9.0 - 0.3 * 2 - 0.45) / 4
    for i, (v, lbl) in enumerate(kpis):
        kx = 0.7 + i * (kw + 0.15)
        ky = re_y + 0.80
        c.text(v, x=kx, y=ky, w=kw, h=0.55,
               size=22, bold=True, color="white",
               align="center", anchor="top")
        c.text(lbl, x=kx, y=ky + 0.60, w=kw, h=0.3,
               size=8, bold=True, color="white",
               align="center", anchor="top")


# ============================================================
# E4. comp_pricing_tiers — Basic / Professional / Enterprise
# ============================================================
def slide_e4_pricing_tiers(prs):
    _, c = make(prs, "Pricing Tiers — 기업 규모별 3단계 구독 패키지",
                "Professional 티어가 중견기업 표준 (Recommended)")

    tiers = [
        {
            "name": "Basic",
            "price": "₩9.9M",
            "unit": "/월 · 최대 50 User",
            "features": [
                "S/4HANA Cloud Essential",
                "표준 Fiori Launchpad",
                "재무 · 구매 · 판매 모듈",
                "월간 패치 적용",
                "비즈니스 아워 지원",
            ],
            "cta": "문의하기",
            "highlight": False,
        },
        {
            "name": "Professional",
            "price": "₩24.9M",
            "unit": "/월 · 최대 200 User",
            "features": [
                "S/4HANA Cloud Advanced",
                "커스텀 Fiori App 5개",
                "전 모듈 (MM·PP·SCM 포함)",
                "실시간 Analytics 대시보드",
                "24/7 Premium 지원",
                "분기별 Health Check",
            ],
            "cta": "가입 신청",
            "highlight": True,
        },
        {
            "name": "Enterprise",
            "price": "Custom",
            "unit": "200+ User · 견적 상담",
            "features": [
                "S/4HANA Private Cloud",
                "무제한 커스터마이징",
                "BTP · AI · IoT 통합",
                "전담 TAM (Technical Account)",
                "SLA 99.95% 보장",
                "연 2회 전략 리뷰",
            ],
            "cta": "영업팀 연결",
            "highlight": False,
        },
    ]

    start_x, start_y = 0.6, 1.4
    card_w, card_h = 2.95, 5.55
    gap = 0.1

    for i, t in enumerate(tiers):
        x = start_x + i * (card_w + gap)
        y = start_y
        hl = t["highlight"]

        # RECOMMENDED tag (above card, only middle)
        if hl:
            c.box(x=x + card_w / 2 - 0.85, y=y - 0.18, w=1.7, h=0.3,
                  fill="accent", border=None)
            c.text("RECOMMENDED", x=x + card_w / 2 - 0.85, y=y - 0.18,
                   w=1.7, h=0.3,
                   size=8, bold=True, color="white",
                   align="center", anchor="middle")

        # card
        border_color = "accent" if hl else "grey_mid"
        border_w = 2.0 if hl else 0.75
        c.box(x=x, y=y, w=card_w, h=card_h,
              fill="white", border=border_w, border_color=border_color)

        # header band
        header_fill = "accent" if hl else "grey_900"
        c.box(x=x, y=y, w=card_w, h=0.55,
              fill=header_fill, border=None)
        c.text(t["name"], x=x, y=y, w=card_w, h=0.55,
               size=14, bold=True, color="white",
               align="center", anchor="middle")

        # price
        c.text(t["price"], x=x, y=y + 0.72, w=card_w, h=0.55,
               size=26, bold=True, color="grey_900",
               align="center", anchor="top")
        c.text(t["unit"], x=x, y=y + 1.32, w=card_w, h=0.25,
               size=8, color="grey_700",
               align="center", anchor="top")

        # divider
        c.line(x1=x + 0.3, y1=y + 1.70, x2=x + card_w - 0.3, y2=y + 1.70,
               color="grey_200", width=0.75)

        # feature list with checkmarks
        fy = y + 1.85
        for feat in t["features"]:
            # check circle
            c.circle(x=x + 0.25, y=fy, d=0.22,
                     fill="accent" if hl else "grey_700",
                     border=None,
                     text="v", text_color="white",
                     text_size=9, text_bold=True)
            c.text(feat, x=x + 0.55, y=fy - 0.02, w=card_w - 0.7, h=0.28,
                   size=9, color="grey_900", anchor="top")
            fy += 0.38

        # CTA button at bottom
        cta_y = y + card_h - 0.55
        cta_fill = "accent" if hl else "grey_900"
        c.box(x=x + 0.3, y=cta_y, w=card_w - 0.6, h=0.4,
              fill=cta_fill, border=None)
        c.text(t["cta"], x=x + 0.3, y=cta_y, w=card_w - 0.6, h=0.4,
               size=10, bold=True, color="white",
               align="center", anchor="middle")


# ============================================================
# E5. comp_testimonial_quote — Large quote + author block
# ============================================================
def slide_e5_testimonial_quote(prs):
    _, c = make(prs, "Customer Testimonial — HD현대 CIO의 직접 인용",
                "18개월 성공적 Go-Live 직후 인터뷰")

    # outer card with accent border
    card_x, card_y = 1.0, 1.55
    card_w, card_h = 8.0, 4.8
    c.box(x=card_x, y=card_y, w=card_w, h=card_h,
          fill="white", border=2.0, border_color="accent")

    # left accent stripe
    c.box(x=card_x, y=card_y, w=0.12, h=card_h,
          fill="accent", border=None)

    # huge opening quote mark
    c.text("\u201C", x=card_x + 0.3, y=card_y + 0.05,
           w=1.5, h=1.8,
           size=120, bold=True, color="accent",
           font="Georgia", anchor="top")

    # quote text (italic-ish via regular, 14pt)
    quote = (
        "이번 S/4HANA 전환은 단순한 시스템 교체가 아니었습니다.  "
        "20년간 쌓인 3,200건의 커스터마이징을 Clean Core로 정리하면서, "
        "전사 프로세스 자체를 재설계한 변혁 프로젝트였습니다.  "
        "월결산이 14일에서 3일로 단축된 순간, 경영진 의사결정 속도가 "
        "완전히 달라졌다는 것을 체감했습니다."
    )
    c.text(quote,
           x=card_x + 1.5, y=card_y + 0.45,
           w=card_w - 1.8, h=2.7,
           size=14, color="grey_900", anchor="top")

    # divider line above author
    c.line(x1=card_x + 0.6, y1=card_y + 3.4,
           x2=card_x + card_w - 0.4, y2=card_y + 3.4,
           color="accent", width=1.5)

    # author row: photo circle + name + title + company
    author_y = card_y + 3.6
    photo_d = 0.95
    photo_x = card_x + 0.6
    c.circle(x=photo_x, y=author_y, d=photo_d,
             fill="grey_700", border=None,
             text="CIO", text_color="white",
             text_size=13, text_bold=True)

    c.text("이 성 호", x=photo_x + photo_d + 0.25,
           y=author_y + 0.02, w=3.5, h=0.35,
           size=14, bold=True, color="grey_900")
    c.text("Chief Information Officer", x=photo_x + photo_d + 0.25,
           y=author_y + 0.40, w=4.5, h=0.3,
           size=10, bold=True, color="accent")
    c.text("HD현대중공업 · Digital Transformation Office",
           x=photo_x + photo_d + 0.25,
           y=author_y + 0.68, w=5.0, h=0.3,
           size=9, color="grey_700")

    # right side small company tag
    c.box(x=card_x + card_w - 1.6, y=author_y + 0.15,
          w=1.2, h=0.7,
          fill="grey_900", border=None)
    c.text("HD HYUNDAI", x=card_x + card_w - 1.6, y=author_y + 0.15,
           w=1.2, h=0.7,
           size=9, bold=True, color="white",
           align="center", anchor="middle")


# ============================================================
# E6. comp_roi_summary — Investment / Benefit split + ROI strip
# ============================================================
def slide_e6_roi_summary(prs):
    _, c = make(prs, "ROI Summary — 3년 누적 287% · 14개월 Payback",
                "Investment ₩500M vs. Benefit ₩1,935M (NPV 기준)")

    # ── Left: Investment table ──
    left_x, left_y = 0.5, 1.35
    left_w, left_h = 4.5, 4.25

    c.box(x=left_x, y=left_y, w=left_w, h=left_h,
          fill="white", border=0.75, border_color="grey_mid")
    # header band
    c.box(x=left_x, y=left_y, w=left_w, h=0.45,
          fill="grey_900", border=None)
    c.text("INVESTMENT BREAKDOWN", x=left_x + 0.15, y=left_y,
           w=left_w - 0.3, h=0.45,
           size=11, bold=True, color="white", anchor="middle")
    c.text("Total: ₩500M", x=left_x, y=left_y,
           w=left_w - 0.15, h=0.45,
           size=10, bold=True, color="accent",
           align="right", anchor="middle")

    inv_rows = [
        ("Phase 1 · Blueprint", "2024 Q1", "₩ 60M"),
        ("Phase 2 · Realization", "2024 Q2-Q3", "₩ 180M"),
        ("Phase 3 · Testing", "2024 Q4", "₩ 90M"),
        ("Phase 4 · Cutover", "2025 Q1", "₩ 80M"),
        ("Phase 5 · Hypercare", "2025 Q2", "₩ 50M"),
        ("Licenses (3yr)", "2024-2026", "₩ 140M"),
    ]
    ry = left_y + 0.55
    row_h = 0.55
    for i, (phase, when, cost) in enumerate(inv_rows):
        fill = "white" if i % 2 == 0 else "grey_100"
        c.box(x=left_x + 0.05, y=ry, w=left_w - 0.1, h=row_h,
              fill=fill, border=None)
        c.text(phase, x=left_x + 0.2, y=ry, w=2.2, h=row_h,
               size=9, bold=True, color="grey_900", anchor="middle")
        c.text(when, x=left_x + 2.3, y=ry, w=1.2, h=row_h,
               size=8, color="grey_700", anchor="middle")
        c.text(cost, x=left_x + 3.4, y=ry, w=left_w - 3.5, h=row_h,
               size=10, bold=True, color="grey_900",
               align="right", anchor="middle")
        ry += row_h

    # ── Right: Benefit calculation ──
    right_x, right_y = 5.2, 1.35
    right_w, right_h = 4.3, 4.25

    c.box(x=right_x, y=right_y, w=right_w, h=right_h,
          fill="white", border=0.75, border_color="grey_mid")
    # header band
    c.box(x=right_x, y=right_y, w=right_w, h=0.45,
          fill="accent", border=None)
    c.text("BENEFIT CALCULATION (3yr NPV)", x=right_x + 0.15, y=right_y,
           w=right_w - 0.3, h=0.45,
           size=11, bold=True, color="white", anchor="middle")

    benefits = [
        ("Cost Savings", "운영비 절감", "₩ 820M", "+","positive"),
        ("Revenue Uplift", "매출 증대 효과", "₩ 680M", "+", "positive"),
        ("Efficiency Gain", "인당 생산성 18%", "₩ 435M", "+", "positive"),
    ]
    by = right_y + 0.65
    for lbl_en, lbl_ko, amount, sign, col in benefits:
        c.box(x=right_x + 0.2, y=by, w=right_w - 0.4, h=0.9,
              fill="grey_100", border=None)
        # left color stripe
        c.box(x=right_x + 0.2, y=by, w=0.08, h=0.9,
              fill=col, border=None)
        c.text(lbl_en, x=right_x + 0.4, y=by + 0.08,
               w=right_w - 1.6, h=0.3,
               size=10, bold=True, color="grey_900")
        c.text(lbl_ko, x=right_x + 0.4, y=by + 0.38,
               w=right_w - 1.6, h=0.25,
               size=8, color="grey_700")
        c.text(amount, x=right_x + right_w - 1.6, y=by,
               w=1.4, h=0.9,
               size=16, bold=True, color=col,
               align="right", anchor="middle")
        by += 1.0

    # Total benefit
    c.box(x=right_x + 0.2, y=by, w=right_w - 0.4, h=0.55,
          fill="accent", border=None)
    c.text("TOTAL BENEFIT", x=right_x + 0.4, y=by,
           w=2.0, h=0.55,
           size=10, bold=True, color="white", anchor="middle")
    c.text("₩ 1,935M", x=right_x + right_w - 1.8, y=by,
           w=1.55, h=0.55,
           size=14, bold=True, color="white",
           align="right", anchor="middle")

    # ── Bottom strip: Big ROI + Payback + NPV ──
    strip_y = 5.75
    strip_h = 1.3
    c.box(x=0.5, y=strip_y, w=9.0, h=strip_h,
          fill="grey_900", border=None)
    # left accent block
    c.box(x=0.5, y=strip_y, w=0.15, h=strip_h,
          fill="accent", border=None)

    # ROI (huge)
    c.text("287%", x=0.8, y=strip_y + 0.15, w=2.8, h=0.9,
           size=46, bold=True, color="accent",
           font="Georgia", anchor="top")
    c.text("3-Year Cumulative ROI", x=0.8, y=strip_y + 0.95,
           w=2.8, h=0.3,
           size=8, bold=True, color="white", anchor="top")

    # divider
    c.line(x1=3.8, y1=strip_y + 0.2, x2=3.8, y2=strip_y + strip_h - 0.2,
           color="grey_700", width=0.75)

    # Payback
    c.text("14 Months", x=4.0, y=strip_y + 0.2, w=2.7, h=0.5,
           size=24, bold=True, color="white",
           font="Georgia", anchor="top")
    c.text("Payback Period", x=4.0, y=strip_y + 0.75,
           w=2.7, h=0.25,
           size=8, bold=True, color="accent", anchor="top")
    c.text("손익분기점 도달", x=4.0, y=strip_y + 0.98,
           w=2.7, h=0.25,
           size=7, color="grey_400", anchor="top")

    # divider
    c.line(x1=6.8, y1=strip_y + 0.2, x2=6.8, y2=strip_y + strip_h - 0.2,
           color="grey_700", width=0.75)

    # NPV
    c.text("₩ 1,435M", x=7.0, y=strip_y + 0.2, w=2.4, h=0.5,
           size=22, bold=True, color="white",
           font="Georgia", anchor="top")
    c.text("Net Present Value", x=7.0, y=strip_y + 0.75,
           w=2.4, h=0.25,
           size=8, bold=True, color="accent", anchor="top")
    c.text("WACC 8% 적용", x=7.0, y=strip_y + 0.98,
           w=2.4, h=0.25,
           size=7, color="grey_400", anchor="top")


# ============================================================
# Export (PDF + PNGs)
# ============================================================
def export_pdf_and_pngs(pptx_path: Path) -> tuple[Path | None, list[Path]]:
    pdf_path: Path | None = None
    png_paths: list[Path] = []

    # PDF
    try:
        from ppt_builder.visual_validate import convert_pptx_to_pdf
        pdf_path = convert_pptx_to_pdf(pptx_path)
        print(f"PDF saved:  {pdf_path}")
    except Exception as e:
        print(f"[WARN] PDF conversion skipped: {e}")

    # PNGs
    try:
        from ppt_builder.track_c.png_export import pptx_to_pngs
        png_dir = pptx_path.parent / pptx_path.stem
        png_paths = pptx_to_pngs(pptx_path, png_dir)
        print(f"PNGs saved: {png_dir} ({len(png_paths)} files)")
    except Exception as e:
        print(f"[WARN] PNG export skipped: {e}")

    return pdf_path, png_paths


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_e1_team_page(prs)
    slide_e2_credentials_wall(prs)
    slide_e3_case_study_card(prs)
    slide_e4_pricing_tiers(prs)
    slide_e5_testimonial_quote(prs)
    slide_e6_roi_summary(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX saved: {pptx_path}")

    export_pdf_and_pngs(pptx_path)

    # Evaluate (best-effort)
    try:
        from ppt_builder.evaluate import evaluate_pptx, print_report
        report = evaluate_pptx(str(pptx_path))
        print_report(report)
    except Exception as e:
        print(f"[WARN] Evaluation skipped: {e}")


if __name__ == "__main__":
    main()
