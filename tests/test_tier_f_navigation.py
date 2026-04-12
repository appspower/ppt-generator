"""Tier F — Navigation / Structure Components (5 slides).

F1. comp_agenda_slide         — Numbered agenda (6 items) with vertical accent line
F2. comp_section_divider      — Full-slide section divider (big number + title)
F3. comp_agenda_tracker       — Horizontal progress tracker (6 segments) + detail
F4. comp_executive_summary_scqa — SCQA 4-section stacked bars
F5. comp_key_findings_summary — 2x2 finding boxes with accent top strip

Output: output/tier_f_navigation.pptx + PDF + PNGs
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "tier_f_navigation"


def make(prs, title_text: str, subtitle: str = ""):
    """Create a blank slide with a standard header."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    # Thin accent top bar
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    # Title
    c.text(title_text, x=0.4, y=0.22, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.66, w=9.2, h=0.25,
               size=9, color="grey_700")
    # Separator under header
    c.line(x1=0.4, y1=0.98, x2=9.6, y2=0.98, color="grey_200", width=1.0)
    return s, c


# ============================================================
# F1. comp_agenda_slide — Numbered agenda, 6 items
# ============================================================
def slide_f1_agenda(prs):
    s, c = make(prs,
                "Agenda — 본 세션의 6가지 논의 주제",
                "각 섹션 페이지 참조와 함께 구성된 전체 목차")

    # Vertical accent line on the left (content column separator)
    line_x = 0.7
    top_y = 1.25
    bottom_y = 7.05
    c.box(x=line_x, y=top_y, w=0.04, h=bottom_y - top_y,
          fill="accent", border=None)

    items = [
        ("01", "경영 환경 진단",
         "거시 지표·산업 트렌드·경쟁 구도 핵심 요약", "p. 03"),
        ("02", "전략 방향성 정의",
         "3~5년 중기 비전·미션·핵심 가치 재정립", "p. 09"),
        ("03", "사업 포트폴리오 재편",
         "성장·유지·철수 축으로 본 사업 우선순위", "p. 17"),
        ("04", "디지털 전환 로드맵",
         "AI/Cloud 기반 Core 혁신 단계별 실행 계획", "p. 25"),
        ("05", "조직 및 인력 운영 모델",
         "Target Operating Model 및 핵심 역량 내재화", "p. 33"),
        ("06", "재무 효과 및 투자 계획",
         "3년 CAPEX/OPEX, NPV, 핵심 KPI 전망", "p. 41"),
    ]

    row_h = (bottom_y - top_y) / len(items)
    for i, (num, title_txt, desc, pg) in enumerate(items):
        y = top_y + i * row_h
        # Row separator (except first)
        if i > 0:
            c.line(x1=line_x + 0.25, y1=y, x2=9.6, y2=y,
                   color="grey_200", width=0.5)

        # Big number (accent, 24pt)
        c.text(num, x=line_x + 0.2, y=y + 0.06, w=0.9, h=0.6,
               size=24, bold=True, color="accent",
               font="Georgia", anchor="top")

        # Section title (bold 14pt)
        c.text(title_txt, x=line_x + 1.25, y=y + 0.08,
               w=6.5, h=0.38,
               size=14, bold=True, color="grey_900", anchor="top")

        # Brief description (8pt, grey_700)
        c.text(desc, x=line_x + 1.25, y=y + 0.48,
               w=6.5, h=0.3,
               size=8, color="grey_700", anchor="top")

        # Page ref (8pt right-aligned)
        c.text(pg, x=8.4, y=y + 0.22, w=1.2, h=0.3,
               size=8, color="grey_400",
               align="right", anchor="top")


# ============================================================
# F2. comp_section_divider — Full-slide section divider
# ============================================================
def slide_f2_section_divider(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)

    # Top grey strip
    c.box(x=0, y=0, w=10, h=0.55, fill="grey_100", border=None)
    # Bottom grey strip
    c.box(x=0, y=6.95, w=10, h=0.55, fill="grey_100", border=None)

    # Small eyebrow label
    c.text("SECTION", x=0.5, y=0.15, w=3.0, h=0.25,
           size=9, bold=True, color="accent")
    c.text("Part 03 of 06", x=6.5, y=0.15, w=3.0, h=0.25,
           size=9, color="grey_700", align="right")

    # Large section number on the left (96pt, accent)
    c.text("03", x=0.6, y=1.9, w=3.6, h=3.6,
           size=96, bold=True, color="accent",
           font="Georgia", anchor="middle")

    # Vertical divider between number and title
    c.line(x1=4.3, y1=2.2, x2=4.3, y2=5.3,
           color="grey_400", width=1.2)

    # Section title on the right (28pt bold grey_900)
    c.text("사업 포트폴리오 재편",
           x=4.6, y=2.5, w=5.2, h=0.8,
           size=28, bold=True, color="grey_900", anchor="top")

    # Small description below (10pt)
    c.text("성장·유지·철수 3축 프레임으로 본 핵심 사업의 우선순위 재정의와 "
           "투자/회수 계획의 구체화. 향후 3년 간의 핵심 자원 배분 원칙을 제시한다.",
           x=4.6, y=3.45, w=5.2, h=1.6,
           size=10, color="grey_700", anchor="top")

    # Accent underline accent for title
    c.box(x=4.6, y=3.35, w=0.6, h=0.05, fill="accent", border=None)

    # Bottom strip text
    c.text("HD현대 · Strategy 2029 · Confidential",
           x=0.5, y=7.08, w=9.0, h=0.3,
           size=8, color="grey_700", anchor="middle")


# ============================================================
# F3. comp_agenda_tracker — Top progress indicator + detail below
# ============================================================
def slide_f3_agenda_tracker(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)

    # --- Horizontal tracker (full width) ---
    sections = [
        ("01", "진단"),
        ("02", "전략"),
        ("03", "포트폴리오"),   # current
        ("04", "디지털"),
        ("05", "조직"),
        ("06", "재무"),
    ]
    current_idx = 2

    tracker_x = 0.3
    tracker_y = 0.3
    tracker_w = 9.4
    tracker_h = 0.7
    gap = 0.08
    seg_w = (tracker_w - gap * (len(sections) - 1)) / len(sections)

    for i, (num, lbl) in enumerate(sections):
        sx = tracker_x + i * (seg_w + gap)
        if i < current_idx:
            fill_c = "grey_400"
            text_c = "white"
            border_c = "grey_400"
            border_w: float | None = None
        elif i == current_idx:
            fill_c = "accent"
            text_c = "white"
            border_c = "accent"
            border_w = None
        else:
            fill_c = "grey_100"
            text_c = "grey_700"
            border_c = "grey_400"
            border_w = 0.75

        c.box(x=sx, y=tracker_y, w=seg_w, h=tracker_h,
              fill=fill_c, border=border_w, border_color=border_c)
        c.text(num, x=sx + 0.1, y=tracker_y + 0.06,
               w=0.7, h=0.3, size=9, bold=True, color=text_c)
        c.text(lbl, x=sx + 0.1, y=tracker_y + 0.34,
               w=seg_w - 0.2, h=0.3, size=10, bold=True,
               color=text_c, anchor="top")

    # Progress bar underline
    total_progress = (current_idx + 1) / len(sections)
    c.box(x=tracker_x, y=tracker_y + tracker_h + 0.08,
          w=tracker_w, h=0.06, fill="grey_200", border=None)
    c.box(x=tracker_x, y=tracker_y + tracker_h + 0.08,
          w=tracker_w * total_progress, h=0.06,
          fill="accent", border=None)

    # --- Current section detail below ---
    detail_y = 1.5
    c.text("SECTION 03 — 사업 포트폴리오 재편",
           x=0.3, y=detail_y, w=9.4, h=0.35,
           size=10, bold=True, color="accent")
    c.text("성장·유지·철수 3축으로 본 핵심 사업의 우선순위",
           x=0.3, y=detail_y + 0.35, w=9.4, h=0.5,
           size=20, bold=True, color="grey_900", anchor="top")
    c.line(x1=0.3, y1=detail_y + 0.95, x2=9.7, y2=detail_y + 0.95,
           color="grey_200", width=0.75)

    # Three detail columns
    cols = [
        ("진단 결과",
         ["• 전체 12개 사업 중 5개가 저수익 구간",
          "• 시장 성장률 < 3% 사업이 매출의 38%",
          "• 핵심 3개 사업에 CAPEX 64% 집중 필요"]),
        ("주요 이슈",
         ["• 성장 사업 자본 배분 병목",
          "• 비핵심 사업 매각/축소 의사결정 지연",
          "• 신사업 Incubation 자원 부족"]),
        ("이 섹션의 결론",
         ["• Growth 3개: 5년 내 매출 2배",
          "• Maintain 4개: 현금 창출 극대화",
          "• Exit 5개: 2027년까지 단계 철수"]),
    ]
    col_x0 = 0.3
    col_w = 3.06
    col_gap = 0.1
    col_y = 2.9
    col_h = 3.9
    for i, (h1, bullets) in enumerate(cols):
        x = col_x0 + i * (col_w + col_gap)
        # Accent top strip (only current column highlighted)
        strip_c = "accent" if i == 2 else "grey_400"
        c.box(x=x, y=col_y, w=col_w, h=0.06,
              fill=strip_c, border=None)
        c.box(x=x, y=col_y + 0.06, w=col_w, h=col_h - 0.06,
              fill="white", border=0.75, border_color="grey_200")
        c.text(h1, x=x + 0.2, y=col_y + 0.2, w=col_w - 0.3, h=0.35,
               size=11, bold=True, color="grey_900")
        c.text("\n".join(bullets),
               x=x + 0.2, y=col_y + 0.65, w=col_w - 0.3, h=col_h - 0.8,
               size=9, color="grey_700", anchor="top")


# ============================================================
# F4. comp_executive_summary_scqa — 4 stacked bars
# ============================================================
def slide_f4_scqa(prs):
    s, c = make(prs,
                "Executive Summary — SCQA 프레임으로 본 핵심 메시지",
                "Situation → Complication → Question → Answer")

    # Layout
    top_y = 1.2
    bar_h = 1.45
    gap = 0.1
    x0, w = 0.4, 9.2
    label_w = 1.8

    sections = [
        ("SITUATION", "현황",
         "글로벌 조선·해양 수요는 회복세이나, 경쟁사 대비 수익성 격차는 "
         "2019년 이후 확대. 원가 구조 및 생산성에서 구조적 불리 영역 존재.",
         "grey_100", "grey_900", False),
        ("COMPLICATION", "문제",
         "원가 경쟁력 회복만으로는 부족 — 디지털 전환·인재 구조 전환까지 "
         "병행되지 않으면 2028년 이후 점유율 하락이 고착화될 위험.",
         "grey_400", "white", False),
        ("QUESTION", "질문",
         "그렇다면 향후 3년 간 어떤 우선순위로 투자를 집중해, "
         "수익성과 성장성을 동시에 회복할 수 있는가?",
         "accent_mid", "white", False),
        ("ANSWER", "답",
         "Core Platform 재구축(40%) + Growth 3대 사업 선택 집중(35%) "
         "+ 조직/인재 Target Operating Model 전환(25%)의 3축 동시 실행.",
         "accent", "white", True),
    ]

    for i, (lbl_en, lbl_ko, body, fill_c, text_c, is_bold) in enumerate(sections):
        y = top_y + i * (bar_h + gap)

        # Full bar
        c.box(x=x0, y=y, w=w, h=bar_h, fill=fill_c, border=None)

        # Left label zone (slightly darker overlay via text only — keep unified bar)
        c.text(lbl_en, x=x0 + 0.25, y=y + 0.2, w=label_w, h=0.3,
               size=10, bold=True, color=text_c)
        c.text(lbl_ko, x=x0 + 0.25, y=y + 0.55, w=label_w, h=0.5,
               size=22, bold=True, color=text_c,
               font="Georgia", anchor="top")

        # Vertical divider between label and body
        c.line(x1=x0 + label_w + 0.2, y1=y + 0.2,
               x2=x0 + label_w + 0.2, y2=y + bar_h - 0.2,
               color=text_c, width=0.5)

        # Body paragraph
        c.text(body,
               x=x0 + label_w + 0.4, y=y + 0.3,
               w=w - label_w - 0.6, h=bar_h - 0.45,
               size=12, bold=is_bold, color=text_c, anchor="top")


# ============================================================
# F5. comp_key_findings_summary — 2x2 grid with accent top strip
# ============================================================
def slide_f5_key_findings(prs):
    s, c = make(prs,
                "Key Findings — 이번 진단의 핵심 결론 4가지",
                "정량 근거 기반 요약, 세부 자료는 Appendix 참조")

    # 2x2 grid layout
    top_y = 1.3
    bottom_y = 7.1
    grid_h = bottom_y - top_y
    gap = 0.2

    box_w = (10 - 0.6 - gap) / 2
    box_h = (grid_h - gap) / 2

    findings = [
        ("1", "원가 경쟁력 격차는 구조적",
         ["선도 경쟁사 대비 톤당 원가 12% 높음 (2025 기준)",
          "노동·물류·소재 3대 항목에서 각 5%p 이상 격차",
          "단순 Cost-out으론 3년 내 회복 불가 — Core 혁신 필수"]),
        ("2", "디지털 전환 ROI는 조기에 가시화",
         ["DT Pilot 3개 사업부 평균 EBIT +2.3%p (12개월)",
          "AI 기반 예측 정비로 비가동 시간 22% 감소",
          "Core Platform 구축 시 전사 확산에 18개월 소요"]),
        ("3", "성장은 선택과 집중에서 온다",
         ["성장 사업 3개에 CAPEX 64% 집중 시 CAGR 9.5%",
          "비핵심 5개 사업 Exit 시 자본 $1.8B 회수",
          "선택 집중 없이는 2028년 점유율 2%p 추가 하락"]),
        ("4", "조직 재설계는 실행의 전제",
         ["현재 조직은 기능 중심 — 제품/고객 중심 재편 필요",
          "디지털·데이터 직군 내재화율 23% → 55% 목표",
          "Target Operating Model 없이는 전략 실행률 < 40%"]),
    ]

    for idx, (num, headline, bullets) in enumerate(findings):
        row = idx // 2
        col = idx % 2
        x = 0.3 + col * (box_w + gap)
        y = top_y + row * (box_h + gap)

        # Accent top strip
        c.box(x=x, y=y, w=box_w, h=0.1, fill="accent", border=None)

        # Main box
        c.box(x=x, y=y + 0.1, w=box_w, h=box_h - 0.1,
              fill="white", border=0.75, border_color="grey_200")

        # Finding number (accent, bold)
        c.text(num, x=x + 0.25, y=y + 0.22, w=0.7, h=0.75,
               size=32, bold=True, color="accent",
               font="Georgia", anchor="top")

        # Headline (bold)
        c.text(headline,
               x=x + 1.0, y=y + 0.3, w=box_w - 1.2, h=0.55,
               size=14, bold=True, color="grey_900", anchor="top")

        # Divider line
        c.line(x1=x + 1.0, y1=y + 0.95,
               x2=x + box_w - 0.2, y2=y + 0.95,
               color="grey_200", width=0.5)

        # Evidence bullets
        bullet_text = "\n".join(f"• {b}" for b in bullets)
        c.text(bullet_text,
               x=x + 1.0, y=y + 1.05,
               w=box_w - 1.2, h=box_h - 1.2,
               size=9, color="grey_700", anchor="top")


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_f1_agenda(prs)
    slide_f2_section_divider(prs)
    slide_f3_agenda_tracker(prs)
    slide_f4_scqa(prs)
    slide_f5_key_findings(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX saved: {pptx_path}")

    # Evaluate (Step 5 auto check)
    try:
        from ppt_builder.evaluate import evaluate_pptx, print_report
        report = evaluate_pptx(str(pptx_path))
        print_report(report)
    except Exception as e:
        print(f"[evaluate skipped] {e}")

    # PDF conversion (PowerPoint COM)
    pdf_path = OUTPUT_DIR / f"{NAME}.pdf"
    try:
        from ppt_builder.visual_validate import convert_pptx_to_pdf
        convert_pptx_to_pdf(pptx_path, pdf_path)
        print(f"PDF saved:  {pdf_path}")
    except Exception as e:
        print(f"[PDF conversion failed] {e}")

    # PNG extraction
    png_dir = OUTPUT_DIR / f"{NAME}_png"
    try:
        from ppt_builder.track_c.png_export import pptx_to_pngs
        pngs = pptx_to_pngs(pptx_path, png_dir)
        print(f"PNGs saved: {png_dir} ({len(pngs)} files)")
    except Exception as e:
        print(f"[PNG extraction failed] {e}")

    print(f"\nDone: {pptx_path}")


if __name__ == "__main__":
    main()
