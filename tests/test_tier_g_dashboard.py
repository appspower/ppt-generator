"""Tier G — Dashboard / KPI Components (5 slides).

G1. comp_kpi_card_sparkline   — 4 KPI cards with sparkline
G2. comp_traffic_light_status — 6 project RAG status table
G3. comp_bullet_chart_target  — 5 horizontal bullet charts
G4. comp_multi_gauge_grid     — 2x2 semi-circle gauges
G5. comp_scoreboard_ranking   — 8-item ranking with trend arrows

Output: output/tier_g_dashboard.pptx + PDF + PNGs
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "tier_g_dashboard"


def make(prs, title_text: str, subtitle: str = ""):
    """Create blank slide with standard header."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    c.text(title_text, x=0.4, y=0.25, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.68, w=9.2, h=0.25,
               size=9, color="grey_700")
    c.line(x1=0.4, y1=0.98, x2=9.6, y2=0.98, color="grey_200", width=1.0)
    return s, c


# ============================================================
# G1: comp_kpi_card_sparkline — 4 KPI cards + mini sparklines
# ============================================================
def slide_g1_kpi_card_sparkline(prs):
    s, c = make(prs,
                "핵심 KPI 대시보드 — 4개 지표 모두 전년 대비 개선",
                "KPI Card with Sparkline: label / big number / delta / 5-point trend")

    # 4 cards, each ~2.2" wide
    cards = [
        # (label, value, delta_text, delta_positive, trend[0..1 normalized 5 pts])
        ("월 매출",        "₩42.8B",  "+12.4%",  True,  [0.35, 0.48, 0.55, 0.72, 0.88]),
        ("신규 고객 수",   "1,284",   "+8.7%",   True,  [0.40, 0.52, 0.50, 0.68, 0.78]),
        ("고객 이탈률",    "2.3%",    "-0.8%p",  True,  [0.85, 0.72, 0.60, 0.48, 0.35]),
        ("평균 응대 시간", "1.8분",   "+0.3분",  False, [0.45, 0.55, 0.62, 0.78, 0.82]),
    ]

    card_w, card_h = 2.2, 2.6
    start_x, start_y = 0.45, 1.5
    gap = 0.15

    for idx, (label, value, delta, positive, trend) in enumerate(cards):
        x = start_x + idx * (card_w + gap)
        y = start_y

        # Card frame
        c.box(x=x, y=y, w=card_w, h=card_h,
              fill="white", border=1.0, border_color="grey_200")
        # Top accent stripe
        c.box(x=x, y=y, w=card_w, h=0.08,
              fill="accent", border=None)

        # Label (top, 8pt grey)
        c.text(label, x=x + 0.15, y=y + 0.18, w=card_w - 0.3, h=0.25,
               size=8, bold=True, color="grey_700")

        # BIG number (28pt bold)
        c.text(value, x=x + 0.15, y=y + 0.5, w=card_w - 0.3, h=0.75,
               size=28, bold=True, color="grey_900", anchor="top")

        # Delta arrow + %
        arrow_char = "▲" if positive else "▼"
        delta_color = "positive" if positive else "negative"
        c.text(f"{arrow_char} {delta}", x=x + 0.15, y=y + 1.25, w=card_w - 0.3, h=0.3,
               size=11, bold=True, color=delta_color)
        c.text("vs 전월", x=x + 0.15, y=y + 1.55, w=card_w - 0.3, h=0.22,
               size=7, color="grey_400")

        # Separator
        c.line(x1=x + 0.15, y1=y + 1.85, x2=x + card_w - 0.15, y2=y + 1.85,
               color="grey_200", width=0.5)

        # Sparkline: 5 points connected, with circles
        spark_x = x + 0.2
        spark_y_base = y + 2.45  # bottom of sparkline area
        spark_height = 0.45
        spark_w = card_w - 0.4
        n_pts = len(trend)

        pts = []
        for i, v in enumerate(trend):
            px = spark_x + (i / (n_pts - 1)) * spark_w
            py = spark_y_base - v * spark_height
            pts.append((px, py))

        # Connecting lines
        line_color = "accent" if positive else "negative"
        for i in range(len(pts) - 1):
            c.line(x1=pts[i][0], y1=pts[i][1],
                   x2=pts[i + 1][0], y2=pts[i + 1][1],
                   color=line_color, width=1.5)

        # Circles on points
        dot_d = 0.09
        for px, py in pts:
            c.circle(x=px - dot_d / 2, y=py - dot_d / 2, d=dot_d,
                     fill=line_color, border=None)
        # Last point emphasis
        last_px, last_py = pts[-1]
        c.circle(x=last_px - 0.07, y=last_py - 0.07, d=0.14,
                 fill="white", border=1.5, border_color=line_color)

    # Footer
    c.text("※ Sparkline = 최근 5개월 추이 | Data: Internal BI (2026-03 기준)",
           x=0.4, y=6.9, w=9.2, h=0.25, size=7, color="grey_400")


# ============================================================
# G2: comp_traffic_light_status — 6 projects RAG
# ============================================================
def slide_g2_traffic_light_status(prs):
    s, c = make(prs,
                "프로젝트 RAG 상태 — 6개 과제 중 2개 적색경보 (긴급 개입 필요)",
                "Traffic Light Status: Schedule / Budget / Risk 3-축 평가")

    projects = [
        # (name, schedule, budget, risk, note)
        ("SAP S/4HANA 전환",       "G", "G", "A", "파일럿 완료, 3Q 확산 예정"),
        ("차세대 MES 구축",         "G", "A", "G", "예산 초과 가능성 (+3.2%)"),
        ("데이터 레이크 마이그",    "A", "G", "A", "일정 2주 지연, 리소스 재배치"),
        ("CRM 통합 플랫폼",         "R", "A", "R", "핵심 이해관계자 이슈 발생"),
        ("RPA Center of Excel.",   "G", "G", "G", "예정 대비 120% 달성"),
        ("글로벌 IAM 통합",         "R", "R", "A", "벤더 이슈, CTO 에스컬레이션"),
    ]

    # Table geometry
    table_x = 0.5
    table_y = 1.4
    row_h = 0.62
    col_widths = [0.35, 2.4, 0.9, 0.9, 0.9, 3.05]  # #, Name, S, B, R, Note
    headers = ["#", "프로젝트명", "Schedule", "Budget", "Risk", "상태 요약"]

    # Header row
    hx = table_x
    c.box(x=table_x, y=table_y, w=sum(col_widths), h=0.45,
          fill="grey_900", border=None)
    for i, (header, cw) in enumerate(zip(headers, col_widths)):
        c.text(header, x=hx + 0.1, y=table_y, w=cw - 0.1, h=0.45,
               size=9, bold=True, color="white",
               align="center" if i != 1 and i != 5 else "left",
               anchor="middle")
        hx += cw

    # Status color map
    status_map = {
        "G": ("positive", "정상"),
        "A": ("warning",  "주의"),
        "R": ("negative", "위험"),
    }

    # Rows
    for r_idx, (name, sch, bud, risk, note) in enumerate(projects):
        ry = table_y + 0.45 + r_idx * row_h
        # Alt row shading
        if r_idx % 2 == 1:
            c.box(x=table_x, y=ry, w=sum(col_widths), h=row_h,
                  fill="grey_100", border=None)

        cx = table_x

        # #
        c.text(str(r_idx + 1), x=cx, y=ry, w=col_widths[0], h=row_h,
               size=9, bold=True, color="grey_700",
               align="center", anchor="middle")
        cx += col_widths[0]

        # Name
        c.text(name, x=cx + 0.1, y=ry, w=col_widths[1] - 0.1, h=row_h,
               size=10, bold=True, color="grey_900", anchor="middle")
        cx += col_widths[1]

        # Status circles (Schedule / Budget / Risk)
        for status in (sch, bud, risk):
            color_name, _ = status_map[status]
            cw = col_widths[2]  # all equal
            # Centered circle in cell
            dot_d = 0.32
            dot_x = cx + (cw - dot_d) / 2
            dot_y = ry + (row_h - dot_d) / 2
            c.circle(x=dot_x, y=dot_y, d=dot_d,
                     fill=color_name, border=None,
                     text=status, text_color="white",
                     text_size=10, text_bold=True)
            cx += cw

        # Note
        c.text(note, x=cx + 0.1, y=ry, w=col_widths[5] - 0.15, h=row_h,
               size=9, color="grey_700", anchor="middle")

        # Row bottom border
        c.line(x1=table_x, y1=ry + row_h,
               x2=table_x + sum(col_widths), y2=ry + row_h,
               color="grey_200", width=0.5)

    # Legend
    ly = table_y + 0.45 + len(projects) * row_h + 0.25
    c.text("범례:", x=table_x, y=ly, w=0.6, h=0.25, size=8, bold=True, color="grey_700")
    legend_items = [
        ("positive", "G", "정상 진행"),
        ("warning",  "A", "주의 필요"),
        ("negative", "R", "위험 / 개입"),
    ]
    lx = table_x + 0.7
    for col, code, label in legend_items:
        c.circle(x=lx, y=ly, d=0.25, fill=col, border=None,
                 text=code, text_color="white", text_size=9, text_bold=True)
        c.text(label, x=lx + 0.32, y=ly, w=1.4, h=0.25,
               size=8, color="grey_700", anchor="middle")
        lx += 1.75


# ============================================================
# G3: comp_bullet_chart_target — 5 horizontal bullet charts
# ============================================================
def slide_g3_bullet_chart_target(prs):
    s, c = make(prs,
                "KPI 목표 대비 실적 — 5개 지표 중 3개 목표 초과 달성",
                "Bullet Chart: 배경 밴드(Poor/Acceptable/Good) + 실적 막대 + 목표선")

    # (label, actual, target, max_val, unit)
    # ranges: poor 0~40%, acceptable 40~70%, good 70~100% of max
    kpis = [
        ("매출 성장률",      11.8, 10.0, 15.0, "%"),
        ("영업이익률",       14.2, 15.0, 20.0, "%"),
        ("고객 만족도 (NPS)", 62,   55,   80,   ""),
        ("직원 몰입도",      78,   75,   100,  "점"),
        ("Time-to-Market",    4.2,  5.0,  8.0,  "개월"),
    ]

    chart_x = 2.6
    chart_w = 5.8
    chart_start_y = 1.4
    row_h = 1.0
    bar_h = 0.38
    bg_h = 0.6  # background band height

    for idx, (label, actual, target, maxv, unit) in enumerate(kpis):
        row_y = chart_start_y + idx * row_h

        # Label on the left
        c.text(label, x=0.4, y=row_y + 0.1, w=2.1, h=0.3,
               size=10, bold=True, color="grey_900")
        # Sub: current / target
        c.text(f"실적 {actual}{unit} / 목표 {target}{unit}",
               x=0.4, y=row_y + 0.38, w=2.1, h=0.25,
               size=8, color="grey_700")

        # Background bands (Poor 0-40%, Acceptable 40-70%, Good 70-100%)
        bg_y = row_y + 0.1
        poor_w = chart_w * 0.4
        acc_w = chart_w * 0.3
        good_w = chart_w * 0.3

        c.box(x=chart_x, y=bg_y, w=poor_w, h=bg_h,
              fill="grey_200", border=None)
        c.box(x=chart_x + poor_w, y=bg_y, w=acc_w, h=bg_h,
              fill="grey_400", border=None)
        c.box(x=chart_x + poor_w + acc_w, y=bg_y, w=good_w, h=bg_h,
              fill="grey_700", border=None)

        # Actual value bar (dark accent), centered vertically in band
        actual_w = chart_w * (actual / maxv)
        bar_y = bg_y + (bg_h - bar_h) / 2
        c.box(x=chart_x, y=bar_y, w=actual_w, h=bar_h,
              fill="accent", border=None)

        # Target line (vertical red)
        target_x = chart_x + chart_w * (target / maxv)
        c.line(x1=target_x, y1=bg_y - 0.05,
               x2=target_x, y2=bg_y + bg_h + 0.05,
               color="negative", width=2.5)

        # Scale min/max labels
        c.text("0", x=chart_x - 0.15, y=bg_y + bg_h + 0.02,
               w=0.3, h=0.2, size=7, color="grey_400", align="center")
        c.text(f"{maxv}{unit}", x=chart_x + chart_w - 0.35,
               y=bg_y + bg_h + 0.02, w=0.7, h=0.2,
               size=7, color="grey_400", align="center")

    # Top legend
    ly = 6.6
    legend = [
        ("grey_200", "Poor"),
        ("grey_400", "Acceptable"),
        ("grey_700", "Good"),
        ("accent",   "실적 (Actual)"),
        ("negative", "목표 (Target)"),
    ]
    lx = 0.5
    for col, lbl in legend:
        if lbl == "목표 (Target)":
            # draw as vertical line marker
            c.line(x1=lx + 0.1, y1=ly, x2=lx + 0.1, y2=ly + 0.25,
                   color="negative", width=2.5)
        else:
            c.box(x=lx, y=ly, w=0.25, h=0.25, fill=col, border=None)
        c.text(lbl, x=lx + 0.32, y=ly, w=1.5, h=0.25,
               size=8, color="grey_700", anchor="middle")
        lx += 1.85


# ============================================================
# G4: comp_multi_gauge_grid — 2x2 semi-circle gauges (approximated)
# ============================================================
def slide_g4_multi_gauge_grid(prs):
    s, c = make(prs,
                "핵심 성과 게이지 — 4개 영역 중 3개 녹색 존 진입",
                "Multi-Gauge Grid: 반원 게이지 (3 color segments + needle + center %)")

    # (label, sub, value_pct, accent_color)
    gauges = [
        ("운영 효율성",   "Operational Eff.",     82, "positive"),
        ("고객 만족도",   "Customer Satisfaction", 71, "positive"),
        ("예산 집행률",   "Budget Execution",     58, "warning"),
        ("혁신 성숙도",   "Innovation Maturity",  39, "negative"),
    ]

    # 2x2 layout
    positions = [
        (1.2, 1.3),   # top-left
        (5.8, 1.3),   # top-right
        (1.2, 4.3),   # bottom-left
        (5.8, 4.3),   # bottom-right
    ]

    gauge_w = 3.0
    gauge_h = 2.6

    for (label, sub, value, _), (gx, gy) in zip(gauges, positions):
        # Frame box
        c.box(x=gx, y=gy, w=gauge_w, h=gauge_h,
              fill="white", border=1.0, border_color="grey_200")

        # Center of semicircle
        cx_mid = gx + gauge_w / 2
        # Semicircle base y (horizontal line)
        base_y = gy + 1.65
        # Outer radius
        R = 1.1
        # Inner (white) radius to approximate donut
        r_inner = 0.55

        # Approximated semicircle by 3 arc-like boxes using chord segments
        # We draw 3 bands as pie-slice-like rectangles rotated.
        # Simpler: use 3 sectors via overlapping circles + mask.
        # Practical approach: 3 thick arcs approximated by bars laid in 3 angular zones.
        # Use MSO_SHAPE.BLOCK_ARC if available; otherwise approximate with colored pie wedges
        # via 3 separate thick curves — we approximate using 3 rotated rectangle chunks.

        # Background semicircle: draw big circle then cover bottom half with white
        c.circle(x=cx_mid - R, y=base_y - R, d=2 * R,
                 fill="grey_200", border=None)
        # Cover bottom half to make semicircle
        c.box(x=cx_mid - R - 0.05, y=base_y,
              w=2 * R + 0.1, h=R + 0.1,
              fill="white", border=None)

        # Colored zones: split semicircle into 3 by drawing sector-like overlays.
        # We'll approximate 3 arc segments by placing 3 rectangles at rotation angles
        # of -60°, 0°, +60° (each sector = 60°). Use thin chord rectangles.
        # Alternative pragmatic: draw 3 thin colored wedges as rotated rectangles.
        zone_colors = ["negative", "warning", "positive"]
        # 3 rotated rectangles simulating sectors by using narrow thick "spokes"
        # covering 60° each. We'll use wide colored arcs drawn as chords (polygons).
        # To keep simple & portable: draw 3 colored slim rectangles as arc-approximations
        # positioned at -60/-30°, -30/+30°, +30/+60° angles.
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches as _Inches, Pt as _Pt
        import math
        for k, col in enumerate(zone_colors):
            # Angle center: -60° (left), 0° (top), +60° (right)
            angle_center_deg = -60 + k * 60
            ang = math.radians(angle_center_deg)
            # Radial segment rectangle: from r_inner to R along direction ang
            seg_len = R - r_inner
            seg_mid_r = (R + r_inner) / 2
            # Segment center coords (relative to cx_mid, base_y)
            sx = cx_mid + seg_mid_r * math.sin(ang)
            sy = base_y - seg_mid_r * math.cos(ang)
            # Rectangle approximating arc chord, thickness seg_len, width ~ chord
            seg_w = 2 * seg_mid_r * math.sin(math.radians(30)) * 1.15
            seg_h_loc = seg_len * 1.15
            rect = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                _Inches(sx - seg_w / 2), _Inches(sy - seg_h_loc / 2),
                _Inches(seg_w), _Inches(seg_h_loc),
            )
            rect.rotation = float(angle_center_deg)
            rect.fill.solid()
            from ppt_builder.primitives import color as _color_fn
            rect.fill.fore_color.rgb = _color_fn(col)
            rect.line.fill.background()

        # Cover inner area to create donut effect
        c.circle(x=cx_mid - r_inner, y=base_y - r_inner, d=2 * r_inner,
                 fill="white", border=None)
        # Re-cover bottom half after inner
        c.box(x=cx_mid - R - 0.05, y=base_y + 0.001,
              w=2 * R + 0.1, h=R + 0.1,
              fill="white", border=None)

        # Needle — line from center to position determined by value
        needle_angle_deg = -90 + (value / 100.0) * 180  # -90 (left) to +90 (right)
        nang = math.radians(needle_angle_deg)
        needle_len = R - 0.08
        nx_end = cx_mid + needle_len * math.sin(nang)
        ny_end = base_y - needle_len * math.cos(nang)
        c.line(x1=cx_mid, y1=base_y, x2=nx_end, y2=ny_end,
               color="grey_900", width=2.5)
        # Needle hub
        c.circle(x=cx_mid - 0.08, y=base_y - 0.08, d=0.16,
                 fill="grey_900", border=None)

        # BIG % in center (below semicircle)
        c.text(f"{value}%", x=gx, y=base_y + 0.05, w=gauge_w, h=0.5,
               size=26, bold=True, color="grey_900",
               align="center", anchor="top")

        # Label
        c.text(label, x=gx, y=base_y + 0.55, w=gauge_w, h=0.28,
               size=11, bold=True, color="grey_900", align="center")
        c.text(sub, x=gx, y=base_y + 0.83, w=gauge_w, h=0.22,
               size=8, color="grey_400", align="center")

    # Scale legend on bottom
    c.text("게이지 구간: 0~33% (위험) | 33~66% (주의) | 66~100% (양호)",
           x=0.4, y=7.0, w=9.2, h=0.25, size=8, color="grey_700", align="center")


# ============================================================
# G5: comp_scoreboard_ranking — 8-item ranking table
# ============================================================
def slide_g5_scoreboard_ranking(prs):
    s, c = make(prs,
                "부문별 성과 순위 — IT/Digital 부문 3기 연속 1위 수성",
                "Scoreboard: Rank / Name / Score / Trend / Status Bar")

    # (name, score, trend ['up'/'down'/'flat'], status_pct 0~1)
    items = [
        ("IT & Digital 본부",       94.2, "up",   0.94),
        ("재무 본부",                91.8, "up",   0.92),
        ("영업 본부 (B2B)",          88.5, "flat", 0.89),
        ("SCM 본부",                 85.1, "up",   0.85),
        ("생산기술 본부",            82.4, "down", 0.82),
        ("인사/조직 본부",           78.9, "flat", 0.79),
        ("마케팅 본부",              74.3, "down", 0.74),
        ("고객서비스 본부",          68.7, "down", 0.69),
    ]

    table_x = 0.5
    table_y = 1.3
    col_widths = [0.7, 3.3, 1.1, 1.0, 2.9]  # Rank, Name, Score, Trend, Bar
    headers = ["순위", "부문", "Score", "Trend", "달성도"]

    # Header
    c.box(x=table_x, y=table_y, w=sum(col_widths), h=0.42,
          fill="grey_900", border=None)
    hx = table_x
    for header, cw in zip(headers, col_widths):
        c.text(header, x=hx, y=table_y, w=cw, h=0.42,
               size=9, bold=True, color="white",
               align="center", anchor="middle")
        hx += cw

    row_h = 0.6

    for r_idx, (name, score, trend, pct) in enumerate(items):
        ry = table_y + 0.42 + r_idx * row_h
        rank = r_idx + 1
        is_top = rank <= 3

        # Top-3 accent band (full row)
        if rank == 1:
            c.box(x=table_x, y=ry, w=sum(col_widths), h=row_h,
                  fill="zone_alert", border=None)
        elif is_top:
            c.box(x=table_x, y=ry, w=sum(col_widths), h=row_h,
                  fill="grey_100", border=None)

        cx = table_x

        # Rank: accent circle for top 3, plain text otherwise
        if is_top:
            rank_color = "accent" if rank == 1 else "accent_mid" if rank == 2 else "grey_700"
            dot_d = 0.38
            dot_x = cx + (col_widths[0] - dot_d) / 2
            dot_y = ry + (row_h - dot_d) / 2
            c.circle(x=dot_x, y=dot_y, d=dot_d,
                     fill=rank_color, border=None,
                     text=str(rank), text_color="white",
                     text_size=12, text_bold=True)
        else:
            c.text(str(rank), x=cx, y=ry, w=col_widths[0], h=row_h,
                   size=12, bold=True, color="grey_700",
                   align="center", anchor="middle")
        cx += col_widths[0]

        # Name
        name_bold = is_top
        c.text(name, x=cx + 0.15, y=ry, w=col_widths[1] - 0.15, h=row_h,
               size=11, bold=name_bold, color="grey_900", anchor="middle")
        cx += col_widths[1]

        # Score
        c.text(f"{score:.1f}", x=cx, y=ry, w=col_widths[2], h=row_h,
               size=13, bold=True, color="grey_900",
               align="center", anchor="middle")
        cx += col_widths[2]

        # Trend arrow
        trend_map = {
            "up":   ("▲", "positive"),
            "down": ("▼", "negative"),
            "flat": ("→", "grey_400"),
        }
        arrow_char, arrow_color = trend_map[trend]
        c.text(arrow_char, x=cx, y=ry, w=col_widths[3], h=row_h,
               size=16, bold=True, color=arrow_color,
               align="center", anchor="middle")
        cx += col_widths[3]

        # Status bar
        bar_margin = 0.15
        bar_full_w = col_widths[4] - bar_margin * 2
        bar_h = 0.22
        bar_y = ry + (row_h - bar_h) / 2
        # Background
        c.box(x=cx + bar_margin, y=bar_y, w=bar_full_w, h=bar_h,
              fill="grey_200", border=None)
        # Fill
        fill_col = "positive" if pct >= 0.85 else "warning" if pct >= 0.75 else "negative"
        c.box(x=cx + bar_margin, y=bar_y, w=bar_full_w * pct, h=bar_h,
              fill=fill_col, border=None)
        # % text on right
        c.text(f"{int(pct * 100)}%",
               x=cx + bar_margin + bar_full_w + 0.05, y=ry,
               w=0.55, h=row_h,
               size=8, bold=True, color="grey_700", anchor="middle")

        # Row divider
        c.line(x1=table_x, y1=ry + row_h,
               x2=table_x + sum(col_widths), y2=ry + row_h,
               color="grey_200", width=0.5)

    # Footer
    c.text("※ Score = 전사 공통 KPI (재무 40% + 고객 30% + 운영 20% + 조직 10%) | "
           "Trend = 전분기 대비 변동",
           x=0.4, y=table_y + 0.42 + len(items) * row_h + 0.2,
           w=9.2, h=0.3, size=7, color="grey_400")


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_g1_kpi_card_sparkline(prs)
    slide_g2_traffic_light_status(prs)
    slide_g3_bullet_chart_target(prs)
    slide_g4_multi_gauge_grid(prs)
    slide_g5_scoreboard_ranking(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX saved: {pptx_path} ({len(prs.slides)} slides)")

    # PDF conversion (Windows + PowerPoint COM)
    try:
        import pythoncom
        import win32com.client
        pythoncom.CoInitialize()
        pdf_path = pptx_path.with_suffix(".pdf")
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        p = ppt.Presentations.Open(str(pptx_path.resolve()), WithWindow=False)
        p.SaveAs(str(pdf_path.resolve()), 32)  # 32 = ppSaveAsPDF
        p.Close()
        print(f"PDF saved:  {pdf_path}")
    except Exception as e:
        print(f"PDF conversion skipped: {e}")

    # PNG export
    try:
        from ppt_builder.track_c.png_export import pptx_to_pngs
        png_dir = OUTPUT_DIR / f"{NAME}_pngs"
        paths = pptx_to_pngs(pptx_path, png_dir)
        print(f"PNGs saved: {png_dir} ({len(paths)} images)")
    except Exception as e:
        print(f"PNG export skipped: {e}")


if __name__ == "__main__":
    main()
