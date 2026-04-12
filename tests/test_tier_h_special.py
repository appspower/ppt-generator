"""Tier H — Special Visualization Components (5 slides).

H1. comp_timeline_with_images       — Horizontal timeline + 5 image events
H2. comp_funnel_with_conversion     — Vertical funnel + conversion arrows
H3. comp_pyramid_maturity           — 4-level maturity pyramid + characteristics
H4. comp_split_comparison_asymmetric — AS-IS / Transition / TO-BE (40/20/40)
H5. comp_venn_3_circle              — 3-circle Venn with intersection labels

Output: output/tier_h_special.pptx (+ PDF + per-slide PNGs).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ppt_builder.primitives import Canvas

OUTPUT_DIR = Path(__file__).parent.parent / "output"
NAME = "tier_h_special"


def make(prs, title_text: str, subtitle: str = ""):
    """Create a blank slide with a consistent header."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    c = Canvas(s)
    c.box(x=0, y=0, w=10, h=0.08, fill="accent", border=None)
    c.text(title_text, x=0.4, y=0.22, w=9.2, h=0.45,
           size=16, bold=True, color="grey_900")
    if subtitle:
        c.text(subtitle, x=0.4, y=0.66, w=9.2, h=0.25,
               size=9, color="grey_700")
    c.line(x1=0.4, y1=0.96, x2=9.6, y2=0.96, color="grey_200", width=1.0)
    return s, c


# ============================================================
# H1. comp_timeline_with_images — 5 events with image placeholders
# ============================================================
def slide_h1_timeline_with_images(prs):
    s, c = make(prs,
                "디지털 전환 5년 로드맵 — 주요 마일스톤 시각화",
                "Horizontal timeline with image placeholders per milestone")

    events = [
        ("2022", "Foundation",
         "기초 인프라 구축\nERP 마이그레이션 착수"),
        ("2023", "Cloud Migration",
         "레거시 전환 완료\nAWS/Azure 이전"),
        ("2024", "Data Platform",
         "데이터 레이크 구축\n실시간 분석 기반"),
        ("2025", "AI Integration",
         "생성형 AI 도입\n업무 자동화 30%"),
        ("2026", "Optimization",
         "풀 디지털 네이티브\nROI 150% 달성"),
    ]

    n = len(events)
    img_w, img_h = 1.55, 1.2
    # Distribute centers evenly across 9.2in width
    left_margin = 0.4
    usable_w = 9.2
    step = usable_w / n
    bar_y = 1.6 + img_h + 0.25  # timeline bar below images

    # Horizontal timeline bar (full span)
    c.box(x=left_margin, y=bar_y, w=usable_w, h=0.06,
          fill="grey_400", border=None)

    for i, (year, title, desc) in enumerate(events):
        cx = left_margin + step * i + step / 2
        img_x = cx - img_w / 2
        img_y = 1.4

        # Image placeholder: grey_400 rectangle with "Image" text
        c.box(x=img_x, y=img_y, w=img_w, h=img_h,
              fill="grey_400", border=0.75, border_color="grey_700")
        c.text("Image", x=img_x, y=img_y, w=img_w, h=img_h,
               size=11, bold=True, color="white",
               align="center", anchor="middle")

        # Connector from image bottom to timeline bar
        c.line(x1=cx, y1=img_y + img_h, x2=cx, y2=bar_y,
               color="grey_700", width=1.0)

        # Year dot on timeline
        c.circle(x=cx - 0.11, y=bar_y - 0.08, d=0.22,
                 fill="accent", border=None)

        # Year label below timeline
        c.text(year, x=cx - 0.6, y=bar_y + 0.12, w=1.2, h=0.28,
               size=11, bold=True, color="accent",
               align="center", anchor="top")

        # Event title (bold)
        c.text(title, x=cx - (step / 2) + 0.05, y=bar_y + 0.45,
               w=step - 0.1, h=0.3,
               size=10, bold=True, color="grey_900",
               align="center", anchor="top")

        # Description (2 lines)
        c.text(desc, x=cx - (step / 2) + 0.05, y=bar_y + 0.78,
               w=step - 0.1, h=0.7,
               size=8, color="grey_700",
               align="center", anchor="top")

    c.text("각 마일스톤은 분기별 체크포인트로 분해되어 월 단위 트래킹 — "
           "지연 시 Steering Committee 에스컬레이션",
           x=0.4, y=6.9, w=9.2, h=0.3,
           size=8, color="grey_400", align="center")


# ============================================================
# H2. comp_funnel_with_conversion — 5-stage vertical funnel
# ============================================================
def slide_h2_funnel_with_conversion(prs):
    s, c = make(prs,
                "Lead → Customer 전환 퍼널 — 단계별 전환율 분석",
                "5-stage funnel with inter-stage conversion rates")

    stages = [
        ("Awareness",   100_000, 100, "grey_400"),
        ("Interest",     45_000,  45, "grey_700"),
        ("Consideration",18_000,  18, "accent_mid"),
        ("Intent",        7_200,   7, "accent"),
        ("Customer",      2_880,   3, "grey_900"),
    ]
    # Conversion percentages between stages (next/prev)
    conv_rates = ["45%↓", "40%↓", "40%↓", "40%↓"]

    # Funnel geometry: decreasing widths, centered on slide
    top_w = 6.0
    bot_w = 3.2
    stage_h = 0.72
    gap = 0.32           # vertical space for conversion arrow
    start_y = 1.25
    cx = 3.0             # left anchor x of funnel center area
    center_line = 5.0    # slide center x

    n = len(stages)
    for i, (name, count, pct, fill_col) in enumerate(stages):
        # Linear interpolation of width from top_w -> bot_w
        t = i / (n - 1)
        w = top_w - (top_w - bot_w) * t
        x = center_line - w / 2
        y = start_y + i * (stage_h + gap)

        # Stage box
        text_color = "white" if fill_col in ("grey_700", "accent_mid",
                                             "accent", "grey_900") else "grey_900"
        c.box(x=x, y=y, w=w, h=stage_h, fill=fill_col, border=None)

        # Stage content: name | count | % of original
        c.text(name, x=x + 0.15, y=y, w=w * 0.45, h=stage_h,
               size=11, bold=True, color=text_color, anchor="middle")
        c.text(f"{count:,}", x=x + w * 0.45, y=y, w=w * 0.3, h=stage_h,
               size=13, bold=True, color=text_color,
               align="center", anchor="middle")
        c.text(f"{pct}%", x=x + w * 0.75, y=y, w=w * 0.22, h=stage_h,
               size=10, bold=True, color=text_color,
               align="right", anchor="middle")

        # Conversion arrow between stages
        if i < n - 1:
            arrow_y = y + stage_h + 0.02
            # Red accent arrow + label
            c.arrow(x1=center_line, y1=arrow_y,
                    x2=center_line, y2=arrow_y + gap - 0.06,
                    color="negative", width=1.5)
            c.text(conv_rates[i],
                   x=center_line + 0.2, y=arrow_y,
                   w=1.2, h=gap,
                   size=10, bold=True, color="negative",
                   align="left", anchor="middle")

    # Side legend
    c.text("단계별 전환율",
           x=7.5, y=1.25, w=2.0, h=0.3,
           size=10, bold=True, color="grey_900")
    c.text("• 100K 리드 → 2.9K 고객\n"
           "• 최종 전환율: 2.88%\n"
           "• 업계 평균(1.5%) 대비 +1.38pp\n"
           "• 최대 병목: Awareness→Interest\n"
           "  (55% 이탈)",
           x=7.5, y=1.6, w=2.2, h=3.5,
           size=8, color="grey_700", anchor="top")

    c.text("Source: CRM Pipeline 2026 Q1, 90-day cohort",
           x=0.4, y=6.9, w=9.2, h=0.3,
           size=7, color="grey_400")


# ============================================================
# H3. comp_pyramid_maturity — 4-level pyramid + characteristics
# ============================================================
def slide_h3_pyramid_maturity(prs):
    s, c = make(prs,
                "디지털 성숙도 모델 — Reactive에서 Optimized까지 4단계 여정",
                "4-level maturity pyramid with per-level characteristics")

    # Pyramid geometry (left side)
    levels = [
        ("Reactive",    "grey_400",
         ["수동 대응 중심", "사일로화된 데이터", "이슈 발생 후 조치"]),
        ("Proactive",   "grey_700",
         ["예측 지표 도입", "부서 간 데이터 공유", "정기 모니터링"]),
        ("Integrated",  "accent_mid",
         ["E2E 통합 플랫폼", "실시간 대시보드", "크로스펑션 KPI"]),
        ("Optimized",   "accent",
         ["AI 기반 자동화", "연속 개선 루프", "생태계 확장"]),
    ]

    # Pyramid base at bottom, apex at top
    base_x_center = 2.6
    base_y = 6.4     # bottom edge of base level
    base_w = 4.4
    apex_w = 1.4
    level_h = 1.05

    n = len(levels)
    for i, (label, fill_col, chars) in enumerate(levels):
        # i=0 is base (Reactive), i=n-1 is apex (Optimized)
        # Interpolate width from base_w (bottom) to apex_w (top)
        t = i / (n - 1)
        w = base_w - (base_w - apex_w) * t
        x = base_x_center - w / 2
        y = base_y - (i + 1) * level_h  # stack upward

        text_color = "white" if fill_col != "grey_400" else "grey_900"
        c.box(x=x, y=y, w=w, h=level_h, fill=fill_col, border=None)
        c.text(label, x=x, y=y, w=w, h=level_h,
               size=13, bold=True, color=text_color,
               align="center", anchor="middle")

        # Characteristics list on the right, aligned with this level
        char_x = 5.5
        c.text(f"Level {i + 1} — {label}",
               x=char_x, y=y + 0.05, w=4.0, h=0.3,
               size=10, bold=True, color=fill_col if fill_col != "grey_400" else "grey_900")
        bullets = "\n".join(f"•  {ch}" for ch in chars)
        c.text(bullets, x=char_x, y=y + 0.35, w=4.0, h=level_h - 0.35,
               size=9, color="grey_700", anchor="top")

    # Upward arrow on the left
    c.arrow(x1=0.8, y1=base_y - 0.1, x2=0.8,
            y2=base_y - n * level_h + 0.1,
            color="accent", width=2.0)
    c.text("성숙도\n↑", x=0.55, y=base_y - n * level_h - 0.4,
           w=0.8, h=0.4, size=9, bold=True, color="accent",
           align="center", anchor="middle")

    c.text("각 단계는 평균 12~18개월의 전환 기간 소요 — "
           "건너뛰기(skip-level) 전략은 실패율 70% 이상",
           x=0.4, y=6.8, w=9.2, h=0.3,
           size=8, color="grey_700", align="center")


# ============================================================
# H4. comp_split_comparison_asymmetric — AS-IS / Transition / TO-BE
# ============================================================
def slide_h4_split_comparison_asymmetric(prs):
    s, c = make(prs,
                "AS-IS vs TO-BE — 전환 전후 핵심 변화 대조",
                "40% / 20% / 40% asymmetric split with transition strip")

    top_y = 1.3
    panel_h = 5.2
    left_w = 3.68   # 40%
    center_w = 1.84 # 20%
    right_w = 3.68  # 40%
    left_x = 0.4
    center_x = left_x + left_w
    right_x = center_x + center_w

    # --- Left panel: AS-IS ---
    c.box(x=left_x, y=top_y, w=left_w, h=panel_h,
          fill="grey_100", border=0.75, border_color="grey_400")
    c.text("AS-IS", x=left_x, y=top_y + 0.15, w=left_w, h=0.45,
           size=16, bold=True, color="grey_900", align="center")
    c.text("현재 운영의 한계",
           x=left_x, y=top_y + 0.6, w=left_w, h=0.3,
           size=10, color="grey_700", align="center")

    issues = [
        "수작업 데이터 집계로 월말 마감 5일 소요",
        "시스템 간 불일치로 재작업률 18%",
        "의사결정 지연 (평균 리드타임 14일)",
        "부서별 KPI 사일로 — 전사 가시성 부재",
    ]
    for i, txt in enumerate(issues):
        y = top_y + 1.15 + i * 0.9
        # ✗ symbol
        c.text("✗", x=left_x + 0.25, y=y, w=0.4, h=0.45,
               size=18, bold=True, color="negative",
               align="center", anchor="middle")
        c.text(txt, x=left_x + 0.7, y=y, w=left_w - 0.85, h=0.8,
               size=9, color="grey_900", anchor="middle")

    # --- Center strip: Transition ---
    c.box(x=center_x, y=top_y, w=center_w, h=panel_h,
          fill="accent", border=None)
    # Big arrow glyph in center
    c.text("▶", x=center_x, y=top_y + 1.6, w=center_w, h=1.6,
           size=60, bold=True, color="white",
           align="center", anchor="middle")
    c.text("Transition",
           x=center_x, y=top_y + 3.3, w=center_w, h=0.4,
           size=14, bold=True, color="white", align="center")
    c.text("18개월\n디지털 전환",
           x=center_x, y=top_y + 3.75, w=center_w, h=0.8,
           size=10, color="white", align="center", anchor="top")

    # --- Right panel: TO-BE ---
    c.box(x=right_x, y=top_y, w=right_w, h=panel_h,
          fill="grey_700", border=None)
    c.text("TO-BE", x=right_x, y=top_y + 0.15, w=right_w, h=0.45,
           size=16, bold=True, color="white", align="center")
    c.text("전환 후 목표 상태",
           x=right_x, y=top_y + 0.6, w=right_w, h=0.3,
           size=10, color="grey_200", align="center")

    benefits = [
        "실시간 자동 집계 — 마감 1일 (80% 단축)",
        "Single Source of Truth — 재작업률 3% 이하",
        "대시보드 기반 의사결정 — 리드타임 2일",
        "전사 통합 KPI Tower — 실시간 가시성 확보",
    ]
    for i, txt in enumerate(benefits):
        y = top_y + 1.15 + i * 0.9
        c.text("✓", x=right_x + 0.25, y=y, w=0.4, h=0.45,
               size=18, bold=True, color="positive",
               align="center", anchor="middle")
        c.text(txt, x=right_x + 0.7, y=y, w=right_w - 0.85, h=0.8,
               size=9, color="white", anchor="middle")

    c.text("※ 전환 비용 $18M, 3년 누적 효익 $52M, Payback 22개월",
           x=0.4, y=6.8, w=9.2, h=0.3,
           size=8, color="grey_700", align="center")


# ============================================================
# H5. comp_venn_3_circle — 3 overlapping circles with labels
# ============================================================
def slide_h5_venn_3_circle(prs):
    s, c = make(prs,
                "Sweet Spot 분석 — Desirability × Feasibility × Viability",
                "3-circle Venn diagram with intersection labels")

    # 3 circles arranged in triangle: A top-left, B top-right, C bottom-center
    d = 3.0  # diameter
    # Centers (cx, cy)
    ax, ay = 3.1, 2.3
    bx, by = 5.9, 2.3
    cx2, cy2 = 4.5, 4.3

    # Circle A (Desirability) — accent with alpha approximation via light fill
    c.circle(x=ax - d / 2, y=ay - d / 2, d=d,
             fill="accent_light", border=2.0, border_color="accent")
    # Circle B (Feasibility)
    c.circle(x=bx - d / 2, y=by - d / 2, d=d,
             fill="grey_200", border=2.0, border_color="grey_700")
    # Circle C (Viability)
    c.circle(x=cx2 - d / 2, y=cy2 - d / 2, d=d,
             fill="zone_positive", border=2.0, border_color="positive")

    # Outer circle labels (non-overlap regions)
    c.text("A · Desirability",
           x=0.6, y=1.2, w=2.2, h=0.3,
           size=11, bold=True, color="accent")
    c.text("고객이 원하는가?\n(사용자 니즈)",
           x=0.6, y=1.5, w=2.2, h=0.6,
           size=8, color="grey_700")

    c.text("B · Feasibility",
           x=7.2, y=1.2, w=2.4, h=0.3,
           size=11, bold=True, color="grey_900")
    c.text("기술로 만들 수 있는가?\n(기술 가능성)",
           x=7.2, y=1.5, w=2.4, h=0.6,
           size=8, color="grey_700")

    c.text("C · Viability",
           x=7.2, y=5.3, w=2.4, h=0.3,
           size=11, bold=True, color="positive")
    c.text("비즈니스로 지속 가능한가?\n(경제성)",
           x=7.2, y=5.6, w=2.4, h=0.6,
           size=8, color="grey_700")

    # Intersection labels (approximated positions)
    # A ∩ B (top center, between A and B)
    c.text("A ∩ B\nUsable", x=3.9, y=1.95, w=1.2, h=0.55,
           size=8, bold=True, color="grey_900",
           align="center", anchor="middle")

    # A ∩ C (left lower)
    c.text("A ∩ C\nEmotional", x=2.4, y=3.55, w=1.3, h=0.55,
           size=8, bold=True, color="grey_900",
           align="center", anchor="middle")

    # B ∩ C (right lower)
    c.text("B ∩ C\nSustainable", x=5.3, y=3.55, w=1.3, h=0.55,
           size=8, bold=True, color="grey_900",
           align="center", anchor="middle")

    # A ∩ B ∩ C — Sweet Spot (center)
    # Draw a filled accent box as the sweet spot emphasis
    c.box(x=4.0, y=3.0, w=1.0, h=0.6,
          fill="accent", border=None)
    c.text("Sweet\nSpot",
           x=4.0, y=3.0, w=1.0, h=0.6,
           size=10, bold=True, color="white",
           align="center", anchor="middle")

    # Bottom note
    c.text("세 영역이 모두 교차하는 Sweet Spot에서 성공 확률이 가장 높은 전략이 도출된다 — "
           "Design Thinking (IDEO) 기반",
           x=0.4, y=6.9, w=9.2, h=0.3,
           size=8, color="grey_700", align="center")


# ============================================================
# Post-processing: PDF + PNG export
# ============================================================
def export_pdf_and_pngs(pptx_path: Path):
    """Convert to PDF and per-slide PNGs (Windows + PowerPoint required)."""
    # PDF
    try:
        from ppt_builder.visual_validate import convert_pptx_to_pdf
        pdf_path = pptx_path.with_suffix(".pdf")
        convert_pptx_to_pdf(pptx_path, pdf_path)
        print(f"PDF saved:  {pdf_path}")
    except Exception as e:
        print(f"[warn] PDF export skipped: {e}")

    # PNGs
    try:
        from ppt_builder.track_c.png_export import pptx_to_pngs
        png_dir = pptx_path.parent / f"{pptx_path.stem}_png"
        pngs = pptx_to_pngs(pptx_path, png_dir)
        print(f"PNGs saved: {png_dir} ({len(pngs)} files)")
    except Exception as e:
        print(f"[warn] PNG export skipped: {e}")


# ============================================================
# Main
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_h1_timeline_with_images(prs)
    slide_h2_funnel_with_conversion(prs)
    slide_h3_pyramid_maturity(prs)
    slide_h4_split_comparison_asymmetric(prs)
    slide_h5_venn_3_circle(prs)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    pptx_path = OUTPUT_DIR / f"{NAME}.pptx"
    prs.save(str(pptx_path))
    print(f"PPTX saved: {pptx_path}")

    export_pdf_and_pngs(pptx_path)

    # Optional evaluation
    try:
        from ppt_builder.evaluate import evaluate_pptx, print_report
        report = evaluate_pptx(str(pptx_path))
        print_report(report)
    except Exception as e:
        print(f"[warn] Evaluation skipped: {e}")


if __name__ == "__main__":
    main()
